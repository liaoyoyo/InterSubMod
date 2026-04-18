#!/usr/bin/env python3
"""
Extract text content from all PPTX files in the meeting PPT directory.

Outputs one Markdown file per PPTX and a consolidated 00_INDEX.md.

Usage:
    python3 scripts/analysis/extract_ppt_text.py
"""

import os
import re
import sys
import traceback
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
PPT_DIR = Path("/big7_disk/liaoyoyo2001/Meeting/PPT/PPT")
OUTPUT_DIR = Path(
    "/big7_disk/liaoyoyo2001/InterSubMod/docs/archive/2026/04/ppt_text_extracts"
)

SKIP_PATTERNS = [
    re.compile(r"^~\$"),           # lock files
    re.compile(r"^YYYYMMDD_"),     # template
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def should_skip(filename: str) -> bool:
    """Return True if the file should be skipped."""
    for pat in SKIP_PATTERNS:
        if pat.search(filename):
            return True
    return False


def extract_date_from_filename(filename: str) -> str:
    """Try to extract YYYYMMDD from the filename."""
    m = re.match(r"(\d{8})", filename)
    if m:
        raw = m.group(1)
        try:
            dt = datetime.strptime(raw, "%Y%m%d")
            return dt.strftime("%Y-%m-%d")
        except ValueError:
            return raw
    return "N/A"


def extract_text_from_shape(shape) -> str:
    """Extract text from a single shape (text frame, table, group, etc.)."""
    lines = []

    # --- Text frame ---
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)

    # --- Table ---
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            line = " | ".join(cells)
            if any(c.strip() for c in cells):
                if row_idx == 0:
                    lines.append(line)
                    lines.append(" | ".join(["---"] * len(cells)))
                else:
                    lines.append(line)

    # --- Group shape (recurse) ---
    if shape.shape_type is not None:
        try:
            # GroupShape has .shapes
            if hasattr(shape, "shapes"):
                for child in shape.shapes:
                    child_text = extract_text_from_shape(child)
                    if child_text:
                        lines.append(child_text)
        except Exception:
            pass

    return "\n".join(lines)


def get_slide_title(slide) -> str:
    """Try to extract slide title from title placeholder or first text shape."""
    # Try title placeholder first
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        if title:
            return title

    # Fallback: first shape with text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) < 200:
                return text[:100]

    return "(no title)"


def extract_notes(slide) -> str:
    """Extract notes text from a slide."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            text = notes_slide.notes_text_frame.text.strip()
            if text:
                return text
    return ""


def process_pptx(filepath: Path) -> dict:
    """
    Process a single PPTX file.

    Returns dict with keys:
        filename, date, slide_count, slides, error
    where slides is a list of dicts with keys:
        number, title, content, notes
    """
    filename = filepath.name
    stem = filepath.stem
    date_str = extract_date_from_filename(filename)

    result = {
        "filename": filename,
        "stem": stem,
        "date": date_str,
        "slide_count": 0,
        "slides": [],
        "error": None,
        "titles": [],
    }

    try:
        prs = Presentation(str(filepath))
        result["slide_count"] = len(prs.slides)

        for idx, slide in enumerate(prs.slides, start=1):
            title = get_slide_title(slide)
            result["titles"].append(title)

            # Collect all text from all shapes
            content_parts = []
            for shape in slide.shapes:
                text = extract_text_from_shape(shape)
                if text:
                    content_parts.append(text)
            content = "\n\n".join(content_parts)

            notes = extract_notes(slide)

            result["slides"].append({
                "number": idx,
                "title": title,
                "content": content,
                "notes": notes,
            })

    except Exception as e:
        result["error"] = f"{type(e).__name__}: {e}"
        traceback.print_exc()

    return result


def write_markdown(result: dict, output_dir: Path) -> Path:
    """Write a Markdown file for one PPTX extraction result."""
    outpath = output_dir / f"{result['stem']}.md"

    lines = []
    lines.append(f"# {result['filename']}")
    lines.append("")
    lines.append(f"- Slides: {result['slide_count']}")
    lines.append(f"- Date: {result['date']}")
    if result["error"]:
        lines.append(f"- **Error**: {result['error']}")
    lines.append("")

    for slide in result["slides"]:
        lines.append(f"## Slide {slide['number']}: {slide['title']}")
        lines.append("")
        if slide["content"]:
            lines.append(slide["content"])
            lines.append("")
        if slide["notes"]:
            lines.append(f"> **Notes:** {slide['notes']}")
            lines.append("")
        lines.append("---")
        lines.append("")

    outpath.write_text("\n".join(lines), encoding="utf-8")
    return outpath


def write_index(all_results: list, output_dir: Path) -> Path:
    """Write 00_INDEX.md with summary table and statistics."""
    outpath = output_dir / "00_INDEX.md"

    total_files = len(all_results)
    total_slides = sum(r["slide_count"] for r in all_results)
    total_errors = sum(1 for r in all_results if r["error"])
    total_success = total_files - total_errors

    lines = []
    lines.append("# PPT Text Extracts Index")
    lines.append("")
    lines.append(f"Extracted: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    lines.append("")
    lines.append("## Statistics")
    lines.append("")
    lines.append(f"- Total files processed: {total_files}")
    lines.append(f"- Successful: {total_success}")
    lines.append(f"- Errors: {total_errors}")
    lines.append(f"- Total slides: {total_slides}")
    lines.append("")
    lines.append("## File Index")
    lines.append("")
    lines.append("| # | Filename | Date | Slides | Key Topics |")
    lines.append("|---|---------|------|--------|------------|")

    for i, r in enumerate(all_results, start=1):
        # Extract key topics from slide titles (first few unique, non-trivial ones)
        topics = []
        seen = set()
        for t in r.get("titles", []):
            t_clean = t.strip()
            if (
                t_clean
                and t_clean != "(no title)"
                and t_clean not in seen
                and len(t_clean) < 100
            ):
                seen.add(t_clean)
                topics.append(t_clean)
            if len(topics) >= 4:
                break
        topics_str = "; ".join(topics) if topics else "(no titles extracted)"

        # Escape pipe characters in topics
        topics_str = topics_str.replace("|", "/")

        err_marker = " **[ERROR]**" if r["error"] else ""
        lines.append(
            f"| {i} | [{r['stem']}.md]({r['stem']}.md){err_marker} "
            f"| {r['date']} | {r['slide_count']} | {topics_str} |"
        )

    lines.append("")

    if total_errors > 0:
        lines.append("## Errors")
        lines.append("")
        for r in all_results:
            if r["error"]:
                lines.append(f"- **{r['filename']}**: {r['error']}")
        lines.append("")

    outpath.write_text("\n".join(lines), encoding="utf-8")
    return outpath


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    print(f"Source directory: {PPT_DIR}")
    print(f"Output directory: {OUTPUT_DIR}")
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Collect PPTX files
    pptx_files = sorted(
        f for f in PPT_DIR.glob("*.pptx") if not should_skip(f.name)
    )

    print(f"\nFound {len(pptx_files)} PPTX files to process (after skipping lock/template files)")
    print()

    all_results = []
    for i, filepath in enumerate(pptx_files, start=1):
        size_mb = filepath.stat().st_size / (1024 * 1024)
        print(f"[{i:2d}/{len(pptx_files)}] Processing: {filepath.name} ({size_mb:.1f} MB) ...", end=" ", flush=True)

        result = process_pptx(filepath)
        all_results.append(result)

        if result["error"]:
            print(f"ERROR: {result['error']}")
        else:
            print(f"OK ({result['slide_count']} slides)")

        # Write individual markdown
        write_markdown(result, OUTPUT_DIR)

    # Write index
    index_path = write_index(all_results, OUTPUT_DIR)

    # Summary
    total_slides = sum(r["slide_count"] for r in all_results)
    total_errors = sum(1 for r in all_results if r["error"])
    print()
    print("=" * 60)
    print(f"DONE")
    print(f"  Files processed : {len(all_results)}")
    print(f"  Total slides    : {total_slides}")
    print(f"  Errors          : {total_errors}")
    print(f"  Output directory: {OUTPUT_DIR}")
    print(f"  Index file      : {index_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()
