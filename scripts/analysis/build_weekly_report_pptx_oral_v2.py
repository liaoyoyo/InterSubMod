#!/usr/bin/env python3
"""
將 2026-03-11 研究主線整合週報重製為較適合口頭報告的新簡報。

使用方式:
python scripts/analysis/build_weekly_report_pptx_oral_v2.py \
  --config docs/references/manual/assets/20260311_weekly_report_pptx_oral_config_01.json \
  --profile docs/references/manual/assets/20260311_liao_research_ppt_profile_01.json
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5


def hex_rgb(value: str) -> RGBColor:
    value = value.replace("#", "").strip()
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)


def add_rect(slide, x, y, w, h, *, fill, line=None, radius=True, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color="17202A",
    font_name="Noto Sans TC",
    bold=False,
    italic=False,
    align="left",
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    line_spacing=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_rgb(color)
    return box


def add_bullets(
    slide,
    bullets,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color="17202A",
    font_name="Noto Sans TC",
    space_after=6,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = hex_rgb(color)
    return box


def wrap_text(text: str, max_len: int = 34, max_lines: int = 4) -> str:
    parts = [token for token in re.split(r"([/ _-])", text) if token]
    lines = []
    current = ""
    for part in parts:
        if not current:
            current = part
        elif len(current) + len(part) <= max_len:
            current += part
        else:
            lines.append(current.rstrip())
            current = part.lstrip()
    if current:
        lines.append(current.rstrip())
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1][: max_len - 1] + "…"
    return "\n".join(lines)


def add_header(slide, zh_title, en_title, zh_subtitle, theme, fonts, *, dark=False):
    fg = theme["light_text"] if dark else theme["text"]
    muted = theme["light_muted"] if dark else theme["muted"]
    add_text(slide, zh_title, 0.7, 0.42, 9.8, 0.58, font_size=25, color=fg, font_name=fonts["zh_title"], bold=True)
    add_text(slide, en_title, 0.72, 0.98, 9.2, 0.22, font_size=11, color=theme["accent"], font_name=fonts["en_body"])
    add_text(slide, zh_subtitle, 0.72, 1.25, 11.4, 0.34, font_size=13, color=muted, font_name=fonts["zh_body"])


def add_footer(slide, theme, fonts, page_num, total_pages, footer_text, *, dark=False):
    color = theme["light_muted"] if dark else "52606D"
    add_text(slide, footer_text, 0.72, 7.01, 8.9, 0.18, font_size=8.6, color=color, font_name=fonts["en_body"], margin=0)
    add_text(slide, f"{page_num}/{total_pages}", 12.05, 6.99, 0.5, 0.18, font_size=10, color=color, font_name=fonts["en_body"], align="right", margin=0)


def make_hero_figure(path: Path, theme: dict) -> None:
    fig, ax = plt.subplots(figsize=(7.6, 5.0), dpi=220)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis("off")
    bg = "#F7F3EC"
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)

    boxes = [
        (0.5, 5.9, 4.4, 1.2, "#1E2A44", "Caller-first", "GQ / QUAL / eligibility"),
        (0.9, 4.25, 4.4, 1.2, "#D8E7E0", "Methylation-support", "Quality / Pairwise / agreement"),
        (1.3, 2.6, 4.4, 1.2, "#EBD8CF", "Annotation / QC / triage", "dataset-aware + phase completeness"),
    ]
    for x, y, w, h, color, title, body in boxes:
        patch = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.08,rounding_size=0.16", facecolor=color, linewidth=0)
        ax.add_patch(patch)
        ax.text(x + 0.24, y + 0.78, title, fontsize=12.5, fontweight="bold", color="white" if color == "#1E2A44" else "#17202A")
        ax.text(x + 0.24, y + 0.36, body, fontsize=9.6, color="white" if color == "#1E2A44" else "#4F5B67")

    ax.text(6.4, 6.75, "Study Matrix", fontsize=13, fontweight="bold", color="#17202A")
    quadrants = [
        (6.2, 5.1, "#F0D8CF", "5kHz\npaired"),
        (8.1, 5.1, "#D7E5DE", "DORADO\npaired"),
        (6.2, 3.2, "#F0D8CF", "5kHz\nTO"),
        (8.1, 3.2, "#D7E5DE", "DORADO\nTO"),
    ]
    for x, y, color, label in quadrants:
        patch = FancyBboxPatch((x, y), 1.55, 1.45, boxstyle="round,pad=0.05,rounding_size=0.12", facecolor=color, linewidth=0)
        ax.add_patch(patch)
        ax.text(x + 0.775, y + 0.73, label, ha="center", va="center", fontsize=11.5, fontweight="bold", color="#17202A")

    ax.text(6.25, 2.2, "Goal", fontsize=11.5, fontweight="bold", color="#C76B50")
    ax.text(6.25, 1.72, "Keep caller as the primary gate,\nuse methylation for support and annotation.", fontsize=9.7, color="#4F5B67")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor=bg)
    plt.close()


def make_grouped_bars(path: Path, categories, series, title, ylabel):
    fig, ax = plt.subplots(figsize=(7.0, 3.8), dpi=220)
    x = range(len(categories))
    width = 0.22
    offsets = [-width, 0, width]
    all_vals = [v for item in series for v in item["values"]]
    ymin = min(all_vals) - 0.005
    ymax = max(all_vals) + 0.004
    for idx, item in enumerate(series):
        positions = [v + offsets[idx] for v in x]
        ax.bar(positions, item["values"], width=width * 0.9, color=f"#{item['color']}", label=item["name"])
        for px, val in zip(positions, item["values"]):
            ax.text(px, val + 0.00055, f"{val:.4f}", ha="center", va="bottom", fontsize=8)
    ax.set_xticks(list(x))
    ax.set_xticklabels(categories, fontsize=9)
    ax.set_ylabel(ylabel, fontsize=10)
    ax.set_title(title, fontsize=12, fontweight="bold")
    ax.set_ylim(ymin, ymax)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.legend(frameon=False, fontsize=8, ncol=3, loc="upper center", bbox_to_anchor=(0.5, 1.14))
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def make_phase2_chart(path: Path, panels):
    fig, axes = plt.subplots(1, 2, figsize=(9.6, 4.1), dpi=220)
    for ax, panel in zip(axes, panels):
        labels = [row["name"] for row in panel["rows"]]
        values = [row["f1"] for row in panel["rows"]]
        colors = [f"#{row['color']}" for row in panel["rows"]]
        ypos = list(range(len(labels)))
        ax.barh(ypos, values, color=colors, height=0.58)
        ax.set_yticks(ypos)
        ax.set_yticklabels(labels, fontsize=8)
        ax.set_xlim(panel["xlim"][0], panel["xlim"][1])
        ax.set_title(panel["title"], fontsize=10.5, fontweight="bold")
        ax.grid(axis="x", linestyle="--", alpha=0.22)
        for y, val in zip(ypos, values):
            ax.text(val + 0.0008, y, f"{val:.4f}", va="center", fontsize=7.6)
    fig.suptitle("Paired raw caller benchmarks", fontsize=12, fontweight="bold")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def make_to_support_chart(path: Path) -> None:
    fig, axes = plt.subplots(1, 2, figsize=(8.2, 3.6), dpi=220)
    specs = [
        ("5kHz TO", ["caller-only", "caller+Pairwise", "caller+agreement"], [0.006824, 0.004219, 0.002163], ["#1E2A44", "#C76B50", "#6E9D8A"]),
        ("DORADO TO", ["caller-only", "Quality-only", "Pairwise-only"], [0.000540, 0.002620, 0.002217], ["#1E2A44", "#C76B50", "#6E9D8A"]),
    ]
    for ax, (title, labels, vals, colors) in zip(axes, specs):
        ax.bar(labels, vals, color=colors, width=0.58)
        ax.set_title(title, fontsize=10.5, fontweight="bold")
        ax.set_ylabel("ΔF1 vs baseline", fontsize=9)
        ax.tick_params(axis="x", labelrotation=20, labelsize=8)
        ax.tick_params(axis="y", labelsize=8)
        ax.grid(axis="y", linestyle="--", alpha=0.22)
        ymax = max(vals) * 1.25
        ax.set_ylim(0, ymax)
        for idx, val in enumerate(vals):
            ax.text(idx, val + ymax * 0.02, f"{val:.4f}", ha="center", va="bottom", fontsize=7.6)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def make_annotation_policy_chart(path: Path, rows):
    labels = [row["dataset"].replace("HCC1395 ", "").replace(" paired", "\npaired").replace(" TO", "\nTO") for row in rows]
    values = [float(row["delta"].replace("+", "")) for row in rows]
    colors = ["#1E2A44", "#C76B50", "#6E9D8A", "#54708C"]
    fig, (ax_top, ax_bot) = plt.subplots(2, 1, figsize=(6.8, 3.8), dpi=220, sharex=True,
                                          gridspec_kw={"height_ratios": [1.6, 1], "hspace": 0.08})
    for ax in (ax_top, ax_bot):
        ax.bar(labels, values, color=colors, width=0.58)
        ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax_top.set_ylim(0.005, max(values) * 1.22)
    ax_bot.set_ylim(0, 0.0015)
    ax_top.spines["bottom"].set_visible(False)
    ax_bot.spines["top"].set_visible(False)
    ax_top.tick_params(bottom=False)
    d = 0.008
    kwargs = dict(transform=ax_top.transAxes, color="k", clip_on=False, linewidth=0.8)
    ax_top.plot((-d, +d), (-d, +d), **kwargs)
    ax_top.plot((1 - d, 1 + d), (-d, +d), **kwargs)
    kwargs.update(transform=ax_bot.transAxes)
    ax_bot.plot((-d, +d), (1 - d, 1 + d), **kwargs)
    ax_bot.plot((1 - d, 1 + d), (1 - d, 1 + d), **kwargs)
    for idx, val in enumerate(values):
        target_ax = ax_top if val > 0.003 else ax_bot
        target_ax.text(idx, val + max(values) * 0.025, f"{val:.4f}", ha="center", va="bottom", fontsize=7.6)
    ax_bot.set_ylabel("ΔF1 vs baseline", fontsize=9)
    ax_top.set_title("Best policy by dataset", fontsize=11, fontweight="bold")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def build_assets(cfg: dict) -> dict:
    assets_dir = Path(cfg["meta"]["assets_dir"])
    ensure_dir(assets_dir)
    hero = assets_dir / "cover_hero.png"
    paired = assets_dir / "paired_rerun_f1.png"
    phase2 = assets_dir / "phase2_raw.png"
    to_chart = assets_dir / "to_support.png"
    anno = assets_dir / "annotation_policy.png"
    make_hero_figure(hero, cfg["theme"])
    make_grouped_bars(paired, cfg["charts"]["paired_rerun"]["categories"], cfg["charts"]["paired_rerun"]["series"], "Paired rerun benchmark", "F1")
    make_phase2_chart(phase2, cfg["charts"]["phase2_raw"])
    make_to_support_chart(to_chart)
    make_annotation_policy_chart(anno, cfg["deck"]["annotation_policy"]["rows"])
    return {
        "hero": hero,
        "paired": paired,
        "phase2": phase2,
        "to_support": to_chart,
        "annotation": anno,
    }


def create_prs(cfg):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = cfg["meta"]["author"]
    prs.core_properties.title = cfg["meta"]["deck_title"]
    return prs


def add_cover(slide, cfg, theme, fonts, assets):
    set_bg(slide, theme["bg"])
    add_rect(slide, 0.45, 0.55, 5.6, 6.1, fill="FFFFFF", line=theme["line"], radius=True)
    slide.shapes.add_picture(str(assets["hero"]), Inches(0.7), Inches(0.85), width=Inches(5.1), height=Inches(5.25))
    add_text(slide, cfg["meta"]["kicker"], 6.45, 0.95, 2.3, 0.22, font_size=11, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, cfg["meta"]["deck_title"], 6.4, 1.35, 6.4, 1.15, font_size=25, color=theme["text"], font_name=fonts["zh_title"], bold=True, line_spacing=1.05)
    add_text(slide, "Weekly Research Review", 6.42, 2.52, 3.2, 0.22, font_size=12, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    add_text(slide, cfg["meta"]["deck_subtitle"], 6.42, 2.84, 5.8, 0.42, font_size=14, color=theme["muted"], font_name=fonts["zh_body"])
    add_rect(slide, 6.4, 3.45, 5.9, 1.2, fill=theme["dark"], line=None, radius=True)
    add_text(slide, cfg["meta"]["summary_line"], 6.68, 3.62, 5.3, 0.72, font_size=14.5, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_text(slide, f"報告整理時間｜{cfg['meta']['generated_date']}", 6.42, 5.25, 2.8, 0.2, font_size=11, color=theme["muted"], font_name=fonts["zh_body"])
    add_text(slide, "Report Date", 6.42, 5.48, 1.6, 0.18, font_size=10, color=theme["accent"], font_name=fonts["en_body"])
    add_text(slide, f"報告人｜{cfg['meta']['author']}", 9.2, 5.25, 2.3, 0.2, font_size=11, color=theme["muted"], font_name=fonts["zh_body"])
    add_text(slide, "Presenter", 9.2, 5.48, 1.4, 0.18, font_size=10, color=theme["accent"], font_name=fonts["en_body"])


def add_agenda(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "報告大綱", "Agenda", "本次簡報共 14 頁，分成 6 個主題，先架構後證據再下一步。", theme, fonts)
    sections = [
        ("01", "破題結論與主架構", "Opening thesis and framework"),
        ("02", "資料四象限與研究推進", "Study matrix and weekly milestones"),
        ("03", "paired / TO 核心結果", "Paired and tumor-only results"),
        ("04", "特徵定位與 phase 2", "Feature roles and phase-2 benchmark"),
        ("05", "annotation 決策", "Annotation-layer decisions"),
        ("06", "缺口、下一步與複查路徑", "Gaps, actions, and traceability"),
    ]
    for idx, (num, zh, en) in enumerate(sections):
        x = 0.85 + (idx % 2) * 6.2
        y = 1.95 + (idx // 2) * 1.55
        add_rect(slide, x, y, 5.45, 1.12, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.2, y + 0.18, 0.62, 0.6, fill=theme["dark"], line=None, radius=True)
        add_text(slide, num, x + 0.33, y + 0.33, 0.28, 0.16, font_size=11, color=theme["light_text"], font_name=fonts["en_title"], bold=True, align="center")
        add_text(slide, zh, x + 1.0, y + 0.18, 3.9, 0.22, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, en, x + 1.0, y + 0.5, 4.1, 0.18, font_size=10, color=theme["accent"], font_name=fonts["en_body"])
    add_rect(slide, 0.85, 6.63, 11.85, 0.35, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "頁數策略：一頁一重點；若觀眾需要更多細節，附錄與原始報告可直接回查。", 1.02, 6.72, 11.3, 0.14, font_size=9.8, color="52606D", font_name=fonts["zh_body"])


def add_thesis(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "本週最重要的結論", "Opening Thesis", "研究主線已從『大量試規則』收斂成固定分工：caller-first、methylation-support、annotation/QC/artifact。", theme, fonts)
    add_rect(slide, 0.8, 1.95, 7.2, 2.3, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "目前最穩的救援仍在 caller 端；甲基訊號已證明有價值，但定位應是 support、annotation 與 triage，而不是直接取代 caller。", 1.08, 2.25, 6.55, 0.78, font_size=20, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_text(slide, "The stable pattern is caller-first. Methylation already shows value, but its current role is support, annotation, and artifact triage.", 1.08, 3.22, 6.35, 0.26, font_size=10.5, color=theme["light_muted"], font_name=fonts["en_body"])
    stats = [
        ("5kHz paired", "F1 +0.0010", "83 個 FP 被額外移除"),
        ("5kHz TO", "caller-only 最佳", "甲基 support 成立但未超過 gate"),
    ]
    for idx, (title, metric, note) in enumerate(stats):
        x = 8.25
        y = 1.95 + idx * 1.53
        add_rect(slide, x, y, 4.25, 1.2, fill="FFFFFF", line=theme["line"], radius=True)
        metric_size = 18 if len(metric) <= 10 else 14.6
        add_text(slide, title, x + 0.2, y + 0.15, 1.8, 0.18, font_size=12.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, metric, x + 0.2, y + 0.43, 2.6, 0.24, font_size=metric_size, color=theme["text"], font_name=fonts["en_title"], bold=True)
        add_text(slide, note, x + 0.2, y + 0.82, 3.7, 0.18, font_size=10.2, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 8.25, 5.03, 4.25, 0.6, fill="EBD8CF", line=None, radius=True)
    add_text(slide, "DORADO TO｜support 非全負，但 Pairwise 方向與 5kHz 相反。", 8.45, 5.2, 3.75, 0.14, font_size=10.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_rect(slide, 0.8, 5.82, 11.7, 0.88, fill="FFFFFF", line=theme["line"], radius=True)
    layer_cards = [
        ("第一層｜caller-first", "用 GQ / QUAL / eligibility 定義可救回的邊緣集合。", theme["dark"], theme["light_text"]),
        ("第二層｜methylation-support", "用 Quality / Pairwise / agreement 補強排序與支持。", theme["accent2"], theme["text"]),
        ("第三層｜annotation / QC / triage", "把 phase completeness、dataset-aware 訊號與 artifact 旁路制度化輸出。", "EBD8CF", theme["text"]),
    ]
    for idx, (title, body, fill, fg) in enumerate(layer_cards):
        x = 1.0 + idx * 3.86
        add_rect(slide, x, 5.94, 3.45, 0.5, fill=fill, line=None, radius=True)
        add_text(slide, title, x + 0.16, 6.03, 3.02, 0.12, font_size=10.8, color=fg, font_name=fonts["zh_title"], bold=True)
        add_text(slide, body, x + 0.16, 6.22, 3.02, 0.12, font_size=8.4, color=fg, font_name=fonts["zh_body"])


def add_matrix(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "資料四象限與比較邊界", "Study matrix and comparability", "四個 dataset 都已跑通，但 snapshot scope 與 candidate-specific coverage 不完全相同，解讀邊界必須分清楚。", theme, fonts)
    add_rect(slide, 0.8, 1.95, 6.15, 4.85, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "四象限資料", 1.0, 2.15, 1.6, 0.18, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    quads = [
        ("HCC1395 5kHz paired", "最 noisy 的主實驗 paired 樣本", "F0D8CF"),
        ("HCC1395 DORADO paired", "跨平台 paired 交叉驗證", "D7E5DE"),
        ("HCC1395 5kHz TO", "最能看見 methylation-support 的 TO 場景", "F0D8CF"),
        ("HCC1395 DORADO TO", "TO 跨平台驗證；本輪用 subset haplotag", "D7E5DE"),
    ]
    positions = [(1.02, 2.6), (4.03, 2.6), (1.02, 4.4), (4.03, 4.4)]
    for (x, y), (title, body, fill) in zip(positions, quads):
        add_rect(slide, x, y, 2.65, 1.35, fill=fill, line=None, radius=True)
        add_text(slide, title, x + 0.14, y + 0.12, 2.25, 0.22, font_size=11.4, color=theme["text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, body, x + 0.14, y + 0.46, 2.25, 0.42, font_size=9.5, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 7.15, 1.95, 5.35, 4.85, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "比較時一定要記住", 7.38, 2.15, 2.2, 0.18, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    bullets = [
        "5kHz TO snapshot 來自 full tagged BAM；DORADO TO 來自 candidate-window subset tagged BAM。",
        "paired candidate-specific coverage 只有 8%~20%，TO 則是 71%~100%。",
        "因此 paired 的甲基無效不能直接解讀成 paired 沒有甲基價值。",
        "PairwiseMedianDist 方向在 5kHz TO 與 DORADO TO 相反，只能做 dataset-aware 解讀。",
    ]
    add_bullets(slide, bullets, 7.35, 2.55, 4.72, 2.3, font_size=11.0, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 7.35, 5.28, 4.65, 0.96, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "結論：四象限可以一起看，但不是所有數值都能硬做跨 dataset 絕對值比較。", 7.62, 5.56, 4.05, 0.3, font_size=11.2, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)


def add_milestones(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "本週推進順序", "Weekly milestones", "先排除 mixed 主線，再跑通 paired / TO、candidate-specific rescue、phase 2 與 annotation 回接。", theme, fonts)
    add_rect(slide, 0.78, 2.0, 11.9, 3.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 1.1, 3.72, 10.95, 0.06, fill=theme["accent"], line=None, radius=False)
    x_positions = [1.2, 4.0, 6.8, 9.6]
    for x, item in zip(x_positions, cfg["deck"]["timeline"]):
        add_rect(slide, x, 2.35, 2.2, 1.55, fill="F7F3EC", line=theme["line"], radius=True)
        add_rect(slide, x + 0.14, 2.48, 0.86, 0.34, fill=theme["dark"], line=None, radius=True)
        add_text(slide, item["date"], x + 0.18, 2.54, 0.78, 0.18, font_size=10.5, color=theme["light_text"], font_name=fonts["en_title"], bold=True, align="center")
        add_text(slide, item["title"], x + 0.12, 2.96, 1.95, 0.22, font_size=11.4, color=theme["text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, item["result"], x + 0.12, 3.3, 1.95, 0.42, font_size=9.4, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 0.78, 5.5, 11.9, 0.55, fill="EBD8CF", line=None, radius=True)
    add_text(slide, "重要收斂：這週不是只多做了幾條規則，而是把『什麼應進主流程、什麼只該留在分析層』這件事正式定型。", 1.0, 5.68, 11.3, 0.18, font_size=11.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_paired_results(slide, cfg, theme, fonts, assets):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "paired pure 正式 rerun", "Paired benchmark", "5kHz 是最有壓力的主樣本；DORADO 是跨平台交叉驗證基準。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 7.0, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    slide.shapes.add_picture(str(assets["paired"]), Inches(1.03), Inches(2.18), width=Inches(6.55), height=Inches(3.45))
    add_text(slide, "圖表解讀｜X 軸：paired dataset；Y 軸：F1。每組柱狀依序代表 ClairS、LongPhase-S、InterSubMod。", 1.03, 5.82, 6.3, 0.24, font_size=9.4, color="52606D", font_name=fonts["zh_body"])
    cards = [
        ("5kHz paired", "0.8522 → 0.8532", "增益主要來自額外移除 83 個 FP，不是大量救回 TP。", theme["dark"], theme["light_text"]),
        ("DORADO paired", "0.8592 → 0.8590", "方向可對照，但同一條規則不保證可攜。", theme["accent"], theme["light_text"]),
        ("paired 高層判讀", "caller-first 仍主導", "InterSubMod 在 paired 有價值，但目前更像精修 precision 的分析層。", "D7E5DE", theme["text"]),
    ]
    for idx, (title, metric, note, fill, fg) in enumerate(cards):
        x = 8.05
        y = 2.0 + idx * 1.52
        add_rect(slide, x, y, 4.35, 1.24, fill=fill, line=None, radius=True)
        metric_size = 15.2 if any(ch.isdigit() for ch in metric) else 13.8
        add_text(slide, title, x + 0.18, y + 0.13, 1.9, 0.16, font_size=11.8, color=fg, font_name=fonts["zh_title"], bold=True)
        add_text(slide, metric, x + 0.18, y + 0.42, 2.65, 0.22, font_size=metric_size, color=fg, font_name=fonts["en_title"], bold=True)
        add_text(slide, note, x + 0.18, y + 0.82, 3.9, 0.2, font_size=9.5, color=fg, font_name=fonts["zh_body"])


def add_to_rescue(slide, cfg, theme, fonts, assets):
    set_bg(slide, theme["bg"])
    add_header(slide, "tumor-only rescue：caller 先、甲基後", "Tumor-only rescue", "TO 下甲基特徵不是全負效果，但目前仍沒有穩定超過最佳 caller gate。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 6.6, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    slide.shapes.add_picture(str(assets["to_support"]), Inches(1.02), Inches(2.18), width=Inches(6.2), height=Inches(3.15))
    add_text(slide, "圖表解讀｜Y 軸為相對 baseline 的 ΔF1。左圖是 5kHz TO，右圖是 DORADO TO。", 1.02, 5.54, 6.1, 0.24, font_size=9.3, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 7.7, 1.92, 4.62, 2.02, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "5kHz TO", 7.95, 2.1, 1.4, 0.18, font_size=13.5, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(
        slide,
        [
            "caller-only 最佳：`QUAL>=10 or GQ>=20`，`491 TP / 118 FP`。",
            "`Pairwise>=0.20` 與 `agreement_positive` 都有 support 價值，但仍低於最佳 caller-only。",
            "結論：適合 `caller-first + methylation-support`。",
        ],
        7.92, 2.42, 4.0, 1.18, font_size=10.4, color=theme["text"], font_name=fonts["zh_body"]
    )
    add_rect(slide, 7.7, 4.12, 4.62, 2.0, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "DORADO TO", 7.95, 4.3, 1.6, 0.18, font_size=13.5, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(
        slide,
        [
            "`GQ>=10` 仍是最佳 caller-only。",
            "`Quality>=60`、`Pairwise<=0.20` 單獨都為正值，但在 gate 後仍只作 support。",
            "本輪使用 `candidate-window subset BAM + subset haplotag`，方向可判讀但不做絕對值硬比。",
        ],
        7.92, 4.62, 4.0, 1.18, font_size=10.2, color=theme["text"], font_name=fonts["zh_body"]
    )
    add_rect(slide, 7.7, 6.06, 4.62, 0.5, fill=theme["dark"], line=None, radius=True)
    add_text(
        slide,
        "結論：甲基 support 成立，但目前最穩的主規則仍是 caller-first。",
        7.93,
        6.2,
        4.08,
        0.18,
        font_size=10.1,
        color=theme["light_text"],
        font_name=fonts["zh_title"],
        bold=True,
    )


def add_feature_architecture(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "特徵分工重新定型", "Feature architecture", "這週已把特徵從『候選規則』重新整理成三種用途：主規則、support、annotation/QC/triage。", theme, fonts)
    cols = [
        ("caller-first", "主規則", ["GQ", "QUAL", "candidate eligibility"], theme["dark"], theme["light_text"]),
        ("methylation-support", "第二層支持", ["Quality_Score", "PairwiseMedianDist", "agreement_positive"], theme["accent"], theme["light_text"]),
        ("annotation / QC / triage", "註記與旁路", ["hp_assign_rate", "artifact_lowvaf_highad_lowcv", "dataset-aware flags"], theme["accent2"], theme["text"]),
    ]
    for idx, (en, zh, items, fill, fg) in enumerate(cols):
        x = 0.85 + idx * 4.16
        add_rect(slide, x, 2.0, 3.6, 3.7, fill=fill, line=None, radius=True)
        add_text(slide, zh, x + 0.18, 2.18, 2.2, 0.18, font_size=15, color=fg, font_name=fonts["zh_title"], bold=True)
        add_text(slide, en, x + 0.18, 2.48, 2.4, 0.16, font_size=10.5, color=fg, font_name=fonts["en_body"])
        add_bullets(slide, items, x + 0.18, 2.88, 3.0, 1.2, font_size=10.8, color=fg, font_name=fonts["en_body"])
        if idx == 0:
            note = "最穩、跨 dataset 最一致，負責決定第一層可救回集合。"
        elif idx == 1:
            note = "在 TO 可提供額外 support，但方向可能依 dataset 改變。"
        else:
            note = "不直接做 keep/filter；負責解讀、QC 與 artifact 診斷。"
        add_text(slide, note, x + 0.18, 4.58, 3.0, 0.46, font_size=10.0, color=fg, font_name=fonts["zh_body"])
    add_rect(slide, 0.85, 6.0, 11.75, 0.52, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "現階段最重要的不是再堆更多條件，而是把有訊息量但不夠穩定的特徵，明確放到正確層級。", 1.05, 6.15, 11.2, 0.16, font_size=11.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_feature_cards(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "關鍵特徵：哪些可升級、哪些不能", "Feature roles", "看的是『目前最合理的定位』，不是單一 dataset 一時有效就直接升級成全域規則。", theme, fonts)
    cards = cfg["deck"]["feature_roles"]
    positions = [(0.82, 1.82), (6.6, 1.82), (0.82, 3.52), (6.6, 3.52), (0.82, 5.22), (6.6, 5.22)]
    for (x, y), row in zip(positions, cards):
        add_rect(slide, x, y, 5.75, 1.48, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, row["feature"], x + 0.18, y + 0.1, 3.8, 0.2, font_size=11.6, color=theme["text"], font_name=fonts["en_title"], bold=True)
        add_text(slide, row["role"], x + 0.18, y + 0.38, 3.6, 0.16, font_size=10.2, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, f"意義：{row['meaning']}", x + 0.18, y + 0.66, 5.1, 0.26, font_size=9.4, color="52606D", font_name=fonts["zh_body"])
        add_text(slide, f"限制：{row['limit']}", x + 0.18, y + 1.0, 5.1, 0.26, font_size=9.4, color="65707C", font_name=fonts["zh_body"])


def add_phase2(slide, cfg, theme, fonts, assets):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "phase 2：raw caller 的邊界", "Paired raw benchmark", "paired raw pileup / full model 有額外召回，但 precision 代價太大，不能直接取代 merged caller。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 7.2, 4.9, fill="FFFFFF", line=theme["line"], radius=True)
    slide.shapes.add_picture(str(assets["phase2"]), Inches(1.02), Inches(2.12), width=Inches(6.82), height=Inches(3.95))
    add_text(slide, "圖表解讀｜每個 panel 都是 paired 場景；X 軸為 F1，Y 軸為 caller 層級來源。", 1.02, 6.0, 6.4, 0.22, font_size=9.2, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 8.25, 1.92, 4.1, 4.9, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "這頁要傳達的重點", 8.48, 2.1, 2.2, 0.18, font_size=14.8, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(
        slide,
        [
            "5kHz paired：raw pileup / full 都比 merged 更高召回，但 precision 崩得更快。",
            "DORADO paired：raw full 優於 raw pileup，但仍低於 merged output。",
            "paired 問題下，raw caller 更適合拿來理解 recall 空間，不適合直接當最終輸出。",
        ],
        8.42, 2.48, 3.45, 1.78, font_size=10.7, color=theme["text"], font_name=fonts["zh_body"]
    )
    add_rect(slide, 8.42, 4.72, 3.48, 0.9, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "正式主線仍應維持\nmerged caller → LongPhase → InterSubMod", 8.68, 4.98, 2.95, 0.34, font_size=12.2, color=theme["light_text"], font_name=fonts["zh_title"], bold=True, align="center")


def add_annotation(slide, cfg, theme, fonts, assets):
    set_bg(slide, theme["bg"])
    add_header(slide, "annotation layer：目前該做的整合層", "Annotation-layer decision", "phase 2 的 finer interval 已足夠回接 annotation，但還不到改 hard keep / hard filter 的時候。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 6.55, 4.85, fill="FFFFFF", line=theme["line"], radius=True)
    slide.shapes.add_picture(str(assets["annotation"]), Inches(1.02), Inches(2.15), width=Inches(6.1), height=Inches(3.35))
    add_text(slide, "圖表解讀｜顯示四個 dataset 在目前最佳 policy 下，相對 baseline 的 ΔF1。", 1.02, 5.72, 5.9, 0.22, font_size=9.2, color="52606D", font_name=fonts["zh_body"])
    add_rect(slide, 7.65, 1.92, 4.67, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "Quality_Score", 7.9, 2.08, 1.8, 0.16, font_size=12.6, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    add_text(slide, "升級為 support annotation，適合排序、報告與人工審閱。", 7.9, 2.42, 4.0, 0.18, font_size=10.3, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 7.65, 3.38, 4.67, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "PairwiseMedianDist", 7.9, 3.54, 2.4, 0.16, font_size=12.6, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    add_text(slide, "升級為 dataset-aware support annotation，不可直接寫成跨樣本全域規則。", 7.9, 3.88, 4.0, 0.18, font_size=10.3, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 7.65, 4.84, 4.67, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "hp_assign_rate", 7.9, 5.0, 2.0, 0.16, font_size=12.6, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    add_text(slide, "明確降階為 phase / QC annotation；反映相位完整度，不直接做 truth keep。", 7.9, 5.34, 4.0, 0.18, font_size=10.3, color=theme["text"], font_name=fonts["zh_body"])


def add_gaps(slide, cfg, theme, fonts):
    set_bg(slide, theme["dark"])
    add_header(slide, "還沒完成的邊界與風險", "Open gaps and caveats", "這頁要交代的是：哪些結論已可定稿、哪些仍需保守解讀。", theme, fonts, dark=True)
    titles = [
        ("coverage ceiling", "coverage 天花板"),
        ("model gap", "模型缺口"),
        ("scope caveat", "解讀限制"),
    ]
    bodies = [
        "paired 的 candidate-specific coverage 只有 8%~20%，因此 paired 的甲基無效不能直接解讀成完全沒有價值。",
        "TO 的 pileup vs full model 尚未補齊雙模型直接對照，因此 TO 模型差異目前只以 pileup 路線為主。",
        "5kHz TO 與 DORADO TO 的 snapshot scope 不同，read-level 圖只能做方向比較；Pool B FN 仍有 caller-side rescue 空間。",
    ]
    fills = ["F4EFE7", "FFFFFF", "E9D8D0"]
    for idx, ((en_title, zh_title), body, fill) in enumerate(zip(titles, bodies, fills)):
        x = 0.82 + idx * 4.06
        add_rect(slide, x, 2.05, 3.62, 2.65, fill=fill, line=None, radius=True)
        add_text(slide, zh_title, x + 0.22, 2.22, 2.32, 0.22, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, en_title, x + 0.22, 2.52, 2.0, 0.16, font_size=10, color=theme["accent"], font_name=fonts["en_body"], bold=True)
        add_text(slide, body, x + 0.22, 2.82, 3.0, 1.6, font_size=12, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 5.25, 11.76, 0.52, fill="F4EFE7", line=None, radius=True)
    add_text(slide, "這些邊界不代表研究失敗，而是記錄現階段不能講太滿的地方，以確保結論可靠。", 1.04, 5.38, 11.2, 0.18, font_size=11.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_next_steps(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "下一步與現在不該做的事", "Next actions", "接下來應優先把 annotation 接到輸出層，而不是急著改核心 hard filter。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 7.1, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    steps = [
        ("1", "把 annotation 接到 round summary / dashboard / diagnostics"),
        ("2", "繼續累積 finer interval 與 orthogonality 的 cross-dataset 證據"),
        ("3", "若要升級規則，先升級成 annotation export，不直接改核心過濾"),
        ("4", "保留 5kHz artifact triage 作 dataset-specific 診斷"),
    ]
    for idx, (num, text) in enumerate(steps):
        y = 2.25 + idx * 1.0
        add_rect(slide, 1.08, y, 0.56, 0.56, fill=theme["dark"], line=None, radius=True)
        add_text(slide, num, 1.25, y + 0.17, 0.18, 0.12, font_size=10.6, color=theme["light_text"], font_name=fonts["en_title"], bold=True, align="center")
        add_text(slide, text, 1.85, y + 0.1, 5.55, 0.26, font_size=11.2, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 8.18, 1.92, 4.15, 4.95, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "現在不該做", 8.45, 2.15, 2.0, 0.18, font_size=15.4, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(
        slide,
        [
            "不要把 Quality 或 Pairwise 直接變成新的 hard keep。",
            "不要把 artifact-veto 直接拿去做 TP rescue 主規則。",
            "不要把 5kHz 單一 dataset 的 Pairwise 方向寫成跨樣本結論。",
        ],
        8.42, 2.55, 3.35, 2.2, font_size=10.8, color=theme["light_text"], font_name=fonts["zh_body"]
    )
    add_rect(slide, 8.42, 5.55, 3.35, 0.82, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "原則：先用 annotation 解決可解釋性，再考慮是否值得改核心決策。", 8.62, 5.83, 2.95, 0.24, font_size=10.4, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)


def add_appendix(slide, cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "複查路徑與再生成方式", "Traceability and rerun", "這份 deck 的來源、版本與重生設定都已獨立整理，方便之後重跑或換版。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 4.2, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 5.2, 1.92, 3.2, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 8.58, 1.92, 3.74, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "核心文件", 1.04, 2.14, 1.6, 0.18, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    for idx, item in enumerate(cfg["deck"]["source_paths"][:3]):
        y = 2.48 + idx * 1.15
        add_text(slide, item["label"], 1.04, y, 1.7, 0.16, font_size=11.0, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, wrap_text(item["path"], 34, 5), 1.04, y + 0.22, 3.75, 0.78, font_size=7.8, color="52606D", font_name=fonts["en_body"])
    add_text(slide, "主要數據", 5.42, 2.14, 1.5, 0.18, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    for idx, item in enumerate(cfg["deck"]["source_paths"][3:5]):
        y = 2.48 + idx * 1.6
        add_text(slide, item["label"], 5.42, y, 2.1, 0.16, font_size=10.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, wrap_text(item["path"], 28, 6), 5.42, y + 0.22, 2.75, 1.1, font_size=7.2, color="52606D", font_name=fonts["en_body"])
    version_name = Path(cfg["meta"]["output_pptx"]).with_suffix(".version.json").name
    add_text(slide, "版本與重生", 8.8, 2.14, 1.8, 0.18, font_size=15, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_text(slide, "version json", 8.8, 2.5, 1.4, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, version_name, 8.8, 2.72, 2.95, 0.38, font_size=8.4, color="52606D", font_name=fonts["en_body"])
    add_text(slide, "deck path", 8.8, 3.35, 1.2, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, Path(cfg["meta"]["output_pptx"]).name, 8.8, 3.57, 3.0, 0.48, font_size=8.4, color="52606D", font_name=fonts["en_body"])
    add_text(slide, "rerun", 8.8, 4.35, 1.0, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_rect(slide, 8.8, 4.6, 2.95, 1.02, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "python build_weekly_report_pptx_oral_v2.py\n--config 20260311_weekly_report_pptx_oral_config_01.json\n--profile 20260311_liao_research_ppt_profile_01.json", 9.02, 4.82, 2.5, 0.72, font_size=7.2, color=theme["light_text"], font_name=fonts["en_body"])
    add_text(slide, "完整絕對路徑仍保留在主週報、設定檔與 version.json；本頁優先保證可口頭說明與快速複查。", 8.8, 5.92, 3.0, 0.34, font_size=8.6, color="52606D", font_name=fonts["zh_body"])


def build_notes(cfg: dict) -> list[str]:
    src = cfg["meta"]["source_report"]
    footer = cfg["meta"]["footer"]
    return [
        f"""頁面目的：快速建立報告主題、時間範圍與高層結論。
口頭重點：
1. 這份簡報整理的是 2026-03-05 到 2026-03-11 的研究主線。
2. 本週最重要收斂是三層分工：caller-first、methylation-support、annotation/QC/artifact triage。
3. 這份 deck 以口頭報告型為主，細節會放在後面與原始週報。
證據來源：
- 主週報：{src}
- deck footer：{footer}
""",
        """頁面目的：先讓聽眾知道簡報會怎麼走，避免在中途迷路。
口頭重點：
1. 先講破題結論，再講四象限資料邊界。
2. 中段聚焦 paired 與 tumor-only 的主要結果。
3. 後段收斂到 phase 2、annotation layer、缺口與下一步。
4. 最後會保留原始檔案與重生路徑，方便複查。
""",
        """頁面目的：用一頁先把本週最重要的結論講完。
口頭重點：
1. 目前最穩的救援訊號仍在 caller 端，不在甲基端。
2. 甲基訊號已經證明有價值，但現階段應定位成 support、annotation 與 triage。
3. 5kHz paired 有小幅正增益；5kHz TO 則最能看見甲基 support。
4. DORADO TO 證明甲基不是全負效果，但 Pairwise 方向不具全域一致性。
""",
        """頁面目的：交代四象限可以一起看，但不是所有數值都能直接硬比。
口頭重點：
1. 四個 dataset 已形成 paired / TO 與 5kHz / DORADO 的矩陣。
2. paired 與 TO 的 candidate-specific coverage 差很多。
3. 5kHz TO 用 full tagged BAM，DORADO TO 用 subset tagged BAM，所以 read-level 只能做方向比較。
4. PairwiseMedianDist 目前明顯是 dataset-aware 訊號。
""",
        """頁面目的：讓聽眾理解本週不是零散試規則，而是有節奏地推進。
口頭重點：
1. 先排除 mixed purity 作主線的干擾。
2. 再跑 paired pure 與 tumor-only 主流程。
3. 接著補 candidate-specific rescue、phase 2 與 annotation layer。
4. 這週真正的收穫是把哪些東西該進主流程、哪些只該留在分析層定型。
""",
        """頁面目的：說明 paired pure 場景的正式 benchmark 結果。
口頭重點：
1. 5kHz paired 是最 noisy 的主樣本，所以小幅提升仍有意義。
2. 5kHz paired 的正增益主要來自多移除 83 個 FP，不是大量救回 TP。
3. DORADO paired 幾乎持平略負，表示同一條規則不具自動可攜性。
4. 因此 paired 目前仍以 caller-first 為主，InterSubMod 較像 precision 精修層。
""",
        """頁面目的：整理 tumor-only rescue 的高層結論。
口頭重點：
1. TO 下甲基特徵不是全負效果，單獨看有正值。
2. 但在同一 caller gate 之後，甲基 support 目前仍沒有超過最佳 caller-only。
3. 5kHz TO 與 DORADO TO 都支持同一個高層分工：caller-first + methylation-support。
4. DORADO TO 目前仍需保留 subset haplotag 的解讀限制。
""",
        """頁面目的：把特徵重新分工，避免再把所有特徵混在同一層解讀。
口頭重點：
1. caller-first 負責定義第一層可救回集合。
2. methylation-support 負責排序、補強與局部支持。
3. annotation/QC/triage 負責 phase 完整度、artifact 與 dataset-aware 註記。
4. 這種分工比單純堆更多條件更重要。
""",
        """頁面目的：逐一說明哪些特徵可升級、哪些不能。
口頭重點：
1. GQ 仍是最穩定的主規則訊號。
2. Quality_Score 可升為 support annotation。
3. PairwiseMedianDist 只能升為 dataset-aware annotation，不能變全域規則。
4. hp_assign_rate 應明確降階為 phase/QC annotation。
5. low VAF + high AlleleDelta + low CramersV 應留在 artifact triage。
""",
        """頁面目的：交代 phase 2 為何重要，以及它排除了什麼方向。
口頭重點：
1. raw pileup / full model 的確揭露了 paired 的召回空間。
2. 但 precision 代價太大，因此不能直接取代 merged caller。
3. phase 2 的價值是理解 recall ceiling，而不是把 raw caller 當正式輸出。
4. paired 主線仍是 merged caller → LongPhase → InterSubMod。
""",
        """頁面目的：把 finer interval 結果正式回接到 annotation layer。
口頭重點：
1. Quality_Score 已足夠穩，適合直接輸出為 support annotation。
2. PairwiseMedianDist 有訊號，但必須標示 dataset-aware。
3. hp_assign_rate 反映的是 phase completeness，不應當成 truth keep。
4. 現在最合理的做法是先接 round summary / dashboard / diagnostics，而不是改 hard filter。
""",
        """頁面目的：主動交代現階段不能講太滿的地方。
口頭重點：
1. paired 的 candidate-specific coverage ceiling 仍偏低。
2. TO 的 pileup vs full model 還沒有完整雙模型對照。
3. 5kHz TO 與 DORADO TO 的 snapshot scope 不同，所以 read-level 解讀要保守。
4. Pool B FN 仍有 caller-side rescue 空間。
""",
        """頁面目的：把下一步說清楚，並避免研究又回到亂試規則。
口頭重點：
1. 下一步優先把 annotation 接到輸出層。
2. 繼續累積 finer interval 與 orthogonality 的 cross-dataset 證據。
3. 若要升級規則，先升級成 annotation export，不直接改核心過濾。
4. 不要把 5kHz 單一方向的 Pairwise 或 artifact triage 寫成全域規則。
""",
        f"""頁面目的：提供複查、再生與版本追蹤入口。
口頭重點：
1. 這份 deck 對應的主週報、phase 2 報告與 annotation 報告都可直接回查。
2. 生成腳本、設定檔與個人 PPT profile 已固定。
3. 若之後需要再版，應優先沿用這套 deck 生成流程。
主要來源：
- 主週報：{src}
- deck path：{cfg["meta"]["output_pptx"]}
""",
    ]


def add_speaker_notes(slide, note_text: str) -> None:
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = note_text


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", required=True, type=Path)
    parser.add_argument("--profile", required=True, type=Path)
    args = parser.parse_args()

    cfg = load_json(args.config)
    profile = load_json(args.profile)
    theme = cfg["theme"]
    fonts = {
        "zh_title": profile["fonts"]["zh_title"],
        "zh_body": profile["fonts"]["zh_body"],
        "en_title": profile["fonts"]["en_title"],
        "en_body": profile["fonts"]["en_body"],
    }

    assets = build_assets(cfg)
    prs = create_prs(cfg)
    notes = build_notes(cfg)
    builders = [
        lambda s: add_cover(s, cfg, theme, fonts, assets),
        lambda s: add_agenda(s, cfg, theme, fonts),
        lambda s: add_thesis(s, cfg, theme, fonts),
        lambda s: add_matrix(s, cfg, theme, fonts),
        lambda s: add_milestones(s, cfg, theme, fonts),
        lambda s: add_paired_results(s, cfg, theme, fonts, assets),
        lambda s: add_to_rescue(s, cfg, theme, fonts, assets),
        lambda s: add_feature_architecture(s, cfg, theme, fonts),
        lambda s: add_feature_cards(s, cfg, theme, fonts),
        lambda s: add_phase2(s, cfg, theme, fonts, assets),
        lambda s: add_annotation(s, cfg, theme, fonts, assets),
        lambda s: add_gaps(s, cfg, theme, fonts),
        lambda s: add_next_steps(s, cfg, theme, fonts),
        lambda s: add_appendix(s, cfg, theme, fonts),
    ]

    total = len(builders)
    dark_pages = {12}
    for idx, builder in enumerate(builders, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder(slide)
        add_speaker_notes(slide, notes[idx - 1])
        add_footer(slide, theme, fonts, idx, total, cfg["meta"]["footer"], dark=idx in dark_pages)

    out = Path(cfg["meta"]["output_pptx"])
    ensure_dir(out.parent)
    prs.save(out)
    print(f"saved_pptx={out}")
    print(f"assets_dir={cfg['meta']['assets_dir']}")


if __name__ == "__main__":
    main()
