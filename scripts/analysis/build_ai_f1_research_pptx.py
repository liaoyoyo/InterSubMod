#!/usr/bin/env python3
"""
將 2026-03-12 AI 自動化 F1 主線整合報告轉成 16 頁技術深度版 PPTX。

使用方式:
python scripts/analysis/build_ai_f1_research_pptx.py \
  --config docs/references/manual/assets/20260312_ai_f1_research_pptx_config_01.json \
  --profile docs/references/manual/assets/20260311_liao_research_ppt_profile_01.json \
  --output-pptx docs/presentations/draft/2026/03/<deck_name>/<deck_name>.pptx
"""

from __future__ import annotations

import argparse
import json
import re
from datetime import datetime
from pathlib import Path
from zoneinfo import ZoneInfo

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5
TZ = ZoneInfo("Asia/Taipei")


def bump_font(font_size: float) -> float:
    if font_size <= 8.5:
        return font_size + 0.7
    if font_size <= 10.0:
        return font_size + 0.6
    if font_size <= 12.0:
        return font_size + 0.5
    if font_size <= 14.0:
        return font_size + 0.4
    return font_size


def hex_rgb(value: str) -> RGBColor:
    value = value.replace("#", "").strip()
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)


def add_rect(slide, x, y, w, h, *, fill, line=None, radius=True, transparency=0):
    use_radius = radius and fill.upper() != "FFFFFF"
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if use_radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color="17202A",
    font_name="Arial",
    bold=False,
    italic=False,
    align="left",
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
    line_spacing=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(bump_font(font_size))
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_rgb(color)
    return box


def add_bullets(
    slide,
    bullets,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    color="17202A",
    font_name="Arial",
    space_after=6,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(bump_font(font_size))
        run.font.name = font_name
        run.font.color.rgb = hex_rgb(color)
    return box


def add_numbered_list(
    slide,
    items,
    x,
    y,
    w,
    h,
    *,
    font_size=18,
    text_color="17202A",
    font_name="Arial",
    number_fill="1E2761",
    number_color="F0F4FB",
    item_fill="F5F8FD",
    line=None,
):
    item_h = h / max(len(items), 1)
    for idx, item in enumerate(items):
        iy = y + idx * item_h
        card_h = max(item_h - 0.08, 0.42)
        add_rect(slide, x, iy, w, card_h, fill=item_fill, line=line, radius=True)
        add_rect(slide, x + 0.10, iy + 0.10, 0.40, 0.40, fill=number_fill, line=None, radius=True)
        add_text(
            slide,
            str(idx + 1),
            x + 0.20,
            iy + 0.18,
            0.18,
            0.12,
            font_size=9.8,
            color=number_color,
            font_name=font_name,
            bold=True,
            align="center",
            margin=0,
        )
        add_text(
            slide,
            item,
            x + 0.62,
            iy + 0.07,
            w - 0.78,
            card_h - 0.12,
            font_size=font_size,
            color=text_color,
            font_name=font_name,
            margin=0.02,
        )


def add_chevron(slide, x, y, w, h, *, fill):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.line.fill.background()
    return shape


def short_rule(rule: str) -> str:
    return (
        rule.replace("PairwiseMedianDist", "Pairwise")
        .replace("Quality_Score", "Quality")
        .replace("agreement_positive", "agreement+")
    )


def wrap_text(text: str, max_len: int = 34, max_lines: int = 5) -> str:
    parts = [token for token in re.split(r"([/ _-])", text) if token]
    lines = []
    current = ""
    for part in parts:
        if not current:
            current = part
        elif len(current) + len(part) <= max_len:
            current += part
        else:
            lines.append(current.rstrip())
            current = part.lstrip()
    if current:
        lines.append(current.rstrip())
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1][: max_len - 1] + "…"
    return "\n".join(lines)


def short_path(text: str, keep_parts: int = 4) -> str:
    parts = [part for part in text.split("/") if part]
    if len(parts) <= keep_parts:
        return text
    return "/.../" + "/".join(parts[-keep_parts:])


def add_header(slide, zh_title, en_title, zh_subtitle, theme, fonts, *, dark=False):
    fg = theme["light_text"] if dark else theme["text"]
    muted = theme["light_muted"] if dark else theme["muted"]
    accent = theme["accent2"] if dark else theme["accent"]
    add_text(slide, zh_title, 0.72, 0.36, 9.8, 0.64, font_size=25.5, color=fg, font_name=fonts["zh_title"], bold=True)
    add_text(slide, en_title, 0.74, 0.98, 9.2, 0.24, font_size=11.2, color=accent, font_name=fonts["en_body"])
    add_text(slide, zh_subtitle, 0.74, 1.24, 11.6, 0.38, font_size=13.0, color=muted, font_name=fonts["zh_body"])


def add_footer(slide, theme, fonts, page_num, total_pages, footer_text, *, dark=False):
    color = theme["light_muted"] if dark else theme["text"]
    add_text(slide, footer_text, 0.72, 7.08, 8.9, 0.18, font_size=10.0, color=color, font_name=fonts["en_body"], margin=0)
    add_text(slide, f"{page_num}/{total_pages}", 11.86, 7.04, 0.66, 0.18, font_size=11.4, color=color, font_name=fonts["en_body"], align="right", margin=0, bold=True)


def style_table(table, theme: dict, fonts: dict, *, header_fill: str, header_font: str, body_fill: str = "FFFFFF", zebra_fill: str | None = None, font_size: float = 11):
    rows = len(table.rows)
    cols = len(table.columns)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = hex_rgb(header_fill)
            elif zebra_fill and r % 2 == 0:
                cell.fill.fore_color.rgb = hex_rgb(zebra_fill)
            else:
                cell.fill.fore_color.rgb = hex_rgb(body_fill)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = fonts["zh_body"]
                    run.font.size = Pt(bump_font(font_size))
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = hex_rgb(header_font)
                    else:
                        run.font.color.rgb = hex_rgb(theme["text"])
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, rows, x, y, w, h, theme, fonts, *, header_fill, body_fill="FFFFFF", zebra_fill=None, font_size=11, col_widths=None):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        total = sum(col_widths)
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(w * (width / total))
    else:
        col_width = Inches(w / len(rows[0]))
        for idx in range(len(rows[0])):
            table.columns[idx].width = col_width
    row_h = Inches(h / len(rows))
    for idx in range(len(rows)):
        table.rows[idx].height = row_h
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = str(value)
    style_table(
        table,
        theme,
        fonts,
        header_fill=header_fill,
        header_font=theme["light_text"],
        body_fill=body_fill,
        zebra_fill=zebra_fill,
        font_size=font_size,
    )
    return table


def add_speaker_notes(slide, note_text: str) -> None:
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = note_text


def make_paired_chart(path: Path, categories, series):
    fig, ax = plt.subplots(figsize=(7.6, 4.2), dpi=240)
    x = range(len(categories))
    width = 0.22
    offsets = [-width, 0, width]
    values = [val for item in series for val in item["values"]]
    ax.set_ylim(min(values) - 0.010, max(values) + 0.007)
    for idx, item in enumerate(series):
        positions = [v + offsets[idx] for v in x]
        ax.bar(positions, item["values"], width=width * 0.9, color=f"#{item['color']}", label=item["name"])
        for px, val in zip(positions, item["values"]):
            ax.text(px, val + 0.0007, f"{val:.4f}", ha="center", va="bottom", fontsize=9.5, fontweight="bold")
    ax.set_xticks(list(x))
    ax.set_xticklabels(categories, fontsize=10.8, fontweight="bold")
    ax.tick_params(axis="y", labelsize=10.8)
    ax.set_ylabel("F1 score", fontsize=12.6, fontweight="bold")
    ax.set_title("Paired Benchmark by Method", fontsize=15.2, fontweight="bold", pad=14)
    ax.grid(axis="y", linestyle="--", alpha=0.22)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(frameon=False, fontsize=10.4, ncol=3, loc="upper center", bbox_to_anchor=(0.5, 1.17))
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def make_to_chart(path: Path, slides_cfg: dict):
    fig, axes = plt.subplots(1, 2, figsize=(8.9, 4.1), dpi=240)
    panel_specs = [
        ("5kHz TO", slides_cfg["to_results"]["5khz"]["rules"]),
        ("DORADO TO", slides_cfg["to_results"]["dorado"]["rules"]),
    ]
    palette = {
        "caller-first": "#1E2761",
        "methyl-support": "#2D5EA8",
        "label-support": "#6DAEE0",
    }
    for ax, (title, rows) in zip(axes, panel_specs):
        labels = [short_rule(row["rule"]) for row in rows]
        values = [float(row["delta_f1"].replace("+", "")) for row in rows]
        colors = [palette.get(row["type"], "#6DAEE0") for row in rows]
        ax.bar(labels, values, color=colors, width=0.58)
        ax.set_title(title, fontsize=13.0, fontweight="bold", pad=10)
        ax.set_ylabel("ΔF1 vs baseline", fontsize=10.8, fontweight="bold")
        ax.tick_params(axis="x", labelrotation=12, labelsize=9.6)
        ax.tick_params(axis="y", labelsize=9.8)
        ax.grid(axis="y", linestyle="--", alpha=0.22)
        ymax = max(values) * 1.34
        ax.set_ylim(0, ymax)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        for idx, val in enumerate(values):
            ax.text(idx, val + ymax * 0.02, f"{val:.4f}", ha="center", va="bottom", fontsize=8.8, fontweight="bold")
    plt.tight_layout(w_pad=2.0)
    plt.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close()


def build_assets(slides_cfg: dict, assets_dir: Path) -> dict:
    ensure_dir(assets_dir)
    paired_path = assets_dir / "paired_results.png"
    to_path = assets_dir / "to_results.png"
    make_paired_chart(paired_path, slides_cfg["paired_results"]["chart_categories"], slides_cfg["paired_results"]["chart_series"])
    make_to_chart(to_path, slides_cfg)
    return {"paired": paired_path, "to": to_path}


def create_prs(deck_title: str, author: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = author
    prs.core_properties.title = deck_title
    return prs


def add_cover(slide, meta, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_rect(slide, 0.48, 0.55, 5.3, 6.1, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "研究主線結構圖", 0.82, 0.86, 2.1, 0.18, font_size=12.4, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
    layer_positions = [(0.82, 1.28), (1.10, 2.45), (1.38, 3.62)]
    for (x, y), layer in zip(layer_positions, slides_cfg["framework"]["layers"]):
        add_rect(slide, x, y, 3.78, 0.9, fill=layer["fill"], line=None, radius=True)
        add_text(slide, layer["title"], x + 0.18, y + 0.12, 3.15, 0.18, font_size=13.2, color=layer["title_color"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, layer["body"], x + 0.18, y + 0.42, 3.25, 0.22, font_size=10.0, color=layer["body_color"], font_name=fonts["zh_body"])
    add_text(slide, "四象限資料", 0.82, 4.86, 1.8, 0.18, font_size=12, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
    quadrants = [
        ("HKU_5kHz\npaired", 1.84, 5.14, "DCE8F8"),
        ("DORADO\npaired", 3.44, 5.14, "CBE0F3"),
        ("HKU_5kHz\nTO", 1.84, 5.82, "DCE8F8"),
        ("DORADO\nTO", 3.44, 5.82, "CBE0F3"),
    ]
    for label, x, y, fill in quadrants:
        add_rect(slide, x, y, 1.42, 0.60, fill=fill, line=None, radius=True)
        add_text(slide, label, x + 0.10, y + 0.10, 1.20, 0.28, font_size=10.2, color=theme["text"], font_name=fonts["zh_title"], bold=True, align="center")

    add_text(slide, meta["kicker"], 6.22, 0.92, 4.9, 0.18, font_size=11.4, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, meta["deck_title"], 6.2, 1.32, 6.4, 1.12, font_size=28.4, color=theme["text"], font_name=fonts["zh_title"], bold=True, line_spacing=1.02)
    add_text(slide, meta["deck_subtitle"], 6.22, 2.48, 5.9, 0.34, font_size=12.8, color=theme["muted"], font_name=fonts["en_body"])
    add_rect(slide, 6.18, 3.00, 5.95, 1.18, fill=theme["dark"], line=None, radius=True)
    add_text(slide, meta["summary_line"], 6.44, 3.24, 5.35, 0.64, font_size=15.0, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_rect(slide, 6.18, 4.52, 2.1, 0.78, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "16 頁技術深度版", 6.44, 4.78, 1.62, 0.18, font_size=11.8, color=theme["light_text"], font_name=fonts["zh_title"], bold=True, align="center")
    add_rect(slide, 8.48, 4.52, 2.1, 0.78, fill=theme["accent2"], line=None, radius=True)
    add_text(slide, meta["audience_mode"], 8.70, 4.70, 1.66, 0.36, font_size=10.1, color=theme["text"], font_name=fonts["zh_title"], bold=True, align="center")
    add_rect(slide, 10.78, 4.52, 1.34, 0.78, fill="DCE8F8", line=None, radius=True)
    add_text(slide, "draft", 11.02, 4.78, 0.82, 0.18, font_size=11.8, color=theme["text"], font_name=fonts["en_title"], bold=True, align="center")
    add_text(slide, f"報告整理時間｜{meta['generated_date']}", 6.22, 5.72, 2.8, 0.18, font_size=11.0, color=theme["muted"], font_name=fonts["zh_body"])
    add_text(slide, f"報告人｜{meta['author']}", 9.18, 5.72, 2.0, 0.18, font_size=11.0, color=theme["muted"], font_name=fonts["zh_body"])
    add_text(slide, "所有細節推論與絕對路徑均嵌入 speaker notes。", 6.22, 6.14, 5.2, 0.18, font_size=10.0, color=theme["accent"], font_name=fonts["zh_body"])


def add_agenda(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "報告總覽", "Agenda", "先講研究問題與主結論，再依序處理 paired、TO、Pool B、annotation 與 mixed sample。", theme, fonts)
    modules = [
        ("01", "主結論與 AI 自動化", "第 3-6 頁"),
        ("02", "樣本策略與 paired/TO", "第 7-9 頁"),
        ("03", "前段 TP rescue：Pool B FN", "第 10 頁"),
        ("04", "特徵分工與 annotation", "第 11-13 頁"),
        ("05", "mixed sample 與下一步", "第 14-15 頁"),
        ("06", "附錄與重生路徑", "第 16 頁"),
    ]
    for idx, (num, zh, page_range) in enumerate(modules):
        x = 0.82 + (idx % 2) * 6.05
        y = 1.92 + (idx // 2) * 1.42
        add_rect(slide, x, y, 5.35, 1.02, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.18, y + 0.18, 0.58, 0.52, fill=theme["dark"], line=None, radius=True)
        add_text(slide, num, x + 0.34, y + 0.33, 0.24, 0.12, font_size=10.4, color=theme["light_text"], font_name=fonts["en_title"], bold=True, align="center")
        add_text(slide, zh, x + 0.95, y + 0.16, 3.45, 0.18, font_size=13.8, color=theme["text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, page_range, x + 0.95, y + 0.48, 1.6, 0.18, font_size=10.2, color=theme["accent"], font_name=fonts["zh_body"], bold=True)
    add_rect(slide, 0.82, 6.36, 11.7, 0.42, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "高層導讀：本 deck 的重點不是再推一條新規則，而是把 caller、甲基 support、annotation 與 mixed 假說放到正確層次。", 1.02, 6.48, 11.1, 0.14, font_size=9.8, color=theme["muted"], font_name=fonts["zh_body"])


def add_thesis(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["dark"])
    add_header(slide, "破題結論", "Opening Thesis", "目前最穩定的結論是 caller-first；甲基訊號有價值，但還沒有穩定升級成第一層主規則。", theme, fonts, dark=True)
    add_rect(slide, 0.82, 1.88, 6.55, 1.72, fill="2D5EA8", line=None, radius=True)
    add_text(slide, "最穩定結論：caller-first 主規則 → methylation-support / annotation → haplotype / SNV / 甲基共分析。", 1.10, 2.14, 5.95, 0.78, font_size=16.5, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_text(slide, "Stable structure: caller-first, then methylation support/annotation, then multi-omics co-analysis.", 1.10, 3.08, 5.75, 0.26, font_size=10.2, color=theme["light_muted"], font_name=fonts["en_body"])
    for idx, layer in enumerate(slides_cfg["framework"]["layers"]):
        x = 0.82 + idx * 4.06
        add_rect(slide, x, 4.06, 3.72, 1.45, fill=layer["fill"], line=None, radius=True)
        add_text(slide, layer["title"], x + 0.16, 4.22, 3.1, 0.18, font_size=12.2, color=layer["title_color"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, layer["title_en"], x + 0.16, 4.48, 3.05, 0.18, font_size=9.2, color=layer["title_color"], font_name=fonts["en_body"])
        add_text(slide, layer["body"], x + 0.16, 4.78, 3.05, 0.34, font_size=9.6, color=layer["body_color"], font_name=fonts["zh_body"])
    add_rect(slide, 8.05, 1.88, 4.46, 1.85, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "三句一定要先講清楚", 8.28, 2.08, 2.8, 0.18, font_size=14.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(slide, slides_cfg["framework"]["conclusions"], 8.20, 2.42, 4.0, 1.00, font_size=10.4, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 6.02, 11.7, 0.48, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "這份簡報之後的所有數據頁，都只是把這個結論拆開驗證，而不是推翻它。", 1.04, 6.16, 11.1, 0.16, font_size=10.6, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_ai_automation(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "為什麼最近改用 AI / Agent 自動化", "Why AI-assisted Automation", "研究策略已從手動觀察零散位點，轉成可持續重跑、回寫、追版本的研究循環。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 5.45, 3.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 6.52, 1.92, 5.95, 3.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "Before", 1.05, 2.10, 1.2, 0.18, font_size=12.5, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    add_text(slide, "After", 6.78, 2.10, 1.2, 0.18, font_size=12.5, color=theme["accent"], font_name=fonts["en_title"], bold=True)
    for idx, row in enumerate(slides_cfg["ai_automation"]["before_after"]):
        y = 2.46 + idx * 0.84
        add_rect(slide, 1.02, y, 4.95, 0.62, fill="F5F8FD", line=None, radius=True)
        add_rect(slide, 6.72, y, 5.45, 0.62, fill="DCE8F8", line=None, radius=True)
        add_text(slide, row["before"], 1.22, y + 0.18, 4.5, 0.14, font_size=11.0, color=theme["text"], font_name=fonts["zh_body"])
        add_text(slide, row["after"], 6.96, y + 0.12, 5.0, 0.34, font_size=10.3, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 6.10, 11.7, 0.48, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "研究價值：" + "；".join(slides_cfg["ai_automation"]["value"]), 1.06, 6.22, 11.1, 0.18, font_size=10.1, color=theme["light_text"], font_name=fonts["zh_body"])


def add_two_problems(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "這輪刻意只限制 AI 先做兩類問題", "Two Restricted Research Questions", "前段救 TP 與後段刪 FP 必須拆開談，因為兩者母體不同。", theme, fonts)
    a = slides_cfg["two_problems"]["problem_a"]
    b = slides_cfg["two_problems"]["problem_b"]
    for idx, block in enumerate([a, b]):
        x = 0.82 + idx * 6.05
        add_rect(slide, x, 1.92, 5.35, 4.38, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.18, 2.10, 2.20, 0.54, fill=theme["dark"] if idx == 0 else theme["accent"], line=None, radius=True)
        add_text(slide, block["title"], x + 0.36, 2.24, 4.4, 0.18, font_size=13.9, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, block["title_en"], x + 0.20, 2.80, 3.0, 0.18, font_size=10.8, color=theme["accent"], font_name=fonts["en_body"], bold=True)
        add_text(slide, block["description"], x + 0.20, 3.10, 4.78, 0.56, font_size=11.4, color=theme["text"], font_name=fonts["zh_body"])
        add_text(slide, "典型研究子集", x + 0.20, 3.86, 1.6, 0.16, font_size=10.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_bullets(slide, block["targets"], x + 0.18, 4.12, 4.65, 0.98, font_size=10.8, color=theme["text"], font_name=fonts["zh_body"])
        add_rect(slide, x + 0.18, 5.34, 4.95, 0.54, fill="DCE8F8" if idx == 0 else "CBE0F3", line=None, radius=True)
        add_text(slide, block["note"], x + 0.36, 5.52, 4.55, 0.16, font_size=10.3, color=theme["text"], font_name=fonts["zh_body"], bold=True)
    add_rect(slide, 0.82, 6.30, 11.7, 0.28, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, slides_cfg["two_problems"]["key_difference"], 1.04, 6.38, 11.1, 0.12, font_size=10.2, color=theme["muted"], font_name=fonts["zh_body"])


def add_workflow(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "自動化研究流程與 tagged BAM 的角色", "Workflow", "這條流程真正的核心，是把報告、資料、LongPhase、InterSubMod 與回寫索引接成一條可重跑鏈。", theme, fonts)
    step_positions = [
        (0.92, 2.08), (3.08, 2.08), (5.24, 2.08), (7.40, 2.08),
        (1.98, 4.18), (4.14, 4.18), (6.30, 4.18),
    ]
    for idx, ((x, y), text) in enumerate(zip(step_positions, slides_cfg["workflow_steps"]), start=1):
        fill = theme["dark"] if idx in {1, 4, 7} else theme["accent"] if idx in {2, 5} else theme["accent2"]
        fg = theme["light_text"] if fill != theme["accent2"] else theme["text"]
        add_rect(slide, x, y, 1.9, 1.18, fill=fill, line=None, radius=True)
        add_rect(slide, x + 0.12, y + 0.10, 0.38, 0.38, fill="FFFFFF", line=None, radius=True)
        add_text(slide, str(idx), x + 0.22, y + 0.18, 0.16, 0.10, font_size=10.0, color=theme["text"], font_name=fonts["en_title"], bold=True, align="center")
        add_text(slide, text, x + 0.14, y + 0.56, 1.58, 0.42, font_size=9.5, color=fg, font_name=fonts["zh_body"], align="center")
    add_rect(slide, 9.72, 2.08, 2.18, 3.28, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "tagged BAM 為什麼重要", 9.94, 2.26, 1.8, 0.18, font_size=13.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(
        slide,
        [
            "InterSubMod 需要 haplotagged BAM。",
            "沒有 tagged BAM，就不能做真正的 HP / label-first 驗證。",
            "因此 paired 用 LongPhase-S，TO 用 LongPhase-TO。",
        ],
        9.88,
        2.64,
        1.82,
        1.70,
        font_size=10.0,
        color=theme["text"],
        font_name=fonts["zh_body"],
    )
    add_rect(slide, 9.88, 4.72, 1.84, 0.46, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "先純樣本，再 mixed", 10.02, 4.86, 1.56, 0.14, font_size=9.8, color=theme["light_text"], font_name=fonts["zh_title"], bold=True, align="center")


def add_sample_strategy(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "樣本策略：為什麼先做 5kHz，再用 DORADO 驗證", "Sample Strategy", "5kHz 提供規則壓力測試；DORADO 負責檢查這些規則是否只是平台特異性。", theme, fonts)
    primary = slides_cfg["sample_strategy"]["primary"]
    secondary = slides_cfg["sample_strategy"]["secondary"]
    for idx, block in enumerate([primary, secondary]):
        x = 0.82 + idx * 6.08
        fill = "FFFFFF"
        banner = theme["dark"] if idx == 0 else theme["accent"]
        add_rect(slide, x, 1.98, 5.45, 4.42, fill=fill, line=theme["line"], radius=True)
        add_rect(slide, x + 0.18, 2.16, 2.44, 0.60, fill=banner, line=None, radius=True)
        add_text(slide, block["label"], x + 0.34, 2.34, 4.6, 0.18, font_size=13.0, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, block["label_en"], x + 0.20, 2.94, 2.8, 0.18, font_size=10.2, color=theme["accent"], font_name=fonts["en_body"], bold=True)
        add_bullets(slide, block["points"], x + 0.18, 3.28, 4.88, 2.10, font_size=10.8, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 6.56, 11.7, 0.34, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, slides_cfg["sample_strategy"]["disk_note"], 1.06, 6.64, 11.1, 0.12, font_size=9.2, color=theme["muted"], font_name=fonts["zh_body"])


def add_paired_results(slide, slides_cfg, assets, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "paired 主結果", "Paired Results", "paired 問題目前最合理的口徑是：caller 與 LongPhase 已做掉主要改善，InterSubMod 在 5kHz 上再做少量 precision 精修。", theme, fonts)
    slide.shapes.add_picture(str(assets["paired"]), Inches(0.78), Inches(1.98), width=Inches(7.52))
    add_rect(slide, 0.78, 5.98, 7.52, 0.48, fill="F5F8FD", line=theme["line"], radius=True)
    add_text(slide, "圖表解讀：兩個 paired dataset 都先經過 ClairS → LongPhase-S → InterSubMod，比較 F1 變化與增益來源。", 0.98, 6.12, 7.10, 0.18, font_size=10.6, color=theme["muted"], font_name=fonts["zh_body"])

    right_cards = [
        {
            "title": "5kHz paired",
            "lines": [
                "ClairS 0.8443 → LongPhase-S 0.8522 → InterSubMod 0.8532",
                "增益來源：FP 627 → 544，額外少 83 個 FP",
                "TP 並未大量救回，反而少 2 個",
            ],
            "fill": "DCE8F8",
        },
        {
            "title": "DORADO paired",
            "lines": [
                "ClairS 0.8565 → LongPhase-S 0.8592 → InterSubMod 0.8590",
                "InterSubMod 與 LongPhase-S 幾乎持平略負",
                "代表規則方向不可直接宣稱跨平台可攜",
            ],
            "fill": "CBE0F3",
        },
    ]
    for idx, card in enumerate(right_cards):
        x = 8.34
        y = 2.04 + idx * 1.86
        add_rect(slide, x, y, 4.04, 1.62, fill=card["fill"], line=None, radius=True)
        add_text(slide, card["title"], x + 0.18, y + 0.16, 2.0, 0.16, font_size=12.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_bullets(slide, card["lines"], x + 0.16, y + 0.42, 3.58, 0.96, font_size=10.4, color=theme["text"], font_name=fonts["zh_body"], space_after=2)
    add_rect(slide, 8.34, 5.96, 4.04, 0.58, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "重點：paired 的正式口徑是『precision 精修』，不是『大量 TP rescue』。", 8.54, 6.14, 3.64, 0.16, font_size=9.6, color=theme["light_text"], font_name=fonts["zh_body"], bold=True)


def add_to_results(slide, slides_cfg, assets, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "tumor-only 主結果", "Tumor-only Results", "TO 下 caller-first 與甲基 support 都成立，但在固定 caller gate 後，甲基尚未穩定超過 caller-only。", theme, fonts)
    slide.shapes.add_picture(str(assets["to"]), Inches(0.76), Inches(1.98), width=Inches(7.30))
    add_rect(slide, 0.76, 5.98, 7.32, 0.48, fill="F5F8FD", line=theme["line"], radius=True)
    add_text(slide, "圖表解讀：左圖為 5kHz TO，右圖為 DORADO TO；Y 軸皆為相對 baseline 的 ΔF1。", 0.96, 6.12, 6.90, 0.18, font_size=10.2, color=theme["muted"], font_name=fonts["zh_body"])
    blocks = [slides_cfg["to_results"]["5khz"], slides_cfg["to_results"]["dorado"]]
    fills = ["DCE8F8", "CBE0F3"]
    for idx, (block, fill) in enumerate(zip(blocks, fills)):
        x = 8.20
        y = 2.02 + idx * 2.04
        add_rect(slide, x, y, 4.16, 1.94, fill=fill, line=None, radius=True)
        add_text(slide, block["label"], x + 0.18, y + 0.14, 2.6, 0.18, font_size=12.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        summary_lines = [
            f"caller-first：{short_rule(block['rules'][0]['rule'])}，{block['rules'][0]['tp_rescued']} TP / {block['rules'][0]['fp_reintro']} FP，ΔF1 {block['rules'][0]['delta_f1']}",
            f"support 1：{short_rule(block['rules'][1]['rule'])}，{block['rules'][1]['tp_rescued']} TP / {block['rules'][1]['fp_reintro']} FP，ΔF1 {block['rules'][1]['delta_f1']}",
            f"support 2：{short_rule(block['rules'][2]['rule'])}，{block['rules'][2]['tp_rescued']} TP / {block['rules'][2]['fp_reintro']} FP，ΔF1 {block['rules'][2]['delta_f1']}",
        ]
        add_bullets(slide, summary_lines, x + 0.16, y + 0.38, 3.70, 1.12, font_size=9.8, color=theme["text"], font_name=fonts["zh_body"], space_after=2)
        add_text(slide, block["note"], x + 0.20, y + 1.50, 3.66, 0.18, font_size=9.3, color=theme["accent"], font_name=fonts["zh_body"], bold=True)


def add_pool_b(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "Pool B FN：前段 TP rescue 問題的正式收尾", "Pool B FN Closeout", "Pool B 再次證明 caller-side rescue 空間存在，但甲基只提供弱補強，沒有翻盤。", theme, fonts)
    add_rect(slide, 0.82, 1.92, 3.05, 4.82, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "Pool B FN", 1.28, 2.22, 2.0, 0.22, font_size=16, color=theme["light_text"], font_name=fonts["zh_title"], bold=True, align="center")
    add_text(slide, str(slides_cfg["pool_b"]["fn_count"]), 1.18, 2.95, 2.3, 0.70, font_size=40, color=theme["light_text"], font_name=fonts["en_title"], bold=True, align="center")
    add_text(slide, "個 non-PASS 但 truth set 包含的位點", 1.04, 3.78, 2.6, 0.30, font_size=10.2, color=theme["light_muted"], font_name=fonts["zh_body"], align="center")
    add_rect(slide, 1.10, 4.68, 2.48, 1.18, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "這是一條『前段救 TP』問題，\n不能和 kept-set artifact triage 混寫。", 1.28, 5.02, 2.08, 0.42, font_size=10.0, color=theme["light_text"], font_name=fonts["zh_title"], bold=True, align="center")
    for idx, row in enumerate(slides_cfg["pool_b"]["results"]):
        x = 4.18
        y = 1.98 + idx * 1.50
        add_rect(slide, x, y, 4.22, 1.24, fill="FFFFFF", line=theme["line"], radius=True)
        title = row["rule"]
        title = title.replace("no_varcluster（caller-only）", "no_varcluster")
        title = title.replace("no_varcluster_and_gq15（最乾淨）", "no_varcluster + gq15")
        title = title.replace("gq15_and_allele_delta_low（combined）", "gq15 + allele_delta_low")
        add_text(slide, title, x + 0.18, y + 0.12, 3.2, 0.18, font_size=11.2, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        metric_parts = []
        if row.get("tp"):
            metric_parts.append(f"TP {row['tp']}")
        if row.get("fp"):
            metric_parts.append(f"FP {row['fp']}")
        if row.get("precision"):
            metric_parts.append(f"precision {row['precision']}")
        metric_parts.append(f"ΔF1 {row['delta_f1']}")
        add_text(slide, " | ".join(metric_parts), x + 0.18, y + 0.44, 3.55, 0.18, font_size=9.8, color=theme["text"], font_name=fonts["zh_body"], bold=True)
        add_text(slide, row["note"], x + 0.18, y + 0.76, 3.7, 0.18, font_size=9.2, color=theme["muted"], font_name=fonts["zh_body"])
    add_rect(slide, 8.66, 1.98, 3.86, 4.24, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "Pool B 要帶出的口徑", 8.90, 2.18, 2.2, 0.18, font_size=13.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(slide, slides_cfg["pool_b"]["conclusions"], 8.88, 2.56, 3.25, 2.15, font_size=10.2, color=theme["text"], font_name=fonts["zh_body"])
    add_rect(slide, 8.88, 5.36, 3.20, 0.48, fill="DCE8F8", line=None, radius=True)
    add_text(slide, "關鍵：AlleleDelta 在 Pool B 與 kept-set / artifact triage 的方向不同。", 9.04, 5.50, 2.92, 0.16, font_size=8.6, color=theme["text"], font_name=fonts["zh_body"], bold=True)


def add_feature_table(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "現在最重要的特徵與分工", "Feature Roles", "有用的特徵不一定能升級成主規則；有些特徵最合理的位置，是 support、annotation 或 artifact 旁路。", theme, fonts)
    columns = [slides_cfg["feature_table"][:4], slides_cfg["feature_table"][4:]]
    for col_idx, items in enumerate(columns):
        x = 0.82 + col_idx * 5.92
        for idx, row in enumerate(items):
            y = 1.92 + idx * 1.06
            add_rect(slide, x, y, 5.48, 0.88, fill="FFFFFF", line=theme["line"], radius=True)
            if "✓" in row["promotion"]:
                badge_fill = theme["dark"]
                badge_fg = theme["light_text"]
            elif "↑" in row["promotion"]:
                badge_fill = theme["accent"]
                badge_fg = theme["light_text"]
            else:
                badge_fill = "DCE8F8"
                badge_fg = theme["text"]
            add_text(slide, row["feature"], x + 0.18, y + 0.12, 2.30, 0.18, font_size=12.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
            add_rect(slide, x + 3.90, y + 0.12, 1.24, 0.34, fill=badge_fill, line=None, radius=True)
            add_text(slide, row["promotion"], x + 4.02, y + 0.20, 1.00, 0.10, font_size=8.8, color=badge_fg, font_name=fonts["zh_body"], bold=True, align="center")
            add_text(slide, f"定位：{row['position']}", x + 0.18, y + 0.42, 2.30, 0.14, font_size=9.8, color=theme["accent"], font_name=fonts["zh_body"], bold=True)
            add_text(slide, f"限制：{row['constraint']}", x + 2.24, y + 0.42, 3.00, 0.20, font_size=9.2, color=theme["muted"], font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 6.20, 11.7, 0.40, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "這張表的重點不是『誰有用』，而是『每個訊號現在應該被擺在哪一層』。", 1.06, 6.32, 11.1, 0.14, font_size=10.4, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_annotation_architecture(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "annotation / 三層架構：現在最合理的整合方式", "Annotation Architecture", "這個三層架構不是抽象口號，而是目前最適合繼續研究與輸出的方法學切分。", theme, fonts)
    layer_specs = [
        (0.68, 3.62, "第一層：caller-first", slides_cfg["framework"]["layers"][0]["body"], slides_cfg["framework"]["layers"][0]["fill"], slides_cfg["framework"]["layers"][0]["title_color"], slides_cfg["framework"]["layers"][0]["body_color"]),
        (4.48, 3.56, "第二層：methylation-support", slides_cfg["framework"]["layers"][1]["body"], slides_cfg["framework"]["layers"][1]["fill"], slides_cfg["framework"]["layers"][1]["title_color"], slides_cfg["framework"]["layers"][1]["body_color"]),
        (8.20, 3.64, "第三層：haplotype / SNV /\n甲基共分析", slides_cfg["framework"]["layers"][2]["body"], slides_cfg["framework"]["layers"][2]["fill"], slides_cfg["framework"]["layers"][2]["title_color"], slides_cfg["framework"]["layers"][2]["body_color"]),
    ]
    for x, w, title, body, fill, title_color, body_color in layer_specs:
        add_rect(slide, x, 1.98, w, 1.60, fill=fill, line=None, radius=True)
        add_text(slide, title, x + 0.18, 2.16, w - 0.34, 0.28, font_size=12.0, color=title_color, font_name=fonts["zh_title"], bold=True)
        add_text(slide, wrap_text(body, 32, 3), x + 0.18, 2.56, w - 0.34, 0.34, font_size=9.8, color=body_color, font_name=fonts["zh_body"])
    add_chevron(slide, 4.18, 2.46, 0.22, 0.54, fill=theme["accent"])
    add_chevron(slide, 7.96, 2.46, 0.22, 0.54, fill=theme["accent2"])
    cards = [
        ("GQ / QUAL / FILTER", "caller-first 主規則", "第一層先避免漏掉太多 TP", theme["dark"], theme["light_text"]),
        ("Quality_Score", "support annotation", "在 TO 兩個資料集都不是負訊號", theme["accent"], theme["light_text"]),
        ("PairwiseMedianDist", "dataset-aware annotation", "方向依 dataset 而變，不可全域化", "DCE8F8", theme["text"]),
        ("hp_assign_rate", "phase / QC annotation", "反映相位完整度，不應直接拿來 truth keep", "CBE0F3", theme["text"]),
    ]
    for idx, (title, role, note, fill, fg) in enumerate(cards):
        x = 0.82 + (idx % 2) * 6.00
        y = 4.12 + (idx // 2) * 1.10
        add_rect(slide, x, y, 5.32, 0.92, fill=fill, line=None, radius=True)
        add_text(slide, title, x + 0.18, y + 0.12, 2.3, 0.16, font_size=11.4, color=fg, font_name=fonts["en_title"], bold=True)
        add_text(slide, role, x + 2.58, y + 0.12, 2.3, 0.16, font_size=10.0, color=fg, font_name=fonts["zh_title"], bold=True)
        add_text(slide, note, x + 0.18, y + 0.44, 4.9, 0.18, font_size=8.9, color=fg, font_name=fonts["zh_body"])
    add_rect(slide, 0.82, 6.24, 11.7, 0.34, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "現階段最合理的策略是：第一層盡量避免漏掉太多 TP，第二層避免帶回太多 FP，第三層負責驗證與可解釋性。", 1.06, 6.32, 11.1, 0.12, font_size=10.2, color=theme["muted"], font_name=fonts["zh_body"])


def add_comparability(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "四象限可比性與解讀邊界", "Comparability Limits", "5kHz 與 DORADO 的差異可用來做方向判讀，但 dataset-dependence 仍需保守解讀。", theme, fonts)
    add_rect(slide, 0.82, 1.96, 5.62, 4.62, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 6.68, 1.96, 5.62, 4.62, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "為什麼 5kHz / DORADO 會分化", 1.06, 2.16, 3.3, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_bullets(slide, slides_cfg["divergence"]["bullets"], 1.06, 2.54, 4.82, 2.08, font_size=10.8, color=theme["text"], font_name=fonts["zh_body"], space_after=4)
    add_rect(slide, 1.06, 5.10, 4.82, 0.64, fill="DCE8F8", line=None, radius=True)
    add_text(slide, "合理推論：noisy sample 對甲基特徵更敏感，但這還不是已證明因果。", 1.22, 5.30, 4.42, 0.18, font_size=9.6, color=theme["text"], font_name=fonts["zh_body"], bold=True)

    add_text(slide, "四象限比較時不能忘記的限制", 6.92, 2.16, 3.5, 0.18, font_size=13.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    limit_bullets = [
        "paired candidate-specific coverage 只有 8%~20%，TO 才到 71%~100%。",
        "5kHz TO snapshot 來自 full tagged BAM；DORADO TO 來自 subset tagged BAM。",
        "PairwiseMedianDist 在 5kHz TO 偏高、在 DORADO TO 偏低，因此只能做 dataset-aware support annotation。",
        "甲基 rescue 在 paired 目前沒看到穩定正向，不代表 paired 一定沒有甲基價值。",
    ]
    add_bullets(slide, limit_bullets, 6.92, 2.54, 4.82, 2.36, font_size=10.5, color=theme["text"], font_name=fonts["zh_body"], space_after=4)
    add_rect(slide, 6.92, 5.16, 4.82, 0.70, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "正確口徑：目前規則在 5kHz 有效，尚未穩定跨平台成立。", 7.10, 5.40, 4.40, 0.18, font_size=9.6, color=theme["light_text"], font_name=fonts["zh_body"], bold=True)


def add_mixed(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg"])
    add_header(slide, "mixed sample 的目前位置：假說，不是定論", "Mixed Sample Hypothesis", "mixed sample 目前只能講到初步觀察與合理推論，正確策略是先把純樣本做完，再回來做校正驗證。", theme, fonts)
    add_rect(slide, 0.82, 1.96, 4.10, 4.74, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "當前觀察", 1.08, 2.18, 1.4, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_rect(slide, 1.12, 2.62, 3.40, 0.96, fill="DCE8F8", line=None, radius=True)
    add_text(slide, "tumor < normal < mixed", 1.28, 2.88, 3.04, 0.18, font_size=15.0, color=theme["text"], font_name=fonts["zh_title"], bold=True, align="center")
    add_text(slide, "PairwiseMedianDist 中位數梯度", 1.28, 3.22, 3.04, 0.16, font_size=10.0, color=theme["muted"], font_name=fonts["zh_body"], align="center")
    add_numbered_list(
        slide,
        [
            "tumor：PairwiseMedianDist 最低。",
            "normal：中位數位於中間。",
            "mixed：目前中位數最高。",
        ],
        1.10,
        4.02,
        3.46,
        1.96,
        font_size=10.0,
        text_color=theme["text"],
        font_name=fonts["zh_body"],
        number_fill=theme["dark"],
        number_color=theme["light_text"],
        item_fill="F5F8FD",
    )

    add_rect(slide, 5.16, 1.96, 3.50, 4.74, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "合理推論", 5.42, 2.18, 1.5, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_numbered_list(
        slide,
        slides_cfg["mixed_sample"]["hypothesis_bullets"],
        5.34,
        2.52,
        2.96,
        2.16,
        font_size=9.9,
        text_color=theme["text"],
        font_name=fonts["zh_body"],
        number_fill=theme["accent"],
        number_color=theme["light_text"],
        item_fill="F5F8FD",
    )
    add_rect(slide, 5.34, 5.02, 2.96, 0.58, fill="DCE8F8", line=None, radius=True)
    add_text(slide, "這條線現在還不是 closeout 結論。", 5.54, 5.22, 2.56, 0.16, font_size=9.4, color=theme["text"], font_name=fonts["zh_body"], bold=True)

    add_rect(slide, 8.90, 1.96, 3.60, 4.74, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "這頁要怎麼講", 9.16, 2.18, 1.7, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_numbered_list(
        slide,
        [
            slides_cfg["mixed_sample"]["caution"],
            slides_cfg["mixed_sample"]["next"],
            "在 deck 中應放在主線後段的一頁 hypothesis，不可拿來當主要證據頁。",
        ],
        9.08,
        2.52,
        3.0,
        2.20,
        font_size=9.9,
        text_color=theme["text"],
        font_name=fonts["zh_body"],
        number_fill=theme["dark"],
        number_color=theme["light_text"],
        item_fill="F5F8FD",
    )


def add_next_steps(slide, slides_cfg, theme, fonts):
    set_bg(slide, theme["dark"])
    add_header(slide, "下一步研究", "Next Steps", "研究主線已清楚：先純樣本、先 annotation，再決定是否有任何規則值得升級成核心流程。", theme, fonts, dark=True)
    add_rect(slide, 0.82, 1.98, 5.70, 4.88, fill="F5F8FD", line=None, radius=True)
    add_text(slide, "下一步要做", 1.08, 2.18, 1.8, 0.18, font_size=14.8, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_numbered_list(
        slide,
        slides_cfg["next_steps"]["actions"],
        1.00,
        2.52,
        5.10,
        3.34,
        font_size=10.0,
        text_color=theme["text"],
        font_name=fonts["zh_body"],
        number_fill=theme["dark"],
        number_color=theme["light_text"],
        item_fill="FFFFFF",
    )
    add_rect(slide, 6.80, 1.98, 5.70, 4.88, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "現在不該做", 7.08, 2.18, 1.8, 0.18, font_size=14.8, color=theme["light_text"], font_name=fonts["zh_title"], bold=True)
    add_numbered_list(
        slide,
        slides_cfg["next_steps"]["not_yet"],
        6.98,
        2.52,
        5.14,
        3.34,
        font_size=10.0,
        text_color=theme["light_text"],
        font_name=fonts["zh_body"],
        number_fill=theme["dark"],
        number_color=theme["light_text"],
        item_fill="507BC3",
    )
    add_rect(slide, 0.82, 6.02, 11.7, 0.42, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "原則：先把 annotation 解決可解釋性與追蹤性，再考慮是否值得改核心 hard keep / hard filter。", 1.08, 6.14, 11.1, 0.14, font_size=10.0, color=theme["text"], font_name=fonts["zh_title"], bold=True)


def add_appendix(slide, meta, runtime, slides_cfg, theme, fonts):
    set_bg(slide, theme["bg_alt"])
    add_header(slide, "附錄與可複查路徑", "Appendix and Traceability", "這份 draft deck 的來源報告、生成腳本、config、profile 與輸出資料夾都能直接回查。", theme, fonts)
    add_rect(slide, 0.82, 1.96, 4.32, 4.86, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 5.34, 1.96, 3.28, 4.86, fill="FFFFFF", line=theme["line"], radius=True)
    add_rect(slide, 8.82, 1.96, 3.50, 4.86, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "核心文件", 1.06, 2.18, 1.4, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    for idx, item in enumerate(slides_cfg["source_paths"][:3]):
        y = 2.52 + idx * 1.14
        add_text(slide, item["label"], 1.06, y, 1.8, 0.16, font_size=11.1, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, Path(item["path"]).name, 1.06, y + 0.24, 3.82, 0.34, font_size=10.2, color=theme["text"], font_name=fonts["zh_body"], bold=True)
    add_text(slide, "設定與資料", 5.58, 2.18, 1.6, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    config_items = slides_cfg["source_paths"][3:]
    for idx, item in enumerate(config_items):
        y = 2.52 + idx * 1.30
        add_text(slide, item["label"], 5.58, y, 2.0, 0.16, font_size=10.8, color=theme["accent"], font_name=fonts["zh_title"], bold=True)
        add_text(slide, Path(item["path"]).name, 5.58, y + 0.24, 2.70, 0.34, font_size=9.8, color=theme["text"], font_name=fonts["zh_body"], bold=True)
    add_text(slide, "版本與重生", 9.06, 2.18, 1.8, 0.18, font_size=13.2, color=theme["text"], font_name=fonts["zh_title"], bold=True)
    add_text(slide, "Deck", 9.06, 2.52, 1.0, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, runtime["output_pptx"].name, 9.06, 2.74, 2.92, 0.42, font_size=8.9, color=theme["muted"], font_name=fonts["en_body"])
    add_text(slide, "Version JSON", 9.06, 3.34, 1.4, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_text(slide, runtime["version_json"].name, 9.06, 3.56, 2.92, 0.42, font_size=8.9, color=theme["muted"], font_name=fonts["en_body"])
    add_text(slide, "Rerun", 9.06, 4.20, 1.0, 0.16, font_size=10.0, color=theme["accent"], font_name=fonts["en_body"], bold=True)
    add_rect(slide, 9.06, 4.44, 2.88, 1.26, fill=theme["dark"], line=None, radius=True)
    rerun_short = (
        "python build_ai_f1_research_pptx.py\n"
        "--config 20260312_ai_f1_research_pptx_config_01.json\n"
        f"--output-pptx {runtime['output_pptx'].name}"
    )
    add_text(slide, rerun_short, 9.22, 4.60, 2.50, 0.96, font_size=8.0, color=theme["light_text"], font_name=fonts["en_body"])
    add_text(slide, "完整絕對路徑與完整指令已保留在 notes / version.json。", 9.06, 5.92, 2.92, 0.18, font_size=8.6, color=theme["muted"], font_name=fonts["zh_body"])
    add_text(slide, f"狀態：{meta['status']} · 版本：{meta['version']} · notes：embedded", 9.06, 6.16, 2.92, 0.16, font_size=8.8, color=theme["text"], font_name=fonts["zh_body"])


def build_notes(meta: dict, slides_cfg: dict, runtime: dict) -> list[str]:
    source_report = meta["source_report"]
    main_weekly = slides_cfg["source_paths"][1]["path"]
    pool_b_path = slides_cfg["source_paths"][2]["path"]
    rescue_path = slides_cfg["source_paths"][3]["path"]
    profile_path = runtime["profile"].as_posix()
    config_path = runtime["config"].as_posix()
    deck_path = runtime["output_pptx"].as_posix()
    notes = [
        f"""這頁目的：
建立簡報主題、時間範圍、輸出定位與高層結論。
口頭講解重點：
1. 這份 deck 的主題是用 AI / Agent 自動化推進 InterSubMod 的 F1 驗證。
2. 這是一份 16 頁技術深度版 draft，speaker notes 為正式產物。
3. 今天不是要推一條新規則，而是要把主線、限制與下一步整理清楚。
完整中文推論：
三層結論已經收斂，因此本 deck 之後的數據只是逐頁驗證 caller-first、methylation-support 與 haplotype 共分析的正確位置。
數據或圖表證據：
- 主整合報告封面與第 11-12 節頁序藍圖
原始檔案或路徑：
- {source_report}
- {config_path}
- {profile_path}
補充背景：
本 deck 採個人固定 PPT 規範：中文主體、英文輔助、一頁一重點、完整 notes。""",
        f"""這頁目的：
先交代簡報結構，讓聽眾知道會依什麼順序理解主線。
口頭講解重點：
1. 前段先講研究問題與主結論。
2. 中段講 paired、TO、Pool B 與特徵分工。
3. 後段講 comparability、mixed 與下一步。
完整中文推論：
這份 deck 並不是把所有實驗流水帳照搬，而是依『研究決策順序』安排章節，避免不同母體、不同問題定義混寫。
數據或圖表證據：
- Agenda 內容來自 20260312 主整合報告的 16 頁規劃
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
報告總覽刻意把 Pool B 和 mixed 留在後段，因為兩者都不應壓過純樣本主線。""",
        f"""這頁目的：
用一頁先講完整份 deck 的破題結論。
口頭講解重點：
1. caller-first 仍是最穩定的主規則。
2. 甲基特徵有價值，但主要位於 support、annotation 與 triage。
3. 5kHz 有效不代表 DORADO 同樣有效。
完整中文推論：
這個三層架構是目前最合理的工作分工，後面各頁都只是把這個高層結論拆成 paired、TO、Pool B 與 mixed 四個面向去驗證。
數據或圖表證據：
- framework.conclusions
- framework.layers
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
這一頁要直接講結論，不要等到中段才讓聽眾知道口徑。""",
        f"""這頁目的：
說明為什麼最近研究方式從手動試規則轉成 AI / Agent 自動化。
口頭講解重點：
1. 自動讀 docs / knowledge。
2. 自動跑 feature sweep 與 candidate-specific rescue。
3. 自動生成報告並回寫索引。
完整中文推論：
研究價值不只在節省時間，而是在於讓每輪實驗都留有可追溯的證據鏈，避免 AI 每次重開就重新踩過失敗方向。
數據或圖表證據：
- ai_automation.before_after
- ai_automation.value
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
這個流程依賴 CURRENT_FOCUS、INDEX 與 docs/README 的三層入口。""",
        f"""這頁目的：
把本輪 AI 被限制只做的兩類問題說清楚。
口頭講解重點：
1. 問題 A 是後段刪 FP。
2. 問題 B 是前段救 TP。
3. 兩者母體不同，不能用同一套規則混寫。
完整中文推論：
像 AlleleDelta 在 kept-set / artifact triage 與 Pool B FN 這兩種母體中方向可以相反，因此若不先分母體，所有結論都會失真。
數據或圖表證據：
- two_problems.problem_a
- two_problems.problem_b
原始檔案或路徑：
- {source_report}
- {pool_b_path}
補充背景：
這頁是整份 deck 的方法學防呆頁。""",
        f"""這頁目的：
把從 docs 到報告的整條自動化研究鏈講清楚。
口頭講解重點：
1. 前段先對齊文件與知識庫。
2. 中段串 caller、LongPhase、tagged BAM 與 InterSubMod。
3. 後段做 feature sweep、候選 rescue、報告與索引回寫。
完整中文推論：
tagged BAM 是 InterSubMod 真正能做 HP / label-first 驗證的必要條件，因此 paired 與 TO 的 LongPhase workflow 都不是可省略細節。
數據或圖表證據：
- workflow_steps
原始檔案或路徑：
- {main_weekly}
- {rescue_path}
補充背景：
空間策略之所以重要，就是因為 full tagged BAM 體積大，會決定研究順序。""",
        f"""這頁目的：
解釋為什麼先做 HCC1395 5kHz，再用 HCC1395 DORADO 做交叉驗證。
口頭講解重點：
1. 5kHz 比較 noisy，適合做規則壓力測試。
2. DORADO 比較乾淨，適合測規則是否只是平台特異性。
3. 其他樣本先暫列後續，不搶主線。
完整中文推論：
這個樣本策略可以同時提供『最容易看見增益的場景』與『最容易看見規則不可攜的場景』，所以是現在最合理的 paired / TO 雙軸設計。
數據或圖表證據：
- sample_strategy.primary
- sample_strategy.secondary
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
big8_disk 目前只保留最有研究價值的主線素材。""",
        f"""這頁目的：
呈現 paired 場景的正式 benchmark 結果與正確口徑。
口頭講解重點：
1. 5kHz paired 有小幅正增益，但主要來自少掉 FP。
2. DORADO paired 幾乎持平略負。
3. paired 目前仍是 caller 與 LongPhase 做掉大部分改善。
完整中文推論：
paired 的正確說法不是『甲基已經能主導救援』，而是『在最 noisy 的主樣本上，InterSubMod 仍能做少量 precision 精修』。
數據或圖表證據：
- paired_results.table
- paired_results.chart_series
原始檔案或路徑：
- {main_weekly}
- {source_report}
補充背景：
paired 與 TO 的結論必須分開講，因為 LongPhase 在 TO pilot 上沒有改變 baseline PASS call set。""",
        f"""這頁目的：
呈現 tumor-only candidate-specific rescue 的主結果。
口頭講解重點：
1. TO 下甲基特徵不是全負效果。
2. 但固定同一 caller gate 後，甲基 support 仍未超過最佳 caller-only。
3. 5kHz 與 DORADO 支持同一個高層結論，但 Pairwise 方向不同。
完整中文推論：
這頁要建立的是『caller-first + methylation-support』的分工，而不是把 support 誤說成已經取代 caller gate。
數據或圖表證據：
- to_results.5khz.rules
- to_results.dorado.rules
原始檔案或路徑：
- {source_report}
- {rescue_path}
補充背景：
DORADO TO 目前的 read-level 素材來自 candidate-window subset haplotag。""",
        f"""這頁目的：
正式關閉 Pool B FN 支線，並把它接回主線敘事。
口頭講解重點：
1. Pool B FN 再次證明 caller-side rescue 空間存在。
2. 最佳結果仍是 caller-only。
3. 甲基只提供弱補強，沒有超過 caller-only 最佳。
完整中文推論：
Pool B 是前段 TP rescue 的問題，因此不能和 kept-set / artifact triage 的規則混成同一條全域規則，特別是 AlleleDelta 的方向會相反。
數據或圖表證據：
- pool_b.fn_count
- pool_b.results
- pool_b.conclusions
原始檔案或路徑：
- {pool_b_path}
- {source_report}
補充背景：
這頁的任務是切清楚『前段救 TP』與『後段刪 FP』。""",
        f"""這頁目的：
整理現在最重要的特徵與它們正確的角色。
口頭講解重點：
1. GQ 已是主規則。
2. Quality_Score 可升 annotation。
3. PairwiseMedianDist 只能升 dataset-aware annotation。
4. artifact triage 不能回頭拿來當 TP rescue 主規則。
完整中文推論：
真正要做的不是再堆更多條件，而是把每個特徵放在正確層級，避免把穩定訊號與不穩定訊號混成同一個決策引擎。
數據或圖表證據：
- feature_table
原始檔案或路徑：
- {source_report}
- {rescue_path}
補充背景：
這張表同時也是後續 annotation export 的需求清單。""",
        f"""這頁目的：
用三層架構把 annotation 與主規則的分工制度化。
口頭講解重點：
1. 第一層是 caller-first。
2. 第二層是 methylation-support。
3. 第三層是 annotation / QC / diagnostics。
完整中文推論：
這個分工的實際用途，是讓我們可以先把 Quality、Pairwise、hp_assign 等特徵接到輸出層，而不是太早把它們硬接進 hard keep / hard filter。
數據或圖表證據：
- framework.layers
- annotation decision cards
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
這頁是之後 round summary、dashboard 與 diagnostics 最重要的設計依據。""",
        f"""這頁目的：
明確交代四象限比較的限制與 dataset-dependence。
口頭講解重點：
1. 5kHz 與 DORADO 的方向差異合理，但不能直接當因果。
2. candidate-specific coverage 與 snapshot scope 不同。
3. paired 的 negative result 仍要保留 coverage ceiling 的 caveat。
完整中文推論：
這頁要阻止過度推論，尤其是把 PairwiseMedianDist 當成跨樣本穩定生物學訊號，或把 paired 的目前無效說成 paired 永遠無效。
數據或圖表證據：
- divergence.bullets
- divergence.caution
原始檔案或路徑：
- {rescue_path}
- {source_report}
補充背景：
這頁等同於簡報版的 methods / limitations。""",
        f"""這頁目的：
把 mixed sample 目前應該被放在哪個位置說清楚。
口頭講解重點：
1. 目前只有 tumor < normal < mixed 的初步觀察。
2. 合理推論是組織來源差異可能大於 haplotype 與癌症趨同的效應。
3. 這條線現在還不是正式 closeout。
完整中文推論：
mixed sample 應放在主線後段的一頁 hypothesis，提醒大家它值得研究，但不能壓過純樣本主線，也不能先講成已驗證結論。
數據或圖表證據：
- mixed_sample.observation
- mixed_sample.hypothesis_bullets
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
策略固定為先完成純樣本，再回頭做 mixed 校正與去噪。""",
        f"""這頁目的：
說明接下來應該做什麼，以及哪些事現在不該做。
口頭講解重點：
1. 先把 annotation 接到輸出層。
2. 繼續累積 cross-dataset 證據。
3. 不要把 5kHz 單一方向直接升級成全域規則。
完整中文推論：
這頁真正的作用是防止研究主線再次回到『不停試規則但缺乏分層與治理』的狀態。
數據或圖表證據：
- next_steps.actions
- next_steps.not_yet
原始檔案或路徑：
- {source_report}
- {main_weekly}
補充背景：
如果後續真的要升級任何規則，也應先從 annotation export 開始，而不是直接改核心流程。""",
        f"""這頁目的：
保留複查、重生與版本追蹤入口。
口頭講解重點：
1. 主報告、週報、Pool B closeout、四象限整合與 config/profile 都可回查。
2. 這份 deck 的生成腳本與 version.json 已固定。
3. 後續再版應沿用這條生成鏈。
完整中文推論：
speaker notes 與 version metadata 是正式產物的一部分，目的是讓後續簡報能被重生、被 QA、被追版本。
數據或圖表證據：
- source_paths
- deck path
原始檔案或路徑：
- {deck_path}
- {config_path}
- {profile_path}
補充背景：
validated 之前，這份 deck 仍屬 draft，但複查鏈已完整。""",
    ]
    return notes


def write_version_json(meta: dict, runtime: dict, qa_status: str, qa_cycles: int, pdf_path: Path | None) -> None:
    payload = {
        "deck_name": runtime["output_pptx"].stem,
        "version": meta["version"],
        "status": meta["status"],
        "generated_at": datetime.now(TZ).isoformat(timespec="seconds"),
        "qa_status": qa_status,
        "qa_cycles": qa_cycles,
        "canonical_deck_path": runtime["output_pptx"].as_posix(),
        "presentation_entry_path": runtime["output_pptx"].parent.as_posix(),
        "pdf_dir": runtime["pdf_dir"].as_posix(),
        "pdf_path": pdf_path.as_posix() if pdf_path else None,
        "source_report": meta["source_report"],
        "generation_script": runtime["script_path"].as_posix(),
        "generation_config": runtime["config"].as_posix(),
        "personal_ppt_profile": runtime["profile"].as_posix(),
        "assets_dir": runtime["assets_dir"].as_posix(),
        "speaker_notes": "embedded",
        "notes": [
            "16 頁研究內部技術深度版 draft deck。",
            "依 20260312 AI 自動化 F1 主線整合報告與個人 PPT 規範生成。",
            "QA metadata 由腳本參數寫入；若已完成 QA，請同步更新 qa_status / qa_cycles / pdf_path。",
        ],
    }
    runtime["version_json"].write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", required=True, type=Path)
    parser.add_argument("--profile", required=True, type=Path)
    parser.add_argument("--output-pptx", required=True, type=Path)
    parser.add_argument("--status", default="draft")
    parser.add_argument("--version", default="v01")
    parser.add_argument("--qa-status", default="generated")
    parser.add_argument("--qa-cycles", default=0, type=int)
    parser.add_argument("--pdf-path", type=Path)
    args = parser.parse_args()

    cfg = load_json(args.config)
    profile = load_json(args.profile)
    meta = dict(cfg["meta"])
    meta["status"] = args.status
    meta["version"] = args.version

    output_pptx = args.output_pptx.resolve()
    deck_dir = output_pptx.parent
    assets_dir = deck_dir / "assets" / f"{output_pptx.stem}_assets"
    pdf_dir = deck_dir / "pdf"
    version_json = output_pptx.with_suffix(".version.json")
    rerun_command = (
        "python scripts/analysis/build_ai_f1_research_pptx.py "
        f"--config {args.config} --profile {args.profile} --output-pptx {args.output_pptx}"
    )
    runtime = {
        "config": args.config.resolve(),
        "profile": args.profile.resolve(),
        "output_pptx": output_pptx,
        "assets_dir": assets_dir,
        "pdf_dir": pdf_dir,
        "version_json": version_json,
        "script_path": Path(__file__).resolve(),
        "rerun_command": rerun_command,
    }

    theme = cfg["theme"]
    fonts = {
        "zh_title": profile["fonts"]["zh_title"],
        "zh_body": profile["fonts"]["zh_body"],
        "en_title": profile["fonts"]["en_title"],
        "en_body": profile["fonts"]["en_body"],
    }
    slides_cfg = cfg["slides"]
    assets = build_assets(slides_cfg, assets_dir)
    notes = build_notes(meta, slides_cfg, runtime)

    prs = create_prs(meta["deck_title"], meta["author"])
    builders = [
        lambda s: add_cover(s, meta, slides_cfg, theme, fonts),
        lambda s: add_agenda(s, slides_cfg, theme, fonts),
        lambda s: add_thesis(s, slides_cfg, theme, fonts),
        lambda s: add_ai_automation(s, slides_cfg, theme, fonts),
        lambda s: add_two_problems(s, slides_cfg, theme, fonts),
        lambda s: add_workflow(s, slides_cfg, theme, fonts),
        lambda s: add_sample_strategy(s, slides_cfg, theme, fonts),
        lambda s: add_paired_results(s, slides_cfg, assets, theme, fonts),
        lambda s: add_to_results(s, slides_cfg, assets, theme, fonts),
        lambda s: add_pool_b(s, slides_cfg, theme, fonts),
        lambda s: add_feature_table(s, slides_cfg, theme, fonts),
        lambda s: add_annotation_architecture(s, slides_cfg, theme, fonts),
        lambda s: add_comparability(s, slides_cfg, theme, fonts),
        lambda s: add_mixed(s, slides_cfg, theme, fonts),
        lambda s: add_next_steps(s, slides_cfg, theme, fonts),
        lambda s: add_appendix(s, meta, runtime, slides_cfg, theme, fonts),
    ]

    total_pages = len(builders)
    dark_pages = {1, 3, 15}
    for idx, builder in enumerate(builders, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder(slide)
        add_speaker_notes(slide, notes[idx - 1])
        add_footer(slide, theme, fonts, idx, total_pages, meta["footer"], dark=idx in dark_pages)

    ensure_dir(deck_dir)
    ensure_dir(pdf_dir)
    prs.save(output_pptx)
    pdf_path = args.pdf_path.resolve() if args.pdf_path else None
    write_version_json(meta, runtime, args.qa_status, args.qa_cycles, pdf_path)
    print(f"saved_pptx={output_pptx}")
    print(f"assets_dir={assets_dir}")
    print(f"version_json={version_json}")


if __name__ == "__main__":
    main()
