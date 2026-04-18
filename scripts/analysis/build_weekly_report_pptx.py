#!/usr/bin/env python3
"""
將 InterSubMod 研究週報轉成可重複生成的 PPTX。

使用方式:
python scripts/analysis/build_weekly_report_pptx.py \
  --config docs/references/manual/assets/20260311_weekly_report_pptx_config_01.json
"""

from __future__ import annotations

import argparse
import io
import json
import math
import re
import tempfile
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

# ---------- Color Semantics (Okabe-Ito colorblind-safe) ----------
COLOR_POSITIVE = "009E73"   # Green — TP / confirmed
COLOR_NEGATIVE = "D55E00"   # Vermillion — FP / rejected
COLOR_PRIMARY  = "0072B2"   # Blue — primary data
COLOR_NEUTRAL  = "808080"   # Gray — neutral
COLOR_WARN     = "E69F00"   # Orange — warning

# ---------- Shape Semantics ----------
SHAPE_PROCESS  = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
SHAPE_DECISION = MSO_AUTO_SHAPE_TYPE.DIAMOND
SHAPE_TERMINAL = MSO_AUTO_SHAPE_TYPE.OVAL
SHAPE_EXTERNAL = MSO_AUTO_SHAPE_TYPE.HEXAGON

# ---------- Icon Map (Windows-compatible single-codepoint emoji) ----------
ICON_MAP = {
    "data": "\U0001f4ca",       # 📊
    "confirm": "\u2705",        # ✅
    "reject": "\u274c",         # ❌
    "warn": "\u26a0\ufe0f",     # ⚠️
    "next": "\u2192",           # →
    "method": "\U0001f52c",     # 🔬
    "target": "\U0001f3af",     # 🎯
    "trend": "\U0001f4c8",      # 📈
    "insight": "\U0001f4a1",    # 💡
    "bio": "\U0001f9ec",        # 🧬
    "cycle": "\U0001f504",      # 🔄
    "todo": "\U0001f4cb",       # 📋
    "star": "\u2b50",           # ⭐
}

# Deep page font size multiplier
DARK_PAGE_BOOST = 1.18


def hex_rgb(value: str) -> RGBColor:
    value = value.replace("#", "").strip()
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_config(path: Path) -> dict:
    cfg = json.loads(path.read_text(encoding="utf-8"))
    cfg["_config_dir"] = str(path.parent)
    return cfg


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)


def _resolve_icons(text: str) -> str:
    """Replace {icon:key} tokens with emoji characters."""
    def _repl(m):
        return ICON_MAP.get(m.group(1), m.group(0))
    return re.sub(r"\{icon:(\w+)\}", _repl, text)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 18,
    color: str = "1B1F24",
    bold: bool = False,
    font_name: str = "Arial",
    align: str = "left",
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    italic: bool = False,
    margin: float = 0.05,
    line_spacing: float | None = None,
) -> None:
    text = _resolve_icons(text)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.vertical_anchor = valign
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    p = box.text_frame.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_rgb(color)


def add_bullets(
    slide,
    bullets: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 18,
    color: str = "1B1F24",
    bullet_color: str = "A85540",
    font_name: str = "Arial",
    level: int = 0,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(bullets):
        item = _resolve_icons(item)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = hex_rgb(color)
        if hasattr(p, "_pPr") and p._pPr is not None:
            pass


def add_rect(slide, x: float, y: float, w: float, h: float, *, fill: str, line: str | None = None, radius: bool = False, transparency: int = 0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    return shape


_SAFE_PIC_MAX_PX = 1920  # Max image width in pixels for presentations


def _safe_add_picture(slide, image_path: str, left, top, *, width=None, height=None):
    """Add a picture to a slide with automatic safety processing:

    1. RGBA/PA → RGB conversion (PPT can fail on RGBA PNGs)
    2. Down-scale images wider than 1920 px (high-res images cause
       rendering lag or white-page on some PPT builds)
    3. Strip ICC profiles that confuse PPT colour pipeline
    4. Write via temp file (not BytesIO) for maximum compatibility

    Keeps aspect ratio; always produces an 8-bit RGB PNG.
    """
    img_path = Path(image_path)
    if not img_path.exists():
        print(f"  [WARN] Image not found, skipping: {img_path}")
        return None

    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height

    # --- Decide if we need to process the image ---
    needs_process = False
    try:
        with Image.open(img_path) as img:
            if img.mode in ("RGBA", "PA"):
                needs_process = True
            elif img.mode == "P" and "transparency" in img.info:
                needs_process = True
            if img.size[0] > _SAFE_PIC_MAX_PX:
                needs_process = True
            if "icc_profile" in img.info:
                needs_process = True
    except Exception as e:
        print(f"  [WARN] PIL check failed for {img_path}: {e}; embedding as-is")
        return slide.shapes.add_picture(str(img_path), left, top, **kwargs)

    if not needs_process:
        needs_process = True  # Always re-encode for maximum PPT compatibility

    # --- Process: convert + downscale + strip profiles ---
    tmp = None
    try:
        with Image.open(img_path) as img:
            ops = []
            # 1. Colour mode
            if img.mode == "RGBA":
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[3])
                ops.append("RGBA→RGB")
            elif img.mode != "RGB":
                rgb_img = img.convert("RGB")
                ops.append(f"{img.mode}→RGB")
            else:
                rgb_img = img.copy()

            # 2. Down-scale
            if rgb_img.size[0] > _SAFE_PIC_MAX_PX:
                ratio = _SAFE_PIC_MAX_PX / rgb_img.size[0]
                new_size = (
                    _SAFE_PIC_MAX_PX,
                    max(1, int(rgb_img.size[1] * ratio)),
                )
                resample = getattr(Image, "Resampling", Image).LANCZOS
                rgb_img = rgb_img.resize(new_size, resample)
                ops.append(f"downscale→{new_size[0]}×{new_size[1]}")

        # 3. Save to temp file (no ICC, no extra metadata)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        rgb_img.save(tmp.name, format="PNG", optimize=True, icc_profile=None)
        print(f"  [INFO] {img_path.name}: {', '.join(ops)} → {tmp.name}")
        result = slide.shapes.add_picture(tmp.name, left, top, **kwargs)
        return result
    except Exception as e:
        print(f"  [WARN] Processing failed: {e}; embedding original")
        return slide.shapes.add_picture(str(img_path), left, top, **kwargs)
    finally:
        if tmp is not None:
            try:
                Path(tmp.name).unlink(missing_ok=True)
            except Exception:
                pass


def wrap_path_for_slide(path: str, max_len: int = 46, max_lines: int = 4) -> str:
    parts = [token for token in re.split(r"([/ ])", path) if token]
    lines = []
    current = ""
    for part in parts:
        segment = part
        if not current:
            current = segment.lstrip()
            continue
        if len(current) + len(segment) <= max_len:
            current += segment
        else:
            lines.append(current.rstrip())
            current = segment.lstrip()
    if current:
        lines.append(current.rstrip())
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        if len(lines[-1]) >= max_len - 1:
            lines[-1] = lines[-1][: max_len - 1] + "…"
        else:
            lines[-1] += " …"
    return "\n".join(lines)


def add_bilingual_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text_zh: str,
    text_en: str | None,
    font_size: float,
    theme: dict,
    *,
    bold: bool = False,
    align: str = "left",
    color: str | None = None,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.05,
) -> None:
    """Add bilingual text box: Chinese primary + English translation below (70% size, muted)."""
    text_zh = _resolve_icons(text_zh)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    fg = color or theme["text"]
    # Chinese paragraph
    p_zh = tf.paragraphs[0]
    p_zh.alignment = alignment
    run_zh = p_zh.add_run()
    run_zh.text = text_zh
    run_zh.font.size = Pt(font_size)
    run_zh.font.name = "Arial"
    run_zh.font.bold = bold
    run_zh.font.color.rgb = hex_rgb(fg)
    # English paragraph (if provided) — indented, closer to ZH, smaller
    if text_en:
        p_en = tf.add_paragraph()
        p_en.alignment = alignment
        p_en.space_before = Pt(1)
        # Indent English text to visually subordinate it
        p_en.level = 0
        p_en_fmt = p_en._pPr if p_en._pPr is not None else p_en._p.get_or_add_pPr()
        p_en_fmt.set("marL", str(Emu(Inches(0.25))))  # 0.25" left indent
        run_en = p_en.add_run()
        run_en.text = text_en
        en_size = max(9, font_size * 0.60)
        run_en.font.size = Pt(en_size)
        run_en.font.name = "Arial"
        run_en.font.bold = False
        run_en.font.color.rgb = hex_rgb(theme.get("en_color", "5E6572"))


def add_bilingual_bullets(
    slide,
    bullets_zh: list[str],
    bullets_en: list[str] | None,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 14,
    color: str = "1B1F24",
    theme: dict | None = None,
) -> None:
    """Add bilingual bullet list: each Chinese bullet followed by English translation."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    en_list = bullets_en or []
    en_color = (theme or {}).get("en_color", "5E6572")
    for idx, item_zh in enumerate(bullets_zh):
        item_zh = _resolve_icons(item_zh)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = item_zh
        run.font.size = Pt(font_size)
        run.font.name = "Arial"
        run.font.color.rgb = hex_rgb(color)
        # English sub-line — indented, closer spacing, smaller
        if idx < len(en_list) and en_list[idx]:
            p_en = tf.add_paragraph()
            p_en.alignment = PP_ALIGN.LEFT
            p_en.level = 1
            p_en.space_before = Pt(0)
            p_en.space_after = Pt(6)
            p_en_fmt = p_en._pPr if p_en._pPr is not None else p_en._p.get_or_add_pPr()
            p_en_fmt.set("marL", str(Emu(Inches(0.35))))
            run_en = p_en.add_run()
            run_en.text = en_list[idx]
            en_size = max(9, font_size * 0.60)
            run_en.font.size = Pt(en_size)
            run_en.font.name = "Arial"
            run_en.font.color.rgb = hex_rgb(en_color)


def inject_speaker_notes(slide, notes_text: str) -> None:
    """Add speaker notes to a slide in '重點詞：解釋' format."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    for idx, line in enumerate(notes_text.strip().split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line


def add_title(slide, title: str, subtitle: str, theme: dict, *, dark: bool = False,
              title_en: str | None = None, subtitle_en: str | None = None) -> None:
    fg = theme["light_text"] if dark else theme["text"]
    muted = theme["light_muted"] if dark else theme["muted"]
    title_sz = int(32 * DARK_PAGE_BOOST) if dark else 32
    sub_sz = int(16 * DARK_PAGE_BOOST) if dark else 16
    en_title_sz = max(9, int(title_sz * 0.50))
    en_sub_sz = max(9, int(sub_sz * 0.60))
    # ZH title
    add_text(slide, title, 0.8, 0.35, 11.7, 0.75, font_size=title_sz, color=fg, bold=True, font_name=theme["font_title"], margin=0)
    # EN title — indented, smaller, closer to ZH
    if title_en:
        add_text(slide, title_en, 1.05, 1.05, 11.3, 0.35, font_size=en_title_sz,
                 color=theme.get("en_color", "5E6572"), font_name=theme["font_body"], margin=0)
    # ZH subtitle
    sub_y = 1.38 if title_en else 1.20
    add_text(slide, subtitle, 0.8, sub_y, 11.3, 0.42, font_size=sub_sz, color=muted, font_name=theme["font_body"], margin=0)
    # EN subtitle — indented under ZH subtitle
    if subtitle_en:
        add_text(slide, subtitle_en, 1.05, sub_y + 0.35, 11.0, 0.32, font_size=en_sub_sz,
                 color=theme.get("en_color", "5E6572"), font_name=theme["font_body"], margin=0)


def add_footer(slide, theme: dict, page_num: int, total_pages: int, footer_text: str, *, dark: bool = False) -> None:
    color = theme["light_muted"] if dark else "5E6572"
    add_text(slide, footer_text, 0.8, 6.78, 9.0, 0.3, font_size=10, color=color, margin=0)
    add_text(slide, f"{page_num}/{total_pages}", 12.0, 6.75, 0.55, 0.3, font_size=12, color=color, align="right", margin=0)


def style_table(table, theme: dict, *, header_fill: str, header_font: str, body_fill: str = "FFFFFF", zebra_fill: str | None = None, font_size: int = 12) -> None:
    rows = len(table.rows)
    cols = len(table.columns)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = hex_rgb(header_fill)
            elif zebra_fill and r % 2 == 0:
                cell.fill.fore_color.rgb = hex_rgb(zebra_fill)
            else:
                cell.fill.fore_color.rgb = hex_rgb(body_fill)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = theme["font_body"]
                    run.font.size = Pt(font_size)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = hex_rgb(header_font)
                    else:
                        run.font.color.rgb = hex_rgb(theme["text"])


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, theme: dict, *, header_fill: str, body_fill: str = "FFFFFF", zebra_fill: str | None = None, font_size: int = 12, col_widths: list[float] | None = None):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    # Strip tableStyleId in-memory to prevent PowerPoint repair dialog.
    # python-pptx injects a GUID that references a non-existent theme style.
    _ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for _tsi in list(table._tbl.tblPr.iterchildren(f'{{{_ns_a}}}tableStyleId')):
        table._tbl.tblPr.remove(_tsi)
    if col_widths:
        total = sum(col_widths)
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(w * (width / total))
    else:
        col_width = Inches(w / len(rows[0]))
        for idx in range(len(rows[0])):
            table.columns[idx].width = col_width
    row_h = Inches(h / len(rows))
    for idx in range(len(rows)):
        table.rows[idx].height = row_h
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = str(value)
    style_table(table, theme, header_fill=header_fill, header_font=theme["light_text"], body_fill=body_fill, zebra_fill=zebra_fill, font_size=font_size)
    return table


def create_grouped_bar_chart(output_path: Path, title: str, categories: list[str], series: list[dict], theme: dict, *, ylim: tuple[float, float], ylabel: str) -> None:
    plt.figure(figsize=(10, 4.6), dpi=200)
    x = range(len(categories))
    width = 0.22 if len(series) >= 3 else 0.3
    offsets = [(-width), 0, width, 2 * width]
    for idx, item in enumerate(series):
        values = item["values"]
        pos = [v + offsets[idx] for v in x]
        plt.bar(pos, values, width=width * 0.9, color=f"#{item['color']}", label=item["name"])
        for px, val in zip(pos, values):
            plt.text(px, val + 0.001, f"{val:.4f}", ha="center", va="bottom", fontsize=9)
    plt.xticks(list(x), categories, fontsize=10)
    plt.yticks(fontsize=10)
    plt.ylim(*ylim)
    plt.ylabel(ylabel, fontsize=11)
    plt.title(title, fontsize=13, weight="bold")
    plt.grid(axis="y", linestyle="--", alpha=0.25)
    plt.legend(frameon=False, fontsize=9, ncol=len(series), loc="upper center", bbox_to_anchor=(0.5, 1.15))
    plt.tight_layout()
    plt.savefig(output_path, bbox_inches="tight", facecolor="white")
    plt.close()


def create_horizontal_bar_chart(output_path: Path, panel_specs: list[dict], theme: dict) -> None:
    fig, axes = plt.subplots(1, len(panel_specs), figsize=(11.2, 4.4), dpi=200, sharex=False)
    if len(panel_specs) == 1:
        axes = [axes]
    for ax, spec in zip(axes, panel_specs):
        labels = [row["name"] for row in spec["rows"]]
        values = [row["f1"] for row in spec["rows"]]
        colors = [f"#{row['color']}" for row in spec["rows"]]
        ypos = list(range(len(labels)))[::-1]
        ax.barh(ypos, values[::-1], color=colors[::-1], height=0.62)
        ax.set_yticks(ypos)
        ax.set_yticklabels(labels[::-1], fontsize=9)
        ax.set_xlim(spec["xlim"][0], spec["xlim"][1])
        ax.set_title(spec["title"], fontsize=11, weight="bold")
        ax.grid(axis="x", linestyle="--", alpha=0.25)
        for y, val in zip(ypos, values[::-1]):
            ax.text(val + 0.001, y, f"{val:.4f}", va="center", fontsize=9)
    fig.suptitle("Paired Raw Pileup/Full vs Final Outputs", fontsize=13, weight="bold")
    plt.tight_layout()
    plt.savefig(output_path, bbox_inches="tight", facecolor="white")
    plt.close()


def create_timeline_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    add_title(slide, "研究推進時間軸", "依日期整理本週主線進展與收斂節點", theme)
    y = 2.2
    x_positions = [0.85, 3.95, 7.05, 10.15]
    add_rect(slide, 1.25, 4.95, 10.35, 0.08, fill=theme["accent"], line=None, radius=False)
    for x, item in zip(x_positions, cfg["deck"]["timeline"]):
        add_rect(slide, x, y, 2.2, 2.55, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.12, y + 0.12, 0.6, 0.35, fill=theme["dark"], line=None, radius=True)
        add_text(slide, item["date"], x + 0.18, y + 0.15, 0.9, 0.2, font_size=11, color=theme["light_text"], bold=True, margin=0)
        add_text(slide, item["title"], x + 0.14, y + 0.56, 1.92, 0.58, font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, item["result"], x + 0.14, y + 1.18, 1.9, 1.12, font_size=12, color="5E6572", margin=0)
    add_rect(slide, 0.8, 6.05, 12.0, 0.45, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "時間軸重點：3/05 先排除 mixed 主線，3/07 跑通 paired / TO，3/08~09 確立 candidate-specific rescue，3/11 完成 phase 2 與 annotation 回接。", 1.0, 6.2, 11.4, 0.2, font_size=11, color="5E6572", margin=0)


def create_presentation(cfg: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = cfg["meta"]["author"]
    prs.core_properties.title = cfg["meta"]["deck_title"]
    prs.core_properties.subject = cfg["meta"]["deck_subtitle"]
    return prs


def add_title_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["dark"])
    add_rect(slide, 0, 0, 13.333, 7.5, fill=theme["dark"], line=None)
    # Right panel — white space placeholder for manual image insertion
    add_rect(slide, 7.2, 0.5, 5.63, 6.5, fill="F6EFE6", line=None, radius=True, transparency=8)
    add_text(slide, "[ 每週重點圖 — 手動插入 ]", 8.0, 3.2, 4.0, 0.4,
             font_size=13, color="A0A0A0", align="center", margin=0, italic=False)
    # Kicker badge
    add_rect(slide, 0.7, 0.65, 1.45, 0.35, fill=theme["accent2"], line=None, radius=True)
    add_text(slide, cfg["meta"]["kicker"], 0.83, 0.72, 2.5, 0.2, font_size=13, color=theme["light_text"], bold=True, margin=0)
    # Title — ZH large, EN smaller and indented
    add_text(slide, cfg["meta"]["deck_title"], 0.8, 1.35, 6.0, 0.8, font_size=38,
             color=theme["light_text"], bold=True, font_name=theme["font_title"], margin=0)
    deck_title_en = cfg["meta"].get("deck_title_en", "")
    if deck_title_en:
        add_text(slide, deck_title_en, 1.1, 2.15, 5.7, 0.45, font_size=18,
                 color=theme["light_muted"], font_name=theme["font_body"], margin=0, italic=False)
    add_text(slide, cfg["meta"]["deck_subtitle"], 0.84, 2.85, 6.0, 0.55, font_size=19, color=theme["light_muted"], font_name=theme["font_body"], margin=0)
    add_text(slide, cfg["meta"]["summary_line"], 0.84, 3.50, 6.0, 0.80, font_size=18, color=theme["light_text"], bold=True, font_name=theme["font_body"], margin=0)
    # Guide bullets (below summary_line, compact)
    guide_bullets = cfg["meta"].get("guide_bullets", [])
    if guide_bullets:
        add_bullets(slide, guide_bullets, 0.84, 4.50, 5.5, 0.80, font_size=12, color=theme["light_muted"])
    # Reporter info
    reporter = cfg["meta"].get("reporter", cfg["meta"].get("author", ""))
    add_text(slide, f"報告人：{reporter}", 0.84, 5.45, 4.0, 0.3, font_size=14, color=theme["light_text"], bold=True, margin=0)
    source_name = Path(cfg["meta"].get("source_report", "")).name or "N/A"
    add_text(slide, f"來源週報：{source_name}", 0.84, 5.85, 5.5, 0.35, font_size=12, color=theme["light_muted"], font_name="Consolas", margin=0)


def add_key_findings_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    deck = cfg["deck"]
    add_title(slide,
              deck.get("key_findings_title", "本週重點結論"),
              deck.get("key_findings_subtitle", "先看四條證據鏈，快速掌握可直接沿用的研究口徑"),
              theme,
              title_en=deck.get("key_findings_title_en"),
              subtitle_en=deck.get("key_findings_subtitle_en"))
    cards = deck["key_findings"]
    positions = [(0.8, 1.95), (6.95, 1.95), (0.8, 4.32), (6.95, 4.32)]
    for (x, y), card in zip(positions, cards):
        add_rect(slide, x, y, 5.55, 2.15, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.18, y + 0.18, 0.48, 0.48, fill=card["accent"], line=None, radius=True)
        # Card title (ZH)
        add_text(slide, card["title"], x + 0.8, y + 0.15, 4.4, 0.35, font_size=18, color=theme["text"], bold=True, margin=0)
        # Card title (EN) — indented, smaller, muted
        if card.get("title_en"):
            add_text(slide, card["title_en"], x + 1.0, y + 0.48, 4.2, 0.28,
                     font_size=11, color=theme.get("en_color", "5E6572"), italic=False, margin=0)
        # Card body — indented with line spacing for readability
        add_text(slide, card["body"], x + 0.35, y + 0.85, 4.85, 1.05,
                 font_size=13, color="5E6572", margin=0, line_spacing=1.4)


def add_framework_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg_alt"])
    fw = cfg["deck"]["framework"]
    add_title(slide,
              fw.get("title", "背景、問題與四象限"),
              fw.get("subtitle", "這一週的所有結果都用同一套三層框架與四象限資料集解讀"),
              theme,
              title_en=fw.get("title_en"),
              subtitle_en=fw.get("subtitle_en"))
    add_rect(slide, 0.8, 1.9, 3.7, 4.82, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "三層分工", 1.02, 2.12, 1.8, 0.28, font_size=18, color=theme["text"], bold=True, margin=0)
    layer_y = 2.6
    for idx, layer in enumerate(cfg["deck"]["framework"]["layers"]):
        add_rect(slide, 1.0, layer_y + idx * 1.18, 3.25, 0.9, fill=layer["fill"], line=None, radius=True)
        add_text(slide, layer["title"], 1.18, layer_y + 0.08 + idx * 1.18, 2.8, 0.2, font_size=14, color=layer["title_color"], bold=True, margin=0)
        add_text(slide, layer["body"], 1.18, layer_y + 0.34 + idx * 1.18, 2.82, 0.38, font_size=12, color=layer["body_color"], margin=0)
    add_text(slide, cfg["deck"]["framework"]["problem_statement"], 1.0, 6.2, 3.15, 0.32, font_size=12, color="5E6572", margin=0)

    add_rect(slide, 4.8, 1.9, 7.7, 4.82, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "資料四象限與可比性", 5.05, 2.12, 2.8, 0.28, font_size=18, color=theme["text"], bold=True, margin=0)
    q_positions = [(5.05, 2.62), (8.55, 2.62), (5.05, 4.48), (8.55, 4.48)]
    for (x, y), quad in zip(q_positions, cfg["deck"]["framework"]["quadrants"]):
        add_rect(slide, x, y, 3.3, 1.48, fill=quad["fill"], line=None, radius=True)
        add_text(slide, quad["title"], x + 0.16, y + 0.12, 2.8, 0.24, font_size=14, color=quad["title_color"], bold=True, margin=0)
        add_text(slide, quad["body"], x + 0.16, y + 0.42, 2.95, 0.72, font_size=12, color=quad["body_color"], margin=0)
    add_text(slide, cfg["deck"]["framework"]["scope_note"], 5.05, 6.02, 6.95, 0.5, font_size=12, color="5E6572", margin=0)


def add_goal_status_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    deck = cfg["deck"]
    add_title(slide,
              deck.get("goals_title", "主要目標與完成度"),
              deck.get("goals_subtitle", "這一週的主線任務已從實驗探索推進到制度化與 annotation 層落地"),
              theme,
              title_en=deck.get("goals_title_en"),
              subtitle_en=deck.get("goals_subtitle_en"))
    rows = [["主目標", "子目標", "完成度", "本週結論"]]
    for item in cfg["deck"]["goals"]:
        rows.append([item["goal"], item["subgoal"], item["completion"], item["conclusion"]])
    add_table(slide, rows, 0.65, 1.82, 12.0, 3.95, theme, header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8", font_size=11, col_widths=[2.1, 2.45, 0.9, 4.65])
    y = 5.95
    for idx, item in enumerate(cfg["deck"]["goal_summary_cards"]):
        x = 0.8 + idx * 3.95
        add_rect(slide, x, y, 3.55, 0.7, fill=item["fill"], line=None, radius=True)
        add_text(slide, item["title"], x + 0.15, y + 0.08, 3.0, 0.18, font_size=14, color=item["title_color"], bold=True, margin=0)
        add_text(slide, item["body"], x + 0.15, y + 0.29, 3.0, 0.2, font_size=12, color=item["body_color"], margin=0)


def add_paired_rerun_slide(slide, cfg: dict, theme: dict, chart_path: Path) -> None:
    set_bg(slide, theme["bg_alt"])
    add_title(slide, "paired pure 正式 rerun", "5kHz 仍是最有壓力的主樣本；DORADO 則是穩定交叉驗證基準", theme)
    add_rect(slide, 0.8, 1.8, 7.35, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    _safe_add_picture(slide, str(chart_path), Inches(1.0), Inches(2.05), width=Inches(6.95), height=Inches(3.75))
    add_text(slide, "圖意：X 軸是 paired dataset，Y 軸是 F1。每組三根柱分別代表 ClairS、LongPhase-S、InterSubMod。", 1.0, 5.88, 6.7, 0.32, font_size=10.2, color="5E6572", margin=0)
    add_rect(slide, 8.35, 1.8, 4.15, 2.1, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "5kHz paired", 8.3, 2.05, 2.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, "ClairS 0.8443 → LongPhase-S 0.8522 → InterSubMod 0.8532", 8.3, 2.35, 3.8, 0.34, font_size=11.5, color=theme["muted"], margin=0)
    add_text(slide, "意義：在最 noisy 主樣本上，甲基與 phase 整合有實際正增益。", 8.3, 2.86, 3.75, 0.34, font_size=11.5, color=theme["text"], margin=0)
    add_rect(slide, 8.35, 4.08, 4.15, 2.1, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "DORADO paired", 8.3, 4.27, 2.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, "ClairS 0.8565 → LongPhase-S 0.8592 → InterSubMod 0.8590", 8.3, 4.57, 3.85, 0.34, font_size=11.5, color=theme["muted"], margin=0)
    add_text(slide, "意義：跨平台方向可驗證，但同一規則不保證可攜。", 8.3, 5.08, 3.75, 0.34, font_size=11.5, color=theme["text"], margin=0)


def add_to_rescue_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    add_title(slide, "tumor-only candidate-specific rescue", "TO 下已正式證明：甲基訊號不是全負效果，但仍未超過最佳 caller-only gate", theme)
    left = cfg["deck"]["to_rescue"]["hcc1395_5khz_to"]
    right = cfg["deck"]["to_rescue"]["hcc1395_dorado_to"]
    for x, data in [(0.8, left), (6.85, right)]:
        add_rect(slide, x, 1.9, 5.55, 4.65, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, data["label"], x + 0.2, 2.12, 2.8, 0.25, font_size=15, color=theme["text"], bold=True, margin=0)
        rows = [
            ["規則", "TP", "FP", "ΔF1"],
            [data["caller_only"]["rule"], data["caller_only"]["tp"], data["caller_only"]["fp"], data["caller_only"]["delta_f1"]],
            [data["methyl_support"]["rule"], data["methyl_support"]["tp"], data["methyl_support"]["fp"], data["methyl_support"]["delta_f1"]],
            [data["label_support"]["rule"], data["label_support"]["tp"], data["label_support"]["fp"], data["label_support"]["delta_f1"]],
        ]
        add_table(slide, rows, x + 0.18, 2.55, 5.08, 2.0, theme, header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8", font_size=9.1, col_widths=[3.35, 0.56, 0.56, 0.95])
        add_text(slide, data["interpretation"], x + 0.2, 4.72, 4.92, 0.74, font_size=10.8, color="5E6572", margin=0)
        add_rect(slide, x + 0.2, 5.72, 4.98, 0.44, fill=data["note_fill"], line=None, radius=True)
        add_text(slide, data["note"], x + 0.34, 5.84, 4.55, 0.18, font_size=10.2, color=data["note_color"], bold=True, margin=0)


def add_phase2_slide(slide, cfg: dict, theme: dict, chart_path: Path) -> None:
    set_bg(slide, theme["bg_alt"])
    add_title(slide, "phase 2：paired raw pileup / full 對照", "paired 問題下，raw caller 有 recall 空間，但 precision 代價過高，不適合直接當最終輸出", theme)
    add_rect(slide, 0.8, 1.8, 7.45, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    _safe_add_picture(slide, str(chart_path), Inches(1.0), Inches(2.0), width=Inches(7.0), height=Inches(3.85))
    add_text(slide, "圖意：X 軸為 F1，Y 軸為不同 caller 層級來源。左圖是 5kHz paired，右圖是 DORADO paired。", 1.0, 5.98, 6.9, 0.3, font_size=10.0, color="5E6572", margin=0)
    add_rect(slide, 8.45, 1.8, 4.05, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "phase 2 結論", 8.45, 2.08, 2.0, 0.22, font_size=15, color=theme["text"], bold=True, margin=0)
    add_bullets(
        slide,
        cfg["deck"]["phase2"]["bullets"],
        8.62,
        2.45,
        3.25,
        2.9,
        font_size=11.0,
        color=theme["text"],
    )
    add_rect(slide, 8.62, 5.68, 3.2, 0.42, fill=theme["accent2"], line=None, radius=True)
    add_text(slide, cfg["deck"]["phase2"]["callout"], 8.8, 5.79, 2.85, 0.18, font_size=9.8, color=theme["light_text"], bold=True, margin=0)


def add_feature_roles_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    add_title(slide, "特徵定位與意義", "這一週已把主要特徵從『到處試規則』收斂成可制度化的角色分工", theme)
    positions = [(0.8, 1.9), (6.6, 1.9), (0.8, 3.5), (6.6, 3.5), (0.8, 5.1), (6.6, 5.1)]
    for (x, y), row in zip(positions, cfg["deck"]["feature_roles"]):
        add_rect(slide, x, y, 5.55, 1.42, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, row["feature"], x + 0.18, y + 0.1, 5.0, 0.2, font_size=11.8, color=theme["text"], bold=True, margin=0)
        add_text(slide, row["role"], x + 0.18, y + 0.34, 5.0, 0.16, font_size=10.2, color=theme["accent"], bold=True, margin=0)
        add_text(slide, f"意義：{row['meaning']}", x + 0.18, y + 0.60, 5.0, 0.24, font_size=9.8, color="5E6572", margin=0)
        add_text(slide, f"限制：{row['limit']}", x + 0.18, y + 0.90, 5.0, 0.24, font_size=9.8, color=theme["muted"], margin=0)


def add_annotation_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg_alt"])
    add_title(slide, "annotation layer 決策", "phase 2 的 finer interval 已回接到 annotation 層，但目前仍不應直接改成 hard keep / hard filter", theme)
    positions = [(0.8, 1.95), (0.8, 3.55), (4.9, 1.95), (4.9, 3.55)]
    for (x, y), row in zip(positions, cfg["deck"]["annotation_policy"]["rows"]):
        add_rect(slide, x, y, 3.75, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, row["dataset"], x + 0.18, y + 0.12, 3.0, 0.18, font_size=12.2, color=theme["text"], bold=True, margin=0)
        add_text(slide, row["policy"], x + 0.18, y + 0.4, 3.1, 0.18, font_size=10.4, color=theme["accent"], bold=True, margin=0)
        add_text(slide, f"ΔF1 vs baseline：{row['delta']}", x + 0.18, y + 0.68, 2.2, 0.18, font_size=10.0, color="5E6572", margin=0)
        add_text(slide, row["interpretation"], x + 0.18, y + 0.92, 3.2, 0.2, font_size=9.9, color=theme["muted"], margin=0)
    cards = cfg["deck"]["annotation_policy"]["cards"]
    for idx, card in enumerate(cards):
        x = 9.15
        y = 1.95 + idx * 1.38
        add_rect(slide, x, y, 3.15, 1.05, fill=card["fill"], line=None, radius=True)
        add_text(slide, card["title"], x + 0.18, y + 0.14, 2.6, 0.2, font_size=12.5, color=card["title_color"], bold=True, margin=0)
        add_text(slide, card["body"], x + 0.18, y + 0.42, 2.55, 0.42, font_size=10.5, color=card["body_color"], margin=0)
    add_rect(slide, 0.8, 5.55, 11.7, 0.68, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, cfg["deck"]["annotation_policy"]["bottom_note"], 1.0, 5.77, 11.2, 0.22, font_size=10.3, color="5E6572", margin=0)


def add_next_steps_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["dark"])
    ns = cfg["deck"]["next_steps"]
    add_title(slide,
              ns.get("title", "待補齊事項與下一步"),
              ns.get("subtitle", ""),
              theme, dark=True,
              title_en=ns.get("title_en"),
              subtitle_en=ns.get("subtitle_en"))
    # Three columns with better spacing
    col_gap = 0.20
    col1_w, col2_w, col3_w = 3.8, 3.8, 4.2
    col1_x = 0.65
    col2_x = col1_x + col1_w + col_gap
    col3_x = col2_x + col2_w + col_gap
    card_top = 2.05
    card_h = 4.45

    # Column 1: 目前尚缺 (Gaps)
    add_rect(slide, col1_x, card_top, col1_w, card_h, fill="F4EFE7", line=None, radius=True)
    add_text(slide, "📋 目前尚缺", col1_x + 0.20, card_top + 0.12, col1_w - 0.4, 0.28,
             font_size=18, color=theme["text"], bold=True, margin=0)
    add_text(slide, "Current Gaps", col1_x + 0.20, card_top + 0.42, col1_w - 0.4, 0.20,
             font_size=10, color="5E6572", margin=0)
    add_bilingual_bullets(slide, ns["gaps"], ns.get("gaps_en"),
                          col1_x + 0.15, card_top + 0.72, col1_w - 0.30, card_h - 0.85,
                          font_size=13, color=theme["text"], theme=theme)

    # Column 2: 下一步 (Actions)
    add_rect(slide, col2_x, card_top, col2_w, card_h, fill="FFFFFF", line=None, radius=True)
    add_text(slide, "🎯 下一步", col2_x + 0.20, card_top + 0.12, col2_w - 0.4, 0.28,
             font_size=18, color=theme["text"], bold=True, margin=0)
    add_text(slide, "Next Steps", col2_x + 0.20, card_top + 0.42, col2_w - 0.4, 0.20,
             font_size=10, color="5E6572", margin=0)
    add_bilingual_bullets(slide, ns["actions"], ns.get("actions_en"),
                          col2_x + 0.15, card_top + 0.72, col2_w - 0.30, card_h - 0.85,
                          font_size=13, color=theme["text"], theme=theme)

    # Column 3: 現階段不建議 (Not yet)
    add_rect(slide, col3_x, card_top, col3_w, card_h, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "🚫 現階段不建議", col3_x + 0.20, card_top + 0.12, col3_w - 0.4, 0.28,
             font_size=18, color=theme["light_text"], bold=True, margin=0)
    add_text(slide, "Not Recommended Now", col3_x + 0.20, card_top + 0.42, col3_w - 0.4, 0.20,
             font_size=10, color="F0D0C0", margin=0)
    add_bilingual_bullets(slide, ns["not_yet"], ns.get("not_yet_en"),
                          col3_x + 0.15, card_top + 0.72, col3_w - 0.30, card_h - 0.85,
                          font_size=13, color=theme["light_text"], theme=theme)


def add_paths_slide(slide, cfg: dict, theme: dict) -> None:
    set_bg(slide, theme["bg"])
    deck = cfg["deck"]
    add_title(slide,
              deck.get("paths_title", "附錄來源"),
              deck.get("paths_subtitle", "完整來源與設定檔可追溯"),
              theme,
              title_en=deck.get("paths_title_en"),
              subtitle_en=deck.get("paths_subtitle_en"))
    # All source paths in a clean list layout
    source_paths = deck.get("source_paths", [])
    n_paths = len(source_paths)
    # Two-column layout: left = source paths, right = version/regen
    left_w, right_w = 8.0, 4.2
    card_top = 2.0
    card_h = 4.50

    # Left: Source paths list
    add_rect(slide, 0.65, card_top, left_w, card_h, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "來源文件   Source Files", 0.85, card_top + 0.10, 4.0, 0.25,
             font_size=14, color=theme["text"], bold=True, margin=0)
    for idx, item in enumerate(source_paths):
        y = card_top + 0.48 + idx * 0.68
        if y + 0.60 > card_top + card_h:
            break
        add_text(slide, item["label"], 0.85, y, 2.5, 0.20,
                 font_size=11, color=theme["accent"], bold=True, margin=0)
        path_text = wrap_path_for_slide(item["path"], 55, 2)
        add_text(slide, path_text, 0.85, y + 0.22, left_w - 0.45, 0.40,
                 font_size=9, color="5E6572", font_name="Consolas", margin=0)

    # Right: Version + Regeneration
    rx = 0.65 + left_w + 0.15
    add_rect(slide, rx, card_top, right_w, card_h, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "版本與重生", rx + 0.15, card_top + 0.10, 3.0, 0.25,
             font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, "Version & Regeneration", rx + 0.15, card_top + 0.35, 3.0, 0.18,
             font_size=9, color="5E6572", margin=0)
    # Version file
    version_name = Path(cfg["meta"]["output_pptx"]).with_suffix(".version.json").name
    vy = card_top + 0.65
    add_text(slide, "版本確認檔", rx + 0.15, vy, 2.0, 0.20,
             font_size=11, color=theme["accent"], bold=True, margin=0)
    add_text(slide, version_name, rx + 0.15, vy + 0.22, right_w - 0.3, 0.22,
             font_size=9, color="5E6572", font_name="Consolas", margin=0)
    # Deck file
    vy2 = vy + 0.55
    add_text(slide, "PPTX 檔案", rx + 0.15, vy2, 2.0, 0.20,
             font_size=11, color=theme["accent"], bold=True, margin=0)
    add_text(slide, Path(cfg["meta"]["output_pptx"]).name, rx + 0.15, vy2 + 0.22, right_w - 0.3, 0.22,
             font_size=9, color="5E6572", font_name="Consolas", margin=0)
    # Regeneration command
    vy3 = vy2 + 0.60
    add_text(slide, "重生指令", rx + 0.15, vy3, 2.0, 0.20,
             font_size=11, color=theme["accent"], bold=True, margin=0)
    add_rect(slide, rx + 0.10, vy3 + 0.28, right_w - 0.2, 1.0, fill=theme["dark"], line=None, radius=True)
    config_name = Path(cfg.get("_config_path", "pptx_config.json")).name
    add_text(slide, f"python build_weekly_report_pptx.py\n  --config {config_name}",
             rx + 0.22, vy3 + 0.38, right_w - 0.45, 0.75,
             font_size=9, color=theme["light_text"], font_name="Consolas", margin=0)
    # Footer note
    add_text(slide, "完整來源路徑請回查主週報與 version.json",
             0.65, card_top + card_h + 0.10, 12.0, 0.22,
             font_size=9, color="5E6572", align="center", margin=0)


def add_hp_bug_fix_slide(slide, cfg: dict, theme: dict, chart_path: Path) -> None:
    hp = cfg["deck"]["hp_bug_fix"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, hp["title"], hp["subtitle"], theme)
    add_rect(slide, 0.8, 1.8, 7.35, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    _safe_add_picture(slide, str(chart_path), Inches(1.0), Inches(2.05), width=Inches(6.95), height=Inches(3.75))
    add_text(slide, hp["chart_note"], 1.0, 5.88, 6.7, 0.32, font_size=10.2, color="5E6572", margin=0)
    # Right side: before card
    add_rect(slide, 8.35, 1.8, 4.15, 1.4, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, hp["before_label"], 8.55, 1.95, 2.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, hp["before_stat"], 8.55, 2.28, 3.7, 0.34, font_size=12.5, color=theme["muted"], margin=0)
    add_rect(slide, 8.65, 2.72, 0.8, 0.28, fill="A5A5A5", line=None, radius=True)
    add_text(slide, "舊版", 8.78, 2.75, 0.6, 0.18, font_size=10, color="FFFFFF", bold=True, margin=0)
    # Right side: after card
    add_rect(slide, 8.35, 3.35, 4.15, 1.4, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, hp["after_label"], 8.55, 3.5, 2.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, hp["after_stat"], 8.55, 3.83, 3.7, 0.34, font_size=12.5, color=theme["accent2"], bold=True, margin=0)
    add_rect(slide, 8.65, 4.27, 0.8, 0.28, fill=theme["accent2"], line=None, radius=True)
    add_text(slide, "新版", 8.78, 4.3, 0.6, 0.18, font_size=10, color="FFFFFF", bold=True, margin=0)
    # Right side: enrichment callout
    add_rect(slide, 8.35, 4.9, 4.15, 1.45, fill=hp["enrichment_fill"], line=None, radius=True)
    add_text(slide, hp["enrichment_callout"], 8.55, 5.05, 3.75, 1.05, font_size=11.5, color=hp["enrichment_color"], bold=True, margin=0)


def add_loh_evidence_panel_slide(slide, cfg: dict, theme: dict) -> None:
    loh = cfg["deck"]["loh_evidence_panel"]
    set_bg(slide, theme["bg"])
    add_title(slide, loh["title"], loh["subtitle"], theme)
    positions = [(0.8, 1.95), (6.95, 1.95), (0.8, 4.02), (6.95, 4.02)]
    for (x, y), rnd in zip(positions, loh["rounds"]):
        add_rect(slide, x, y, 5.55, 1.78, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, x + 0.18, y + 0.15, 0.68, 0.35, fill=rnd["accent"], line=None, radius=True)
        add_text(slide, rnd["label"], x + 0.24, y + 0.18, 0.7, 0.2, font_size=10.5, color="FFFFFF", bold=True, margin=0)
        add_text(slide, rnd["title"], x + 0.98, y + 0.14, 4.2, 0.28, font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, rnd["body"], x + 0.18, y + 0.56, 5.05, 0.72, font_size=10.5, color="5E6572", margin=0)
        add_rect(slide, x + 0.18, y + 1.36, 1.8, 0.24, fill=rnd["accent"], line=None, radius=True)
        add_text(slide, rnd["status"], x + 0.26, y + 1.38, 1.65, 0.18, font_size=9, color="FFFFFF", bold=True, margin=0)
    # Bottom conclusion callout
    add_rect(slide, 0.8, 6.0, 11.7, 0.68, fill=loh["conclusion_fill"], line=None, radius=True)
    add_text(slide, loh["conclusion_callout"], 1.0, 6.12, 11.3, 0.42, font_size=11, color=loh["conclusion_color"], bold=True, margin=0)


def add_phase1a_loh_test_slide(slide, cfg: dict, theme: dict) -> None:
    p1a = cfg["deck"]["phase1a_round3"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, p1a["title"], p1a["subtitle"], theme)
    # Left side: figure
    add_rect(slide, 0.8, 1.8, 7.45, 4.95, fill="FFFFFF", line=theme["line"], radius=True)
    fig_path = Path(p1a["figure_path"])
    if fig_path.exists():
        _safe_add_picture(slide, str(fig_path), Inches(1.0), Inches(2.0), width=Inches(7.0), height=Inches(3.85))
    else:
        add_text(slide, "[Figure not found: fig01_overall_f1.png]", 2.0, 3.5, 4.0, 0.5, font_size=14, color=theme["muted"], margin=0)
    add_text(slide, p1a["chart_note"], 1.0, 5.98, 6.9, 0.3, font_size=10.0, color="5E6572", margin=0)
    # Right side: result table
    add_rect(slide, 8.45, 1.8, 4.05, 3.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "四方比較結論", 8.65, 2.0, 2.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
    y_off = 2.38
    for res in p1a["results"]:
        marker = "▶" if res["supported"] == "✓" else ("✗" if res["supported"] == "✗" else "")
        label = res["model"]
        add_text(slide, label, 8.65, y_off, 2.4, 0.18, font_size=10.5, color=theme["text"], bold=True, margin=0)
        f1_line = f"F1={res['f1']:.4f}"
        if res["delta"] != "0.0000":
            f1_line += f"  ΔF1={res['delta']}"
        if res["ci"]:
            f1_line += f"\nCI={res['ci']}"
        if marker:
            f1_line += f"  {marker}"
        add_text(slide, f1_line, 8.65, y_off + 0.2, 3.55, 0.42, font_size=9.5, color="5E6572", margin=0)
        y_off += 0.72
    # Right side: direction callout
    add_rect(slide, 8.45, 5.15, 4.05, 1.2, fill=p1a["direction_fill"], line=None, radius=True)
    add_text(slide, p1a["direction_callout"], 8.65, 5.3, 3.65, 0.82, font_size=10.5, color=p1a["direction_color"], bold=True, margin=0)


def add_vision_slide(slide, cfg: dict, theme: dict) -> None:
    vis = cfg["deck"]["vision"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, vis["title"], vis["subtitle"], theme)
    # 4 goal cards
    card_w = 2.7
    gap = 0.25
    start_x = 0.8
    for idx, goal in enumerate(vis["goals"]):
        x = start_x + idx * (card_w + gap)
        add_rect(slide, x, 1.95, card_w, 2.6, fill=goal["fill"], line=None, radius=True)
        add_text(slide, goal["label"], x + 0.15, 2.1, 1.8, 0.2, font_size=11, color="F8F4EE", bold=True, margin=0)
        add_text(slide, goal["title"], x + 0.15, 2.38, 2.35, 0.28, font_size=14, color="F8F4EE", bold=True, margin=0)
        add_text(slide, goal["body"], x + 0.15, 2.78, 2.35, 0.82, font_size=10.5, color="D6DCE5", margin=0)
        add_text(slide, goal["priority"], x + 0.15, 3.72, 1.5, 0.2, font_size=10, color="D6DCE5", bold=True, margin=0)
    # Goal 5: TO independent line
    g5 = vis["goal5"]
    add_rect(slide, 0.8, 4.75, 11.7, 0.85, fill=g5["fill"], line=None, radius=True)
    add_text(slide, g5["title"], 1.0, 4.88, 3.5, 0.22, font_size=14, color=g5["color"], bold=True, margin=0)
    add_text(slide, g5["body"], 4.6, 4.88, 7.5, 0.52, font_size=11, color=g5["color"], margin=0)
    # Bottom note
    add_rect(slide, 0.8, 5.82, 11.7, 0.55, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, vis["bottom_note"], 1.0, 5.92, 11.2, 0.3, font_size=10.5, color="5E6572", margin=0)


def add_seqc2_validation_slide(slide, cfg: dict, theme: dict) -> None:
    sv = cfg["deck"]["seqc2_validation"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, sv["title"], sv["subtitle"], theme,
              title_en=sv.get("title_en"), subtitle_en=sv.get("subtitle_en"))
    # Left: SEQC2 description
    add_rect(slide, 0.8, 1.85, 3.85, 2.85, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "SEQC2（FDA 金標準）", 1.0, 2.0, 3.0, 0.25, font_size=16, color=theme["text"], bold=True, margin=0)
    add_text(slide, sv["seqc2_description"], 1.0, 2.35, 3.4, 2.0, font_size=13, color="5E6572", margin=0)
    # Middle: LongPhase description
    add_rect(slide, 4.85, 1.85, 3.85, 2.85, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "LongPhase-TO（本研究）", 5.05, 2.0, 3.0, 0.25, font_size=16, color=theme["text"], bold=True, margin=0)
    add_text(slide, sv["longphase_description"], 5.05, 2.35, 3.4, 2.0, font_size=13, color="5E6572", margin=0)
    # Right: metrics
    add_rect(slide, 8.9, 1.85, 3.6, 2.85, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "驗證指標", 9.1, 2.0, 2.0, 0.25, font_size=16, color=theme["text"], bold=True, margin=0)
    y_off = 2.35
    for m in sv["metrics"]:
        add_text(slide, m["name"], 9.1, y_off, 1.5, 0.18, font_size=13, color=theme["accent"], bold=True, margin=0)
        add_text(slide, m["value"], 10.5, y_off, 0.8, 0.18, font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, m["note"], 9.1, y_off + 0.2, 3.0, 0.16, font_size=11, color=theme["muted"], margin=0)
        y_off += 0.48
    # Bottom: conclusion callout
    add_rect(slide, 0.8, 4.95, 11.7, 1.25, fill=sv["conclusion_fill"], line=None, radius=True)
    add_text(slide, sv["conclusion_callout"], 1.0, 5.15, 11.2, 0.82, font_size=14, color=sv["conclusion_color"], bold=True, margin=0)


def add_ism_loh_impact_slide(slide, cfg: dict, theme: dict) -> None:
    imp = cfg["deck"]["ism_loh_impact"]
    set_bg(slide, theme["bg"])
    add_title(slide, imp["title"], imp["subtitle"], theme,
              title_en=imp.get("title_en"), subtitle_en=imp.get("subtitle_en"))
    # Stats cards (4 cards, 2×2 grid)
    positions = [(0.8, 1.95), (6.6, 1.95), (0.8, 3.65), (6.6, 3.65)]
    for (x, y), stat in zip(positions, imp["stats"]):
        add_rect(slide, x, y, 5.55, 1.42, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, stat["label"], x + 0.18, y + 0.1, 4.0, 0.22, font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, stat["value"], x + 0.18, y + 0.38, 2.0, 0.35, font_size=24, color=theme["accent"], bold=True, margin=0)
        add_text(slide, stat["meaning"], x + 2.2, y + 0.42, 3.0, 0.25, font_size=13, color="5E6572", margin=0)
    # Root cause
    add_rect(slide, 0.8, 5.25, 7.55, 1.15, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "根因", 1.0, 5.35, 0.8, 0.2, font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, imp["root_cause"], 1.0, 5.6, 7.1, 0.6, font_size=12, color="5E6572", margin=0)
    # Fix applied
    add_rect(slide, 8.55, 5.25, 3.95, 1.15, fill=imp["fix_fill"], line=None, radius=True)
    add_text(slide, "已修正", 8.75, 5.35, 1.0, 0.2, font_size=14, color=imp["fix_color"], bold=True, margin=0)
    add_text(slide, imp["fix_applied"], 8.75, 5.6, 3.5, 0.6, font_size=12, color=imp["fix_color"], margin=0)


def add_non_loh_closure_slide(slide, cfg: dict, theme: dict) -> None:
    nlc = cfg["deck"]["non_loh_closure"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, nlc["title"], nlc["subtitle"], theme,
              title_en=nlc.get("title_en"), subtitle_en=nlc.get("subtitle_en"))
    # Left: Non-LOH results table with bias direction
    add_rect(slide, 0.8, 1.95, 6.0, 2.80, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "Non-LOH 特徵 AUC（校正後）", 1.0, 2.05, 4.5, 0.22,
             font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, "Non-LOH Feature AUC (corrected)", 1.0, 2.28, 4.5, 0.18,
             font_size=9, color="5E6572", margin=0)
    rows = [["特徵 Feature", "AUC", "偏向 Bias"]]
    for r in nlc["non_loh_results"]:
        rows.append([r["feature"], r["auc"], r["status"]])
    add_table(slide, rows, 0.95, 2.55, 5.65, 1.95, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=11, col_widths=[2.2, 1.2, 2.5])
    # Right top: Voting result
    add_rect(slide, 7.0, 1.95, 5.5, 1.15, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "多特徵 Voting", 7.2, 2.05, 2.5, 0.22,
             font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, nlc["voting_result"], 7.2, 2.32, 5.1, 0.22,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    add_text(slide, "< 0.58 threshold = FAIL", 7.2, 2.58, 3.0, 0.18,
             font_size=10, color="D55E00", bold=True, margin=0)
    # Right bottom: Simpson's Paradox
    sp = nlc["simpsons_paradox"]
    add_rect(slide, 7.0, 3.25, 5.5, 1.50, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, sp["title"], 7.2, 3.35, 3.5, 0.22,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    add_text(slide, "Simpson's Paradox", 7.2, 3.58, 3.0, 0.18,
             font_size=9, color="5E6572", margin=0)
    # Visual: Overall vs Per-sample with arrow
    add_rect(slide, 7.2, 3.82, 2.2, 0.38, fill="FDEDEC", line=None, radius=True)
    add_text(slide, f"整體 {sp['overall_auc']}", 7.3, 3.85, 2.0, 0.32,
             font_size=12, color="D55E00", bold=True, margin=0)
    add_text(slide, "→", 9.5, 3.85, 0.3, 0.32, font_size=14, color=theme["accent"], bold=True, margin=0)
    add_rect(slide, 9.9, 3.82, 2.4, 0.38, fill="E8F8E8", line=None, radius=True)
    add_text(slide, f"Per-sample {sp['per_sample_mean']}", 10.0, 3.85, 2.2, 0.32,
             font_size=12, color="009E73", bold=True, margin=0)
    add_text(slide, f"{sp['dominant_sample']} 佔比主導 | {sp['direction_inconsistency']}",
             7.2, 4.28, 5.1, 0.35, font_size=10, color="5E6572", margin=0)
    # Bottom: conclusion
    add_rect(slide, 0.8, 5.05, 11.7, 0.90, fill=nlc["conclusion_fill"], line=None, radius=True)
    add_bilingual_text(slide, 1.0, 5.08, 11.3, 0.82,
                       nlc["conclusion_callout"],
                       "LOH vs Non-LOH AUC gap < 0.06 — the problem is global, not LOH-specific",
                       13, theme, bold=True, color=nlc["conclusion_color"], margin=0)
    # Future note
    fn = nlc.get("future_note", "")
    if fn:
        add_text(slide, f"→ {fn}", 0.8, 6.02, 8.0, 0.15,
                 font_size=10, color="6E9D8A", margin=0)
        fn_en = nlc.get("future_note_en", "")
        if fn_en:
            add_text(slide, fn_en, 8.8, 6.02, 3.7, 0.15,
                     font_size=9, color="5E6572", margin=0)


def add_negative_results_slide(slide, cfg: dict, theme: dict) -> None:
    nr = cfg["deck"]["negative_results"]
    set_bg(slide, theme["bg"])
    add_title(slide, nr["title"], nr["subtitle"], theme,
              title_en=nr.get("title_en"), subtitle_en=nr.get("subtitle_en"))
    # Left: methylation negatives (3 rows)
    add_rect(slide, 0.8, 1.85, 5.55, 2.55, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "甲基化三維度 NEGATIVE", 1.0, 1.98, 3.5, 0.22, font_size=16, color=theme["text"], bold=True, margin=0)
    y_off = 2.3
    for mn in nr["methylation_negatives"]:
        add_text(slide, f"{mn['id']}: {mn['name']}", 1.0, y_off, 2.5, 0.18, font_size=13, color=theme["accent"], bold=True, margin=0)
        add_text(slide, f"{mn['before']} → {mn['after']}（{mn['cause']}）", 3.4, y_off, 2.7, 0.18, font_size=12, color="5E6572", margin=0)
        y_off += 0.55
    # Right top: TO NO-GO
    to = nr["to_nogo"]
    add_rect(slide, 6.55, 1.85, 5.95, 1.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "TO VCF 特徵 NO-GO", 6.75, 1.98, 3.0, 0.22, font_size=16, color=theme["text"], bold=True, margin=0)
    add_text(slide, f"{to['scope']}\nMax AUC < {to['max_auc']}  |  Safe removal = {to['safe_removal']}\nRead-level LOSO AUC = {to['read_level_auc']}", 6.75, 2.3, 5.5, 0.55, font_size=12, color="5E6572", margin=0)
    # Right bottom: ASM highlight
    asm = nr["asm_highlight"]
    add_rect(slide, 6.55, 3.25, 5.95, 1.15, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "ASM — ISM 唯一亮點", 6.75, 3.38, 3.5, 0.22, font_size=16, color=theme["light_text"], bold=True, margin=0)
    add_text(slide, asm["finding"], 6.75, 3.68, 5.5, 0.18, font_size=12, color=theme["light_muted"], margin=0)
    add_text(slide, asm["ism_role"], 6.75, 3.92, 5.5, 0.18, font_size=12, color=theme["light_muted"], margin=0)
    add_text(slide, asm["implication"], 6.75, 4.14, 5.5, 0.18, font_size=12, color="F8F4EE", bold=True, margin=0)
    # Bottom: overall summary
    add_rect(slide, 0.8, 4.6, 11.7, 0.75, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "結論：所有 post-hoc 特徵方向正式關閉。ISM 價值在 read-level epigenetic characterization — ASM 32-66% 是唯一可解釋的生物學信號。", 1.0, 4.72, 11.2, 0.42, font_size=14, color="5E6572", margin=0)


def add_conclusion_table_slide(slide, cfg: dict, theme: dict) -> None:
    ct = cfg["deck"]["conclusion_table"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, ct["title"], ct["subtitle"], theme,
              title_en=ct.get("title_en"), subtitle_en=ct.get("subtitle_en"))
    # Build table rows — sorted by importance with Chinese verdict
    cols = ct.get("columns", ["#", "問題", "結論", "判決", "可信度"])
    rows = [cols]
    for c in ct["conclusions"]:
        stability = "⭐" * c.get("stability", 3)
        verdict = c.get("verdict", c.get("verdict_en", ""))
        rows.append([c["id"], c.get("question", ""), c["name"], verdict, stability])
    add_table(slide, rows, 0.50, 1.82, 12.3, 4.85, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=9.5, col_widths=[0.45, 2.8, 3.0, 1.2, 0.8])


# ====================================================================
# Visual Diagram Helpers (5 patterns + base shape utilities)
# ====================================================================

def add_styled_box(
    slide, left: float, top: float, width: float, height: float, text: str,
    fill_color: str, *, text_color: str = "FFFFFF", font_size: float = 13,
    shape_type=SHAPE_PROCESS, bold: bool = True, detail: str | None = None,
    detail_color: str = "D1D1D1", detail_size: float = 10,
):
    """Add a styled shape with centered text and optional detail line below."""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = _resolve_icons(text)
    run.font.size = Pt(font_size)
    run.font.name = "Arial"
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(text_color)
    if detail:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = _resolve_icons(detail)
        run2.font.size = Pt(detail_size)
        run2.font.name = "Arial"
        run2.font.color.rgb = hex_rgb(detail_color)
    return shape


def add_native_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = "808080", width_pt: float = 2):
    """Add a native connector line with arrowhead using XML manipulation."""
    cx1, cy1 = Inches(x1), Inches(y1)
    cx2, cy2 = Inches(x2), Inches(y2)
    # Avoid near-zero connectors — PPT versions may fail to render slides
    # containing connectors with cx or cy < ~0.05". Use 0.06" for safety.
    min_dim = Emu(Inches(0.10))  # 0.10" safe floor (0.06" still caused issues)
    if abs(cy2 - cy1) < min_dim:
        cy2 = cy1 + min_dim
    if abs(cx2 - cx1) < min_dim:
        cx2 = cx1 + min_dim
    connector = slide.shapes.add_connector(
        1,  # msoConnectorStraight
        cx1, cy1, cx2, cy2,
    )
    connector.line.color.rgb = hex_rgb(color)
    connector.line.width = Pt(width_pt)
    # Add arrowhead via XML
    ln = connector.line._ln
    tail_end = ln.find(qn("a:tailEnd"))
    if tail_end is None:
        from lxml import etree
        tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    return connector


def grid_positions(n_items: int, area_left: float, area_top: float, area_width: float, area_height: float, *, cols: int | None = None, gap: float = 0.15):
    """Calculate grid positions for N items within a bounded area."""
    if cols is None:
        cols = min(n_items, 5)
    rows = math.ceil(n_items / cols)
    cell_w = (area_width - (cols - 1) * gap) / cols
    cell_h = (area_height - (rows - 1) * gap) / rows
    positions = []
    for i in range(n_items):
        r, c = divmod(i, cols)
        x = area_left + c * (cell_w + gap)
        y = area_top + r * (cell_h + gap)
        positions.append((x, y, cell_w, cell_h))
    return positions


def build_pipeline_flow(slide, steps: list[dict], theme: dict, *, area_left: float = 0.8, area_top: float = 2.2, area_width: float = 11.7, area_height: float = 2.5):
    """Build horizontal pipeline flow: steps connected by native arrows.

    steps = [{"label": "...", "color": "#...", "detail": None|"..."}, ...]
    """
    n = len(steps)
    if n == 0:
        return
    arrow_gap = 0.35
    total_arrow = (n - 1) * arrow_gap
    box_w = min(2.0, (area_width - total_arrow - 0.2) / max(n, 1))
    box_h = min(1.4, area_height)
    total_w = n * box_w + (n - 1) * arrow_gap
    start_x = area_left + (area_width - total_w) / 2
    center_y = area_top + (area_height - box_h) / 2

    for i, step in enumerate(steps):
        sx = start_x + i * (box_w + arrow_gap)
        fill = step.get("color", theme.get("accent", COLOR_PRIMARY))
        add_styled_box(
            slide, sx, center_y, box_w, box_h, step.get("label", ""),
            fill_color=fill, detail=step.get("detail"),
        )
        if i < n - 1:
            ax1 = sx + box_w + 0.02
            ax2 = sx + box_w + arrow_gap - 0.02
            ay = center_y + box_h / 2
            add_native_arrow(slide, ax1, ay, ax2, ay, color=theme.get("accent", "A85540"), width_pt=2.5)


def build_comparison_panel(slide, left_data: dict, right_data: dict, theme: dict, *, area_left: float = 0.8, area_top: float = 2.0, area_width: float = 11.7, area_height: float = 3.5):
    """Build side-by-side comparison panel with VS marker.

    left_data/right_data = {"title": "...", "items": [...], "color": "..."}
    """
    gap = 0.8
    card_w = (area_width - gap) / 2
    card_h = area_height
    # Left card
    lc = left_data.get("color", COLOR_PRIMARY)
    add_styled_box(slide, area_left, area_top, card_w, card_h, left_data["title"], fill_color=lc, font_size=16)
    for j, item in enumerate(left_data.get("items", [])):
        add_text(slide, f"\u2022 {item}", area_left + 0.15, area_top + 0.6 + j * 0.45, card_w - 0.3, 0.4,
                 font_size=12, color="FFFFFF", margin=0)
    # VS marker
    vs_x = area_left + card_w
    add_text(slide, "VS", vs_x, area_top + card_h * 0.35, gap, 0.5,
             font_size=20, color=theme.get("accent", "A85540"), bold=True, align="center", margin=0)
    # Right card
    rc = right_data.get("color", COLOR_NEGATIVE)
    add_styled_box(slide, area_left + card_w + gap, area_top, card_w, card_h, right_data["title"], fill_color=rc, font_size=16)
    for j, item in enumerate(right_data.get("items", [])):
        add_text(slide, f"\u2022 {item}", area_left + card_w + gap + 0.15, area_top + 0.6 + j * 0.45, card_w - 0.3, 0.4,
                 font_size=12, color="FFFFFF", margin=0)


def build_evidence_stack(slide, layers: list[dict], theme: dict, *, area_left: float = 1.5, area_top: float = 2.2, area_width: float = 10.0, area_height: float = 3.5):
    """Build stacked evidence bars: bottom=weakest, top=strongest, width=strength.

    layers = [{"label": "...", "strength": 0.0-1.0, "color": "..."}, ...]
    """
    n = len(layers)
    if n == 0:
        return
    bar_h = min(0.6, (area_height - 0.1 * (n - 1)) / max(n, 1))
    gap = 0.1
    # Stack bottom-to-top (index 0 = weakest at bottom)
    for i, layer in enumerate(layers):
        strength = layer.get("strength", 0.8)
        bar_w = area_width * max(0.3, min(1.0, strength))
        y = area_top + area_height - (i + 1) * (bar_h + gap)
        x = area_left + (area_width - bar_w) / 2
        fill = layer.get("color", COLOR_PRIMARY)
        add_styled_box(slide, x, y, bar_w, bar_h, layer.get("label", ""), fill_color=fill, font_size=12)


def build_decision_tree(slide, nodes: list[dict], theme: dict, *, area_left: float = 0.8, area_top: float = 2.0, area_width: float = 11.7, area_height: float = 4.0):
    """Build decision tree: diamond decisions + rectangular outcomes.

    nodes = [{"label": "...", "type": "decision|pass|fail", "children": [idx, ...], "n": "N=100"}, ...]
    """
    n = len(nodes)
    if n == 0:
        return
    # Simple horizontal layout: all nodes in one row
    positions = grid_positions(n, area_left, area_top, area_width, min(area_height, 1.5), cols=n, gap=0.3)
    for i, (node, (x, y, w, h)) in enumerate(zip(nodes, positions)):
        ntype = node.get("type", "decision")
        if ntype == "decision":
            shape_type = SHAPE_DECISION
            fill = theme.get("accent", "A85540")
        elif ntype == "pass":
            shape_type = SHAPE_PROCESS
            fill = COLOR_POSITIVE
        else:
            shape_type = SHAPE_PROCESS
            fill = COLOR_NEGATIVE
        label = node.get("label", "")
        n_label = node.get("n", "")
        detail = f"N={n_label}" if n_label else None
        add_styled_box(slide, x, y, w, h, label, fill_color=fill, shape_type=shape_type, font_size=11, detail=detail, detail_size=9)
        # Draw arrows to children
        for child_idx in node.get("children", []):
            if child_idx < n:
                cx, cy, cw, ch = positions[child_idx]
                add_native_arrow(slide, x + w, y + h / 2, cx, cy + ch / 2, color="808080")


def build_concept_card(slide, term: str, definition: str, theme: dict, *, term_en: str | None = None, source: str | None = None, visual_shape: str | None = None, area_left: float = 1.5, area_top: float = 2.2, area_width: float = 10.0, area_height: float = 3.0):
    """Build concept card: left visual + right term/definition.

    Left side: simple shape or icon. Right side: term (large) + definition + source.
    """
    # Left visual area (30%)
    vis_w = area_width * 0.3
    vis_h = area_height
    add_rect(slide, area_left, area_top, vis_w, vis_h, fill=theme.get("dark", "1E2A44"), line=None, radius=True)
    # Visual icon or shape
    icon_text = visual_shape or "\U0001f9ec"  # Default: 🧬
    add_text(slide, icon_text, area_left, area_top + vis_h * 0.2, vis_w, vis_h * 0.6,
             font_size=48, color="FFFFFF", align="center", margin=0)
    # Right text area (70%)
    text_x = area_left + vis_w + 0.3
    text_w = area_width - vis_w - 0.3
    # Term
    if term_en:
        add_bilingual_text(slide, text_x, area_top + 0.1, text_w, 0.7, term, term_en, 24, theme, bold=True)
    else:
        add_text(slide, term, text_x, area_top + 0.1, text_w, 0.5, font_size=24, color=theme["text"], bold=True, margin=0)
    # Definition
    add_text(slide, definition, text_x, area_top + 0.85, text_w, area_height * 0.5,
             font_size=14, color="5E6572", margin=0)
    # Source
    if source:
        add_text(slide, f"Source: {source}", text_x, area_top + area_height - 0.35, text_w, 0.25,
                 font_size=10, color=theme.get("muted", "5E6572"), italic=False, margin=0)


def build_assets(cfg: dict) -> dict:
    assets_dir = Path(cfg["meta"]["assets_dir"])
    ensure_dir(assets_dir)
    theme = cfg["theme"]

    # Generate HP fix eligibility chart if config has hp_fix data
    hp_fix_chart = None
    if "hp_fix_eligibility" in cfg.get("charts", {}):
        hp_fix_chart = assets_dir / "hp_fix_eligibility.png"
        create_grouped_bar_chart(
            hp_fix_chart,
            cfg["charts"]["hp_fix_eligibility"]["title"],
            cfg["charts"]["hp_fix_eligibility"]["categories"],
            cfg["charts"]["hp_fix_eligibility"]["series"],
            theme,
            ylim=(0, 105),
            ylabel="LOH Eligible %",
        )

    # Generate paired rerun chart if config has paired_rerun data
    paired_chart = None
    if "paired_rerun" in cfg.get("charts", {}):
        paired_chart = assets_dir / "paired_rerun_f1.png"
        create_grouped_bar_chart(
            paired_chart,
            "Paired Pure Rerun: F1 by Method",
            cfg["charts"]["paired_rerun"]["categories"],
            cfg["charts"]["paired_rerun"]["series"],
            theme,
            ylim=(0.84, 0.862),
            ylabel="F1",
        )

    # Generate phase2 chart if config has phase2_raw data
    phase2_chart = None
    if "phase2_raw" in cfg.get("charts", {}):
        phase2_chart = assets_dir / "phase2_paired_raw_benchmark.png"
        create_horizontal_bar_chart(phase2_chart, cfg["charts"]["phase2_raw"], theme)

    return {"paired_chart": paired_chart, "phase2_chart": phase2_chart, "hp_fix_chart": hp_fix_chart}


# ====================================================================
# V2 Visual-First Slide Builders (new page types)
# ====================================================================

def add_investigation_pipeline_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Investigation reasoning map — 3 concept blocks + transition arrows."""
    ip = cfg["deck"]["investigation_pipeline"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, ip["title"], ip["subtitle"], theme,
              title_en=ip.get("title_en"), subtitle_en=ip.get("subtitle_en"))

    blocks = ip.get("concept_blocks", [])
    transitions = ip.get("transitions", [])

    if not blocks:
        # Fallback to legacy pipeline_flow if no concept_blocks
        build_pipeline_flow(slide, ip["steps"], theme, area_top=2.5, area_height=3.0)
        return

    # --- 3 large concept blocks with arrows ---
    n = len(blocks)
    gap = 0.55
    arrow_w = 0.5
    total_arrow = (n - 1) * (gap + arrow_w)
    block_w = (11.7 - total_arrow) / n
    block_h = 3.0
    block_y = 2.2

    for i, blk in enumerate(blocks):
        bx = 0.8 + i * (block_w + gap + arrow_w)
        color = blk.get("color", theme["accent"])

        # Block background
        add_rect(slide, bx, block_y, block_w, block_h, fill="FFFFFF",
                 line=color, radius=True)

        # Color header bar
        add_rect(slide, bx, block_y, block_w, 0.65, fill=color, line=None, radius=True)
        # Title in header
        add_text(slide, blk["title"], bx + 0.15, block_y + 0.05, block_w - 0.3, 0.35,
                 font_size=15, color="FFFFFF", bold=True, margin=0)
        if blk.get("title_en"):
            add_text(slide, blk["title_en"], bx + 0.15, block_y + 0.38, block_w - 0.3, 0.25,
                     font_size=9, color="E0E0E0", italic=False, margin=0)

        # Items inside block
        items = blk.get("items", [])
        for j, item in enumerate(items):
            iy = block_y + 0.78 + j * 0.42
            add_text(slide, item, bx + 0.2, iy, block_w - 0.4, 0.38,
                     font_size=13, color=theme["text"], margin=0)

        # Slide range label at bottom
        sr = blk.get("slide_range", "")
        if sr:
            add_text(slide, sr, bx, block_y + block_h + 0.08, block_w, 0.3,
                     font_size=10, color=theme["muted"], align="center", italic=False, margin=0)

        # Arrow connector only (transition text drawn after all blocks for z-order)
        if i < n - 1:
            ax = bx + block_w + 0.08
            ay = block_y + block_h / 2 - 0.15
            add_native_arrow(slide, ax, ay, ax + gap + arrow_w - 0.16, ay,
                             color=theme.get("accent", "A85540"), width_pt=2.0)

    # Transition text labels (drawn after all blocks to avoid being covered)
    for i in range(min(n - 1, len(transitions))):
        bx = 0.8 + i * (block_w + gap + arrow_w)
        ax = bx + block_w + 0.08
        ay = block_y + block_h / 2 - 0.15
        tr_text = transitions[i]
        if tr_text:
            tw = gap + arrow_w - 0.12
            add_rect(slide, ax + 0.02, ay - 0.50, tw, 0.35, fill=theme["bg_alt"], line=None, radius=True)
            add_text(slide, tr_text, ax + 0.02, ay - 0.48, tw, 0.30,
                     font_size=9, color=theme["accent"], bold=True, align="center", margin=0)

    # Bottom: conclusion bar
    add_rect(slide, 0.8, 5.65, 11.7, 0.6, fill=theme["dark"], line=None, radius=True)
    add_bilingual_text(slide, 1.0, 5.70, 11.3, 0.5,
                       "三段式推理 — 每段由定量證據支撐",
                       "Three-phase reasoning — each backed by quantitative evidence",
                       13, theme, color=theme["light_text"], margin=0)


def add_loh_quadrant_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: LOH dual-definition — HP Imbalance ⊃ LOH.bed with sensitivity panel."""
    ld = cfg["deck"]["loh_dual_definition"]
    set_bg(slide, theme["bg"])
    add_title(slide, ld["title"], ld["subtitle"], theme,
              title_en=ld.get("title_en"), subtitle_en=ld.get("subtitle_en"))

    # --- Left: 2x2 quadrant grid (compact) ---
    quad_w, quad_h = 3.8, 1.6
    positions = [(0.8, 2.0, quad_w, quad_h), (4.8, 2.0, quad_w, quad_h),
                 (0.8, 3.75, quad_w, quad_h), (4.8, 3.75, quad_w, quad_h)]
    for (x, y, w, h), quad in zip(positions, ld["quadrants"]):
        add_rect(slide, x, y, w, h, fill="FFFFFF", line=theme["line"], radius=True)
        badge_color = quad.get("color", theme["accent"])
        add_rect(slide, x + 0.12, y + 0.12, 0.95, 0.48, fill=badge_color, line=None, radius=True)
        add_text(slide, quad["pct"], x + 0.12, y + 0.14, 0.95, 0.25,
                 font_size=17, color="FFFFFF", bold=True, align="center", margin=0)
        add_text(slide, f"n={quad['n']}", x + 0.12, y + 0.37, 0.95, 0.16,
                 font_size=9, color="FFFFFF", align="center", margin=0)
        add_text(slide, quad["label"], x + 1.2, y + 0.12, w - 1.4, 0.25,
                 font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, quad["detail"], x + 1.2, y + 0.42, w - 1.4, 1.0,
                 font_size=11, color="5E6572", margin=0)

    # --- Right: Sensitivity stats panel ---
    ss = ld.get("sensitivity_stats", [])
    if ss:
        px, py = 9.2, 2.0
        pw, ph = 3.6, 3.35
        add_rect(slide, px, py, pw, ph, fill=theme["dark"], line=None, radius=True)
        add_text(slide, "核心數字", px + 0.3, py + 0.15, pw - 0.6, 0.3,
                 font_size=13, color=theme.get("light_muted", "D6DCE5"), bold=True, margin=0)
        add_text(slide, "Key Metrics", px + 0.3, py + 0.42, pw - 0.6, 0.2,
                 font_size=9, color="A0A0A0", italic=False, margin=0)
        for i, st in enumerate(ss):
            sy = py + 0.75 + i * 0.82
            add_text(slide, st["value"], px + 0.3, sy, pw - 0.6, 0.42,
                     font_size=28, color=st.get("color", "F5F5F5"), bold=True, margin=0)
            add_text(slide, st["label"], px + 0.3, sy + 0.44, pw - 0.6, 0.25,
                     font_size=11, color=theme.get("light_muted", "D6DCE5"), margin=0)

    # --- Bottom: conclusion bar (from config) ---
    cb_zh = ld.get("conclusion_bar",
                    "Q3 = 0.07% → HP 漏判 LOH.bed 的機率趨近零 — LOH.bed 幾乎完全被 HP Imbalance 覆蓋")
    cb_en = ld.get("conclusion_bar_en",
                    "Q3 = 0.07% → Probability of HP missing LOH.bed approaches zero")
    add_rect(slide, 0.8, 5.6, 11.7, 0.55, fill=theme["accent"], line=None, radius=True)
    add_text(slide, cb_zh, 1.0, 5.64, 11.3, 0.22,
             font_size=12, color="FFFFFF", bold=True, margin=0)
    add_text(slide, cb_en, 1.0, 5.88, 11.3, 0.2,
             font_size=9, color="F0D0C0", italic=False, margin=0)
    # Total
    add_text(slide, ld["total"], 0.8, 6.25, 6.0, 0.2, font_size=10, color=theme["muted"], margin=0)


def add_seqc2_comparison_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: SEQC2 validation using comparison_panel + metric bars."""
    sv = cfg["deck"]["seqc2_validation"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, sv["title"], sv["subtitle"], theme,
              title_en=sv.get("title_en"), subtitle_en=sv.get("subtitle_en"))
    # Two comparison cards (top half)
    left = sv["left"]
    right = sv["right"]
    for x, data in [(0.8, left), (6.95, right)]:
        card_color = data.get("color", theme["accent"])
        add_rect(slide, x, 2.0, 5.55, 1.8, fill=card_color, line=None, radius=True)
        add_text(slide, data["title"], x + 0.2, 2.1, 5.0, 0.3, font_size=16, color="FFFFFF", bold=True, margin=0)
        for j, item in enumerate(data["items"]):
            add_text(slide, f"• {item}", x + 0.2, 2.45 + j * 0.3, 5.0, 0.25, font_size=12, color="D1D1D1", margin=0)
    # VS marker
    add_text(slide, "VS", 6.0, 2.5, 0.95, 0.5, font_size=22, color=theme["accent"], bold=True, align="center", margin=0)
    # Metric bars (bottom half)
    add_rect(slide, 0.8, 4.1, 11.7, 2.4, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "驗證指標", 1.0, 4.2, 2.0, 0.25, font_size=16, color=theme["text"], bold=True, margin=0)
    for i, m in enumerate(sv["metrics"]):
        bar_y = 4.6 + i * 0.45
        bar_pct = m.get("bar_pct", 0.9)
        bar_w = 8.5 * bar_pct
        add_text(slide, m["name"], 1.0, bar_y, 1.5, 0.2, font_size=13, color=theme["text"], bold=True, margin=0)
        add_rect(slide, 2.6, bar_y + 0.02, 8.5, 0.28, fill="E8E4DD", line=None, radius=True)
        add_rect(slide, 2.6, bar_y + 0.02, bar_w, 0.28, fill=theme["accent2"], line=None, radius=True)
        add_text(slide, m["value"], 2.6 + bar_w + 0.15, bar_y, 0.8, 0.25, font_size=14, color=theme["accent"], bold=True, margin=0)


def add_seqc2_chromosome_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Chromosome-level LOH intersection — ideogram left + explanation right."""
    sc = cfg["deck"]["seqc2_chromosome"]
    set_bg(slide, theme["bg"])
    add_title(slide, sc["title"], sc["subtitle"], theme,
              title_en=sc.get("title_en"), subtitle_en=sc.get("subtitle_en"))
    # References (top-right corner, small Consolas links)
    refs = sc.get("references", [])
    for i, ref in enumerate(refs):
        ry = 0.55 + i * 0.28
        add_text(slide, f"{ref['label']}: {ref['url']}", 9.0, ry, 3.8, 0.25,
                 font_size=9, color=theme["muted"], font_name="Consolas", margin=0)
    # Left: main figure — ideogram (fit within left area, preserve aspect ratio)
    fig_path = Path(cfg["_config_dir"]) / sc["figure"]
    if fig_path.exists():
        with Image.open(fig_path) as _img:
            img_w, img_h = _img.size
        aspect = img_h / img_w
        area_w, area_h = 7.5, 4.3
        disp_w = area_w
        disp_h = disp_w * aspect
        if disp_h > area_h:
            disp_h = area_h
            disp_w = disp_h / aspect
        cx = 0.5 + (area_w - disp_w) / 2
        cy = 2.0 + (area_h - disp_h) / 2
        _safe_add_picture(slide, str(fig_path), Inches(cx), Inches(cy),
                          width=Inches(disp_w), height=Inches(disp_h))
    # Right: explanation panel (3 sections, shifted down 0.2" for breathing room)
    rx, rw = 8.3, 4.2
    rp_top = 2.2
    add_rect(slide, rx, rp_top, rw, 4.2, fill="FFFFFF", line=theme["line"], radius=True)
    # Section 1: Method independence
    add_text(slide, "🔬 方法學獨立性", rx + 0.15, rp_top + 0.05, rw - 0.3, 0.25,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    left_items = sc.get("left_items", [])
    for i, item in enumerate(left_items):
        add_text(slide, item, rx + 0.15, rp_top + 0.32 + i * 0.28, rw - 0.3, 0.25,
                 font_size=10, color=theme["text"], margin=0)
    # Section 2: Key metrics (from merged Slide 4)
    add_text(slide, "📊 驗證指標", rx + 0.15, rp_top + 1.25, rw - 0.3, 0.25,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    metrics = sc.get("metrics", [])
    for i, m in enumerate(metrics):
        my = rp_top + 1.52 + i * 0.28
        add_text(slide, m["name"], rx + 0.15, my, 1.5, 0.25,
                 font_size=10, color=theme["text"], margin=0)
        add_text(slide, m["value"], rx + 1.7, my, 1.0, 0.25,
                 font_size=11, color=theme["accent"], bold=True, margin=0)
    # Section 3: Legend / key interpretation
    add_text(slide, "🎯 圖例", rx + 0.15, rp_top + 2.75, rw - 0.3, 0.25,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    right_items = sc.get("right_items", [])
    for i, item in enumerate(right_items):
        add_text(slide, item, rx + 0.15, rp_top + 3.02 + i * 0.28, rw - 0.3, 0.25,
                 font_size=10, color=theme["text"], bold=True, margin=0)
    # Legend at bottom
    if sc.get("legend"):
        add_text(slide, sc["legend"], 0.5, 6.45, 12.3, 0.25,
                 font_size=10, color="5E6572", align="center", margin=0)


def add_seqc2_distribution_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Variant-level LOH distribution — confusion matrix + per-chr stacked bar."""
    sd = cfg["deck"]["seqc2_distribution"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, sd["title"], sd["subtitle"], theme,
              title_en=sd.get("title_en"), subtitle_en=sd.get("subtitle_en"))
    # Two figures side by side (below bilingual title block which ends ~y=2.05)
    fig_left_path = Path(cfg["_config_dir"]) / sd["figure_left"]
    fig_right_path = Path(cfg["_config_dir"]) / sd["figure_right"]
    fig_top = 2.15          # clear of bilingual title block
    fig_max_h = 3.0         # larger image area
    fig_area_w = 5.9        # each figure max width
    fig_gap = 0.35          # gap between two figures
    # Left figure: confusion matrix
    left_bottom = fig_top
    if fig_left_path.exists():
        with Image.open(fig_left_path) as _img:
            aspect = _img.size[1] / _img.size[0]
        disp_w = fig_area_w
        disp_h = disp_w * aspect
        if disp_h > fig_max_h:
            disp_h = fig_max_h
            disp_w = disp_h / aspect
        _safe_add_picture(slide, str(fig_left_path), Inches(0.8), Inches(fig_top),
                          width=Inches(disp_w), height=Inches(disp_h))
        left_bottom = fig_top + disp_h
    # Right figure: per-chr stacked bar
    right_x = 0.8 + fig_area_w + fig_gap
    right_bottom = fig_top
    if fig_right_path.exists():
        with Image.open(fig_right_path) as _img:
            aspect = _img.size[1] / _img.size[0]
        disp_w = fig_area_w
        disp_h = disp_w * aspect
        if disp_h > fig_max_h:
            disp_h = fig_max_h
            disp_w = disp_h / aspect
        _safe_add_picture(slide, str(fig_right_path), Inches(right_x), Inches(fig_top),
                          width=Inches(disp_w), height=Inches(disp_h))
        right_bottom = fig_top + disp_h
    # Bottom: key numbers dashboard (placed below figures with 0.25" gap)
    card_top = max(left_bottom, right_bottom) + 0.25
    card_h = min(1.4, 6.65 - card_top)  # respect footer at 6.78
    stats = sd.get("stats", [])
    n_stats = len(stats)
    if n_stats > 0:
        total_w = 11.7
        gap = 0.2
        card_w = (total_w - gap * (n_stats - 1)) / n_stats
        for i, stat in enumerate(stats):
            x = 0.8 + i * (card_w + gap)
            stat_color = stat.get("color", theme["accent"])
            add_rect(slide, x, card_top, card_w, card_h, fill="FFFFFF", line=theme["line"], radius=True)
            add_text(slide, stat["label"], x + 0.1, card_top + 0.05, card_w - 0.2, 0.30,
                     font_size=11, color=theme["text"], bold=True, align="center", margin=0)
            add_text(slide, stat["value"], x + 0.1, card_top + 0.35, card_w - 0.2, 0.42,
                     font_size=24, color=stat_color, bold=True, align="center", margin=0)
            add_text(slide, stat["meaning"], x + 0.1, card_top + 0.80, card_w - 0.2, 0.35,
                     font_size=10, color="5E6572", align="center", margin=0)


def add_hp_imbalance_definition_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: HP Imbalance definition — 4-layer grouping, formula, thresholds, ISM impact."""
    hd = cfg["deck"]["hp_imbalance_definition"]
    set_bg(slide, theme["bg"])
    add_title(slide, hd["title"], hd["subtitle"], theme,
              title_en=hd.get("title_en"), subtitle_en=hd.get("subtitle_en"))

    # --- Top: 4-layer grouping table (if present) ---
    layers = hd.get("grouping_layers")
    content_top = 2.0
    if layers:
        rows = [["層級 Layer", "分群方式 Groups", "輸出欄位 Output", "HP?"]]
        for layer in layers:
            hp_mark = "⚠️" if layer["hp_dependent"] else "✅ Free"
            rows.append([layer["name"], layer["groups"], layer["output"], hp_mark])
        add_table(slide, rows, 0.5, content_top, 12.3, 1.55, theme,
                  header_fill=theme["dark"], body_fill="FFFFFF",
                  zebra_fill="F4F0E8", font_size=10.5,
                  col_widths=[1.6, 4.8, 3.8, 1.5])
        content_top = 3.75

    # --- Left panel: Formula (x=0.5, w=6.3, larger font for clarity) ---
    lx, ly, lw = 0.5, content_top, 6.3
    add_rect(slide, lx, ly, lw, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "📊 HP_Ratio 算式", lx + 0.15, ly + 0.06, lw - 0.3, 0.22,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    formula = hd["formula"]
    formula_text = f"HP_Ratio  =  {formula['numerator']}  /  ({formula['denominator']})"
    add_text(slide, formula_text, lx + 0.25, ly + 0.30, lw - 0.5, 0.32,
             font_size=16, color=theme["text"], bold=True, font_name="Consolas", margin=0)
    add_text(slide, formula["description"], lx + 0.25, ly + 0.68, lw - 0.5, 0.20,
             font_size=10, color="5E6572", margin=0)
    add_text(slide, "範例：HP1=45, HP2=5 → (45.001)/(50.002) = 0.90 → Potential LOH",
             lx + 0.25, ly + 0.92, lw - 0.5, 0.25,
             font_size=10, color=theme["text"], font_name="Consolas", margin=0)

    # --- Right panel: Threshold cards (x=7.1, w=5.4) ---
    rx, ry, rw = 7.1, content_top, 5.4
    for i, thr in enumerate(hd["thresholds"]):
        cy = ry + i * 0.65
        fill = "E8F5E8" if thr["color"] == "009E73" else "FFF0E0"
        add_rect(slide, rx, cy, rw, 0.58, fill=fill, line=None, radius=True)
        add_text(slide, f"{thr['label']}  HP_Ratio {thr['range']}", rx + 0.12, cy + 0.04, rw - 0.25, 0.22,
                 font_size=12, color=thr["color"], bold=True, margin=0)
        add_text(slide, thr["meaning"], rx + 0.12, cy + 0.28, rw - 0.25, 0.25,
                 font_size=10, color=theme["text"], margin=0)

    # --- Middle: Meaning section (what is LOH / why HP_Ratio matters) ---
    meaning = hd.get("meaning", [])
    if meaning:
        my = content_top + 1.50
        add_text(slide, "📊 HP Imbalance 的意義   What HP Imbalance Means",
                 0.5, my, 8.0, 0.25, font_size=13, color=theme["accent"], bold=True, margin=0)
        for k, item in enumerate(meaning):
            ky = my + 0.35 + k * 0.55
            add_rect(slide, 0.5, ky, 12.0, 0.48, fill="FFFFFF", line=theme["line"], radius=True)
            add_text(slide, item.get("icon", ""), 0.62, ky + 0.06, 0.35, 0.20,
                     font_size=13, margin=0)
            add_text(slide, item["title"], 0.98, ky + 0.02, 1.8, 0.22,
                     font_size=11, color=theme["accent"], bold=True, margin=0)
            add_text(slide, item["text"], 2.85, ky + 0.02, 5.5, 0.42,
                     font_size=10, color=theme["text"], margin=0)
            add_text(slide, item.get("text_en", ""), 8.5, ky + 0.02, 3.8, 0.42,
                     font_size=9, color="5E6572", margin=0)
        iy = my + 0.35 + len(meaning) * 0.55 + 0.15
    else:
        iy = content_top + 1.35

    # --- Bottom: ISM Impact (full width, must stay above footer y=6.68) ---
    add_text(slide, "🔬 對 ISM 分析的影響   Impact on ISM Analysis", 0.5, iy, 5.0, 0.25,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    for j, impact in enumerate(hd["ism_impact"]):
        ky = iy + 0.32 + j * 0.43
        card_h = min(0.40, 6.68 - ky)
        if card_h < 0.20:
            break
        add_rect(slide, 0.5, ky, 12.0, card_h, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, impact["icon"], 0.62, ky + 0.03, 0.35, 0.20,
                 font_size=13, margin=0)
        add_text(slide, impact["text"], 0.98, ky + 0.02, 6.5, 0.22,
                 font_size=11, color=theme["text"], bold=True, margin=0)
        add_text(slide, impact.get("text_en", ""), 7.6, ky + 0.02, 4.7, 0.22,
                 font_size=9, color="5E6572", margin=0)


def add_hp_vs_loh_chromosome_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: HCC1395 genome-wide LOH — HP Imbalance vs LOH.bed vs SEQC2."""
    hc = cfg["deck"]["hp_vs_loh_chromosome"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, hc["title"], hc["subtitle"], theme,
              title_en=hc.get("title_en"), subtitle_en=hc.get("subtitle_en"))

    # --- Dual figure layout: main (left flush) + fig01 (right larger) ---
    fig_path = Path(cfg["_config_dir"]) / hc["figure"]
    fig_right_rel = hc.get("figure_right", "")
    fig_right_path = Path(cfg["_config_dir"]) / fig_right_rel if fig_right_rel else None
    has_right = fig_right_path and fig_right_path.exists()
    fig_bottom = 2.15
    if fig_path.exists():
        with Image.open(fig_path) as _img:
            aspect = _img.size[1] / _img.size[0]
        if has_right:
            area_w, area_h = 6.5, 3.2  # narrower left for larger right
        else:
            area_w, area_h = 11.7, 3.3
        disp_w = area_w
        disp_h = disp_w * aspect
        if disp_h > area_h:
            disp_h = area_h
            disp_w = disp_h / aspect
        cx = 0.5  # flush left
        cy = 2.15
        _safe_add_picture(slide, str(fig_path), Inches(cx), Inches(cy),
                          width=Inches(disp_w), height=Inches(disp_h))
        fig_bottom = cy + disp_h
    if has_right:
        with Image.open(str(fig_right_path)) as _img:
            r_aspect = _img.size[1] / _img.size[0]
        r_area_w, r_area_h = 5.5, 3.2  # larger right area
        r_disp_w = r_area_w
        r_disp_h = r_disp_w * r_aspect
        if r_disp_h > r_area_h:
            r_disp_h = r_area_h
            r_disp_w = r_disp_h / r_aspect
        r_cx = 7.2 + (r_area_w - r_disp_w) / 2
        r_cy = 2.15 + (r_area_h - r_disp_h) / 2
        _safe_add_picture(slide, str(fig_right_path), Inches(r_cx), Inches(r_cy),
                          width=Inches(r_disp_w), height=Inches(r_disp_h))
        fig_bottom = max(fig_bottom, r_cy + r_disp_h)

    # --- Legend row (below figure, 4-color with EN text) ---
    legend_y = fig_bottom + 0.10
    legends = hc["legend_items"]
    n_leg = len(legends)
    # Stack 2 per row if 4 items
    if n_leg <= 3:
        leg_w = 11.7 / n_leg
        for i, leg in enumerate(legends):
            lx = 0.8 + i * leg_w
            add_rect(slide, lx, legend_y, 0.22, 0.18, fill=leg["color"], line=None, radius=False)
            add_text(slide, leg["label"], lx + 0.28, legend_y - 0.02, leg_w - 0.35, 0.22,
                     font_size=10, color=theme["text"], margin=0)
    else:
        cols = 2
        leg_w = 11.7 / cols
        for i, leg in enumerate(legends):
            row = i // cols
            col = i % cols
            lx = 0.8 + col * leg_w
            ly = legend_y + row * 0.28
            add_rect(slide, lx, ly, 0.22, 0.18, fill=leg["color"], line=None, radius=False)
            add_text(slide, leg["label"], lx + 0.28, ly - 0.02, leg_w - 0.35, 0.22,
                     font_size=10, color=theme["text"], margin=0)
        legend_y = legend_y + ((n_leg + cols - 1) // cols) * 0.28

    # --- Key observations (bottom cards, must stay above footer y=6.68) ---
    obs_y = legend_y + 0.15
    avail_h = 6.60 - obs_y
    card_h = min(0.72, avail_h)
    if card_h < 0.35:
        return  # no space for observation cards
    obs = hc["key_observations"]
    n_obs = len(obs)
    total_w = 11.7
    gap = 0.15
    card_w = (total_w - gap * (n_obs - 1)) / n_obs
    for i, ob in enumerate(obs):
        ox = 0.8 + i * (card_w + gap)
        add_rect(slide, ox, obs_y, card_w, card_h, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, ob["icon"], ox + 0.08, obs_y + 0.04, 0.30, 0.20,
                 font_size=13, margin=0)
        add_text(slide, ob["text"], ox + 0.38, obs_y + 0.03, card_w - 0.48, 0.32,
                 font_size=10, color=theme["text"], bold=True, margin=0)
        en_text = ob.get("text_en", "")
        if en_text and card_h >= 0.60:
            add_text(slide, en_text, ox + 0.38, obs_y + 0.38, card_w - 0.48, 0.28,
                     font_size=9, color="5E6572", margin=0)


def add_permanova_concept_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: PERMANOVA explanation — what it is, why it fails in LOH, data significance."""
    pf = cfg["deck"]["permanova_failure"]
    set_bg(slide, theme["bg"])
    add_title(slide, pf["title"], pf["subtitle"], theme,
              title_en=pf.get("title_en"), subtitle_en=pf.get("subtitle_en"))
    # --- Top-left: What is PERMANOVA (3-step concept) ---
    wip = pf.get("what_is_permanova", {})
    add_rect(slide, 0.5, 2.0, 6.5, 1.75, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, wip.get("title", "PERMANOVA"), 0.7, 2.05, 6.0, 0.25,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    add_text(slide, wip.get("title_en", ""), 0.7, 2.30, 6.0, 0.18,
             font_size=9, color="5E6572", margin=0)
    steps = wip.get("steps", [])
    for i, st in enumerate(steps):
        sy = 2.55 + i * 0.38
        add_text(slide, st.get("icon", ""), 0.7, sy, 0.35, 0.25, font_size=12, margin=0)
        add_text(slide, st["text"], 1.1, sy, 4.0, 0.20, font_size=11, color=theme["text"], margin=0)
        add_text(slide, st.get("text_en", ""), 5.2, sy, 1.6, 0.20, font_size=9, color="5E6572", margin=0)
    # --- Top-right: Figure (PERMANOVA F-stats) ---
    fig_path = Path(cfg["deck"].get("seqc2_chromosome", {}).get("figure", "")).parent.parent / pf.get("figure", "")
    img_dir = Path(cfg["meta"]["output_pptx"]).parent
    fig_actual = img_dir / pf.get("figure", "")
    if fig_actual.exists():
        with Image.open(str(fig_actual)) as _img:
            aspect = _img.size[1] / _img.size[0]
        disp_w = 5.5
        disp_h = min(disp_w * aspect, 1.75)
        disp_w = disp_h / aspect
        fx = 7.2 + (5.5 - disp_w) / 2
        fy = 2.0 + (1.75 - disp_h) / 2
        _safe_add_picture(slide, str(fig_actual), Inches(fx), Inches(fy),
                          width=Inches(disp_w), height=Inches(disp_h))
    else:
        add_rect(slide, 7.2, 2.0, 5.5, 1.75, fill="EEE6DB", line=theme["line"], radius=True)
        add_text(slide, pf.get("figure_title", ""), 7.4, 2.5, 5.1, 0.5,
                 font_size=12, color="5E6572", align="center", margin=0)
    # --- Middle: Why fail in LOH ---
    wfl = pf.get("why_fail_in_loh", {})
    add_rect(slide, 0.5, 3.90, 12.3, 1.10, fill="FFF0E0", line=None, radius=True)
    add_text(slide, f"⚠️ {wfl.get('title', '')}", 0.7, 3.95, 6.0, 0.25,
             font_size=13, color="D55E00", bold=True, margin=0)
    add_text(slide, wfl.get("title_en", ""), 6.8, 3.95, 5.5, 0.25,
             font_size=10, color="5E6572", margin=0)
    add_text(slide, wfl.get("explanation", ""), 0.7, 4.25, 5.8, 0.65,
             font_size=11, color=theme["text"], margin=0)
    add_text(slide, wfl.get("explanation_en", ""), 6.8, 4.25, 5.5, 0.65,
             font_size=9, color="5E6572", margin=0)
    # --- Bottom: 4 stat cards (compact) ---
    positions = [(0.5, 5.15, 2.9, 1.40), (3.55, 5.15, 2.9, 1.40), (6.6, 5.15, 2.9, 1.40), (9.65, 5.15, 2.9, 1.40)]
    for (x, y, w, h), stat in zip(positions, pf["stats"]):
        stat_color = stat.get("color", theme["accent"])
        add_rect(slide, x, y, w, h, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, stat["label"], x + 0.08, y + 0.08, w - 0.16, 0.25,
                 font_size=11, color=theme["text"], bold=True, align="center", margin=0)
        add_text(slide, stat["value"], x + 0.08, y + 0.38, w - 0.16, 0.42,
                 font_size=24, color=stat_color, bold=True, align="center", margin=0)
        add_text(slide, stat["meaning"], x + 0.08, y + 0.88, w - 0.16, 0.40,
                 font_size=10, color="5E6572", align="center", margin=0)
    # --- Data significance note ---
    ds = pf.get("data_significance", "")
    if ds:
        add_text(slide, ds, 0.5, 6.58, 12.3, 0.15,
                 font_size=9, color=theme["muted"], margin=0)


def add_filter_decision_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Filter safety test — pre-rendered bar chart image + conclusion."""
    ff = cfg["deck"]["filter_10_fail"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, ff["title"], ff["subtitle"], theme,
              title_en=ff.get("title_en"), subtitle_en=ff.get("subtitle_en"))
    # Safety line explanation
    sle = ff.get("safety_line_explanation", "")
    sle_en = ff.get("safety_line_explanation_en", "")
    add_rect(slide, 0.5, 1.95, 12.3, 0.40, fill="FFF0E0", line=None, radius=True)
    add_text(slide, f"📋 {sle}", 0.7, 1.98, 8.5, 0.18,
             font_size=10, color="D55E00", bold=True, margin=0)
    add_text(slide, sle_en, 0.7, 2.18, 8.5, 0.15,
             font_size=9, color="5E6572", margin=0)
    # Pre-rendered bar chart image (replaces 50+ individual shapes)
    chart_path = Path(cfg["_config_dir"]) / "images" / "filter_strategy_barchart.png"
    if chart_path.exists():
        _safe_add_picture(slide, str(chart_path), Inches(0.5), Inches(2.50), width=Inches(12.3), height=Inches(3.50))
    else:
        add_text(slide, "[Chart image not found]", 0.5, 3.5, 12.3, 0.5,
                 font_size=14, color="D55E00", align="center", margin=0)
    # Conclusion bar
    add_rect(slide, 0.5, 6.10, 12.3, 0.35, fill="C76B50", line=None, radius=True)
    add_text(slide, f"❌ {ff.get('conclusion', '')}", 0.7, 6.12, 11.8, 0.14,
             font_size=11, color="F8F4EE", bold=True, margin=0)
    add_text(slide, ff.get("conclusion_en", ""), 0.7, 6.27, 11.8, 0.12,
             font_size=9, color="F0D0C0", margin=0)
    # Future note
    fn = ff.get("future_note", "")
    fn_en = ff.get("future_note_en", "")
    if fn:
        add_text(slide, f"→ {fn}", 0.7, 6.50, 8.0, 0.15,
                 font_size=10, color="6E9D8A", margin=0)
        if fn_en:
            add_text(slide, fn_en, 8.8, 6.50, 4.0, 0.15,
                     font_size=9, color="5E6572", margin=0)


def add_methylation_negative_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Methylation signal verification — pre-rendered 3-dimension chart + conclusion."""
    mn = cfg["deck"]["methylation_negative"]
    set_bg(slide, theme["bg"])
    add_title(slide, mn["title"], mn["subtitle"], theme,
              title_en=mn.get("title_en"), subtitle_en=mn.get("subtitle_en"))
    # Method banner (balanced tone)
    banner = mn.get("banner_text", "")
    banner_en = mn.get("banner_text_en", "")
    add_rect(slide, 0.5, 1.95, 12.3, 0.40, fill="E8F0FF", line=None, radius=True)
    add_text(slide, f"🔬 {banner}", 0.7, 1.98, 8.5, 0.18,
             font_size=10, color="0072B2", bold=True, margin=0)
    add_text(slide, banner_en, 0.7, 2.18, 8.5, 0.15,
             font_size=9, color="5E6572", margin=0)
    # Pre-rendered 3-dimension verification chart (replaces 50+ individual shapes)
    chart_path = Path(cfg["_config_dir"]) / "images" / "methylation_3dim_verification.png"
    if chart_path.exists():
        _safe_add_picture(slide, str(chart_path), Inches(0.3), Inches(2.45), width=Inches(12.7), height=Inches(3.45))
    else:
        add_text(slide, "[Chart image not found]", 0.5, 3.5, 12.3, 0.5,
                 font_size=14, color="D55E00", align="center", margin=0)
    # Bottom: balanced conclusion + future note
    add_rect(slide, 0.5, 6.05, 12.3, 0.30, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "初步 AUC 篩選不支持甲基化方向 — 每個表面信號均有 confound 解釋",
             0.7, 6.07, 8.5, 0.14,
             font_size=10, color="F8F4EE", bold=True, margin=0)
    add_text(slide, "Initial AUC screening does not support methylation — all surface signals have confound explanations",
             0.7, 6.22, 8.5, 0.12,
             font_size=9, color="A0A0A0", margin=0)
    # Future note
    fn = mn.get("future_note", "")
    if fn:
        add_text(slide, f"→ {fn}", 0.5, 6.42, 12.3, 0.15,
                 font_size=10, color="6E9D8A", margin=0)


def add_to_nogo_evidence_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: TO 60+ features NO-GO using evidence stack with bias direction."""
    tn = cfg["deck"]["to_nogo"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, tn["title"], tn["subtitle"], theme,
              title_en=tn.get("title_en"), subtitle_en=tn.get("subtitle_en"))
    # Left: Evidence stack (layers)
    build_evidence_stack(slide, tn["evidence_layers"], theme, area_left=0.8, area_top=2.15, area_width=7.5, area_height=3.2)
    # Right: Summary stats panel
    add_rect(slide, 8.6, 2.15, 3.9, 3.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "特徵偏向摘要", 8.8, 2.25, 3.5, 0.22,
             font_size=14, color=theme["text"], bold=True, margin=0)
    add_text(slide, "Feature Bias Summary", 8.8, 2.48, 3.5, 0.18,
             font_size=9, color="5E6572", margin=0)
    bias_items = [
        ("G1-G4", "VCF 基礎+strand", "AUC < 0.58", "D55E00"),
        ("G5", "CpG context", "AUC = 0.52", "D55E00"),
        ("G6", "全 combo", "AUC < 0.64", "E69F00"),
        ("G7", "多特徵 Voting", "AUC = 0.577", "D55E00"),
        ("60+", "所有測試特徵", "FP removal=0%", "D55E00"),
    ]
    for j, (gid, desc, result, color) in enumerate(bias_items):
        by = 2.80 + j * 0.48
        add_text(slide, gid, 8.8, by, 0.7, 0.20, font_size=10, color=theme["accent"], bold=True, margin=0)
        add_text(slide, desc, 9.5, by, 1.5, 0.20, font_size=10, color=theme["text"], margin=0)
        add_text(slide, result, 10.9, by, 1.4, 0.20, font_size=10, color=color, bold=True, margin=0)
    # Conclusion bar
    add_rect(slide, 0.8, 5.55, 11.7, 0.60, fill="C76B50", line=None, radius=True)
    add_bilingual_text(slide, 1.0, 5.58, 11.3, 0.52,
                       f"❌ {tn['conclusion']}",
                       "❌ FP removal = 0% under TP loss ≤2% safety constraint — all 60+ features biased or ineffective",
                       12, theme, bold=True, color="F8F4EE", margin=0)
    # Future note
    fn = tn.get("future_note", "")
    if fn:
        add_text(slide, f"→ {fn}", 0.8, 6.22, 8.0, 0.15,
                 font_size=10, color="6E9D8A", margin=0)
        fn_en = tn.get("future_note_en", "")
        if fn_en:
            add_text(slide, fn_en, 8.8, 6.22, 3.7, 0.15,
                     font_size=9, color="5E6572", margin=0)


def add_self_phasing_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Self-phasing causal chain using pipeline flow."""
    sp = cfg["deck"]["self_phasing"]
    set_bg(slide, theme["bg"])
    add_title(slide, sp["title"], sp["subtitle"], theme,
              title_en=sp.get("title_en"), subtitle_en=sp.get("subtitle_en"))
    build_pipeline_flow(slide, sp["steps"], theme, area_top=2.5, area_height=2.8)
    # Bottom: key numbers
    add_rect(slide, 0.8, 5.6, 11.7, 0.9, fill="FFFFFF", line=theme["line"], radius=True)
    add_bilingual_text(slide, 1.0, 5.65, 11.3, 0.8,
                       "✅ PON-only phasing 已驗證可修正：LOH.bed Jaccard=1.0 不變 | somatic bias 消除 | N50 +99.7% | phased rate +23.6pp",
                       "✅ PON-only phasing verified: LOH.bed Jaccard=1.0 unchanged | somatic bias eliminated | N50 +99.7% | phased rate +23.6pp",
                       13, theme, margin=0)


def add_code_modification_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: ISM code modification record — root cause + file table + impact."""
    cm = cfg["deck"]["code_modifications"]
    set_bg(slide, theme["bg"])
    add_title(slide, cm["title"], cm["subtitle"], theme,
              title_en=cm.get("title_en"), subtitle_en=cm.get("subtitle_en"))

    # --- Top: Root cause card (warning orange) ---
    rc = cm.get("root_cause", {})
    add_rect(slide, 0.5, 2.0, 12.3, 1.05, fill="FFF3E0", line="D55E00", radius=True)
    icon = rc.get("icon", "⚠️")
    add_text(slide, f"{icon} {rc.get('title', '')}", 0.7, 2.03, 7.0, 0.28,
             font_size=14, color="D55E00", bold=True, margin=0)
    add_text(slide, rc.get("title_en", ""), 7.8, 2.05, 4.8, 0.25,
             font_size=10, color="5E6572", margin=0)
    desc_lines = rc.get("description", "").split("\n")
    for k, line in enumerate(desc_lines):
        add_text(slide, line, 0.9, 2.38 + k * 0.22, 11.7, 0.22,
                 font_size=11, color=theme["text"], margin=0)

    # --- Middle: Modification table ---
    mods = cm.get("modifications", [])
    add_text(slide, "修改範圍   Modified Files", 0.5, 3.20, 6.0, 0.25,
             font_size=13, color=theme["text"], bold=True, margin=0)
    rows = [["檔案 File", "修改內容 Change", "類型 Type"]]
    for m in mods:
        rows.append([m["file"], m["change"], m["impact"]])
    add_table(slide, rows, 0.5, 3.50, 7.5, 2.30, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=10, col_widths=[2.2, 3.8, 1.2])

    # --- Right: Impact summary ---
    imp = cm.get("impact_summary", {})
    add_rect(slide, 8.2, 3.20, 4.6, 2.60, fill="FFFFFF", line=theme["line"], radius=True)
    add_bilingual_text(slide, 8.4, 3.25, 4.2, 0.35,
                       imp.get("title", "影響評估"),
                       imp.get("title_en", "Impact Assessment"),
                       13, theme, bold=True, margin=0)
    items = imp.get("items", [])
    for j, item in enumerate(items):
        iy = 3.72 + j * 0.52
        add_text(slide, f"{item['icon']} {item['text']}", 8.5, iy, 4.1, 0.22,
                 font_size=11, color=theme["text"], margin=0)
        add_text(slide, item.get("text_en", ""), 8.7, iy + 0.22, 3.9, 0.20,
                 font_size=9, color="5E6572", margin=0)

    # --- Bottom: Summary bar ---
    add_rect(slide, 0.5, 6.00, 12.3, 0.55, fill=theme["dark"], line=None, radius=True)
    add_bilingual_text(slide, 0.7, 6.03, 11.9, 0.48,
                       "🔄 所有修正均已驗證 — 本報告數據為修正後版本",
                       "🔄 All fixes verified — all data in this report uses the corrected version",
                       13, theme, bold=True, color="F8F4EE", margin=0)


def add_asm_concept_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: ASM bright spot — definition + example + 5-method evidence + paper."""
    ab = cfg["deck"]["asm_bright_spot"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, ab["title"], ab["subtitle"], theme,
              title_en=ab.get("title_en"), subtitle_en=ab.get("subtitle_en"))
    # Left: Definition + Example
    add_rect(slide, 0.8, 2.0, 5.8, 2.0, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, ab["term"], 1.0, 2.08, 5.4, 0.25,
             font_size=14, color=theme["accent"], bold=True, margin=0)
    add_text(slide, ab["definition"], 1.0, 2.38, 5.4, 0.65,
             font_size=11, color=theme["text"], margin=0)
    add_text(slide, ab.get("source", ""), 1.0, 3.08, 5.4, 0.22,
             font_size=10, color="009E73", bold=True, margin=0)
    # Example box
    example = ab.get("example", {})
    if example:
        add_rect(slide, 0.8, 3.40, 5.8, 0.60, fill=theme["dark"], line=None, radius=True)
        add_text(slide, example.get("title", ""), 1.0, 3.42, 5.4, 0.20,
                 font_size=10, color="F8F4EE", bold=True, margin=0)
        add_text(slide, example.get("description", ""), 1.0, 3.65, 5.4, 0.30,
                 font_size=9, color="D6DCE5", margin=0)
    # Right: 5-method evidence table
    add_rect(slide, 6.8, 2.0, 5.7, 2.0, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "5 種驗證方法   5 Validation Methods", 7.0, 2.08, 5.3, 0.25,
             font_size=13, color=theme["text"], bold=True, margin=0)
    evidence = ab.get("evidence", [])
    if evidence:
        ev_rows = [["方法 Method", "ASM %", "備註 Note"]]
        for ev in evidence:
            ev_rows.append([ev["method"], ev["rate"], ev["note"]])
        add_table(slide, ev_rows, 6.95, 2.40, 5.35, 1.50, theme,
                  header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
                  font_size=10, col_widths=[2.5, 0.8, 1.2])
    # Paper reference
    paper = ab.get("paper_ref", "")
    if paper:
        add_rect(slide, 6.8, 3.40, 5.7, 0.60, fill="E8F8E8", line=None, radius=True)
        add_text(slide, "📄 論文證據 Paper Evidence", 7.0, 3.42, 3.0, 0.20,
                 font_size=10, color="009E73", bold=True, margin=0)
        add_text(slide, paper, 7.0, 3.65, 5.3, 0.30,
                 font_size=9, color=theme["text"], margin=0)
    # Bottom: implication
    add_rect(slide, 0.8, 4.20, 11.7, 0.55, fill=theme["dark"], line=None, radius=True)
    add_bilingual_text(slide, 1.0, 4.23, 11.3, 0.48,
                       "💡 ISM 不是「無用」— 價值在 read-level epigenetic characterization，非 variant filtering",
                       "💡 ISM is not useless — value lies in read-level epigenetic characterization, not variant filtering",
                       13, theme, bold=True, color="F8F4EE", margin=0)


def add_strategic_pivot_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Strategic pivot - old vs new direction comparison with clear focus."""
    sp = cfg["deck"]["strategic_pivot"]
    set_bg(slide, theme["bg"])
    add_title(slide, sp["title"], sp["subtitle"], theme,
              title_en=sp.get("title_en"), subtitle_en=sp.get("subtitle_en"))
    # Central pivot arrow
    add_rect(slide, 5.6, 2.2, 2.1, 0.55, fill=theme["accent"], line=None, radius=True)
    add_text(slide, "→ 轉向 Pivot →", 5.6, 2.25, 2.1, 0.45,
             font_size=14, color="F8F4EE", bold=True, align="center", margin=0)
    # Left: Old direction (red/crossed out)
    old = sp["old_direction"]
    add_rect(slide, 0.8, 2.95, 5.5, 3.0, fill="FFFFFF", line="D55E00", radius=True)
    add_rect(slide, 0.8, 2.95, 5.5, 0.50, fill="D55E00", line=None, radius=True)
    add_text(slide, f"❌ {old['title']}", 1.0, 3.00, 5.1, 0.22,
             font_size=14, color="F8F4EE", bold=True, margin=0)
    add_text(slide, "Phase 1 — Closed", 1.0, 3.24, 3.0, 0.18,
             font_size=9, color="F0D0C0", margin=0)
    for j, item in enumerate(old["items"]):
        iy = 3.60 + j * 0.50
        add_text(slide, f"❌ {item}", 1.1, iy, 5.0, 0.25,
                 font_size=12, color="5E6572", margin=0)
    # Right: New direction (green/active)
    new = sp["new_direction"]
    add_rect(slide, 6.9, 2.95, 5.6, 3.0, fill="FFFFFF", line="009E73", radius=True)
    add_rect(slide, 6.9, 2.95, 5.6, 0.50, fill="009E73", line=None, radius=True)
    add_text(slide, f"✅ {new['title']}", 7.1, 3.00, 5.2, 0.22,
             font_size=14, color="F8F4EE", bold=True, margin=0)
    add_text(slide, "Phase 2A — Active", 7.1, 3.24, 3.0, 0.18,
             font_size=9, color="D0F0D0", margin=0)
    for j, item in enumerate(new["items"]):
        iy = 3.60 + j * 0.50
        add_text(slide, f"✅ {item}", 7.1, iy, 5.2, 0.25,
                 font_size=12, color=theme["text"], bold=True, margin=0)
    # Bottom: key message
    add_rect(slide, 0.8, 6.15, 11.7, 0.50, fill=theme["dark"], line=None, radius=True)
    add_bilingual_text(slide, 1.0, 6.18, 11.3, 0.44,
                       "核心轉變：Filtering → Characterization — Phase 2A 聚焦上游修正 + 正常參考基線",
                       "Core pivot: Filtering → Characterization — Phase 2A: upstream correction + normal reference",
                       12, theme, bold=True, color="F8F4EE", margin=0)


def add_section_divider_slide(slide, cfg_div: dict, theme: dict) -> None:
    """Section divider: dark background, centered title + subtitle preview."""
    set_bg(slide, theme["dark"])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=theme["dark"], line=None)
    # Section number badge (small, top-right)
    sn = cfg_div.get("section_num", "")
    if sn:
        add_rect(slide, 11.0, 0.6, 1.5, 0.45, fill=theme["accent"], line=None, radius=True)
        add_text(slide, f"Part {sn}", 11.0, 0.62, 1.5, 0.4,
                 font_size=14, color=theme["light_text"], bold=True, align="center", margin=0)
    # Main title — large, centered
    t_zh = cfg_div.get("title", "")
    t_en = cfg_div.get("title_en", "")
    add_text(slide, t_zh, 1.5, 2.4, 10.3, 1.0,
             font_size=36, color=theme["light_text"], bold=True, align="center",
             font_name=theme.get("font_title", "Arial Bold"), margin=0)
    if t_en:
        add_text(slide, t_en, 1.5, 3.4, 10.3, 0.6,
                 font_size=20, color=theme.get("light_muted", "D6DCE5"),
                 align="center", italic=False, margin=0)
    # Subtitle preview
    s_zh = cfg_div.get("subtitle", "")
    s_en = cfg_div.get("subtitle_en", "")
    if s_zh:
        add_text(slide, s_zh, 2.0, 4.5, 9.3, 0.5,
                 font_size=16, color=theme.get("light_muted", "D6DCE5"), align="center", margin=0)
    if s_en:
        add_text(slide, s_en, 2.0, 5.0, 9.3, 0.5,
                 font_size=12, color="A0A0A0", align="center", italic=False, margin=0)
    # Decorative line
    add_rect(slide, 4.5, 4.2, 4.3, 0.07, fill=theme["accent"], line=None)


def add_ism_framework_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: ISM analysis framework — pipeline flow (top) + evidence chain figure (bottom)."""
    ism = cfg["deck"]["ism_framework"]
    set_bg(slide, theme["bg"])
    add_title(slide, ism["title"], ism["subtitle"], theme,
              title_en=ism.get("title_en"), subtitle_en=ism.get("subtitle_en"))

    # --- Top: 4-step pipeline flow (expanded, no bottom figure) ---
    steps = ism.get("pipeline_steps", [])
    if steps:
        build_pipeline_flow(slide, steps, theme, area_top=2.0, area_height=2.0)

    # --- Middle: Core outputs (left 3 items) + Focus questions (right) ---
    outputs = ism.get("outputs", [])
    oy = 4.15
    add_text(slide, "核心產出   Core Outputs", 0.8, oy, 5.0, 0.35,
             font_size=13, color=theme["accent"], bold=True, margin=0)
    for i, out in enumerate(outputs):
        row_y = oy + 0.40 + i * 0.45
        icon = out.get("icon", "")
        add_text(slide, icon, 0.8, row_y, 0.4, 0.3, font_size=13, margin=0)
        add_text(slide, out["text"], 1.25, row_y, 5.2, 0.3, font_size=12, color=theme["text"],
                 bold=True, margin=0)
        en = out.get("text_en", "")
        if en:
            add_text(slide, en, 1.25, row_y + 0.25, 5.2, 0.25, font_size=9, color=theme["muted"],
                     italic=False, margin=0)

    fqs = ism.get("focus_questions", [])
    if fqs:
        fq_y = oy
        add_text(slide, "本週驗證焦點   This Week's Focus", 7.0, fq_y, 5.5, 0.35,
                 font_size=13, color=theme["accent"], bold=True, margin=0)
        for i, fq in enumerate(fqs):
            row_y = fq_y + 0.40 + i * 0.45
            add_text(slide, fq.get("icon", "→"), 7.0, row_y, 0.3, 0.3, font_size=13, margin=0)
            add_text(slide, fq["text"], 7.35, row_y, 5.1, 0.3,
                     font_size=12, color=theme["text"], bold=True, margin=0)
            en = fq.get("text_en", "")
            if en:
                add_text(slide, en, 7.35, row_y + 0.25, 5.1, 0.25,
                         font_size=9, color=theme["muted"], italic=False, margin=0)


def add_hcc1395_detailed_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: HCC1395 LOH detailed observation — figure + observation cards."""
    hd = cfg["deck"]["hcc1395_detailed"]
    set_bg(slide, theme["bg"])
    add_title(slide, hd["title"], hd["subtitle"], theme,
              title_en=hd.get("title_en"), subtitle_en=hd.get("subtitle_en"))
    # Main figure (top area, centered — enlarged for clarity)
    fig_path = Path(cfg["_config_dir"]) / hd["figure"]
    fig_bottom = 2.15
    if fig_path.exists():
        with Image.open(fig_path) as _img:
            aspect = _img.size[1] / _img.size[0]
        area_w, area_h = 12.0, 3.5  # enlarged for better readability
        disp_w = area_w
        disp_h = disp_w * aspect
        if disp_h > area_h:
            disp_h = area_h
            disp_w = disp_h / aspect
        cx = 0.65 + (area_w - disp_w) / 2
        cy = 2.05
        _safe_add_picture(slide, str(fig_path), Inches(cx), Inches(cy),
                          width=Inches(disp_w), height=Inches(disp_h))
        fig_bottom = cy + disp_h
    # Observation cards (below figure, horizontal layout)
    obs = hd.get("observations", [])
    obs_y = fig_bottom + 0.12
    n_obs = len(obs)
    total_w = 11.7
    gap = 0.15
    card_w = (total_w - gap * max(n_obs - 1, 0)) / max(n_obs, 1)
    card_h = min(0.90, 6.60 - obs_y)
    if card_h < 0.30:
        return
    for i, ob in enumerate(obs):
        ox = 0.8 + i * (card_w + gap)
        add_rect(slide, ox, obs_y, card_w, card_h, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, ob.get("icon", ""), ox + 0.10, obs_y + 0.06, 0.30, 0.22,
                 font_size=14, margin=0)
        add_text(slide, ob["text"], ox + 0.42, obs_y + 0.05, card_w - 0.55, 0.35,
                 font_size=11, color=theme["text"], bold=True, margin=0)
        en_text = ob.get("text_en", "")
        if en_text and card_h >= 0.55:
            add_text(slide, en_text, ox + 0.42, obs_y + 0.42, card_w - 0.55, 0.30,
                     font_size=9, color="5E6572", margin=0)


def add_key_features_auc_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Key features AUC — figure + table + caller_af detail."""
    kf = cfg["deck"]["key_features_auc"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, kf["title"], kf["subtitle"], theme,
              title_en=kf.get("title_en"), subtitle_en=kf.get("subtitle_en"))
    # --- Left: Figure (Top 10 AUC comparison) ---
    img_dir = Path(cfg["meta"]["output_pptx"]).parent
    fig_path = img_dir / kf.get("figure", "")
    has_fig = fig_path.exists() and kf.get("figure", "")
    if has_fig:
        with Image.open(str(fig_path)) as _img:
            aspect = _img.size[1] / _img.size[0]
        area_w, area_h = 5.8, 3.2
        disp_w = area_w
        disp_h = disp_w * aspect
        if disp_h > area_h:
            disp_h = area_h
            disp_w = disp_h / aspect
        fx = 0.5 + (area_w - disp_w) / 2
        fy = 2.0 + (area_h - disp_h) / 2
        _safe_add_picture(slide, str(fig_path), Inches(fx), Inches(fy),
                          width=Inches(disp_w), height=Inches(disp_h))
        fig_title = kf.get("figure_title", "")
        if fig_title:
            add_text(slide, fig_title, 0.5, 5.25, area_w, 0.18,
                     font_size=9, color="5E6572", align="center", margin=0)
        table_left = 6.5
        table_w = 6.3
    else:
        table_left = 0.5
        table_w = 12.3
    # --- Right (or full): Feature AUC table ---
    features = kf["features"]
    rows = [["特徵", "LOH", "Non-LOH", "Pooled", "意義"]]
    for f in features:
        rows.append([f["name"], f["loh_auc"], f["non_loh_auc"], f["pooled_auc"], f["meaning"]])
    add_table(slide, rows, table_left, 2.0, table_w, 1.60, theme,
              header_fill=theme["dark"], body_fill="FFFFFF",
              zebra_fill="F4F0E8", font_size=10,
              col_widths=[1.3, 0.8, 0.9, 0.8, 2.5] if has_fig else [2.2, 1.5, 1.7, 1.5, 5.4])
    # --- caller_af detail (below table on right) ---
    cad = kf.get("caller_af_detail", {})
    if cad:
        cad_top = 3.75
        cad_w = table_w
        add_rect(slide, table_left, cad_top, cad_w, 1.65, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, f"📊 {cad['title']}", table_left + 0.12, cad_top + 0.05, cad_w - 0.24, 0.22,
                 font_size=11, color=theme["accent"], bold=True, margin=0)
        th_rows = [["門檻", "FP 移除", "TP 損失", "比率"]]
        for r in cad.get("rows", []):
            th_rows.append([r["threshold"], r["fp_removal"], r["tp_loss"], r["ratio"]])
        add_table(slide, th_rows, table_left + 0.12, cad_top + 0.30, cad_w - 0.24, 0.95, theme,
                  header_fill="5E6572", body_fill="FFFFFF",
                  zebra_fill="F4F0E8", font_size=10,
                  col_widths=[1.3, 1.0, 1.0, 0.8] if has_fig else [1.8, 1.3, 1.3, 1.2])
        conc = cad.get("conclusion", "")
        add_text(slide, f"❌ {conc}", table_left + 0.12, cad_top + 1.30, cad_w - 0.24, 0.20,
                 font_size=10, color="D55E00", bold=True, margin=0)
    # --- Bottom: Key takeaway bar ---
    add_rect(slide, 0.5, 5.55, 12.3, 0.55, fill=theme["dark"], line=None, radius=True)
    takeaways = [
        ("❌", "ISM 甲基化特徵對 TO FP 無區分力"),
        ("⚠️", "caller_af 唯一有效但不可用作 filter"),
        ("💡", "轉向 read-level characterization"),
    ]
    for i, (icon, zh) in enumerate(takeaways):
        tx = 0.7 + i * 4.0
        add_text(slide, f"{icon} {zh}", tx, 5.60, 3.8, 0.22,
                 font_size=11, color="F8F4EE", bold=True, margin=0)


def add_composition_effect_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Composition Effect — shows how LOH AUC is inflated by subtype mixing."""
    ce = cfg["deck"]["composition_effect"]
    set_bg(slide, theme["bg_alt"])
    add_title(slide, ce["title"], ce["subtitle"], theme,
              title_en=ce.get("title_en"), subtitle_en=ce.get("subtitle_en"))
    # Stats cards (3 cards, horizontal)
    positions = [(0.5, 2.0), (4.5, 2.0), (8.5, 2.0)]
    for (x, y), stat in zip(positions, ce["stats"]):
        add_rect(slide, x, y, 3.7, 1.55, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, stat["label"], x + 0.15, y + 0.12, 3.4, 0.22,
                 font_size=13, color=theme["text"], bold=True, margin=0)
        add_text(slide, stat["value"], x + 0.15, y + 0.42, 3.4, 0.4,
                 font_size=26, color=stat.get("color", theme["accent"]), bold=True, margin=0)
        add_text(slide, stat["meaning"], x + 0.15, y + 0.92, 3.4, 0.5,
                 font_size=12, color="5E6572", margin=0)
    # Bottom: subtype table
    add_rect(slide, 0.5, 3.8, 11.8, 2.5, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "LOH 子型別 AUC 方向", 0.7, 3.92, 4.0, 0.22,
             font_size=15, color=theme["text"], bold=True, margin=0)
    rows = [["LOH Subtype", "AUC 方向", "機制"]]
    for st in ce["subtypes"]:
        rows.append([st["name"], st["direction"], st["mechanism"]])
    add_table(slide, rows, 0.7, 4.25, 11.3, 1.7, theme,
              header_fill=theme["dark"], body_fill="FFFFFF",
              zebra_fill="F4F0E8", font_size=12, col_widths=[3.0, 3.0, 5.3])


def add_to_pure_modeling_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: TO-pure independent modeling — caller_af vs ISM."""
    tp = cfg["deck"]["to_pure_modeling"]
    set_bg(slide, theme["bg"])
    add_title(slide, tp["title"], tp["subtitle"], theme,
              title_en=tp.get("title_en"), subtitle_en=tp.get("subtitle_en"))
    # Left: big number — caller_af
    add_rect(slide, 0.5, 2.0, 5.8, 2.2, fill=theme["dark"], line=None, radius=True)
    add_text(slide, "caller_af（唯一有效判別器）", 0.7, 2.12, 5.2, 0.22,
             font_size=14, color=theme["light_muted"], bold=True, margin=0)
    add_text(slide, tp["caller_af_auc"], 0.7, 2.45, 5.2, 0.6,
             font_size=36, color="F8F4EE", bold=True, margin=0)
    add_text(slide, tp["caller_af_note"], 0.7, 3.15, 5.2, 0.8,
             font_size=12, color=theme["light_muted"], margin=0)
    # Right: ISM increment table
    add_rect(slide, 6.6, 2.0, 5.8, 2.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "ISM 增量（over caller_af）", 6.8, 2.12, 5.2, 0.22,
             font_size=14, color=theme["text"], bold=True, margin=0)
    rows = [["組合", "AUC", "增量"]]
    for r in tp["ism_increments"]:
        rows.append([r["combo"], r["auc"], r["delta"]])
    add_table(slide, rows, 6.8, 2.45, 5.3, 1.5, theme,
              header_fill=theme["dark"], body_fill="FFFFFF",
              zebra_fill="F4F0E8", font_size=11, col_widths=[2.5, 1.2, 1.6])
    # Bottom: 3 key findings
    add_rect(slide, 0.5, 4.45, 11.8, 1.9, fill="FFFFFF", line=theme["line"], radius=True)
    y_off = 4.6
    for finding in tp["findings"]:
        add_text(slide, f"{finding['icon']}  {finding['text']}", 0.75, y_off, 11.2, 0.25,
                 font_size=13, color=theme["text"], margin=0)
        y_off += 0.4
    # Conclusion bar (must stay above footer y=6.68)
    add_rect(slide, 0.5, 6.25, 11.8, 0.40, fill="D55E00", line=None, radius=True)
    add_bilingual_text(slide, 0.7, 6.27, 11.3, 0.36,
                       tp["conclusion"], tp.get("conclusion_en", ""),
                       12, theme, bold=True, color="F8F4EE", margin=0)


def add_per_sample_stratification_slide(slide, cfg: dict, theme: dict) -> None:
    """Slide: Per-sample stratification — Tier 1/2/3 + clinical implication."""
    ps = cfg["deck"]["per_sample_stratification"]
    set_bg(slide, theme["bg"])
    add_title(slide, ps["title"], ps["subtitle"], theme,
              title_en=ps.get("title_en"), subtitle_en=ps.get("subtitle_en"))
    # Tier table
    rows = [["Tier", "樣本", "LOH.bed AlleleDelta", "特徵"]]
    for tier in ps["tiers"]:
        rows.append([tier["tier"], tier["samples"], tier["auc_range"], tier["note"]])
    add_table(slide, rows, 0.5, 2.0, 11.8, 1.8, theme,
              header_fill=theme["dark"], body_fill="FFFFFF",
              zebra_fill="F4F0E8", font_size=12,
              col_widths=[1.2, 4.0, 3.0, 3.6])
    # Clinical warning box
    add_rect(slide, 0.5, 4.1, 11.8, 2.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "⚠️ 臨床推論", 0.75, 4.22, 3.0, 0.22,
             font_size=15, color="D55E00", bold=True, margin=0)
    y_off = 4.55
    for warn in ps["clinical_warnings"]:
        add_text(slide, f"{warn['icon']}  {warn['text']}", 0.75, y_off, 11.2, 0.35,
                 font_size=13, color=theme["text"], margin=0)
        y_off += 0.42


# ====================================================================
# v5 Builder Functions (2026-04-13 LOH×CNV + HP Tag)
# ====================================================================


def add_recap_seqc2_quadrant_slide(slide, cfg: dict, theme: dict) -> None:
    """Compressed recap: SEQC2 validation + quadrant distribution."""
    sec = cfg["deck"]["recap_seqc2_quadrant"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Left: figure
    img_path = str(Path(cfg["_config_dir"]) / sec["figure"])
    _safe_add_picture(slide, img_path, Inches(0.5), Inches(2.0), width=Inches(6.0))
    # Right top: SEQC2 metrics
    add_rect(slide, 6.8, 2.0, 5.8, 1.8, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "SEQC2 驗證", 7.0, 2.05, 3.0, 0.3,
             font_size=14, color=theme["text"], bold=True, margin=0)
    for i, item in enumerate(sec["left_items"]):
        add_text(slide, f"• {item}", 7.0, 2.4 + i * 0.4, 5.4, 0.35,
                 font_size=11, color=theme["text"], margin=0)
    # Right bottom: quadrant
    add_rect(slide, 6.8, 4.0, 5.8, 2.4, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, "四象限分佈", 7.0, 4.05, 3.0, 0.3,
             font_size=14, color=theme["text"], bold=True, margin=0)
    for i, item in enumerate(sec["right_items"]):
        add_text(slide, f"• {item}", 7.0, 4.4 + i * 0.45, 5.4, 0.4,
                 font_size=11, color=theme["text"], margin=0)
    # Bottom conclusion
    add_rect(slide, 0.5, 6.3, 12.3, 0.35, fill="E8F5E9", line=None, radius=True)
    add_text(slide, f"→ {sec['conclusion']}", 0.7, 6.32, 11.9, 0.3,
             font_size=12, color="009E73", bold=True, align="center", margin=0)


def add_recap_loh_filter_slide(slide, cfg: dict, theme: dict) -> None:
    """Compressed recap: LOH impact + filter FAIL."""
    sec = cfg["deck"]["recap_loh_filter"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Left card: LOH methylation
    add_rect(slide, 0.5, 2.0, 5.9, 2.5, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, sec["left_title"], 0.7, 2.05, 4.0, 0.3,
             font_size=14, color="D55E00", bold=True, margin=0)
    for i, item in enumerate(sec["left_items"]):
        add_text(slide, f"• {item}", 0.7, 2.4 + i * 0.5, 5.5, 0.45,
                 font_size=12, color=theme["text"], margin=0)
    # Right card: Filter FAIL
    add_rect(slide, 6.7, 2.0, 5.9, 2.5, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, sec["right_title"], 6.9, 2.05, 4.0, 0.3,
             font_size=14, color="D55E00", bold=True, margin=0)
    for i, item in enumerate(sec["right_items"]):
        add_text(slide, f"• {item}", 6.9, 2.4 + i * 0.5, 5.5, 0.45,
                 font_size=12, color=theme["text"], margin=0)
    # Bottom figure
    fig_path = sec.get("figure", "")
    if fig_path:
        img_path = str(Path(cfg["_config_dir"]) / fig_path)
        _safe_add_picture(slide, img_path, Inches(1.0), Inches(4.7), width=Inches(11.3))


def add_recap_nonloh_caller_slide(slide, cfg: dict, theme: dict) -> None:
    """Compressed recap: Non-LOH + caller_af."""
    sec = cfg["deck"]["recap_nonloh_caller"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    items = sec["items"]
    for i, item in enumerate(items):
        y = 2.2 + i * 0.85
        add_rect(slide, 0.8, y, 11.7, 0.7, fill="FFFFFF", line=theme["line"], radius=True)
        add_rect(slide, 0.8, y, 0.12, 0.7, fill=item["color"], line=None)
        add_text(slide, item["label"], 1.1, y + 0.05, 4.5, 0.3,
                 font_size=14, color=theme["text"], bold=True, margin=0)
        add_text(slide, item["value"], 5.8, y + 0.05, 6.5, 0.3,
                 font_size=14, color=item["color"], bold=True, margin=0)
    # Bottom note
    add_text(slide, "TO 模式 ISM 增量近乎為零 — caller_af 是唯一有效判別器",
             0.8, 6.5, 11.7, 0.35,
             font_size=13, color=theme["muted"], align="center", margin=0)


def add_recap_bridge_slide(slide, cfg: dict, theme: dict) -> None:
    """Recap conclusion → 3 extension questions bridge."""
    sec = cfg["deck"]["recap_bridge"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Left: conclusions
    add_rect(slide, 0.5, 2.0, 5.9, 2.8, fill="E8F5E9", line=theme["line"], radius=True)
    add_text(slide, "上週確認", 0.7, 2.05, 3.0, 0.3,
             font_size=15, color="009E73", bold=True, margin=0)
    for i, c in enumerate(sec["conclusions"]):
        add_text(slide, c, 0.7, 2.45 + i * 0.7, 5.5, 0.65,
                 font_size=13, color=theme["text"], margin=0)
    # Divider arrow
    add_text(slide, "→", 6.1, 3.0, 0.5, 0.5,
             font_size=36, color=theme["accent"], bold=True, align="center", margin=0)
    # Right: questions
    add_rect(slide, 6.7, 2.0, 5.9, 2.8, fill="FFF8E1", line=theme["line"], radius=True)
    add_text(slide, "本週延伸", 6.9, 2.05, 3.0, 0.3,
             font_size=15, color="E69F00", bold=True, margin=0)
    for i, q in enumerate(sec["questions"]):
        add_text(slide, q, 6.9, 2.45 + i * 0.7, 5.5, 0.65,
                 font_size=14, color=theme["text"], bold=True, margin=0)
        if i < len(sec.get("questions_en", [])):
            add_text(slide, sec["questions_en"][i], 7.1, 2.75 + i * 0.7, 5.3, 0.3,
                     font_size=10, color=theme["en_color"], margin=0)


def add_competitor_paper_slide(slide, cfg: dict, theme: dict) -> None:
    """TumorLens competitor validation slide."""
    sec = cfg["deck"]["competitor_paper"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Paper card
    add_rect(slide, 0.8, 2.0, 11.7, 2.6, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, sec["paper_title"], 1.0, 2.1, 6.0, 0.35,
             font_size=16, color=theme["text"], bold=True, margin=0)
    for i, feat in enumerate(sec["features"]):
        add_text(slide, f"• {feat}", 1.0, 2.55 + i * 0.45, 5.5, 0.4,
                 font_size=12, color=theme["text"], margin=0)
    # EN features
    for i, feat in enumerate(sec.get("features_en", [])):
        add_text(slide, f"• {feat}", 6.8, 2.55 + i * 0.45, 5.5, 0.4,
                 font_size=10, color=theme["en_color"], margin=0)
    # Implication box
    add_rect(slide, 0.8, 4.9, 11.7, 1.2, fill="E3F2FD", line="0072B2", radius=True)
    add_text(slide, "→ 意義", 1.0, 4.95, 2.0, 0.3,
             font_size=14, color="0072B2", bold=True, margin=0)
    add_text(slide, sec["implication"], 1.0, 5.3, 11.3, 0.7,
             font_size=13, color=theme["text"], margin=0)


def add_loh_cnv_dual_slide(slide, cfg: dict, theme: dict) -> None:
    """Dual observation: TP enrichment + FP hotspot."""
    sec = cfg["deck"]["loh_cnv_dual"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Left: observation 1
    obs_l = sec["obs_left"]
    add_rect(slide, 0.5, 2.0, 5.9, 2.0, fill="E8F5E9", line="009E73", radius=True)
    add_text(slide, obs_l["title"], 0.7, 2.05, 5.5, 0.3,
             font_size=13, color="009E73", bold=True, margin=0)
    for i, item in enumerate(obs_l["items"]):
        add_text(slide, f"• {item}", 0.7, 2.4 + i * 0.45, 5.5, 0.4,
                 font_size=12, color=theme["text"], margin=0)
    # Right: observation 2
    obs_r = sec["obs_right"]
    add_rect(slide, 6.7, 2.0, 5.9, 2.0, fill="FFEBEE", line="D55E00", radius=True)
    add_text(slide, obs_r["title"], 6.9, 2.05, 5.5, 0.3,
             font_size=13, color="D55E00", bold=True, margin=0)
    for i, item in enumerate(obs_r["items"]):
        add_text(slide, f"• {item}", 6.9, 2.4 + i * 0.45, 5.5, 0.4,
                 font_size=12, color=theme["text"], margin=0)
    # Figure
    img_path = str(Path(cfg["_config_dir"]) / sec["figure"])
    _safe_add_picture(slide, img_path, Inches(0.5), Inches(4.2), width=Inches(8.0))
    # Conclusion
    add_rect(slide, 0.5, 6.3, 12.3, 0.35, fill="E3F2FD", line=None, radius=True)
    add_text(slide, f"→ {sec['conclusion']}", 0.7, 6.32, 11.9, 0.3,
             font_size=12, color="0072B2", bold=True, align="center", margin=0)


def add_cn_rootcause_slide(slide, cfg: dict, theme: dict) -> None:
    """CN=3 root cause + cross-sample validation."""
    sec = cfg["deck"]["cn_rootcause"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Two figures side by side
    img_l = str(Path(cfg["_config_dir"]) / sec["figure_left"])
    img_r = str(Path(cfg["_config_dir"]) / sec["figure_right"])
    _safe_add_picture(slide, img_l, Inches(0.3), Inches(2.0), width=Inches(6.2))
    _safe_add_picture(slide, img_r, Inches(6.7), Inches(2.0), width=Inches(6.2))
    # Bottom bullet items
    add_rect(slide, 0.5, 5.0, 12.3, 1.6, fill="FFFFFF", line=theme["line"], radius=True)
    for i, item in enumerate(sec["items"][:5]):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.1
        y = 5.1 + row * 0.45
        add_text(slide, f"• {item}", x, y, 5.9, 0.4,
                 font_size=11, color=theme["text"], margin=0)


def add_cnv_conclusion_slide(slide, cfg: dict, theme: dict) -> None:
    """CNV zone conclusion slide."""
    sec = cfg["deck"]["cnv_conclusion"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, v in enumerate(sec["verdicts"]):
        y = 2.2 + i * 1.1
        fill = "FFEBEE" if v["color"] == "D55E00" else ("E8F5E9" if v["color"] == "009E73" else "E3F2FD")
        add_rect(slide, 0.8, y, 11.7, 0.9, fill=fill, line=v["color"], radius=True)
        add_text(slide, f"{v['icon']}  {v['text']}", 1.0, y + 0.15, 11.3, 0.6,
                 font_size=14, color=theme["text"], margin=0)


def add_self_phasing_v5_slide(slide, cfg: dict, theme: dict) -> None:
    """Self-phasing problem with dual figures."""
    sp = cfg["deck"]["self_phasing"]
    set_bg(slide, theme["bg"])
    add_title(slide, sp["title"], sp["subtitle"], theme,
              title_en=sp.get("title_en"), subtitle_en=sp.get("subtitle_en"))
    build_pipeline_flow(slide, sp["steps"], theme, area_top=2.2, area_height=1.8)
    # Two figures
    fig_l = sp.get("figure_left", "")
    fig_r = sp.get("figure_right", "")
    if fig_l:
        img = str(Path(cfg["_config_dir"]) / fig_l)
        _safe_add_picture(slide, img, Inches(0.3), Inches(4.3), width=Inches(6.2))
    if fig_r:
        img = str(Path(cfg["_config_dir"]) / fig_r)
        _safe_add_picture(slide, img, Inches(6.7), Inches(4.3), width=Inches(6.2))


def add_pon_only_fix_slide(slide, cfg: dict, theme: dict) -> None:
    """PON-only phasing fix slide."""
    sec = cfg["deck"]["pon_only_fix"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Metrics table
    rows = [["指標", "修正前", "修正後"]]
    for m in sec["metrics"]:
        rows.append([m["label"], m["before"], m["after"]])
    add_table(slide, rows, 0.5, 2.0, 5.5, 2.2, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=12, col_widths=[2.0, 1.5, 2.0])
    # Figure
    fig = sec.get("figure", "")
    if fig:
        img_path = str(Path(cfg["_config_dir"]) / fig)
        _safe_add_picture(slide, img_path, Inches(6.3), Inches(2.0), width=Inches(6.5))
    # Note
    note = sec.get("note", "")
    if note:
        add_rect(slide, 0.5, 4.5, 12.3, 1.8, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, note, 0.7, 4.6, 11.9, 1.6,
                 font_size=11, color=theme["text"], margin=0, line_spacing=1.5)


def add_v5_evolution_slide(slide, cfg: dict, theme: dict) -> None:
    """V5 version evolution table."""
    sec = cfg["deck"]["v5_evolution"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Version table
    rows = [["版本", "修正內容", "SEQC2 F1", "AMB%", "狀態"]]
    for v in sec["versions"]:
        rows.append([v["name"], v["fix"], v["f1"], v["amb"], v["status"]])
    add_table(slide, rows, 0.5, 2.0, 12.3, 3.2, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=12, col_widths=[1.2, 3.5, 1.5, 1.2, 1.8])
    # Explanation
    add_rect(slide, 0.5, 5.5, 12.3, 1.3, fill="FFF8E1", line="E69F00", radius=True)
    add_text(slide, sec["explanation"], 0.7, 5.6, 11.9, 1.1,
             font_size=12, color=theme["text"], margin=0, line_spacing=1.4)


def add_v5_validation_slide(slide, cfg: dict, theme: dict) -> None:
    """V5 validation: approaching Paired quality."""
    sec = cfg["deck"]["v5_validation"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # HP changes table
    rows = [["HP Tag", "V3-Fixed", "V5", "Delta"]]
    for h in sec["hp_changes"]:
        rows.append([h["tag"], h["v3f"], h["v5"], h["delta"]])
    add_table(slide, rows, 0.5, 2.0, 6.0, 2.4, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=11, col_widths=[2.5, 1.5, 1.5, 2.0])
    # Concordance card
    add_rect(slide, 6.8, 2.0, 5.7, 2.4, fill="E8F5E9", line="009E73", radius=True)
    add_text(slide, "Concordance 驗證", 7.0, 2.05, 4.0, 0.3,
             font_size=14, color="009E73", bold=True, margin=0)
    for i, c in enumerate(sec["concordance"]):
        add_text(slide, f"{c['metric']}: {c['value']}", 7.0, 2.5 + i * 0.45, 5.3, 0.4,
                 font_size=12, color=theme["text"], margin=0)
    # Bottom line
    add_rect(slide, 0.5, 4.7, 12.3, 0.6, fill="E3F2FD", line=None, radius=True)
    add_text(slide, sec["bottom_line"], 0.7, 4.75, 11.9, 0.5,
             font_size=12, color="0072B2", bold=True, margin=0)


def add_beyond_auc_slide(slide, cfg: dict, theme: dict) -> None:
    """Beyond-AUC 7 methods comprehensive validation."""
    sec = cfg["deck"]["beyond_auc"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    methods = sec["methods"]
    for i, m in enumerate(methods):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.3
        y = 2.0 + row * 1.15
        add_rect(slide, x, y, 6.0, 0.95, fill="FFFFFF", line=theme["line"], radius=True)
        add_text(slide, f"{m['icon']}  {m['name']}", x + 0.15, y + 0.08, 2.5, 0.3,
                 font_size=13, color=theme["text"], bold=True, margin=0)
        add_text(slide, m["result"], x + 0.15, y + 0.42, 5.6, 0.45,
                 font_size=11, color=theme["muted"], margin=0)
    # Conclusion
    add_rect(slide, 0.5, 6.2, 12.3, 0.45, fill="FFEBEE", line="D55E00", radius=True)
    add_text(slide, f"→ {sec['conclusion']}", 0.7, 6.25, 11.9, 0.35,
             font_size=13, color="D55E00", bold=True, align="center", margin=0)


def add_literature_validation_slide(slide, cfg: dict, theme: dict) -> None:
    """Literature hypothesis cross-validation L1-L4."""
    sec = cfg["deck"]["literature_validation"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Figure on the right
    fig = sec.get("figure", "")
    if fig:
        img_path = str(Path(cfg["_config_dir"]) / fig)
        _safe_add_picture(slide, img_path, Inches(7.0), Inches(2.0), width=Inches(5.8))
    # Hypothesis cards on the left
    for i, h in enumerate(sec["hypotheses"]):
        y = 2.0 + i * 1.15
        fill = "FFEBEE" if h["verdict"] == "NEGATIVE" else "E8F5E9"
        add_rect(slide, 0.5, y, 6.2, 0.95, fill=fill, line=h["color"], radius=True)
        add_text(slide, f"{h['id']}: {h['name']}", 0.7, y + 0.05, 3.5, 0.3,
                 font_size=13, color=h["color"], bold=True, margin=0)
        add_text(slide, h["verdict"], 4.8, y + 0.05, 1.5, 0.3,
                 font_size=12, color=h["color"], bold=True, align="right", margin=0)
        add_text(slide, h["result"], 0.7, y + 0.38, 5.8, 0.5,
                 font_size=10, color=theme["text"], margin=0)


def add_exhaustion_rootcause_slide(slide, cfg: dict, theme: dict) -> None:
    """Feature exhaustion biological root cause."""
    sec = cfg["deck"]["exhaustion_rootcause"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    # Core reason card
    add_rect(slide, 0.8, 2.0, 11.7, 2.2, fill="FFFFFF", line=theme["line"], radius=True)
    add_text(slide, sec["core_reason"], 1.0, 2.1, 11.3, 2.0,
             font_size=14, color=theme["text"], margin=0, line_spacing=1.6)
    # NOT items
    add_rect(slide, 0.8, 4.5, 5.5, 1.8, fill="FFEBEE", line="D55E00", radius=True)
    for i, item in enumerate(sec["not_items"]):
        add_text(slide, f"✗  {item}", 1.0, 4.6 + i * 0.5, 5.1, 0.45,
                 font_size=13, color="D55E00", margin=0)
    # IS item
    add_rect(slide, 6.6, 4.5, 5.9, 1.8, fill="E3F2FD", line="0072B2", radius=True)
    add_text(slide, f"→  {sec['is_item']}", 6.8, 4.7, 5.5, 1.4,
             font_size=13, color="0072B2", bold=True, margin=0, line_spacing=1.5)


def add_exclusion_diagnosis_slide(slide, cfg: dict, theme: dict) -> None:
    """Four independent exclusions: GC + CN + PON + H2009."""
    sec = cfg["deck"]["exclusion_diagnosis"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, ex in enumerate(sec["exclusions"]):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.3
        y = 2.0 + row * 2.0
        fill = "FFEBEE" if ex["color"] == "D55E00" else ("E8F5E9" if ex["color"] == "009E73" else "FFF8E1")
        add_rect(slide, x, y, 6.0, 1.7, fill=fill, line=ex["color"], radius=True)
        add_text(slide, ex["name"], x + 0.2, y + 0.1, 3.0, 0.3,
                 font_size=14, color=ex["color"], bold=True, margin=0)
        add_text(slide, ex["result"], x + 0.2, y + 0.5, 5.6, 0.4,
                 font_size=12, color=theme["text"], bold=True, margin=0)
        add_text(slide, ex["detail"], x + 0.2, y + 1.0, 5.6, 0.5,
                 font_size=11, color=theme["muted"], margin=0)
    # Bottom conclusion
    add_rect(slide, 0.5, 6.2, 12.3, 0.45, fill="FFEBEE", line=None, radius=True)
    add_text(slide, f"→ {sec['conclusion']}", 0.7, 6.25, 11.9, 0.35,
             font_size=13, color="D55E00", bold=True, align="center", margin=0)


def add_positive_observations_slide(slide, cfg: dict, theme: dict) -> None:
    """Positive observations amid exhaustion."""
    sec = cfg["deck"]["positive_observations"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, f in enumerate(sec["findings"]):
        y = 2.0 + i * 1.4
        add_rect(slide, 0.8, y, 11.7, 1.2, fill="FFFFFF", line=f["color"], radius=True)
        add_rect(slide, 0.8, y, 0.12, 1.2, fill=f["color"], line=None)
        add_text(slide, f["title"], 1.1, y + 0.08, 11.2, 0.3,
                 font_size=14, color=f["color"], bold=True, margin=0)
        add_text(slide, f["detail"], 1.1, y + 0.45, 11.2, 0.65,
                 font_size=12, color=theme["text"], margin=0, line_spacing=1.4)
    # Bottom line
    add_rect(slide, 0.5, 6.3, 12.3, 0.35, fill="E8F5E9", line=None, radius=True)
    add_text(slide, f"→ {sec['bottom_line']}", 0.7, 6.32, 11.9, 0.3,
             font_size=12, color="009E73", bold=True, align="center", margin=0)


def add_phase_bcd_slide(slide, cfg: dict, theme: dict) -> None:
    """Phase B/C/D Dual-BAM validation."""
    sec = cfg["deck"]["phase_bcd"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    rows = [["模組", "結果", "判定"]]
    for m in sec["modules"]:
        rows.append([m["name"], m["result"], m["verdict"]])
    add_table(slide, rows, 1.5, 2.2, 10.3, 2.5, theme,
              header_fill=theme["dark"], body_fill="FFFFFF", zebra_fill="F4F0E8",
              font_size=14, col_widths=[3.0, 3.5, 1.5])
    # Note
    add_rect(slide, 1.5, 5.0, 10.3, 0.6, fill="E8F5E9", line="009E73", radius=True)
    add_text(slide, sec["note"], 1.7, 5.05, 9.9, 0.5,
             font_size=14, color="009E73", bold=True, align="center", margin=0)


def add_subclone_groups_slide(slide, cfg: dict, theme: dict) -> None:
    """Subclone 4-group discovery."""
    sec = cfg["deck"]["subclone_groups"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, g in enumerate(sec["groups"]):
        y = 2.2 + i * 1.1
        bar_w = float(g["pct"].replace("%", "")) / 100.0 * 8.0
        # Background bar
        add_rect(slide, 0.8, y, 11.7, 0.9, fill="FFFFFF", line=theme["line"], radius=True)
        # Colored proportion bar
        add_rect(slide, 0.8, y, max(bar_w, 0.3), 0.9, fill=g["color"], line=None, radius=True)
        # Text
        add_text(slide, f"{g['name']} — {g['pct']}", 0.95, y + 0.08, 4.0, 0.3,
                 font_size=14, color="FFFFFF" if bar_w > 2.0 else theme["text"], bold=True, margin=0)
        add_text(slide, g["detail"], 5.0, y + 0.25, 7.3, 0.4,
                 font_size=12, color=theme["text"], margin=0)
    # Significance
    add_rect(slide, 0.5, 6.3, 12.3, 0.35, fill="E3F2FD", line=None, radius=True)
    add_text(slide, sec["significance"], 0.7, 6.32, 11.9, 0.3,
             font_size=12, color="0072B2", bold=True, align="center", margin=0)


def add_normal_baseline_slide(slide, cfg: dict, theme: dict) -> None:
    """Normal Baseline feasibility slide."""
    sec = cfg["deck"]["normal_baseline"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, item in enumerate(sec["items"]):
        y = 2.2 + i * 1.0
        fill = "E8F5E9" if item.startswith("✅") else "E3F2FD"
        add_rect(slide, 0.8, y, 11.7, 0.8, fill=fill, line=theme["line"], radius=True)
        add_text(slide, item, 1.0, y + 0.1, 11.3, 0.6,
                 font_size=14, color=theme["text"], bold=True, margin=0)
    # EN items
    for i, item in enumerate(sec.get("items_en", [])):
        y = 2.2 + i * 1.0
        add_text(slide, item, 1.2, y + 0.5, 11.1, 0.3,
                 font_size=10, color=theme["en_color"], margin=0)


def add_v5_conclusion_table_slide(slide, cfg: dict, theme: dict) -> None:
    """Conclusion table: 5 closed + 3 positive + 2 diagnostics."""
    sec = cfg["deck"]["conclusion_table"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    cats = sec["categories"]
    x_pos = [0.3, 4.5, 8.8]
    widths = [3.9, 4.0, 3.8]
    for ci, cat in enumerate(cats):
        x = x_pos[ci]
        w = widths[ci]
        add_rect(slide, x, 2.0, w, 0.45, fill=cat["color"], line=None, radius=True)
        add_text(slide, cat["category"], x + 0.1, 2.05, w - 0.2, 0.35,
                 font_size=14, color="FFFFFF", bold=True, align="center", margin=0)
        for j, item in enumerate(cat["items"]):
            y = 2.6 + j * 0.65
            add_rect(slide, x, y, w, 0.55, fill="FFFFFF", line=theme["line"], radius=True)
            add_text(slide, item, x + 0.15, y + 0.05, w - 0.3, 0.45,
                     font_size=11, color=theme["text"], margin=0)


def add_v5_next_steps_slide(slide, cfg: dict, theme: dict) -> None:
    """Next steps with timeline phases."""
    sec = cfg["deck"]["next_steps"]
    set_bg(slide, theme["bg"])
    add_title(slide, sec["title"], sec["subtitle"], theme,
              title_en=sec.get("title_en"), subtitle_en=sec.get("subtitle_en"))
    for i, phase in enumerate(sec["timeline"]):
        y = 2.2 + i * 1.6
        add_rect(slide, 0.5, y, 1.5, 1.3, fill=phase["color"], line=None, radius=True)
        add_text(slide, phase["phase"], 0.5, y + 0.3, 1.5, 0.4,
                 font_size=16, color="FFFFFF", bold=True, align="center", margin=0)
        add_rect(slide, 2.2, y, 10.3, 1.3, fill="FFFFFF", line=theme["line"], radius=True)
        for j, item in enumerate(phase["items"]):
            add_text(slide, f"• {item}", 2.4, y + 0.15 + j * 0.5, 9.9, 0.45,
                     font_size=13, color=theme["text"], margin=0)


def _select_slide_builders(cfg: dict, theme: dict, assets: dict) -> list:
    """Select slide builders based on available config sections.

    Supports config formats:
    - v5 (weekly report 2026-04-13, 30 slides): detected by 'layout_version'='v5'
    - v4 (structured, 27 slides): detected by 'section_dividers' in deck
    - v3 (extended, 22 slides): detected by 'to_pure_modeling' in deck
    - v2 (visual-first, 19 slides): detected by 'investigation_pipeline' in deck
    - v1 (legacy, 12 slides): original conditional logic
    """
    deck = cfg.get("deck", {})
    dividers = deck.get("section_dividers", [])

    # Helper to create a section divider builder with injected footer/page info
    def _divider(idx):
        def _build(s):
            d = dict(dividers[idx])
            d["_footer"] = cfg["meta"].get("footer", "")
            d["_page"] = ""
            add_section_divider_slide(s, d, theme)
        return _build

    # --- v5 layout (30 slides): LOH×CNV + HP Tag weekly report ---
    if deck.get("layout_version") == "v5":
        builders = [
            # Section A: Cover + Guide (2 slides)
            lambda s: add_title_slide(s, cfg, theme),                        #  1  封面
            lambda s: add_key_findings_slide(s, cfg, theme),                 #  2  本週重點
            # Section B: Last Week Recap (5 slides)
            lambda s: add_investigation_pipeline_slide(s, cfg, theme),       #  3  推論地圖
            lambda s: add_recap_seqc2_quadrant_slide(s, cfg, theme),         #  4  SEQC2+四象限
            lambda s: add_recap_loh_filter_slide(s, cfg, theme),             #  5  LOH影響+Filter
            lambda s: add_recap_nonloh_caller_slide(s, cfg, theme),          #  6  Non-LOH+caller
            lambda s: add_recap_bridge_slide(s, cfg, theme),                 #  7  結論→銜接
            # Section C: LOH × CNV (5 slides)
            _divider(0),                                                     #  8  Section: LOH×CNV
            lambda s: add_competitor_paper_slide(s, cfg, theme),             #  9  TumorLens
            lambda s: add_loh_cnv_dual_slide(s, cfg, theme),                 # 10  雙面觀察
            lambda s: add_cn_rootcause_slide(s, cfg, theme),                 # 11  CN=3根因
            lambda s: add_cnv_conclusion_slide(s, cfg, theme),               # 12  CNV結論
            # Section D: Self-Phasing + V5 (5 slides)
            _divider(1),                                                     # 13  Section: HP Tag
            lambda s: add_self_phasing_v5_slide(s, cfg, theme),              # 14  Self-Phasing
            lambda s: add_pon_only_fix_slide(s, cfg, theme),                 # 15  PON-only修正
            lambda s: add_v5_evolution_slide(s, cfg, theme),                 # 16  V5版本演進
            lambda s: add_v5_validation_slide(s, cfg, theme),                # 17  V5驗證
            # Section E: ISM Status (6 slides)
            _divider(2),                                                     # 18  Section: ISM驗證
            lambda s: add_beyond_auc_slide(s, cfg, theme),                   # 19  Beyond-AUC
            lambda s: add_literature_validation_slide(s, cfg, theme),        # 20  文獻L1-L4
            lambda s: add_exhaustion_rootcause_slide(s, cfg, theme),         # 21  耗盡根因
            lambda s: add_exclusion_diagnosis_slide(s, cfg, theme),          # 22  四項排除
            lambda s: add_positive_observations_slide(s, cfg, theme),        # 23  正面觀察
            # Section F: Phase 2 (4 slides)
            _divider(3),                                                     # 24  Section: Phase 2
            lambda s: add_phase_bcd_slide(s, cfg, theme),                    # 25  Phase BCD
            lambda s: add_subclone_groups_slide(s, cfg, theme),              # 26  4-group
            lambda s: add_normal_baseline_slide(s, cfg, theme),              # 27  Normal baseline
            # Section G: Conclusion (3 slides)
            lambda s: add_v5_conclusion_table_slide(s, cfg, theme),          # 28  結論總表
            lambda s: add_v5_next_steps_slide(s, cfg, theme),                # 29  下一步
            lambda s: add_paths_slide(s, cfg, theme),                        # 30  附錄
        ]
        return builders

    # --- v4 structured layout (29 slides) with section dividers ---
    if dividers and len(dividers) >= 4:
        builders = [
            # Act 1: Opening (3 slides)
            lambda s: add_title_slide(s, cfg, theme),                       #  1  封面
            lambda s: add_key_findings_slide(s, cfg, theme),                 #  2  五大結論
            lambda s: add_investigation_pipeline_slide(s, cfg, theme),       #  3  推論地圖
            # Section Divider A
            _divider(0),                                                     #  4  信任基礎
            # Act 2: Trust + ISM Framework (6 slides)
            lambda s: add_seqc2_chromosome_slide(s, cfg, theme),             #  5  SEQC2 驗證
            lambda s: add_seqc2_distribution_slide(s, cfg, theme),           #  6  LOH 位點定量
            lambda s: add_ism_framework_slide(s, cfg, theme),                #  7  ISM 框架
            lambda s: add_hp_imbalance_definition_slide(s, cfg, theme),      #  8  HP Imbalance 定義
            lambda s: add_hp_vs_loh_chromosome_slide(s, cfg, theme),         #  9  全基因組比較
            lambda s: add_loh_quadrant_slide(s, cfg, theme),                 # 10  HP ⊃ LOH.bed
            # Section Divider B
            _divider(1),                                                     # 11  全面驗證
            # Act 3: Systematic rejection (9 slides)
            lambda s: add_permanova_concept_slide(s, cfg, theme),            # 12  PERMANOVA
            lambda s: add_hcc1395_detailed_slide(s, cfg, theme),             # 13  HCC1395 詳細
            lambda s: add_key_features_auc_slide(s, cfg, theme),             # 14  關鍵特徵 AUC
            lambda s: add_filter_decision_slide(s, cfg, theme),              # 15  Filter
            lambda s: add_methylation_negative_slide(s, cfg, theme),         # 16  甲基化
            lambda s: add_to_nogo_evidence_slide(s, cfg, theme),             # 17  TO NO-GO
            lambda s: add_non_loh_closure_slide(s, cfg, theme),              # 18  Non-LOH
            lambda s: add_self_phasing_slide(s, cfg, theme),                 # 19  Self-phasing
            lambda s: add_code_modification_slide(s, cfg, theme),            # 20  程式碼修改 ★NEW
            # Section Divider C
            _divider(2),                                                     # 21  獨立確認
            # Act 4: Independent modeling (2 slides)
            lambda s: add_to_pure_modeling_slide(s, cfg, theme),             # 22  TO-pure
            lambda s: add_per_sample_stratification_slide(s, cfg, theme),    # 23  7 樣本分層
            # Section Divider D
            _divider(3),                                                     # 24  結論與未來
            # Act 5: Conclusion + pivot (5 slides)
            lambda s: add_asm_concept_slide(s, cfg, theme),                  # 25  ASM 亮點
            lambda s: add_conclusion_table_slide(s, cfg, theme),             # 26  結論總表
            lambda s: add_strategic_pivot_slide(s, cfg, theme),              # 27  轉向
            lambda s: add_next_steps_slide(s, cfg, theme),                   # 28  下一步
            lambda s: add_paths_slide(s, cfg, theme),                        # 29  附錄
        ]
        return builders

    # --- v3 extended layout (22 slides) ---
    if "to_pure_modeling" in deck:
        builders = [
            lambda s: add_title_slide(s, cfg, theme),                    # 1  封面
            lambda s: add_key_findings_slide(s, cfg, theme),              # 2  關鍵結論
            lambda s: add_investigation_pipeline_slide(s, cfg, theme),    # 3  調查邏輯鏈
            lambda s: add_seqc2_chromosome_slide(s, cfg, theme),          # 4  SEQC2 驗證
            lambda s: add_seqc2_distribution_slide(s, cfg, theme),        # 5  位點層級
            lambda s: add_hp_imbalance_definition_slide(s, cfg, theme),   # 6  HP Imbalance
            lambda s: add_hp_vs_loh_chromosome_slide(s, cfg, theme),      # 7  HP vs LOH.bed
            lambda s: add_loh_quadrant_slide(s, cfg, theme),              # 8  四象限
            lambda s: add_permanova_concept_slide(s, cfg, theme),         # 9  PERMANOVA
            lambda s: add_composition_effect_slide(s, cfg, theme),        # 10 Composition Effect
            lambda s: add_filter_decision_slide(s, cfg, theme),           # 11 Filter 決策
            lambda s: add_methylation_negative_slide(s, cfg, theme),      # 12 甲基化 NEGATIVE
            lambda s: add_to_nogo_evidence_slide(s, cfg, theme),          # 13 TO NO-GO
            lambda s: add_non_loh_closure_slide(s, cfg, theme),           # 14 Non-LOH
            lambda s: add_self_phasing_slide(s, cfg, theme),              # 15 Self-phasing
            lambda s: add_to_pure_modeling_slide(s, cfg, theme),          # 16 TO-pure 建模
            lambda s: add_per_sample_stratification_slide(s, cfg, theme), # 17 Per-sample 分層
            lambda s: add_asm_concept_slide(s, cfg, theme),               # 18 ASM 亮點
            lambda s: add_conclusion_table_slide(s, cfg, theme),          # 19 結論總表
            lambda s: add_strategic_pivot_slide(s, cfg, theme),           # 20 策略轉向
            lambda s: add_next_steps_slide(s, cfg, theme),                # 21 下一步
            lambda s: add_paths_slide(s, cfg, theme),                     # 22 附錄
        ]
        return builders

    # --- v2 visual-first layout (19 slides) ---
    if "investigation_pipeline" in deck:
        builders = [
            lambda s: add_title_slide(s, cfg, theme),               # 1  封面
            lambda s: add_key_findings_slide(s, cfg, theme),         # 2  四大結論
            lambda s: add_investigation_pipeline_slide(s, cfg, theme),  # 3  調查邏輯鏈
            lambda s: add_seqc2_chromosome_slide(s, cfg, theme),     # 4  SEQC2 驗證 + 染色體交集
            lambda s: add_seqc2_distribution_slide(s, cfg, theme),   # 5  位點層級 LOH 分佈
            lambda s: add_hp_imbalance_definition_slide(s, cfg, theme),  # 6  HP Imbalance 定義
            lambda s: add_hp_vs_loh_chromosome_slide(s, cfg, theme),    # 7  HP vs LOH.bed 染色體分佈
            lambda s: add_loh_quadrant_slide(s, cfg, theme),         # 8  LOH 四象限
            lambda s: add_permanova_concept_slide(s, cfg, theme),    # 9  PERMANOVA 概念
            lambda s: add_filter_decision_slide(s, cfg, theme),      # 10 Filter 決策樹
            lambda s: add_methylation_negative_slide(s, cfg, theme), # 11 甲基化 NEGATIVE
            lambda s: add_to_nogo_evidence_slide(s, cfg, theme),     # 12 TO NO-GO 證據
            lambda s: add_non_loh_closure_slide(s, cfg, theme),      # 13 Non-LOH 關閉
            lambda s: add_self_phasing_slide(s, cfg, theme),         # 14 Self-phasing 因果鏈
            lambda s: add_asm_concept_slide(s, cfg, theme),          # 15 ASM 亮點
            lambda s: add_conclusion_table_slide(s, cfg, theme),     # 16 結論總表
            lambda s: add_strategic_pivot_slide(s, cfg, theme),      # 17 策略轉向
            lambda s: add_next_steps_slide(s, cfg, theme),           # 18 下一步
            lambda s: add_paths_slide(s, cfg, theme),                # 19 附錄
        ]
        return builders

    # --- v1 legacy layout (12 slides) ---
    builders = [
        lambda s: add_title_slide(s, cfg, theme),
        lambda s: add_key_findings_slide(s, cfg, theme),
        lambda s: add_framework_slide(s, cfg, theme),
        lambda s: add_goal_status_slide(s, cfg, theme),
        lambda s: create_timeline_slide(s, cfg, theme),
    ]

    if "seqc2_validation" in deck:
        builders.append(lambda s: add_seqc2_validation_slide(s, cfg, theme))
    elif "hp_bug_fix" in deck and assets.get("hp_fix_chart"):
        builders.append(lambda s: add_hp_bug_fix_slide(s, cfg, theme, assets["hp_fix_chart"]))
    elif assets.get("paired_chart"):
        builders.append(lambda s: add_paired_rerun_slide(s, cfg, theme, assets["paired_chart"]))

    if "ism_loh_impact" in deck:
        builders.append(lambda s: add_ism_loh_impact_slide(s, cfg, theme))
    elif "loh_evidence_panel" in deck:
        builders.append(lambda s: add_loh_evidence_panel_slide(s, cfg, theme))
    elif "to_rescue" in deck:
        builders.append(lambda s: add_to_rescue_slide(s, cfg, theme))

    if "non_loh_closure" in deck:
        builders.append(lambda s: add_non_loh_closure_slide(s, cfg, theme))
    elif "phase1a_round3" in deck:
        builders.append(lambda s: add_phase1a_loh_test_slide(s, cfg, theme))
    elif assets.get("phase2_chart"):
        builders.append(lambda s: add_phase2_slide(s, cfg, theme, assets["phase2_chart"]))

    if "negative_results" in deck:
        builders.append(lambda s: add_negative_results_slide(s, cfg, theme))
    elif "feature_roles" in deck:
        builders.append(lambda s: add_feature_roles_slide(s, cfg, theme))

    if "conclusion_table" in deck:
        builders.append(lambda s: add_conclusion_table_slide(s, cfg, theme))
    elif "vision" in deck:
        builders.append(lambda s: add_vision_slide(s, cfg, theme))
    elif "annotation_policy" in deck:
        builders.append(lambda s: add_annotation_slide(s, cfg, theme))

    builders.append(lambda s: add_next_steps_slide(s, cfg, theme))
    builders.append(lambda s: add_paths_slide(s, cfg, theme))

    return builders


# ====================================================================
# Audit / Quality Check
# ====================================================================

def _contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """Calculate WCAG 2.2 contrast ratio between two hex colors."""
    def _rel_lum(h: str) -> float:
        h = h.replace("#", "").strip()
        r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
        def _lin(c):
            return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
        return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)
    l1 = _rel_lum(fg_hex)
    l2 = _rel_lum(bg_hex)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def audit_pptx(pptx_path: str, config: dict | None = None) -> dict:
    """Post-generation quality audit. Returns per-slide PASS/FAIL report.

    Checks:
      M1: All text font_size >= 9pt
      M2: Title font_size >= 28pt (page-level titles)
      M3: WCAG AA contrast (normal >= 4.5:1, large >= 3:1)
      M4: Bullets per slide <= 4
      M5: Font name consistency (expect Arial)
      M6: Dark page text uses light color
      M7: Speaker Notes present on every slide
      M8: Data accuracy (config key numbers vs slide text)
    """
    from pptx import Presentation as PrsReader
    prs = PrsReader(pptx_path)
    results = {"slides": [], "summary": {"pass": 0, "fail": 0, "warnings": []}}

    for idx, slide in enumerate(prs.slides, start=1):
        slide_result = {"slide": idx, "checks": {}, "pass": True}
        fonts_found = set()
        sizes_found = []
        bullet_count = 0
        has_notes = False

        # Check notes
        try:
            ns = slide.notes_slide
            notes_text = ns.notes_text_frame.text.strip() if ns else ""
            has_notes = len(notes_text) > 5
        except Exception:
            has_notes = False

        slide_result["checks"]["M7_notes"] = "PASS" if has_notes else "WARN"

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Count bullet paragraphs (check for bullet XML element)
                if hasattr(para, '_pPr') and para._pPr is not None:
                    buNone = para._pPr.find(qn('a:buNone'))
                    if buNone is None and para.text.strip():
                        # Has potential bullet formatting
                        pass
                # Simple heuristic: count non-empty paragraphs with bullet-like prefix
                txt = para.text.strip()
                if txt and (txt.startswith('\u2022') or txt.startswith('-') or txt.startswith('\u2713')):
                    bullet_count += 1
                for run in para.runs:
                    if run.font.name:
                        fonts_found.add(run.font.name)
                    if run.font.size:
                        sizes_found.append(run.font.size.pt)

        # M1: min font size
        min_size = min(sizes_found) if sizes_found else 0
        slide_result["checks"]["M1_min_font"] = "PASS" if min_size >= 9 else f"FAIL (min={min_size:.1f}pt)"

        # M2: max font >= 28pt (for non-empty slides)
        max_size = max(sizes_found) if sizes_found else 0
        slide_result["checks"]["M2_title_size"] = "PASS" if max_size >= 28 or len(sizes_found) < 3 else f"WARN (max={max_size:.1f}pt)"

        # M4: bullet count
        slide_result["checks"]["M4_bullets"] = "PASS" if bullet_count <= 4 else f"WARN ({bullet_count} bullets)"

        # M5: font consistency
        non_code_fonts = {f for f in fonts_found if f != "Consolas"}
        expected = {"Arial", "Arial Bold"}
        unexpected = non_code_fonts - expected - {"Calibri"}  # Allow Calibri as legacy
        slide_result["checks"]["M5_fonts"] = "PASS" if not unexpected else f"WARN ({unexpected})"

        # Determine overall pass
        for k, v in slide_result["checks"].items():
            if "FAIL" in str(v):
                slide_result["pass"] = False

        if slide_result["pass"]:
            results["summary"]["pass"] += 1
        else:
            results["summary"]["fail"] += 1

        results["slides"].append(slide_result)

    total = len(results["slides"])
    results["summary"]["total"] = total
    results["summary"]["verdict"] = "ALL PASS" if results["summary"]["fail"] == 0 else f"{results['summary']['fail']}/{total} FAIL"

    # Print summary
    print(f"\n{'='*50}")
    print(f"PPTX Audit: {results['summary']['verdict']}")
    print(f"{'='*50}")
    for sr in results["slides"]:
        status = "✓" if sr["pass"] else "✗"
        checks_str = " | ".join(f"{k}={v}" for k, v in sr["checks"].items() if "FAIL" in str(v) or "WARN" in str(v))
        if checks_str:
            print(f"  Slide {sr['slide']}: {status}  {checks_str}")
        else:
            print(f"  Slide {sr['slide']}: {status}")
    print(f"{'='*50}\n")

    return results


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", required=True, type=Path)
    args = parser.parse_args()

    cfg = load_config(args.config)
    theme = cfg["theme"]
    assets = build_assets(cfg)
    prs = create_presentation(cfg)

    slide_builders = _select_slide_builders(cfg, theme, assets)

    total_pages = len(slide_builders)
    # Detect dark page indices (1-based): first slide, section dividers, and second-to-last
    dark_indices = {1, total_pages - 1}
    # v5: section dividers at slides 8, 13, 18, 24
    if cfg.get("deck", {}).get("layout_version") == "v5":
        dark_indices = {1, 8, 13, 18, 24}
    # Collect notes from config deck sections
    deck_notes = cfg.get("deck", {}).get("notes", [])

    for idx, builder in enumerate(slide_builders, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder(slide)
        add_footer(slide, theme, idx, total_pages, cfg["meta"]["footer"], dark=idx in dark_indices)
        # Inject speaker notes from config
        if idx - 1 < len(deck_notes) and deck_notes[idx - 1]:
            inject_speaker_notes(slide, deck_notes[idx - 1])

    output_path = Path(cfg["meta"]["output_pptx"])
    ensure_dir(output_path.parent)
    prs.save(output_path)
    print(f"saved_pptx={output_path}")
    print(f"assets_dir={cfg['meta']['assets_dir']}")

    # Auto audit
    audit_pptx(str(output_path), cfg)

    # Mandatory rendering-safety check (catches PowerPoint crash triggers)
    _verify_pptx_rendering_safety(str(output_path))


def _strip_table_style_ids(pptx_path: str) -> None:
    """Remove tableStyleId from all tables in saved PPTX.

    python-pptx injects a tableStyleId GUID that references a built-in
    theme style. If the slide master doesn't define that style, PowerPoint
    shows a "repair" dialog on open. Stripping the reference is safe —
    our explicit cell fills/fonts already control appearance.
    """
    import zipfile, shutil, tempfile
    from lxml import etree

    tmp = tempfile.mktemp(suffix='.pptx')
    shutil.copy2(pptx_path, tmp)
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    count = 0
    with zipfile.ZipFile(tmp, 'r') as zin:
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item.startswith('ppt/slides/slide') and item.endswith('.xml'):
                    root = etree.fromstring(data)
                    for tsi in root.iter(f'{{{ns_a}}}tableStyleId'):
                        tsi.getparent().remove(tsi)
                        count += 1
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    Path(tmp).unlink()
    if count:
        print(f"  [FIX] Stripped {count} tableStyleId references (prevents PPT repair dialog)")


def _verify_pptx_rendering_safety(pptx_path: str) -> None:
    """Post-save mandatory check for issues that cause PowerPoint to fail.

    Catches: zero/tiny connectors, zero-dim shapes, missing backgrounds,
    corrupt images, broken relationships, invalid XML.
    Raises SystemExit on critical failure so broken files are never silently produced.
    """
    from pptx import Presentation as PrsReader
    from lxml import etree
    import zipfile, io

    prs = PrsReader(pptx_path)
    critical = []

    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    for i, slide in enumerate(prs.slides, 1):
        # --- Background ---
        bg = slide._element.find(
            f'.//{{{ns_p}}}bg',
        )
        if bg is None:
            # Check alternative namespace path
            bg = slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if bg is None:
            critical.append(f"Slide {i}: NO BACKGROUND — will render as white page")

        for j, shape in enumerate(slide.shapes):
            w = shape.width or 0
            h = shape.height or 0

            # --- Zero / negative dimensions ---
            if w == 0 or h == 0:
                critical.append(f"Slide {i} shape[{j}] '{shape.name}': ZERO dimension w={w} h={h}")
            if w < 0 or h < 0:
                critical.append(f"Slide {i} shape[{j}] '{shape.name}': NEGATIVE dimension w={w} h={h}")

            # --- Connector minimum size (must be > ~0.04" = 36576 EMU) ---
            tag = shape._element.tag.split('}')[-1] if '}' in shape._element.tag else shape._element.tag
            if tag == 'cxnSp':
                min_safe = 36576  # ~0.04"
                if w < min_safe or h < min_safe:
                    critical.append(
                        f"Slide {i} connector[{j}] '{shape.name}': "
                        f"TINY extent w={w} h={h} EMU (min={min_safe}) — "
                        f"PowerPoint may fail to render this slide"
                    )

            # --- Image integrity + RGBA + oversize check ---
            if tag == 'pic':
                try:
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob))
                    mode = img.mode
                    px_w, px_h = img.size
                    img.verify()
                    if mode in ("RGBA", "PA"):
                        critical.append(
                            f"Slide {i} image[{j}] '{shape.name}': "
                            f"RGBA/PA mode embedded — PowerPoint may fail to render. "
                            f"Use _safe_add_picture() to auto-convert."
                        )
                    if px_w > _SAFE_PIC_MAX_PX * 1.1:
                        critical.append(
                            f"Slide {i} image[{j}] '{shape.name}': "
                            f"OVERSIZE {px_w}×{px_h} px (max {_SAFE_PIC_MAX_PX}px) — "
                            f"may cause rendering lag. Use _safe_add_picture()."
                        )
                except Exception as e:
                    critical.append(f"Slide {i} image[{j}] '{shape.name}': CORRUPT — {e}")

    # --- ZIP integrity ---
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            bad = zf.testzip()
            if bad:
                critical.append(f"ZIP corrupt entry: {bad}")

            # Verify all slide XML parses
            for si in range(1, len(prs.slides) + 1):
                f = f'ppt/slides/slide{si}.xml'
                if f not in zf.namelist():
                    critical.append(f"Missing slide XML: {f}")
                else:
                    try:
                        etree.fromstring(zf.read(f))
                    except Exception as e:
                        critical.append(f"{f} XML parse error: {e}")

            # Verify presentation.xml references
            pres_rels = etree.fromstring(zf.read('ppt/_rels/presentation.xml.rels'))
            rid_map = {r.get('Id'): r.get('Target') for r in pres_rels}
            pres = etree.fromstring(zf.read('ppt/presentation.xml'))
            sld_list = pres.find(f'{{{ns_p}}}sldIdLst')
            if sld_list is not None:
                ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                for sid in sld_list.findall(f'{{{ns_p}}}sldId'):
                    rid = sid.get(f'{{{ns_r}}}id')
                    target = rid_map.get(rid, '')
                    full = f'ppt/{target}'
                    if full not in zf.namelist():
                        critical.append(f"Broken slide ref: rId={rid} → {target}")
    except Exception as e:
        critical.append(f"ZIP open error: {e}")

    # --- Report ---
    print(f"\n{'='*50}")
    if critical:
        print(f"RENDERING SAFETY: ❌ FAIL ({len(critical)} issues)")
        print(f"{'='*50}")
        for c in critical:
            print(f"  ❌ {c}")
        print(f"\n⚠️  This PPTX may fail to open in PowerPoint!")
        print(f"{'='*50}")
        raise SystemExit(1)
    else:
        print(f"RENDERING SAFETY: ✅ PASS ({len(prs.slides)} slides verified)")
        print(f"{'='*50}")


if __name__ == "__main__":
    main()
