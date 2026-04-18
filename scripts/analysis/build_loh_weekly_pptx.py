#!/usr/bin/env python3
"""
LOH Weekly Report PPTX Generator — Dark Theme Detective Story
Generates a 16-slide PPTX from a JSON config file.

Usage:
    python scripts/analysis/build_loh_weekly_pptx.py \
        docs/presentations/validated/2026/04/20260401_LOH_weekly_report_draft/pptx_config.json \
        docs/presentations/validated/2026/04/20260401_LOH_weekly_report_draft/LOH_weekly_report_v1.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.6
MARGIN_R = 0.6
MARGIN_T = 0.4
TITLE_H = 0.85
FOOTER_Y = 7.05
FOOTER_H = 0.35
MIN_FONT = 10  # 任何元素的絕對最小字體 (pt)


# ── Helpers ──────────────────────────────────────────────────────────────

def hex_rgb(value: str) -> RGBColor:
    v = value.replace("#", "").strip()
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)


def add_text(
    slide, text: str, x: float, y: float, w: float, h: float, *,
    font_size: float = 18, color: str = "FFFFFF", bold: bool = False,
    font_name: str = "Calibri", align: str = "left",
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP, italic: bool = False,
    margin: float = 0.05, line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_rgb(color)


def add_rich_text(slide, parts: list[dict], x: float, y: float, w: float, h: float, *,
                  align: str = "left", font_name: str = "Calibri", line_spacing: float | None = None):
    """Add text box with mixed formatting (bold/color segments)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing:
        p.line_spacing = line_spacing
    for part in parts:
        run = p.add_run()
        run.text = part.get("text", "")
        run.font.size = Pt(part.get("size", 18))
        run.font.name = part.get("font", font_name)
        run.font.bold = part.get("bold", False)
        run.font.italic = part.get("italic", False)
        run.font.color.rgb = hex_rgb(part.get("color", "E0E0E0"))


def add_bullets(
    slide, items: list[str], x: float, y: float, w: float, h: float, *,
    font_size: float = 18, color: str = "E0E0E0", font_name: str = "Calibri",
    spacing: float = 6, items_en: list[str] | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = hex_rgb(color)
        # English translation (smaller, italic, muted)
        if items_en and idx < len(items_en) and items_en[idx]:
            en_run = p.add_run()
            en_run.text = "\n" + items_en[idx]
            en_size = max(MIN_FONT, font_size - 4)
            en_run.font.size = Pt(en_size)
            en_run.font.name = font_name
            en_run.font.italic = True
            en_run.font.color.rgb = hex_rgb("8888AA")


def add_rect(slide, x: float, y: float, w: float, h: float, *,
             fill: str, line: str | None = None, radius: bool = False, transparency: int = 0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_image(slide, img_path: Path, x: float, y: float, w: float, h: float):
    """Add image preserving original aspect ratio, fit-within bounding box, centered."""
    if img_path.exists():
        from PIL import Image
        with Image.open(img_path) as img:
            orig_w, orig_h = img.size
        aspect = orig_w / orig_h
        box_aspect = w / h
        if aspect > box_aspect:
            final_w = w
            final_h = w / aspect
        else:
            final_h = h
            final_w = h * aspect
        final_x = x + (w - final_w) / 2
        final_y = y + (h - final_h) / 2
        slide.shapes.add_picture(str(img_path), Inches(final_x), Inches(final_y),
                                 width=Inches(final_w), height=Inches(final_h))
    else:
        add_rect(slide, x, y, w, h, fill="333355", line="555577", radius=True)
        add_text(slide, f"[Missing: {img_path.name}]", x + 0.1, y + h / 2 - 0.15, w - 0.2, 0.3,
                 font_size=11, color="FF6B6B", align="center")


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, theme: dict, *,
              col_widths: list[float] | None = None, font_size: int = 12):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        total = sum(col_widths)
        for idx, cw in enumerate(col_widths):
            tbl.columns[idx].width = Inches(w * (cw / total))
    row_h = Inches(h / len(rows))
    for idx in range(len(rows)):
        tbl.rows[idx].height = row_h
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = hex_rgb(theme.get("table_header", "2a2a4e"))
            else:
                cell.fill.fore_color.rgb = hex_rgb(theme.get("table_body", "22223a") if r % 2 == 1 else theme.get("table_alt", "1e1e36"))
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = theme.get("font_body", "Calibri")
                    run.font.size = Pt(font_size)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = hex_rgb(theme.get("accent", "00d4ff"))
                    else:
                        run.font.color.rgb = hex_rgb(theme.get("body_color", "E0E0E0"))
    return tbl


def add_slide_title(slide, title: str, theme: dict, *, subtitle: str | None = None, title_en: str | None = None):
    add_text(slide, title, MARGIN_L, MARGIN_T, SLIDE_W - MARGIN_L - MARGIN_R, TITLE_H,
             font_size=28, color=theme["title_color"], bold=True,
             font_name=theme.get("font_title", "Microsoft JhengHei"), margin=0)
    next_y = MARGIN_T + 0.6
    if title_en:
        add_text(slide, title_en, MARGIN_L, next_y, SLIDE_W - MARGIN_L - MARGIN_R, 0.25,
                 font_size=11, color="8888AA", italic=True, margin=0)
        next_y += 0.22
    if subtitle:
        add_text(slide, subtitle, MARGIN_L, next_y, SLIDE_W - MARGIN_L - MARGIN_R, 0.4,
                 font_size=14, color=theme.get("subtitle_color", "8888AA"), margin=0)


def add_footer(slide, theme: dict, page_num: int, total: int):
    add_text(slide, f"{page_num} / {total}", SLIDE_W - 1.2, FOOTER_Y, 0.8, FOOTER_H,
             font_size=10, color=theme.get("footer_color", "666688"), align="right", margin=0)


def add_notes(slide, text: str):
    if text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text


def resolve_img(img_dir: Path, name: str) -> Path:
    return img_dir / name


def _get_aspect(img_path: Path) -> float:
    """Return width/height aspect ratio. Returns 1.5 if file missing."""
    if img_path and img_path.exists():
        from PIL import Image
        with Image.open(img_path) as img:
            w, h = img.size
        return w / h if h > 0 else 1.5
    return 1.5


def add_comparison_diagram(slide, comp: dict, x: float, y: float, w: float, h: float, theme: dict):
    """Draw a two-column comparison diagram with paired rows."""
    col_w = (w - 0.4) / 2
    rows = comp.get("rows", [])
    n = len(rows)

    # Column headers
    left_title = comp.get("left_title", "")
    right_title = comp.get("right_title", "")
    add_text(slide, left_title, x, y, col_w, 0.4,
             font_size=14, color=theme["accent"], bold=True, align="center", margin=0)
    add_text(slide, right_title, x + col_w + 0.4, y, col_w, 0.4,
             font_size=14, color=theme.get("accent_red", "ff6b6b"), bold=True, align="center", margin=0)

    # Rows
    avail_h = h - 0.6
    row_h = min(0.85, (avail_h - (n - 1) * 0.12) / max(n, 1))
    gap = 0.12
    for i, row in enumerate(rows):
        ry = y + 0.5 + i * (row_h + gap)
        # Left box
        left_color = row.get("left_color", theme["accent"])
        add_rect(slide, x, ry, col_w, row_h, fill="22224a", line=left_color, radius=True)
        add_text(slide, row["left"], x + 0.08, ry + 0.04, col_w - 0.16, row_h - 0.08,
                 font_size=11, color=theme["body_color"], align="center",
                 valign=MSO_ANCHOR.MIDDLE, margin=0.02)
        # VS label
        add_text(slide, "vs", x + col_w + 0.08, ry + row_h * 0.3, 0.24, 0.3,
                 font_size=10, color="666688", align="center", margin=0)
        # Right box
        right_color = row.get("right_color", theme.get("accent_red", "ff6b6b"))
        add_rect(slide, x + col_w + 0.4, ry, col_w, row_h, fill="22224a", line=right_color, radius=True)
        add_text(slide, row["right"], x + col_w + 0.48, ry + 0.04, col_w - 0.16, row_h - 0.08,
                 font_size=11, color=theme["body_color"], align="center",
                 valign=MSO_ANCHOR.MIDDLE, margin=0.02)


# ── Slide Builders ───────────────────────────────────────────────────────

def build_cover(slide, cfg: dict, theme: dict, img_dir: Path):
    set_bg(slide, theme["bg"])
    # Decorative accent strip on right
    add_rect(slide, 10.0, 0, 3.333, SLIDE_H, fill=theme["accent"], transparency=85)

    # Main title (Chinese)
    add_text(slide, cfg["title"], 0.8, 1.0, 9.0, 1.2,
             font_size=32, color=theme["title_color"], bold=True,
             font_name=theme.get("font_title", "Microsoft JhengHei"), margin=0, line_spacing=1.15)
    # English title
    title_en = cfg.get("title_en", "")
    if title_en:
        add_text(slide, title_en, 0.8, 2.15, 9.0, 0.4,
                 font_size=14, color="8888AA", italic=True, margin=0)

    # Subtitle
    add_text(slide, cfg.get("subtitle", ""), 0.8, 2.6, 9.0, 0.5,
             font_size=16, color=theme.get("subtitle_color", "8888AA"), margin=0)

    # Reporter + Date
    reporter = cfg.get("reporter", "")
    date_display = cfg.get("date_display", "")
    info_parts = []
    if reporter:
        info_parts.append(f"報告人：{reporter}")
    if date_display:
        info_parts.append(f"日期：{date_display}")
    info = cfg.get("info", "")
    if info_parts:
        info = " | ".join(info_parts) + "\n" + info
    add_text(slide, info, 0.8, 3.3, 9.0, 0.8,
             font_size=13, color="AAAACC", margin=0)

    # Summary box
    summary = cfg.get("summary", "")
    if summary:
        add_rect(slide, 0.8, 4.6, 8.5, 1.2, fill="22224a", line=theme["accent"], radius=True)
        add_text(slide, summary, 1.0, 4.75, 8.1, 0.9,
                 font_size=14, color="E0E0E0", margin=0.05)

    # Cover image (right side, within accent strip area)
    cover_image = cfg.get("cover_image")
    if cover_image:
        img_path = resolve_img(img_dir, cover_image)
        add_image(slide, img_path, 10.2, 1.0, 2.8, 4.5)


def build_text_image(slide, cfg: dict, theme: dict, img_dir: Path):
    """Adaptive layout: stacked for wide images, side-by-side for normal, full-width for no image."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    content_y = 1.4
    usable_w = SLIDE_W - MARGIN_L - MARGIN_R

    has_image = bool(cfg.get("image"))
    img_path = resolve_img(img_dir, cfg["image"]) if has_image else None
    aspect = _get_aspect(img_path) if has_image else None

    # Reserve bottom space for alert / footnote
    bottom_limit = 6.2
    if "alert" in cfg:
        bottom_limit = 6.1
    elif "footnote" in cfg:
        bottom_limit = 6.35

    bullets_en = cfg.get("bullets_en")

    if not has_image:
        # === NO IMAGE: full-width text + optional comparison diagram ===
        text_h = bottom_limit - content_y
        if "comparison" in cfg:
            text_w = 5.5
            if "bullets" in cfg:
                add_bullets(slide, cfg["bullets"], MARGIN_L, content_y, text_w, text_h,
                            font_size=15, color=theme["body_color"], items_en=bullets_en)
            add_comparison_diagram(slide, cfg["comparison"],
                                  MARGIN_L + text_w + 0.3, content_y, usable_w - text_w - 0.3, text_h, theme)
        else:
            if "bullets" in cfg:
                add_bullets(slide, cfg["bullets"], MARGIN_L, content_y, usable_w, text_h,
                            font_size=16, color=theme["body_color"], items_en=bullets_en)

    elif aspect and aspect > 1.8:
        # === WIDE IMAGE: stacked (bullets top, full-width image bottom) ===
        bullet_count = len(cfg.get("bullets", []))
        # Bilingual uses compact 12pt + 10pt to maximize image area
        per_bullet = 0.35 if bullets_en else 0.42
        text_h = min(2.8, max(1.2, bullet_count * per_bullet))
        bfont = 12 if bullets_en else 15
        if "bullets" in cfg:
            add_bullets(slide, cfg["bullets"], MARGIN_L, content_y, usable_w, text_h,
                        font_size=bfont, color=theme["body_color"], items_en=bullets_en)
        img_top = content_y + text_h + 0.10
        img_avail_h = bottom_limit - img_top - 0.05
        if img_avail_h > 0.5:
            add_image(slide, img_path, MARGIN_L, img_top, usable_w, img_avail_h)

    else:
        # === STANDARD: side-by-side ===
        content_h = bottom_limit - content_y
        text_w = 5.8
        img_x = 6.8
        img_w = 5.8
        if "bullets" in cfg:
            add_bullets(slide, cfg["bullets"], MARGIN_L, content_y, text_w, content_h,
                        font_size=16, color=theme["body_color"], items_en=bullets_en)
        if img_path:
            add_image(slide, img_path, img_x, content_y, img_w, content_h)

    # Bottom: alert box
    if "alert" in cfg:
        add_rect(slide, MARGIN_L, 6.3, usable_w, 0.55,
                 fill="3a1a1a", line=theme["accent_red"], radius=True)
        add_text(slide, cfg["alert"], MARGIN_L + 0.15, 6.38, usable_w - 0.3, 0.4,
                 font_size=13, color=theme["accent_red"], bold=True, margin=0)

    # Bottom: footnote
    if "footnote" in cfg:
        fn_y = 6.5 if "alert" not in cfg else 6.9
        add_text(slide, cfg["footnote"], MARGIN_L, fn_y, usable_w, 0.4,
                 font_size=11, color="8888AA", margin=0)


def build_text_image_wide(slide, cfg: dict, theme: dict, img_dir: Path):
    """Adaptive: side-by-side for portrait images, stacked for wide images."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    has_image = "image" in cfg
    img_path = resolve_img(img_dir, cfg["image"]) if has_image else None
    aspect = _get_aspect(img_path) if has_image else None

    if aspect and aspect < 1.2:
        # === PORTRAIT IMAGE: left table+numbers, right tall image ===
        left_w = 6.0
        img_x = 6.8
        img_w = SLIDE_W - MARGIN_R - img_x
        content_y = 1.4

        if "table" in cfg:
            add_table(slide, cfg["table"], MARGIN_L, content_y, left_w, 1.8, theme,
                      font_size=13, col_widths=cfg.get("col_widths"))

        if "key_numbers" in cfg:
            for i, kn in enumerate(cfg["key_numbers"]):
                y = 3.5 + i * 0.8
                add_text(slide, kn["label"], MARGIN_L, y, left_w, 0.3,
                         font_size=13, color="8888AA", margin=0)
                add_text(slide, kn["value"], MARGIN_L, y + 0.28, left_w, 0.4,
                         font_size=22, color=theme["accent"], bold=True, margin=0)

        img_h = 5.3
        add_image(slide, img_path, img_x, content_y, img_w, img_h)
    else:
        # === WIDE/NORMAL IMAGE: top table+numbers, bottom wide image ===
        if "table" in cfg:
            add_table(slide, cfg["table"], MARGIN_L, 1.4, 5.5, 1.8, theme,
                      font_size=13, col_widths=cfg.get("col_widths"))

        if "key_numbers" in cfg:
            for i, kn in enumerate(cfg["key_numbers"]):
                y = 1.4 + i * 0.7
                add_text(slide, kn["label"], 7.0, y, 3.5, 0.3,
                         font_size=13, color="8888AA", margin=0)
                add_text(slide, kn["value"], 7.0, y + 0.25, 3.5, 0.4,
                         font_size=24, color=theme["accent"], bold=True, margin=0)

        if has_image:
            img_y = cfg.get("image_y", 3.4)
            img_h = cfg.get("image_h", 3.3)
            add_image(slide, img_path,
                      MARGIN_L, img_y, SLIDE_W - MARGIN_L - MARGIN_R, img_h)

    if "footnote" in cfg:
        add_text(slide, cfg["footnote"], MARGIN_L, 6.75, SLIDE_W - 1.2, 0.3,
                 font_size=10, color="8888AA", margin=0)


def build_two_image(slide, cfg: dict, theme: dict, img_dir: Path):
    """Title + two images side by side + optional bottom bullets."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    img_y = 1.4
    img_h = cfg.get("image_h", 4.0)
    gap = 0.3
    img_w = (SLIDE_W - MARGIN_L - MARGIN_R - gap) / 2

    images = cfg.get("images", [])
    if len(images) >= 1:
        add_image(slide, resolve_img(img_dir, images[0]), MARGIN_L, img_y, img_w, img_h)
    if len(images) >= 2:
        add_image(slide, resolve_img(img_dir, images[1]), MARGIN_L + img_w + gap, img_y, img_w, img_h)

    # Captions under each image
    captions = cfg.get("captions", [])
    cap_y = img_y + img_h + 0.05
    for i, cap in enumerate(captions[:2]):
        cx = MARGIN_L + i * (img_w + gap)
        add_text(slide, cap, cx, cap_y, img_w, 0.3,
                 font_size=10, color="8888AA", align="center", margin=0)

    if "bullets" in cfg:
        add_bullets(slide, cfg["bullets"], MARGIN_L, 5.7, SLIDE_W - 1.2, 1.0,
                    font_size=13, color=theme["body_color"], items_en=cfg.get("bullets_en"))


def build_text_two_image(slide, cfg: dict, theme: dict, img_dir: Path):
    """Left: text/flow chart, Right: two stacked images."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, subtitle=cfg.get("subtitle"), title_en=cfg.get("title_en"))

    content_y = 1.5
    text_w = 6.0
    img_x = 6.8
    img_w = 5.6

    # Left: flow steps or bullets
    if "steps" in cfg:
        for i, step in enumerate(cfg["steps"]):
            sy = content_y + i * 1.05
            # Step number circle
            add_rect(slide, MARGIN_L, sy, 0.4, 0.4,
                     fill=theme["accent"], radius=True)
            add_text(slide, str(i + 1), MARGIN_L + 0.08, sy + 0.05, 0.25, 0.25,
                     font_size=14, color=theme["bg"], bold=True, align="center", margin=0)
            # Step text
            add_text(slide, step, MARGIN_L + 0.55, sy, text_w - 0.8, 0.9,
                     font_size=14, color=theme["body_color"], margin=0)

    if "bullets" in cfg:
        add_bullets(slide, cfg["bullets"], MARGIN_L, content_y, text_w, 3.5,
                    font_size=14, color=theme["body_color"], items_en=cfg.get("bullets_en"))

    # Right: two stacked images with optional captions
    images = cfg.get("images", [])
    captions = cfg.get("captions", [])
    cap_h = 0.18 if captions else 0
    img_h_each = 2.1 if captions else 2.3  # preserve v4 sizing when no captions
    for i, img_name in enumerate(images[:2]):
        slot_y = content_y + i * (img_h_each + cap_h + 0.15)
        if i < len(captions) and captions[i]:
            add_text(slide, captions[i], img_x, slot_y, img_w, cap_h,
                     font_size=10, color="8888AA", align="center", margin=0)
        add_image(slide, resolve_img(img_dir, img_name),
                  img_x, slot_y + cap_h, img_w, img_h_each)

    # Bottom quantification
    if "bottom_text" in cfg:
        add_text(slide, cfg["bottom_text"], MARGIN_L, 6.2, SLIDE_W - 1.2, 0.6,
                 font_size=12, color="AAAACC", margin=0)


def build_cards(slide, cfg: dict, theme: dict, img_dir: Path):
    """2x2 card layout for hypotheses/findings."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, subtitle=cfg.get("subtitle"), title_en=cfg.get("title_en"))

    cards = cfg.get("cards", [])
    positions = [(MARGIN_L, 1.5), (6.6, 1.5), (MARGIN_L, 3.8), (6.6, 3.8)]
    card_w = 5.7
    card_h = 2.0

    for (cx, cy), card in zip(positions, cards):
        # Card background
        status_color = card.get("status_color", theme["accent_red"])
        add_rect(slide, cx, cy, card_w, card_h, fill="22224a", line="333366", radius=True)
        # Status marker
        marker = card.get("marker", "✗")
        add_rect(slide, cx + 0.15, cy + 0.15, 0.35, 0.35, fill=status_color, radius=True)
        add_text(slide, marker, cx + 0.18, cy + 0.15, 0.3, 0.35,
                 font_size=16, color="FFFFFF", bold=True, align="center", margin=0)
        # Card title
        add_text(slide, card["title"], cx + 0.6, cy + 0.15, card_w - 0.9, 0.35,
                 font_size=15, color=theme["title_color"], bold=True, margin=0)
        # Card body
        add_text(slide, card["body"], cx + 0.15, cy + 0.6, card_w - 0.3, card_h - 0.8,
                 font_size=12, color=theme["body_color"], margin=0.05)

    # Bottom lesson box
    if "lesson" in cfg:
        add_rect(slide, MARGIN_L, 6.0, SLIDE_W - MARGIN_L - MARGIN_R, 0.65,
                 fill="1a2a3a", line=theme.get("accent_teal", "4ecdc4"), radius=True)
        add_text(slide, cfg["lesson"], MARGIN_L + 0.15, 6.1, SLIDE_W - 1.5, 0.45,
                 font_size=12, color=theme.get("accent_teal", "4ecdc4"), margin=0)


def build_table_text(slide, cfg: dict, theme: dict, img_dir: Path):
    """Table on top/left + explanatory text/numbers."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    # Table
    if "table" in cfg:
        tbl_rows = cfg["table"]
        tbl_y = cfg.get("table_y", 1.4)
        tbl_h = cfg.get("table_h", 1.6)
        add_table(slide, tbl_rows, MARGIN_L, tbl_y,
                  SLIDE_W - MARGIN_L - MARGIN_R, tbl_h, theme,
                  font_size=cfg.get("table_font_size", 13),
                  col_widths=cfg.get("col_widths"))

    # Below table: explanation block
    if "explanation" in cfg:
        exp = cfg["explanation"]
        ey = cfg.get("explanation_y", 3.3)
        add_rect(slide, MARGIN_L, ey, SLIDE_W - MARGIN_L - MARGIN_R, 2.0,
                 fill="22224a", line="333366", radius=True)
        for i, line in enumerate(exp):
            add_text(slide, line, MARGIN_L + 0.2, ey + 0.15 + i * 0.45,
                     SLIDE_W - 1.8, 0.4,
                     font_size=14, color=theme["body_color"], margin=0)

    # Bottom info
    if "bottom_text" in cfg:
        add_text(slide, cfg["bottom_text"], MARGIN_L, 5.8, SLIDE_W - 1.2, 0.8,
                 font_size=12, color="AAAACC", margin=0)


def build_numbered_list(slide, cfg: dict, theme: dict, img_dir: Path):
    """Numbered items with descriptions (for future directions etc.)."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    items = cfg.get("items", [])
    start_y = 1.4
    item_h = 0.85
    for i, item in enumerate(items):
        y = start_y + i * item_h
        # Number badge
        add_rect(slide, MARGIN_L, y, 0.45, 0.45,
                 fill=theme["accent"], radius=True)
        add_text(slide, str(i + 1), MARGIN_L + 0.1, y + 0.06, 0.25, 0.3,
                 font_size=16, color=theme["bg"], bold=True, align="center", margin=0)
        # Item title
        add_text(slide, item["title"], MARGIN_L + 0.6, y, 11.0, 0.35,
                 font_size=16, color=theme["title_color"], bold=True, margin=0)
        # Item description
        add_text(slide, item["desc"], MARGIN_L + 0.6, y + 0.35, 11.0, 0.5,
                 font_size=13, color=theme["body_color"], margin=0)

    # Bottom section
    if "bottom_section" in cfg:
        bs = cfg["bottom_section"]
        by = start_y + len(items) * item_h + 0.15
        by = min(by, 5.9)  # 不超出安全區域
        available_h = FOOTER_Y - by - 0.1
        box_h = min(1.2, available_h)
        add_rect(slide, MARGIN_L, by, SLIDE_W - MARGIN_L - MARGIN_R, box_h,
                 fill="1a2a3a", line=theme.get("accent_teal", "4ecdc4"), radius=True)
        add_text(slide, bs.get("title", ""), MARGIN_L + 0.15, by + 0.08, 5.0, 0.25,
                 font_size=13, color=theme.get("accent_teal", "4ecdc4"), bold=True, margin=0)
        if "bullets" in bs:
            add_bullets(slide, bs["bullets"], MARGIN_L + 0.15, by + 0.35, SLIDE_W - 1.5, box_h - 0.45,
                        font_size=11, color=theme["body_color"], items_en=bs.get("bullets_en"))


def build_summary(slide, cfg: dict, theme: dict, img_dir: Path):
    """Left: comparison table, Right: layered conclusions."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    # Left: table
    if "table" in cfg:
        add_table(slide, cfg["table"], MARGIN_L, 1.4, 5.8, 3.5, theme,
                  font_size=13, col_widths=cfg.get("col_widths"))

    # Right: conclusion layers
    conclusions = cfg.get("conclusions", [])
    for i, c in enumerate(conclusions):
        y = 1.4 + i * 0.68
        icon = c.get("icon", "")
        text = f"{icon} {c['text']}" if icon else c["text"]
        clr = c.get("color", theme["body_color"])
        add_text(slide, text, 6.8, y, 5.5, 0.6,
                 font_size=14, color=clr, bold=c.get("bold", False), margin=0)

    # Limitations
    if "limitations" in cfg:
        add_rect(slide, MARGIN_L, 5.5, SLIDE_W - MARGIN_L - MARGIN_R, 1.0,
                 fill="22224a", line="555577", radius=True)
        add_text(slide, "Limitations", MARGIN_L + 0.15, 5.55, 2.0, 0.3,
                 font_size=13, color=theme["accent_red"], bold=True, margin=0)
        add_text(slide, cfg["limitations"], MARGIN_L + 0.15, 5.85, SLIDE_W - 1.5, 0.55,
                 font_size=11, color="CCCCDD", margin=0)


def build_two_column(slide, cfg: dict, theme: dict, img_dir: Path):
    """Left column: table/list, Right column: discussion items."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    col_w = 5.8
    gap = 0.5
    col1_x = MARGIN_L
    col2_x = MARGIN_L + col_w + gap

    # Left column: action table
    if "left_table" in cfg:
        add_text(slide, cfg.get("left_title", ""), col1_x, 1.3, col_w, 0.35,
                 font_size=16, color=theme["accent"], bold=True, margin=0)
        add_table(slide, cfg["left_table"], col1_x, 1.7, col_w, 2.8, theme,
                  font_size=12, col_widths=cfg.get("left_col_widths"))

    # Right column: discussion questions
    if "right_items" in cfg:
        add_text(slide, cfg.get("right_title", ""), col2_x, 1.3, col_w, 0.35,
                 font_size=16, color=theme["accent"], bold=True, margin=0)
        for i, item in enumerate(cfg["right_items"]):
            y = 1.8 + i * 1.3
            add_rect(slide, col2_x, y, col_w, 1.1,
                     fill="22224a", line="333366", radius=True)
            add_text(slide, f"Q{i + 1}", col2_x + 0.15, y + 0.1, 0.4, 0.3,
                     font_size=14, color=theme["accent"], bold=True, margin=0)
            add_text(slide, item, col2_x + 0.5, y + 0.1, col_w - 0.7, 0.9,
                     font_size=13, color=theme["body_color"], margin=0)

    # Bottom thanks
    if "footer_text" in cfg:
        add_text(slide, cfg["footer_text"], MARGIN_L, 6.5, SLIDE_W - 1.2, 0.4,
                 font_size=18, color=theme["accent"], bold=True, align="center", margin=0)


def build_key_numbers(slide, cfg: dict, theme: dict, img_dir: Path):
    """Large key numbers with image on right — for ruling out hypotheses."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    # Left: key number cards
    numbers = cfg.get("numbers", [])
    for i, num in enumerate(numbers):
        y = 1.4 + i * 1.4
        add_rect(slide, MARGIN_L, y, 5.8, 1.2, fill="22224a", line="333366", radius=True)
        add_text(slide, num["label"], MARGIN_L + 0.2, y + 0.1, 3.5, 0.35,
                 font_size=13, color="8888AA", margin=0)
        add_text(slide, num["value"], MARGIN_L + 0.2, y + 0.4, 3.5, 0.5,
                 font_size=28, color=num.get("color", theme["accent"]), bold=True, margin=0)
        if "detail" in num:
            add_text(slide, num["detail"], MARGIN_L + 0.2, y + 0.85, 5.2, 0.3,
                     font_size=11, color="AAAACC", margin=0)

    # Right: image
    if "image" in cfg:
        add_image(slide, resolve_img(img_dir, cfg["image"]), 6.8, 1.4, 5.8, 4.8)

    if "bottom_text" in cfg:
        add_text(slide, cfg["bottom_text"], MARGIN_L, 6.3, SLIDE_W - 1.2, 0.5,
                 font_size=12, color="AAAACC", margin=0)


def build_section_divider(slide, cfg: dict, theme: dict, img_dir: Path):
    """Section divider page: dark background, centered title, subtitle, progress dots."""
    set_bg(slide, "#151528")

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R

    # Title — 36pt, white, bold, vertically centered (upper)
    add_text(slide, cfg["title"], MARGIN_L, 2.2, usable_w, 1.2,
             font_size=36, color="FFFFFF", bold=True, align="center",
             font_name=theme.get("font_title", "Microsoft JhengHei"),
             valign=MSO_ANCHOR.MIDDLE, margin=0)

    # English title
    title_en = cfg.get("title_en", "")
    if title_en:
        add_text(slide, title_en, MARGIN_L, 3.2, usable_w, 0.4,
                 font_size=14, color="666688", italic=True, align="center", margin=0)

    # Subtitle — 18pt, muted
    subtitle = cfg.get("subtitle", "")
    if subtitle:
        add_text(slide, subtitle, MARGIN_L, 3.5, usable_w, 0.6,
                 font_size=18, color="8888AA", align="center", margin=0)

    # Progress dots at bottom
    act = cfg.get("act", 1)
    total_acts = cfg.get("total_acts", 3)
    dot_r = 0.18
    dot_gap = 0.55
    total_dot_w = total_acts * dot_r * 2 + (total_acts - 1) * (dot_gap - dot_r * 2)
    dot_start_x = (SLIDE_W - total_dot_w) / 2
    dot_y = 6.5
    for i in range(total_acts):
        dx = dot_start_x + i * dot_gap
        dot_color = theme.get("accent", "00d4ff") if (i + 1) == act else "333366"
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(dx), Inches(dot_y), Inches(dot_r * 2), Inches(dot_r * 2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_rgb(dot_color)
        shape.line.fill.background()


def build_big_number(slide, cfg: dict, theme: dict, img_dir: Path):
    """Big number emphasis page: large number left, optional image right."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    has_image = bool(cfg.get("image"))
    left_w = 7.0 if has_image else SLIDE_W - MARGIN_L - MARGIN_R
    left_x = MARGIN_L

    # Main number — 56pt accent bold, centered in left area
    number_color = cfg.get("number_color", theme.get("accent", "00d4ff"))
    add_text(slide, cfg["number"], left_x, 2.0, left_w, 1.4,
             font_size=56, color=number_color, bold=True, align="center",
             valign=MSO_ANCHOR.MIDDLE, margin=0)

    # Number label — 18pt below number
    label = cfg.get("number_label", "")
    if label:
        add_text(slide, label, left_x, 3.4, left_w, 0.5,
                 font_size=18, color="AAAACC", align="center", margin=0)

    # Sub-numbers — horizontal layout below label
    sub_numbers = cfg.get("sub_numbers", [])
    if sub_numbers:
        n_sub = len(sub_numbers)
        sub_w = left_w / max(n_sub, 1)
        for i, sn in enumerate(sub_numbers):
            sx = left_x + i * sub_w
            sn_color = sn.get("color", theme.get("accent", "00d4ff"))
            add_text(slide, sn["value"], sx, 4.2, sub_w, 0.7,
                     font_size=28, color=sn_color, bold=True, align="center",
                     valign=MSO_ANCHOR.MIDDLE, margin=0)
            add_text(slide, sn.get("label", ""), sx, 4.9, sub_w, 0.4,
                     font_size=13, color="8888AA", align="center", margin=0)

    # Right image (if present)
    if has_image:
        img_path = resolve_img(img_dir, cfg["image"])
        img_x = MARGIN_L + left_w + 0.3
        img_w = SLIDE_W - MARGIN_R - img_x
        add_image(slide, img_path, img_x, 1.4, img_w, 5.0)


def build_full_image(slide, cfg: dict, theme: dict, img_dir: Path):
    """Full image page: title on top, image fills content area, optional caption."""
    set_bg(slide, theme["bg"])

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R

    # Title — 28pt bold
    add_text(slide, cfg["title"], MARGIN_L, MARGIN_T, usable_w, TITLE_H,
             font_size=28, color=theme["title_color"], bold=True,
             font_name=theme.get("font_title", "Microsoft JhengHei"), margin=0)

    # English title
    title_en = cfg.get("title_en", "")
    if title_en:
        add_text(slide, title_en, MARGIN_L, MARGIN_T + 0.6, usable_w, 0.25,
                 font_size=11, color="8888AA", italic=True, margin=0)

    # Image — fills content area
    img_path = resolve_img(img_dir, cfg.get("image", ""))
    add_image(slide, img_path, MARGIN_L, 1.4, usable_w, 4.8)

    # Caption (if present) — 11pt muted, below image
    caption = cfg.get("caption", "")
    if caption:
        add_text(slide, caption, MARGIN_L, 6.3, usable_w, 0.4,
                 font_size=11, color="8888AA", align="center", margin=0)


def build_process_flow(slide, cfg: dict, theme: dict, img_dir: Path):
    """Process flow page: horizontal steps with arrows."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    steps = cfg.get("steps", [])
    n = len(steps)
    if n == 0:
        return

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R
    # Layout: each step is a box, arrows between them
    arrow_w = 0.4
    total_arrow_w = (n - 1) * arrow_w
    step_w = min(1.8, (usable_w - total_arrow_w - 0.2) / max(n, 1))
    step_h = 1.3
    total_w = n * step_w + (n - 1) * arrow_w
    start_x = MARGIN_L + (usable_w - total_w) / 2
    # Vertically center in content area (y: 1.4 to 6.2)
    center_y = 1.4 + (6.2 - 1.4 - step_h) / 2

    accent = theme.get("accent", "00d4ff")

    for i, step in enumerate(steps):
        sx = start_x + i * (step_w + arrow_w)

        # Rounded rectangle for step
        add_rect(slide, sx, center_y, step_w, step_h,
                 fill="22224a", line=accent, radius=True)

        # Step text — 13pt centered
        step_label = step.get("label", "") if isinstance(step, dict) else str(step)
        add_text(slide, step_label, sx + 0.05, center_y + 0.1, step_w - 0.1, step_h - 0.2,
                 font_size=13, color=theme.get("body_color", "E0E0E0"),
                 align="center", valign=MSO_ANCHOR.MIDDLE, margin=0.02)

        # Arrow between steps
        if i < n - 1:
            ax = sx + step_w
            add_text(slide, "\u2192", ax, center_y + step_h * 0.25, arrow_w, step_h * 0.5,
                     font_size=18, color=accent, bold=True, align="center",
                     valign=MSO_ANCHOR.MIDDLE, margin=0)

    # Bottom text (if present)
    if "bottom_text" in cfg:
        add_text(slide, cfg["bottom_text"], MARGIN_L, 6.2, usable_w, 0.6,
                 font_size=12, color="AAAACC", margin=0)


def build_before_after(slide, cfg: dict, theme: dict, img_dir: Path):
    """Before/after comparison page: two cards with arrow in between."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R
    card_w = (usable_w - 1.2) / 2  # space for arrow in the middle
    card_h = 3.8
    card_y = 1.5
    accent_red = theme.get("accent_red", "ff6b6b")
    accent_teal = theme.get("accent_teal", "4ecdc4")

    # ---- Before card (left) ----
    before_x = MARGIN_L
    add_rect(slide, before_x, card_y, card_w, card_h,
             fill="22224a", line=accent_red, radius=True)
    before_title = cfg.get("before_title", "Before")
    add_text(slide, before_title, before_x + 0.15, card_y + 0.2, card_w - 0.3, 0.4,
             font_size=16, color=accent_red, bold=True, margin=0)
    before_items = cfg.get("before_items", [])
    for j, item in enumerate(before_items):
        add_text(slide, f"\u2022 {item}", before_x + 0.2, card_y + 0.7 + j * 0.55,
                 card_w - 0.4, 0.5,
                 font_size=14, color=theme.get("body_color", "E0E0E0"), margin=0)

    # ---- Arrow in middle ----
    arrow_x = before_x + card_w
    arrow_w = 1.2
    add_text(slide, "\u2192", arrow_x, card_y + card_h * 0.3, arrow_w, 0.8,
             font_size=36, color=theme.get("accent", "00d4ff"), bold=True,
             align="center", valign=MSO_ANCHOR.MIDDLE, margin=0)
    arrow_text = cfg.get("arrow_text", "")
    if arrow_text:
        add_text(slide, arrow_text, arrow_x, card_y + card_h * 0.55, arrow_w, 0.4,
                 font_size=12, color="8888AA", align="center", margin=0)

    # ---- After card (right) ----
    after_x = before_x + card_w + arrow_w
    add_rect(slide, after_x, card_y, card_w, card_h,
             fill="22224a", line=accent_teal, radius=True)
    after_title = cfg.get("after_title", "After")
    add_text(slide, after_title, after_x + 0.15, card_y + 0.2, card_w - 0.3, 0.4,
             font_size=16, color=accent_teal, bold=True, margin=0)
    after_items = cfg.get("after_items", [])
    for j, item in enumerate(after_items):
        add_text(slide, f"\u2022 {item}", after_x + 0.2, card_y + 0.7 + j * 0.55,
                 card_w - 0.4, 0.5,
                 font_size=14, color=theme.get("body_color", "E0E0E0"), margin=0)

    # ---- Bottom text ----
    if "bottom_text" in cfg:
        add_rect(slide, MARGIN_L, 5.9, usable_w, 0.6,
                 fill="1a2a3a", line=accent_teal, radius=True)
        add_text(slide, cfg["bottom_text"], MARGIN_L + 0.15, 5.98, usable_w - 0.3, 0.45,
                 font_size=13, color=accent_teal, margin=0)


def build_funnel(slide, cfg: dict, theme: dict, img_dir: Path):
    """Funnel / hypothesis rejection page: narrowing bars from top to bottom."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, subtitle=cfg.get("subtitle"), title_en=cfg.get("title_en"))

    items = cfg.get("items", [])
    n = len(items)
    if n == 0:
        return

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R
    accent_red = theme.get("accent_red", "ff6b6b")
    accent_teal = theme.get("accent_teal", "4ecdc4")

    # Funnel bars: widest at top, narrowest at bottom
    start_y = 1.8 if not cfg.get("subtitle") else 2.1
    bar_h = 0.75
    bar_gap = 0.12
    avail_h = 5.5 - start_y  # reserve space for conclusion
    if n > 0:
        bar_h = min(bar_h, (avail_h - (n - 1) * bar_gap) / n)
        bar_h = max(bar_h, 0.5)

    max_w = usable_w
    min_w = usable_w * 0.55  # narrowest bar is 55% of widest
    for i, item in enumerate(items):
        # Funnel effect: linearly shrink width
        frac = i / max(n - 1, 1)
        bar_w = max_w - frac * (max_w - min_w)
        bx = MARGIN_L + (usable_w - bar_w) / 2
        by = start_y + i * (bar_h + bar_gap)

        status = item.get("status", "rejected")
        bar_fill = "2a1a1a" if status == "rejected" else "1a2a2a"
        bar_line = accent_red if status == "rejected" else accent_teal
        add_rect(slide, bx, by, bar_w, bar_h,
                 fill=bar_fill, line=bar_line, radius=True)

        # Label (left portion)
        label = item.get("label", "")
        add_text(slide, label, bx + 0.15, by + 0.04, bar_w * 0.35, bar_h - 0.08,
                 font_size=12, color=theme.get("body_color", "E0E0E0"), bold=True,
                 valign=MSO_ANCHOR.MIDDLE, margin=0.02)

        # Initial value (teal) -> final value (red)
        initial = item.get("initial", "")
        final = item.get("final", "")
        mid_text = f"{initial} {final}" if initial and final else initial or final
        add_text(slide, mid_text, bx + bar_w * 0.36, by + 0.04, bar_w * 0.3, bar_h - 0.08,
                 font_size=12, color="AAAACC",
                 valign=MSO_ANCHOR.MIDDLE, margin=0.02)

        # Reason
        reason = item.get("reason", "")
        if reason:
            add_text(slide, reason, bx + bar_w * 0.67, by + 0.04, bar_w * 0.23, bar_h - 0.08,
                     font_size=11, color="8888AA",
                     valign=MSO_ANCHOR.MIDDLE, margin=0.02)

        # Status marker
        marker = "\u2717" if status == "rejected" else "\u2713"
        marker_color = accent_red if status == "rejected" else accent_teal
        add_text(slide, marker, bx + bar_w - 0.45, by + 0.04, 0.35, bar_h - 0.08,
                 font_size=16, color=marker_color, bold=True, align="center",
                 valign=MSO_ANCHOR.MIDDLE, margin=0)

    # Conclusion box at bottom
    conclusion = cfg.get("conclusion", "")
    if conclusion:
        con_y = start_y + n * (bar_h + bar_gap) + 0.15
        con_y = max(con_y, 5.6)
        add_rect(slide, MARGIN_L, con_y, usable_w, 0.65,
                 fill="1a2a3a", line=accent_teal, radius=True)
        add_text(slide, conclusion, MARGIN_L + 0.15, con_y + 0.08, usable_w - 0.3, 0.5,
                 font_size=12, color=accent_teal, margin=0)


def build_definition(slide, cfg: dict, theme: dict, img_dir: Path):
    """Definition page: term definition with visual example and formula."""
    set_bg(slide, theme["bg"])
    add_slide_title(slide, cfg["title"], theme, title_en=cfg.get("title_en"))

    usable_w = SLIDE_W - MARGIN_L - MARGIN_R
    content_y = 1.4

    # Left side: definitions (compact layout to fit 3 definitions + criteria in 4.8")
    left_w = 5.5
    definitions = cfg.get("definitions", [])
    dy = content_y
    for defn in definitions:
        term = defn.get("term", "")
        desc = defn.get("desc", "")
        formula = defn.get("formula", "")

        # Term in accent color
        add_text(slide, term, MARGIN_L, dy, left_w, 0.28,
                 font_size=14, color=theme.get("accent", "00d4ff"), bold=True, margin=0)
        dy += 0.28
        # Description (compact: 12pt, reduced height)
        if desc:
            add_text(slide, desc, MARGIN_L + 0.1, dy, left_w - 0.1, 0.45,
                     font_size=11, color=theme.get("body_color", "E0E0E0"), margin=0)
            dy += 0.40
        # Formula box (compact)
        if formula:
            add_rect(slide, MARGIN_L + 0.1, dy, left_w - 0.2, 0.35,
                     fill="22224a", line=theme.get("accent_gold", "ffd93d"), radius=True)
            add_text(slide, formula, MARGIN_L + 0.25, dy + 0.03, left_w - 0.5, 0.29,
                     font_size=12, color=theme.get("accent_gold", "ffd93d"), bold=True,
                     align="center", valign=MSO_ANCHOR.MIDDLE, margin=0)
            dy += 0.40
        dy += 0.08  # compact gap between definitions

    # Criteria box (if present)
    criteria = cfg.get("criteria", [])
    if criteria:
        crit_y = max(dy, 4.0)
        crit_h = 0.30 + len(criteria) * 0.26
        add_rect(slide, MARGIN_L, crit_y, left_w, crit_h,
                 fill="1a2a3a", line=theme.get("accent_teal", "4ecdc4"), radius=True)
        add_text(slide, "判定標準 Classification Criteria", MARGIN_L + 0.15, crit_y + 0.04, left_w - 0.3, 0.22,
                 font_size=11, color=theme.get("accent_teal", "4ecdc4"), bold=True, margin=0)
        for i, crit in enumerate(criteria):
            add_text(slide, f"• {crit}", MARGIN_L + 0.15, crit_y + 0.28 + i * 0.26, left_w - 0.3, 0.24,
                     font_size=10, color=theme.get("body_color", "E0E0E0"), margin=0)

    # Right side: image
    if cfg.get("image"):
        img_path = resolve_img(img_dir, cfg["image"])
        add_image(slide, img_path, 6.3, content_y, 6.0, 4.8)

    # Bottom note
    if "bottom_text" in cfg:
        add_text(slide, cfg["bottom_text"], MARGIN_L, 6.3, usable_w, 0.5,
                 font_size=11, color="AAAACC", margin=0)


# ── Builder Dispatch ─────────────────────────────────────────────────────

BUILDERS = {
    "cover": build_cover,
    "text_image": build_text_image,
    "text_image_wide": build_text_image_wide,
    "two_image": build_two_image,
    "text_two_image": build_text_two_image,
    "cards": build_cards,
    "table_text": build_table_text,
    "numbered_list": build_numbered_list,
    "summary": build_summary,
    "two_column": build_two_column,
    "key_numbers": build_key_numbers,
    "section_divider": build_section_divider,
    "big_number": build_big_number,
    "full_image": build_full_image,
    "process_flow": build_process_flow,
    "before_after": build_before_after,
    "funnel": build_funnel,
    "definition": build_definition,
}


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Generate LOH Weekly Report PPTX")
    parser.add_argument("config", type=Path, help="JSON config file")
    parser.add_argument("output", type=Path, help="Output PPTX path")
    args = parser.parse_args()

    cfg = json.loads(args.config.read_text(encoding="utf-8"))
    theme = cfg["theme"]
    img_dir = args.config.parent / "images"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = cfg["meta"]["title"]
    prs.core_properties.author = cfg["meta"].get("author", "")

    deck = cfg["deck"]
    total = len(deck)

    for idx, slide_cfg in enumerate(deck, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide_type = slide_cfg["type"]
        builder = BUILDERS.get(slide_type)
        if builder is None:
            print(f"WARNING: Unknown slide type '{slide_type}' for slide {idx}, skipping")
            set_bg(slide, theme["bg"])
            add_text(slide, f"[Unknown type: {slide_type}]", 2, 3, 8, 1,
                     font_size=24, color="FF6B6B", align="center")
        else:
            builder(slide, slide_cfg, theme, img_dir)

        add_footer(slide, theme, idx, total)
        add_notes(slide, slide_cfg.get("notes", ""))

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"Generated: {args.output} ({total} slides)")


if __name__ == "__main__":
    main()
