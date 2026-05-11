"""assertion_title.py — Enforce thesis-style slide titles (myPPT playbook §5).

PLOS Ten Simple Rules + Alley Assertion-Evidence: each slide heading must be
a *full thesis sentence* (a claim) — not a generic label like "Results",
"Discussion", "Background". This helper wraps add_title() with a guard that
rejects titles failing a heuristic assertion check.

Usage:

    >>> from ppt_toolkit.assertion_title import assertion_title
    >>> assertion_title(slide, "V5 三層投票把 17.3:1 降回 1:1")  # OK
    >>> assertion_title(slide, "結果")                            # raises
"""

from __future__ import annotations

import re
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from ppt_toolkit.text_helpers import (
    DEFAULT_CJK_FONT,
    DEFAULT_LATIN_FONT,
    PALETTE,
    add_text_with_fallback,
)

# Generic labels that fail the assertion check (English + 繁/簡中文).
_GENERIC_LABELS = {
    # English
    "results", "discussion", "background", "introduction", "conclusion",
    "method", "methods", "summary", "overview", "analysis", "data",
    "approach", "challenges", "future work", "next steps", "agenda",
    # Chinese
    "結果", "討論", "背景", "簡介", "結論", "方法", "總結", "概述",
    "分析", "資料", "數據", "做法", "挑戰", "未來工作", "後續步驟",
    "議程", "目錄", "前言", "本文", "正文",
}

# Verbs / claim markers that strengthen an assertion sentence.
_ASSERTION_HINTS_EN = re.compile(
    r"\b(is|are|was|were|shows?|demonstrate(s|d)?|reveal(s|ed)?|"
    r"prove(s|d|n)?|fix(es|ed)?|reduce(s|d)?|increase(s|d)?|"
    r"confirm(s|ed)?|achieve(s|d)?|enable(s|d)?|cause(s|d)?|"
    r"trigger(s|ed)?|expose(s|d)?|OK|NG|PASS|FAIL|"
    r"unlock(s|ed)?|drop(s|ped)?|raise(s|d)?)\b",
    re.IGNORECASE,
)
_ASSERTION_HINTS_CN = re.compile(
    r"(是|為|不是|降至|降為|升至|升為|提升|降低|減少|增加|"
    r"修補|修正|解決|觸發|暴露|證實|證明|顯示|揭示|表示|"
    r"通過|未通過|不可信|可信|失衡|平衡|翻轉|不變|相同)"
)


def is_assertion_sentence(text: str, min_len: int = 6) -> bool:
    """Return True if *text* looks like a thesis-style assertion.

    Heuristic, not perfect — designed to catch the worst offenders
    (single-word labels, "Results", etc.). Caller can override by passing
    ``allow_label=True`` to :func:`assertion_title`.

    Rules:
        1. Must be longer than ``min_len`` chars (after strip).
        2. Must NOT match the _GENERIC_LABELS set (case-insensitive).
        3. Must contain at least one assertion verb (English or CJK marker)
           OR a colon / arrow / equals sign suggesting a claim:
           "17.3:1 → 1:1", "AMB% = 8.0%", "purity ≤ 0.95".
    """
    s = text.strip()
    if len(s) < min_len:
        return False
    low = s.lower().strip("。.、，, ")
    if low in _GENERIC_LABELS:
        return False
    if _ASSERTION_HINTS_EN.search(s) or _ASSERTION_HINTS_CN.search(s):
        return True
    # Numerical / comparison style assertion (e.g. "17.3 : 1 → 1 : 1")
    if re.search(r"[→⟶=≤≥<>%‰]", s) or re.search(r"\d", s):
        return True
    return False


def assertion_title(
    slide,
    thesis_sentence: str,
    subtitle: Optional[str] = None,
    bar_height_in: float = 0.85,
    bar_color: Optional[RGBColor] = None,
    title_color: Optional[RGBColor] = None,
    title_size_pt: int = 26,
    subtitle_color: Optional[RGBColor] = None,
    subtitle_size_pt: int = 14,
    allow_label: bool = False,
):
    """Render a thesis-style title bar.

    Parameters
    ----------
    slide : pptx.slide.Slide
    thesis_sentence : str
        Full assertion sentence (e.g. "V5 把 17.3:1 降回 1:1").
        Must NOT be a generic label unless ``allow_label=True``.
    subtitle : str, optional
        Soft secondary line (e.g. caveat, citation). Rendered smaller below.
    bar_height_in, bar_color : layout knobs.
    title_color, title_size_pt : title typography.
    subtitle_color, subtitle_size_pt : subtitle typography.
    allow_label : bool
        Override the assertion check. Use sparingly (cover slide acknowledgements
        page, etc.). Logs a warning rather than raising.

    Raises
    ------
    ValueError
        Title fails the assertion check and ``allow_label`` is False.
    """
    if not allow_label and not is_assertion_sentence(thesis_sentence):
        raise ValueError(
            f"Title is not a thesis-style assertion: {thesis_sentence!r}. "
            f"Replace with a full claim (e.g. 'V5 把 17.3:1 降回 1:1' instead "
            f"of '結果'). Pass allow_label=True to bypass for cover slides."
        )

    # Local import so this module does not require a global slide width.
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Pt

    bar_color = bar_color or PALETTE["title_bg"]
    title_color = title_color or PALETTE["title_fg"]
    subtitle_color = subtitle_color or RGBColor(0xCF, 0xE2, 0xF3)

    bar_h = Inches(bar_height_in)
    # Use slide width from presentation if accessible; fall back to 13.333in.
    try:
        slide_w = slide.part.package.presentation.slide_width
    except Exception:
        slide_w = Inches(13.333)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, bar_h)
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = ""

    add_text_with_fallback(
        tf,
        thesis_sentence,
        font_size=title_size_pt,
        bold=True,
        color=title_color,
        paragraph_index=0,
        latin_font=DEFAULT_LATIN_FONT,
        cjk_font=DEFAULT_CJK_FONT,
    )
    if subtitle:
        add_text_with_fallback(
            tf,
            subtitle,
            font_size=subtitle_size_pt,
            bold=False,
            color=subtitle_color,
        )
    return bar
