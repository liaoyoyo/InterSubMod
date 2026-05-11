"""text_helpers.py — Reusable PPTX text/image/note primitives.

Extracted from v2 build_pptx.py (lines 109-271) without altering behaviour.
Usage:

    from ppt_toolkit.text_helpers import (
        add_text_with_fallback, fit_image_within, set_speaker_note,
        DEFAULT_LATIN_FONT, DEFAULT_CJK_FONT,
    )

All functions are pure helpers operating on python-pptx Slide / TextFrame
objects. They are independent of any specific deck and can be reused across
projects (v3 build, future weekly reports, etc.).
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Optional, Tuple

from PIL import Image
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Defaults & constants
# ---------------------------------------------------------------------------
DEFAULT_LATIN_FONT = "Arial"
DEFAULT_CJK_FONT = "Droid Sans Fallback"
DEFAULT_CJK_FONT_PATH = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"

#: Palette tokens shared by myPPT skill style_library/colors/palette.yaml.
#: Kept as RGBColor for direct python-pptx use; YAML loader returns hex.
PALETTE = {
    "title_bg": RGBColor(0x1F, 0x4E, 0x79),
    "title_fg": RGBColor(0xFF, 0xFF, 0xFF),
    "text": RGBColor(0x1F, 0x1F, 0x1F),
    "accent": RGBColor(0xC0, 0x39, 0x2B),
    "green": RGBColor(0x27, 0xAE, 0x60),
    "amber": RGBColor(0xE6, 0x7E, 0x22),
    "gray": RGBColor(0x7F, 0x8C, 0x8D),
    "light": RGBColor(0xEC, 0xF0, 0xF1),
    "border": RGBColor(0xBD, 0xC3, 0xC7),
    "blue": RGBColor(0x29, 0x80, 0xB9),
    "purple": RGBColor(0x8E, 0x44, 0xAD),
}

#: Unicode blocks treated as CJK for per-character font fallback.
#: Source: v2 build_pptx.py CJK_RANGES (verified against UTR-50).
CJK_RANGES: Tuple[Tuple[int, int], ...] = (
    (0x3000, 0x303F),  # CJK Symbols and Punctuation
    (0x3040, 0x30FF),  # Hiragana + Katakana
    (0x3400, 0x4DBF),  # CJK Unified Ideographs Extension A
    (0x4E00, 0x9FFF),  # CJK Unified Ideographs
    (0xF900, 0xFAFF),  # CJK Compatibility Ideographs
    (0xFE30, 0xFE4F),  # CJK Compatibility Forms
    (0xFF00, 0xFFEF),  # Halfwidth and Fullwidth Forms
)


# ---------------------------------------------------------------------------
# CJK detection (private)
# ---------------------------------------------------------------------------
def _is_cjk(ch: str) -> bool:
    """Return True if codepoint falls in any CJK_RANGES block."""
    cp = ord(ch)
    for lo, hi in CJK_RANGES:
        if lo <= cp <= hi:
            return True
    return False


def _segments(text: str) -> Iterable[Tuple[str, bool]]:
    """Split text into (chunk, is_cjk) segments by alternating script.

    Mirrors v2 build_pptx.py._segments — minimum-segment algorithm so that
    runs of consecutive Latin or CJK characters share one font fallback.
    """
    if not text:
        return []
    out = []
    buf = [text[0]]
    cur_is_cjk = _is_cjk(text[0])
    for ch in text[1:]:
        if _is_cjk(ch) == cur_is_cjk:
            buf.append(ch)
        else:
            out.append(("".join(buf), cur_is_cjk))
            buf = [ch]
            cur_is_cjk = _is_cjk(ch)
    out.append(("".join(buf), cur_is_cjk))
    return out


# ---------------------------------------------------------------------------
# Public helpers
# ---------------------------------------------------------------------------
def add_text_with_fallback(
    text_frame,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    latin_font: str = DEFAULT_LATIN_FONT,
    cjk_font: str = DEFAULT_CJK_FONT,
    align=PP_ALIGN.LEFT,
    paragraph_index: Optional[int] = None,
    italic: bool = False,
):
    """Append a paragraph with per-character Latin / CJK font fallback.

    The XML manipulation injects ``<a:latin>`` and ``<a:ea>`` runs so PowerPoint
    selects the Latin face for ASCII glyphs and the East-Asian face for CJK
    glyphs *within the same run*. This avoids the "tofu / mojibake" failure
    mode where a CJK-only typeface paints English as tofu boxes.

    Parameters
    ----------
    text_frame : pptx.text.text.TextFrame
        Target text frame (slide.shapes.add_textbox(...).text_frame).
    text : str
        Paragraph content (Latin + CJK mix allowed).
    font_size : int
        Pt size for all runs in this paragraph.
    bold, italic : bool
        Run-level style.
    color : RGBColor, optional
        Run colour. Falls back to inheriting theme colour when None.
    latin_font, cjk_font : str
        Typeface names. Defaults to Arial / Droid Sans Fallback.
    align : PP_ALIGN
        Paragraph alignment.
    paragraph_index : int, optional
        When 0 and the first paragraph is empty, write into it instead of
        creating a new paragraph (avoids leading blank line).

    Returns
    -------
    pptx.text.text._Paragraph
        The paragraph object that was populated.
    """
    if (
        paragraph_index is not None
        and paragraph_index == 0
        and text_frame.paragraphs[0].text == ""
    ):
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = align

    for seg, _is_cjk_seg in _segments(text):
        run = p.add_run()
        run.text = seg
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        run.font.italic = bool(italic)
        if color is not None:
            run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        # Remove any pre-existing latin/ea elements before injecting ours.
        for tag in ("a:latin", "a:ea"):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        latin_el = etree.SubElement(rPr, qn("a:latin"))
        latin_el.set("typeface", latin_font)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", cjk_font)
    return p


def fit_image_within(
    slide,
    image_path,
    x,
    y,
    max_w,
    max_h,
    border: bool = False,
    fallback_label: Optional[str] = None,
):
    """Insert a picture preserving aspect ratio inside (max_w, max_h).

    Equal-ratio scaling: ``ratio = min(max_w / image_w, max_h / image_h)``.
    Never forces both width and height — that is the failure mode v2/v3
    explicitly forbid (causes squashed IGV screenshots).

    When the source file is missing, draws a light-grey placeholder rectangle
    with a red border + the filename, so wireframe screenshots still surface
    the missing asset without crashing the build.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    image_path : str | Path
        Source image. PNG / JPG / SVG-via-PIL etc.
    x, y, max_w, max_h : EMU (use Inches() to construct)
        Bounding box.
    border : bool
        When True, draw a thin grey rectangle outlining the placed picture.
    fallback_label : str, optional
        Override the auto-generated "[圖片缺失] {name}" placeholder text.

    Returns
    -------
    pptx.shapes.picture.Picture | None
        The inserted picture, or None if the file was missing (placeholder
        rendered instead).
    """
    image_path = Path(image_path)
    if not image_path.exists():
        _draw_missing_placeholder(slide, x, y, max_w, max_h, image_path.name, fallback_label)
        return None

    img = Image.open(image_path)
    iw, ih = img.size
    img.close()

    # python-pptx EMU per pixel (assuming 96 DPI source). 9525 EMU = 1/9525 in.
    iw_emu = iw * 9525
    ih_emu = ih * 9525
    ratio = min(max_w / iw_emu, max_h / ih_emu)
    final_w = int(iw_emu * ratio)
    final_h = int(ih_emu * ratio)
    cx = x + (max_w - final_w) // 2
    cy = y + (max_h - final_h) // 2
    pic = slide.shapes.add_picture(str(image_path), cx, cy, width=final_w, height=final_h)
    if border:
        _add_border(slide, cx, cy, final_w, final_h)
    return pic


def set_speaker_note(slide, text: str, min_chars: int = 350, max_chars: Optional[int] = None) -> None:
    """Set the speaker note for a slide with mandatory length bounds.

    v3 enhancement over v2: also enforces an *upper* bound. Long speaker notes
    (> 360 CJK chars ≈ 90 sec) overshoot the 1 min/slide PLOS rule — the
    myPPT playbook §4 requires Tier 3 oral-optional content to be split out.

    Parameters
    ----------
    slide : pptx.slide.Slide
    text : str
        Speaker note content (≥ min_chars, ≤ max_chars).
    min_chars : int
        Lower bound. Default 350 matches v2/v3 storyboard requirement.
    max_chars : int, optional
        Upper bound. When set, raises if exceeded — caller should split into
        Tier 2 (must-say) + Tier 3 (oral-optional, marked [ORAL-OPTIONAL]).

    Raises
    ------
    ValueError
        Empty note, length < min_chars, or length > max_chars.
    """
    if not text or not text.strip():
        raise ValueError(f"Empty speaker note for slide {slide.slide_id}")
    if len(text) < min_chars:
        raise ValueError(
            f"Speaker note too short ({len(text)} chars) for slide {slide.slide_id}; "
            f"need >= {min_chars}"
        )
    if max_chars is not None and len(text) > max_chars:
        raise ValueError(
            f"Speaker note too long ({len(text)} chars) for slide {slide.slide_id}; "
            f"max {max_chars}. Split tier-3 oral-optional content into a separate "
            f"[ORAL-OPTIONAL] section or move to backup slide."
        )
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    add_text_with_fallback(notes, text, font_size=14, paragraph_index=0)


# ---------------------------------------------------------------------------
# Internal helpers (placeholder + border)
# ---------------------------------------------------------------------------
def _draw_missing_placeholder(slide, x, y, max_w, max_h, name: str, label_override: Optional[str]) -> None:
    """Draw a light-grey rectangle with a red border and filename caption."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, max_w, max_h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = PALETTE["light"]
    shp.line.color.rgb = PALETTE["accent"]
    shp.line.width = Pt(0.75)

    # Caption box centered vertically.
    caption_h = Inches(0.4)
    caption_y = y + max_h // 2 - caption_h // 2
    cap = slide.shapes.add_textbox(x, caption_y, max_w, caption_h)
    cap.text_frame.margin_left = Emu(60000)
    cap.text_frame.margin_right = Emu(60000)
    cap.text_frame.word_wrap = True
    cap.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cap.text_frame.paragraphs[0].text = ""
    label = label_override if label_override else f"[圖片缺失] {name}"
    add_text_with_fallback(
        cap.text_frame,
        label,
        font_size=14,
        color=PALETTE["accent"],
        align=PP_ALIGN.CENTER,
        paragraph_index=0,
    )


def _add_border(slide, x, y, w, h) -> None:
    """Add a thin rectangle border around a placed image."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.background()
    shp.line.color.rgb = PALETTE["border"]
    shp.line.width = Pt(0.75)
