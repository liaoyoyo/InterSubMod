"""style_library_loader.py — YAML-driven object & layout deserializer.

Loads object / layout / palette definitions from the myPPT style_library
(see InterSubMod/.claude/skills/myPPT/style_library/) into Python dataclasses.

Design constraints (myPPT plan §13, §17):

- Object YAML files describe a reusable visual element (caveat strip, KPI
  card, citation footnote, etc.) with required_fields + content_template.
- Layout YAML files describe a slide skeleton (zones with x/y/w/h/object).
- The renderer never inlines colour or alignment — those flow from
  ``style_library/colors/palette.yaml`` and ``alignment_rules`` blocks.
- Loader discovers the library root from one of:
    1. ``$MYPPT_STYLE_LIBRARY`` environment variable (testing override),
    2. relative path: ``<repo>/.claude/skills/myPPT/style_library/``,
       searched upward from CWD.

Public API:

    >>> from ppt_toolkit.style_library_loader import load_object, load_layout
    >>> spec = load_object("caveat_red_strip")
    >>> spec.render(slide, x=Inches(1.0), y=Inches(5.5),
    ...             content={"caveat_text": "Source mis-reported 38%, real 16.5%"})
    >>> layout = load_layout("before_after_split")
    >>> layout.populate(slide, {"title": "Baseline vs V5", ...})

Both render() and populate() are best-effort: missing optional fields fall
back to the YAML defaults; missing required fields raise ValueError early
with a helpful message naming the spec + missing key.
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

import yaml
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ppt_toolkit.text_helpers import (
    DEFAULT_CJK_FONT,
    DEFAULT_LATIN_FONT,
    PALETTE,
    add_text_with_fallback,
    fit_image_within,
)

# ---------------------------------------------------------------------------
# Library root resolution
# ---------------------------------------------------------------------------
_ENV_VAR = "MYPPT_STYLE_LIBRARY"
_DEFAULT_REL = ".claude/skills/myPPT/style_library"


def _find_style_library_root() -> Path:
    """Return absolute path to style_library/ root.

    Resolution order:
        1. $MYPPT_STYLE_LIBRARY (absolute path)
        2. Walk upward from CWD looking for .claude/skills/myPPT/style_library/
        3. Walk upward from this file's location (works when toolkit is
           installed at <repo>/tools/ppt_toolkit/ as designed in plan §26).
    """
    env = os.environ.get(_ENV_VAR)
    if env:
        p = Path(env).expanduser().resolve()
        if p.is_dir():
            return p
        raise FileNotFoundError(
            f"${_ENV_VAR}={env!r} does not point to an existing directory."
        )

    candidates: List[Path] = []
    candidates.append(Path.cwd())
    candidates.append(Path(__file__).resolve().parent.parent.parent)  # repo root from tools/ppt_toolkit/

    for start in candidates:
        cur = start
        for _ in range(8):  # walk up at most 8 levels
            target = cur / _DEFAULT_REL
            if target.is_dir():
                return target.resolve()
            if cur.parent == cur:
                break
            cur = cur.parent

    raise FileNotFoundError(
        f"Could not locate style_library. Set ${_ENV_VAR} or run from a directory "
        f"under InterSubMod/."
    )


# ---------------------------------------------------------------------------
# Helpers: YAML-side colour / measurement parsing
# ---------------------------------------------------------------------------
_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{6})$")


def _parse_color(value: Any) -> Optional[RGBColor]:
    """Convert YAML colour value to RGBColor.

    Accepts: ``"#D32F2F"``, ``"D32F2F"``, palette token name (e.g. ``"red"``),
    or an existing RGBColor passthrough.
    """
    if value is None:
        return None
    if isinstance(value, RGBColor):
        return value
    if isinstance(value, str):
        m = _HEX_RE.match(value.strip())
        if m:
            h = m.group(1)
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        # Fallback: palette token lookup
        token = value.strip().lower()
        if token in PALETTE:
            return PALETTE[token]
    raise ValueError(f"Unrecognised colour value: {value!r}")


def _parse_inches(value: Any) -> int:
    """Convert YAML measurement (e.g. ``"5.8in"``, ``5.8``, ``"0.5"``) to EMU."""
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return Inches(float(value))
    if isinstance(value, str):
        s = value.strip().lower()
        if s.endswith("in"):
            s = s[:-2].strip()
        if s.endswith('"'):
            s = s[:-1].strip()
        try:
            return Inches(float(s))
        except ValueError as e:
            raise ValueError(f"Cannot parse inches from {value!r}") from e
    raise ValueError(f"Cannot parse measurement: {value!r}")


def _format_template(template: str, content: Dict[str, Any]) -> str:
    """Best-effort ``str.format_map`` with a default substitution for missing keys."""

    class _Defaulting(dict):
        def __missing__(self, key):
            return "{" + key + "}"

    return template.format_map(_Defaulting(content))


# ---------------------------------------------------------------------------
# ObjectSpec
# ---------------------------------------------------------------------------
@dataclass
class ObjectSpec:
    """Parsed YAML object specification.

    Renders a single visual element (e.g. caveat_red_strip) at a given
    position. The caller supplies content dictionary matching required_fields.
    """

    name: str
    category: str = "generic"
    purpose: str = ""
    visual: Dict[str, Any] = field(default_factory=dict)
    content_template: str = "{text}"
    required_fields: List[str] = field(default_factory=list)
    optional_fields: List[str] = field(default_factory=list)
    tier_recommendation: str = "A"  # myPPT playbook §20 Tier S/A/B/C/D
    raw: Dict[str, Any] = field(default_factory=dict)
    source_path: Optional[Path] = None

    @classmethod
    def from_yaml(cls, path: Path) -> "ObjectSpec":
        with open(path, "r", encoding="utf-8") as fh:
            data = yaml.safe_load(fh) or {}
        return cls(
            name=data.get("name", path.stem),
            category=data.get("category", "generic"),
            purpose=data.get("purpose", ""),
            visual=data.get("visual", {}) or {},
            content_template=data.get("content_template", "{text}"),
            required_fields=list(data.get("required_fields", []) or []),
            optional_fields=list(data.get("optional_fields", []) or []),
            tier_recommendation=str(data.get("tier_recommendation", "A")),
            raw=data,
            source_path=path,
        )

    def validate(self, content: Dict[str, Any]) -> None:
        """Ensure all required_fields are present in content (raises early)."""
        # required_fields entries may be ``"name"`` or ``"name (constraint...)"``
        keys = [re.split(r"\s*\(", f, 1)[0].strip() for f in self.required_fields]
        missing = [k for k in keys if k and k not in content]
        if missing:
            raise ValueError(
                f"ObjectSpec[{self.name}] missing required fields: {missing}. "
                f"Provided keys: {list(content.keys())}"
            )

    def render(
        self,
        slide,
        x,
        y,
        content: Dict[str, Any],
        w=None,
        h=None,
    ):
        """Render this object onto ``slide`` at (x, y).

        Parameters
        ----------
        slide : pptx.slide.Slide
        x, y : EMU
            Top-left position. Use Inches() to construct.
        content : dict
            Substitutions for content_template (must include required_fields).
        w, h : EMU, optional
            Bounding box. Defaults to the YAML ``visual.size`` block or
            (10in × 0.8in) if absent.
        """
        self.validate(content)
        text = _format_template(self.content_template, content).strip()

        size = self.visual.get("size", {}) or {}
        w = w if w is not None else _parse_inches(size.get("w", 10.0))
        h = h if h is not None else _parse_inches(size.get("h", 0.8))

        bg = _parse_color(self.visual.get("background"))
        text_color = _parse_color(self.visual.get("text_color")) or PALETTE["text"]
        border_color = _parse_color(self.visual.get("border_color"))
        border_width_pt = float(self.visual.get("border_width_pt", 0.75))
        font_size_pt = int(self.visual.get("font_size_pt", 14))
        bold = bool(self.visual.get("bold", False))
        align = self._parse_align(self.visual.get("align", "left"))

        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.shadow.inherit = False
        if bg is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = bg
        else:
            shp.fill.background()
        if border_color is not None:
            shp.line.color.rgb = border_color
            shp.line.width = Pt(border_width_pt)
        else:
            shp.line.fill.background()

        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(60000)
        tf.margin_right = Emu(60000)
        tf.margin_top = Emu(40000)
        tf.margin_bottom = Emu(40000)
        tf.paragraphs[0].text = ""
        for i, line in enumerate(text.splitlines()):
            add_text_with_fallback(
                tf,
                line,
                font_size=font_size_pt,
                bold=bold,
                color=text_color,
                align=align,
                paragraph_index=i,
                latin_font=self.visual.get("latin_font", DEFAULT_LATIN_FONT),
                cjk_font=self.visual.get("cjk_font", DEFAULT_CJK_FONT),
            )
        return shp

    @staticmethod
    def _parse_align(value: Any):
        if isinstance(value, str):
            v = value.strip().lower()
            return {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "centre": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(v, PP_ALIGN.LEFT)
        return PP_ALIGN.LEFT


# ---------------------------------------------------------------------------
# LayoutSpec
# ---------------------------------------------------------------------------
@dataclass
class LayoutSpec:
    """Parsed YAML layout specification.

    A layout is a collection of zones; each zone has a position (x/y/w/h) and
    optionally references an ObjectSpec by name. populate() walks zones in
    declaration order, resolving content_dict keys via dotted notation
    ("left.heading", "delta_strip.content", etc.).
    """

    name: str
    template_type: List[str] = field(default_factory=list)
    purpose: str = ""
    canvas_w_in: float = 13.333
    canvas_h_in: float = 7.5
    zones: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    focal_point_zone: Optional[str] = None
    raw: Dict[str, Any] = field(default_factory=dict)
    source_path: Optional[Path] = None

    @classmethod
    def from_yaml(cls, path: Path) -> "LayoutSpec":
        with open(path, "r", encoding="utf-8") as fh:
            data = yaml.safe_load(fh) or {}
        canvas = data.get("canvas", "13.333in x 7.5in")
        cw, ch = cls._parse_canvas(canvas)
        focal = data.get("focal_point_zone")
        if isinstance(focal, dict):
            focal = focal.get("position") or focal.get("zone")
        return cls(
            name=data.get("name", path.stem),
            template_type=list(data.get("template_type", []) or []),
            purpose=data.get("purpose", ""),
            canvas_w_in=cw,
            canvas_h_in=ch,
            zones=dict(data.get("zones", {}) or {}),
            focal_point_zone=focal,
            raw=data,
            source_path=path,
        )

    @staticmethod
    def _parse_canvas(value: str):
        if isinstance(value, str):
            parts = re.split(r"[x×]", value.lower())
            if len(parts) == 2:
                try:
                    return float(parts[0].replace("in", "").strip()), float(
                        parts[1].replace("in", "").strip()
                    )
                except ValueError:
                    pass
        return 13.333, 7.5

    def populate(
        self,
        slide,
        content_dict: Dict[str, Any],
        object_resolver=None,
    ) -> None:
        """Populate zones with content_dict.

        Content keys use dotted notation: ``"<zone>.<field>"``.
        Bare ``"title"`` is treated as ``"title.text"`` for ergonomics.
        Object resolution defaults to :func:`load_object`.
        """
        if object_resolver is None:
            object_resolver = load_object

        # Group by zone.
        per_zone: Dict[str, Dict[str, Any]] = {}
        for k, v in content_dict.items():
            if "." in k:
                zone, key = k.split(".", 1)
            else:
                zone, key = k, "text"
            per_zone.setdefault(zone, {})[key] = v

        for zone_name, zone_def in self.zones.items():
            x = _parse_inches(zone_def.get("x", 0.0))
            y = _parse_inches(zone_def.get("y", 0.0))
            w = _parse_inches(zone_def.get("w", self.canvas_w_in - 1.0))
            h = _parse_inches(zone_def.get("h", 1.0))
            obj_name = zone_def.get("object")
            zone_content = per_zone.get(zone_name, {})

            # Image zone shortcut
            image_path = zone_content.get("image") or zone_def.get("image")
            if image_path:
                fit_image_within(slide, image_path, x, y, w, h)
                continue

            if obj_name:
                spec = object_resolver(obj_name)
                spec.render(slide, x, y, zone_content, w=w, h=h)
                continue

            # Plain text zone fallback
            text = zone_content.get("text") or zone_def.get("text")
            if text:
                _render_plain_text(slide, x, y, w, h, str(text))


def _render_plain_text(slide, x, y, w, h, text: str) -> None:
    """Render a textbox without any object styling (last-resort zone path)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for i, line in enumerate(text.splitlines() or [""]):
        add_text_with_fallback(tf, line, font_size=14, paragraph_index=i)


# ---------------------------------------------------------------------------
# Public loaders (cached)
# ---------------------------------------------------------------------------
_OBJECT_CACHE: Dict[str, ObjectSpec] = {}
_LAYOUT_CACHE: Dict[str, LayoutSpec] = {}


def load_object(name: str, root: Optional[Path] = None) -> ObjectSpec:
    """Load and cache an ObjectSpec by name (without ``.yaml`` extension)."""
    cache_key = name
    if cache_key in _OBJECT_CACHE:
        return _OBJECT_CACHE[cache_key]
    base = root or _find_style_library_root()
    path = base / "objects" / f"{name}.yaml"
    if not path.exists():
        raise FileNotFoundError(f"ObjectSpec not found: {path}")
    spec = ObjectSpec.from_yaml(path)
    _OBJECT_CACHE[cache_key] = spec
    return spec


def load_layout(name: str, root: Optional[Path] = None) -> LayoutSpec:
    """Load and cache a LayoutSpec by name (without ``.yaml`` extension)."""
    if name in _LAYOUT_CACHE:
        return _LAYOUT_CACHE[name]
    base = root or _find_style_library_root()
    path = base / "layouts" / f"{name}.yaml"
    if not path.exists():
        raise FileNotFoundError(f"LayoutSpec not found: {path}")
    spec = LayoutSpec.from_yaml(path)
    _LAYOUT_CACHE[name] = spec
    return spec


def load_palette(root: Optional[Path] = None) -> Dict[str, RGBColor]:
    """Load palette.yaml and return {token: RGBColor} mapping.

    Returns the in-memory PALETTE dict if the YAML file is missing — this
    keeps ppt_toolkit usable when style_library has not yet been authored
    (Stream B not yet complete).
    """
    base = root or _find_style_library_root()
    path = base / "colors" / "palette.yaml"
    if not path.exists():
        return dict(PALETTE)
    with open(path, "r", encoding="utf-8") as fh:
        data = yaml.safe_load(fh) or {}
    out: Dict[str, RGBColor] = {}
    for tok, body in (data.get("primary_palette") or {}).items():
        hex_value = body.get("hex") if isinstance(body, dict) else body
        col = _parse_color(hex_value)
        if col is not None:
            out[tok] = col
    return out or dict(PALETTE)


def clear_cache() -> None:
    """Clear loader caches (test-only helper)."""
    _OBJECT_CACHE.clear()
    _LAYOUT_CACHE.clear()
