"""Weekly Report v2 — Assertion-Evidence build script.

Design principles:
  * 16 slides (vs v1's 41)
  * Every title is a complete claim sentence
  * Okabe-Ito colorblind-safe palette (from theme_okabe_ito)
  * 15/70/15 vertical layout (headline/evidence/footer)
  * Bilingual: Chinese primary (28pt) + English subtitle (17pt, 60% grey)

Usage:
    python3 build_pptx_v2.py

Output: outputs/20260417_週報v2_AE.pptx + outputs/_snapshot.json
"""
import os
import sys
import datetime as _dt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from theme_okabe_ito import OKABE_ITO, ROLES, FONTS, CANVAS, LAYOUT, FINDING_COLORS

# ---------- Setup ----------
HERE = os.path.dirname(os.path.abspath(__file__))
FIGS_DIR = os.path.join(HERE, "figures")
OUT_DIR = os.path.join(HERE, "outputs")
os.makedirs(OUT_DIR, exist_ok=True)

TOTAL_SLIDES = 16
TODAY = "2026-04-17"
TITLE_SHORT = "LOH Subclone v2 • AE"


def _rgb(hex_code):
    return RGBColor.from_string(hex_code.lstrip("#"))


def _set_run_font(run, size_pt, *, color_hex=None, bold=False, name="Helvetica"):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color_hex:
        run.font.color.rgb = _rgb(color_hex)


# ---------- Helper functions (8) ----------
def add_claim_header(slide, zh_claim, en_subtitle, finding_color=None):
    """Add a complete-sentence claim headline + English subtitle + accent rule."""
    margin = LAYOUT["margin_x"]
    width = CANVAS["width_in"] - 2 * margin

    # ZH claim
    zh_box = slide.shapes.add_textbox(
        Inches(margin), Inches(LAYOUT["headline_top"]),
        Inches(width), Inches(0.58))
    tf = zh_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = zh_claim
    _set_run_font(r, FONTS["title_zh"], color_hex=ROLES["text"], bold=True, name="PingFang TC")

    # EN subtitle
    en_box = slide.shapes.add_textbox(
        Inches(margin), Inches(LAYOUT["headline_top"] + LAYOUT["en_offset"] + 0.10),
        Inches(width), Inches(0.45))
    tf = en_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = en_subtitle
    _set_run_font(r, FONTS["title_en"], color_hex=ROLES["text_light"], name="Helvetica")

    # Accent rule (2pt horizontal line under header)
    rule_color = finding_color or ROLES["accent"]
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(margin), Inches(LAYOUT["headline_top"] + LAYOUT["headline_height"] - 0.08),
        Inches(0.80), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(rule_color)
    rule.line.fill.background()


def add_evidence_image(slide, image_path, caption=None):
    """Center-fit an image within the evidence region, with optional caption."""
    margin = LAYOUT["margin_x"]
    region_w = CANVAS["width_in"] - 2 * margin
    region_h = LAYOUT["evidence_height"] - (0.35 if caption else 0)
    # Add image at full region width; let pptx keep aspect.
    pic = slide.shapes.add_picture(
        image_path,
        Inches(margin), Inches(LAYOUT["evidence_top"]),
        width=Inches(region_w))
    # If picture taller than region, rescale by height.
    if pic.height > Inches(region_h):
        scale = Inches(region_h) / pic.height
        pic.height = Inches(region_h)
        pic.width = int(pic.width * scale)
    # Horizontal center
    pic.left = Inches(margin + (region_w - pic.width / 914400) / 2)

    if caption:
        cap = slide.shapes.add_textbox(
            Inches(margin), Inches(LAYOUT["evidence_top"] + region_h + 0.05),
            Inches(region_w), Inches(0.30))
        tf = cap.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        _set_run_font(r, FONTS["caption"], color_hex=ROLES["text_light"])


def add_bullet_list(slide, items, left_in, top_in, width_in, height_in,
                     font_size=None, line_spacing=1.25):
    """Add a bulleted list with Okabe-Ito accent bullets."""
    font_size = font_size or FONTS["body"]
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                    Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r_bullet = p.add_run()
        r_bullet.text = "●  "
        _set_run_font(r_bullet, font_size, color_hex=ROLES["accent"], bold=True)
        r = p.add_run()
        r.text = item
        _set_run_font(r, font_size, color_hex=ROLES["text"], name="PingFang TC")


def add_two_column(slide, left_title, left_items, right_title, right_items, top_in=None):
    """Two-column layout for comparisons."""
    top = top_in if top_in is not None else LAYOUT["evidence_top"] + 0.10
    margin = LAYOUT["margin_x"]
    col_w = (CANVAS["width_in"] - 2 * margin - 0.50) / 2

    # Left column
    _add_column_header(slide, left_title, margin, top, col_w, ROLES["primary"])
    add_bullet_list(slide, left_items, margin, top + 0.55, col_w, 4.2, font_size=15)

    # Right column
    right_x = margin + col_w + 0.50
    _add_column_header(slide, right_title, right_x, top, col_w, ROLES["secondary"])
    add_bullet_list(slide, right_items, right_x, top + 0.55, col_w, 4.2, font_size=15)


def _add_column_header(slide, title, x, y, w, color):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.45))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, FONTS["label"] + 2, color_hex=color, bold=True, name="PingFang TC")
    # underline bar
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y + 0.42),
                                   Inches(w), Inches(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(color)
    rule.line.fill.background()


def add_three_column(slide, items, top_in=None):
    """Three-column layout for Q1/Q2/Q3 on slide 2."""
    top = top_in if top_in is not None else LAYOUT["evidence_top"] + 0.20
    margin = LAYOUT["margin_x"]
    col_w = (CANVAS["width_in"] - 2 * margin - 0.80) / 3
    colors = [ROLES["primary"], ROLES["success"], ROLES["accent"]]
    for i, (label, question, sub) in enumerate(items):
        x = margin + i * (col_w + 0.40)
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(x), Inches(top),
                                        Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(colors[i])
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_left = btf.margin_right = 0
        btf.margin_top = btf.margin_bottom = 0
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = label
        _set_run_font(br, 18, color_hex="#FFFFFF", bold=True)
        # Question text
        qbox = slide.shapes.add_textbox(Inches(x), Inches(top + 0.75),
                                         Inches(col_w), Inches(1.2))
        qtf = qbox.text_frame
        qtf.word_wrap = True
        qtf.margin_left = qtf.margin_right = 0
        qp = qtf.paragraphs[0]
        qr = qp.add_run()
        qr.text = question
        _set_run_font(qr, 17, color_hex=ROLES["text"], bold=True, name="PingFang TC")
        # Sub description
        sbox = slide.shapes.add_textbox(Inches(x), Inches(top + 2.10),
                                         Inches(col_w), Inches(2.4))
        stf = sbox.text_frame
        stf.word_wrap = True
        stf.margin_left = stf.margin_right = 0
        sp = stf.paragraphs[0]
        sp.line_spacing = 1.35
        sr = sp.add_run()
        sr.text = sub
        _set_run_font(sr, 13, color_hex=ROLES["text_light"], name="PingFang TC")


def add_stat_callout(slide, left_in, top_in, width_in, stat_value, stat_label, color):
    """Large statistic with small descriptive label below."""
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                    Inches(width_in), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = stat_value
    _set_run_font(r, FONTS["stat"] + 14, color_hex=color, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = stat_label
    _set_run_font(r2, FONTS["label"] - 2, color_hex=ROLES["text_light"], name="PingFang TC")


def add_footer(slide, slide_n, total=TOTAL_SLIDES, date=TODAY, title=TITLE_SHORT):
    """Footer with slide number, date, and short title."""
    margin = LAYOUT["margin_x"]
    y = LAYOUT["footer_top"]
    w = CANVAS["width_in"] - 2 * margin
    # Rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(margin), Inches(y),
                                   Inches(w), Inches(0.015))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(ROLES["rule"])
    rule.line.fill.background()
    # Left: title
    lbox = slide.shapes.add_textbox(Inches(margin), Inches(y + 0.10),
                                     Inches(w / 2), Inches(0.30))
    ltf = lbox.text_frame
    ltf.margin_left = ltf.margin_right = 0
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = title
    _set_run_font(lr, FONTS["footer"], color_hex=ROLES["text_muted"], name="PingFang TC")
    # Right: slide N / total + date
    rbox = slide.shapes.add_textbox(Inches(margin + w / 2), Inches(y + 0.10),
                                     Inches(w / 2), Inches(0.30))
    rtf = rbox.text_frame
    rtf.margin_left = rtf.margin_right = 0
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run()
    rr.text = f"{slide_n:02d} / {total}   •   {date}"
    _set_run_font(rr, FONTS["footer"], color_hex=ROLES["text_muted"])


def add_blank_slide(prs):
    """Create a blank slide (layout 6 = blank)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ---------- Icon helpers (shape-based pictograms) ----------
def _add_rect(slide, x, y, w, h, color, line_bg=True):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = _rgb(color)
    if line_bg:
        r.line.fill.background()
    return r


def _add_oval(slide, x, y, w, h, color, line_color=None, line_width=0.5):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    o.fill.solid()
    o.fill.fore_color.rgb = _rgb(color)
    if line_color:
        o.line.color.rgb = _rgb(line_color)
        o.line.width = Pt(line_width)
    else:
        o.line.fill.background()
    return o


def draw_icon(slide, x, y, size, icon_type):
    """Draw a simple pictogram icon at (x, y) with a square bounding box `size`.

    Icon types:
      af       : Two stacked bars (ref vs alt reads)
      loh      : Two chromosome bars (maternal preserved, paternal lost with ✕)
      ngroups  : Three clustered color dots
      hp       : Two parallel haplotype lines with filled/hollow markers
      segment  : Segmented genomic bar (4 colored portions)
      dual     : Two overlapping circles (dual evidence)
    """
    s = size
    if icon_type == "af":
        _add_rect(slide, x, y + s * 0.40, s, s * 0.18, ROLES["rule"])
        _add_rect(slide, x, y + s * 0.40, s * 0.45, s * 0.18, OKABE_ITO["vermillion"])
        # Label ticks
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + s * 0.62),
                                       Inches(s), Inches(s * 0.22))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "alt"
        _set_run_font(r, 9, color_hex=OKABE_ITO["vermillion"], bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = ""  # spacer
        _set_run_font(r2, 1, color_hex=ROLES["text"])

    elif icon_type == "loh":
        _add_rect(slide, x, y + s * 0.18, s, s * 0.22, OKABE_ITO["blue"])
        lost = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y + s * 0.56),
                                       Inches(s), Inches(s * 0.22))
        lost.fill.solid()
        lost.fill.fore_color.rgb = _rgb(ROLES["rule"])
        lost.line.fill.background()
        # Red X across lost chromosome
        xtb = slide.shapes.add_textbox(Inches(x), Inches(y + s * 0.48),
                                        Inches(s), Inches(s * 0.40))
        xtf = xtb.text_frame
        xtf.margin_left = xtf.margin_right = 0
        xtf.margin_top = xtf.margin_bottom = 0
        xp = xtf.paragraphs[0]
        xp.alignment = PP_ALIGN.CENTER
        xr = xp.add_run()
        xr.text = "✕"
        _set_run_font(xr, int(s * 28), color_hex=OKABE_ITO["vermillion"], bold=True)

    elif icon_type == "ngroups":
        # 3 clusters of 3 dots each in different colors
        clusters = [(0.18, 0.22, OKABE_ITO["blue"]),
                    (0.68, 0.20, OKABE_ITO["vermillion"]),
                    (0.40, 0.68, OKABE_ITO["blue_green"])]
        dot_size = s * 0.10
        for (cx, cy, col) in clusters:
            for (dx, dy) in [(0, 0), (0.10, 0.06), (-0.04, 0.10)]:
                _add_oval(slide, x + s * (cx + dx) - dot_size / 2,
                          y + s * (cy + dy) - dot_size / 2,
                          dot_size, dot_size, col)

    elif icon_type == "hp":
        # Two horizontal haplotype lines
        _add_rect(slide, x, y + s * 0.30, s, s * 0.04, OKABE_ITO["blue"])
        _add_rect(slide, x, y + s * 0.66, s, s * 0.04, OKABE_ITO["vermillion"])
        # Markers (filled for HP1, hollow for HP2)
        for mx in [0.20, 0.50, 0.80]:
            _add_oval(slide, x + s * mx - s * 0.05, y + s * 0.32 - s * 0.05,
                      s * 0.10, s * 0.10, OKABE_ITO["blue"])
            _add_oval(slide, x + s * mx - s * 0.05, y + s * 0.68 - s * 0.05,
                      s * 0.10, s * 0.10, "#FFFFFF",
                      line_color=OKABE_ITO["vermillion"], line_width=1.4)
        # Labels HP1/HP2
        for (txt, yy, col) in [("HP1", 0.10, OKABE_ITO["blue"]),
                                ("HP2", 0.80, OKABE_ITO["vermillion"])]:
            tb = slide.shapes.add_textbox(Inches(x), Inches(y + s * yy),
                                           Inches(s * 0.40), Inches(s * 0.20))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = txt
            _set_run_font(r, 9, color_hex=col, bold=True)

    elif icon_type == "segment":
        parts = [(0.0, 0.30, OKABE_ITO["blue"]),
                 (0.30, 0.55, OKABE_ITO["orange"]),
                 (0.55, 0.78, ROLES["rule"]),
                 (0.78, 1.0, OKABE_ITO["blue_green"])]
        for (a, b, col) in parts:
            seg = _add_rect(slide, x + s * a, y + s * 0.42,
                            s * (b - a), s * 0.20, col, line_bg=False)
            seg.line.color.rgb = _rgb("#FFFFFF")
            seg.line.width = Pt(1.5)
        # Tick marks on top
        for tx in [0.30, 0.55, 0.78]:
            _add_rect(slide, x + s * tx - s * 0.005, y + s * 0.36,
                      s * 0.010, s * 0.08, ROLES["text"])

    elif icon_type == "dual":
        _add_oval(slide, x + s * 0.05, y + s * 0.20, s * 0.55, s * 0.55,
                  OKABE_ITO["blue"])
        # Overlapping circle
        ol = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + s * 0.40), Inches(y + s * 0.20),
                                     Inches(s * 0.55), Inches(s * 0.55))
        ol.fill.solid()
        ol.fill.fore_color.rgb = _rgb(OKABE_ITO["vermillion"])
        ol.fill.transparency = 0.3  # may not render; fallback OK
        ol.line.fill.background()


def add_concept_badge(slide, x, y, w, h, icon_type, term_zh, term_en, definition):
    """Composite: rounded card with icon + bilingual term + one-line definition."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(ROLES["bg_panel"])
    card.line.color.rgb = _rgb(ROLES["rule"])
    card.line.width = Pt(0.5)

    # Icon in upper-left (size relative to card height)
    icon_size = min(h * 0.45, 0.70)
    icon_x = x + 0.15
    icon_y = y + 0.15
    draw_icon(slide, icon_x, icon_y, icon_size, icon_type)

    # Term (right of icon)
    term_x = icon_x + icon_size + 0.15
    tbox = slide.shapes.add_textbox(Inches(term_x), Inches(y + 0.10),
                                     Inches(w - (term_x - x) - 0.10), Inches(0.40))
    tf = tbox.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = term_zh
    _set_run_font(r, 15, color_hex=ROLES["primary"], bold=True, name="PingFang TC")

    enbox = slide.shapes.add_textbox(Inches(term_x), Inches(y + 0.45),
                                      Inches(w - (term_x - x) - 0.10), Inches(0.30))
    enf = enbox.text_frame
    enf.margin_left = enf.margin_right = 0
    ep = enf.paragraphs[0]
    er = ep.add_run()
    er.text = term_en
    _set_run_font(er, 10, color_hex=ROLES["text_light"], name="Helvetica")

    # Definition (full width, below icon/term row)
    def_top = y + max(icon_size + 0.25, 0.90)
    dbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(def_top),
                                     Inches(w - 0.30), Inches(h - (def_top - y) - 0.10))
    dtf = dbox.text_frame
    dtf.word_wrap = True
    dtf.margin_left = dtf.margin_right = 0
    dp = dtf.paragraphs[0]
    dp.line_spacing = 1.30
    dr = dp.add_run()
    dr.text = definition
    _set_run_font(dr, 11, color_hex=ROLES["text"], name="PingFang TC")


def add_mini_concept_tag(slide, x, y, w, icon_type, term, definition):
    """Compact inline concept tag (one-row) for inline references in findings slides."""
    # Small icon
    draw_icon(slide, x, y, 0.35, icon_type)
    # Term + def on right
    tbox = slide.shapes.add_textbox(Inches(x + 0.45), Inches(y),
                                     Inches(w - 0.45), Inches(0.40))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = term + "  "
    _set_run_font(r1, 10, color_hex=ROLES["primary"], bold=True, name="PingFang TC")
    r2 = p.add_run()
    r2.text = definition
    _set_run_font(r2, 10, color_hex=ROLES["text_light"], name="PingFang TC")


# ---------- 16 Slide builders ----------
def slide_01_cover(prs):
    """Cover: LOH 區域發現 subclonal 表觀遺傳結構."""
    s = add_blank_slide(prs)
    margin = LAYOUT["margin_x"]

    # Top tag
    tag = s.shapes.add_textbox(Inches(margin), Inches(0.45),
                                Inches(6.0), Inches(0.35))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "WEEKLY RESEARCH UPDATE  •  2026-04-17"
    _set_run_font(r, 11, color_hex=ROLES["accent"], bold=True)

    # Main title (ZH)
    tbox = s.shapes.add_textbox(Inches(margin), Inches(1.10),
                                 Inches(CANVAS["width_in"] - 2 * margin), Inches(1.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "LOH 區域發現 subclonal 表觀遺傳結構"
    _set_run_font(r, 40, color_hex=ROLES["text"], bold=True, name="PingFang TC")

    # EN title
    ebox = s.shapes.add_textbox(Inches(margin), Inches(2.55),
                                 Inches(CANVAS["width_in"] - 2 * margin), Inches(0.7))
    tf = ebox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "LOH regions reveal subclonal epigenetic signatures in long-read data"
    _set_run_font(r, 20, color_hex=ROLES["text_light"], name="Helvetica")

    # Key statistics row
    stats_y = 3.80
    stats = [
        ("6.0×", "Intermediate-AF enrichment"),
        ("+0.705", "ΔNGroups (p<10⁻³⁹)"),
        ("d = 0.72", "AlleleDelta effect size"),
        ("ρ = 0.27", "Segment spatial correlation"),
    ]
    col_w = (CANVAS["width_in"] - 2 * margin) / 4
    stat_colors = [OKABE_ITO["blue"], OKABE_ITO["vermillion"],
                   OKABE_ITO["blue_green"], OKABE_ITO["orange"]]
    for i, (val, lab) in enumerate(stats):
        add_stat_callout(s, margin + i * col_w, stats_y, col_w, val, lab, stat_colors[i])

    # Author / context
    abox = s.shapes.add_textbox(Inches(margin), Inches(6.10),
                                 Inches(CANVAS["width_in"] - 2 * margin), Inches(0.5))
    tf = abox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "InterSubMod Research  •  Long-read ONT  •  HCC1395, HCC1937, H1975, H2009, H838, H1437, H2228"
    _set_run_font(r, 11, color_hex=ROLES["text_muted"], name="Helvetica")

    add_footer(s, 1)
    return s


def slide_02_three_questions(prs):
    """Concept primer (4 key terms) + three framing questions."""
    s = add_blank_slide(prs)
    add_claim_header(s,
        "先認識 4 個關鍵術語，再談三個研究問題",
        "Meet 4 key terms, then three framing questions")

    margin = LAYOUT["margin_x"]

    # ---- Row 1: 4 concept badges ----
    row1_y = LAYOUT["evidence_top"] + 0.05
    badge_h = 1.75
    badge_gap = 0.18
    badge_w = (CANVAS["width_in"] - 2 * margin - 3 * badge_gap) / 4

    concepts = [
        ("af",      "AF  等位基因頻率",
         "Allele Frequency",
         "變異 read / 總 read。0.5 代表雜合，中間 AF 常源自 subclone。"),
        ("loh",     "LOH  雜合性喪失",
         "Loss of Heterozygosity",
         "一條染色體區段失去；剩下的 allele 主導 variant 呈現。"),
        ("ngroups", "NGroups  甲基化群集數",
         "Methylation cluster count",
         "Reads 依 CpG 模式分群後的群集數；高 NGroups = 異質性。"),
        ("hp",      "HP  單倍體相位標籤",
         "Haplotype phase tag",
         "將 read 歸入 HP1 / HP2；讓我們分別觀察兩條染色體的甲基化。"),
    ]
    for i, (icon, zh, en, df) in enumerate(concepts):
        x = margin + i * (badge_w + badge_gap)
        add_concept_badge(s, x, row1_y, badge_w, badge_h, icon, zh, en, df)

    # ---- Divider ----
    div_y = row1_y + badge_h + 0.18
    _add_rect(s, margin, div_y, CANVAS["width_in"] - 2 * margin, 0.012,
              ROLES["rule_light"])

    # ---- Row 2: 3 questions (compact) ----
    row2_y = div_y + 0.15
    q_gap = 0.40
    q_w = (CANVAS["width_in"] - 2 * margin - 2 * q_gap) / 3
    questions = [
        ("Q1", "LOH 區域是否富含 subclonal AF？",
         "中間 AF（0.4-0.6）富集度是檢驗 subclonal 假說的第一個訊號。",
         ROLES["primary"]),
        ("Q2", "甲基化 NGroups 能否區辨 LOH？",
         "ΔNGroups 跨 7 樣本一致 > 0 代表生物訊號穩健。",
         ROLES["success"]),
        ("Q3", "AF × methylation 空間是否一致？",
         "Segment-level ρ 與 4-group 分層驗證 non-confound。",
         ROLES["accent"]),
    ]
    for i, (tag, q, sub, col) in enumerate(questions):
        x = margin + i * (q_w + q_gap)
        # Number badge
        badge = slide_badge_circle(s, x, row2_y, 0.50, tag, col)
        # Question
        qbox = s.shapes.add_textbox(Inches(x + 0.60), Inches(row2_y - 0.02),
                                     Inches(q_w - 0.60), Inches(0.55))
        qtf = qbox.text_frame
        qtf.word_wrap = True
        qtf.margin_left = qtf.margin_right = 0
        qp = qtf.paragraphs[0]
        qp.line_spacing = 1.15
        qr = qp.add_run()
        qr.text = q
        _set_run_font(qr, 14, color_hex=ROLES["text"], bold=True, name="PingFang TC")
        # Sub
        sbox = s.shapes.add_textbox(Inches(x), Inches(row2_y + 0.65),
                                     Inches(q_w), Inches(1.0))
        stf = sbox.text_frame
        stf.word_wrap = True
        stf.margin_left = stf.margin_right = 0
        sp = stf.paragraphs[0]
        sp.line_spacing = 1.30
        sr = sp.add_run()
        sr.text = sub
        _set_run_font(sr, 11, color_hex=ROLES["text_light"], name="PingFang TC")

    add_footer(s, 2)
    return s


def slide_badge_circle(slide, x, y, size, text, color):
    """Small filled circle with white text (used for Q1/Q2/Q3 tags)."""
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x), Inches(y), Inches(size), Inches(size))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(color)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, 14, color_hex="#FFFFFF", bold=True)
    return badge


def slide_03_intermediate_af(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "中間 AF 區間富含 LOH，是 subclonal 主要訊號區",
        "Intermediate AF (0.4–0.6) is enriched 6× in LOH regions",
        FINDING_COLORS["F1_af_enrichment"])
    add_evidence_image(s, os.path.join(FIGS_DIR, "01_intermediate_af_enrichment.png"),
                       "Left: AF density by LOH status (n=7 samples).  "
                       "Right: Proportion of variants with AF ∈ [0.4, 0.6].")
    add_footer(s, 3)
    return s


def slide_04_delta_ngroups(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "LOH 區域 NGroups 每樣本皆顯著較高",
        "ΔNGroups = +0.705 across all 7 samples (p < 10⁻³⁹)",
        FINDING_COLORS["F2_delta_ngroups"])
    add_evidence_image(s, os.path.join(FIGS_DIR, "02_delta_ngroups_7samples.png"),
                       "Per-sample ΔNGroups (LOH − Non-LOH) with 95% CI. "
                       "Pooled effect: +0.705, p < 10⁻³⁹ (Mann–Whitney).")
    add_footer(s, 4)
    return s


def slide_05_allele_delta(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "AlleleDelta 在 LOH 呈大效應分離",
        "AlleleDelta effect size Cohen's d = 0.724 (large)",
        FINDING_COLORS["F3_allele_delta"])
    add_evidence_image(s, os.path.join(FIGS_DIR, "03_allele_delta_effect.png"),
                       "Violin distribution with median; Cohen's d interprets as 'large effect' per Cohen (1988).")
    add_footer(s, 5)
    return s


def slide_06_segment_rho(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "Segment-level 空間相關性中等但顯著",
        "Segment-level Spearman ρ = 0.270 between AF SD and NGroups",
        FINDING_COLORS["F4_segment_rho"])
    add_evidence_image(s, os.path.join(FIGS_DIR, "04_segment_spatial_rho.png"),
                       "450 LOH segments pooled across 7 samples. "
                       "Linear fit shown; significance by permutation test.")
    add_footer(s, 6)
    return s


def slide_07_4group(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "四分層策略分離 high-TP 亞群",
        "Four-group (AF × NGroups) stratification isolates TP-enriched cluster",
        FINDING_COLORS["F5_4group_strat"])
    add_evidence_image(s, os.path.join(FIGS_DIR, "05_4group_stratification.png"),
                       "High-AF × High-NGroups 亞群 TP rate 89.1% (n=295)，遠高於隨機預期 50%。")
    add_footer(s, 7)
    return s


def slide_08_hpfine_marker(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "HPFineNGroups 作為 somatic heterogeneity 標記",
        "HPFineNGroups marks residual somatic diversity beyond HP phasing")
    margin = LAYOUT["margin_x"]
    top = LAYOUT["evidence_top"] + 0.10

    # ===== LEFT: Visual schematic (3-level nested clustering) =====
    panel_w = 7.0
    panel_h = 4.55
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(margin), Inches(top),
                                Inches(panel_w), Inches(panel_h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(ROLES["bg_panel"])
    panel.line.color.rgb = _rgb(ROLES["rule"])
    panel.line.width = Pt(0.5)

    # Panel title
    pt_box = s.shapes.add_textbox(Inches(margin + 0.25), Inches(top + 0.15),
                                   Inches(panel_w - 0.50), Inches(0.35))
    ptf = pt_box.text_frame
    ptf.margin_left = ptf.margin_right = 0
    pp = ptf.paragraphs[0]
    pr = pp.add_run()
    pr.text = "HPFineNGroups = 在 HP 層下再做 fine-grained 分群"
    _set_run_font(pr, 14, color_hex=ROLES["primary"], bold=True, name="PingFang TC")

    # ---- Level 1: Read pool (top center) ----
    pool_cx = margin + panel_w / 2
    pool_w = 1.8
    pool_h = 0.55
    pool_y = top + 0.65
    pool = _add_rect(s, pool_cx - pool_w / 2, pool_y, pool_w, pool_h, "#FFFFFF")
    pool.line.color.rgb = _rgb(ROLES["rule"])
    pool.line.width = Pt(0.8)
    # Dots inside pool (mixed colors)
    dot_colors = [OKABE_ITO["blue"], OKABE_ITO["blue"], OKABE_ITO["vermillion"],
                  OKABE_ITO["vermillion"], OKABE_ITO["blue"], OKABE_ITO["vermillion"]]
    for i, col in enumerate(dot_colors):
        dx = pool_cx - pool_w / 2 + 0.22 + i * 0.24
        _add_oval(s, dx, pool_y + 0.18, 0.18, 0.18, col)
    # Label left of pool
    plbl = s.shapes.add_textbox(Inches(pool_cx - pool_w / 2 - 1.35), Inches(pool_y + 0.10),
                                 Inches(1.30), Inches(0.35))
    plf = plbl.text_frame
    plf.margin_left = plf.margin_right = 0
    lp = plf.paragraphs[0]
    lp.alignment = PP_ALIGN.RIGHT
    lr = lp.add_run()
    lr.text = "① Read pool"
    _set_run_font(lr, 12, color_hex=ROLES["text"], bold=True, name="PingFang TC")

    # ---- Arrow: HP phasing ----
    arr1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                               Inches(pool_cx - 0.15), Inches(pool_y + pool_h + 0.05),
                               Inches(0.30), Inches(0.30))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = _rgb(ROLES["accent"])
    arr1.line.fill.background()
    arr1_lbl = s.shapes.add_textbox(Inches(pool_cx + 0.25), Inches(pool_y + pool_h + 0.08),
                                     Inches(2.0), Inches(0.30))
    al1 = arr1_lbl.text_frame
    al1.margin_left = al1.margin_right = 0
    ap1 = al1.paragraphs[0]
    ar1 = ap1.add_run()
    ar1.text = "HP phasing"
    _set_run_font(ar1, 10, color_hex=ROLES["accent"], bold=True)

    # ---- Level 2: Two HP groups ----
    hp_y = pool_y + pool_h + 0.55
    hp_h = 0.55
    hp_w = 2.0
    hp1_cx = pool_cx - 1.6
    hp2_cx = pool_cx + 1.6
    # HP1
    hp1 = _add_rect(s, hp1_cx - hp_w / 2, hp_y, hp_w, hp_h, OKABE_ITO["blue"])
    hp1.fill.transparency = 0.0
    for i in range(3):
        _add_oval(s, hp1_cx - hp_w / 2 + 0.35 + i * 0.45, hp_y + 0.18,
                  0.18, 0.18, "#FFFFFF")
    h1lbl = s.shapes.add_textbox(Inches(hp1_cx - hp_w / 2), Inches(hp_y - 0.35),
                                  Inches(hp_w), Inches(0.30))
    h1tf = h1lbl.text_frame
    h1tf.margin_left = h1tf.margin_right = 0
    h1p = h1tf.paragraphs[0]
    h1p.alignment = PP_ALIGN.CENTER
    h1r = h1p.add_run()
    h1r.text = "② HP1 reads"
    _set_run_font(h1r, 11, color_hex=OKABE_ITO["blue"], bold=True, name="PingFang TC")
    # HP2
    hp2 = _add_rect(s, hp2_cx - hp_w / 2, hp_y, hp_w, hp_h, OKABE_ITO["vermillion"])
    for i in range(3):
        _add_oval(s, hp2_cx - hp_w / 2 + 0.35 + i * 0.45, hp_y + 0.18,
                  0.18, 0.18, "#FFFFFF")
    h2lbl = s.shapes.add_textbox(Inches(hp2_cx - hp_w / 2), Inches(hp_y - 0.35),
                                  Inches(hp_w), Inches(0.30))
    h2tf = h2lbl.text_frame
    h2tf.margin_left = h2tf.margin_right = 0
    h2p = h2tf.paragraphs[0]
    h2p.alignment = PP_ALIGN.CENTER
    h2r = h2p.add_run()
    h2r.text = "② HP2 reads"
    _set_run_font(h2r, 11, color_hex=OKABE_ITO["vermillion"], bold=True, name="PingFang TC")

    # ---- Arrows: fine-grained clustering ----
    for cx in [hp1_cx, hp2_cx]:
        arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  Inches(cx - 0.12), Inches(hp_y + hp_h + 0.05),
                                  Inches(0.24), Inches(0.25))
        arr.fill.solid()
        arr.fill.fore_color.rgb = _rgb(ROLES["accent"])
        arr.line.fill.background()
    # Central label for step 3
    s3lbl = s.shapes.add_textbox(Inches(pool_cx - 1.8), Inches(hp_y + hp_h + 0.30),
                                  Inches(3.6), Inches(0.28))
    s3tf = s3lbl.text_frame
    s3tf.margin_left = s3tf.margin_right = 0
    s3p = s3tf.paragraphs[0]
    s3p.alignment = PP_ALIGN.CENTER
    s3r = s3p.add_run()
    s3r.text = "fine-grained methylation clustering"
    _set_run_font(s3r, 10, color_hex=ROLES["accent"], bold=True)

    # ---- Level 3: sub-clusters ----
    sub_y = hp_y + hp_h + 0.65
    # HP1 sub-clusters (2 small circles — NGroups=2 for HP1 in this example)
    for i, col_shade in enumerate([OKABE_ITO["blue"], OKABE_ITO["sky_blue"]]):
        cx = hp1_cx - 0.45 + i * 0.90
        _add_oval(s, cx - 0.20, sub_y, 0.40, 0.40, col_shade)
        # dots inside
        for (dx, dy) in [(-0.08, -0.04), (0.06, 0.02)]:
            _add_oval(s, cx + dx - 0.04, sub_y + 0.20 + dy - 0.04, 0.08, 0.08, "#FFFFFF")
    # HP2 sub-clusters (3 small circles — NGroups=3 for HP2)
    for i, col_shade in enumerate([OKABE_ITO["vermillion"], OKABE_ITO["orange"], OKABE_ITO["pink"]]):
        cx = hp2_cx - 0.65 + i * 0.65
        _add_oval(s, cx - 0.18, sub_y, 0.36, 0.36, col_shade)
        for (dx, dy) in [(-0.06, -0.02), (0.04, 0.04)]:
            _add_oval(s, cx + dx - 0.04, sub_y + 0.18 + dy - 0.04, 0.07, 0.07, "#FFFFFF")

    # Level-3 labels
    s_lb1 = s.shapes.add_textbox(Inches(hp1_cx - 1.0), Inches(sub_y + 0.45),
                                  Inches(2.0), Inches(0.28))
    lb1f = s_lb1.text_frame
    lb1f.margin_left = lb1f.margin_right = 0
    lb1p = lb1f.paragraphs[0]
    lb1p.alignment = PP_ALIGN.CENTER
    lb1r = lb1p.add_run()
    lb1r.text = "③ HP1 fine groups = 2"
    _set_run_font(lb1r, 10, color_hex=ROLES["text"], bold=True, name="PingFang TC")

    s_lb2 = s.shapes.add_textbox(Inches(hp2_cx - 1.0), Inches(sub_y + 0.45),
                                  Inches(2.0), Inches(0.28))
    lb2f = s_lb2.text_frame
    lb2f.margin_left = lb2f.margin_right = 0
    lb2p = lb2f.paragraphs[0]
    lb2p.alignment = PP_ALIGN.CENTER
    lb2r = lb2p.add_run()
    lb2r.text = "③ HP2 fine groups = 3"
    _set_run_font(lb2r, 10, color_hex=ROLES["text"], bold=True, name="PingFang TC")

    # Formula at bottom
    formula_y = sub_y + 0.85
    formula_box = _add_rect(s, margin + 0.35, formula_y, panel_w - 0.70, 0.40,
                             ROLES["text"])
    ftf = formula_box.text_frame
    ftf.margin_left = ftf.margin_right = 0
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "HPFineNGroups = Σ fine groups (across HPs)  =  2 + 3  =  5"
    _set_run_font(fr, 12, color_hex="#FFFFFF", bold=True, name="Helvetica")

    # ===== RIGHT: metric comparison (stat callouts vertically) =====
    right_x = margin + panel_w + 0.35
    right_w = CANVAS["width_in"] - margin - right_x
    metrics = [
        ("89.1%", "TP rate (N≥4, NR≥80)", ROLES["success"]),
        ("2.3×", "vs overall baseline", ROLES["accent"]),
        ("7 / 7", "samples consistent direction", ROLES["primary"]),
    ]
    for i, (val, lab, col) in enumerate(metrics):
        add_stat_callout(s, right_x, top + 0.20 + i * 1.45, right_w, val, lab, col)

    add_footer(s, 8)
    return s


def slide_09_dual_evidence(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "雙證據鏈：AF × methylation 共同支持 subclonal 假說",
        "Dual evidence chain: AF and methylation co-support the subclonal hypothesis")

    margin = LAYOUT["margin_x"]
    top = LAYOUT["evidence_top"] + 0.20

    # Three boxes: AF evidence | Methylation evidence | Convergence
    box_w = (CANVAS["width_in"] - 2 * margin - 1.00) / 3
    specs = [
        ("AF 證據", "af", ROLES["primary"], [
            "中間 AF 富集 6.0×",
            "Segment AF SD 與 NGroups ρ=0.27",
            "4-group 中 High-AF 亞群 TP+",
        ]),
        ("Methylation 證據", "ngroups", FINDING_COLORS["F2_delta_ngroups"], [
            "ΔNGroups=+0.705, p<10⁻³⁹",
            "AlleleDelta d=0.724 (大效應)",
            "HPFineNGroups somatic marker 89.1%",
        ]),
        ("匯流 Convergence", "dual", ROLES["success"], [
            "兩證據鏈獨立成立",
            "方向一致 (positive)",
            "跨 7 樣本皆穩健",
        ]),
    ]
    for i, (title, icon_type, col, items) in enumerate(specs):
        x = margin + i * (box_w + 0.50)
        # Header strip
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(top),
                                    Inches(box_w), Inches(0.45))
        strip.fill.solid()
        strip.fill.fore_color.rgb = _rgb(col)
        strip.line.fill.background()
        stf = strip.text_frame
        stf.margin_left = Inches(0.15)
        stf.margin_top = Inches(0.05)
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = title
        _set_run_font(sr, 15, color_hex="#FFFFFF", bold=True, name="PingFang TC")
        # Body
        body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(top + 0.45),
                                   Inches(box_w), Inches(3.5))
        body.fill.solid()
        body.fill.fore_color.rgb = _rgb(ROLES["bg_panel"])
        body.line.color.rgb = _rgb(ROLES["rule"])
        body.line.width = Pt(0.5)
        # Icon at top-center of body
        icon_size = 0.70
        draw_icon(s, x + box_w / 2 - icon_size / 2, top + 0.60, icon_size, icon_type)
        # Items below icon
        items_box = s.shapes.add_textbox(Inches(x + 0.20),
                                          Inches(top + 0.60 + icon_size + 0.15),
                                          Inches(box_w - 0.40), Inches(2.3))
        btf = items_box.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = 0
        for j, line in enumerate(items):
            pp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            pp.line_spacing = 1.40
            rr = pp.add_run()
            rr.text = "• " + line
            _set_run_font(rr, 12, color_hex=ROLES["text"], name="PingFang TC")

    # Arrows between boxes (aligned with icon row vertically)
    for i in range(2):
        arrow_x = margin + (i + 1) * box_w + (i + 0.5) * 0.50 - 0.07
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(arrow_x), Inches(top + 0.85),
                                    Inches(0.40), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = _rgb(ROLES["accent"])
        arrow.line.fill.background()

    add_footer(s, 9)
    return s


def slide_10_positioning(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "研究定位：read-level epigenetic characterization，非 variant filter",
        "Positioning: read-level epigenetic context — not a variant filter")
    add_two_column(s,
        "不做的（Variant Filter）", [
            "❌ 從 ISM 決定變異 TP/FP",
            "❌ TO-mode methylation-based 分類",
            "❌ germline FP discrimination（AUC<0.64）",
            "❌ pure methylation AUC > 0.58",
        ],
        "要做的（Epigenetic Characterization）", [
            "✓ 描繪 LOH 區 subclonal 表觀結構",
            "✓ HP 下的 fine-grained 異質性",
            "✓ AF × methylation 共同 subclone 證據",
            "✓ 跨樣本可重現的 biology",
        ])
    add_footer(s, 10)
    return s


def slide_11_closed(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "已關閉方向 prevent re-investigation",
        "Closed directions — documented to prevent loops",
        ROLES["warn"])
    add_two_column(s,
        "Closed (NEGATIVE / NO-GO)", [
            "TO Germline FP：G1-G7 AUC<0.64",
            "O11 Heterogeneity：n_reads confound",
            "O12 LOH scenarios：AlleleDelta=AF confound",
            "O13 Cross-region：shared reads confound",
            "Wave 3 Non-LOH：Simpson's Paradox",
            "Fine-Pairwise：6 distances 全無效",
        ],
        "Lesson learned", [
            "AF 與 methylation 經常 confound，必須 stratify",
            "HP-free AUC<0.53 證明信號來自 HP phasing",
            "pure methylation ≤0.58 限制已確認",
            "TO-pure caller_af=0.654 超越所有 ISM",
            "跨區域訊號源自 shared reads 而非 biology",
            "2×2 CramersV 93% 零值：特徵空間耗盡",
        ],
        top_in=LAYOUT["evidence_top"] + 0.10)
    add_footer(s, 11)
    return s


def slide_12_phase2_priority(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "Phase 2 方向 A+D 為當前最高優先",
        "Phase 2 (A: Normal BAM + D: LOH BED annotation) is top priority",
        ROLES["success"])
    margin = LAYOUT["margin_x"]
    top = LAYOUT["evidence_top"] + 0.20

    # Horizontal roadmap: 4 milestones
    milestones = [
        ("方向 A", "Normal BAM 整合", "讀入 normal read，建立 per-sample ASM baseline", ROLES["primary"]),
        ("方向 D", "LOH BED 標註", "將 LOH 區域標註至輸出，支援 per-segment 分析", ROLES["accent"]),
        ("方向 B", "Sample-level ASM", "分離 normal 與 tumor 甲基化趨勢", ROLES["neutral"]),
        ("方向 C", "Cross-region subclone", "跨 region subclone 聯合建模", FINDING_COLORS["F5_4group_strat"]),
    ]
    col_w = (CANVAS["width_in"] - 2 * margin - 1.5) / 4
    for i, (tag, title, desc, col) in enumerate(milestones):
        x = margin + i * (col_w + 0.50)
        # Tag circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + col_w / 2 - 0.35), Inches(top),
                                     Inches(0.70), Inches(0.70))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(col)
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.margin_left = ctf.margin_right = 0
        ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = tag
        _set_run_font(cr, 13, color_hex="#FFFFFF", bold=True, name="PingFang TC")
        # Title
        tbox = s.shapes.add_textbox(Inches(x), Inches(top + 0.85),
                                     Inches(col_w), Inches(0.45))
        tbtf = tbox.text_frame; tbtf.word_wrap = True
        tbtf.margin_left = tbtf.margin_right = 0
        tp = tbtf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = title
        _set_run_font(tr, 15, color_hex=ROLES["text"], bold=True, name="PingFang TC")
        # Description
        dbox = s.shapes.add_textbox(Inches(x), Inches(top + 1.40),
                                     Inches(col_w), Inches(1.8))
        dtf = dbox.text_frame; dtf.word_wrap = True
        dtf.margin_left = dtf.margin_right = 0
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        dp.line_spacing = 1.30
        dr = dp.add_run()
        dr.text = desc
        _set_run_font(dr, 12, color_hex=ROLES["text_light"], name="PingFang TC")
        # Arrow to next
        if i < 3:
            ax = x + col_w + 0.10
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        Inches(ax), Inches(top + 0.22),
                                        Inches(0.30), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = _rgb(ROLES["rule"])
            arrow.line.fill.background()

    # Priority bar at bottom
    bar_y = top + 3.6
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(margin), Inches(bar_y),
                              Inches(CANVAS["width_in"] - 2 * margin), Inches(0.70))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(ROLES["success"])
    bar.line.fill.background()
    btf = bar.text_frame
    btf.margin_left = Inches(0.30)
    btf.margin_top = Inches(0.15)
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = "當前執行：方向 A+D 並行開發    |    下一里程碑：ReadParser Normal BAM 支援 + LOH.bed 寫入 (本月內)"
    _set_run_font(br, 13, color_hex="#FFFFFF", bold=True, name="PingFang TC")

    add_footer(s, 12)
    return s


def slide_13_limitations(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "限制：7 樣本 × 有限 coverage 的統計與生物邊界",
        "Limits: 7 samples, limited coverage — statistical and biological boundaries",
        ROLES["warn"])
    margin = LAYOUT["margin_x"]
    items = [
        "樣本規模 n = 7 細胞株，外推至 primary tumor 需驗證",
        "ONT coverage 中位數 ~30×，低 AF 區統計功率受限",
        "LOH 範圍取自 ClairS + Nanocaller 整合，邊界定義仍有模糊",
        "NGroups 受 read count confound，需 within-group 回歸殘差化",
        "Methylation 訊號 L2 殘差化有 collider bias 風險",
        "缺 normal sample 平行對照，Phase 2A 即將補齊",
    ]
    add_bullet_list(s, items, margin + 0.30, LAYOUT["evidence_top"] + 0.30,
                    CANVAS["width_in"] - 2 * margin - 0.60, 4.5,
                    font_size=16, line_spacing=1.55)
    add_footer(s, 13)
    return s


def slide_14_next_week(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "下週 3 個可驗證交付物",
        "Next week: 3 verifiable deliverables",
        ROLES["success"])
    margin = LAYOUT["margin_x"]
    top = LAYOUT["evidence_top"] + 0.20

    deliverables = [
        ("D1", "Normal BAM ReadParser",
         "✓ ReadParser 讀入 normal BAM 路徑\n✓ 單元測試 normal-only / tumor-only / mixed 三情境\n✓ F1 差異 < 0.01"),
        ("D2", "LOH.bed 寫入 output",
         "✓ 每 run 輸出 LOH.bed\n✓ bedtools intersect 驗證範圍正確\n✓ 7 樣本跑通"),
        ("D3", "Subclone stratified 分析報告",
         "✓ 4-group stratification 全 7 樣本\n✓ AF × methylation 交叉表\n✓ docs/experiments 更新"),
    ]
    col_w = (CANVAS["width_in"] - 2 * margin - 1.0) / 3
    for i, (tag, title, body) in enumerate(deliverables):
        x = margin + i * (col_w + 0.50)
        # Badge
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(top),
                                    Inches(col_w), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(ROLES["success"])
        badge.line.fill.background()
        bf = badge.text_frame
        bf.margin_left = Inches(0.20); bf.margin_top = Inches(0.08)
        bp = bf.paragraphs[0]
        br = bp.add_run()
        br.text = f"{tag}   {title}"
        _set_run_font(br, 15, color_hex="#FFFFFF", bold=True, name="PingFang TC")
        # Body card
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(top + 0.55),
                                   Inches(col_w), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(ROLES["bg_panel"])
        card.line.color.rgb = _rgb(ROLES["rule"])
        card.line.width = Pt(0.5)
        cf = card.text_frame
        cf.word_wrap = True
        cf.margin_left = Inches(0.22); cf.margin_right = Inches(0.22)
        cf.margin_top = Inches(0.20)
        for j, line in enumerate(body.split("\n")):
            pp = cf.paragraphs[0] if j == 0 else cf.add_paragraph()
            pp.line_spacing = 1.45
            rr = pp.add_run()
            rr.text = line
            _set_run_font(rr, 13, color_hex=ROLES["text"], name="PingFang TC")

    add_footer(s, 14)
    return s


def slide_15_qa(prs):
    s = add_blank_slide(prs)
    margin = LAYOUT["margin_x"]
    # Large centered text
    box = s.shapes.add_textbox(Inches(margin), Inches(2.8),
                                Inches(CANVAS["width_in"] - 2 * margin), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "問題與討論"
    _set_run_font(r, 56, color_hex=ROLES["text"], bold=True, name="PingFang TC")
    # EN
    box2 = s.shapes.add_textbox(Inches(margin), Inches(4.1),
                                 Inches(CANVAS["width_in"] - 2 * margin), Inches(0.7))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Questions & Discussion"
    _set_run_font(r2, 24, color_hex=ROLES["text_light"], name="Helvetica")

    # Accent line centered
    line_w = 1.5
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches((CANVAS["width_in"] - line_w) / 2),
                               Inches(4.80),
                               Inches(line_w), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(ROLES["accent"])
    rule.line.fill.background()

    add_footer(s, 15)
    return s


def slide_16_appendix(prs):
    s = add_blank_slide(prs)
    add_claim_header(s,
        "Appendix：完整方法與原始數據連結",
        "Methods & data reference table")

    margin = LAYOUT["margin_x"]
    top = LAYOUT["evidence_top"] + 0.10
    table_w = CANVAS["width_in"] - 2 * margin

    # Header row
    header_y = top
    col_widths = [3.0, 5.0, table_w - 8.0]
    headers = ["類別", "內容", "位置"]
    x_cursor = margin
    for i, h in enumerate(headers):
        hbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x_cursor), Inches(header_y),
                                   Inches(col_widths[i]), Inches(0.45))
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = _rgb(ROLES["text"])
        hbox.line.fill.background()
        htf = hbox.text_frame
        htf.margin_left = Inches(0.15); htf.margin_top = Inches(0.08)
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = h
        _set_run_font(hr, 13, color_hex="#FFFFFF", bold=True, name="PingFang TC")
        x_cursor += col_widths[i]

    rows = [
        ("方法：AF 分布", "7 樣本合併 AF density + intermediate AF proportion",
         "docs/research/loh_subclone_af_methylation/"),
        ("方法：NGroups", "HP-aware methylation clustering (k=auto)",
         "src/core/methylation_clustering.cpp"),
        ("方法：AlleleDelta", "per-variant methylation delta between HP1/HP2",
         "src/core/allele_delta.cpp"),
        ("統計：效應量", "Cohen's d; Mann–Whitney U; Spearman ρ",
         "docs/methodology/20260315_effect_size_guide.md"),
        ("資料：LOH BED", "ClairS ∩ Nanocaller, 7 samples",
         "data/loh/*.bed"),
        ("參考：closed dirs", "MEMORY 14 conclusions, 8 evidence chains",
         "docs/reports/research_landscape/00_INDEX.md"),
    ]
    row_y = header_y + 0.45
    for j, row in enumerate(rows):
        x_cursor = margin
        fill_col = ROLES["bg"] if j % 2 == 0 else ROLES["bg_panel"]
        for i, cell in enumerate(row):
            rbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x_cursor), Inches(row_y + j * 0.55),
                                       Inches(col_widths[i]), Inches(0.55))
            rbox.fill.solid()
            rbox.fill.fore_color.rgb = _rgb(fill_col)
            rbox.line.color.rgb = _rgb(ROLES["rule_light"])
            rbox.line.width = Pt(0.25)
            rtf = rbox.text_frame
            rtf.word_wrap = True
            rtf.margin_left = Inches(0.15); rtf.margin_top = Inches(0.10)
            rp = rtf.paragraphs[0]
            rr = rp.add_run()
            rr.text = cell
            _set_run_font(rr, 11, color_hex=ROLES["text"], name="PingFang TC")
            x_cursor += col_widths[i]

    add_footer(s, 16)
    return s


# ---------- Main ----------
SLIDE_BUILDERS = [
    slide_01_cover,
    slide_02_three_questions,
    slide_03_intermediate_af,
    slide_04_delta_ngroups,
    slide_05_allele_delta,
    slide_06_segment_rho,
    slide_07_4group,
    slide_08_hpfine_marker,
    slide_09_dual_evidence,
    slide_10_positioning,
    slide_11_closed,
    slide_12_phase2_priority,
    slide_13_limitations,
    slide_14_next_week,
    slide_15_qa,
    slide_16_appendix,
]


def main():
    print("=" * 60)
    print("Building Weekly Report v2 PPTX (Assertion-Evidence)")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(CANVAS["width_in"])
    prs.slide_height = Inches(CANVAS["height_in"])

    assert len(SLIDE_BUILDERS) == TOTAL_SLIDES, \
        f"Expected {TOTAL_SLIDES} builders, got {len(SLIDE_BUILDERS)}"

    for i, builder in enumerate(SLIDE_BUILDERS, 1):
        print(f"  Slide {i:02d}: {builder.__name__} ...", end=" ")
        try:
            builder(prs)
            print("OK")
        except Exception as e:
            print(f"FAIL: {e}")
            raise

    out_path = os.path.join(OUT_DIR, "20260417_週報v2_AE.pptx")
    prs.save(out_path)
    size_kb = os.path.getsize(out_path) / 1024
    print("=" * 60)
    print(f"Saved: {os.path.relpath(out_path)} ({size_kb:.1f} KB, {len(prs.slides)} slides)")

    # Auto-snapshot (reuse pptx_snapshot.py from v1)
    try:
        v1_dir = os.path.abspath(os.path.join(HERE, "..", "..", "..", "..",
                                               "validated", "2026", "04",
                                               "20260414_研究週報_LOH_Subclone_AF_Methylation"))
        sys.path.insert(0, v1_dir)
        from pptx_snapshot import take_snapshot
        snap_path = os.path.join(OUT_DIR, "_snapshot.json")
        slide_index_map = {i: b.__name__ for i, b in enumerate(SLIDE_BUILDERS)}
        take_snapshot(out_path, snap_path, slide_index_map=slide_index_map)
        print(f"Snapshot: {os.path.relpath(snap_path)}")
    except Exception as e:
        print(f"(Snapshot skipped: {e})")

    return 0


if __name__ == "__main__":
    sys.exit(main())
