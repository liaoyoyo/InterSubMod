#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_pptx.py — Self-Phasing V3 精簡教授報告 (12 slides, 16:9).

V3 = "高密度精簡版"
    - one-slide-one-message
    - 大圖示低認知負荷
    - speaker note >= 500 字 (因 slide 變少, note 比 v2 350 下限更長)

Helpers/THEME 直接複用 v2 的 build_pptx.py — 不修改 helper 行為。
所有自繪 schematic 改寫成 V3 視覺強化版 (大字、簡潔、單一焦點)。

Slide structure (12 張):
    S1  17.3:1 -> 1:1 大字衝擊 banner
    S2  Pipeline 一覽 (fig1 為主)
    S3  HP tag 5 整數值 + 三層證據 (fig17 + fig01d 並排)
    S4  Phasing OK, Tag 有問題 (左右兩欄)
    S5  ISM 影響 29/14/42 (03_self_phasing_impact)
    S6  ★ 根因樹 (Purity 0.95 主者 + 三 bug)
    S7  ★ V5 三層投票流程
    S8  Sanity 15/15 + 5 證據鏈
    S9  量化指標 AMB/HP33/Concordance (fig18)
    S10 V5max1 climax (39 reads 翻轉)
    S11 業界家族樹 + 五大目標
    S12 Take-home + 3 P0 + Q&A

Strict rendering rules (繼承 v2 PPT_RENDERING_RULES.md):
    1. Per-character Latin + CJK font fallback
    2. fit_image_within (等比, 不雙軸強制)
    3. Speaker note 強制 >= 500 字 (V3 標準, 比 v2 嚴)
    4. Slide title <= 15 CJK chars
    5. 禁原生 emoji (用 ASCII / box symbols)
"""

from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FIG_DIR = ROOT / "figures"          # symlink to ../v2/figures
OUT_DIR = ROOT / "output"
AUTO_FIG_DIR = OUT_DIR / "auto_figs"
OUT_DIR.mkdir(exist_ok=True)
AUTO_FIG_DIR.mkdir(exist_ok=True)

OUTPUT_PPTX = OUT_DIR / "20260429_Self_Phasing_精簡教授報告_v3.pptx"
TOTAL_SLIDES = 12

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------
CJK_FONT_PATH = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
LATIN_FONT = "Arial"
CJK_FONT_NAME = "Droid Sans Fallback"

if Path(CJK_FONT_PATH).exists():
    fm.fontManager.addfont(CJK_FONT_PATH)
    _font = fm.FontProperties(fname=CJK_FONT_PATH).get_name()
    plt.rcParams.update({
        "font.family": [_font, "DejaVu Sans"],
        "axes.unicode_minus": False,
    })
else:
    print(f"[WARN] CJK font not found at {CJK_FONT_PATH}; CJK in matplotlib may render as boxes.")

# ---------------------------------------------------------------------------
# Slide geometry (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Theme colours (繼承 v2)
C_TITLE_BG = RGBColor(0x1F, 0x4E, 0x79)
C_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
C_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_AMBER = RGBColor(0xE6, 0x7E, 0x22)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
C_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_BORDER = RGBColor(0xBD, 0xC3, 0xC7)
C_BLUE = RGBColor(0x29, 0x80, 0xB9)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)

# ---------------------------------------------------------------------------
# Font fallback helpers (verbatim copy from v2)
# ---------------------------------------------------------------------------
CJK_RANGES = (
    (0x3000, 0x303F),
    (0x3040, 0x30FF),
    (0x3400, 0x4DBF),
    (0x4E00, 0x9FFF),
    (0xF900, 0xFAFF),
    (0xFE30, 0xFE4F),
    (0xFF00, 0xFFEF),
)


def _is_cjk(ch: str) -> bool:
    cp = ord(ch)
    for lo, hi in CJK_RANGES:
        if lo <= cp <= hi:
            return True
    return False


def _segments(text: str):
    if not text:
        return []
    out = []
    buf = [text[0]]
    cur_is_cjk = _is_cjk(text[0])
    for ch in text[1:]:
        if _is_cjk(ch) == cur_is_cjk:
            buf.append(ch)
        else:
            out.append(("".join(buf), cur_is_cjk))
            buf = [ch]
            cur_is_cjk = _is_cjk(ch)
    out.append(("".join(buf), cur_is_cjk))
    return out


def add_text_with_fallback(text_frame, text, font_size=18, bold=False,
                           color=None, latin=LATIN_FONT, cjk=CJK_FONT_NAME,
                           align=PP_ALIGN.LEFT, paragraph_index=None,
                           italic=False):
    """Append paragraph with per-char Latin/CJK font fallback."""
    if paragraph_index is not None and paragraph_index == 0 and text_frame.paragraphs[0].text == "":
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = align

    for seg, is_cjk in _segments(text):
        run = p.add_run()
        run.text = seg
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        run.font.italic = bool(italic)
        if color is not None:
            run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea"):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        latin_el = etree.SubElement(rPr, qn("a:latin"))
        latin_el.set("typeface", latin)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", cjk)
    return p


# ---------------------------------------------------------------------------
# Generic shape helpers (verbatim copy from v2)
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w_pt=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)
    shp.text_frame.margin_left = Emu(60000)
    shp.text_frame.margin_right = Emu(60000)
    shp.text_frame.margin_top = Emu(40000)
    shp.text_frame.margin_bottom = Emu(40000)
    return shp


def add_text_block(slide, x, y, w, h, lines, font_size=14, bold=False,
                   color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   fill=None, line=None, italic=False):
    if fill is not None or line is not None:
        shp = add_rect(slide, x, y, w, h, fill=fill, line=line)
    else:
        shp = slide.shapes.add_textbox(x, y, w, h)
        shp.text_frame.margin_left = Emu(60000)
        shp.text_frame.margin_right = Emu(60000)
        shp.text_frame.margin_top = Emu(40000)
        shp.text_frame.margin_bottom = Emu(40000)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        b = bold if not isinstance(bold, list) else bold[i] if i < len(bold) else False
        sz = font_size if not isinstance(font_size, list) else font_size[i] if i < len(font_size) else font_size[-1]
        col = color if not isinstance(color, list) else color[i] if i < len(color) else color[-1]
        add_text_with_fallback(tf, ln, font_size=sz, bold=b, color=col,
                               align=align, paragraph_index=i, italic=italic)
    return shp


def add_title(slide, text, subtitle=None):
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=C_TITLE_BG)
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = ""
    add_text_with_fallback(tf, text, font_size=26, bold=True,
                           color=C_TITLE_FG, paragraph_index=0)
    if subtitle:
        add_text_with_fallback(tf, subtitle, font_size=14, bold=False,
                               color=RGBColor(0xCF, 0xE2, 0xF3))


def fit_image_within(slide, path, x, y, max_w, max_h, border=False):
    """Add picture preserving aspect ratio inside (max_w, max_h)."""
    if not Path(path).exists():
        add_rect(slide, x, y, max_w, max_h, fill=C_LIGHT, line=C_ACCENT)
        add_text_block(slide, x, y + max_h // 2 - Inches(0.2), max_w, Inches(0.4),
                       [f"[圖片缺失] {Path(path).name}"],
                       font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
        return None
    img = Image.open(path)
    iw, ih = img.size
    img.close()
    iw_emu = iw * 9525
    ih_emu = ih * 9525
    ratio = min(max_w / iw_emu, max_h / ih_emu)
    final_w = int(iw_emu * ratio)
    final_h = int(ih_emu * ratio)
    cx = x + (max_w - final_w) // 2
    cy = y + (max_h - final_h) // 2
    pic = slide.shapes.add_picture(path, cx, cy, width=final_w, height=final_h)
    if border:
        add_rect(slide, cx, cy, final_w, final_h, fill=None, line=C_BORDER)
    return pic


def set_speaker_note(slide, text, min_chars=500):
    """Set speaker note (mandatory; >= min_chars). V3 default 500 chars."""
    if not text or not text.strip():
        raise ValueError(f"Empty speaker note for slide {slide.slide_id}")
    if len(text) < min_chars:
        raise ValueError(f"Speaker note too short ({len(text)} chars) for slide {slide.slide_id}; need >= {min_chars}")
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    add_text_with_fallback(notes, text, font_size=14, paragraph_index=0)


def add_footer(slide, page_no, total=TOTAL_SLIDES):
    add_text_block(slide, Inches(0.3), Inches(7.05), Inches(9), Inches(0.35),
                   ["Self-Phasing V3 精簡 · 2026-04-30 · longphase-to-mod V5 (TO + PON)"],
                   font_size=10, color=C_GRAY)
    add_text_block(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.35),
                   [f"{page_no} / {total}"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


def _save_fig(name, fig, dpi=160):
    out = AUTO_FIG_DIR / name
    fig.savefig(out, dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


# ---------------------------------------------------------------------------
# V3 自繪 schematic figures
# ---------------------------------------------------------------------------
def make_v3_s1_banner() -> Path:
    """S1 大字衝擊 banner — 17.3:1 -> 1:1 + 副標."""
    fig, ax = plt.subplots(figsize=(13.2, 6.8))
    ax.set_xlim(0, 13.2)
    ax.set_ylim(0, 6.8)
    ax.axis("off")

    # 大字: 左 17.3:1 紅, 中 ->, 右 1:1 綠
    ax.text(2.8, 4.2, "17.3 : 1", ha="center", va="center",
            fontsize=110, fontweight="bold", color="#c0392b")
    ax.text(6.6, 4.2, "→", ha="center", va="center",
            fontsize=120, fontweight="bold", color="#27ae60")
    ax.text(10.4, 4.2, "1 : 1", ha="center", va="center",
            fontsize=110, fontweight="bold", color="#27ae60")

    # 副標
    ax.text(6.6, 2.2,
            "Tumor-only somatic HP tag bias 已修補",
            ha="center", va="center",
            fontsize=22, color="#1f4e79", fontweight="bold")
    ax.text(6.6, 1.55,
            "read-level concordance +8.3 pp (clean PS blocks, 全基因組)",
            ha="center", va="center",
            fontsize=18, color="#27ae60", fontweight="bold")

    # 底部 caveat strip
    ax.add_patch(plt.Rectangle((1.0, 0.3), 11.2, 0.75,
                               facecolor="#fdf6e3", edgecolor="#a04000", lw=1.5))
    ax.text(6.6, 0.67,
            "修補 self-phasing 是解鎖 InterSubMod 五大研究目標的前提",
            ha="center", va="center",
            fontsize=14, color="#a04000", fontweight="bold")

    # 上方標籤
    ax.text(2.8, 5.55, "Baseline (HP1 family : HP2 family)",
            ha="center", va="center", fontsize=12, color="#7f8c8d")
    ax.text(10.4, 5.55, "V5 (生物學期望)",
            ha="center", va="center", fontsize=12, color="#7f8c8d")

    return _save_fig("v3_s1_banner.png", fig)


def make_v3_s4_phasing_vs_tag() -> Path:
    """S4 左右兩欄拆鎖 — phasing 層 OK (綠) vs tag 層 FAIL (紅)."""
    fig, ax = plt.subplots(figsize=(13.2, 6.4))
    ax.set_xlim(0, 13.2)
    ax.set_ylim(0, 6.4)
    ax.axis("off")

    # 左欄: phasing 層 (綠)
    ax.add_patch(plt.Rectangle((0.3, 0.5), 6.0, 5.6,
                               facecolor="#e8f8f0", edgecolor="#27ae60", lw=3.0))
    ax.text(3.3, 5.65, "Phasing 層", ha="center", va="center",
            fontsize=22, fontweight="bold", color="#27ae60")
    ax.text(3.3, 5.20, "(LOH.bed, region-level)",
            ha="center", va="center", fontsize=12, color="#444")
    ax.text(3.3, 4.75, "✓  PASS", ha="center", va="center",
            fontsize=28, fontweight="bold", color="#27ae60")
    rows_left = [
        ("Jaccard = 1.0", "V5 前後完全不變"),
        ("LOH.bed 不變", "VCF AD 來源穩定"),
        ("r = 0.001", "TO vs paired 同位點"),
        ("Cohen's d = -1.20", "巨大效應量"),
    ]
    cy = 3.85
    for k, v in rows_left:
        ax.text(0.6, cy, "•  " + k, fontsize=14, fontweight="bold", color="#1f4e79")
        ax.text(0.9, cy - 0.32, v, fontsize=11, color="#555")
        cy -= 0.78

    # 右欄: tag 層 (紅)
    ax.add_patch(plt.Rectangle((6.9, 0.5), 6.0, 5.6,
                               facecolor="#fdedec", edgecolor="#c0392b", lw=3.0))
    ax.text(9.9, 5.65, "Tag 層", ha="center", va="center",
            fontsize=22, fontweight="bold", color="#c0392b")
    ax.text(9.9, 5.20, "(BAM HP_Ratio, read-level)",
            ha="center", va="center", fontsize=12, color="#444")
    ax.text(9.9, 4.75, "●  FAIL", ha="center", va="center",
            fontsize=28, fontweight="bold", color="#c0392b")
    rows_right = [
        ("HP_Ratio 17.3 : 1", "baseline 偏向 HP1"),
        ("62% LOH = artifact", "ISM HP_Ratio LOH"),
        ("kappa = 0.67", "兩 LOH 系統不一致"),
        ("AF 0.1-0.8 近 100% 假", "中等 AF 最嚴重"),
    ]
    cy = 3.85
    for k, v in rows_right:
        ax.text(7.2, cy, "•  " + k, fontsize=14, fontweight="bold", color="#922b21")
        ax.text(7.5, cy - 0.32, v, fontsize=11, color="#555")
        cy -= 0.78

    # 底部結論
    ax.text(6.6, 0.18,
            "本 PPT 焦點 = tag 層;  phase 層的 self-phasing scaffold 已由 V2b PON-only 處理",
            ha="center", va="center",
            fontsize=12, color="#1f4e79", fontweight="bold")
    return _save_fig("v3_s4_phasing_vs_tag.png", fig)


def make_v3_s6_root_cause_tree() -> Path:
    """S6 根因樹 — Purity 0.95 觸發鎖 (root) + 三 bug (leaves)."""
    fig, ax = plt.subplots(figsize=(13.4, 6.6))
    ax.set_xlim(0, 13.4)
    ax.set_ylim(0, 6.6)
    ax.axis("off")

    # Root: 觸發條件 (大紅框)
    ax.add_patch(plt.Rectangle((2.4, 5.30), 8.6, 1.10,
                               facecolor="#fdedec", edgecolor="#c0392b", lw=3.0))
    ax.text(6.7, 6.10, "★  觸發條件 (主者)",
            ha="center", va="center",
            fontsize=14, fontweight="bold", color="#c0392b")
    ax.text(6.7, 5.78, "Baseline 估 purity = 0.927 ≤ 0.95 (純樣本標準)",
            ha="center", va="center",
            fontsize=13, color="#222", fontweight="bold")
    ax.text(6.7, 5.45, "→ 未觸發 Two-Pass → 走 baseline 標準三條路流程",
            ha="center", va="center",
            fontsize=12, color="#222")

    # Three branches (arrows)
    ax.annotate("", xy=(2.3, 4.55), xytext=(4.3, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(6.7, 4.55), xytext=(6.7, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(11.1, 4.55), xytext=(9.1, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))

    # Bug 1 (紅, getVote priority)
    ax.add_patch(plt.Rectangle((0.3, 1.5), 4.0, 3.0,
                               facecolor="#fdedec", edgecolor="#c0392b", lw=2.5))
    ax.text(2.3, 4.20, "Bug 1", ha="center",
            fontsize=14, fontweight="bold", color="#c0392b")
    ax.text(2.3, 3.85, "getVote priority", ha="center",
            fontsize=12, fontweight="bold", color="#c0392b")
    ax.text(2.3, 3.45, "somatic-first 投票順序",
            ha="center", fontsize=11, color="#222")
    ax.text(2.3, 3.10, "{HP1_1, HP2_1} 排第一",
            ha="center", fontsize=10, color="#222", family="monospace")
    ax.text(2.3, 2.75, "任一 somatic > 0 即 break",
            ha="center", fontsize=10, color="#222")
    ax.text(2.3, 2.40, "germline 完全跳過",
            ha="center", fontsize=10, color="#222")
    ax.text(2.3, 1.85, "✓ V3F 41ff147 修",
            ha="center", fontsize=11, fontweight="bold", color="#c0392b")

    # Bug 2 (橘, enum vs int)
    ax.add_patch(plt.Rectangle((4.7, 1.5), 4.0, 3.0,
                               facecolor="#fef5e7", edgecolor="#e67e22", lw=2.5))
    ax.text(6.7, 4.20, "Bug 2", ha="center",
            fontsize=14, fontweight="bold", color="#e67e22")
    ax.text(6.7, 3.85, "enum vs int mismatch",
            ha="center", fontsize=12, fontweight="bold", color="#e67e22")
    ax.text(6.7, 3.45, "Util.h:20-25",
            ha="center", fontsize=10, color="#222", family="monospace")
    ax.text(6.7, 3.10, "HAPLOTYPE1_1 = 3 (enum)",
            ha="center", fontsize=10, color="#222", family="monospace")
    ax.text(6.7, 2.75, "HP tag integer = 11",
            ha="center", fontsize=10, color="#222", family="monospace")
    ax.text(6.7, 2.40, "→ HP:i:33 永不出現",
            ha="center", fontsize=10, color="#222")
    ax.text(6.7, 1.85, "✓ V3F 41ff147 修",
            ha="center", fontsize=11, fontweight="bold", color="#e67e22")

    # Bug 3 (灰, scaffold V2b 已解)
    ax.add_patch(plt.Rectangle((9.1, 1.5), 4.0, 3.0,
                               facecolor="#ecf0f1", edgecolor="#7f8c8d", lw=1.8))
    ax.text(11.1, 4.20, "Bug 3", ha="center",
            fontsize=14, fontweight="bold", color="#7f8c8d")
    ax.text(11.1, 3.85, "phase scaffold (已解)",
            ha="center", fontsize=12, fontweight="bold", color="#7f8c8d")
    ax.text(11.1, 3.45, "PhasingProcess.cpp",
            ha="center", fontsize=10, color="#444", family="monospace")
    ax.text(11.1, 3.10, "convertNonGermlineToSomatic",
            ha="center", fontsize=9, color="#444", family="monospace")
    ax.text(11.1, 2.75, "未觸發 → somatic 當 anchor",
            ha="center", fontsize=10, color="#444")
    ax.text(11.1, 2.40, "(已由 V2b 8b8c1fd 解)",
            ha="center", fontsize=10, color="#444", style="italic")
    ax.text(11.1, 1.85, "scaffold 已解, 焦點在 1+2",
            ha="center", fontsize=10, fontweight="bold", color="#7f8c8d")

    # Bottom 結果
    ax.text(6.7, 1.05,
            "→ 暴露 tag 層 somatic-first 投票 bias",
            ha="center", fontsize=12, color="#1f4e79", fontweight="bold")
    ax.text(6.7, 0.55,
            "結果:  self-phasing  17.3 : 1  artifact",
            ha="center", fontsize=14, fontweight="bold", color="#c0392b")
    return _save_fig("v3_s6_root_cause_tree.png", fig)


def make_v3_s7_three_layer_voting() -> Path:
    """S7 V5 三層投票流程 — 垂直堆疊 layer 1 / 1.5 / 2."""
    fig, ax = plt.subplots(figsize=(13.2, 6.6))
    ax.set_xlim(0, 13.2)
    ax.set_ylim(0, 6.6)
    ax.axis("off")

    # Top label
    ax.text(6.6, 6.30, "V5 getVote() — 三層投票流程",
            ha="center", va="center", fontsize=18, fontweight="bold", color="#1f4e79")

    layers = [
        # (y_top, color, name, body, ring_text)
        (5.55, "#3498db", "Layer 1:  germline first",
         "if (germlineHP1 > germlineHP2)  → 1\nelif (germlineHP2 > germlineHP1)  → 2\nelse  → 進入 Layer 1.5",
         "★ 翻轉 baseline 的 somatic-first 投票順序"),
        (3.85, "#c0392b", "Layer 1.5:  somatic fallback  (V5 新增)",
         "if (sTotal == 0)  → 0\nelif (s1 > 0.6 × tot)  → align HP1\nelif (s2 > 0.6 × tot)  → align HP2\nelse  → 33 (true ambiguous)",
         "★ Confidence 0.6 = mid-low purity 防禦層"),
        (2.15, "#27ae60", "Layer 2:  encode hpResult",
         "1 + somatic  → HP:i:11    |    2 + somatic  → HP:i:21\n0 + somatic  → HP:i:33    |    germline only  → HP:i:1 / 2\nnothing  → HP:i:0",
         "對應 5 種 BAM HP integer 編碼"),
    ]
    for y_top, color, name, body, ring in layers:
        # Layer box
        ax.add_patch(plt.Rectangle((1.0, y_top - 1.50), 11.2, 1.50,
                                   facecolor=color, alpha=0.13,
                                   edgecolor=color, lw=2.2))
        ax.text(1.2, y_top - 0.30, name, fontsize=14, fontweight="bold", color=color)
        ax.text(1.2, y_top - 0.55, body, fontsize=10.5, color="#222",
                family="monospace", va="top")
        # Ring marker (key idea, right side)
        ax.add_patch(plt.Rectangle((8.6, y_top - 1.40), 3.5, 0.45,
                                   facecolor="#fdedec", edgecolor="#c0392b", lw=1.5))
        ax.text(10.35, y_top - 1.18, ring,
                ha="center", va="center", fontsize=9, fontweight="bold", color="#c0392b")

    # Arrows between layers
    ax.annotate("", xy=(6.6, 4.05), xytext=(6.6, 4.20),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(6.6, 2.35), xytext=(6.6, 2.50),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))

    # Bottom: 量化效果
    ax.add_patch(plt.Rectangle((1.0, 0.20), 11.2, 0.85,
                               facecolor="#fef5e7", edgecolor="#a04000", lw=1.5))
    ax.text(6.6, 0.78, "效果:  AMB%  17.5% → 8.0% (-9.5pp)   ·   HP:i:33 reads -54%",
            ha="center", fontsize=12, fontweight="bold", color="#a04000")
    ax.text(6.6, 0.42, "介面契約零變動 (HaplotagProcess.h:66-68 一字未變, drop-in replace)",
            ha="center", fontsize=10, color="#555")
    return _save_fig("v3_s7_three_layer_voting.png", fig)


def make_v3_s8_sanity_card() -> Path:
    """S8 守恆律公式卡 + 5 證據鏈."""
    fig, ax = plt.subplots(figsize=(13.2, 6.4))
    ax.set_xlim(0, 13.2)
    ax.set_ylim(0, 6.4)
    ax.axis("off")

    # 上半: 守恆律公式 (兩個方框)
    ax.text(6.6, 6.15, "4 項硬性 sanity check  —  15 / 15 PASS",
            ha="center", va="center", fontsize=15, fontweight="bold", color="#1f4e79")

    # 公式 A
    ax.add_patch(plt.Rectangle((0.5, 4.10), 6.0, 1.65,
                               facecolor="#e8f8f0", edgecolor="#27ae60", lw=2.5))
    ax.text(3.5, 5.40, "公式 A · Δ-consistency",
            ha="center", fontsize=12, fontweight="bold", color="#27ae60")
    ax.text(3.5, 4.85, "ΔHP33 + (ΔHP11 + ΔHP21) = 0",
            ha="center", fontsize=14, color="#222", family="monospace", fontweight="bold")
    ax.text(3.5, 4.40, "tag 移轉量平衡  ·  in = out",
            ha="center", fontsize=10, color="#444")
    ax.text(3.5, 4.20, "✓  15 / 15 PASS",
            ha="center", fontsize=12, fontweight="bold", color="#27ae60")

    # 公式 B
    ax.add_patch(plt.Rectangle((6.7, 4.10), 6.0, 1.65,
                               facecolor="#e8f8f0", edgecolor="#27ae60", lw=2.5))
    ax.text(9.7, 5.40, "公式 B · Germline 不變",
            ha="center", fontsize=12, fontweight="bold", color="#27ae60")
    ax.text(9.7, 4.85, "HP1 / HP2 reads:  V3F == V5",
            ha="center", fontsize=14, color="#222", family="monospace", fontweight="bold")
    ax.text(9.7, 4.40, "0 reads 漂移  ·  逐 site 比對",
            ha="center", fontsize=10, color="#444")
    ax.text(9.7, 4.20, "✓  15 / 15 PASS",
            ha="center", fontsize=12, fontweight="bold", color="#27ae60")

    # 下半: 5 證據鏈 (一行 5 列)
    ax.text(6.6, 3.55, "5 條獨立證據鏈同步收斂",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#c0392b")

    chains = [
        (0.5, "1.  理論層", "germline het 隨機\n→ 預期 1:1", "#3498db"),
        (3.0, "2.  全基因組層", "HP1:HP2 = 17.3:1\n跨 23 chr 一致", "#c0392b"),
        (5.5, "3.  個別位點層", "SP1 chr19:17565944\nbaseline 113:0", "#a04000"),
        (8.0, "4.  Sanity check", "4 項全 PASS\n0 violation", "#27ae60"),
        (10.5, "5.  程式碼最小必要", "+68 / -36 行\n3 函式集中", "#8e44ad"),
    ]
    for x, head, body, color in chains:
        ax.add_patch(plt.Rectangle((x, 0.70), 2.4, 2.55,
                                   facecolor=color, alpha=0.12,
                                   edgecolor=color, lw=1.8))
        ax.text(x + 1.2, 2.85, head, ha="center", fontsize=11, fontweight="bold", color=color)
        ax.text(x + 1.2, 1.85, body, ha="center", fontsize=10, color="#222")
        ax.add_patch(plt.Circle((x + 1.2, 1.10), 0.18, facecolor=color, edgecolor="white", lw=1.5))
        ax.text(x + 1.2, 1.10, "✓", ha="center", va="center",
                fontsize=14, color="white", fontweight="bold")

    ax.text(6.6, 0.30,
            "5 條獨立路徑同步收斂於同一結論  ·  self-phasing 修補必要且充分",
            ha="center", fontsize=11, color="#1f4e79", fontweight="bold")
    return _save_fig("v3_s8_sanity_card.png", fig)


def make_v3_s11_family_and_goals() -> Path:
    """S11 業界家族樹 (上半) + 五大目標卡 (下半)."""
    fig, ax = plt.subplots(figsize=(13.2, 6.6))
    ax.set_xlim(0, 13.2)
    ax.set_ylim(0, 6.6)
    ax.axis("off")

    # 上半: family tree
    ax.text(6.6, 6.30, "業界家族樹 (同實驗室相鄰工作)",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#1f4e79")

    # Root: LongPhase (2020)
    ax.add_patch(plt.Rectangle((4.7, 5.30), 3.8, 0.75,
                               facecolor="#3498db", alpha=0.25, edgecolor="#1f6391", lw=2.2))
    ax.text(6.6, 5.85, "LongPhase  (Lin 2022)", ha="center",
            fontsize=12, fontweight="bold", color="#1f6391")
    ax.text(6.6, 5.50, "germline phasing 基礎", ha="center",
            fontsize=10, color="#444")

    # Branch L: LongPhase-S (2025 paired)
    ax.add_patch(plt.Rectangle((0.5, 3.60), 5.5, 1.05,
                               facecolor="#27ae60", alpha=0.22, edgecolor="#27ae60", lw=2.0))
    ax.text(3.25, 4.40, "LongPhase-S  (bioRxiv 2025)",
            ha="center", fontsize=11.5, fontweight="bold", color="#27ae60")
    ax.text(3.25, 4.05, "Tumor-Normal Paired",
            ha="center", fontsize=10, color="#1f4e79")
    ax.text(3.25, 3.78, "ClairS SNV F1 +4.5%  (paired)",
            ha="center", fontsize=9, color="#444")

    # Branch R: longphase-to-mod V5 (本實作, 紅框 emphasis)
    ax.add_patch(plt.Rectangle((7.2, 3.60), 5.5, 1.05,
                               facecolor="#c0392b", alpha=0.22, edgecolor="#c0392b", lw=3.0))
    ax.text(9.95, 4.40, "longphase-to-mod V5  (本實作)",
            ha="center", fontsize=11.5, fontweight="bold", color="#c0392b")
    ax.text(9.95, 4.05, "Tumor-Only + PON",
            ha="center", fontsize=10, color="#1f4e79")
    ax.text(9.95, 3.78, "+8.3 pp clean PS concordance",
            ha="center", fontsize=9, color="#444")

    # Arrows
    ax.annotate("", xy=(3.25, 4.65), xytext=(5.4, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(9.95, 4.65), xytext=(7.8, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))

    # 下半: 五大目標卡 + 銜接箭頭
    ax.text(6.6, 3.20, "Self-phasing 修補  →  五大目標解鎖",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#c0392b")

    goals = [
        (0.3, "目標 1", "per-CpG\n關聯", "#c0392b", True),
        (2.85, "目標 2", "clone\n結構", "#e67e22", True),
        (5.4, "目標 3", "二次打擊\n順序", "#f39c12", False),
        (7.95, "目標 4", "TO normal\n補強", "#27ae60", True),
        (10.5, "目標 5", "F1 panel", "#8e44ad", False),
    ]
    for x, head, body, color, direct in goals:
        ax.add_patch(plt.Rectangle((x, 0.70), 2.4, 2.20,
                                   facecolor=color, alpha=0.18,
                                   edgecolor=color, lw=2.0 if direct else 1.0))
        ax.text(x + 1.2, 2.55, head, ha="center", fontsize=12, fontweight="bold", color=color)
        ax.text(x + 1.2, 1.85, body, ha="center", fontsize=11, color="#222")
        if direct:
            ax.text(x + 1.2, 1.05, "★ 直接依賴",
                    ha="center", fontsize=9, fontweight="bold", color="#c0392b")
        else:
            ax.text(x + 1.2, 1.05, "部分依賴",
                    ha="center", fontsize=9, color="#7f8c8d")

    # Bottom note
    ax.text(6.6, 0.30,
            "目標 1 / 2 / 4 直接依賴 HP tag 正確  ·  V5 是研究可信度前提",
            ha="center", fontsize=11, fontweight="bold", color="#1f4e79")
    return _save_fig("v3_s11_family_and_goals.png", fig)


# ---------------------------------------------------------------------------
# V3 Slide builders (12 張)
# ---------------------------------------------------------------------------
def slide_01_impact(prs):
    """S1: 17.3:1 -> 1:1 大字衝擊."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    img = make_v3_s1_banner()
    fit_image_within(s, str(img), Inches(0.2), Inches(0.4), Inches(13.0), Inches(6.7))
    add_footer(s, 1)
    set_speaker_note(s,
        "60 秒鎖定核心結論。本 PPT 是 V3 12-slide 高密度精簡版。"
        "**核心訊息**: 大字 17.3 : 1 → 1 : 1 是 HP1 family : HP2 family ratio, "
        "baseline 17.3 倍偏在 HP1 (614,000 vs 35,500 reads, 跨 23 染色體一致), "
        "經 4-commit 漸進修補後回到生物學期望 1:1。"
        "**第二行 +8.3 pp**: 全基因組 clean PS blocks 對 paired ground truth 的吻合度提升 (V5 90.5% vs Baseline 82.2%, "
        "PI 報告 4 §3.7), 必須加 caveat 「(clean PS blocks, 全基因組)」, 不是全基因組 raw。"
        "**底部動機 strip**: 修補 self-phasing 是解鎖 InterSubMod 五大研究目標的前提 — 目標 1 (per-CpG 關聯)、目標 2 (clone 結構)、"
        "目標 4 (TO normal 補強) 直接依賴 HP tag 正確; 目標 3 (二次打擊)、目標 5 (F1 panel) 部分依賴。"
        "**邊界澄清**: 本 PPT 焦點在 longphase-to-mod fork (TO + PON 條件下的本地實作) 4-commit 漸進修補, "
        "InterSubMod 本 repo 無 C++ 改動。修補位置在 HaplotagProcess.cpp (longphase-to-mod fork, "
        "獨立 git repo @ /big7_disk/liaoyoyo2001/longphase-to-mod/), 不在 ISM 內。"
        "**重要分層**: self-phasing 問題鏈是被分層處理的 — V2b 解 phase scaffold (--pon-only-phasing flag), "
        "V3F/V5 解 tag 層 (priority bug + enum mismatch + V3F 過度保守的 directional reassignment)。"
        "V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold)。"
        "**Q&A 預備**: "
        "Q1 (「17.3 怎麼算的?」) — HP1 family = 614,000 reads (germline + somatic on HP1), "
        "HP2 family = 35,500, 614000 / 35500 = 17.3, 跨 23 染色體一致, 94.6% 集中於 HP1; "
        "生物學上 germline het 隨機分配應 ~1:1, 17.3 倍偏差不可能是統計噪音。"
        "Q2 (「+8.3 pp 是怎麼測的?」) — 取 clean PS blocks (germline accuracy ≥ 70% 的 phase block), "
        "對全基因組統計 read-level haplotype assignment 與 paired tumor BAM 的 concordance, "
        "V5 = 90.5% vs Baseline = 82.2%, 差 +8.3 pp。"
        "**過場**: 接下來 S2 用 pipeline 圖澄清「修在哪個位置」, 再進入觀察與根因。")
    return s


def slide_02_pipeline(prs):
    """S2: 工作流程一覽 (fig1_pipeline_comparison)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "工作流程一覽")
    fit_image_within(s, str(FIG_DIR / "fig1_pipeline_comparison.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.7))
    # Bottom 邊界澄清
    add_rect(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.30),
             fill=RGBColor(0xEC, 0xF0, 0xF1), line=C_BORDER, line_w_pt=1.0)
    add_text_block(s, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.28),
                   ["修補在 longphase-to-mod (fork)  ·  本 repo (InterSubMod) 無 C++ 改動  ·  ISM 是下游受惠者"],
                   font_size=11, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 2)
    set_speaker_note(s,
        "邊界澄清。本張用 fig1_pipeline_comparison 展示三系統 (Tumor BAM → ClairS-TO → longphase-to-mod V5 → tagged BAM → InterSubMod) 的 pipeline 與本工作的位置。"
        "**核心訊息**: 修補位置只在中間 longphase-to-mod fork 的 HaplotagProcess.cpp, "
        "ClairS-TO caller 與 InterSubMod ISM 都無 C++ 改動。InterSubMod 是讀新版 BAM 後下游分析自動受惠。"
        "**Pipeline 各層職責**: "
        "(1) Tumor BAM — Nanopore long-read, 含 MM/ML 甲基化 tag。"
        "(2) ClairS-TO — caller, 輸出 tumor-only somatic SNV/INDEL VCF, V5 完全不影響此層 raw F1=0.7166 三版本相同。"
        "(3) longphase-to-mod V5 — phasing + haplotag, V5 修補位置, getVote() / countSNPHaplotype() / countINDELHaplotype() 三函式 +68/-36 行。"
        "(4) tumor_tagged.bam — 輸出 BAM, 每條 read 一個 HP:i 整數 tag (0/1/2/11/21/33)。"
        "(5) InterSubMod — read-level ISM 分析, 用 HP tag 做 4-bucket 分群 (HP1/HP1-1/HP2/HP2-1) 計算特徵。"
        "**為何 ISM 不直接修?**: 違反單一職責, phasing/haplotag 該在 longphase-to-mod (fork), ISM 重複實作會介面割裂; "
        "ISM 對 TO germline FP 的 F1 增益只 0.0124, 無法取代上游修補。修補在 fork 也讓 V5 BAM 可被其他工具直接消費。"
        "**4 階段工作對應 PPT 段落**: (1) 機制定位 (S3-S5 觀察 17.3:1 與根因) → (2) 4-commit 修補 (S6-S7 根因樹與 V5 三層投票) → "
        "(3) 驗證 (S8-S9 sanity + 量化指標) → (4) 影響評估 (S10-S11 climax 與五大目標銜接)。"
        "**Q&A 預備**: "
        "Q (「為什麼用 fork 而不是上游 PR?」) — 上游 LongPhase 不支援 TO + PON 模式, longphase-to-mod 是同實驗室相鄰工作的 fork, "
        "本實作填補 TO + PON 場景 gap, 後續可考慮回饋上游。"
        "**過場**: 下一張 S3 補完 HP tag 的 5 種整數值定義, 並用三層證據展示 self-phasing 的全貌。")
    return s


def slide_03_definition_evidence(prs):
    """S3: HP tag 5 整數值 + 三層證據 (fig17 + fig01d 並排)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "HP tag 五值 + 三層證據")
    # 左半 fig17
    fit_image_within(s, str(FIG_DIR / "fig17_hp_tag_5versions.png"),
                     Inches(0.3), Inches(1.0), Inches(6.4), Inches(5.7), border=True)
    add_text_block(s, Inches(0.3), Inches(6.75), Inches(6.4), Inches(0.30),
                   ["HP tag 5 整數值定義  (HP:i:0 / 1 / 2 / 11 / 21 / 33)"],
                   font_size=11, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    # 右半 fig01d
    fit_image_within(s, str(FIG_DIR / "fig01d_somatic_bias_explanation.png"),
                     Inches(6.85), Inches(1.0), Inches(6.2), Inches(5.7), border=True)
    add_text_block(s, Inches(6.85), Inches(6.75), Inches(6.2), Inches(0.30),
                   ["全基因組 17.3 : 1 證據  (HP1=614K vs HP2=35.5K reads)"],
                   font_size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 3)
    set_speaker_note(s,
        "兩張圖並排 — 左定義, 右證據。讓教授 30 秒內理解「HP tag 是什麼」與「為什麼 17.3:1 是 artifact」。"
        "**左半 (fig17 HP tag 5 整數值)**: BAM 內每條 read 一個整數 HP tag, 6 種可能值 — "
        "HP:i:0 (unphased)、HP:i:1 (germline HP1)、HP:i:2 (germline HP2)、"
        "HP:i:11 (somatic on HP1)、HP:i:21 (somatic on HP2)、HP:i:33 (somatic ambiguous)。"
        "這個編碼是 longphase-to-mod 自訂, 不在 LongPhase 官方 spec, 與 C++ enum (Util.h:20-25) "
        "HAPLOTYPE_UNDEFINED=-1 / HAPLOTYPE1=1 / HAPLOTYPE2=2 / HAPLOTYPE1_1=3 / HAPLOTYPE2_1=4 / HAPLOTYPE3=5 不同 — "
        "這個型別語意失配是後面 S6 enum bug (Bug 2) 的伏筆。"
        "**右半 (fig01d 全基因組 17.3:1 證據)**: HCC1395 5kHz baseline 硬數據 — "
        "HP1 family (HP1 + HP1_1 reads) = 614,000, HP2 family = 35,500, 比例 17.3:1; 94.6% 集中於 HP1; 跨 23 染色體一致。"
        "fig01d Panel D 的 6 個 sub-panel 證明這是全基因組層的系統性偏差, 不是個別 hotspot 也不是統計噪音。"
        "**三層證據彙整 (融合在 speaker note)**: "
        "(1) 理論層 — germline het 隨機 → 預期 1:1; "
        "(2) 全基因組層 (右圖) — 17.3:1 跨 23 chr; "
        "(3) 個別位點層 — SP1 chr19:17565944 baseline HP2:HP1 = 113:0 完全失衡, paired 真實方向是 HP2 (110:2)。"
        "三層獨立證據彙整於同一結論: self-phasing artifact 真實存在。"
        "**LOH 兩層精確化 (放這裡避免後面誤解)**: ISM HP_Ratio LOH (read-level, BAM HP tag 來源, 受 self-phasing 影響 62% 是 artifact) "
        "與 LOH.bed (region-level, VCF AD 來源, 不受影響, V5 前後 Jaccard = 1.0) 是兩套不同數據源的系統; "
        "kappa = 0.670 不是 bug 而是兩獨立衡量的自然差異。"
        "**Q&A 預備**: "
        "Q (「LOH.bed Jaccard=1.0 怎麼證明?」) — baseline 與 V2b PON-only LOH.bed 的 region-level Jaccard 相似度 = 1.0, "
        "兩 BED 完全相同, 證實 phasing 層不受 self-phasing 影響。"
        "**過場**: 下一張 S4 把 phasing 層 vs tag 層拆鎖 — phasing 沒問題, 是 tag 的問題。")
    return s


def slide_04_phasing_vs_tag(prs):
    """S4: 拆鎖 phasing OK vs tag FAIL."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Phasing OK,  Tag 有問題")
    img = make_v3_s4_phasing_vs_tag()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.8))
    add_text_block(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.3),
                   ["LOH.bed 由 VCF AD  ·  HP_Ratio 由 BAM HP tag  ·  本 PPT 焦點 = tag 層"],
                   font_size=11, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 4)
    set_speaker_note(s,
        "拆鎖關鍵 — self-phasing 影響的是哪一層? 答案是 tag 層, 不是 phasing 層。"
        "**左欄 phasing 層 (綠色 PASS)**: "
        "(1) 路徑 — VCF allele depth (AD) → LongPhase region 偵測, 輸出 BED 格式 region-level LOH。"
        "(2) Jaccard = 1.0 — V5 前後 LOH.bed 完全不變 (baseline vs V2b PON-only LOH.bed Jaccard = 1.0000)。"
        "(3) r = 0.001 — 同位點 TO vs paired HP_Ratio 相關係數, n=288K pairs, 幾乎完全無相關。"
        "(4) Cohen's d = -1.20 — 巨大效應量, paired 與 TO HP_Ratio 跨樣本差異。"
        "(5) V2b PON-only 已解 phase scaffold self-phasing (commit 8b8c1fd)。"
        "**右欄 tag 層 (紅色 FAIL)**: "
        "(1) 路徑 — BAM HP:i tag → ISM 統計 HP_Ratio per region。"
        "(2) HP_Ratio 17.3:1 — baseline 嚴重偏向 HP1 (61.4 萬 vs 3.55 萬 reads)。"
        "(3) 62% LOH = artifact — ISM HP_Ratio LOH 在 paired 比對下其實是 artifact, "
        "TO TP 中 86.5% 在 paired 下平衡 (HP_Ratio 0.4-0.6)。"
        "(4) kappa = 0.67 — ISM HP_Ratio LOH 與 LOH.bed 兩套系統不一致, "
        "由「不同數據源」自然差異解釋, 不是 bug。"
        "(5) AF 0.1-0.8 中等 AF 範圍近 100% 假 LOH, 是 self-phasing 最嚴重的區段。"
        "**重要區分 (避免誤解)**: 兩套 LOH 系統使用不同數據源 (BAM HP tag vs VCF AD), 過去把 ISM HP_Ratio LOH 當 LOH.bed 真值是常見誤解; "
        "本工作 self-phasing 修補只影響 ISM HP_Ratio LOH (左欄 read-level), 不改變 LOH.bed (右欄 region-level)。"
        "**Q&A 預備**: "
        "Q (「62% 是怎麼算的?」) — TO mode ISM-level 標為 LOH 的位點 (HP_Ratio < 0.1 或 > 0.9), 在同樣本 paired BAM 上重算 HP_Ratio, "
        "62% 落在 0.4-0.6 平衡區間 → baseline TO 的 LOH 標記 62% 是 self-phasing artifact。"
        "**過場**: 下一張 S5 量化影響 — 哪些 ISM 特徵嚴重 / 中度 / 不受影響, 重跑工作可聚焦在 43 個而非全 85。")
    return s


def slide_05_ism_impact(prs):
    """S5: 29/85 features 受影響 (03_self_phasing_impact + 註解)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "ISM 影響範圍  29 / 14 / 42")
    fit_image_within(s, str(FIG_DIR / "03_self_phasing_impact.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.35),
                   ["嚴重 29 個 (HP-依賴)  ·  中度 14 個 (間接污染)  ·  不受影響 42 個  ·  總計 85 features"],
                   font_size=12, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.30),
                   ["ISM 4-bucket 分群 (HP1 / HP1-1 / HP2 / HP2-1) 全部依賴 HP tag 正確"],
                   font_size=10, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 5)
    set_speaker_note(s,
        "影響量化。本張用 03_self_phasing_impact 3-tier 金字塔展示 V5 修補對下游 ISM 特徵的影響範圍 (來源 source 03 §4.3)。"
        "**嚴重影響 (HP-依賴) — 29 個**: "
        "代表特徵: HP_Ratio (62% 假 LOH)、Potential_LOH、HPMergedDelta / HPMergedSig (方向反轉)、"
        "HPFineNGroups (含 NG=2 LOH-constrained phasing)、HPSig、HP_Ratio_norm。"
        "這 29 個必須在 V5 BAM 上重跑才能信任, 舊結論需加註版本 (haplotag_version = baseline / V5)。"
        "**中度影響 (間接污染) — 14 個**: "
        "代表特徵: QualityScore (AUC 0.497, 已移除 LOH penalty)、GlobalP (取 HP/Allele 最小值)、"
        "CramersV (取 HP/Allele 最大值)、VerificationClass (label_sig 含 HP 成分)。"
        "重評多數影響微弱, 但仍建議在 V5 BAM 上 cross-check。"
        "**不受影響 (無 HP 依賴) — 42 個**: "
        "代表特徵: PairwiseMean / MedianDist (無 HP 分組)、AlleleDelta / AlleleP、Caller 特徵 (AF / GQ / DP / SB)、"
        "甲基化矩陣本身 (BAM MM/ML tag, 不依賴 HP)、CpG 座標、region_methyl_mean。"
        "這 42 個結論穩固不需重測。"
        "**為什麼 ISM 4-bucket 分群依賴 HP tag**: "
        "ISM 把 reads 分成 HP1 (germline HP1)、HP1-1 (somatic on HP1)、HP2 (germline HP2)、HP2-1 (somatic on HP2) 四個 bucket, "
        "計算 region-level 特徵。如果 HP tag 偏差 (17.3:1), 4-bucket 分群錯誤, 所有依賴特徵的結論都不可信。"
        "**口徑校準**: slide 上只列 count 不列 %。來源報告寫 14/85 = 7% 但 14/85 實際 16.5%, 來源百分比有誤。"
        "為避免傳遞錯誤數字, 本 slide 全部以 count 表達。"
        "**Q&A 預備**: "
        "Q (「重跑 29 + 14 = 43 個要多久?」) — 7 樣本 V5 BAM 全量重跑預估 ~10 hr parallel (依 archive TO rerun 規劃), "
        "聚焦 43 個 features 而非全 85, 大量節省計算資源。"
        "**過場**: 下一張 S6 進入根因 — 為什麼 baseline 會出現 17.3:1 self-phasing? 答案是觸發條件 + 三 bug 的集中暴露。")
    return s


def slide_06_root_cause(prs):
    """S6: ★ 根因樹 (Purity 0.95 主者 + 三 bug)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "★ 根因樹  —  觸發條件 + 三 bug")
    img = make_v3_s6_root_cause_tree()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(6.0))
    add_footer(s, 6)
    set_speaker_note(s,
        "★ 核心新發現 — Self-phasing 不是單一 bug, 而是「特定觸發條件下三層獨立故障的集中暴露」。"
        "**根 (觸發條件主者, 紅色框)**: Baseline 估 purity = 0.927 ≤ 0.95 (純樣本標準) → "
        "未觸發 Two-Pass → 走 baseline 標準三條路流程。"
        "**口徑校準關鍵 (避免誤解)**: 不是「baseline 誤判 purity」(0.927 是正確估計, 是邊緣 case), "
        "而是 baseline 設計上的硬編碼閾值 0.95 (PhasingProcess.cpp:197) 對近似純樣本不夠保守。"
        "HCC1395 5kHz 真實 purity > 0.95 (純 tumor, 無 normal contamination), 但 baseline 的 polynomial calculator 估出 0.927, "
        "這個 0.927 數值本身正確, 只是落在 0.95 閾值的邊緣下方就觸發了 baseline 標準流程而非 Two-Pass 路徑。"
        "**Bug 1 (紅色, getVote priority)**: variantKeys 順序 {HP1_1, HP2_1} 排第一, 任一 somatic count > 0 即 break, "
        "germline 完全跳過。為何 paired 沒事? paired germline votes 多但 priority bug 仍存在; "
        "TO PON-only 後 germline votes 急減立刻顯形 (HCC1395 99.9% reads 拿 HP21)。V3F commit 41ff147 修。"
        "**Bug 2 (橘色, enum vs int literal mismatch)**: C++ enum 定義在 Util.h:20-25 (注意是 20-25 不是 21-25, "
        "HAPLOTYPE_UNDEFINED=-1 起於 line 20), HAPLOTYPE1_1 = 3。但 BAM HP tag integer = 11。"
        "caller 端 if(hpResult != HAPLOTYPE1_1) 用 enum=3 比較 hpResult=11/21/33, fallback 死分支永不執行, "
        "HP:i:33 永不出現 (baseline 0/15 sites)。C++ 把 enum 隱式轉 int 不會 warn。V3F 同 commit 改為 integer literal 比較。"
        "**Bug 3 (灰色, scaffold, 已由 V2b 解)**: PhasingProcess.cpp:154-157 convertNonGermlineToSomatic() 不被觸發, "
        "somatic 當 phasing anchor。已由 V2b commit 8b8c1fd 解決 (透過 --pon-only-phasing flag), "
        "本 PPT 焦點不在此, 視覺只佔 20%。"
        "**整體因果鏈**: 因樣本實為純無 normal contamination 而流程假設有 → 暴露 tag 層 somatic-first 投票 bias → "
        "結果 self-phasing 17.3:1 artifact。「修一個 bug 暴露另兩個 bug」: PON-only 啟用後 germline votes 急減, "
        "原本被 paired germline 訊號掩蓋的 tag 層 priority bug + enum bug 立刻顯形。"
        "**驗證來源**: source_materials/04_purity_calculator_failure_root_cause.md (v5_audit_suite/18 複製)。"
        "**Q&A 預備**: "
        "Q (「PON 有資料還會 self-phasing?」) — PON classification (VCF 層) 與 phasing anchor (read 連結層) 是獨立步驟。"
        "Baseline bug: PON 雖標出 somatic 但 phasing 階段仍把 somatic 當 anchor (--pon-only-phasing=false 預設)。"
        "V2b 啟用後 phasing 才只用 PON-confirmed germline 做 anchor。"
        "**過場**: 下一張 S7 講 V5 三層投票流程如何修補 Bug 1 與補 Layer 1.5 mid-low purity 防禦。")
    return s


def slide_07_v5_voting(prs):
    """S7: ★ V5 三層投票流程."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "★ V5 修法  —  三層投票流程")
    img = make_v3_s7_three_layer_voting()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(6.0))
    add_footer(s, 7)
    set_speaker_note(s,
        "★ 修法核心。V5 對 getVote() 的具體改動: 從 V3F 兩層升級為三層投票邏輯。"
        "**Layer 1 (germline first, 藍色)**: "
        "if (germlineHP1 > germlineHP2) → 1; elif (germlineHP2 > germlineHP1) → 2; else → 進入 Layer 1.5。"
        "**這裡是 baseline 的 somatic-first 投票順序的關鍵翻轉**: baseline getVote() 用 variantKeys loop 把 {HP1_1, HP2_1} 排第一, "
        "somatic 投票優先; V5 直接看 germline 投票數, 翻轉投票優先序為 germline-first。"
        "保持 germline first 核心原則 — germline 是穩定的 zygosity 訊號, somatic 是後天突變不應主導 phasing。"
        "**Layer 1.5 (somatic fallback, V5 新增, 紅色 mid-low purity 防禦層)**: 當 germline 平手時參考 somatic 投票。"
        "if (somaticTotal == 0) → hpResult = 0 (unphased); "
        "elif (somaticHP1 > 0.6 × somaticTotal) → 對齊 HP1 方向 (germlineResult = 1); "
        "elif (somaticHP2 > 0.6 × somaticTotal) → 對齊 HP2 方向; "
        "else → hpResult = 33 (真正 ambiguous)。"
        "Confidence threshold 0.6 是經驗值 (R4 caveat), **設計上是 mid-low purity 的防禦層** — "
        "在 0.6-0.93 中等純度範圍, 弱 directional 訊號標 33 ambiguous (保守 +10pp HP33 比例) 避免錯誤分配。"
        "**Layer 2 (encode hpResult, 綠色)**: germline=1 + somatic_total>0 → HP:i:11; germline=2 + somatic_total>0 → HP:i:21; "
        "germline=0 + somatic_total>0 → HP:i:33; germline only → HP:i:1 or 2; nothing → HP:i:0。"
        "對應 5 種 BAM HP integer 編碼 (見 S3)。"
        "**效果量化 (底部 strip)**: "
        "AMB% (HP:i:33 比例) 從 17.5% 降到 8.0% (-9.5 pp), V5 解了過半 ambiguous reads; "
        "HP:i:33 reads 數從 239,679 降到 110,197, 減 54%; "
        "介面契約 HaplotagProcess.h:66-68 一字未變 (drop-in replace), 上下游無需配合修改。"
        "**Q&A 預備**: "
        "Q (「V5 vs baseline 在 mid-low purity 表現?」) — V5 更好。"
        "純樣本 (HCC1395 5kHz): baseline somatic-first 在 PON-only 啟用後 germline votes 急減 → somatic 主導 → 17.3:1 artifact; "
        "V5 germline-first 不受此影響。"
        "Mid-low purity 場景 (0.6 sample, t30_n20): baseline 仍 somatic-first 但 self-phasing 自然弱化; "
        "V5 加入 confidence 0.6 threshold → 對「弱 directional 訊號」標 HP33 ambiguous → "
        "09 報告: 0.6 sample HP33 比例 12.4% vs Baseline 2% (+10pp 保守) → 避免錯誤分配。"
        "**Caveat**: V5 在 SP-extreme self-phasing 位點仍不修 (V5 設計就不修 phase 層, 那是 V2b PON-only 的責任)。"
        "**過場**: 下一張 S8 用 sanity 4 項 + 5 證據鏈展示驗證強度。")
    return s


def slide_08_sanity(prs):
    """S8: Sanity 15/15 + 5 證據鏈."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Sanity 15 / 15 PASS  +  5 證據鏈")
    img = make_v3_s8_sanity_card()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(5.9))
    add_text_block(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.30),
                   ["守恆律 (in = out)  +  5 條獨立路徑同步收斂  ·  self-phasing 修補必要且充分"],
                   font_size=11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_footer(s, 8)
    set_speaker_note(s,
        "驗證強度。本張用守恆律公式卡 + 5 證據鏈展示 V5 修補通過嚴格驗證。"
        "**4 項硬性 sanity check (HCC1395 5kHz 15 sites, v5_audit_suite/06 報告)**: "
        "(1) 守恆律 A · Δ-consistency: ΔHP33 + (ΔHP11 + ΔHP21) = 0, tag 移轉量平衡, in = out, 15/15 PASS。"
        "意義: V5 reassign 39 reads 從 HP33 到 HP11, ΔHP33 = -39, ΔHP11 = +39, 總和為 0; "
        "沒有 read 憑空消失或新增, 所有 tag 變動都是內部 reassignment。"
        "(2) 守恆律 B · Germline 不變: V3F 與 V5 的 HP1 / HP2 reads 逐 site 比對, 0 reads 漂移, 15/15 PASS。"
        "意義: V5 的 Layer 1.5 只動 germline 平手的 reads (ambiguous bucket), 不影響 germline 已決定的 reads。"
        "(3) Layer 1.5 期望 1 · 33→directional 精確守恆: ΔHP11 == n(V3F=33→V5=11)、ΔHP21 == n(V3F=33→V5=21), "
        "15/15 PASS (V5max1 chr19:4639528 39 reads 精確守恆)。"
        "(4) Layer 1.5 期望 2 · 無 germline → HP33: 跨 15 sites pool transition table, germlineResult ≠ 0 不出 HP33, 0 violation。"
        "**5 條獨立證據鏈 (下半 5 卡)**: "
        "(1) 理論層 — germline het 隨機 → 預期 1:1 (生物學第一性原理); "
        "(2) 全基因組層 — HP1:HP2 = 17.3:1 跨 23 chr 一致 (硬數據); "
        "(3) 個別位點層 — SP1 chr19:17565944 baseline 113:0 全失衡, paired 確認 HP2 為真; "
        "(4) Sanity check 4 項 15/15 PASS (左欄, 守恆律 + Layer 1.5 預期); "
        "(5) 程式碼最小必要 — +68/-36 行集中於 3 函式 (getVote / countSNPHaplotype / countINDELHaplotype), "
        "介面契約零變動, 改動範圍可審計。"
        "5 條獨立路徑 (理論 / 觀察 / 個案 / 守恆律 / 程式碼) 同步收斂於同一結論: self-phasing 修補必要且充分。"
        "**Q&A 預備**: "
        "Q (「sanity check 是誰寫的? 覆蓋率多少?」) — 本工作新增 (v5_audit_suite/06 Agent D), "
        "覆蓋限制是 15 sites cherry-picked (R3), 需 7 樣本擴展 (F3) + 100 隨機位點 cross-validate (F8); "
        "整體 sanity check 必要不充分但已是業界水準以上。"
        "Q (「PON-only Two-Pass 在不同 purity 都有效嗎?」) — "
        "0.93 純樣本: PASS, Aggregate paired GT concordance V5 78.85% vs Baseline 72.20% (+6.65pp); "
        "0.6 sample: PASS, V5 HP33 比例 12.4% vs Baseline 2% (+10pp 保守 tagging); 兩端都有改善。"
        "**過場**: 下一張 S9 用 fig18 量化指標可視化呈現 AMB / HP33 / Concordance 三大改善。")
    return s


def slide_09_metrics(prs):
    """S9: 量化指標 fig18 + 數字大字註記."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "量化指標  —  AMB ↓  HP33 ↓  Concordance ↑")
    # 上方數字大字
    add_text_block(s, Inches(0.4), Inches(1.0), Inches(4.1), Inches(0.7),
                   ["−9.5 pp"], font_size=42, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.4), Inches(1.65), Inches(4.1), Inches(0.30),
                   ["AMB%  17.5 → 8.0"], font_size=12, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

    add_text_block(s, Inches(4.6), Inches(1.0), Inches(4.1), Inches(0.7),
                   ["−54 %"], font_size=42, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(4.6), Inches(1.65), Inches(4.1), Inches(0.30),
                   ["HP:i:33 reads  239K → 110K"], font_size=12, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

    add_text_block(s, Inches(8.8), Inches(1.0), Inches(4.1), Inches(0.7),
                   ["+8.3 pp"], font_size=42, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(8.8), Inches(1.65), Inches(4.1), Inches(0.30),
                   ["Concordance  82.2 → 90.5"], font_size=12, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

    # fig18 主圖佔 80%
    fit_image_within(s, str(FIG_DIR / "fig18_concordance_amb_f1.png"),
                     Inches(0.4), Inches(2.1), Inches(12.5), Inches(4.6))
    add_text_block(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.30),
                   ["clean PS blocks, 全基因組  ·  非全基因組 raw  ·  非 cherry-picked"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 9)
    set_speaker_note(s,
        "量化結果。本張上方三個大字指標 + 主圖 fig18 並列, 一眼看出三項改善。"
        "**指標 1: AMB% (HP:i:33 比例) 17.5% → 8.0% (-9.5 pp)**: "
        "ambiguous reads 比例減半。意義是 V5 解了過半「baseline 與 V3F 都沒辦法決定 hap 方向」的 reads, "
        "Layer 1.5 confidence 0.6 threshold 把這些 reads 透過 somatic directional 訊號重新分配到具體 hap。"
        "**指標 2: HP:i:33 reads 239,679 → 110,197 (-54%)**: "
        "原始 read 數降 54%。HCC1395 5kHz 全基因組 baseline 有 23.97 萬條 reads 標為 ambiguous, V5 後降到 11.02 萬, 改善幅度 12.97 萬 reads。"
        "**指標 3: 全基因組 clean PS blocks paired GT concordance V5 90.5% vs Baseline 82.2% (+8.3 pp)**: "
        "**口徑校準關鍵**: +8.3 pp 必須加 caveat 「(clean PS blocks, 全基因組)」, 不是全基因組 raw, 不是 cherry-picked。"
        "clean PS blocks 篩選 germline accuracy ≥ 70% 的 phase block; problem PS blocks 上 V5 與 Baseline 接近隨機 "
        "(read-level orientation 不穩定 — 那是 phase 層 V2b 的責任, 不是 V5 修補對象)。"
        "**其他補充指標**: "
        "(a) 15-site cherry-picked 上 clean PS concordance: V5 = 88.2% vs BL = 74.9%, +13.3 pp。"
        "(b) 15-site aggregate (all 15 pooled): V5 78.85% vs BL 72.20%, +6.65 pp。"
        "(c) TP Balanced% 從 V2b 13.0% 改善 (V5 同 V3F 基準)。"
        "**注意**: 這些都是 read-level metric, 不是 caller F1 metric — F1 不變是預期 (V5 不改 caller, 改 BAM HP tag 編碼)。"
        "ClairS-TO raw F1 = 0.7166 對所有版本 (Baseline / V3F / V5) 完全相同。"
        "**Q&A 預備**: "
        "Q (「N 多少? 信賴區間? 跨樣本一致?」) — 15-site cherry-picked N=11 clean PS; "
        "全基因組 clean PS N 在數萬 - 數十萬 reads 量級 (PI 報告 4 §3.7); "
        "跨樣本擴展未做 (F3 待辦, 7 樣本 V5 BAM 全量重跑); 信賴區間 source 03 未直接給, 推估 Wilson 95% CI ±1 pp 內。"
        "Q (「為什麼 F1 不變仍要做?」) — 真實價值在 read-level concordance +8.3 pp, ISM 29 個 HP-依賴特徵必須在 V5 BAM 才可信; "
        "跨樣本研究可信度 (7/7 同方向) 是論文 contribution。"
        "**過場**: 下一張 S10 climax — 用 V5max1 chr19:4639528 的 IGV 大圖展示 39 reads 100% reassigned 的視覺證據。")
    return s


def slide_10_climax_v5max1(prs):
    """S10: 🎯 V5max1 climax — 39 reads 翻轉 (用 ★★★ 替代 emoji)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "★★★  V5max1  —  39 reads 翻轉")
    fit_image_within(s, str(FIG_DIR / "C_V5max1_chr19_4639528.png"),
                     Inches(0.6), Inches(1.0), Inches(12.0), Inches(5.4))
    # Caveat strip
    add_rect(s, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.55),
             fill=RGBColor(0xFE, 0xF5, 0xE7), line=C_AMBER, line_w_pt=1.5)
    add_text_block(s, Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.30),
                   ["chr19:4639528  ·  V3-F 紫色 HP33 群 (39 reads) → V5 全部正確重分配為 HP11"],
                   font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.22),
                   ["但 V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold)  ·  3/3 SP-extreme 一致翻轉"],
                   font_size=10, color=C_AMBER, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 10)
    set_speaker_note(s,
        "演講 CLIMAX。chr19:4639528 V5max1 是最戲劇化的視覺證據, 用 IGV 大圖佔 80% 寬度展示。"
        "**視覺敘事**: 在 V3-Fixed panel 看到一大群紫色 reads 標為 HP:i:33 (ambiguous), 共 39 條; "
        "V5 修補後這 39 條全部正確重分配為 HP:i:11 (somatic on HP1)。"
        "守恆律驗證: V3F 39 reads HP33 → V5 39 reads HP11, in = out 完全平衡, sanity check (前一張 S8 的公式 A) PASS。"
        "這是 Layer 1.5 somatic fallback 的具體效果 — 原本因為 V3F 過於保守卡在 ambiguous bucket 的 reads, "
        "現在透過 confidence 0.6 threshold 被正確歸到具體 hap (somatic 投票 39/39 都指向 HP1, 符合 0.6 threshold 即 align HP1)。"
        "看 IGV 4-BAM 並列圖 (從上到下: baseline / V2b / V3F / V5), 從 baseline → V2b → V3F → V5 的視覺差異很明顯。"
        "**底部 caveat 防誤解 (重要!)**: 「但 V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold)」。"
        "**口徑校準關鍵**: 高潮圖容易讓教授誤以為 V5 全面解決 self-phasing, 但 self-phasing 問題鏈是被分層處理的: "
        "(1) phase 層的 self-phasing scaffold 由 V2b PON-only phasing 處理 (改 phasing graph anchor 來源, "
        "從含 somatic 改為僅 PON-confirmed germline); "
        "(2) V3F 解 priority bug + enum mismatch (Bug 1+2); "
        "(3) V5 解 V3F 過於保守 (germline 平手時 fallback 到 33), "
        "新增 Layer 1.5 directional reassignment 把卡在 33 的 reads 解到 11/21。"
        "所以 V5 的「39 reads → 100% reassigned」是 directional reassignment 不是 phasing graph 重建。"
        "**3/3 SP-extreme 一致翻轉證據**: 除 V5max1 (chr19:4639528 39 reads), 還有 V5max2 (chr19:2235521 26 reads) "
        "與 V5max3 (chr19:7405500 16 reads), 全部 100% reassigned, 跨位點一致性。"
        "另外 SP1/SP2/SP3 三個 SP-extreme 位點 (baseline 113:0 / 109:1 / 108:0) V2b 已翻轉至 paired 方向, V5 在這些位點 Δ ≈ 0 (V5 不修 phase 層)。"
        "speaker 必須在 60-90 sec 內一句說清這個分層, 避免教授誤以為「V5 = 解所有」。"
        "**Q&A 預備**: "
        "Q (「為什麼這 39 條 baseline 標 HP33 而 V5 知道是 HP11?」) — "
        "baseline 因為 enum=3 vs HP integer=11 mismatch, fallback 死分支永不執行, HP33 永不出現 (39 條應分類但分類失敗); "
        "V3F 修了 enum 但 priority 邏輯讓這些 reads 卡在 ambiguous; "
        "V5 加 Layer 1.5 看 somatic 投票方向 (39 條 somatic 全指向 HP1, 比例 100% > 0.6 threshold), 確認後重 align 為 HP11。"
        "**過場**: 下一張 S11 把本工作放回業界家族樹, 並銜接到 InterSubMod 五大目標。")
    return s


def slide_11_industry_goals(prs):
    """S11: 業界家族樹 + 五大目標卡."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "業界家族樹  +  五大目標銜接")
    img = make_v3_s11_family_and_goals()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(6.0))
    add_footer(s, 11)
    set_speaker_note(s,
        "銜接段。本張上半業界家族樹 + 下半五大目標卡, 把本工作放回研究脈絡。"
        "**上半 業界家族樹 (同實驗室相鄰工作)**: "
        "上游基底是 LongPhase (Lin 2022 Bioinformatics) 處理 germline phasing。由此分出兩個分支: "
        "(a) LongPhase-S (bioRxiv 2025.11.20.689492v1) 是同實驗室的 paired 版, 在 ClairS 上 SNV F1 +4.5% (indel +7.1%), "
        "提供「somatic 錨在 germline scaffold」的 anchoring 概念。"
        "(b) longphase-to-mod V5 是本實作, 是 TO + PON 條件下的本地實作, +8.3 pp clean PS concordance。"
        "**口徑校準關鍵**: 本實作填補 tumor-only 在公開工具中的 gap (WhatsHap 不支援 TO), "
        "**不是「業界共識」「標準替代」**, 而是「同實驗室相鄰工作」。"
        "LongPhase-S 為 paired 場景參考 (注意非直接可比, 因本工作 F1 不變是預期行為), 不重疊本工作 TO 場景。"
        "**Contribution 區隔**: 4-commit 漸進修補在 longphase-to-mod fork (獨立 git repo); "
        "修補對象是 LongPhase-TO 在 PON-only 啟用後集中暴露的 3 層 tag-side bug; "
        "4 項 sanity check 是新貢獻; ISM 下游影響 3-tier 分類是新貢獻; "
        "ISM 跨樣本一致性 7/7 + Cohen's d=-1.20 是新貢獻。"
        "**下半 五大目標銜接**: "
        "InterSubMod 五大研究目標 (來源 InterSubMod/docs/architecture/20260327_InterSubMod研究願景定錨_01.md): "
        "目標 1 — per-CpG 甲基位點多標籤關聯性評分 (每個 CpG 與 ALT/REF/HP1/HP2 關聯); "
        "目標 2 — clone 結構分析 (sub-clone + 共演化); "
        "目標 3 — 二次打擊與事件順序推論 (依目標 1/2 + LOH/CNV/HP); "
        "目標 4 — TO normal 資訊補強 (無配對 normal 下估計 normal 背景); "
        "目標 5 — 整合 evidence panel 提升 F1 (補充 caller, 保留 TP, 過濾 FP)。"
        "**目標 1/2/4 直接依賴 (粗框星號)**: per-CpG 關聯需 hap 分層、clone 結構需 4-bucket、TO normal 補強需 phasing 正確。"
        "**目標 3/5 部分依賴 (細框)**: 3 間接依賴 1/2; 5 部分依賴 + 整合 evidence。"
        "Self-phasing 修補是 5 大目標解鎖前提, 特別是目標 1/2/4 (V5 修補完才能信任分析結果)。"
        "本張展示研究藍圖, 讓教授看到本工作不只是一次性修補, 而是支撐整個 InterSubMod 的根基。"
        "**Q&A 預備**: "
        "Q (「跟 LongPhase-S 重疊多少? 算 contribution 嗎?」) — LongPhase-S 是 paired, 本實作是 tumor-only (PON 替代 normal); "
        "設計哲學一致, 本工作填補 TO gap; 4 項新 sanity check + ISM 3-tier 影響分類 + 7/7 跨樣本一致性都是新貢獻。"
        "Q (「Thread D NG=2 6 樣本可信嗎?」) — Wilcoxon 6 樣本配對最強顯著性是 6/6 同方向 W=21 p=0.0156, 本實驗達此上限; "
        "evidence grade B (待 Phase 2B 升 A); 多重比較校正未做 (只 1 假設)。"
        "**過場**: 最後一張 S12 take-home + 3 P0 行動 + Q&A 結語。")
    return s


def slide_12_takehome_qa(prs):
    """S12: Take-home + 3 P0 + Q&A."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))

    # 上半: take-home (大字)
    add_rect(s, 0, 0, SLIDE_W, Inches(2.7), fill=C_TITLE_BG)
    add_text_block(s, Inches(0.5), Inches(0.30), Inches(12.5), Inches(0.45),
                   ["Take-home"],
                   font_size=18, bold=True, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(0.85), Inches(12.5), Inches(0.7),
                   ["TO 模式可用  ·  Tag 層問題已解  ·  五大目標解鎖"],
                   font_size=28, bold=True, color=C_TITLE_FG,
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(1.75), Inches(12.5), Inches(0.45),
                   ["read-level concordance +8.3 pp (clean PS, 全基因組)"],
                   font_size=16, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(2.20), Inches(12.5), Inches(0.40),
                   ["V5 不修 self-phasing 本身 (V2b 已處理)  ·  分層處理 phase 與 tag 兩層"],
                   font_size=11, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER, italic=True)

    # 中段: 3 P0 卡片
    add_text_block(s, Inches(0.5), Inches(2.95), Inches(12.5), Inches(0.40),
                   ["3 P0 行動  (下次 meeting 教授會看到的進度)"],
                   font_size=15, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    items = [
        ("F1", "Commit V5 working tree", "切 2 獨立 commits + tag v5.0", C_ACCENT),
        ("F3", "7 樣本 V5 BAM 全量重跑", "HCC1395 / DORADO / HCC1937 / HCC1954 / H1437 / H2009 / COLO829", C_AMBER),
        ("F4", "HPFineNGroups master × flag 重驗", "判定 phasing signature 機制本質", C_GREEN),
    ]
    cy = Inches(3.55)
    rh = Inches(0.95)
    for fid, head, sub, color in items:
        add_rect(s, Inches(0.6), cy, Inches(12.2), rh,
                 fill=C_LIGHT, line=color, line_w_pt=2.5)
        add_text_block(s, Inches(0.8), cy + Inches(0.20), Inches(0.9), Inches(0.6),
                       [fid], font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_block(s, Inches(1.9), cy + Inches(0.10), Inches(7.8), Inches(0.45),
                       [head], font_size=15, bold=True, color=C_TEXT)
        add_text_block(s, Inches(1.9), cy + Inches(0.55), Inches(10.2), Inches(0.40),
                       [sub], font_size=11, color=C_GRAY)
        cy += rh + Inches(0.10)

    # 底部: Q&A 感謝
    add_rect(s, Inches(0.5), Inches(6.85), Inches(12.4), Inches(0.55),
             fill=RGBColor(0xFD, 0xF6, 0xE3), line=C_TITLE_BG, line_w_pt=2.0)
    add_text_block(s, Inches(0.5), Inches(6.92), Inches(12.4), Inches(0.40),
                   ["Q & A 歡迎隨時打斷  ·  Thank you"],
                   font_size=16, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    set_speaker_note(s,
        "結語。倒過來收尾: take-home + 3 P0 行動承諾 + Q&A 感謝, 不重複 S1 TL;DR。"
        "**Take-home (上半大字)**: TO 模式可用、Tag 層問題已解、五大目標解鎖。"
        "三個重點: "
        "(1) TO 模式可用 — 透過 V5 修補, TO + PON 條件下 tag 品質接近 paired 水準, "
        "read-level concordance +8.3 pp (clean PS blocks, 全基因組)。"
        "(2) Tag 層問題已解 — V3F 解 priority bug + enum mismatch (Bug 1+2), V5 補 Layer 1.5 directional reassignment, "
        "AMB% 17.5%→8.0%, HP:i:33 reads -54%, 介面契約零變動。"
        "(3) 五大目標解鎖 — InterSubMod 目標 1/2/4 直接依賴 HP tag 正確, V5 修補是研究可信度前提。"
        "**口徑校準分層 caveat**: V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold); "
        "self-phasing 問題鏈是被分層處理的 (phase 層由 V2b PON-only, tag 層由 V3F/V5)。"
        "**3 P0 行動承諾 (下次 meeting 教授會看到的進度)**: "
        "**F1 (高) Commit V5 working tree**: 切 2 獨立 commits — Commit A getVote() Layer 1.5 somatic fallback (行 512-563); "
        "Commit B countSNPHaplotype() alt guard 對稱化 (行 489-494)。完成後 tag v5.0, "
        "在 InterSubMod manifest.yaml 加 haplotag_version = v5.0 讓下游可追溯。"
        "**F3 (高) 7 樣本 V5 BAM 全量重跑**: 目前只在 HCC1395 5kHz 完整驗證, 跨樣本擴展 R6 未做; "
        "樣本清單 (KDE-corrected): HCC1395_5kHz (已驗證)、HCC1395_DORADO、HCC1937、HCC1954、H1437、H2009、COLO829; "
        "預估 ~10 hr parallel (依 archive TO rerun 規劃)。"
        "**F4 (高) master × flag 重驗 HPFineNGroups**: 對應 S11 機制本質判定, "
        "確認 HPFineNGroups 是 phasing signature 還是 methylation bimodality; "
        "subclone marker ⭐4→⭐3 已降級因 HCC1395 TO ClairS-TO raw split 無法重現 master 89.1%。"
        "以承諾收尾比口號收尾更有力, 教授知道下次會看到具體進度。"
        "**F5-F8 中低優先 (口頭帶過)**: F5 Manifest 加 haplotag_version 欄位; F6 cnLOH 區獨立評估方案 (R1 caveat 應對); "
        "F7 Trio-phased 第二 ground truth (R9 應對); F8 100 隨機位點 cross-validate (R3 cherry-pick 風險應對)。"
        "**底部 Q&A 條**: 歡迎隨時打斷, 詳細數據在 v2/source_materials 三份報告 (01 IGV visual audit、02 V5 audit、"
        "03 longphase TO vs V5 技術報告) 與 v2/notes 6 份 onboarding 文件 (00 background、qa 11+v3 Q12 Q13、"
        "key metrics、code references、industry references、glossary) + 04 Purity calculator failure root cause。"
        "Thank you。")
    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_impact,                # S1 17.3:1 -> 1:1 banner
        slide_02_pipeline,              # S2 工作流程一覽
        slide_03_definition_evidence,   # S3 HP tag 五值 + 三層證據
        slide_04_phasing_vs_tag,        # S4 拆鎖 phasing OK vs tag FAIL
        slide_05_ism_impact,            # S5 ISM 影響 29/14/42
        slide_06_root_cause,            # S6 ★ 根因樹 (Purity 0.95 + 三 bug)
        slide_07_v5_voting,             # S7 ★ V5 三層投票流程
        slide_08_sanity,                # S8 Sanity 15/15 + 5 證據鏈
        slide_09_metrics,               # S9 量化指標 fig18
        slide_10_climax_v5max1,         # S10 V5max1 climax 39 reads
        slide_11_industry_goals,        # S11 業界家族樹 + 五大目標
        slide_12_takehome_qa,           # S12 Take-home + 3 P0 + Q&A
    ]

    assert len(builders) == TOTAL_SLIDES, f"Expected {TOTAL_SLIDES} slides, got {len(builders)}"

    for i, fn in enumerate(builders, 1):
        print(f"  [{i:2d}/{TOTAL_SLIDES}] building {fn.__name__} ...")
        fn(prs)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"\n[OK] PPTX saved to {OUTPUT_PPTX}")
    print(f"     Slides: {len(prs.slides)}")
    return OUTPUT_PPTX


if __name__ == "__main__":
    build()
