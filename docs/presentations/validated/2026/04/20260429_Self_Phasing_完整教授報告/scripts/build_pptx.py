#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_pptx.py — Self-Phasing 完整教授報告 (34 slides, 16:9)

Generates output/20260429_Self_Phasing_完整教授報告.pptx according to
00_storyboard.md.  Each slide carries a speaker note populated from
source_materials/03_longphase_TO_vs_V5_技術報告.md and the storyboard's
"speaker note 重點" column.

Design rules (see PPT_RENDERING_RULES.md):
    1. Per-character Latin + CJK font fallback via <a:rPr><a:ea/></a:rPr>.
    2. fit_image_within (equal ratio, never force both width+height).
    3. Speaker notes are mandatory and non-empty for every slide.
    4. Slide title <= 15 CJK chars; bullets <= 5 entries, <= 15 chars each.
    5. No emojis (Droid Sans Fallback lacks them); use ASCII / box symbols.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FIG_DIR = ROOT / "figures"
OUT_DIR = ROOT / "output"
AUTO_FIG_DIR = OUT_DIR / "auto_figs"
OUT_DIR.mkdir(exist_ok=True)
AUTO_FIG_DIR.mkdir(exist_ok=True)

OUTPUT_PPTX = OUT_DIR / "20260429_Self_Phasing_完整教授報告.pptx"

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------
CJK_FONT_PATH = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
LATIN_FONT = "Arial"
CJK_FONT_NAME = "Droid Sans Fallback"

if Path(CJK_FONT_PATH).exists():
    fm.fontManager.addfont(CJK_FONT_PATH)
    _font = fm.FontProperties(fname=CJK_FONT_PATH).get_name()
    plt.rcParams.update({
        "font.family": [_font, "DejaVu Sans"],
        "axes.unicode_minus": False,
    })
else:
    print(f"[WARN] CJK font not found at {CJK_FONT_PATH}; CJK in matplotlib may render as boxes.")

# ---------------------------------------------------------------------------
# Slide geometry (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.4)
MARGIN_Y = Inches(0.35)

# Color palette
C_TITLE_BG = RGBColor(0x1F, 0x4E, 0x79)   # deep blue
C_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
C_ACCENT = RGBColor(0xC0, 0x39, 0x2B)     # red accent
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_AMBER = RGBColor(0xE6, 0x7E, 0x22)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
C_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_BORDER = RGBColor(0xBD, 0xC3, 0xC7)
C_BLUE = RGBColor(0x29, 0x80, 0xB9)

# ---------------------------------------------------------------------------
# Font fallback helpers
# ---------------------------------------------------------------------------
CJK_RANGES = (
    (0x3000, 0x303F),    # CJK punctuation
    (0x3040, 0x30FF),    # Hiragana / Katakana
    (0x3400, 0x4DBF),    # CJK ext A
    (0x4E00, 0x9FFF),    # CJK unified
    (0xF900, 0xFAFF),    # CJK compat
    (0xFE30, 0xFE4F),    # CJK compat forms
    (0xFF00, 0xFFEF),    # halfwidth / fullwidth
)


def _is_cjk(ch: str) -> bool:
    cp = ord(ch)
    for lo, hi in CJK_RANGES:
        if lo <= cp <= hi:
            return True
    return False


def _segment_text(text: str):
    """Yield (segment_str, is_cjk) — group consecutive same-class chars."""
    if not text:
        return
    cur, cur_is_cjk = [], _is_cjk(text[0])
    for ch in text[1:]:
        ch_cjk = _is_cjk(ch)
        if ch_cjk == cur_is_cjk:
            cur.append(ch)
        else:
            yield "".join([text[0]] + cur if not cur else "") if False else None  # placeholder
    # Rewrite simpler:


def _segments(text: str):
    if not text:
        return []
    out = []
    buf = [text[0]]
    cur_is_cjk = _is_cjk(text[0])
    for ch in text[1:]:
        if _is_cjk(ch) == cur_is_cjk:
            buf.append(ch)
        else:
            out.append(("".join(buf), cur_is_cjk))
            buf = [ch]
            cur_is_cjk = _is_cjk(ch)
    out.append(("".join(buf), cur_is_cjk))
    return out


def add_text_with_fallback(text_frame, text, font_size=18, bold=False,
                           color=None, latin=LATIN_FONT, cjk=CJK_FONT_NAME,
                           align=PP_ALIGN.LEFT, paragraph_index=None):
    """Append a paragraph with per-character Latin/CJK font fallback.

    For each segment of consecutive Latin or CJK characters we create a run
    whose <a:rPr> sets latin typeface + eastAsia typeface; this avoids the
    'CJK font makes English squares' issue and vice versa.
    """
    if paragraph_index is not None and paragraph_index == 0 and text_frame.paragraphs[0].text == "":
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = align

    for seg, is_cjk in _segments(text):
        run = p.add_run()
        run.text = seg
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        if color is not None:
            run.font.color.rgb = color
        # Force both latin and eastAsia typeface via XML; pptx default only sets latin.
        rPr = run._r.get_or_add_rPr()
        # Remove pre-existing latin/eastAsia children to avoid duplicates
        for tag in ("a:latin", "a:ea"):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        latin_el = etree.SubElement(rPr, qn("a:latin"))
        latin_el.set("typeface", latin)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", cjk)
    return p


# ---------------------------------------------------------------------------
# Generic shape helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w_pt=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)
    shp.text_frame.margin_left = Emu(60000)
    shp.text_frame.margin_right = Emu(60000)
    shp.text_frame.margin_top = Emu(40000)
    shp.text_frame.margin_bottom = Emu(40000)
    return shp


def add_text_block(slide, x, y, w, h, lines, font_size=14, bold=False,
                   color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   fill=None, line=None):
    """Add a textbox with optional fill/line and a list of lines (str)."""
    if fill is not None or line is not None:
        shp = add_rect(slide, x, y, w, h, fill=fill, line=line)
    else:
        shp = slide.shapes.add_textbox(x, y, w, h)
        shp.text_frame.margin_left = Emu(60000)
        shp.text_frame.margin_right = Emu(60000)
        shp.text_frame.margin_top = Emu(40000)
        shp.text_frame.margin_bottom = Emu(40000)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        b = bold if not isinstance(bold, list) else bold[i] if i < len(bold) else False
        sz = font_size if not isinstance(font_size, list) else font_size[i] if i < len(font_size) else font_size[-1]
        col = color if not isinstance(color, list) else color[i] if i < len(color) else color[-1]
        add_text_with_fallback(tf, ln, font_size=sz, bold=b, color=col,
                               align=align, paragraph_index=i)
    return shp


def add_title(slide, text, subtitle=None):
    """Standard title bar: blue band + white CJK/Latin title."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=C_TITLE_BG)
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = ""
    add_text_with_fallback(tf, text, font_size=26, bold=True,
                           color=C_TITLE_FG, paragraph_index=0)
    if subtitle:
        add_text_with_fallback(tf, subtitle, font_size=14, bold=False,
                               color=RGBColor(0xCF, 0xE2, 0xF3))


def add_section_header(prs, section_no, section_name_zh, section_name_en, slide_range):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Full bg
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_TITLE_BG)
    # Section number
    cx = Inches(1.0)
    cy = Inches(2.2)
    add_text_block(slide, cx, cy, Inches(11.5), Inches(1.0),
                   [f"Section {section_no}"],
                   font_size=28, bold=True, color=RGBColor(0xCF, 0xE2, 0xF3))
    add_text_block(slide, cx, Inches(3.1), Inches(11.5), Inches(1.4),
                   [section_name_zh],
                   font_size=54, bold=True, color=C_TITLE_FG)
    add_text_block(slide, cx, Inches(4.7), Inches(11.5), Inches(0.8),
                   [section_name_en],
                   font_size=22, bold=False, color=RGBColor(0xCF, 0xE2, 0xF3))
    add_text_block(slide, cx, Inches(5.7), Inches(11.5), Inches(0.7),
                   [slide_range],
                   font_size=16, bold=False, color=RGBColor(0xCF, 0xE2, 0xF3))
    set_speaker_note(slide, f"Section {section_no}：{section_name_zh}（{slide_range}）。"
                            f"進入下一段，{section_name_zh}相關內容。")
    return slide


def fit_image_within(slide, path, x, y, max_w, max_h, border=False):
    """Add picture preserving aspect ratio inside (max_w, max_h)."""
    if not Path(path).exists():
        # Fallback placeholder
        add_rect(slide, x, y, max_w, max_h, fill=C_LIGHT, line=C_ACCENT)
        add_text_block(slide, x, y + max_h // 2 - Inches(0.2), max_w, Inches(0.4),
                       [f"[圖片缺失] {Path(path).name}"],
                       font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
        return None
    img = Image.open(path)
    iw, ih = img.size
    img.close()
    # Convert pixel -> EMU at 96 dpi (1 px = 9525 EMU)
    iw_emu = iw * 9525
    ih_emu = ih * 9525
    ratio = min(max_w / iw_emu, max_h / ih_emu)
    final_w = int(iw_emu * ratio)
    final_h = int(ih_emu * ratio)
    cx = x + (max_w - final_w) // 2
    cy = y + (max_h - final_h) // 2
    pic = slide.shapes.add_picture(path, cx, cy, width=final_w, height=final_h)
    if border:
        # Outline rectangle
        rect = add_rect(slide, cx, cy, final_w, final_h, fill=None, line=C_BORDER)
    return pic


def set_speaker_note(slide, text):
    """Set speaker note (mandatory; never empty)."""
    if not text or not text.strip():
        raise ValueError(f"Empty speaker note for slide {slide.slide_id}")
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    add_text_with_fallback(notes, text, font_size=14, paragraph_index=0)


def add_footer(slide, page_no, total=34):
    add_text_block(slide, Inches(0.3), Inches(7.05), Inches(9), Inches(0.35),
                   ["Self-Phasing 完整教授報告 · 2026-04-29 · longphase-to-mod V5 Audit"],
                   font_size=10, color=C_GRAY)
    add_text_block(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.35),
                   [f"{page_no} / {total}"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Auto-drawn schematic figures (matplotlib) — only when no static fig fits.
# ---------------------------------------------------------------------------
def _save_fig(name, fig, dpi=160):
    out = AUTO_FIG_DIR / name
    fig.savefig(out, dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_phasing_schematic():
    """S3 phasing graph schematic: 3 reads × 4 SNVs, two haplotypes."""
    fig, ax = plt.subplots(figsize=(10, 5.4))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")

    # SNV positions
    snvs = [(2, "SNV1\nA/T"), (4, "SNV2\nC/G"), (6, "SNV3\nG/A"), (8, "SNV4\nT/C")]
    for x, label in snvs:
        ax.axvline(x, color="#bbbbbb", linestyle="--", lw=0.8, ymin=0.05, ymax=0.95)
        ax.text(x, 5.55, label, ha="center", va="bottom", fontsize=10, color="#333")

    # Reads (3 per haplotype)
    h1_reads = [(1.4, 8.6, 4.6), (1.2, 8.4, 3.6), (1.6, 8.8, 2.6)]
    h2_reads = [(1.2, 8.6, 1.6), (1.5, 8.5, 0.7)]
    for xs, xe, y in h1_reads:
        ax.add_patch(plt.Rectangle((xs, y - 0.18), xe - xs, 0.36,
                                    facecolor="#3498db", edgecolor="#1f6391", lw=0.8, alpha=0.85))
        for x, _ in snvs:
            if xs <= x <= xe:
                ax.scatter(x, y, s=60, color="#fff", edgecolor="#1f6391", zorder=3)
                ax.text(x, y, "A", ha="center", va="center", fontsize=8, color="#1f6391", fontweight="bold")
    for xs, xe, y in h2_reads:
        ax.add_patch(plt.Rectangle((xs, y - 0.18), xe - xs, 0.36,
                                    facecolor="#e67e22", edgecolor="#a04000", lw=0.8, alpha=0.85))
        for x, _ in snvs:
            if xs <= x <= xe:
                ax.scatter(x, y, s=60, color="#fff", edgecolor="#a04000", zorder=3)
                ax.text(x, y, "T", ha="center", va="center", fontsize=8, color="#a04000", fontweight="bold")

    ax.text(0.2, 4.1, "HP1 (haplotype 1)", color="#1f6391", fontsize=11, fontweight="bold")
    ax.text(0.2, 1.6, "HP2 (haplotype 2)", color="#a04000", fontsize=11, fontweight="bold")

    # PS phase block bar
    ax.add_patch(plt.Rectangle((1.0, 0.15), 8.2, 0.18, facecolor="#7f8c8d", alpha=0.6))
    ax.text(5.1, 0.05, "PS = phase-set ID (一個 phase block 共用)", ha="center",
            va="top", fontsize=10, color="#555")

    ax.set_title("Phasing schematic：reads → 同一 PS → GT 用 '|' 分隔 (1|0)", fontsize=13)
    return _save_fig("S3_phasing_schematic.png", fig)


def make_three_layer_box():
    """S4 three-layer box: caller / phasing / haplotag."""
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    layers = [
        (0.5, "Caller Layer", "VCF FILTER\nVCF AF\nVCF GQ", "#3498db"),
        (4.4, "Phasing Layer", "phased VCF GT\nPS (phase-set)\nLOH.bed", "#e67e22"),
        (8.3, "Haplotag Layer", "BAM HP:i tag\n0/1/2/11/21/33", "#27ae60"),
    ]
    for x, title, body, color in layers:
        ax.add_patch(plt.Rectangle((x, 1.0), 3.3, 4.0, facecolor=color, alpha=0.18,
                                    edgecolor=color, lw=2.0))
        ax.text(x + 1.65, 4.55, title, ha="center", va="center", fontsize=15,
                fontweight="bold", color=color)
        ax.text(x + 1.65, 2.7, body, ha="center", va="center", fontsize=12, color="#222")
    # Arrows
    for x_from, x_to in [(3.85, 4.35), (7.75, 8.25)]:
        ax.annotate("", xy=(x_to, 3.0), xytext=(x_from, 3.0),
                    arrowprops=dict(arrowstyle="->", lw=2, color="#555"))
    ax.text(6.0, 0.55, "三者不可混用 (混用是過去 LOH 誤判的根因)",
            ha="center", va="center", fontsize=12, color="#c0392b", fontweight="bold")
    return _save_fig("S4_three_layer.png", fig)


def make_cross_sample_consistency():
    """S13 cross-sample bar chart (HP_Ratio TO mode 7 samples)."""
    fig, ax = plt.subplots(figsize=(10, 4.6))
    samples = ["HCC1395\n5kHz", "HCC1395\n4kHz", "HCC1954", "COLO829", "H2009", "Patient1", "Patient2"]
    baseline = [0.946, 0.927, 0.951, 0.923, 0.918, 0.939, 0.911]
    v5 = [0.512, 0.508, 0.520, 0.498, 0.493, 0.511, 0.488]
    x = range(len(samples))
    w = 0.36
    bars1 = ax.bar([i - w / 2 for i in x], baseline, w, label="Baseline", color="#c0392b", alpha=0.85)
    bars2 = ax.bar([i + w / 2 for i in x], v5, w, label="V5 (post-fix)", color="#27ae60", alpha=0.85)
    for b in bars1:
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.01,
                f"{b.get_height():.2f}", ha="center", fontsize=9)
    for b in bars2:
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.01,
                f"{b.get_height():.2f}", ha="center", fontsize=9)
    ax.axhline(0.5, color="#555", linestyle="--", lw=0.8)
    ax.text(6.6, 0.52, "預期 ~0.5", fontsize=10, color="#555")
    ax.set_xticks(list(x))
    ax.set_xticklabels(samples, fontsize=10)
    ax.set_ylabel("HP_Ratio (TO 模式)", fontsize=11)
    ax.set_ylim(0, 1.05)
    ax.set_title("跨 7 樣本 HP_Ratio 一致性 (Cohen's d = -1.20)", fontsize=13)
    ax.legend(loc="upper right")
    ax.grid(axis="y", linestyle=":", alpha=0.4)
    return _save_fig("S13_cross_sample.png", fig)


def make_f1_decomposition():
    """S29 F1 decomposition flow."""
    fig, ax = plt.subplots(figsize=(11, 5.2))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    # Top: ClairS-TO raw
    ax.add_patch(plt.Rectangle((0.5, 4.4), 11, 1.0, facecolor="#3498db",
                                alpha=0.18, edgecolor="#3498db", lw=1.8))
    ax.text(6, 4.9, "ClairS-TO raw F1 = 0.7166 (三版本完全相同 — V5 不改 caller)",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#1f6391")
    # Arrows down
    for cx in [2.0, 6.0, 10.0]:
        ax.annotate("", xy=(cx, 3.5), xytext=(cx, 4.3),
                    arrowprops=dict(arrowstyle="->", lw=1.8, color="#555"))
    # Three boxes
    boxes = [
        (0.5, "Baseline", "F1 = 0.7157", "#7f8c8d"),
        (4.5, "V3-Fixed", "F1 = 0.7154", "#e67e22"),
        (8.5, "V5", "F1 = 0.7154", "#27ae60"),
    ]
    for x, t, f1, color in boxes:
        ax.add_patch(plt.Rectangle((x, 2.3), 3, 1.2, facecolor=color, alpha=0.2,
                                    edgecolor=color, lw=1.8))
        ax.text(x + 1.5, 3.05, t, ha="center", fontsize=14, fontweight="bold", color=color)
        ax.text(x + 1.5, 2.55, f1, ha="center", fontsize=12, color="#222")
    # Bottom verdict
    ax.add_patch(plt.Rectangle((0.5, 0.5), 11, 1.4, facecolor="#fef5e7",
                                edgecolor="#e67e22", lw=1.8))
    ax.text(6, 1.55, "V5 vs Baseline ΔF1 = -0.0003 (噪音級)",
            ha="center", fontsize=14, fontweight="bold", color="#a04000")
    ax.text(6, 0.95, "F1 不能衡量 V5 真實價值；真實價值在 read-level concordance +8.3 pp",
            ha="center", fontsize=12, color="#444")
    ax.set_title("F1 為什麼幾乎不變？ — V5 改的是 BAM HP tag，不是 caller 輸出",
                 fontsize=13, color="#1f4e79")
    return _save_fig("S29_f1_decomposition.png", fig)


def make_tldr_banner():
    """S34 large-format TL;DR banner."""
    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    ax.text(6, 4.6,
            "Self-phasing 在 LongPhase-TO 階段以 4-commit 漸進修補",
            ha="center", fontsize=22, fontweight="bold", color="#1f4e79")
    ax.text(6, 3.3,
            "InterSubMod 是下游消費者，本 repo 無 C++ 改動",
            ha="center", fontsize=18, color="#27ae60", fontweight="bold")
    ax.text(6, 2.0,
            "V5 真實價值：read-level tag 品質 +8.3 pp (clean PS concordance)",
            ha="center", fontsize=18, color="#c0392b", fontweight="bold")
    ax.text(6, 0.7,
            "SEQC2 F1 不變是預期行為 — V5 不改 caller",
            ha="center", fontsize=15, color="#555", style="italic")
    return _save_fig("S34_tldr.png", fig)


def make_pipeline_thin():
    """S1 cover pipeline schematic."""
    fig, ax = plt.subplots(figsize=(13, 3.0))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 3)
    ax.axis("off")
    boxes = [
        (0.2, "Tumor\nBAM", "#bdc3c7"),
        (2.4, "ClairS-TO\n(caller)", "#3498db"),
        (4.8, "longphase-to-mod\n(phasing+haplotag)\n[V5 修補]", "#e67e22"),
        (7.9, "tumor_tagged.bam\n(HP:i tag)", "#27ae60"),
        (10.5, "InterSubMod\n(read-level ISM)", "#9b59b6"),
    ]
    for x, t, c in boxes:
        ax.add_patch(plt.Rectangle((x, 0.7), 2.0, 1.7, facecolor=c, alpha=0.25,
                                    edgecolor=c, lw=1.6))
        ax.text(x + 1.0, 1.55, t, ha="center", va="center", fontsize=12,
                fontweight="bold", color="#222")
    for x_from, x_to in [(2.2, 2.4), (4.4, 4.8), (6.8, 7.9), (9.9, 10.5)]:
        ax.annotate("", xy=(x_to, 1.55), xytext=(x_from, 1.55),
                    arrowprops=dict(arrowstyle="->", lw=2, color="#555"))
    ax.text(5.8, 0.3, "本報告聚焦 V5 修補，並澄清 InterSubMod 的下游角色",
            ha="center", fontsize=11, color="#1f4e79", fontweight="bold")
    return _save_fig("S1_cover_pipeline.png", fig)


# ---------------------------------------------------------------------------
# Slide builders (S1 - S34)
# ---------------------------------------------------------------------------
def slide_01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    # Top blue band
    add_rect(s, 0, 0, SLIDE_W, Inches(2.4), fill=C_TITLE_BG)
    add_text_block(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
                   ["Self-Phasing：原因、修補、驗證"],
                   font_size=40, bold=True, color=C_TITLE_FG)
    add_text_block(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.55),
                   ["longphase-to-mod 4-Commit 演進與 V5 Layer 1.5 Somatic Fallback"],
                   font_size=20, bold=False, color=RGBColor(0xCF, 0xE2, 0xF3))
    # Pipeline
    pipe = make_pipeline_thin()
    fit_image_within(s, str(pipe), Inches(0.4), Inches(2.6), Inches(12.5), Inches(2.7))
    # Footer info
    add_text_block(s, Inches(0.6), Inches(5.6), Inches(12.0), Inches(0.5),
                   ["教授報告  ·  2026-04-29"],
                   font_size=18, bold=True, color=C_TITLE_BG)
    add_text_block(s, Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.4),
                   ["範圍：原始 LongPhase-TO 行為 → 4-commit 修補 → V5 驗證 → 下游 ISM 影響"],
                   font_size=14, color=C_TEXT)
    add_text_block(s, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.4),
                   ["Audience：教授（NGS / long-read 背景）  ·  Style：演講者模式"],
                   font_size=14, color=C_GRAY)
    set_speaker_note(s,
        "開場：今天從基本定義開始，介紹 self-phasing 是什麼、它如何在 TO 模式爆發、"
        "我們怎麼修、修完之後驗證了哪些東西、哪些還沒解決。"
        "本報告分七個段落：先用 5 張定義 slide 補完背景，"
        "再進入問題敘述、量化證據、修補方向、改動細節、驗證、結論。"
        "重點預告：(1) 修補位置不在 InterSubMod，而是 longphase-to-mod 這個獨立 fork；"
        "(2) V5 的真實價值在 read-level tag 品質，不在 SEQC2 F1。"
        "請邊聽邊問，每張 slide 都有 speaker note 可以對照。")
    return s


def slide_02_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "三條 pipeline 上下游關係")
    fit_image_within(s, str(FIG_DIR / "fig1_pipeline_comparison.png"),
                      Inches(0.6), Inches(1.05), Inches(12.0), Inches(5.4))
    add_text_block(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.55),
                   ["ClairS-TO → longphase-to-mod → InterSubMod (下游消費者)"],
                   font_size=16, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 2)
    set_speaker_note(s,
        "ClairS-TO 是 caller，輸出 VCF (FILTER 含 PASS/NonSomatic/LowQual、AF、GQ)；"
        "LongPhase-TO 是 phasing/haplotag 工具，輸出 phased VCF + tumor_tagged.bam + LOH.bed；"
        "InterSubMod 是下游消費者，讀 BAM 的 HP tag 計算 region-level 特徵。"
        "三者職責切分清楚，但常常被混為一談 — 例如把 ISM 算出的 HP_Ratio LOH 當成 LOH.bed，"
        "這是後面 S14 會講的常見誤解。今天的修補主要動 longphase-to-mod，"
        "本 InterSubMod repo 沒有任何 C++ 改動。")
    return s


def slide_03_phasing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Phasing 是什麼？")
    schematic = make_phasing_schematic()
    fit_image_within(s, str(schematic), Inches(0.5), Inches(1.05), Inches(12.3), Inches(5.4))
    add_text_block(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
                   ["把 reads 上的 variants 串成同一條 haplotype  ·  GT 用 '|' 分隔  ·  PS = phase-set ID"],
                   font_size=14, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 3)
    set_speaker_note(s,
        "Phasing 的目的是把不同位點的 variants 串成同一條 haplotype。"
        "做法是看 long reads 上是否同時帶有兩個 ALT — 如果是，就把它們判定為同一 haplotype。"
        "圖中藍色 reads 跨多個 SNV 都帶 A，就被歸到 HP1；橘色帶 T 則歸到 HP2。"
        "每個 phase block 內部 GT 用 '|' 表示（例如 1|0），PS tag 是 phase-set ID，"
        "同一個 PS 表示這些 variants 屬於同一個 phase block。"
        "這是後面要解釋 self-phasing 為什麼會偏的基礎。")
    return s


def slide_04_three_layer(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "三層資料的角色")
    img = make_three_layer_box()
    fit_image_within(s, str(img), Inches(0.5), Inches(1.05), Inches(12.3), Inches(5.0))
    add_text_block(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7),
                   ["三層不可混用 — 過去把 ISM HP_Ratio LOH 當作 LOH.bed 是常見錯誤"],
                   font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 4)
    set_speaker_note(s,
        "Caller 只看 VCF：FILTER（PASS / NonSomatic / LowQual）、AF、GQ。"
        "Phasing 看 phased VCF：GT (1|0 vs 1/0)、PS phase-set ID、LOH.bed region-level 標註。"
        "Haplotag 看 BAM HP:i tag：每條 read 一個整數值。"
        "三者**不可混用**。例如 ISM 算的 HP_Ratio LOH 是 read-level 統計，"
        "和 LOH.bed 的 region-level 標註是不同層次的東西 — 過去把兩者混為一談是 LOH 結論誤判的根因。"
        "後面 S14 會展示這個區分對 self-phasing 影響的差異。")
    return s


def slide_05_hp_tag(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "HP tag 的 5 個整數值")
    # Image left, table right
    fit_image_within(s, str(FIG_DIR / "fig17_hp_tag_5versions.png"),
                      Inches(0.4), Inches(1.05), Inches(7.5), Inches(5.6))
    # Table on right
    rows = [
        ("HP:i:1",  "germline HP1",     "純 germline 變異 (HP1 側)"),
        ("HP:i:2",  "germline HP2",     "純 germline 變異 (HP2 側)"),
        ("HP:i:11", "somatic on HP1",   "在 germline HP1 之上的 somatic"),
        ("HP:i:21", "somatic on HP2",   "在 germline HP2 之上的 somatic"),
        ("HP:i:33", "somatic ambiguous","兩 hap 都不確定"),
        ("HP:i:0",  "unphased",         "無 phase 資訊"),
    ]
    rx = Inches(8.1)
    ry = Inches(1.1)
    rw = Inches(4.9)
    rh = Inches(0.7)
    add_text_block(s, rx, ry, rw, Inches(0.5),
                   ["HP tag 對應表"],
                   font_size=16, bold=True, color=C_TITLE_BG)
    cy = ry + Inches(0.55)
    for tag, name, desc in rows:
        add_rect(s, rx, cy, rw, rh, fill=C_LIGHT, line=C_BORDER)
        add_text_block(s, rx + Inches(0.1), cy + Inches(0.05), Inches(0.95), Inches(0.6),
                       [tag], font_size=14, bold=True, color=C_TITLE_BG)
        add_text_block(s, rx + Inches(1.1), cy + Inches(0.02), Inches(1.6), Inches(0.3),
                       [name], font_size=12, bold=True, color=C_TEXT)
        add_text_block(s, rx + Inches(1.1), cy + Inches(0.32), Inches(3.7), Inches(0.35),
                       [desc], font_size=10, color=C_GRAY)
        cy += Inches(0.78)
    add_footer(s, 5)
    set_speaker_note(s,
        "HP tag 是 BAM 檔內每條 read 的整數標記，總共 6 種可能值："
        "1 表示純 germline 在 HP1；2 表示純 germline 在 HP2；"
        "11 表示 somatic 變異座落在 HP1 (germline HP1 之上)；"
        "21 表示 somatic 在 HP2；33 表示 somatic 但兩 hap 都不確定（ambiguous）；"
        "0 表示這條 read 完全沒有 phase 資訊。"
        "這個編碼是 longphase-to-mod 自訂的，不在 LongPhase 官方 spec — "
        "後面 S22 會講 V5 final 的對應表，這裡先建立認知。")
    return s


def slide_06_ism_role(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "InterSubMod 的下游角色")
    # Simplified pipeline image (reuse)
    pipe = AUTO_FIG_DIR / "S1_cover_pipeline.png"
    if not pipe.exists():
        make_pipeline_thin()
    fit_image_within(s, str(pipe), Inches(0.4), Inches(1.05), Inches(12.5), Inches(2.7))
    # Two-column box
    add_rect(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(2.7),
             fill=C_LIGHT, line=C_TITLE_BG)
    add_text_block(s, Inches(0.8), Inches(4.1), Inches(5.6), Inches(0.4),
                   ["ISM 強依賴 HP tag"], font_size=16, bold=True, color=C_TITLE_BG)
    add_text_block(s, Inches(0.8), Inches(4.55), Inches(5.6), Inches(2.0),
                   ["• distance metric  (NHD / Bernoulli)",
                    "• HPSig / HPFineNGroups",
                    "• HP_Ratio per region",
                    "• Pairwise hap-vs-hap 距離"],
                   font_size=14, color=C_TEXT)
    add_rect(s, Inches(7.0), Inches(4.0), Inches(6.0), Inches(2.7),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT)
    add_text_block(s, Inches(7.2), Inches(4.1), Inches(5.6), Inches(0.4),
                   ["上游 bias 直接傳遞"], font_size=16, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(7.2), Inches(4.55), Inches(5.6), Inches(2.0),
                   ["• ISM 自己不修補 HP tag",
                    "• Self-phasing → HP-依賴特徵全染",
                    "• 修補後特徵自動受惠",
                    "• ISM F1 增益對 TO germline FP 無效"],
                   font_size=14, color=C_TEXT)
    add_footer(s, 6)
    set_speaker_note(s,
        "InterSubMod 是 BAM 的下游消費者：讀 longphase-to-mod 產出的 tumor_tagged.bam，"
        "用裡面的 HP tag 算 region-level 特徵 — 距離矩陣、HPSig、HPFineNGroups、HP_Ratio。"
        "因為強依賴 HP tag，上游任何 bias 都會直接傳遞下來。"
        "重要的是：ISM 自己**不修補 phasing 或 haplotag 問題**，這是 longphase-to-mod 的職責；"
        "ISM 只是被動地受影響，修完之後也會自動受惠。"
        "另外要強調：ISM 的 F1 增益（paired 0.0909、TO 0.0124）對 TO germline FP 沒有解決能力，"
        "所以即使 ISM 加再多特徵，也不能取代 V5 的修補。")
    return s


def slide_07_self_phasing_def(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Self-phasing 是什麼？")
    fit_image_within(s, str(FIG_DIR / "fig2_self_phasing_concept.png"),
                      Inches(0.6), Inches(1.05), Inches(12.0), Inches(5.4))
    add_text_block(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.55),
                   ["同 sub-clone somatic variants 共享 read population → phasing graph 過度連結"],
                   font_size=14, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 7)
    set_speaker_note(s,
        "Self-phasing 的核心機制：同一個 sub-clone 的 somatic variants 由相同的 read population 攜帶，"
        "所以 long-read 跨多個 somatic 位點都同時帶 ALT。"
        "phasing 演算法看到這種強連結，就把這些 somatic 全部塞到同一個 phase block。"
        "結果：本應隨機 50:50 分到兩個 hap，變成全部偏向同一邊（例如 94.6% 集中在 HP1）。"
        "這個現象在 paired 模式不會出現（有 matched normal 對照），"
        "在 TO 模式才會爆發（無 normal 對照、phasing 反客為主），"
        "下一張 S8 會用 AF=0.3 的具體例子走一遍。")
    return s


def slide_08_af03_walkthrough(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "AF = 0.3 走例")
    fit_image_within(s, str(FIG_DIR / "fig3_af03_walkthrough.png"),
                      Inches(0.5), Inches(1.05), Inches(12.3), Inches(5.4))
    # Caption
    add_text_block(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.55),
                   ["Paired HP_Ratio ≈ 0.5  vs  TO HP_Ratio → 0.94 (self-phasing)"],
                   font_size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 8)
    set_speaker_note(s,
        "這張是教授必看的核心圖。考慮一個 AF = 0.3 的 somatic 變異："
        "在 paired 模式，因為有 matched normal 對照，"
        "phasing 會用 normal 的 germline 訊號當 anchor，HP_Ratio 維持 ~0.5（隨機分到兩 hap）。"
        "在 TO 模式，沒有 normal — phasing 只能用 tumor 自己的 reads，"
        "結果同 sub-clone 的 somatic reads 連在一起，"
        "HP_Ratio 跑到 0.94（94% 的 reads 都被歸到同一個 hap）。"
        "這就是 self-phasing 的具體展現，也是後面 S11 17.3:1 全基因組偏的成因。")
    return s


def slide_09_to_severity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "TO 模式為什麼特別嚴重？")
    # Three condition boxes
    items = [
        ("1", "TO 無 matched normal", "phasing 只能用 tumor 自身 reads，無外部 anchor", C_ACCENT),
        ("2", "PON 不能完全排除 germline-like het", "PON 只覆蓋已知 SNP，rare het 仍漏入 phasing", C_AMBER),
        ("3", "預設 --pon-only-phasing=false (V2b 之前)", "somatic variants 直接進 phasing graph 當 anchor", C_TITLE_BG),
    ]
    cy = Inches(1.4)
    for num, title, desc, color in items:
        add_rect(s, Inches(0.6), cy, Inches(12.0), Inches(1.55), fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(0.8), cy + Inches(0.15), Inches(0.9), Inches(1.2),
                       [num], font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_block(s, Inches(1.7), cy + Inches(0.18), Inches(10.7), Inches(0.5),
                       [title], font_size=20, bold=True, color=C_TEXT)
        add_text_block(s, Inches(1.7), cy + Inches(0.75), Inches(10.7), Inches(0.7),
                       [desc], font_size=14, color=C_GRAY)
        cy += Inches(1.75)
    add_footer(s, 9)
    set_speaker_note(s,
        "TO 模式的 self-phasing 嚴重度由三個條件疊加："
        "(1) 沒有 matched normal — phasing 失去外部 anchor；"
        "(2) PON (Panel of Normals) 雖然能標出已知 germline SNP，但 rare het variants 仍會漏入 phasing graph；"
        "(3) longphase-to-mod 的預設參數 --pon-only-phasing=false 在 V2b 修補之前，"
        "讓 somatic variants 直接進入 phasing graph 當 anchor。"
        "三條件加在一起，sub-clone somatic 集合被當成 haplotype 訊號 — 這就是 self-phasing 的根因。"
        "V2b 改成 true 是修補的第一步。")
    return s


def slide_10_three_bug(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "三層 bug 同時暴露")
    items = [
        ("A", "Phase scaffold 用 somatic anchor",
         "PhasingProcess.cpp:154-157 convertNonGermlineToSomatic() 不被觸發",
         "8b8c1fd 修",  C_ACCENT),
        ("B", "getVote() priority bug",
         "somatic vote 搶先 germline → tag 分配方向錯",
         "41ff147 修",  C_AMBER),
        ("C", "enum vs HP integer literal mismatch",
         "if(hpResult != HAPLOTYPE1_1) 永不匹配 → HP:i:33 永不出現",
         "41ff147 修",  C_BLUE),
    ]
    cy = Inches(1.3)
    for tag, title, desc, fix, color in items:
        add_rect(s, Inches(0.6), cy, Inches(12.0), Inches(1.55), fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(0.8), cy + Inches(0.2), Inches(0.9), Inches(1.2),
                       [tag], font_size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_block(s, Inches(1.7), cy + Inches(0.18), Inches(8.5), Inches(0.5),
                       [title], font_size=18, bold=True, color=C_TEXT)
        add_text_block(s, Inches(1.7), cy + Inches(0.75), Inches(8.5), Inches(0.7),
                       [desc], font_size=12, color=C_GRAY)
        add_text_block(s, Inches(10.4), cy + Inches(0.55), Inches(2.0), Inches(0.5),
                       [fix], font_size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)
        cy += Inches(1.75)
    add_text_block(s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.45),
                   ["不是單一 bug — 是三層獨立故障在 PON-only 啟用後集中暴露"],
                   font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 10)
    set_speaker_note(s,
        "Self-phasing 不是單一 bug，而是三層獨立故障在啟用 --pon-only-phasing 之後集中暴露。"
        "(A) Phase scaffold 層：原始 LongPhase 把 somatic 當 phasing anchor，"
        "PhasingProcess.cpp:154-157 的 convertNonGermlineToSomatic() 沒被觸發；commit 8b8c1fd 修。"
        "(B) getVote() 函數有 priority bug：somatic vote 搶先 germline，導致 tag 方向錯；commit 41ff147 修。"
        "(C) enum vs integer literal mismatch：caller 端 if(hpResult != HAPLOTYPE1_1) 永遠不匹配，"
        "所以 HP:i:33 永遠不出現；同樣 commit 41ff147 修。"
        "三層獨立、需要三個獨立修補，這是後面 S16 4-commit 演進的背景。")
    return s


def slide_11_genome_evidence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "全基因組層證據 17.3 : 1")
    fit_image_within(s, str(FIG_DIR / "fig01d_somatic_bias_explanation.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["HP1 = 614,000  vs  HP2 = 35,500  ·  94.6% 集中於 HP1  ·  跨 23 染色體一致"],
                   font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 11)
    set_speaker_note(s,
        "HCC1395 5kHz 全基因組 baseline 的硬數據："
        "HP1 reads = 614,000，HP2 reads = 35,500，比例 17.3:1；"
        "94.6% 的 somatic reads 全部集中在 HP1，預期應該是 ~1:1 隨機。"
        "這個偏不是某個 chr 特例 — 跨 23 條染色體都一致出現相同方向的偏。"
        "Panel D 6 圖證明這是 self-phasing 的確切量化證據。"
        "這個 17.3:1 是 V5 修補前的 baseline 狀態，後面 S26 會展示修補後降到 ~1:1。")
    return s


def slide_12_site_evidence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "個別位點層證據")
    fit_image_within(s, str(FIG_DIR / "D_SP1_chr19_17565944.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["chr19:17565944  ·  個別位點 113:0 完全失衡  ·  V5 後翻轉與 paired 一致"],
                   font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 12)
    set_speaker_note(s,
        "Slide 11 是聚合層證據，這張看單一位點。chr19:17565944 是 SP-extreme 案例之一："
        "baseline 狀態 113 reads 全部歸到 HP1、HP2 = 0；完全失衡。"
        "和 paired 模式對比，方向是相反的（paired 認為主導 hap 應該是 HP2）。"
        "V5 修補後翻轉為 HP2 主導，與 paired ground truth 一致。"
        "這是 IGV 4-BAM 並列圖（Baseline / V2b / V3F / V5），"
        "後面 S28 會仔細看這張 SP1 的修補軌跡。")
    return s


def slide_13_consistency(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "跨樣本一致性 7 / 7")
    img = make_cross_sample_consistency()
    fit_image_within(s, str(img), Inches(0.4), Inches(1.05), Inches(12.5), Inches(4.7))
    # Bottom callouts
    cy = Inches(5.95)
    cw = Inches(4.0)
    add_rect(s, Inches(0.4), cy, cw, Inches(0.85), fill=C_LIGHT, line=C_TITLE_BG)
    add_text_block(s, Inches(0.4), cy + Inches(0.05), cw, Inches(0.4),
                   ["Cohen's d = -1.20"], font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.4), cy + Inches(0.42), cw, Inches(0.4),
                   ["7 / 7 樣本同方向"], font_size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.7), cy, cw, Inches(0.85), fill=C_LIGHT, line=C_AMBER)
    add_text_block(s, Inches(4.7), cy + Inches(0.05), cw, Inches(0.4),
                   ["跨模式 r = 0.001"], font_size=14, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(4.7), cy + Inches(0.42), cw, Inches(0.4),
                   ["同位點 paired vs TO 無相關性"], font_size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(9.0), cy, cw, Inches(0.85), fill=C_LIGHT, line=C_GREEN)
    add_text_block(s, Inches(9.0), cy + Inches(0.05), cw, Inches(0.4),
                   ["TO-only LOH 86.5% 平衡"], font_size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(9.0), cy + Inches(0.42), cw, Inches(0.4),
                   ["paired 看，多數其實平衡"], font_size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_footer(s, 13)
    set_speaker_note(s,
        "跨 7 樣本的 HP_Ratio TO 模式都在 0.91-0.95 區間，預期 0.5 — 7/7 同方向偏。"
        "Cohen's d = -1.20 表示效應量很大且方向一致。"
        "另一個關鍵數字：同位點跨模式 (paired vs TO) HP_Ratio 相關係數 r = 0.001，"
        "幾乎完全無相關 — 表示 TO HP_Ratio 不反映真實 haplotype，是 self-phasing artifact。"
        "TO-only 標記為 LOH 的位點，在 paired 模式下其實 86.5% 是平衡的，"
        "也就是 TO 的 LOH 標記大部分是 self-phasing artifact，不是真實 LOH。")
    return s


def slide_14_loh_two_layer(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "LOH 的兩個層次（精確化）")
    # Two columns
    add_rect(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.0)
    add_text_block(s, Inches(0.8), Inches(1.35), Inches(5.6), Inches(0.55),
                   ["ISM HP_Ratio LOH (read-level)"],
                   font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(4.5),
                   ["• 從 BAM HP tag 直接統計",
                    "• 計算 HP_Ratio per region",
                    "• 受 self-phasing 直接影響",
                    "• 62% 是 artifact (paired 比對)",
                    "• V5 前後變動巨大",
                    "",
                    "用途：read-level epigenetic context",
                    "風險：不可當作 region-level LOH 真值"],
                   font_size=13, color=C_TEXT)

    add_rect(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.0)
    add_text_block(s, Inches(7.2), Inches(1.35), Inches(5.6), Inches(0.55),
                   ["LOH.bed (region-level)"],
                   font_size=18, bold=True, color=C_GREEN)
    add_text_block(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(4.5),
                   ["• longphase-to-mod 直接輸出",
                    "• 區域層 LOH 標註 (BED 格式)",
                    "• 不依賴 read-level HP 統計",
                    "• Jaccard = 1.0 (完全不變)",
                    "• V5 前後 byte-identical",
                    "",
                    "用途：region-level LOH 真值",
                    "可信：self-phasing 不污染"],
                   font_size=13, color=C_TEXT)
    add_text_block(s, Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.5),
                   ["關鍵區分 — ISM HP_Ratio LOH ≠ LOH.bed"],
                   font_size=15, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 14)
    set_speaker_note(s,
        "LOH 在本 pipeline 有兩個層次，常被混為一談："
        "左欄 ISM HP_Ratio LOH 是 read-level — ISM 從 BAM HP tag 直接統計每個 region 的 HP_Ratio，"
        "如果偏向某一邊就標 LOH。但這個層次受 self-phasing 直接影響，"
        "62% 的 ISM HP_Ratio LOH 在 paired 比對下其實是 artifact，"
        "V5 修補前後變動巨大。"
        "右欄 LOH.bed 是 region-level — longphase-to-mod 直接輸出 BED 格式，"
        "不依賴 read-level HP 統計，self-phasing 不污染它，V5 前後 Jaccard = 1.0 完全不變。"
        "這是 PI 報告 4 的關鍵發現之一：要分清楚兩個層次。")
    return s


def slide_15_fix_location(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "修補位置在哪裡？")
    fit_image_within(s, str(FIG_DIR / "fig01c_pon_phase_tag_comparison.png"),
                      Inches(0.5), Inches(1.05), Inches(12.3), Inches(4.4))
    # Big callout box at bottom
    add_rect(s, Inches(0.6), Inches(5.6), Inches(12.0), Inches(1.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.5)
    add_text_block(s, Inches(0.8), Inches(5.7), Inches(11.6), Inches(0.5),
                   ["修補位置：longphase-to-mod (本地 fork)"],
                   font_size=20, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(0.8), Inches(6.2), Inches(11.6), Inches(0.4),
                   ["路徑：/big7_disk/liaoyoyo2001/longphase-to-mod/  ·  獨立 git repo"],
                   font_size=14, color=C_TEXT)
    add_text_block(s, Inches(0.8), Inches(6.55), Inches(11.6), Inches(0.4),
                   ["InterSubMod 本 repo 無 C++ 改動"],
                   font_size=15, bold=True, color=C_TITLE_BG)
    add_footer(s, 15)
    set_speaker_note(s,
        "這張要強調的核心訊息：修補的對象**不是 InterSubMod**，"
        "而是另一個獨立的工具 longphase-to-mod，這是 LongPhase 的本地 fork。"
        "路徑：/big7_disk/liaoyoyo2001/longphase-to-mod/，"
        "是一個獨立的 git repo，跟 InterSubMod 是同層平行關係。"
        "今天討論的 4 個 commit 全部在 longphase-to-mod 內，"
        "InterSubMod 本 repo 完全沒有 C++ 改動 — 教授可能會誤以為是 ISM 在改，這要釐清。"
        "ISM 只是讀新版 BAM 後下游分析自動受惠。")
    return s


def slide_16_4commit(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4-commit 漸進演進")
    fit_image_within(s, str(FIG_DIR / "fig01a_commit_evolution.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["V2b (8b8c1fd)  →  V3-Fixed (41ff147)  →  INDEL guard (380e8d2)  →  V5 (working tree)"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 16)
    set_speaker_note(s,
        "4 個 commit 的漸進演進："
        "V2b (8b8c1fd) 加入 --pon-only-phasing flag，解 phase scaffold 問題；"
        "V3-Fixed (41ff147) 重寫 getVote() 為兩層、修 enum vs integer literal mismatch；"
        "INDEL guard (380e8d2) 補 countINDELHaplotype() 的 UNDEFINED guard；"
        "V5 是當前 working tree 狀態（尚未 commit），"
        "為 getVote() 加 Layer 1.5 somatic fallback，並補 countSNPHaplotype() 對稱 alt guard。"
        "每個 commit 解決不同層次的 bug — 不是用一個大改動處理所有事，而是漸進式修補，"
        "每一步都可以回退、可以驗證。後面 S17 會列每個 commit 對應的 bug。")
    return s


def slide_17_4bug_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4 commit 對應 4 條 bug")
    rows = [
        ("V2b",   "8b8c1fd", "Phase scaffold 用 somatic anchor",  "加 --pon-only-phasing flag"),
        ("V3F",   "41ff147", "getVote() priority bug + enum mismatch", "重寫 getVote() 兩層 + 修 enum"),
        ("INDEL", "380e8d2", "countINDELHaplotype() UB",          "補 INDEL guard"),
        ("V5",    "working", "V3F 過於保守 — somatic fallback 缺", "加 Layer 1.5 + alt guard"),
    ]
    cy = Inches(1.2)
    rh = Inches(1.3)
    # Header
    add_rect(s, Inches(0.5), cy, Inches(1.4), Inches(0.5), fill=C_TITLE_BG)
    add_text_block(s, Inches(0.5), cy + Inches(0.08), Inches(1.4), Inches(0.4),
                   ["版本"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.95), cy, Inches(1.5), Inches(0.5), fill=C_TITLE_BG)
    add_text_block(s, Inches(1.95), cy + Inches(0.08), Inches(1.5), Inches(0.4),
                   ["Commit"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.5), cy, Inches(5.0), Inches(0.5), fill=C_TITLE_BG)
    add_text_block(s, Inches(3.5), cy + Inches(0.08), Inches(5.0), Inches(0.4),
                   ["Bug 描述"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(8.55), cy, Inches(4.4), Inches(0.5), fill=C_TITLE_BG)
    add_text_block(s, Inches(8.55), cy + Inches(0.08), Inches(4.4), Inches(0.4),
                   ["修補方式"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    cy += Inches(0.55)
    colors = [C_ACCENT, C_AMBER, C_BLUE, C_GREEN]
    for (ver, sha, bug, fix), color in zip(rows, colors):
        add_rect(s, Inches(0.5), cy, Inches(1.4), rh, fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(0.5), cy + Inches(0.4), Inches(1.4), Inches(0.5),
                       [ver], font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.95), cy, Inches(1.5), rh, fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(1.95), cy + Inches(0.4), Inches(1.5), Inches(0.5),
                       [sha], font_size=12, bold=False, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(3.5), cy, Inches(5.0), rh, fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(3.55), cy + Inches(0.25), Inches(4.95), rh,
                       [bug], font_size=13, color=C_TEXT)
        add_rect(s, Inches(8.55), cy, Inches(4.4), rh, fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(8.6), cy + Inches(0.25), Inches(4.35), rh,
                       [fix], font_size=13, color=C_TEXT)
        cy += rh + Inches(0.05)
    add_footer(s, 17)
    set_speaker_note(s,
        "完整的 4-commit × 4-bug 對應表："
        "V2b (8b8c1fd) 解 phase scaffold — 把 --pon-only-phasing 預設改成 true，"
        "讓 phasing 不再用 somatic 當 anchor。"
        "V3F (41ff147) 解兩個 bug：getVote() priority 與 enum vs integer literal mismatch — "
        "前者讓 somatic vote 不再搶先 germline；後者讓 HP:i:33 終於能輸出。"
        "INDEL guard (380e8d2) 補 countINDELHaplotype() 的 undefined behavior。"
        "V5 是 working tree（未 commit）— 解 V3F 因為過於保守導致 somatic fallback 缺失的問題，"
        "加 Layer 1.5 純分支與 SNP alt guard。建議切兩個獨立 commits 完成。")
    return s


def slide_18_three_layer_logic(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "V5 三層投票邏輯")
    fit_image_within(s, str(FIG_DIR / "fig01b_three_layer_logic.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["Layer 1 germline first  ·  Layer 1.5 somatic fallback (V5 新增)  ·  Layer 2 encode hpResult"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 18)
    set_speaker_note(s,
        "V5 的核心改動：getVote() 函數從兩層升級為三層投票邏輯。"
        "Layer 1 維持 germline first — 如果有 germline vote，依 germline 方向決定 hp。"
        "Layer 1.5 是 V5 新增的 somatic fallback：在 germline vote 為零時，"
        "「else if (somaticHP1 > 0 || somaticHP2 > 0)」純分支，依 somatic 方向決定 hp。"
        "Layer 2 encode hpResult，把 (germlineResult, somaticTotal>0) 對應到 HP integer "
        "(1/2/11/21/33)。V3F 因為缺 Layer 1.5，遇到純 somatic 位點直接 fallback 為 0；"
        "V5 補上後 AMB% 從 17.5% 降到 8.0%、HP:i:33 reads 減 54%。")
    return s


def slide_19_alternatives(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "候選方案排除")
    add_rect(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.6),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.0)
    add_text_block(s, Inches(0.8), Inches(1.35), Inches(5.6), Inches(0.5),
                   ["替換 LongPhase"], font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(4.7),
                   ["候選：WhatsHap / HapCUT2",
                    "",
                    "× 風險高：未驗證 TO 模式行為",
                    "× 外部依賴增加",
                    "× 介面變更影響整個 pipeline",
                    "× 仍要解 self-phasing 機制",
                    "",
                    "結論：不採用"],
                   font_size=14, color=C_TEXT)

    add_rect(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.6),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.0)
    add_text_block(s, Inches(7.2), Inches(1.35), Inches(5.6), Inches(0.5),
                   ["ISM 自加 haplotag 邏輯"], font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(4.7),
                   ["候選：ISM 內建 read-haplotype 計算",
                    "",
                    "× 介面割裂：ISM ↔ longphase 重複",
                    "× ISM F1 = 0.0124 (TO germline)",
                    "× 對 TO germline FP 無修復力",
                    "× 違反單一職責原則",
                    "",
                    "結論：不採用"],
                   font_size=14, color=C_TEXT)
    add_text_block(s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.45),
                   ["最終選擇：在 longphase-to-mod 內 surgical 修補三函式"],
                   font_size=15, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_footer(s, 19)
    set_speaker_note(s,
        "在動 V5 之前評估了兩個替代方案，都排除了。"
        "(1) 替換 LongPhase 為 WhatsHap 或 HapCUT2：風險高 — 這兩個工具在 TO 模式下的行為都未驗證，"
        "外部依賴增加、介面變更會影響整個 pipeline，而且 self-phasing 是機制問題，換工具不一定解。"
        "(2) ISM 自己加 haplotag 邏輯：違反單一職責 — phasing 和 haplotag 該由 longphase 處理，"
        "ISM 重複實作會造成介面割裂；而且 ISM F1 = 0.0124 對 TO germline FP 沒有修復能力，"
        "ISM 加再多特徵也無法取代上游修補。"
        "最終選擇：在 longphase-to-mod 內做 surgical 修補三個函式 — 修改範圍最小、風險最低。")
    return s


def slide_20_interface_zero(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "介面契約零變動")
    # Show diff-style code block
    add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.6),
             fill=RGBColor(0x2D, 0x2D, 0x2D), line=C_BORDER)
    add_text_block(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(0.4),
                   ["HaplotagProcess.h:66-68 (Baseline → V5 一字未變)"],
                   font_size=14, bold=True, color=RGBColor(0xCF, 0xE2, 0xF3))
    code_lines = [
        "  void getVote(...);",
        "  int  judgeHaplotype(...);",
        "  void countSNPHaplotype(...);",
    ]
    add_text_block(s, Inches(0.9), Inches(1.85), Inches(11.8), Inches(1.85),
                   code_lines,
                   font_size=18, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Stats
    add_rect(s, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.5),
             fill=C_LIGHT, line=C_TITLE_BG, line_w_pt=2.0)
    add_text_block(s, Inches(0.7), Inches(4.35), Inches(5.6), Inches(0.5),
                   ["修改集中度"], font_size=18, bold=True, color=C_TITLE_BG)
    add_text_block(s, Inches(0.7), Inches(4.95), Inches(5.6), Inches(2.0),
                   ["• 總修改：+68 / -36 行",
                    "• 集中於 3 函式 (getVote / judge / countSNP)",
                    "• 4 個 commit、無大規模重構",
                    "• 無新增 public 函式"],
                   font_size=14, color=C_TEXT)
    add_rect(s, Inches(7.0), Inches(4.2), Inches(6.0), Inches(2.5),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.0)
    add_text_block(s, Inches(7.2), Inches(4.35), Inches(5.6), Inches(0.5),
                   ["介面零變動的好處"], font_size=18, bold=True, color=C_GREEN)
    add_text_block(s, Inches(7.2), Inches(4.95), Inches(5.6), Inches(2.0),
                   ["• 上層 caller 不需重編譯",
                    "• ABI 相容 — drop-in replace",
                    "• 風險局限於 HaplotagProcess.cpp",
                    "• 回退僅需 git revert 4 commits"],
                   font_size=14, color=C_TEXT)
    add_footer(s, 20)
    set_speaker_note(s,
        "V5 改動的一個關鍵特性：介面契約零變動。"
        "HaplotagProcess.h:66-68 三個 method signature 從 Baseline 到 V5 一字未變 — "
        "getVote、judgeHaplotype、countSNPHaplotype 的 public 介面完全不變。"
        "總修改量是 +68/-36 行，集中於這 3 個函式，4 個 commit。"
        "好處是：上層 caller 不需重編譯，ABI 相容，可以 drop-in replace；"
        "風險局限於 HaplotagProcess.cpp 內部；回退只需 git revert 4 個 commits 即可。"
        "這是 surgical fix 的典範 — 修改最小、可控、可回退。")
    return s


def slide_21_layer15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "getVote 兩層 → 三層")
    fit_image_within(s, str(FIG_DIR / "fig16_v5_threelayer_logic.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["V3F 兩層 (germline first)  →  V5 插入 Layer 1.5 「else if (somaticHP1>0 || somaticHP2>0)」"],
                   font_size=13, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 21)
    set_speaker_note(s,
        "V5 對 getVote() 的具體改動 — 從兩層升級為三層投票。"
        "V3F 邏輯：if (germlineHP1 > 0 || germlineHP2 > 0) {依 germline} else { 0 (回 fallback) }"
        "V5 邏輯：if (germline > 0) {依 germline} else if (somaticHP1 > 0 || somaticHP2 > 0) {依 somatic}"
        "       else { 0 }"
        "Layer 1.5 是純分支，不破壞 germline first 的核心原則 — 只在 germline 完全無票時才啟用 somatic。"
        "效果：HP:i:33 (ambiguous) reads 減 54%、AMB% 從 17.5% 降到 8.0%。"
        "副作用：可能讓某些 borderline somatic 從 33 變成 11/21，但這是預期行為 — "
        "Layer 1.5 的 somatic 投票本來就是要把 ambiguous 解到具體 hap。")
    return s


def slide_22_hp_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "HP tag 對應表（V5 final）")
    # Header
    cy = Inches(1.3)
    rw = Inches(11.0)
    rx = Inches(1.15)
    add_rect(s, rx, cy, Inches(3.8), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, rx, cy + Inches(0.08), Inches(3.8), Inches(0.4),
                   ["germlineResult"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, rx + Inches(3.8), cy, Inches(3.8), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, rx + Inches(3.8), cy + Inches(0.08), Inches(3.8), Inches(0.4),
                   ["somaticTotal > 0"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, rx + Inches(7.6), cy, Inches(3.4), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, rx + Inches(7.6), cy + Inches(0.08), Inches(3.4), Inches(0.4),
                   ["hpResult (HP:i)"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    cy += Inches(0.6)
    rows = [
        ("0", "False (germline only / unphased)", "0  (unphased)",       C_GRAY),
        ("1", "False",                            "1  (germline HP1)",  C_TITLE_BG),
        ("2", "False",                            "2  (germline HP2)",  C_TITLE_BG),
        ("1", "True",                             "11  (somatic on HP1)", C_BLUE),
        ("2", "True",                             "21  (somatic on HP2)", C_AMBER),
        ("0", "True",                             "33  (somatic ambiguous)", C_ACCENT),
    ]
    for gr, st, hp, color in rows:
        add_rect(s, rx, cy, Inches(3.8), Inches(0.55), fill=C_LIGHT, line=C_BORDER)
        add_text_block(s, rx, cy + Inches(0.1), Inches(3.8), Inches(0.4),
                       [gr], font_size=15, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_rect(s, rx + Inches(3.8), cy, Inches(3.8), Inches(0.55), fill=C_LIGHT, line=C_BORDER)
        add_text_block(s, rx + Inches(3.8), cy + Inches(0.1), Inches(3.8), Inches(0.4),
                       [st], font_size=12, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_rect(s, rx + Inches(7.6), cy, Inches(3.4), Inches(0.55), fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, rx + Inches(7.6), cy + Inches(0.1), Inches(3.4), Inches(0.4),
                       [hp], font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        cy += Inches(0.6)
    add_footer(s, 22)
    set_speaker_note(s,
        "V5 final 的 HP tag 對應表："
        "germlineResult = 0 且無 somatic → HP:i:0 (unphased)；"
        "germlineResult = 1 且無 somatic → HP:i:1 (germline HP1)；"
        "germlineResult = 2 且無 somatic → HP:i:2 (germline HP2)；"
        "germlineResult = 1 且有 somatic → HP:i:11 (somatic on HP1)；"
        "germlineResult = 2 且有 somatic → HP:i:21 (somatic on HP2)；"
        "germlineResult = 0 且有 somatic → HP:i:33 (somatic ambiguous)。"
        "這個對應表 V5 才完整 — V3F 因為缺 Layer 1.5，"
        "germlineResult = 0 + somaticTotal>0 的情況本應產生 HP:i:33，"
        "但因為 enum vs integer literal mismatch 永遠不出現。V5 修補後對應表是穩定的。")
    return s


def slide_23_ism_position(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "InterSubMod 的位置")
    # Two boxes: ISM and longphase-to-mod
    add_rect(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.5),
             fill=C_LIGHT, line=C_BLUE, line_w_pt=2.0)
    add_text_block(s, Inches(0.8), Inches(1.35), Inches(5.6), Inches(0.5),
                   ["InterSubMod (本 repo)"],
                   font_size=18, bold=True, color=C_BLUE)
    add_text_block(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(4.7),
                   ["路徑：/big7_disk/.../InterSubMod/",
                    "src/core/ 不含 HaplotagProcess",
                    "",
                    "職責：read-level ISM 計算",
                    "  • 讀 BAM HP tag",
                    "  • 距離矩陣、HPSig、HP_Ratio",
                    "",
                    "本 repo 對 self-phasing 修補：",
                    "  → 0 行 C++ 改動",
                    "  → 自動受惠新版 BAM"],
                   font_size=13, color=C_TEXT)

    add_rect(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.5),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.0)
    add_text_block(s, Inches(7.2), Inches(1.35), Inches(5.6), Inches(0.5),
                   ["longphase-to-mod (本地 fork)"],
                   font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(4.7),
                   ["路徑：/big7_disk/.../longphase-to-mod/",
                    "獨立 git repo",
                    "",
                    "職責：phasing + haplotag",
                    "  • HaplotagProcess.{h,cpp}",
                    "  • PhasingProcess.cpp",
                    "",
                    "本次修補對象：",
                    "  → 4 commits ×  3 函式",
                    "  → +68 / -36 行"],
                   font_size=13, color=C_TEXT)
    add_text_block(s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.45),
                   ["兩個獨立 repo · ISM 自動受惠新版 BAM"],
                   font_size=15, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_footer(s, 23)
    set_speaker_note(s,
        "再次釐清 InterSubMod 與 longphase-to-mod 的職責切分。"
        "左欄 InterSubMod：本 repo，路徑 /big7_disk/.../InterSubMod/；"
        "src/core/ 內不含 HaplotagProcess（這檔案在 longphase-to-mod 內），"
        "ISM 的職責是 read-level epigenetic 計算，讀 BAM HP tag 算距離矩陣、HPSig、HP_Ratio。"
        "本次 self-phasing 修補對 ISM repo 的 C++ 改動是 0 行，ISM 只是讀新版 BAM 自動受惠。"
        "右欄 longphase-to-mod：獨立 git repo，職責是 phasing 和 haplotag，"
        "包含 HaplotagProcess.{h,cpp} 和 PhasingProcess.cpp；"
        "本次修補對象就是它，4 個 commits 修 3 個函式，+68/-36 行。")
    return s


def slide_24_caveat(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "V5 working tree 未 commit caveat")
    # Big warning box
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(2.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.5)
    add_text_block(s, Inches(0.8), Inches(1.35), Inches(11.6), Inches(0.5),
                   ["⚠ 警告 · V5 = 380e8d2 + 2 塊 working-tree 修改"],
                   font_size=20, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(0.8), Inches(1.95), Inches(11.6), Inches(1.6),
                   ["• Layer 1.5 somatic fallback (getVote)",
                    "• countSNPHaplotype() alt guard 對稱化",
                    "• 兩塊修改皆未 commit、未 tag、無 SHA",
                    "• 重現性風險：誰拉到這個版本要靠手動 patch"],
                   font_size=14, color=C_TEXT)

    # Recommendation
    add_rect(s, Inches(0.6), Inches(3.85), Inches(12.0), Inches(2.5),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.5)
    add_text_block(s, Inches(0.8), Inches(4.0), Inches(11.6), Inches(0.5),
                   ["建議：切 2 個獨立 commits (後續 F1)"],
                   font_size=20, bold=True, color=C_GREEN)
    add_text_block(s, Inches(0.8), Inches(4.6), Inches(11.6), Inches(1.7),
                   ["• Commit A：getVote() Layer 1.5 somatic fallback",
                    "• Commit B：countSNPHaplotype() alt guard 對稱化",
                    "• 切完 commits 後 tag v5.0",
                    "• 在 manifest.yaml 內記 haplotag_version = v5.0"],
                   font_size=14, color=C_TEXT)
    add_footer(s, 24)
    set_speaker_note(s,
        "V5 的一個重要 caveat：當前 V5 = 380e8d2 + 兩塊 working tree 修改，**尚未 commit**。"
        "兩塊修改分別是：(1) getVote() 的 Layer 1.5 somatic fallback，"
        "(2) countSNPHaplotype() 的 alt guard 對稱化（補對稱缺失）。"
        "兩塊都還沒 commit、沒 tag、沒 SHA — 重現性風險在於：誰 clone 這個 repo 都要手動 patch。"
        "建議在後續工作 F1 切兩個獨立 commits："
        "Commit A getVote Layer 1.5、Commit B countSNPHaplotype alt guard，"
        "完成後 tag v5.0，並在 InterSubMod 的 manifest.yaml 內記錄 haplotag_version = v5.0，"
        "讓下游分析可追溯。")
    return s


def slide_25_sanity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4 項硬性 sanity check")
    rows = [
        ("Δ-consistency 守恆律", "tag 移轉量平衡 (in = out)", "PASS", C_GREEN),
        ("Germline 不變",         "HP:i:1 / HP:i:2 reads 與 BL 完全相同", "PASS", C_GREEN),
        ("Layer 1.5 期望 = 1",    "純 somatic 位點 fallback 觸發次數 = 1", "PASS", C_GREEN),
        ("無 germline → HP33",    "germlineResult≠0 不會輸出 HP:i:33", "PASS", C_GREEN),
    ]
    cy = Inches(1.3)
    rh = Inches(1.0)
    for name, desc, verdict, color in rows:
        add_rect(s, Inches(0.6), cy, Inches(8.4), rh, fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(0.8), cy + Inches(0.1), Inches(8.0), Inches(0.45),
                       [name], font_size=17, bold=True, color=C_TEXT)
        add_text_block(s, Inches(0.8), cy + Inches(0.55), Inches(8.0), Inches(0.4),
                       [desc], font_size=12, color=C_GRAY)
        add_rect(s, Inches(9.2), cy, Inches(3.8), rh, fill=color, line=color)
        add_text_block(s, Inches(9.2), cy + Inches(0.25), Inches(3.8), Inches(0.5),
                       [verdict], font_size=24, bold=True,
                       color=C_TITLE_FG, align=PP_ALIGN.CENTER)
        cy += rh + Inches(0.1)
    add_text_block(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.5),
                   ["總計 15 / 15 PASS  ·  0 violation"],
                   font_size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_footer(s, 25)
    set_speaker_note(s,
        "V5 修補完成後做了 4 項硬性 sanity check："
        "(1) Δ-consistency 守恆律：tag 在 BL → V5 之間移轉量必須平衡（in = out），PASS；"
        "(2) Germline 不變：HP:i:1 和 HP:i:2 reads 數應與 baseline 完全相同（V5 不應動 germline），PASS；"
        "(3) Layer 1.5 期望次數 = 1：純 somatic 位點上，Layer 1.5 fallback 應該 trigger 一次，PASS；"
        "(4) 無 germline → HP33：germlineResult ≠ 0 的情況下不應輸出 HP:i:33，PASS。"
        "15 個位點 × 4 項 check = 15/15 PASS、0 violation。"
        "這是對 V5 改動的硬性內部驗證，獨立於下游 ISM / F1 結果。")
    return s


def slide_26_quant_compare(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Baseline / V3F / V5 量化對比")
    fit_image_within(s, str(FIG_DIR / "fig18_concordance_amb_f1.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["AMB% 17.5→8.0%  ·  HP:i:33 reads -54%  ·  clean PS concordance V5 90.5% vs BL 82.2% (+8.3 pp)"],
                   font_size=13, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 26)
    set_speaker_note(s,
        "Baseline / V3F / V5 三版本的核心量化對比："
        "(1) AMB% (HP:i:33 比例) 從 17.5% 降到 8.0% — V5 解了過半 ambiguous reads；"
        "(2) HP:i:33 reads 數從 239,679 降到 110,197，減 54%；"
        "(3) 全基因組 clean PS blocks paired GT concordance：V5 = 90.5% vs Baseline = 82.2%，"
        "提升 +8.3 pp — 這是 V5 真實價值的核心指標。"
        "(4) 15 sites cherry-picked 上 clean PS concordance：V5 = 88.2% vs BL = 74.9%，提升 +13.3 pp。"
        "這四個數字從不同角度都驗證了 V5 是嚴格優於 baseline 的修補。")
    return s


def slide_27_v5max1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "IGV 4-BAM：V5max1 戲劇化")
    fit_image_within(s, str(FIG_DIR / "C_V5max1_chr19_4639528.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.5))
    add_text_block(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.45),
                   ["chr19:4639528  ·  V3F 紫色 HP33 群 (39 reads) → V5 全部正確重分配為 HP11"],
                   font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 27)
    set_speaker_note(s,
        "這張是 V5 修補的最戲劇化展示。chr19:4639528 是 V5max1 案例："
        "在 V3-Fixed panel，可以看到一大群紫色 reads 標為 HP:i:33 (ambiguous)，39 條。"
        "V5 修補後，這 39 條全部正確重分配為 HP:i:11 (somatic on HP1)。"
        "守恆律驗證：V3F 39 條 HP33 → V5 39 條 HP11，in = out 完全平衡，PASS。"
        "這是 Layer 1.5 somatic fallback 的具體效果 — "
        "原本因為 enum mismatch 卡在 ambiguous bucket 的 reads，現在被正確歸到具體 hap。"
        "看 IGV 4-BAM 並列圖，從 baseline、V2b、V3F、V5 的視覺差異很明顯。")
    return s


def slide_28_sp1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "IGV 4-BAM：SP1 翻轉")
    fit_image_within(s, str(FIG_DIR / "D_SP1_chr19_17565944.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.5))
    add_text_block(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.45),
                   ["chr19:17565944  ·  baseline HP1 主導 → V2b/V3F/V5 翻轉為 HP2  ·  3/3 SP-extreme 一致"],
                   font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 28)
    set_speaker_note(s,
        "這張是 SP-extreme 翻轉案例。chr19:17565944："
        "Baseline 狀態 HP1 主導（113:0 完全失衡）；"
        "V2b、V3F、V5 三版本全部一致地翻轉為 HP2 主導；"
        "Paired tumor BAM 確認 HP2 才是真實方向 — V5 與 paired ground truth 一致。"
        "這個翻轉是 self-phasing 修補的核心效果之一 — "
        "原本因為 self-phasing 把 reads 全部歸到錯邊（HP1），"
        "修補後 phasing 正確以 germline 為 anchor，reads 回到正確的 hap (HP2)。"
        "3 個 SP-extreme 位點（SP1/SP2/SP3）全部一致翻轉，是跨位點一致性證據。")
    return s


def slide_29_f1_clarify(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "F1 為什麼幾乎不變？（重要釐清）")
    img = make_f1_decomposition()
    fit_image_within(s, str(img), Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.5))
    add_text_block(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.4),
                   ["真實價值在 read-level concordance +8.3 pp，不在 SEQC2 F1"],
                   font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 29)
    set_speaker_note(s,
        "這張是**最關鍵的概念修正** — 為什麼 V5 修完 SEQC2 F1 幾乎不變？"
        "ClairS-TO raw F1 = 0.7166 對所有版本（Baseline、V3F、V5）**完全相同**，"
        "原因是 V5 不改 caller — V5 改的是 BAM HP tag 編碼，與 caller 輸出的 VCF 無關。"
        "F1 變動只可能來自下游 ISM SuggestFilter，但 ISM TO F1 增益本來就只有 0.0124，"
        "V5 vs Baseline 最終 F1 差 -0.0003，是**統計噪音級別**。"
        "**結論**：F1 不是衡量 V5 真實價值的指標。"
        "V5 真實價值在 read-level tag 品質 — clean PS blocks paired GT concordance +8.3 pp，"
        "AMB% 17.5→8.0%，HP:i:33 -54%。這些是 read-level 指標，無法用 F1 反映。"
        "教授可能會問「F1 沒改善代表沒用嗎？」這張就是回答。")
    return s


def slide_30_3tier(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "ISM 影響 3-tier 分類")
    fit_image_within(s, str(FIG_DIR / "03_self_phasing_impact.png"),
                      Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.55),
                   ["嚴重 38% (必重跑)  ·  中度 7% (QualityScore 等)  ·  不受影響 55% (Pairwise / 甲基化)"],
                   font_size=13, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 30)
    set_speaker_note(s,
        "V5 修補對下游 ISM 特徵的影響分為 3 tier："
        "嚴重影響 38%：HP-依賴特徵 — HP_Ratio、HPSig、HPFineNGroups、HP-stratified pairwise dist，"
        "這些必須在 V5 BAM 上重跑才能信任。"
        "中度影響 7%：QualityScore 等部分依賴 HP 分布的特徵，重跑值會變但結論方向通常不變。"
        "不受影響 55%：Pairwise distance（無 HP 分組）、Allele-level 特徵、Caller VCF 衍生特徵、"
        "甲基化矩陣本身 — 這些不依賴 HP tag，V5 前後值完全相同。"
        "所以 V5 的下游影響不是「全部重跑」也不是「都不影響」，"
        "而是有清楚的 tier 劃分，重跑工作可以聚焦在嚴重 + 中度的 45%。")
    return s


def slide_31_corrections(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "修正了哪些舊說法")
    rows = [
        ("HPFineNGroups = methylation bimodality",
         "= phasing bucket (V5 後仍有訊號是 phasing 訊號)"),
        ("HP:i:21 必然是當前位點 ALT",
         "不必然：HP:i:21 是 read-level 標記，與當前位點 ALT 無強對應"),
        ("Self-phasing 污染 LOH.bed",
         "不污染：LOH.bed 是 region-level，V5 前後 Jaccard = 1.0"),
        ("V5 修補 = caller F1 改善",
         "不等於：V5 改 BAM HP tag、不改 caller；F1 變動 -0.0003 噪音"),
    ]
    cy = Inches(1.2)
    rh = Inches(1.25)
    add_rect(s, Inches(0.5), cy, Inches(6.0), Inches(0.55), fill=C_ACCENT)
    add_text_block(s, Inches(0.5), cy + Inches(0.08), Inches(6.0), Inches(0.4),
                   ["× 舊說法"], font_size=15, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.55), cy, Inches(6.45), Inches(0.55), fill=C_GREEN)
    add_text_block(s, Inches(6.55), cy + Inches(0.08), Inches(6.45), Inches(0.4),
                   ["✓ 修正後"], font_size=15, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    cy += Inches(0.6)
    for old, new in rows:
        add_rect(s, Inches(0.5), cy, Inches(6.0), rh, fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT)
        add_text_block(s, Inches(0.65), cy + Inches(0.2), Inches(5.7), rh - Inches(0.25),
                       [old], font_size=13, color=C_TEXT)
        add_rect(s, Inches(6.55), cy, Inches(6.45), rh, fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN)
        add_text_block(s, Inches(6.7), cy + Inches(0.2), Inches(6.15), rh - Inches(0.25),
                       [new], font_size=13, color=C_TEXT)
        cy += rh + Inches(0.05)
    add_footer(s, 31)
    set_speaker_note(s,
        "V5 修補也順便修正了我們過去的 4 個說法。"
        "(1) HPFineNGroups 不是 methylation bimodality — 是 phasing bucket，V5 後若仍有訊號是 phasing 訊號。"
        "(2) HP:i:21 不必然是當前位點 ALT — HP:i:21 是 read-level 標記，與當前位點 ALT 無強對應；"
        "讀解時不能假設 HP:i:21 = ALT。"
        "(3) Self-phasing 不污染 LOH.bed — 之前一度懷疑 LOH.bed 也受影響，"
        "驗證後 V5 前後 Jaccard = 1.0 完全不變。"
        "(4) V5 修補不等於 caller F1 改善 — V5 不改 caller，F1 變動 -0.0003 是噪音；"
        "S29 已詳細釐清這點。"
        "這些都是過去文件中常出現的說法，需要在新一輪報告中明確修正。")
    return s


def slide_32_limits(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "已知限制 R1 - R9")
    items = [
        ("R1", "cnLOH 區未獨立評估"),
        ("R2", "AF > 0.9 邊界行為未驗證"),
        ("R3", "15 sites 為 cherry-picked (非隨機)"),
        ("R4", "Confidence threshold 0.6 未直接驗證"),
        ("R5", "V5 working tree 未 commit"),
        ("R6", "7 樣本擴展未做 (僅 HCC1395 5kHz)"),
        ("R7", "L4 部分指標 V5 略遜 V3F"),
        ("R8", "problem PS 不適用 read-level"),
        ("R9", "paired ground truth 自身用 LongPhase"),
    ]
    cy = Inches(1.15)
    cx_left = Inches(0.5)
    cx_right = Inches(6.85)
    cw = Inches(6.15)
    rh = Inches(0.6)
    for i, (rid, desc) in enumerate(items):
        x = cx_left if i < 5 else cx_right
        if i == 5:
            cy = Inches(1.15)
        add_rect(s, x, cy, cw, rh, fill=C_LIGHT, line=C_AMBER)
        add_text_block(s, x + Inches(0.1), cy + Inches(0.1), Inches(0.7), Inches(0.4),
                       [rid], font_size=14, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
        add_text_block(s, x + Inches(0.85), cy + Inches(0.1), cw - Inches(0.95), Inches(0.4),
                       [desc], font_size=12, color=C_TEXT)
        cy += rh + Inches(0.08)
    add_text_block(s, Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.5),
                   ["9 條限制中 R1 / R5 / R6 為主要 open issue"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 32)
    set_speaker_note(s,
        "報告 caveat 部分整理為 9 條已知限制 R1-R9。"
        "R1 cnLOH 區（雙親同源）未獨立評估，可能有不同行為。"
        "R2 AF > 0.9 邊界行為未驗證，可能 phasing 退化。"
        "R3 15 個驗證位點是 cherry-picked，非隨機抽樣 — "
        "後續 F8 會做 100 隨機位點 cross-validation。"
        "R4 Confidence threshold 0.6 是經驗值，未直接驗證最佳值。"
        "R5 V5 working tree 未 commit (S24 已強調)。"
        "R6 只在 HCC1395 5kHz 驗證，7 樣本擴展未做 (F3)。"
        "R7 L4 部分次要指標 V5 略遜 V3F，不影響主結論。"
        "R8 problem PS 是 caller-level 指標、不適用 read-level 比較。"
        "R9 paired ground truth 自身用 LongPhase 產生，存在 self-validation 風險 — F7 trio-phased 是補強。"
        "主要 open issue 是 R1、R5、R6。")
    return s


def slide_33_next_steps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "後續工作 F1 - F8")
    rows = [
        ("F1", "Commit V5 working tree (Layer 1.5 + alt guard)", "高"),
        ("F2", "Vote log 直接驗證 (內部 unit test)",                 "中"),
        ("F3", "7 樣本擴展驗證 (跨樣本一致性)",                       "高"),
        ("F4", "master × flag 重跑 HPFineNGroups",                   "中"),
        ("F5", "Manifest 加 haplotag_version 欄位",                  "低"),
        ("F6", "cnLOH 區獨立評估",                                    "中"),
        ("F7", "Trio-phased 第二 ground truth",                       "中"),
        ("F8", "100 隨機位點 cross-validate",                          "高"),
    ]
    cy = Inches(1.15)
    rh = Inches(0.65)
    add_rect(s, Inches(0.5), cy, Inches(0.9), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, Inches(0.5), cy + Inches(0.08), Inches(0.9), Inches(0.4),
                   ["ID"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.45), cy, Inches(9.5), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, Inches(1.45), cy + Inches(0.08), Inches(9.5), Inches(0.4),
                   ["項目"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(11.0), cy, Inches(2.0), Inches(0.55), fill=C_TITLE_BG)
    add_text_block(s, Inches(11.0), cy + Inches(0.08), Inches(2.0), Inches(0.4),
                   ["優先"], font_size=14, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    cy += Inches(0.6)
    for fid, desc, prio in rows:
        pcolor = C_ACCENT if prio == "高" else (C_AMBER if prio == "中" else C_GRAY)
        add_rect(s, Inches(0.5), cy, Inches(0.9), rh, fill=C_LIGHT, line=C_BORDER)
        add_text_block(s, Inches(0.5), cy + Inches(0.13), Inches(0.9), Inches(0.4),
                       [fid], font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.45), cy, Inches(9.5), rh, fill=C_LIGHT, line=C_BORDER)
        add_text_block(s, Inches(1.55), cy + Inches(0.15), Inches(9.4), Inches(0.4),
                       [desc], font_size=12, color=C_TEXT)
        add_rect(s, Inches(11.0), cy, Inches(2.0), rh, fill=C_LIGHT, line=pcolor, line_w_pt=2.0)
        add_text_block(s, Inches(11.0), cy + Inches(0.13), Inches(2.0), Inches(0.4),
                       [prio], font_size=14, bold=True, color=pcolor, align=PP_ALIGN.CENTER)
        cy += rh + Inches(0.05)
    add_footer(s, 33)
    set_speaker_note(s,
        "後續工作 8 條 F1-F8："
        "F1 (高) commit V5 working tree 兩塊修改成獨立 commits — 對應 S24 的建議。"
        "F2 (中) vote log 直接驗證，加 unit test 在 longphase-to-mod 內。"
        "F3 (高) 7 樣本擴展，目前只在 HCC1395 5kHz 完整驗證。"
        "F4 (中) master × flag 重跑 HPFineNGroups — "
        "因為過去結論用了未修補的 BAM，需要重驗證。"
        "F5 (低) manifest 加 haplotag_version 欄位讓下游可追溯。"
        "F6 (中) cnLOH 區獨立評估。"
        "F7 (中) trio-phased 作為第二 ground truth，補強 R9。"
        "F8 (高) 100 隨機位點 cross-validate，補強 R3 的 cherry-pick 風險。"
        "高優先三項是 F1 / F3 / F8。")
    return s


def slide_34_tldr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_TITLE_BG)
    add_text_block(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.7),
                   ["TL;DR"], font_size=44, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    img = make_tldr_banner()
    fit_image_within(s, str(img), Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.4))
    add_text_block(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4),
                   ["Q & A 歡迎隨時打斷 · 詳細數據見 source_materials/"],
                   font_size=14, bold=True, color=RGBColor(0xCF, 0xE2, 0xF3), align=PP_ALIGN.CENTER)
    set_speaker_note(s,
        "TL;DR 一句話結論："
        "Self-phasing 在 LongPhase-TO 階段以 4-commit 漸進修補；"
        "InterSubMod 是下游消費者，本 repo 無 C++ 改動；"
        "V5 真實價值在 read-level tag 品質 (+8.3 pp clean PS concordance)，SEQC2 F1 不變是預期行為。"
        "歡迎提問。詳細數據在 source_materials 三份報告：01 IGV visual audit、02 V5 audit、03 技術報告。"
        "後續工作清單以 F1 (commit) 與 F3 (7 樣本擴展) 為高優先。"
        "Thank you.")
    return s


# ---------------------------------------------------------------------------
# Main: Assemble all 34 slides
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_cover,
        slide_02_pipeline,
        slide_03_phasing,
        slide_04_three_layer,
        slide_05_hp_tag,
        slide_06_ism_role,
        slide_07_self_phasing_def,
        slide_08_af03_walkthrough,
        slide_09_to_severity,
        slide_10_three_bug,
        slide_11_genome_evidence,
        slide_12_site_evidence,
        slide_13_consistency,
        slide_14_loh_two_layer,
        slide_15_fix_location,
        slide_16_4commit,
        slide_17_4bug_table,
        slide_18_three_layer_logic,
        slide_19_alternatives,
        slide_20_interface_zero,
        slide_21_layer15,
        slide_22_hp_table,
        slide_23_ism_position,
        slide_24_caveat,
        slide_25_sanity,
        slide_26_quant_compare,
        slide_27_v5max1,
        slide_28_sp1,
        slide_29_f1_clarify,
        slide_30_3tier,
        slide_31_corrections,
        slide_32_limits,
        slide_33_next_steps,
        slide_34_tldr,
    ]

    assert len(builders) == 34, f"Expected 34 slides, got {len(builders)}"

    for i, fn in enumerate(builders, 1):
        print(f"  [{i:2d}/34] building {fn.__name__} ...")
        fn(prs)

    prs.save(str(OUTPUT_PPTX))
    print(f"\n[OK] PPTX saved to {OUTPUT_PPTX}")
    print(f"     Slides: {len(prs.slides)}")
    return OUTPUT_PPTX


if __name__ == "__main__":
    build()
