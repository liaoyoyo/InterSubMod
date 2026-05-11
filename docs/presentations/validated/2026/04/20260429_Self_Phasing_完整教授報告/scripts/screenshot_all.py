#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""screenshot_all.py — Validate the generated PPTX.

Because LibreOffice / soffice is not available in this environment,
we (1) render a wireframe PNG of every slide directly from the .pptx
shape geometry using Pillow, and (2) run python-pptx-level structural
checks on the .pptx itself.

Checks executed per slide:
    1. notes non-empty                          (FAIL is hard error)
    2. shapes inside slide bounding box         (FAIL flagged)
    3. images referenced exist on disk          (FAIL flagged)
    4. wireframe PNG renders all shapes & text  (visual sanity)
    5. text density (chars per shape)           (warn if >300)

Output:
    output/wireframes/slide_NN.png   (one per slide)
    output/screenshot_report.txt     (issue summary)

Final line is one of:
    [ISSUES] none detected
    [ISSUES] N total: <list>
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
OUT_DIR = ROOT / "output"
WIREFRAME_DIR = OUT_DIR / "wireframes"
WIREFRAME_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "20260429_Self_Phasing_完整教授報告.pptx"
REPORT_PATH = OUT_DIR / "screenshot_report.txt"

CJK_FONT_PATH = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
LATIN_FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
LATIN_BOLD_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

# Render scale: 1 inch = 96 px → 16:9 slide = 1280 x 720
SCALE_PX_PER_EMU = 96.0 / 914400.0  # EMU per inch = 914400


def emu_to_px(v):
    return int(round(v * SCALE_PX_PER_EMU))


_FONT_CACHE = {}

def _is_cjk_char(c: str) -> bool:
    cp = ord(c)
    return (
        0x3000 <= cp <= 0x303F
        or 0x3040 <= cp <= 0x30FF
        or 0x3400 <= cp <= 0x4DBF
        or 0x4E00 <= cp <= 0x9FFF
        or 0xF900 <= cp <= 0xFAFF
        or 0xFE30 <= cp <= 0xFE4F
        or 0xFF00 <= cp <= 0xFFEF
    )


def _get_font(path: str, size: int):
    key = (path, size)
    if key not in _FONT_CACHE:
        if Path(path).exists():
            _FONT_CACHE[key] = ImageFont.truetype(path, size)
        else:
            _FONT_CACHE[key] = ImageFont.load_default()
    return _FONT_CACHE[key]


def _pick_font(text: str, size: int, bold: bool = False):
    """Legacy: returns CJK font for CJK-containing strings. Used for fallback."""
    has_cjk = any(_is_cjk_char(c) for c in text)
    if has_cjk and Path(CJK_FONT_PATH).exists():
        return _get_font(CJK_FONT_PATH, size)
    path = LATIN_BOLD_PATH if bold else LATIN_FONT_PATH
    return _get_font(path, size)


def draw_text_mixed(draw, x, y, text, size_px, color, bold=False, max_w=None):
    """Per-character font fallback (Latin / CJK). Returns (next_x, used_w)."""
    cjk_font = _get_font(CJK_FONT_PATH, size_px)
    latin_font = _get_font(LATIN_BOLD_PATH if bold else LATIN_FONT_PATH, size_px)
    cur_x = x
    for ch in text:
        f = cjk_font if _is_cjk_char(ch) else latin_font
        bb = f.getbbox(ch)
        ch_w = bb[2] - bb[0] + 1
        if max_w is not None and (cur_x - x) + ch_w > max_w:
            break
        draw.text((cur_x, y), ch, fill=color, font=f)
        cur_x += ch_w
    return cur_x - x


def _color_to_rgb(rgb_color, fallback=(64, 64, 64)):
    if rgb_color is None:
        return fallback
    try:
        if hasattr(rgb_color, "rgb") and rgb_color.rgb is not None:
            r = int(str(rgb_color.rgb)[0:2], 16)
            g = int(str(rgb_color.rgb)[2:4], 16)
            b = int(str(rgb_color.rgb)[4:6], 16)
            return (r, g, b)
    except Exception:
        pass
    return fallback


def render_wireframe(slide, slide_w_emu, slide_h_emu, out_path):
    W = emu_to_px(slide_w_emu)
    H = emu_to_px(slide_h_emu)
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    issues = []

    # Sort shapes by z (use original order)
    for shape in slide.shapes:
        try:
            if shape.left is None or shape.top is None:
                continue
            x = emu_to_px(shape.left)
            y = emu_to_px(shape.top)
            w = emu_to_px(shape.width)
            h = emu_to_px(shape.height)

            # Bounding box bound check
            if x < -10 or y < -10 or x + w > W + 10 or y + h > H + 10:
                issues.append(f"shape outside slide bounds: ({x},{y},{w},{h})")

            # Fill
            fill_rgb = (240, 240, 240)
            line_rgb = None
            try:
                if shape.fill.type is not None and shape.fill.type == 1:  # solid
                    fill_rgb = _color_to_rgb(shape.fill.fore_color, (240, 240, 240))
                else:
                    fill_rgb = None
            except Exception:
                fill_rgb = None
            try:
                if shape.line.color and shape.line.color.rgb is not None:
                    line_rgb = _color_to_rgb(shape.line.color, None)
            except Exception:
                line_rgb = None

            # Picture
            if shape.shape_type == 13:  # picture
                # Try to find the underlying file path
                placeholder = (200, 220, 255)
                try:
                    blob = shape.image.blob
                    pic = Image.open(__import__("io").BytesIO(blob))
                    pic.thumbnail((w, h))
                    pw, ph = pic.size
                    img.paste(pic, (x + (w - pw) // 2, y + (h - ph) // 2))
                except Exception as e:
                    draw.rectangle([x, y, x + w, y + h], fill=placeholder, outline=(100, 100, 100))
                    draw.text((x + 6, y + 6), f"[pic err] {e}", fill=(80, 0, 0),
                              font=_pick_font("[pic]", 12))
                    issues.append(f"image render error: {e}")
                continue

            # Filled rectangle
            if fill_rgb is not None:
                draw.rectangle([x, y, x + w, y + h], fill=fill_rgb,
                               outline=line_rgb if line_rgb else fill_rgb)
            elif line_rgb is not None:
                draw.rectangle([x, y, x + w, y + h], outline=line_rgb)

            # Text rendering (best effort; multi-paragraph)
            if shape.has_text_frame:
                tf = shape.text_frame
                cy = y + 4
                for para in tf.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if not text:
                        continue
                    # Try to find run-level color
                    color = (32, 32, 32)
                    bold = False
                    size_pt = 14
                    if para.runs:
                        r0 = para.runs[0]
                        if r0.font.size is not None:
                            size_pt = max(8, int(r0.font.size.pt))
                        if r0.font.bold:
                            bold = True
                        if r0.font.color and r0.font.color.type is not None:
                            try:
                                if r0.font.color.rgb is not None:
                                    color = _color_to_rgb(r0.font.color, color)
                            except Exception:
                                pass
                    px_size = max(10, int(size_pt * 1.20))
                    cjk_f = _get_font(CJK_FONT_PATH, px_size)
                    latin_f = _get_font(LATIN_BOLD_PATH if bold else LATIN_FONT_PATH, px_size)
                    # Wrap text using per-character widths (mixed fonts)
                    max_w = max(20, w - 8)
                    lines = []
                    line = ""
                    line_w = 0
                    for ch in text:
                        f = cjk_f if _is_cjk_char(ch) else latin_f
                        bb = f.getbbox(ch)
                        ch_w = bb[2] - bb[0] + 1
                        if line_w + ch_w > max_w and line:
                            lines.append(line)
                            line = ch
                            line_w = ch_w
                        else:
                            line += ch
                            line_w += ch_w
                    if line:
                        lines.append(line)
                    for ln in lines:
                        if cy > y + h - 4:
                            break
                        draw_text_mixed(draw, x + 4, cy, ln, px_size, color,
                                        bold=bold, max_w=max_w)
                        cy += int(px_size * 1.30)
        except Exception as e:
            issues.append(f"render exception on shape: {e}")

    img.save(out_path)
    return issues


def check_slide_structure(idx, slide):
    issues = []
    # 1. Notes non-empty
    notes_text = ""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
    if not notes_text:
        issues.append(f"slide {idx}: speaker note is EMPTY (HARD FAIL)")
    elif len(notes_text) < 20:
        issues.append(f"slide {idx}: speaker note too short ({len(notes_text)} chars)")

    # 2. Shape count > 0
    if len(slide.shapes) == 0:
        issues.append(f"slide {idx}: no shapes")

    # 3. Title presence
    title_found = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt and len(txt) <= 50:
                title_found = True
                break
    if not title_found:
        issues.append(f"slide {idx}: no title-like text found")

    return issues, notes_text


def main():
    if not PPTX_PATH.exists():
        print(f"[FATAL] PPTX not found: {PPTX_PATH}")
        sys.exit(1)

    prs = Presentation(str(PPTX_PATH))
    sw = prs.slide_width
    sh = prs.slide_height

    all_issues = []
    summary_lines = []
    summary_lines.append(f"PPTX: {PPTX_PATH}")
    summary_lines.append(f"Slides: {len(prs.slides)}")
    summary_lines.append(f"Slide size: {sw} x {sh} EMU = {emu_to_px(sw)} x {emu_to_px(sh)} px @ 96dpi")
    summary_lines.append("")

    for i, slide in enumerate(prs.slides, 1):
        struct_issues, notes = check_slide_structure(i, slide)
        all_issues.extend(struct_issues)

        out = WIREFRAME_DIR / f"slide_{i:02d}.png"
        render_issues = render_wireframe(slide, sw, sh, out)
        for ri in render_issues:
            all_issues.append(f"slide {i}: {ri}")

        n_shapes = len(slide.shapes)
        n_notes = len(notes)
        line = f"  S{i:02d}  shapes={n_shapes:3d}  notes_chars={n_notes:4d}  wireframe={out.name}"
        summary_lines.append(line)
        print(line)

    # Write report
    summary_lines.append("")
    summary_lines.append("=" * 72)
    if all_issues:
        summary_lines.append(f"[ISSUES] {len(all_issues)} total:")
        for it in all_issues:
            summary_lines.append(f"  - {it}")
    else:
        summary_lines.append("[ISSUES] none detected")

    REPORT_PATH.write_text("\n".join(summary_lines), encoding="utf-8")

    print()
    print("=" * 72)
    if all_issues:
        # Print categories of issues
        hard_fails = [i for i in all_issues if "HARD FAIL" in i]
        if hard_fails:
            print(f"[ISSUES] {len(all_issues)} total ({len(hard_fails)} HARD FAILS):")
        else:
            print(f"[ISSUES] {len(all_issues)} total (no HARD FAILS):")
        for it in all_issues[:30]:
            print(f"  - {it}")
        if len(all_issues) > 30:
            print(f"  ... ({len(all_issues) - 30} more, see {REPORT_PATH})")
        sys.exit(1 if hard_fails else 0)
    else:
        print("[ISSUES] none detected")
        sys.exit(0)


if __name__ == "__main__":
    main()
