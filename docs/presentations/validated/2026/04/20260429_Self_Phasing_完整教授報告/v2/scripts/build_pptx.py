#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_pptx.py — Self-Phasing 完整教授報告 v3 (24 slides, 16:9)

Generates output/20260429_Self_Phasing_完整教授報告_v2.pptx according to
00_storyboard_v2.md (24 slides v3). Each slide carries a speaker note >= 350
chars (S11/S12/S14 >= 500 chars including Q1/Q2 + Purity 0.95 root cause
integration).

v3 changes vs v2-26:
    - merge S5+S6 -> S4 (3-layer evidence); S11+S12 -> S10 (prerequisite);
      S22+S24 -> S21; S2 industry tree absorbed into S20
    - split original S8 -> S6 unlock + S7 LOH two layers; S17 -> S15 metrics
      + S16 timeline
    - new: S11 root-cause tree adds Purity 0.927 <= 0.95 trigger lock root
    - speaker note: S10 adds Purity 0.95 threshold; S12 adds mid-low purity
      defence + Q2 answer; S13 baseline column adds "somatic-first" red mark;
      S14 adds Q1 PON dual-path effectiveness
    - tightened ending: S24 adds Q&A invitation + thanks bottom strip

Strict rendering rules (see scripts/PPT_RENDERING_RULES.md):
    1. Per-character Latin + CJK font fallback via <a:rPr><a:latin/><a:ea/></a:rPr>.
    2. fit_image_within (equal ratio, never force both width+height).
    3. Speaker notes mandatory (>=350; S11/S12/S14 >=500 chars).
    4. Slide title <= 15 CJK chars; bullets <= 5 entries.
    5. No emojis (Droid Sans Fallback lacks them); use ASCII / box symbols.

Verbiage calibration (6 must-follow under v3):
    - V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold)
    - 同實驗室相鄰工作 (NOT 業界共識)
    - ISM 影響只列 count (29 / 14 / 42 / 共 85)
    - Util.h:20-25 (NOT 21-25)
    - +8.3 pp (clean PS blocks, 全基因組)
    - Baseline somatic-first 投票順序 (not 「有 HP11/HP21 即 somatic」)
    - Purity 0.927 <= 0.95 未觸發 Two-Pass -> 走三條路 (not 「baseline 誤判 purity」)
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FIG_DIR = ROOT / "figures"
OUT_DIR = ROOT / "output"
AUTO_FIG_DIR = OUT_DIR / "auto_figs"
OUT_DIR.mkdir(exist_ok=True)
AUTO_FIG_DIR.mkdir(exist_ok=True)

OUTPUT_PPTX = OUT_DIR / "20260429_Self_Phasing_完整教授報告_v2.pptx"
TOTAL_SLIDES = 24

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------
CJK_FONT_PATH = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
LATIN_FONT = "Arial"
CJK_FONT_NAME = "Droid Sans Fallback"

if Path(CJK_FONT_PATH).exists():
    fm.fontManager.addfont(CJK_FONT_PATH)
    _font = fm.FontProperties(fname=CJK_FONT_PATH).get_name()
    plt.rcParams.update({
        "font.family": [_font, "DejaVu Sans"],
        "axes.unicode_minus": False,
    })
else:
    print(f"[WARN] CJK font not found at {CJK_FONT_PATH}; CJK in matplotlib may render as boxes.")

# ---------------------------------------------------------------------------
# Slide geometry (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color palette
C_TITLE_BG = RGBColor(0x1F, 0x4E, 0x79)
C_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
C_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_AMBER = RGBColor(0xE6, 0x7E, 0x22)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
C_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_BORDER = RGBColor(0xBD, 0xC3, 0xC7)
C_BLUE = RGBColor(0x29, 0x80, 0xB9)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)

# ---------------------------------------------------------------------------
# Font fallback helpers
# ---------------------------------------------------------------------------
CJK_RANGES = (
    (0x3000, 0x303F),
    (0x3040, 0x30FF),
    (0x3400, 0x4DBF),
    (0x4E00, 0x9FFF),
    (0xF900, 0xFAFF),
    (0xFE30, 0xFE4F),
    (0xFF00, 0xFFEF),
)


def _is_cjk(ch: str) -> bool:
    cp = ord(ch)
    for lo, hi in CJK_RANGES:
        if lo <= cp <= hi:
            return True
    return False


def _segments(text: str):
    if not text:
        return []
    out = []
    buf = [text[0]]
    cur_is_cjk = _is_cjk(text[0])
    for ch in text[1:]:
        if _is_cjk(ch) == cur_is_cjk:
            buf.append(ch)
        else:
            out.append(("".join(buf), cur_is_cjk))
            buf = [ch]
            cur_is_cjk = _is_cjk(ch)
    out.append(("".join(buf), cur_is_cjk))
    return out


def add_text_with_fallback(text_frame, text, font_size=18, bold=False,
                           color=None, latin=LATIN_FONT, cjk=CJK_FONT_NAME,
                           align=PP_ALIGN.LEFT, paragraph_index=None,
                           italic=False):
    """Append paragraph with per-char Latin/CJK font fallback."""
    if paragraph_index is not None and paragraph_index == 0 and text_frame.paragraphs[0].text == "":
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.alignment = align

    for seg, is_cjk in _segments(text):
        run = p.add_run()
        run.text = seg
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold)
        run.font.italic = bool(italic)
        if color is not None:
            run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea"):
            for el in rPr.findall(qn(tag)):
                rPr.remove(el)
        latin_el = etree.SubElement(rPr, qn("a:latin"))
        latin_el.set("typeface", latin)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", cjk)
    return p


# ---------------------------------------------------------------------------
# Generic shape helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w_pt=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)
    shp.text_frame.margin_left = Emu(60000)
    shp.text_frame.margin_right = Emu(60000)
    shp.text_frame.margin_top = Emu(40000)
    shp.text_frame.margin_bottom = Emu(40000)
    return shp


def add_text_block(slide, x, y, w, h, lines, font_size=14, bold=False,
                   color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   fill=None, line=None, italic=False):
    if fill is not None or line is not None:
        shp = add_rect(slide, x, y, w, h, fill=fill, line=line)
    else:
        shp = slide.shapes.add_textbox(x, y, w, h)
        shp.text_frame.margin_left = Emu(60000)
        shp.text_frame.margin_right = Emu(60000)
        shp.text_frame.margin_top = Emu(40000)
        shp.text_frame.margin_bottom = Emu(40000)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        b = bold if not isinstance(bold, list) else bold[i] if i < len(bold) else False
        sz = font_size if not isinstance(font_size, list) else font_size[i] if i < len(font_size) else font_size[-1]
        col = color if not isinstance(color, list) else color[i] if i < len(color) else color[-1]
        add_text_with_fallback(tf, ln, font_size=sz, bold=b, color=col,
                               align=align, paragraph_index=i, italic=italic)
    return shp


def add_title(slide, text, subtitle=None):
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=C_TITLE_BG)
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].text = ""
    add_text_with_fallback(tf, text, font_size=26, bold=True,
                           color=C_TITLE_FG, paragraph_index=0)
    if subtitle:
        add_text_with_fallback(tf, subtitle, font_size=14, bold=False,
                               color=RGBColor(0xCF, 0xE2, 0xF3))


def fit_image_within(slide, path, x, y, max_w, max_h, border=False):
    """Add picture preserving aspect ratio inside (max_w, max_h)."""
    if not Path(path).exists():
        add_rect(slide, x, y, max_w, max_h, fill=C_LIGHT, line=C_ACCENT)
        add_text_block(slide, x, y + max_h // 2 - Inches(0.2), max_w, Inches(0.4),
                       [f"[圖片缺失] {Path(path).name}"],
                       font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
        return None
    img = Image.open(path)
    iw, ih = img.size
    img.close()
    iw_emu = iw * 9525
    ih_emu = ih * 9525
    ratio = min(max_w / iw_emu, max_h / ih_emu)
    final_w = int(iw_emu * ratio)
    final_h = int(ih_emu * ratio)
    cx = x + (max_w - final_w) // 2
    cy = y + (max_h - final_h) // 2
    pic = slide.shapes.add_picture(path, cx, cy, width=final_w, height=final_h)
    if border:
        add_rect(slide, cx, cy, final_w, final_h, fill=None, line=C_BORDER)
    return pic


def set_speaker_note(slide, text, min_chars=350):
    """Set speaker note (mandatory; >= min_chars)."""
    if not text or not text.strip():
        raise ValueError(f"Empty speaker note for slide {slide.slide_id}")
    if len(text) < min_chars:
        raise ValueError(f"Speaker note too short ({len(text)} chars) for slide {slide.slide_id}; need >= {min_chars}")
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    add_text_with_fallback(notes, text, font_size=14, paragraph_index=0)


def add_footer(slide, page_no, total=TOTAL_SLIDES):
    add_text_block(slide, Inches(0.3), Inches(7.05), Inches(9), Inches(0.35),
                   ["Self-Phasing v3 · 2026-04-30 · longphase-to-mod V5 (TO + PON)"],
                   font_size=10, color=C_GRAY)
    add_text_block(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.35),
                   [f"{page_no} / {total}"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Auto-drawn schematic figures (matplotlib)
# ---------------------------------------------------------------------------
def _save_fig(name, fig, dpi=160):
    out = AUTO_FIG_DIR / name
    fig.savefig(out, dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_s1_impact_banner():
    """S1 衝擊 TL;DR — 大字 17.3:1 -> 1:1 + concordance + 分層 caveat + 動機小字."""
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.5)
    ax.axis("off")
    # 大字主結論
    ax.text(6.5, 5.2, "17.3 : 1   →   1 : 1",
            ha="center", va="center", fontsize=68, fontweight="bold", color="#1f4e79")
    # 中字綠色
    ax.text(6.5, 3.5,
            "read-level concordance +8.3 pp (clean PS blocks, 全基因組)",
            ha="center", va="center", fontsize=22, color="#27ae60", fontweight="bold")
    # 灰字 caveat (口徑校準)
    ax.text(6.5, 2.3,
            "self-phasing 問題鏈分層處理:  V2b 解 phase scaffold |  V3F/V5 解 tag 層",
            ha="center", va="center", fontsize=15, color="#555")
    ax.text(6.5, 1.5,
            "本 PPT 焦點:  longphase-to-mod (TO + PON) 4-commit 漸進修補   ·   InterSubMod 無 C++ 改動",
            ha="center", va="center", fontsize=13, color="#7f8c8d", style="italic")
    # 底部小字 (v3 新增動機 strip)
    ax.add_patch(plt.Rectangle((1.0, 0.2), 11.0, 0.7, facecolor="#fdf6e3",
                               edgecolor="#a04000", lw=1.5))
    ax.text(6.5, 0.55,
            "修補 self-phasing 是解鎖 InterSubMod 五大研究目標的前提",
            ha="center", va="center", fontsize=14, color="#a04000", fontweight="bold")
    return _save_fig("S1_impact.png", fig)


def make_s2_pipeline_4stage():
    """S2 pipeline + 4 階段 (原 v2-26 S3)."""
    fig, ax = plt.subplots(figsize=(13, 6.0))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6)
    ax.axis("off")
    # Top: pipeline horizontal
    boxes = [
        (0.2, "Tumor\nBAM", "#bdc3c7"),
        (2.4, "ClairS-TO\n(caller)", "#3498db"),
        (4.8, "longphase-\nto-mod\n[V5]", "#c0392b"),
        (7.7, "tumor_\ntagged.bam\n(HP:i)", "#27ae60"),
        (10.4, "InterSubMod\n(read-level\nISM)", "#9b59b6"),
    ]
    for x, t, c in boxes:
        ax.add_patch(plt.Rectangle((x, 4.0), 2.0, 1.6, facecolor=c, alpha=0.25,
                                   edgecolor=c, lw=1.6))
        ax.text(x + 1.0, 4.8, t, ha="center", va="center", fontsize=11,
                fontweight="bold", color="#222")
    for x_from, x_to in [(2.2, 2.4), (4.4, 4.8), (6.8, 7.7), (9.7, 10.4)]:
        ax.annotate("", xy=(x_to, 4.8), xytext=(x_from, 4.8),
                    arrowprops=dict(arrowstyle="->", lw=2, color="#555"))
    ax.text(6.5, 3.55,
            "InterSubMod 無 C++ 改動  ·  ISM 是下游受惠者",
            ha="center", va="center", fontsize=12, color="#1f4e79", fontweight="bold")
    # Bottom: 4-stage strip
    stages = [
        (0.5, "1. 機制定位", "理論 1:1 → 觀察 17.3:1", "#3498db"),
        (3.6, "2. 4-commit 修補", "V2b → V3F → INDEL → V5", "#e67e22"),
        (6.7, "3. 驗證", "15/15 PASS + 5 證據鏈", "#27ae60"),
        (9.8, "4. 影響評估", "ISM 3-tier × 五大目標", "#8e44ad"),
    ]
    for x, t, sub, c in stages:
        ax.add_patch(plt.Rectangle((x, 1.4), 2.7, 1.6, facecolor=c, alpha=0.18,
                                   edgecolor=c, lw=1.5))
        ax.text(x + 1.35, 2.55, t, ha="center", va="center", fontsize=12,
                fontweight="bold", color=c)
        ax.text(x + 1.35, 1.95, sub, ha="center", va="center", fontsize=10, color="#333")
    ax.text(6.5, 0.7, "後續 PPT 段落依此 4 階段展開",
            ha="center", va="center", fontsize=12, color="#1f4e79", fontweight="bold")
    return _save_fig("S2_pipeline_4stage.png", fig)


def make_s4_three_layer_evidence():
    """S4 三層證據並列 (原 v2-26 S5+S6 合併, 理論+全基因組+SP1)."""
    fig, ax = plt.subplots(figsize=(13, 6.2))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.2)
    ax.axis("off")
    # Three column header
    cols = [
        (0.3, 4.3, "理論層", "germline het 隨機 → 1:1", "#3498db"),
        (4.7, 4.3, "全基因組層", "HP1:HP2 = 17.3:1", "#c0392b"),
        (9.1, 4.3, "個別位點層", "SP1 113:0 完全失衡", "#a04000"),
    ]
    for x, y, t, sub, c in cols:
        ax.add_patch(plt.Rectangle((x, y), 3.8, 1.5, facecolor=c, alpha=0.20,
                                   edgecolor=c, lw=2.0))
        ax.text(x + 1.9, y + 1.10, t, ha="center", va="center",
                fontsize=15, fontweight="bold", color=c)
        ax.text(x + 1.9, y + 0.55, sub, ha="center", va="center",
                fontsize=11, color="#222")
    # Detail boxes below each col
    details = [
        (0.3, "long-read 跨 phase block",
         "兩 hap 同概率 (~50%/50%)",
         "預期 HP1 : HP2 = 1 : 1"),
        (4.7, "HP1 = 614,000 reads",
         "HP2 = 35,500 reads",
         "94.6% 集中於 HP1\n跨 23 chr 一致"),
        (9.1, "SP1 chr19:17565944",
         "Baseline HP2:HP1 = 113:0",
         "Paired 顯示真實方向 HP2\n→ baseline 完全反向"),
    ]
    for (x, l1, l2, l3), (_, _, _, _, c) in zip(details, cols):
        ax.add_patch(plt.Rectangle((x, 0.6), 3.8, 3.4, facecolor="#fafafa",
                                   edgecolor=c, lw=1.5))
        ax.text(x + 0.2, 3.5, "·  " + l1, fontsize=11, color="#222")
        ax.text(x + 0.2, 2.95, "·  " + l2, fontsize=11, color="#222")
        ax.text(x + 0.2, 2.10, l3, fontsize=11, color="#222", va="top")
    # Bottom note
    ax.text(6.5, 0.25,
            "三層獨立證據彙整於同一結論: self-phasing artifact 真實存在",
            ha="center", va="center", fontsize=12, color="#1f4e79", fontweight="bold")
    return _save_fig("S4_three_layer.png", fig)


def make_s5_paired_vs_to():
    """S5 Paired vs TO bar (原 v2-26 S7)."""
    fig, ax = plt.subplots(figsize=(11, 5.0))
    samples = ["HCC1395\n5kHz", "HCC1395\nDORADO", "HCC1954", "COLO829", "H2009", "HCC1937", "H1437"]
    paired = [0.50, 0.51, 0.49, 0.51, 0.50, 0.50, 0.49]
    to = [0.946, 0.927, 0.951, 0.923, 0.918, 0.939, 0.911]
    x = list(range(len(samples)))
    w = 0.36
    ax.bar([i - w / 2 for i in x], paired, w, label="Paired (有 normal)",
           color="#27ae60", alpha=0.85)
    ax.bar([i + w / 2 for i in x], to, w, label="TO baseline (無 normal)",
           color="#c0392b", alpha=0.85)
    ax.axhline(0.5, color="#555", linestyle="--", lw=0.8)
    ax.text(6.6, 0.53, "預期 ~0.5", fontsize=10, color="#555")
    ax.set_xticks(x)
    ax.set_xticklabels(samples, fontsize=10)
    ax.set_ylabel("HP_Ratio", fontsize=11)
    ax.set_ylim(0, 1.05)
    ax.set_title("Paired ~0.5 vs TO 0.91-0.95  ·  Cohen's d = -1.20  ·  同位點 r = 0.001",
                 fontsize=13)
    ax.legend(loc="center right")
    ax.grid(axis="y", linestyle=":", alpha=0.4)
    return _save_fig("S5_paired_vs_to.png", fig)


def make_s10_prereq_combined():
    """S10 prerequisite — 4-row function table + PON schematic (兩欄合併).

    v3: 合併原 S11 4 函數表 + S12 PON 概念為單一 slide; speaker note 補
    Purity 0.95 閾值概念。
    """
    fig, ax = plt.subplots(figsize=(13, 6.2))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.2)
    ax.axis("off")
    # Left: 4-row function table
    ax.add_patch(plt.Rectangle((0.2, 0.4), 7.0, 5.6, facecolor="#fafafa",
                               edgecolor="#bdc3c7", lw=1.2))
    ax.text(3.7, 5.65, "longphase-to-mod 4 函數架構",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#1f4e79")
    rows = [
        ("getVote()", "統計 reads 對 HP1/HP2/somatic 投票",
         "HaplotagProcess.cpp:512-563", "#c0392b"),
        ("judgeHaplotype()", "vote → hpResult integer",
         "HaplotagProcess.cpp:697", "#e67e22"),
        ("countSNPHaplotype()", "SNP 位點各 hap read 計數",
         "HaplotagProcess.cpp:489-494", "#3498db"),
        ("countINDELHaplotype()", "INDEL 版本",
         "HaplotagProcess.cpp:497-510", "#27ae60"),
    ]
    cy = 5.05
    for fn, desc, loc, color in rows:
        ax.add_patch(plt.Rectangle((0.4, cy - 0.85), 6.6, 0.95, facecolor=color,
                                   alpha=0.10, edgecolor=color, lw=1.2))
        ax.text(0.6, cy - 0.30, fn, fontsize=12, fontweight="bold", color=color)
        ax.text(0.6, cy - 0.56, desc, fontsize=10, color="#222")
        ax.text(0.6, cy - 0.78, loc, fontsize=9, color="#555", style="italic")
        cy -= 1.05

    # Right: PON schematic
    ax.add_patch(plt.Rectangle((7.4, 0.4), 5.4, 5.6, facecolor="#fafafa",
                               edgecolor="#bdc3c7", lw=1.2))
    ax.text(10.1, 5.65, "PON 替代 normal 概念",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#1f4e79")
    # Tumor only box
    ax.add_patch(plt.Rectangle((7.6, 4.0), 1.8, 1.0, facecolor="#bdc3c7",
                               alpha=0.4, edgecolor="#7f8c8d", lw=1.5))
    ax.text(8.5, 4.6, "Tumor BAM", ha="center", fontsize=11, fontweight="bold")
    ax.text(8.5, 4.25, "(無 normal)", ha="center", fontsize=9, color="#444")
    # Arrow
    ax.annotate("", xy=(9.7, 4.5), xytext=(9.4, 4.5),
                arrowprops=dict(arrowstyle="->", lw=1.5, color="#555"))
    # PON box
    ax.add_patch(plt.Rectangle((9.8, 3.5), 2.9, 2.0, facecolor="#3498db",
                               alpha=0.20, edgecolor="#3498db", lw=2.0))
    ax.text(11.25, 5.10, "PON", ha="center", fontsize=12, fontweight="bold", color="#1f6391")
    ax.text(11.25, 4.75, "群體 germline DB", ha="center", fontsize=9, color="#444")
    ax.text(11.25, 4.40, "1000g, dbSNP, gnomAD", ha="center", fontsize=8, color="#333")
    ax.text(11.25, 4.05, "CoLoRSdb", ha="center", fontsize=8, color="#333")
    ax.text(11.25, 3.65, "→ phasing anchor", ha="center", fontsize=9, color="#c0392b",
            fontweight="bold")
    # V2b box (under)
    ax.add_patch(plt.Rectangle((7.6, 2.4), 5.0, 0.7, facecolor="#27ae60",
                               alpha=0.20, edgecolor="#27ae60", lw=1.5))
    ax.text(10.1, 2.75, "V2b 啟用 --pon-only-phasing  →  解 phase scaffold",
            ha="center", fontsize=11, color="#27ae60", fontweight="bold")
    # Purity threshold box (V3 NEW)
    ax.add_patch(plt.Rectangle((7.6, 0.7), 5.0, 1.5, facecolor="#fef5e7",
                               alpha=0.85, edgecolor="#a04000", lw=2.0))
    ax.text(10.1, 1.85, "★ Purity 0.95 觸發鎖 (S11 主者)",
            ha="center", fontsize=11, fontweight="bold", color="#a04000")
    ax.text(10.1, 1.45, "PhasingProcess.cpp:197 hardcoded",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(10.1, 1.10, "purity > 0.95 → Two-Pass | else → 三條路",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(10.1, 0.85, "(HCC1395 baseline 估 0.927 → 走三條路)",
            ha="center", fontsize=9, color="#a04000", style="italic")
    return _save_fig("S10_prereq_combined.png", fig)


def make_s11_root_cause_with_trigger():
    """S11 root-cause tree (v3 KEY) — 觸發條件主者 (root) + 三 bug (leaves)."""
    fig, ax = plt.subplots(figsize=(13.5, 6.5))
    ax.set_xlim(0, 13.5)
    ax.set_ylim(0, 6.5)
    ax.axis("off")

    # Root: 觸發條件主者 (RED, 大框, v3 NEW)
    ax.add_patch(plt.Rectangle((2.7, 5.4), 8.1, 1.0, facecolor="#fdedec",
                               edgecolor="#c0392b", lw=3.0))
    ax.text(6.75, 6.10, "★ 觸發條件 (主者)",
            ha="center", va="center", fontsize=12, fontweight="bold", color="#c0392b")
    ax.text(6.75, 5.80, "Baseline 估 purity = 0.927 ≤ 0.95 (純樣本標準)",
            ha="center", va="center", fontsize=12, color="#222")
    ax.text(6.75, 5.55, "→ 未觸發 Two-Pass → 走 baseline 標準三條路流程",
            ha="center", va="center", fontsize=11, color="#222", fontweight="bold")

    # Three branches
    ax.annotate("", xy=(2.3, 4.6), xytext=(4.5, 5.4),
                arrowprops=dict(arrowstyle="->", lw=1.8, color="#555"))
    ax.annotate("", xy=(6.75, 4.6), xytext=(6.75, 5.4),
                arrowprops=dict(arrowstyle="->", lw=1.8, color="#555"))
    ax.annotate("", xy=(11.2, 4.6), xytext=(9.0, 5.4),
                arrowprops=dict(arrowstyle="->", lw=1.8, color="#555"))

    # Bug 1: priority (red)
    ax.add_patch(plt.Rectangle((0.3, 1.8), 4.0, 2.8, facecolor="#fdedec",
                               edgecolor="#c0392b", lw=2.5))
    ax.text(2.3, 4.30, "Bug 1: getVote priority",
            ha="center", fontsize=12, fontweight="bold", color="#c0392b")
    ax.text(2.3, 3.85, "somatic-first 投票順序",
            ha="center", fontsize=10, color="#222")
    ax.text(2.3, 3.55, "variantKeys 排序",
            ha="center", fontsize=9, color="#222", style="italic")
    ax.text(2.3, 3.20, "{HP1_1, HP2_1} 排第一",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(2.3, 2.85, "任一 somatic > 0 即 break",
            ha="center", fontsize=9, color="#222")
    ax.text(2.3, 2.55, "germline 跳過",
            ha="center", fontsize=9, color="#222")
    ax.text(2.3, 2.10, "V3F 41ff147 修",
            ha="center", fontsize=10, fontweight="bold", color="#c0392b")

    # Bug 2: enum (orange)
    ax.add_patch(plt.Rectangle((4.75, 1.8), 4.0, 2.8, facecolor="#fef5e7",
                               edgecolor="#e67e22", lw=2.5))
    ax.text(6.75, 4.30, "Bug 2: enum vs int literal",
            ha="center", fontsize=12, fontweight="bold", color="#e67e22")
    ax.text(6.75, 3.95, "Util.h:20-25",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(6.75, 3.65, "HAPLOTYPE1_1 = 3 (enum)",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(6.75, 3.35, "HP tag integer = 11",
            ha="center", fontsize=9, color="#222", family="monospace")
    ax.text(6.75, 3.00, "if(hpResult != HAPLOTYPE1_1)",
            ha="center", fontsize=8, color="#222", family="monospace")
    ax.text(6.75, 2.70, "永不匹配 → HP:i:33 不出現",
            ha="center", fontsize=9, color="#222")
    ax.text(6.75, 2.10, "V3F 41ff147 修",
            ha="center", fontsize=10, fontweight="bold", color="#e67e22")

    # Bug 3: scaffold (gray, V2b 已解)
    ax.add_patch(plt.Rectangle((9.2, 1.8), 4.0, 2.8, facecolor="#ecf0f1",
                               edgecolor="#7f8c8d", lw=1.5))
    ax.text(11.2, 4.30, "Bug 3: phase scaffold",
            ha="center", fontsize=12, fontweight="bold", color="#7f8c8d")
    ax.text(11.2, 3.95, "PhasingProcess.cpp",
            ha="center", fontsize=9, color="#444", family="monospace")
    ax.text(11.2, 3.65, "convertNonGermlineToSomatic",
            ha="center", fontsize=8, color="#444", family="monospace")
    ax.text(11.2, 3.30, "未觸發 → somatic 當 anchor",
            ha="center", fontsize=9, color="#444")
    ax.text(11.2, 2.85, "(已由 V2b 8b8c1fd 解決)",
            ha="center", fontsize=9, color="#444", style="italic")
    ax.text(11.2, 2.10, "scaffold 已解, 焦點在 1+2",
            ha="center", fontsize=10, color="#7f8c8d", fontweight="bold")

    # Bottom note
    ax.text(6.75, 1.25,
            "因樣本實為純無 normal contamination 而流程假設有 → 暴露 tag 層 somatic-first 投票 bias",
            ha="center", fontsize=11, color="#1f4e79", fontweight="bold")
    ax.text(6.75, 0.75,
            "結果:  self-phasing 17.3 : 1 artifact",
            ha="center", fontsize=12, color="#c0392b", fontweight="bold")
    ax.text(6.75, 0.35,
            "驗證來源: source_materials/04_purity_calculator_failure_root_cause.md (v5_audit_suite/18 複製)",
            ha="center", fontsize=8, color="#7f8c8d", style="italic")
    return _save_fig("S11_root_cause_v3.png", fig)


def make_s12_v5_three_layer():
    """S12 V5 三層投票流程 (原 v2-26 S14, 加紅圈標 mid-low purity 防禦層)."""
    fig, ax = plt.subplots(figsize=(13, 6.2))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.2)
    ax.axis("off")
    layers = [
        (0.4, 4.0, "Layer 1: germline first",
         "if (g1 > g2) → 1\nelif (g2 > g1) → 2\nelse → Layer 1.5",
         "#3498db", "翻轉 baseline 的\nsomatic-first 投票順序"),
        (4.6, 4.0, "Layer 1.5: somatic fallback (V5 新增)",
         "if (sTotal == 0) → 0\nelif s1 > 0.6*tot → align HP1\nelif s2 > 0.6*tot → align HP2\nelse → 33 (true ambiguous)",
         "#c0392b", "Confidence 0.6 =\nmid-low purity 防禦層"),
        (9.0, 4.0, "Layer 2: encode hpResult",
         "1 + somatic → HP:i:11\n2 + somatic → HP:i:21\n0 + somatic → HP:i:33\ngermline only → 1 / 2\nnothing → 0",
         "#27ae60", "對應 5 種 BAM HP\n integer 編碼"),
    ]
    for x, y, title, body, color, marker in layers:
        ax.add_patch(plt.Rectangle((x, y - 2.6), 4.0, 2.6, facecolor=color,
                                   alpha=0.15, edgecolor=color, lw=2.0))
        ax.text(x + 2.0, y - 0.25, title, ha="center", va="center",
                fontsize=12, fontweight="bold", color=color)
        ax.text(x + 0.2, y - 0.7, body, va="top", fontsize=10, color="#222",
                family="monospace")
        # Red ring marker (key idea per layer)
        ax.add_patch(plt.Rectangle((x + 0.2, y - 2.55), 3.6, 0.55,
                                   facecolor="#fdedec", edgecolor="#c0392b", lw=1.5))
        ax.text(x + 2.0, y - 2.27, marker, ha="center", va="center",
                fontsize=9, color="#c0392b", fontweight="bold")

    # Top label
    ax.text(6.5, 5.85, "V5 getVote() 三層投票流程",
            ha="center", va="center", fontsize=15, fontweight="bold", color="#1f4e79")
    ax.text(6.5, 5.45, "對應 S13 程式碼 diff 三處紅標",
            ha="center", va="center", fontsize=11, color="#555")
    # Arrows between layers
    ax.annotate("", xy=(4.6, 2.7), xytext=(4.4, 2.7),
                arrowprops=dict(arrowstyle="->", lw=2, color="#555"))
    ax.annotate("", xy=(9.0, 2.7), xytext=(8.6, 2.7),
                arrowprops=dict(arrowstyle="->", lw=2, color="#555"))
    # Bottom Q2 hint
    ax.add_patch(plt.Rectangle((0.4, 0.2), 12.4, 0.8, facecolor="#fef5e7",
                               edgecolor="#a04000", lw=1.5))
    ax.text(6.5, 0.70, "Q2 答: V5 在 mid-low purity 比 baseline 更好",
            ha="center", fontsize=10, color="#a04000", fontweight="bold")
    ax.text(6.5, 0.40, "0.6 sample HP33 比例 12.4% vs Baseline 2% (+10pp 保守)  ·  避免錯誤分配",
            ha="center", fontsize=10, color="#444")
    return _save_fig("S12_v5_layers_v3.png", fig)


def make_s13_code_diff():
    """S13 程式碼 diff (baseline vs V5, 兩欄 + 紅標 baseline somatic-first caveat)."""
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.5)
    ax.axis("off")
    # Headers
    ax.add_patch(plt.Rectangle((0.2, 5.85), 6.3, 0.5, facecolor="#c0392b",
                               alpha=0.85, edgecolor="#c0392b", lw=1.0))
    ax.text(3.35, 6.10, "Baseline (有 bug)", ha="center", va="center",
            fontsize=14, fontweight="bold", color="#fff")
    ax.add_patch(plt.Rectangle((6.7, 5.85), 6.1, 0.5, facecolor="#27ae60",
                               alpha=0.85, edgecolor="#27ae60", lw=1.0))
    ax.text(9.75, 6.10, "V5 (修補後)", ha="center", va="center",
            fontsize=14, fontweight="bold", color="#fff")

    # ★ V3 新增: baseline 欄上方紅標 caveat (somatic-first 投票順序)
    ax.add_patch(plt.Rectangle((0.2, 5.30), 6.3, 0.45, facecolor="#fdedec",
                               edgecolor="#c0392b", lw=1.5))
    ax.text(3.35, 5.52, "★ caveat: 投票邏輯 = somatic-first 優先序 (非「有 HP11/HP21 即 somatic」這個簡化敘述)",
            ha="center", va="center", fontsize=9, color="#c0392b", fontweight="bold")

    # Code blocks (dark bg)
    ax.add_patch(plt.Rectangle((0.2, 0.9), 6.3, 4.3, facecolor="#1e1e1e",
                               edgecolor="#444", lw=1.0))
    ax.add_patch(plt.Rectangle((6.7, 0.9), 6.1, 4.3, facecolor="#1e1e1e",
                               edgecolor="#444", lw=1.0))

    # Baseline code
    base_lines = [
        ("for (auto& vk : variantKeys) {", "#fff", False),
        ("  // {HP1_1,HP2_1} listed first", "#888", False),
        ("  if (count[vk.first] > 0", "#fff", False),
        ("      || count[vk.second] > 0)", "#ff7777", True),
        ("    break;  // somatic preempts", "#ff7777", True),
        ("}", "#fff", False),
        ("if (hpResult != HAPLOTYPE1_1)", "#ff7777", True),
        ("  hpResult = 0;  // enum=3", "#ff7777", True),
        ("  // never matches HP tag=11", "#888", False),
    ]
    cy = 4.95
    for line, color, hl in base_lines:
        ax.text(0.35, cy, line, fontsize=11, family="monospace",
                color=color, va="top",
                fontweight="bold" if hl else "normal")
        cy -= 0.42
    # V5 code
    v5_lines = [
        ("// Layer 1: germline first", "#888", False),
        ("if (g1 > g2) germlineRes = 1;", "#7eff7e", True),
        ("else if (g2 > g1) germlineRes = 2;", "#7eff7e", True),
        ("// Layer 1.5: somatic fallback", "#888", False),
        ("else if (s1 + s2 > 0) {", "#7eff7e", True),
        ("  if (s1 > 0.6*tot) align HP1;", "#7eff7e", True),
        ("  else if (s2 > 0.6*tot) align HP2;", "#7eff7e", True),
        ("}", "#fff", False),
        ("if (hpResult != 11)  // int literal", "#7eff7e", True),
    ]
    cy = 4.95
    for line, color, hl in v5_lines:
        ax.text(6.85, cy, line, fontsize=11, family="monospace",
                color=color, va="top",
                fontweight="bold" if hl else "normal")
        cy -= 0.42
    # Bottom annotations
    ax.text(6.5, 0.45,
            "三處紅標: (1) priority loop 拆解 (somatic-first → germline-first)  (2) Layer 1.5 純分支  (3) enum 改 int literal 11",
            ha="center", fontsize=10, color="#c0392b", fontweight="bold")
    return _save_fig("S13_code_diff_v3.png", fig)


def make_s17_hp_flip_bar():
    """S18 HP flip bar (3 SP-extreme sites)."""
    fig, ax = plt.subplots(figsize=(11, 5.0))
    sites = ["SP1\nchr19:17565944", "SP2\nchr19:12452332", "SP3\nchr19:12467180"]
    base_hp1 = [113, 109, 108]
    base_hp2 = [0, 1, 0]
    v5_hp1 = [0, 1, 0]
    v5_hp2 = [113, 109, 108]
    pair_hp1 = [2, 3, 1]
    pair_hp2 = [110, 105, 100]

    x = list(range(len(sites)))
    w = 0.25
    ax.bar([i - w for i in x], base_hp1, w, color="#c0392b", alpha=0.85, label="HP1")
    ax.bar([i - w for i in x], base_hp2, w, bottom=base_hp1, color="#3498db", alpha=0.5, label="HP2")
    ax.bar(x, v5_hp1, w, color="#c0392b", alpha=0.85)
    ax.bar(x, v5_hp2, w, bottom=v5_hp1, color="#3498db", alpha=0.5)
    ax.bar([i + w for i in x], pair_hp1, w, color="#c0392b", alpha=0.85)
    ax.bar([i + w for i in x], pair_hp2, w, bottom=pair_hp1, color="#3498db", alpha=0.5)
    for i in x:
        ax.text(i - w, -8, "Base", ha="center", fontsize=9, color="#444")
        ax.text(i, -8, "V5", ha="center", fontsize=9, color="#444")
        ax.text(i + w, -8, "Pair", ha="center", fontsize=9, color="#444")
    ax.set_xticks(x)
    ax.set_xticklabels(sites, fontsize=10)
    ax.set_ylabel("read count", fontsize=11)
    ax.set_ylim(-15, 130)
    ax.set_title("HP 翻轉: Baseline → V5 → Paired ground truth (3/3 SP-extreme 一致)",
                 fontsize=13)
    ax.legend(loc="upper right", fontsize=9)
    ax.grid(axis="y", linestyle=":", alpha=0.4)
    return _save_fig("S18_hp_flip_v3.png", fig)


def make_s19_f1_decomposition():
    """S19 F1 流程圖 (caller raw 不變 → ISM SuggestFilter ΔF1 噪音)."""
    fig, ax = plt.subplots(figsize=(13, 5.6))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6)
    ax.axis("off")
    ax.add_patch(plt.Rectangle((0.5, 4.6), 12.0, 0.9, facecolor="#3498db",
                               alpha=0.18, edgecolor="#3498db", lw=1.8))
    ax.text(6.5, 5.05, "ClairS-TO raw F1 = 0.7166  (Baseline / V3F / V5 三版本完全相同)",
            ha="center", va="center", fontsize=14, fontweight="bold", color="#1f6391")
    ax.text(6.5, 4.75, "→ V5 不改 caller, 不改 VCF",
            ha="center", va="center", fontsize=11, color="#555")
    boxes = [
        (0.5, "Baseline", "F1 = 0.7157", "#7f8c8d"),
        (4.7, "V3-Fixed", "F1 = 0.7154", "#e67e22"),
        (8.9, "V5", "F1 = 0.7154", "#27ae60"),
    ]
    for x, t, f1, color in boxes:
        ax.add_patch(plt.Rectangle((x, 2.6), 4.0, 1.4, facecolor=color, alpha=0.20,
                                   edgecolor=color, lw=1.8))
        ax.text(x + 2.0, 3.55, t, ha="center", fontsize=14, fontweight="bold", color=color)
        ax.text(x + 2.0, 3.0, f1, ha="center", fontsize=12, color="#222")
    for cx in [2.5, 6.7, 10.9]:
        ax.annotate("", xy=(cx, 4.0), xytext=(cx, 4.55),
                    arrowprops=dict(arrowstyle="->", lw=1.6, color="#555"))
    ax.add_patch(plt.Rectangle((0.5, 0.4), 12.0, 1.7, facecolor="#fef5e7",
                               edgecolor="#e67e22", lw=1.8))
    ax.text(6.5, 1.78, "ΔF1 = -0.0003 (V5 vs Baseline) 噪音級",
            ha="center", fontsize=14, fontweight="bold", color="#a04000")
    ax.text(6.5, 1.30,
            "F1 不能衡量本實作品質  ·  真實價值在 read-level concordance +8.3 pp",
            ha="center", fontsize=11, color="#444")
    ax.text(6.5, 0.85,
            "caveat: cnLOH 區未解 (R1)  ·  AF>0.9 邊界 (R2)  ·  V5 working tree 未 commit (R5)",
            ha="center", fontsize=10, color="#555", style="italic")
    return _save_fig("S19_f1_decomp.png", fig)


def make_s20_industry_combined():
    """S20 業界家族樹 + 2x2 合併 (上半樹 + 下半 matrix)."""
    fig, ax = plt.subplots(figsize=(13, 7.0))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 7)
    ax.axis("off")

    # Top half: family tree
    ax.text(6.5, 6.7, "業界家族樹 (同實驗室相鄰工作)",
            ha="center", fontsize=13, fontweight="bold", color="#1f4e79")
    # Root
    ax.add_patch(plt.Rectangle((4.8, 5.65), 3.4, 0.7, facecolor="#3498db",
                               alpha=0.25, edgecolor="#1f6391", lw=2.0))
    ax.text(6.5, 5.97, "LongPhase (Lin 2022)",
            ha="center", va="center", fontsize=12, fontweight="bold", color="#1f6391")
    # Two branches
    ax.add_patch(plt.Rectangle((1.0, 4.10), 4.4, 1.05, facecolor="#27ae60",
                               alpha=0.20, edgecolor="#27ae60", lw=1.8))
    ax.text(3.2, 4.85, "LongPhase-S (bioRxiv 2025)",
            ha="center", fontsize=11, fontweight="bold", color="#27ae60")
    ax.text(3.2, 4.50, "Tumor-Normal Paired",
            ha="center", fontsize=10, color="#27ae60")
    ax.text(3.2, 4.20, "ClairS SNV F1 +4.5% (paired)",
            ha="center", fontsize=9, color="#444")
    ax.add_patch(plt.Rectangle((7.6, 4.10), 4.4, 1.05, facecolor="#c0392b",
                               alpha=0.20, edgecolor="#c0392b", lw=2.5))
    ax.text(9.8, 4.85, "longphase-to-mod V5 (本實作)",
            ha="center", fontsize=11, fontweight="bold", color="#c0392b")
    ax.text(9.8, 4.50, "Tumor-Only + PON",
            ha="center", fontsize=10, color="#c0392b")
    ax.text(9.8, 4.20, "+8.3 pp clean PS concordance",
            ha="center", fontsize=9, color="#444")
    # Arrows
    ax.annotate("", xy=(3.2, 5.18), xytext=(5.6, 5.65),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(9.8, 5.18), xytext=(7.4, 5.65),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))

    # Bottom half: 2x2 matrix
    ax.text(6.5, 3.50, "2x2 比較表 (Source × Mode)",
            ha="center", fontsize=13, fontweight="bold", color="#1f4e79")
    # Headers
    ax.add_patch(plt.Rectangle((2.2, 2.95), 4.5, 0.4, facecolor="#3498db",
                               alpha=0.85, edgecolor="#3498db"))
    ax.text(4.45, 3.15, "Germline", ha="center", va="center", fontsize=11,
            fontweight="bold", color="#fff")
    ax.add_patch(plt.Rectangle((7.0, 2.95), 4.5, 0.4, facecolor="#c0392b",
                               alpha=0.85, edgecolor="#c0392b"))
    ax.text(9.25, 3.15, "Tumor", ha="center", va="center", fontsize=11,
            fontweight="bold", color="#fff")
    # Row labels
    ax.add_patch(plt.Rectangle((1.4, 2.0), 0.7, 0.9, facecolor="#7f8c8d", alpha=0.85))
    ax.text(1.75, 2.45, "公開\n工具", ha="center", va="center",
            fontsize=10, fontweight="bold", color="#fff")
    ax.add_patch(plt.Rectangle((1.4, 0.6), 0.7, 1.3, facecolor="#27ae60", alpha=0.85))
    ax.text(1.75, 1.25, "本\n實作", ha="center", va="center",
            fontsize=10, fontweight="bold", color="#fff")
    # Cells
    ax.add_patch(plt.Rectangle((2.2, 2.0), 4.5, 0.9, facecolor="#ecf0f1",
                               edgecolor="#bdc3c7", lw=1.0))
    ax.text(4.45, 2.55, "LongPhase / WhatsHap / HapCUT2",
            ha="center", fontsize=10, color="#222")
    ax.text(4.45, 2.25, "成熟 germline phasing", ha="center", fontsize=9, color="#27ae60")

    ax.add_patch(plt.Rectangle((7.0, 2.0), 4.5, 0.9, facecolor="#fef5e7",
                               edgecolor="#e67e22", lw=1.5))
    ax.text(9.25, 2.55, "LongPhase-S (bioRxiv 2025)",
            ha="center", fontsize=10, fontweight="bold", color="#e67e22")
    ax.text(9.25, 2.25, "Paired only (需 normal)", ha="center", fontsize=9, color="#444")

    ax.add_patch(plt.Rectangle((2.2, 0.6), 4.5, 1.3, facecolor="#ecf0f1",
                               edgecolor="#bdc3c7", lw=1.0))
    ax.text(4.45, 1.45, "(本實作不做 germline)", ha="center", fontsize=10,
            color="#888", style="italic")
    ax.text(4.45, 1.05, "上游用 ClairS-TO + LongPhase",
            ha="center", fontsize=9, color="#888")

    ax.add_patch(plt.Rectangle((7.0, 0.6), 4.5, 1.3, facecolor="#fdedec",
                               edgecolor="#c0392b", lw=2.5))
    ax.text(9.25, 1.65, "longphase-to-mod V5",
            ha="center", fontsize=12, fontweight="bold", color="#c0392b")
    ax.text(9.25, 1.30, "Tumor-Only + PON (本工作)",
            ha="center", fontsize=10, color="#222")
    ax.text(9.25, 0.95, "TO + PON 條件下的本地實作",
            ha="center", fontsize=9, color="#444", style="italic")
    ax.text(9.25, 0.70, "填補 TO 場景 gap", ha="center", fontsize=9, color="#27ae60",
            fontweight="bold")
    # Bottom note
    ax.text(6.5, 0.20,
            "口徑: 同實驗室相鄰工作 (非「業界共識」「標準替代」)  ·  LongPhase-S +4.5% 為 paired 參考非直接可比",
            ha="center", fontsize=9, color="#7f8c8d", style="italic")
    return _save_fig("S20_industry_combined.png", fig)


def make_s22_research_tree():
    """S22 五大目標 + 研究發展樹."""
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.5)
    ax.axis("off")
    # Top label
    ax.text(6.5, 6.20, "Self-phasing 修補 → 五大目標解鎖",
            ha="center", fontsize=13, fontweight="bold", color="#1f4e79")
    # Top: self-phasing 修補
    ax.add_patch(plt.Rectangle((4.5, 5.30), 4.0, 0.7, facecolor="#27ae60",
                               alpha=0.25, edgecolor="#27ae60", lw=2.0))
    ax.text(6.5, 5.65, "Self-phasing 修補完成 (V5)",
            ha="center", va="center", fontsize=12, fontweight="bold", color="#27ae60")
    # Middle: 4-bucket 可信
    ax.add_patch(plt.Rectangle((4.5, 4.10), 4.0, 0.7, facecolor="#3498db",
                               alpha=0.25, edgecolor="#3498db", lw=2.0))
    ax.text(6.5, 4.45, "4-bucket 分群可信", ha="center", va="center",
            fontsize=11, fontweight="bold", color="#1f6391")
    # NG analysis
    ax.add_patch(plt.Rectangle((4.5, 3.00), 4.0, 0.6, facecolor="#3498db",
                               alpha=0.18, edgecolor="#3498db", lw=1.5))
    ax.text(6.5, 3.30, "NGroups / HPSig / HP_Ratio 可信",
            ha="center", va="center", fontsize=10, fontweight="bold", color="#1f6391")
    # 5 goals
    goals = [
        (0.3, "目標 1\nper-CpG\n關聯", "#c0392b"),
        (2.9, "目標 2\nclone\n結構", "#e67e22"),
        (5.5, "目標 3\n二次打擊\n順序", "#f39c12"),
        (8.1, "目標 4\nTO normal\n補強", "#27ae60"),
        (10.7, "目標 5\nF1 panel", "#8e44ad"),
    ]
    for x, label, color in goals:
        ax.add_patch(plt.Rectangle((x, 0.6), 2.2, 1.7, facecolor=color,
                                   alpha=0.20, edgecolor=color, lw=1.8))
        ax.text(x + 1.1, 1.45, label, ha="center", va="center",
                fontsize=10, fontweight="bold", color=color)
    # Arrows
    ax.annotate("", xy=(6.5, 4.80), xytext=(6.5, 5.30),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    ax.annotate("", xy=(6.5, 3.60), xytext=(6.5, 4.10),
                arrowprops=dict(arrowstyle="->", lw=2.0, color="#555"))
    for x, _, _ in goals:
        ax.annotate("", xy=(x + 1.1, 2.30), xytext=(6.5, 3.00),
                    arrowprops=dict(arrowstyle="->", lw=1.0, color="#888", alpha=0.5))
    ax.text(11.5, 5.0, "目標 1/2/4\n直接依賴", ha="center", va="center",
            fontsize=9, color="#c0392b", fontweight="bold")
    ax.text(0.5, 5.0, "目標 3/5\n部分依賴", ha="center", va="center",
            fontsize=9, color="#e67e22", fontweight="bold")
    ax.text(6.5, 0.30,
            "來源: InterSubMod/docs/architecture/20260327_InterSubMod研究願景定錨_01.md",
            ha="center", fontsize=8, color="#7f8c8d", style="italic")
    return _save_fig("S22_research_tree_v3.png", fig)


# ---------------------------------------------------------------------------
# Slide builders (S1 - S24, v3 24-slide structure)
# ---------------------------------------------------------------------------
def slide_01_impact(prs):
    """S1: 17.3:1 → 1:1 + 動機小字."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    img = make_s1_impact_banner()
    fit_image_within(s, str(img), Inches(0.2), Inches(0.4), Inches(13.0), Inches(6.7))
    add_footer(s, 1)
    set_speaker_note(s,
        "開場 60 秒鎖定核心結論。本 PPT 是 v3 24-slide 版本, 主軸是 longphase-to-mod fork "
        "(TO + PON 條件下的本地實作) 4-commit 漸進修補, InterSubMod 本 repo 無 C++ 改動。"
        "三行核心訊息: (1) 大字 17.3 : 1 → 1 : 1 是 HP1 family : HP2 family ratio, "
        "baseline 17.3 倍偏在 HP1, 修補後回到生物學期望 1:1。"
        "(2) read-level concordance +8.3 pp 是全基因組 clean PS blocks (PI 報告 4 §3.7) 對 paired ground truth "
        "的吻合度提升 (V5 90.5% vs Baseline 82.2%)。"
        "(3) 重要分層: self-phasing 問題鏈是被分層處理的 — V2b 解 phase scaffold (--pon-only-phasing flag), "
        "V3F/V5 解 tag 層 (priority bug + enum mismatch + V3F 過度保守的 directional reassignment)。"
        "**V5 不修 self-phasing 本身** (V2b 已處理 phase scaffold)。"
        "底部小字 (v3 新增動機 strip): 修補 self-phasing 是解鎖 InterSubMod 五大研究目標的前提。"
        "ISM 五大目標包括: 目標 1 per-CpG 多標籤關聯、目標 2 clone 結構、目標 3 二次打擊、"
        "目標 4 TO normal 補強、目標 5 F1 panel; 目標 1/2/4 直接依賴 HP tag 正確。"
        "本 PPT 24 slides 分 6 段: 高層次重點 (S1-S2) → 觀察問題 (S3-S7) → 為何重要 (S8-S9) → "
        "解釋與原因 (S10-S13) → 改動驗證 (S14-S20) → 未來規劃結語 (S21-S24)。"
        "Q1 預備: HP1 family = 614,000 reads, HP2 family = 35,500, 比值 17.3, 跨 23 染色體一致, "
        "94.6% 集中於 HP1; 生物學上 germline het 隨機分配應 ~1:1。")
    return s


def slide_02_pipeline(prs):
    """S2: Pipeline + 4 階段 (原 v2-26 S3 內容)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Pipeline 與 4 階段工作")
    img = make_s2_pipeline_4stage()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.8))
    add_footer(s, 2)
    set_speaker_note(s,
        "本張展示三個系統的 pipeline 切分與本工作 4 階段的關係。"
        "Pipeline: tumor BAM → ClairS-TO (caller) → longphase-to-mod (phasing + haplotag, V5 修補位置) → "
        "tumor_tagged.bam (HP:i tag) → InterSubMod (read-level ISM)。"
        "**強調**: InterSubMod 不是修補位置, 本 repo 無 C++ 改動; HaplotagProcess.cpp 在 longphase-to-mod fork "
        "(獨立 git repo, /big7_disk/liaoyoyo2001/longphase-to-mod/), 不在 ISM 內。"
        "ISM 只是讀新版 BAM 後下游分析自動受惠。"
        "4 階段工作: (1) 機制定位 — 從理論預期 1:1 對照觀察到的 17.3:1, 推導 self-phasing 機制。"
        "(2) 4-commit 修補 — V2b → V3F → INDEL guard → V5, 每個 commit 解一層 bug。"
        "(3) 驗證 — 4 項硬性 sanity check 15/15 PASS + 5 條獨立證據鏈 "
        "(理論 / 全基因組 / 個別位點 / sanity / 程式碼最小必要)。"
        "(4) 影響評估 — ISM 3-tier (29 嚴重 / 14 中度 / 42 不受影響, 共 85 features) 與五大研究目標銜接。"
        "後續 PPT 段落依此 4 階段順序展開, 教授可以對照這張作為 navigation map。"
        "v2-26 原 S2 業界家族樹已吸收進 v3 S20 (S20 業界家族樹 + 2x2 合併)。"
        "Q 預備: 「ISM 為什麼不直接修?」答: 違反單一職責, phasing/haplotag 該在 longphase-to-mod, "
        "ISM 重複實作會介面割裂; ISM F1 = 0.0124 對 TO germline FP 無修復力, 無法取代上游修補。")
    return s


def slide_03_hp_tag_def(prs):
    """S3: HP tag 五整數值 + PS/LOH 兩層."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "HP tag 定義 (5 整數值 + PS / LOH 兩層)")
    fit_image_within(s, str(FIG_DIR / "fig17_hp_tag_5versions.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.6))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["HP:i:0 / 1 / 2 / 11 / 21 / 33  ·  PS = phase-set ID  ·  LOH 兩層 (HP_Ratio vs LOH.bed)"],
                   font_size=12, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 3)
    set_speaker_note(s,
        "建立 HP tag 認知。BAM 內每條 read 一個整數 HP tag, 6 種可能值: "
        "HP:i:0 (unphased), HP:i:1 (germline HP1), HP:i:2 (germline HP2), "
        "HP:i:11 (somatic on HP1), HP:i:21 (somatic on HP2), HP:i:33 (somatic ambiguous)。"
        "這個編碼是 longphase-to-mod 自訂, 不在 LongPhase 官方 spec, 與 C++ enum (Util.h:20-25) "
        "HAPLOTYPE_UNDEFINED=-1, HAPLOTYPE1=1, HAPLOTYPE2=2, HAPLOTYPE1_1=3, HAPLOTYPE2_1=4, HAPLOTYPE3=5 不同 — "
        "這個型別語意失配是後面 S11 enum bug (Bug 2) 的伏筆。"
        "PS = phase-set ID: 同一個 PS 表示這些 reads/variants 屬於同一個 phase block。"
        "LOH 兩層必須區分: ISM HP_Ratio LOH 是 read-level 從 BAM HP tag 統計每 region 的 HP_Ratio "
        "(受 self-phasing 影響, 62% artifact); LOH.bed 是 region-level 由 longphase-to-mod 直接從 VCF AD 偵測 "
        "(不受 self-phasing 影響, V5 前後 Jaccard=1.0)。"
        "兩套 LOH 系統使用**不同數據源** (BAM HP tag vs VCF AD), 過去把 ISM HP_Ratio LOH 當作 LOH.bed 真值是常見誤解。"
        "後面 S6 (拆鎖) + S7 (LOH 兩層精確化) 會深入這個區分。"
        "Q3 預備 (「LOH.bed Jaccard=1.0 怎麼證明?」): baseline LOH.bed 與 V2b PON-only LOH.bed 的 region-level "
        "Jaccard 相似度 = 1.0, 兩 BED 完全相同, phasing 層 LOH region 偵測不受 self-phasing 影響。")
    return s


def slide_04_three_layer_evidence(prs):
    """S4: 證據三層 (理論+全基因組+SP1, 原 v2-26 S5+S6 合併)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "證據三層 — 理論 / 全基因組 / 個別位點")
    img = make_s4_three_layer_evidence()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.7))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["3 層獨立證據彙整於同一結論  ·  self-phasing 是真實 artifact 非統計噪音"],
                   font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 4)
    set_speaker_note(s,
        "證據三層並列 (v3 合併原 v2-26 S5 觀察 17.3:1 + S6 理論 1:1 為單一 slide)。"
        "**第一層 理論**: 為什麼 17.3:1 是 artifact 而非真實生物? "
        "Germline het 是個體 zygosity 層的差異 (一個位點兩個 allele), 在染色體上**隨機分散** — "
        "沒有方向性偏好, 兩 hap 的 het variants 數量相同。"
        "Long-read (Nanopore / PacBio) 平均 5-30 kb, 跨多個 phase block, 理論上 HP1 / HP2 應各分到一半 reads。"
        "正確預期 HP1 : HP2 = 1 : 1 (隨機 50% / 50% 分配)。"
        "**第二層 全基因組**: HCC1395 5kHz baseline 硬數據: HP1 family (HP1 + HP1_1 reads) = 614,000, "
        "HP2 family = 35,500, 比例 17.3:1; 94.6% 集中於 HP1; 跨 23 染色體一致。"
        "fig01d Panel D 的 6 圖證明這是全基因組層的系統性偏差, 不是個別 hotspot 也不是統計噪音。"
        "**第三層 個別位點 (SP1 chr19:17565944)**: baseline HP2:HP1 = 113:0 完全失衡 (113 reads 全 HP1, HP2 = 0); "
        "paired tumor BAM 確認真實方向是 HP2 (110:2) — baseline 完全反向。"
        "三層獨立證據彙整於同一結論: self-phasing artifact 真實存在, 不可能是隨機誤差。"
        "Q1 預備: 17.3 = 614000 / 35500, 跨 23 染色體一致, 94.6% = HP1 family / total。"
        "下一張 S5 用 paired 對照印證理論預期 1:1。")
    return s


def slide_05_paired_vs_to(prs):
    """S5: Paired ~1:1 vs TO 0.94 (原 v2-26 S7)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Paired 對照確實 ~1 : 1")
    img = make_s5_paired_vs_to()
    fit_image_within(s, str(img), Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.6))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["7/7 樣本同方向  ·  Cohen's d = -1.20  ·  同位點跨模式 r = 0.001"],
                   font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 5)
    set_speaker_note(s,
        "核心對照證據。同樣本 paired 與 TO 的 HP_Ratio 對照: "
        "Paired (有 normal) 全 7 樣本 HP_Ratio 都在 0.49-0.51 之間, 接近預期 0.5; "
        "TO baseline (無 normal) 全 7 樣本 HP_Ratio 在 0.91-0.95, 嚴重偏向 0.94 平均, "
        "Cohen's d = -1.20 (巨大效應量)。"
        "**關鍵統計**: 同位點跨模式 (paired vs TO) HP_Ratio 相關係數 r = 0.001, n = 288K pairs — "
        "幾乎完全無相關, 表示 TO HP_Ratio 完全不反映真實 haplotype, 而是 self-phasing artifact。"
        "TO-only 標記為 LOH 的位點, 在 paired 模式下 86.5% 是平衡的, "
        "也就是 TO 的 LOH 標記大部分是 self-phasing artifact, 不是真實 LOH。"
        "這個對照排除了「TO 樣本本身就 LOH 多」的解釋。"
        "**Q2 預備** (「那為什麼不直接用 paired 就好?」): normal BAM 不一定可得 (臨床、archive、cell-free DNA 等), "
        "TO 是必要研究方向。本實作目標就是讓 TO 模式達到接近 paired 的 tag 品質。"
        "**Q2 機制**: paired normal 提供 germline het 的 ground truth scaffold; "
        "TO 模式無 normal, phasing 必須從 tumor 自己推 germline, 但 tumor 含 somatic, "
        "somatic 反客為主進 phasing graph → self-phasing。"
        "下一張 S6 進入拆鎖: phasing 層其實沒問題 (LOH.bed Jaccard = 1.0), 是 tag 層 (BAM HP_Ratio) 的 17.3:1 偏差。")
    return s


def slide_06_unlock_phasing_vs_tag(prs):
    """S6: 拆鎖 phasing vs tag (原 v2-26 S8 拆出)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "拆鎖: phasing 沒問題, 是 tag 的問題")
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.5)
    add_text_block(s, Inches(0.7), Inches(1.25), Inches(5.6), Inches(0.55),
                   ["Phasing 層 (LOH.bed)"],
                   font_size=18, bold=True, color=C_GREEN)
    add_text_block(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(4.3),
                   ["路徑:  VCF AD → LongPhase region 偵測",
                    "輸出:  LOH.bed (region-level)",
                    "",
                    "V5 前後 Jaccard = 1.0",
                    "完全不變  (PASS)",
                    "",
                    "self-phasing 不污染此層",
                    "  → V2b PON-only 已解 phase scaffold",
                    "",
                    "結論:  region-level LOH 真值"],
                   font_size=13, color=C_TEXT)
    add_rect(s, Inches(6.85), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.5)
    add_text_block(s, Inches(7.05), Inches(1.25), Inches(5.6), Inches(0.55),
                   ["Tag 層 (BAM HP_Ratio)"],
                   font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(7.05), Inches(1.85), Inches(5.6), Inches(4.3),
                   ["路徑:  BAM HP:i tag → ISM 統計",
                    "輸出:  HP_Ratio per region",
                    "",
                    "Baseline 17.3 : 1  (FAIL)",
                    "62% ISM HP_Ratio LOH 是 artifact",
                    "",
                    "self-phasing 嚴重污染",
                    "  → V3F + V5 修補對象",
                    "",
                    "結論:  本 PPT 焦點在此層"],
                   font_size=13, color=C_TEXT)
    add_text_block(s, Inches(0.5), Inches(6.65), Inches(12.5), Inches(0.5),
                   ["LOH.bed 由 VCF AD 產生  ·  HP_Ratio 由 BAM HP tag 產生  ·  本 PPT 焦點 = tag 層"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 6)
    set_speaker_note(s,
        "拆鎖關鍵: self-phasing 影響的是哪一層? 答案是 tag 層, 不是 phasing 層 (v3 從原 v2-26 S8 拆出, 與 S7 LOH 兩層精確化分頁)。"
        "左欄 phasing 層 (LOH.bed): 路徑是 VCF allele depth (AD) → LongPhase region 偵測, 輸出 BED 格式 region-level LOH。"
        "V5 前後 Jaccard = 1.0 完全不變, V2b PON-only phasing 已解 phase scaffold self-phasing。"
        "右欄 tag 層 (BAM HP_Ratio): 路徑是 BAM HP:i tag → ISM per-region HP_Ratio 統計。"
        "Baseline 17.3:1 嚴重偏差, 62% ISM HP_Ratio LOH 在 paired 比對下其實是 artifact, V5 修補對象就是這層。"
        "兩套 LOH 系統使用**不同數據源** (VCF AD vs BAM HP tag), kappa = 0.670 的不一致由此解釋。"
        "Q3 預備 (「Jaccard=1.0 怎麼證明?」): baseline LOH.bed 與 V2b PON-only LOH.bed 的 region-level Jaccard 相似度 = 1.0, "
        "兩個 BED 完全相同 → phasing 層 LOH region 偵測不受 self-phasing 影響。"
        "**重要區分**: 本工作影響的是 BAM HP tag → ISM HP_Ratio LOH (62% artifact), 不是 LOH.bed; "
        "LOH.bed 即使在 V5 後也不變。"
        "下一張 S7 進一步把 LOH 兩層的精確細節展開 (kappa = 0.670 來自不同數據源解釋)。")
    return s


def slide_07_loh_two_layers(prs):
    """S7: LOH 兩層精確化 (v3 從 S8 拆出)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "LOH 兩層精確化 — kappa = 0.670 解釋")
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.5)
    add_text_block(s, Inches(0.7), Inches(1.25), Inches(5.6), Inches(0.55),
                   ["ISM HP_Ratio LOH (read-level)"],
                   font_size=16, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(4.3),
                   ["數據源:  BAM HP:i tag",
                    "計算:  HP_Ratio < 0.1 or > 0.9",
                    "",
                    "受 self-phasing 影響嚴重:",
                    "  62% 是 artifact",
                    "  AF 0.1-0.8 近 100% artifact",
                    "  TO TP 中 86.5% paired 平衡",
                    "",
                    "本工作修補後改善"],
                   font_size=12, color=C_TEXT)

    add_rect(s, Inches(6.85), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.5)
    add_text_block(s, Inches(7.05), Inches(1.25), Inches(5.6), Inches(0.55),
                   ["LOH.bed (region-level)"],
                   font_size=16, bold=True, color=C_GREEN)
    add_text_block(s, Inches(7.05), Inches(1.85), Inches(5.6), Inches(4.3),
                   ["數據源:  VCF allele depth (AD)",
                    "計算:  LongPhase region 偵測",
                    "",
                    "不受 self-phasing 影響:",
                    "  V5 前後 Jaccard = 1.0",
                    "  region-level LOH 真值穩固",
                    "",
                    "本工作不影響此層"],
                   font_size=12, color=C_TEXT)
    add_text_block(s, Inches(0.5), Inches(6.55), Inches(12.5), Inches(0.5),
                   ["kappa = 0.670 不一致 = 兩套系統「不同數據源」自然差異, 非 bug"],
                   font_size=13, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 7)
    set_speaker_note(s,
        "LOH 兩層精確化 (v3 從 S8 拆出, 為避免常見誤解)。"
        "兩層的精確區分: 左欄 ISM HP_Ratio LOH (read-level) 與 右欄 LOH.bed (region-level)。"
        "**ISM HP_Ratio LOH (左欄)**: 數據源 = BAM HP:i tag; 計算方式 = 統計 region 中各 read 的 HP_Ratio, "
        "若 < 0.1 or > 0.9 標為 LOH。受 self-phasing 影響嚴重: 62% 是 artifact, "
        "AF 0.1-0.8 範圍近 100% artifact, TO TP 中 86.5% 在 paired 下其實平衡。本工作修補後改善。"
        "**LOH.bed (右欄)**: 數據源 = VCF allele depth (AD); 計算方式 = LongPhase 直接從 AD 偵測 region-level LOH。"
        "不受 self-phasing 影響: V5 前後 Jaccard = 1.0, region-level LOH 真值穩固。本工作不影響此層。"
        "**kappa = 0.670 解釋**: 兩套 LOH 系統使用不同數據源, 本就期待中度一致性 (kappa 0.4-0.7); "
        "0.670 不是 bug, 而是兩個獨立衡量的自然差異。"
        "**本工作 self-phasing 修補只影響 ISM HP_Ratio LOH (左欄), 不改變 LOH.bed (右欄)**。"
        "教授若混淆兩層, 會誤以為「修補後 LOH.bed 也變」, 這是錯的。"
        "下一張 S8 進入「為何重要」段, 解釋 TO 是研究根基 + ISM 4-bucket 強依賴 HP tag。")
    return s


def slide_08_why_important(prs):
    """S8: TO 根基 + 4-bucket (原 v2-26 S9)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "為何重要 — TO 根基 + ISM 強依賴 HP tag")
    add_rect(s, Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.4),
             fill=C_LIGHT, line=C_TITLE_BG, line_w_pt=2.0)
    add_text_block(s, Inches(0.7), Inches(1.25), Inches(5.1), Inches(0.5),
                   ["TO 模式是研究根基"],
                   font_size=18, bold=True, color=C_TITLE_BG)
    add_text_block(s, Inches(0.7), Inches(1.85), Inches(5.1), Inches(4.5),
                   ["臨床 / archive 樣本 normal 不可得",
                    "  → TO 模式是必要研究方向",
                    "",
                    "TO 純數據是 TP / FP 觀察根基",
                    "  → 不能依賴 paired 比對外推",
                    "",
                    "本實作:  TO + PON 條件下",
                    "  讓 tag 品質接近 paired 水準"],
                   font_size=13, color=C_TEXT)
    add_rect(s, Inches(6.3), Inches(1.1), Inches(6.5), Inches(5.4),
             fill=RGBColor(0xFD, 0xED, 0xEC), line=C_ACCENT, line_w_pt=2.0)
    add_text_block(s, Inches(6.5), Inches(1.25), Inches(6.1), Inches(0.5),
                   ["ISM 4-bucket 分群核心"],
                   font_size=18, bold=True, color=C_ACCENT)
    add_text_block(s, Inches(6.5), Inches(1.85), Inches(6.1), Inches(4.5),
                   ["·  HP1   (germline HP1 reads)",
                    "·  HP1-1 (somatic on HP1 reads)",
                    "·  HP2   (germline HP2 reads)",
                    "·  HP2-1 (somatic on HP2 reads)",
                    "",
                    "依賴於 4-bucket 的特徵:",
                    "  NGroups, HPSig, HP_Ratio_norm,",
                    "  HPMergedDelta, HPFineNGroups",
                    "",
                    "全部 → 依賴 HP tag 正確"],
                   font_size=13, color=C_TEXT)
    add_text_block(s, Inches(0.5), Inches(6.65), Inches(12.5), Inches(0.5),
                   ["整個 ISM 研究都建立在 TO 模式 + 4-bucket 分群"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 8)
    set_speaker_note(s,
        "為何重要這個 slide 解釋影響的學術意義。左欄: TO 模式為何是研究根基。"
        "(1) 臨床、archive、cell-free DNA、舊樣本等情境下 normal BAM 不可得, TO 是必要研究方向。"
        "(2) TO 純數據是 TP/FP 觀察的根基, 不能依賴 paired 比對外推 (paired 比對引入 normal 訊號污染獨立性)。"
        "(3) 本實作目標就是讓 TO + PON 條件下 tag 品質接近 paired 水準, 讓下游 ISM 研究有可信基礎。"
        "右欄: ISM 4-bucket 分群是核心架構。"
        "ISM 把 reads 分成 HP1 (germline HP1)、HP1-1 (somatic on HP1)、HP2 (germline HP2)、HP2-1 (somatic on HP2) "
        "四個 bucket, 計算 region-level 特徵。"
        "依賴 4-bucket 的關鍵特徵: NGroups (subclone count)、HPSig (HP signature)、HP_Ratio_norm、"
        "HPMergedDelta (HP family delta)、HPFineNGroups (subclone marker)。"
        "如果 HP tag 偏差 (17.3:1), 4-bucket 分群錯誤, 所有依賴特徵的結論都不可信。"
        "整個 ISM 研究都建立在這個基礎上, 因此 self-phasing 修補是研究可信度的前提。"
        "下一張 S9 量化影響範圍 — 哪些特徵嚴重 / 中度 / 不受影響, 重跑工作可聚焦在 43 個而非全 85 features。")
    return s


def slide_09_3tier_count(prs):
    """S9: ISM 3-tier 影響範圍 29/14/42 (原 v2-26 S10)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "影響範圍 — 29 / 14 / 42 個 (共 85)")
    fit_image_within(s, str(FIG_DIR / "03_self_phasing_impact.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.5))
    add_text_block(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4),
                   ["嚴重 29 個  ·  中度 14 個  ·  不受影響 42 個  ·  總計 85 features"],
                   font_size=13, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4),
                   ["註: 來源報告誤寫 14/85 = 7%, 實際 16.5%; slide 以 count 為準"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 9)
    set_speaker_note(s,
        "V5 修補對下游 ISM 特徵的影響量化為 3 tier (來源 source 03 §4.3)。"
        "嚴重影響 (HP-依賴) **29 個**: HP_Ratio (62% 假 LOH)、Potential_LOH、HPMergedDelta/Sig (方向反轉)、"
        "HPFineNGroups (含 NG=2 LOH-constrained phasing)。這 29 個必須在 V5 BAM 上重跑才能信任, "
        "舊結論需加註版本 (haplotag_version = baseline / V5)。"
        "中度影響 (間接污染) **14 個**: QualityScore (AUC 0.497, 已移除 LOH penalty)、GlobalP (取 HP/Allele 最小值)、"
        "CramersV (取 HP/Allele 最大值)、VerificationClass (label_sig 含 HP 成分)。重評多數影響微弱。"
        "不受影響 (無 HP 依賴) **42 個**: PairwiseMean/MedianDist (無 HP 分組)、AlleleDelta/AlleleP、"
        "Caller 特徵 (AF / GQ / DP / SB)、甲基化矩陣本身 (BAM MM/ML tag, 不依賴 HP)、CpG 座標、region_methyl_mean。"
        "這 42 個結論穩固不需重測。"
        "**口徑校準關鍵**: slide 上只列 count 不列 %。來源報告寫 14/85 = 7% 但 14/85 實際 16.5%, 來源百分比有誤。"
        "為避免傳遞錯誤數字, 本 slide 與 speaker note 全部以 count 表達, 並註明來源報告誤寫。"
        "總數 85 features。重跑工作可聚焦在嚴重 + 中度的 43 個, 不需全 85 重跑, 大量節省計算資源。"
        "下一張 S10 進入解釋段, 介紹 4 函數架構 + PON 概念 + Purity 0.95 觸發鎖。")
    return s


def slide_10_prereq_combined(prs):
    """S10: 4 函數 + PON + Purity 0.95 觸發鎖 (v3 合併原 S11+S12)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Prerequisite — 4 函數 + PON + Purity 觸發鎖")
    img = make_s10_prereq_combined()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.7))
    add_footer(s, 10)
    set_speaker_note(s,
        "v3 合併原 v2-26 S11 (4 函數表) + S12 (PON 概念) 為單一 prerequisite slide, 並補上 Purity 0.95 閾值概念。"
        "教授不熟 longphase-to-mod 內部架構, 本張為必補的 prerequisite。"
        "**左半 longphase-to-mod 4 函數**: "
        "(1) getVote() — 統計位點 reads 對 HP1/HP2/somatic 的投票, V3F + V5 修補主戰場 (HaplotagProcess.cpp:512-563)。"
        "(2) judgeHaplotype() — caller 端, 將 vote → hpResult integer, V3F enum 修在這 (line 697)。"
        "(3) countSNPHaplotype() — 計算 SNP 位點各 hap 的 read 計數, V5 補 alt path UNDEFINED guard (line 489-494)。"
        "(4) countINDELHaplotype() — INDEL 版本, 380e8d2 commit 補 UB guard (line 497-510)。"
        "**右半 PON 概念**: PON (Panel of Normals) 是群體 germline 變異資料庫, 包含 1000 Genomes、CoLoRSdb、dbSNP、gnomAD; "
        "在 TO 模式無 normal BAM 時 PON 替代 normal 作 phasing anchor。"
        "V2b 啟用 --pon-only-phasing 解 phase scaffold self-phasing。"
        "**★ Purity 0.95 觸發鎖 (v3 補充, S11 主者預告)**: PhasingProcess.cpp:197 硬編碼 `purity > 0.95` 閾值, "
        "決定是否啟動 baseline 內建的 Two-Pass 路徑。HCC1395 5kHz 真實 purity > 0.95 但 baseline 估計 0.927 "
        "(四捨五入 0.93, 非誤判, 是邊緣 case), 0.927 ≤ 0.95 → 未觸發 Two-Pass → 走「三條路」標準流程。"
        "Q6 預備 (「PON 有資料還會 self-phasing?」): PON classification (VCF 層) 與 phasing anchor (read 連結層) 是獨立步驟。"
        "Baseline bug: PON 雖標出 somatic 但 phasing 階段仍把 somatic 當 anchor (--pon-only-phasing=false 預設)。"
        "V2b 啟用後 phasing 才只用 PON-confirmed germline 做 anchor。"
        "「修一個 bug 暴露另兩個 bug」: PON-only 啟用後 germline votes 急減, 原本被 paired germline 訊號掩蓋的 "
        "tag 層 priority bug + enum bug 立刻顯形。下一張 S11 用 root-cause tree 視覺化這個觸發條件 + 三 bug。")
    return s


def slide_11_root_cause_v3(prs):
    """S11 (v3 ★ KEY): Baseline 根因樹含觸發條件主者."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "★ Baseline 根因樹 (觸發條件主者 + 三 bug)")
    img = make_s11_root_cause_with_trigger()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(5.7))
    add_footer(s, 11)
    set_speaker_note(s,
        "v3 ★ KEY slide — 加入 Purity 0.95 觸發條件主者作為 root-cause tree 的根。"
        "Self-phasing 不是單一 bug, 而是「特定觸發條件下三層獨立故障的集中暴露」。"
        "**根 (觸發條件主者, 紅色框)**: Baseline 估 purity = 0.927 ≤ 0.95 (純樣本標準) → "
        "未觸發 Two-Pass → 走 baseline 標準三條路流程。"
        "**口徑校準關鍵**: 不是「baseline 誤判 purity」(0.927 是正確估計, 是邊緣 case), 而是 baseline 設計上的硬編碼閾值 "
        "0.95 對近似純樣本不夠保守。HCC1395 5kHz 真實 purity > 0.95 (純 tumor, 無 normal contamination), "
        "但 baseline 的 polynomial calculator 估出 0.927, 這個 0.927 數值本身正確, "
        "只是落在 0.95 閾值的邊緣下方就觸發了 baseline 標準流程而非 Two-Pass 路徑。"
        "**三葉 (bug)**: "
        "Bug 1 (紅色, getVote priority): variantKeys 順序 {HP1_1, HP2_1} 排第一, 任一 somatic count > 0 即 break, "
        "germline 完全跳過。為何 paired 沒事? paired germline votes 多但 priority bug 仍存在; "
        "TO PON-only 後 germline votes 急減立刻顯形 (HCC1395 99.9% reads 拿 HP21)。V3F commit 41ff147 修。"
        "Bug 2 (橘色, enum vs int literal mismatch): C++ enum 定義在 Util.h:20-25 (注意是 20-25 不是 21-25, "
        "HAPLOTYPE_UNDEFINED=-1 起於 line 20), HAPLOTYPE1_1 = 3。但 BAM HP tag integer = 11。"
        "caller 端 if(hpResult != HAPLOTYPE1_1) 用 enum=3 比較 hpResult=11/21/33, fallback 死分支永不執行, "
        "HP:i:33 永不出現 (baseline 0/15 sites)。C++ 把 enum 隱式轉 int 不會 warn。V3F 同 commit 改為 integer literal 比較。"
        "Bug 3 (灰色, scaffold, 已由 V2b 解): PhasingProcess.cpp:154-157 convertNonGermlineToSomatic() 不被觸發, "
        "somatic 當 phasing anchor。已由 V2b commit 8b8c1fd 解決 (透過 --pon-only-phasing flag), 本 PPT 焦點不在此, 視覺只佔 20%。"
        "**因樣本實為純無 normal contamination 而流程假設有 → 暴露 tag 層 somatic-first 投票 bias → 結果 self-phasing 17.3:1 artifact**。"
        "驗證來源: source_materials/04_purity_calculator_failure_root_cause.md (v5_audit_suite/18 複製)。"
        "下一張 S12 講 V5 三層投票流程如何修補 Bug 1, 並整合 Q2 答 (V5 在 mid-low purity 比 baseline 更好)。",
        min_chars=500)
    return s


def slide_12_v5_layers_v3(prs):
    """S12 (v3 ★ KEY): V5 三層投票流程 + mid-low purity 防禦層."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "改動方式 — V5 三層投票流程")
    img = make_s12_v5_three_layer()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.7))
    add_footer(s, 12)
    set_speaker_note(s,
        "v3 ★ 強化版 — speaker note 整合 Q2 答 (V5 在 mid-low purity 比 baseline 更好)。"
        "V5 對 getVote() 的具體改動: 從 V3F 兩層升級為三層投票邏輯。"
        "**Layer 1 (germline first)**: if (germlineHP1 > germlineHP2) → 1; elif (germlineHP2 > germlineHP1) → 2; else → 進入 Layer 1.5。"
        "保持 germline first 核心原則。**這裡是 baseline 的 somatic-first 投票順序的關鍵翻轉**: "
        "baseline getVote() 用 variantKeys loop 把 {HP1_1, HP2_1} 排第一, somatic 投票優先; "
        "V5 直接看 germline 投票數, 翻轉投票優先序。"
        "**Layer 1.5 (somatic fallback, V5 新增, mid-low purity 防禦層)**: 當 germline 平手時參考 somatic 投票。"
        "if (somaticTotal == 0) → hpResult = 0 (unphased); "
        "elif (somaticHP1 > 0.6 * somaticTotal) → 對齊 HP1 方向 (germlineResult = 1); "
        "elif (somaticHP2 > 0.6 * somaticTotal) → 對齊 HP2 方向; "
        "else → hpResult = 33 (真正 ambiguous)。"
        "Confidence threshold 0.6 是經驗值 (R4 caveat), **設計上是 mid-low purity 的防禦層** — "
        "在 0.6-0.93 中等純度範圍, 弱 directional 訊號標 33 ambiguous (保守 +10pp HP33 比例) 避免錯誤分配。"
        "**Layer 2 (encode hpResult)**: germline=1 + somatic_total>0 → HP:i:11; germline=2 + somatic_total>0 → HP:i:21; "
        "germline=0 + somatic_total>0 → HP:i:33; germline only → HP:i:1 or 2; nothing → HP:i:0。"
        "**★ Q2 答 (V5 vs baseline 在 mid-low purity)**: V5 更好。"
        "純樣本 (HCC1395 5kHz): baseline somatic-first 在 PON-only 啟用後 germline votes 急減 → somatic 主導 → 17.3:1 artifact; "
        "V5 germline-first 不受此影響。"
        "**Mid-low purity 場景 (0.6 sample, t30_n20)**: baseline 仍 somatic-first 但 self-phasing 自然弱化 (normal contamination 平衡 hap); "
        "V5 加入 confidence 0.6 threshold → 對「弱 directional 訊號」標 HP33 ambiguous → "
        "09 報告: 0.6 sample HP33 比例 12.4% vs Baseline 2% (+10pp 保守) → 避免錯誤分配。"
        "效果: AMB% 17.5%→8.0% (-9.5 pp); HP:i:33 reads 239,679→110,197 (-54%); "
        "介面契約零變動 (HaplotagProcess.h:66-68 一字未變, drop-in replace)。"
        "Caveat: V5 在 SP-extreme self-phasing 位點仍不修 (V5 設計就不修 phase 層, 那是 V2b PON-only 的責任)。"
        "下一張 S13 用程式碼 diff 對照 baseline 與 V5 兩欄, 標出 baseline somatic-first 紅標 caveat。",
        min_chars=500)
    return s


def slide_13_code_diff_v3(prs):
    """S13 (v3): 程式碼 diff baseline somatic-first 紅標."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "程式碼 diff (Baseline vs V5)")
    img = make_s13_code_diff()
    fit_image_within(s, str(img), Inches(0.1), Inches(1.0), Inches(13.2), Inches(5.6))
    add_text_block(s, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.4),
                   ["三處紅標: priority loop 拆解 (somatic-first → germline-first)  ·  Layer 1.5 純分支  ·  enum → int literal 11"],
                   font_size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 13)
    set_speaker_note(s,
        "v3 修正: baseline 欄上方加紅標 caveat 「投票邏輯 = somatic-first 優先序 (非「有 HP11/HP21 即 somatic」這個簡化敘述)」。"
        "程式碼 diff 兩欄並列。"
        "**左欄 Baseline (有 bug)**: "
        "for (auto& vk : variantKeys) — variantKeys 順序 {HP1_1, HP2_1} 排第一; "
        "if (count[vk.first] > 0 || count[vk.second] > 0) break — 任一 somatic 即 break, germline 跳過 (Bug 1 priority); "
        "if (hpResult != HAPLOTYPE1_1) hpResult = 0 — 用 enum=3 比較 hpResult=11/21/33, 永不匹配 (Bug 2 enum)。"
        "**★ V3 紅標 caveat**: baseline 投票邏輯不是「有 HP11/HP21 即 somatic」這個簡化敘述, 而是「somatic-first 投票優先序設計缺陷」 — "
        "variantKeys 排序把 somatic key 排在前, 任一 somatic 投票觸發 break, 結果 germline 投票被忽略。"
        "這個區分對教授很重要: 簡化敘述會讓人誤以為「baseline 一旦看到 somatic tag 就分類為 somatic」, "
        "但實際是「投票順序決定優先權」, V5 修法是「翻轉投票優先序為 germline-first」, 不是「忽略 somatic」。"
        "**右欄 V5 (修補後)**: "
        "// Layer 1: germline first — if (g1 > g2) germlineRes = 1; else if (g2 > g1) germlineRes = 2; "
        "// Layer 1.5: somatic fallback — else if (s1 + s2 > 0) { 0.6 threshold align HP1/HP2 }; "
        "if (hpResult != 11) — 改 integer literal 比較。"
        "**三處紅標對應前面 S11 root-cause tree**: "
        "(1) priority loop 拆解 (somatic-first → germline-first 翻轉, Bug 1); "
        "(2) Layer 1.5 純分支新增 (V5 新工作, mid-low purity 防禦層); "
        "(3) enum HAPLOTYPE1_1 改為 integer literal 11 (Bug 2)。"
        "Q13 預備 (「程式碼 diff 看不太懂」): 請看這三處紅標, 每行對應前面 S11 一條 bug。"
        "monospace 24 pt 是用戶指定字型大小, 雙欄 fit-within 等比縮放。"
        "完整程式碼參考見 notes/code_references.md, V5 commit hash 待 working tree commit 後補 (R5)。"
        "下一張 S14 進入驗證段, 4 項 sanity check 15/15 PASS + 5 條獨立證據鏈 + Q1 PON 雙路徑有效性整合。")
    return s


def slide_14_sanity_5evidence_v3(prs):
    """S14 (v3 ★ KEY): Sanity 4 項 + 5 證據鏈 + Q1 PON 雙路徑."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Sanity check 4 項 PASS + 5 條獨立證據鏈")
    add_rect(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=RGBColor(0xE8, 0xF8, 0xF0), line=C_GREEN, line_w_pt=2.0)
    add_text_block(s, Inches(0.7), Inches(1.25), Inches(5.6), Inches(0.5),
                   ["4 項硬性 sanity check"],
                   font_size=18, bold=True, color=C_GREEN)
    sanity_rows = [
        ("守恆律 A · Δ-consistency", "ΔHP33 + (ΔHP11 + ΔHP21) = 0", "15/15 PASS"),
        ("守恆律 B · Germline 不變", "HP1/HP2 reads V3F vs V5 比對", "15/15 PASS"),
        ("Layer 1.5 期望 1 · 33→directional 守恆", "ΔHP11 == n(33→11)", "15/15 PASS"),
        ("Layer 1.5 期望 2 · 無 germline → HP33", "germlineResult ≠ 0 不出 HP33", "0 violation"),
    ]
    cy = Inches(1.85)
    for name, desc, verdict in sanity_rows:
        add_text_block(s, Inches(0.7), cy, Inches(5.6), Inches(0.35),
                       [name], font_size=12, bold=True, color=C_TEXT)
        add_text_block(s, Inches(0.7), cy + Inches(0.30), Inches(4.0), Inches(0.3),
                       [desc], font_size=10, color=C_GRAY)
        add_text_block(s, Inches(4.7), cy + Inches(0.10), Inches(1.6), Inches(0.4),
                       [verdict], font_size=11, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)
        cy += Inches(1.10)

    add_rect(s, Inches(6.85), Inches(1.1), Inches(6.0), Inches(5.4),
             fill=C_LIGHT, line=C_TITLE_BG, line_w_pt=2.0)
    add_text_block(s, Inches(7.05), Inches(1.25), Inches(5.6), Inches(0.5),
                   ["5 條獨立證據鏈"],
                   font_size=18, bold=True, color=C_TITLE_BG)
    chains = [
        ("1. 理論層", "germline het 隨機 → 預期 1:1"),
        ("2. 全基因組層", "HP1:HP2 = 17.3:1 跨 23 chr 一致"),
        ("3. 個別位點層", "SP1 chr19:17565944 113:0"),
        ("4. Sanity check", "15/15 PASS, 0 violation"),
        ("5. 程式碼最小必要", "+68/-36 行 / 3 函式集中"),
    ]
    cy = Inches(1.85)
    for tag, desc in chains:
        add_text_block(s, Inches(7.05), cy, Inches(5.6), Inches(0.35),
                       [tag], font_size=12, bold=True, color=C_TITLE_BG)
        add_text_block(s, Inches(7.05), cy + Inches(0.30), Inches(5.6), Inches(0.3),
                       [desc], font_size=10, color=C_TEXT)
        cy += Inches(0.85)
    add_text_block(s, Inches(0.5), Inches(6.65), Inches(12.5), Inches(0.5),
                   ["總計 15 / 15 PASS  ·  5 條獨立路徑同步收斂  ·  PON-only Two-Pass 在 0.93 與 0.6 都有效"],
                   font_size=13, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_footer(s, 14)
    set_speaker_note(s,
        "v3 ★ 強化版 — speaker note 整合 Q1 答 (PON-only Two-Pass 在 0.93 與 0.6 sample 都有效)。"
        "驗證段第一張: sanity check + 5 證據鏈彙整。"
        "**左欄 4 項硬性 sanity check** (HCC1395 5kHz 15 sites, v5_audit_suite/06 報告): "
        "(1) 守恆律 A · Δ-consistency: ΔHP33 + (ΔHP11 + ΔHP21) = 0, tag 移轉量平衡, in = out, 15/15 PASS。"
        "(2) 守恆律 B · Germline 不變: V3F 與 V5 的 HP1 / HP2 reads 逐 site 比對, 0 reads 漂移, 15/15 PASS。"
        "(3) Layer 1.5 期望 1 · 33→directional 精確守恆: ΔHP11 == n(V3F=33→V5=11)、ΔHP21 == n(V3F=33→V5=21), "
        "15/15 PASS (V5max1 chr19:4639528 39 reads 精確守恆)。"
        "(4) Layer 1.5 期望 2 · 無 germline → HP33: 跨 15 sites pool transition table, germlineResult ≠ 0 不出 HP33, 0 violation。"
        "**右欄 5 條獨立證據鏈**: "
        "(1) 理論層 — germline het 隨機 → 預期 1:1; "
        "(2) 全基因組層 — HP1:HP2 = 17.3:1 跨 23 chr 一致; "
        "(3) 個別位點層 — SP1 chr19:17565944 baseline 113:0 全失衡; "
        "(4) Sanity check 4 項 15/15 PASS (左欄); "
        "(5) 程式碼最小必要 — +68/-36 行集中於 3 函式 (getVote/countSNPHaplotype/countINDELHaplotype)。"
        "5 條獨立路徑同步收斂於同一結論: self-phasing 修補必要且充分。"
        "**★ Q1 答 (PON-only Two-Pass 在不同 purity 下都有效)**: "
        "(a) 0.93 純樣本 (HCC1395 5kHz): PON-only Two-Pass 可執行, sanity 15/15 PASS, "
        "Aggregate paired GT concordance V5 78.85% vs Baseline 72.20% (+6.65pp), "
        "Clean PS paired GT V5 88.2% vs Baseline 74.9% (+13.3pp), 全基因組 clean PS V5 90.5% vs BL 82.2% (+8.3pp)。"
        "(b) 0.6 sample (t30_n20): PON-only Two-Pass 可執行, 09 報告 V5 HP33 比例 12.4% vs Baseline 2% "
        "(+10pp 保守 tagging) → 在 mid-low purity 也有改善, 設計上的 conservative tagging 防禦層。"
        "Caveat: 0.6 sample 沒有 paired reference (synthetic mix 無 ground truth) 無法直接做 concordance, "
        "但 HP33 比例分析顯示 V5 確實在 mid-low purity 提供保守 tagging 防禦。"
        "Q7 預備 (「sanity check 是誰寫的? 覆蓋率多少?」): 本工作新增 (v5_audit_suite/06 Agent D), "
        "覆蓋限制是 15 sites cherry-picked (R3), 需 7 樣本擴展 (F3) + 100 隨機位點 cross-validate (F8); "
        "整體 sanity check 必要不充分。下一張 S15 進入量化指標 fig18。",
        min_chars=500)
    return s


def slide_15_quant_metrics(prs):
    """S15 (v3): 量化指標 (原 v2-26 S17 拆出上半)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "量化指標 — AMB / HP33 / Concordance")
    fit_image_within(s, str(FIG_DIR / "fig18_concordance_amb_f1.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.7))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["AMB% 17.5→8.0% (-9.5pp)  ·  HP:i:33 reads -54%  ·  +8.3 pp clean PS, 全基因組"],
                   font_size=12, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 15)
    set_speaker_note(s,
        "v3 從原 v2-26 S17 拆出, 專注量化指標。下一張 S16 才講 4-commit timeline。"
        "量化指標彙整: "
        "(1) AMB% (HP:i:33 比例) 從 17.5% 降到 8.0% (-9.5 pp), V5 解了過半 ambiguous reads。"
        "(2) HP:i:33 reads 數從 239,679 降到 110,197, 減 54%。"
        "(3) 全基因組 clean PS blocks paired GT concordance: V5 = 90.5% vs Baseline = 82.2%, **+8.3 pp**。"
        "**口徑校準關鍵**: +8.3 pp 必須加 caveat 「(clean PS blocks, 全基因組)」, "
        "不是全基因組 raw, 不是 cherry-picked。"
        "(4) 15 sites cherry-picked 上 clean PS concordance: V5 = 88.2% vs BL = 74.9%, +13.3 pp。"
        "**Q8 預備** (「N 多少? 信賴區間? 跨樣本一致?」): 15-site cherry-picked N=11 clean PS; "
        "全基因組 clean PS N 在數萬 - 數十萬 reads 量級 (具體 N 在 PI 報告 4 §3.7); "
        "跨樣本擴展未做 (F3 待辦); 信賴區間 source 03 未直接給, 推估 Wilson 95% CI ±1 pp 內。"
        "Caveat: clean PS blocks 篩選 germline accuracy ≥ 70%; problem PS blocks 上 V5 與 Baseline 接近隨機 "
        "(read-level orientation 不穩定)。"
        "TP Balanced% 從 V2b 13.0% 改善 (V5 同 V3F 基準)。"
        "**注意**: 這些是 read-level metric, 不是 caller F1 metric — F1 不變是預期 (見 S19)。"
        "下一張 S16 用 fig01a 展示 4-commit timeline (V2b → V3F → INDEL guard → V5)。")
    return s


def slide_16_commit_timeline(prs):
    """S16: 4-commit timeline (原 v2-26 S17 拆出下半)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "4-commit 漸進演進 timeline")
    fit_image_within(s, str(FIG_DIR / "fig01a_commit_evolution.png"),
                     Inches(0.4), Inches(1.0), Inches(12.5), Inches(5.7))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["V2b (PON-only) → V3F (priority + enum) → INDEL guard → V5 (Layer 1.5)"],
                   font_size=12, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 16)
    set_speaker_note(s,
        "v3 從原 v2-26 S17 拆出, 專注 4-commit timeline 視覺化。"
        "**4-commit 漸進演進** (路徑: longphase-to-mod fork @ /big7_disk/liaoyoyo2001/longphase-to-mod/, 獨立 git repo): "
        "(1) **V2b commit 8b8c1fd** (PON-only) — 解 phase 層 self-phasing scaffold, "
        "改 Phasing.cpp +9/-2、PhasingGraph.cpp +34/-0、PhasingProcess.cpp +25/-3; HaplotagProcess.cpp 未動。"
        "啟用 --pon-only-phasing flag 讓 phasing graph 只用 PON-confirmed germline 做 anchor。"
        "(2) **V3F commit 41ff147** (priority + enum) — 解 tag 層 Bug 1 (priority) + Bug 2 (enum mismatch), "
        "HaplotagProcess.cpp:506-541 (getVote 重寫) + :697 (caller 端 enum → integer literal); +36/-25。"
        "(3) **INDEL guard commit 380e8d2** — 補 UB 漏洞 (countINDELHaplotype), HaplotagProcess.cpp:497-510; +8/-4。"
        "(4) **V5 working tree (未 commit)** — 解 V3F 過於保守, HaplotagProcess.cpp:489-494 (countSNPHaplotype alt guard) "
        "+ :512-563 (getVote Layer 1.5 三層); +24/-7。"
        "**累計**: +68/-36 行集中於 3 函式 (getVote / countSNPHaplotype / countINDELHaplotype); "
        "介面契約 HaplotagProcess.h:66-68 一字未變 (drop-in replace)。"
        "**V5 working tree 未 commit (R5 caveat)**, 後續 P0 行動 F1 切 2 獨立 commits 完成 + tag v5.0。"
        "F1 切 2 commits: (A) Layer 1.5 somatic fallback (行 512-563); (B) countSNPHaplotype alt guard (行 489-494)。"
        "下一張 S17 進入演講高潮 V5max1 + caveat 防誤解。")
    return s


def slide_17_v5max1_climax(prs):
    """S17 (v3): V5max1 CLIMAX + caveat (原 v2-26 S18)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "V5max1 — 39 reads, 100% reassigned")
    fit_image_within(s, str(FIG_DIR / "C_V5max1_chr19_4639528.png"),
                     Inches(0.6), Inches(1.0), Inches(12.0), Inches(5.0))
    # Caveat strip
    add_rect(s, Inches(0.4), Inches(6.10), Inches(12.5), Inches(0.85),
             fill=RGBColor(0xFE, 0xF5, 0xE7), line=C_AMBER, line_w_pt=1.5)
    add_text_block(s, Inches(0.5), Inches(6.18), Inches(12.3), Inches(0.35),
                   ["chr19:4639528  ·  V3-F 紫色 HP33 群 (39 reads) → V5 全部正確重分配為 HP11"],
                   font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                   ["但 V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold) — V5 解的是 tag 層 directional reassignment"],
                   font_size=10, color=C_AMBER, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 17)
    set_speaker_note(s,
        "演講 CLIMAX。chr19:4639528 V5max1 是最戲劇化的視覺證據。"
        "在 V3-Fixed panel 看到一大群紫色 reads 標為 HP:i:33 (ambiguous), 39 條; "
        "V5 修補後這 39 條全部正確重分配為 HP:i:11 (somatic on HP1)。"
        "守恆律驗證: V3F 39 reads HP33 → V5 39 reads HP11, in = out 完全平衡, sanity check PASS。"
        "這是 Layer 1.5 somatic fallback 的具體效果 — 原本因為 V3F 過於保守卡在 ambiguous bucket 的 reads, "
        "現在透過 confidence 0.6 threshold 被正確歸到具體 hap。"
        "看 IGV 4-BAM 並列圖, 從 baseline、V2b、V3F、V5 的視覺差異很明顯。"
        "**底部 caveat 防誤解 (reviewer R2 採納)**: 「但 V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold)」。"
        "**口徑校準關鍵**: 高潮圖容易讓教授誤以為 V5 全面解決 self-phasing, 但 self-phasing 問題鏈是被分層處理的: "
        "phase 層的 self-phasing scaffold 由 V2b PON-only phasing 處理 (改 phasing graph anchor 來源, "
        "從含 somatic 改為僅 PON-confirmed germline); "
        "V3F/V5 解的是 tag 層: V3F 解 priority bug + enum mismatch (Bug 1+2); "
        "V5 解 V3F 過於保守 (germline 平手時 fallback 到 33), 新增 Layer 1.5 directional reassignment "
        "把卡在 33 的 reads 解到 11/21。"
        "所以 V5 的「39 reads → 100% reassigned」是 directional reassignment 不是 phasing graph 重建。"
        "speaker 必須在 60-90 sec 內一句說清這個分層, 避免教授誤以為「V5 = 解所有」。"
        "下一張 S18 用 bar chart 打破 IGV 視覺單調, 展示 SP-extreme 翻轉證據。")
    return s


def slide_18_hp_flip_bar(prs):
    """S18 (v3): HP 翻轉 bar (3 SP-extreme, 原 v2-26 S19)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "HP 翻轉證據 — Bar chart + IGV 縮圖")
    img = make_s17_hp_flip_bar()
    fit_image_within(s, str(img), Inches(0.3), Inches(1.0), Inches(8.5), Inches(5.7))
    fit_image_within(s, str(FIG_DIR / "D_SP1_chr19_17565944.png"),
                     Inches(9.0), Inches(1.0), Inches(4.0), Inches(5.7), border=True)
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["3/3 SP-extreme 一致翻轉  ·  V5 與 paired ground truth 一致"],
                   font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 18)
    set_speaker_note(s,
        "S17 climax 用單一 IGV 圖, 視覺單調; S18 用 bar chart 加上 IGV 縮圖打破單調 (climax 延伸)。"
        "Bar chart (左): 3 個 SP-extreme 位點, 每位點 3 組 bar (Baseline / V5 / Paired)。"
        "SP1 chr19:17565944 baseline 113 reads 全 HP1 (HP2 = 0), V5/V2b 翻轉為 HP2 主導 (113:0 → 0:113), "
        "paired tumor BAM 確認 HP2 為真實方向 (110 HP2 vs 2 HP1) — V5 與 paired ground truth 一致。"
        "SP2 chr19:12452332: baseline 109:1 → V5 1:109 → paired 105:3 確認 HP2。"
        "SP3 chr19:12467180: baseline 108:0 → V5 0:108 → paired 100:1 確認 HP2。"
        "3 個 SP-extreme 位點 (SP1/SP2/SP3) 全部一致地翻轉至 paired 方向, 是跨位點一致性證據。"
        "右側小 IGV 縮圖 (D_SP1) 展示 visual ground truth, 從 baseline → V2b → V3F → V5 的 4-BAM 並列。"
        "**重要**: 為什麼 baseline orientation 跟 paired 完全相反? 因為 self-phasing 把 reads 歸到錯邊 "
        "(sub-clone somatic 集合反客為主), 修補後 phasing 正確以 germline 為 anchor, reads 回到正確 hap。"
        "這是 V2b 解 phase scaffold 的效果, V3F/V5 確保 tag 層 directional 正確。"
        "speaker note 提醒: SP1-3 的修補主要是 V2b 主導 (phase 層), V5 在這些位點不再變動 (Δ ≈ 0); "
        "V5 主戰場是 V3F 過度保守的 33 → 11/21 (S17 V5max1)。"
        "下一張 S19 釐清為什麼這些大改善 F1 還是不變。")
    return s


def slide_19_f1_clarify(prs):
    """S19 (v3): F1 釐清 + cnLOH 邊界 (原 v2-26 S20)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "矛盾或盲點 — F1 釐清 + cnLOH 邊界")
    img = make_s19_f1_decomposition()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.7))
    add_footer(s, 19)
    set_speaker_note(s,
        "教授必問 (Q9): F1 不變是預期 — 那為什麼還要做這修補? 對發論文有意義嗎?"
        "**為什麼 F1 不變**: ClairS-TO raw F1 = 0.7166 對所有版本 (Baseline / V3F / V5) **完全相同**, "
        "因為 V5 不改 caller, 改的是 BAM HP tag 編碼, 與 caller 輸出 VCF 無關。"
        "F1 變動只可能來自下游 ISM SuggestFilter (RegionProcessor.cpp:1120, 1269), "
        "但 ISM TO F1 增益本來只有 0.0124 (對 caller germline FP 修復力先天不足), V5 vs Baseline 最終 F1 差 -0.0003 噪音。"
        "**結論**: F1 不能衡量本實作品質, 真實價值在 read-level concordance +8.3 pp (clean PS blocks, 全基因組)。"
        "**為什麼仍要做**: "
        "(1) 下游 ISM 影響 — 29 個 HP-依賴特徵必須在 V5 BAM 上才可信, 舊結論需重跑或加註版本。"
        "(2) 生物學詮釋正確性 — HPFineNGroups marker 過去解讀為 methylation bimodality, "
        "V5 修補後重詮釋為 phasing bucket signature (Thread D 主軸候選)。"
        "(3) 跨樣本研究可信度 — 7/7 樣本同方向 self-phasing 影響全部研究方向 (五大目標 1/2/4)。"
        "**對發論文意義**: 本工作不是新 caller / 新 phaser, 是填補 LongPhase-TO 在 PON-only 啟用後的 tag 層 bug; "
        "與 LongPhase-S 2025 paired 形成姐妹工作 (同實驗室哲學一致); "
        "真實 contribution 是 read-level tag 品質 +8.3 pp 全基因組 clean PS, 對下游 epigenetic analysis 是 enabling step; "
        "學術定位是 InterSubMod 五大目標的解鎖前提 (特別是目標 1/2/4)。"
        "**Caveat 4 條**: cnLOH 區未獨立評估 (R1, 雙親同源無 het, V5 fallback 也無法 anchor); "
        "AF > 0.9 邊界 (R2); Confidence threshold 0.6 未直接驗證 (R4); V5 working tree 未 commit (R5)。"
        "下一張 S20 業界家族樹 + 2x2 合併 (吸收原 v2-26 S2 業界內容)。")
    return s


def slide_20_industry_combined(prs):
    """S20 (v3): 業界家族樹 + 2x2 合併 (原 S2 + S21 合併)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "業界家族樹 + 2x2 比較表")
    img = make_s20_industry_combined()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.8))
    add_text_block(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
                   ["填補 TO + PON 場景 gap  ·  非「業界共識」  ·  LongPhase-S +4.5% 為 paired 場景參考"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 20)
    set_speaker_note(s,
        "v3 合併 — 上半業界家族樹 (原 v2-26 S2) + 下半 2x2 比較表 (原 v2-26 S21) 為單一 slide。"
        "**上半業界家族樹**: 上游基底是 LongPhase (Lin 2022 Bioinformatics) 處理 germline phasing。"
        "由此分出兩個分支: "
        "(a) LongPhase-S (bioRxiv 2025.11.20.689492v1) 是同實驗室的 paired 版, 在 ClairS 上 SNV F1 +4.5% "
        "(indel +7.1%), 提供「somatic 錨在 germline scaffold」的 anchoring 概念。"
        "(b) longphase-to-mod V5 是本實作, 是 TO + PON 條件下的本地實作。"
        "**下半 2x2 matrix** (Source × Mode): "
        "公開工具 + Germline (左上): LongPhase / WhatsHap / HapCUT2 — 成熟 germline phasing。"
        "公開工具 + Tumor (右上): LongPhase-S — Paired only (需 normal), ClairS SNV F1 +4.5% paired 場景。"
        "本實作 + Germline (左下): 不做 germline, 上游用 ClairS-TO + LongPhase。"
        "本實作 + Tumor (右下): longphase-to-mod V5 — Tumor-Only + PON, +8.3 pp clean PS concordance, "
        "TO + PON 條件下的本地實作。"
        "**口徑校準關鍵**: 本實作填補 tumor-only 在公開工具中的 gap (WhatsHap 不支援 TO), "
        "**不是「業界共識」「標準替代」**, 而是「同實驗室相鄰工作」。"
        "LongPhase-S 為 paired 場景參考 (注意非直接可比, 因本工作 F1 不變是預期行為), 不重疊本工作 TO 場景。"
        "Q10 預備 (「跟 LongPhase-S 重疊多少? 算 contribution 嗎?」): "
        "LongPhase-S 是 paired, 本實作是 tumor-only (PON 替代 normal); 設計哲學一致, 本工作填補 TO gap; "
        "Contribution 區隔: 4-commit 漸進修補在 longphase-to-mod fork (獨立 git repo), "
        "修補對象是 LongPhase-TO 在 PON-only 啟用後集中暴露的 3 層 tag-side bug, "
        "4 項 sanity check 是新貢獻, ISM 下游影響 3-tier 分類是新貢獻, ISM 跨樣本一致性 7/7 + Cohen's d=-1.20 是新貢獻。"
        "下一張 S21 進入未來規劃段, 後續可動 / 已初步現方向合併。")
    return s


def slide_21_followup_combined(prs):
    """S21 (v3): 後續可動 / 已初步現 合併 (原 v2-26 S22+S24 合併)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "後續可動 + 已初步發現的方向")
    # Top half: 4-row 後續可動
    add_text_block(s, Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.4),
                   ["可一起做的事 (依 V5 完成後)"],
                   font_size=15, bold=True, color=C_TITLE_BG)
    items = [
        ("Phase 2A — Normal methylation",
         "依賴 V5 BAM; HP tag 正確後 normal-only methylation 可用 (目標 4 解鎖)"),
        ("HPFineNGroups marker 重驗",
         "subclone marker ⭐4→⭐3 已降級; V5 BAM 上 master × flag 重驗"),
        ("Archive haplotag_version 標記",
         "manifest.yaml 加 haplotag_version 欄位; baseline / V5 區分可追溯"),
        ("Thread D — LOH-constrained phasing 主軸",
         "NG=2 cross-sample 6/6 POSITIVE, B1 Wilcoxon p=0.0156 (論文主軸候選)"),
    ]
    cy = Inches(1.45)
    rh = Inches(0.55)
    for title, desc in items:
        add_rect(s, Inches(0.5), cy, Inches(12.4), rh,
                 fill=C_LIGHT, line=C_TITLE_BG, line_w_pt=1.0)
        add_text_block(s, Inches(0.7), cy + Inches(0.05), Inches(5.0), Inches(0.45),
                       [title], font_size=12, bold=True, color=C_TEXT)
        add_text_block(s, Inches(5.8), cy + Inches(0.10), Inches(7.0), Inches(0.4),
                       [desc], font_size=10, color=C_GRAY)
        cy += rh + Inches(0.05)
    # Bottom half: 3 cards 已初步現
    add_text_block(s, Inches(0.5), Inches(4.30), Inches(12.5), Inches(0.4),
                   ["已初步發現 (全部依賴 V5 BAM 可信)"],
                   font_size=15, bold=True, color=C_ACCENT)
    cards = [
        ("Thread D NG=2",
         "6/6 POSITIVE",
         "Wilcoxon p=0.0156\nmedian gap 0.365",
         C_ACCENT),
        ("HPFineNGroups",
         "重驗中 (⭐3)",
         "Phase 2B master ×\nflag 機制本質判定",
         C_AMBER),
        ("LOH-constrained phasing",
         "論文主軸候選",
         "Inner same-hap\n93-99% obs18 6 樣本",
         C_GREEN),
    ]
    cw = Inches(4.05)
    cy = Inches(4.85)
    cx = Inches(0.5)
    for title, sub1, sub2, color in cards:
        add_rect(s, cx, cy, cw, Inches(1.95),
                 fill=RGBColor(0xEC, 0xF0, 0xF1), line=color, line_w_pt=2.0)
        add_text_block(s, cx + Inches(0.15), cy + Inches(0.10), cw - Inches(0.3), Inches(0.4),
                       [title], font_size=12, bold=True, color=color)
        add_text_block(s, cx + Inches(0.15), cy + Inches(0.55), cw - Inches(0.3), Inches(0.4),
                       [sub1], font_size=11, bold=True, color=C_TEXT)
        add_text_block(s, cx + Inches(0.15), cy + Inches(1.00), cw - Inches(0.3), Inches(0.9),
                       [sub2], font_size=10, color=C_TEXT)
        cx += cw + Inches(0.1)
    add_footer(s, 21)
    set_speaker_note(s,
        "v3 合併 — 上半「後續可動」(原 S22 4 row) + 下半「已初步現」(原 S24 3 cards) 為單一 slide。"
        "**上半 4 row 可一起做的事**: "
        "(1) Phase 2A — Normal methylation reference: 依賴 V5 BAM; HP tag 正確後 normal-only methylation reference 可用, "
        "對應五大目標 4 (TO normal 補強) 的解鎖。"
        "(2) HPFineNGroups marker — master × flag 重驗: subclone marker ⭐4→⭐3 已降級因 HCC1395 TO ClairS-TO raw split "
        "無法重現 master 89.1% (Fisher odds 0.913 反向 p=3.5e-3) + flag=on NG≥3=0; 機制重詮釋為 phasing signature。"
        "(3) Archive haplotag_version 標記: manifest.yaml 加 haplotag_version 欄位 (baseline / V5), "
        "後續 cross-version comparison 可追溯。對應 P0 行動 F5。"
        "(4) Thread D — LOH-constrained phasing (論文主軸候選): NG=2 cross-sample 6/6 POSITIVE "
        "(Wilcoxon p=0.0156, median gap 0.365), Inner same-hap 93-99% (obs18 6 樣本); "
        "需獨立 phasing-vs-methylation 驗證確認非 phasing artifact; V5 BAM 是前提。"
        "**下半 3 卡片 已初步現**: "
        "(a) Thread D NG=2 6/6 POSITIVE: B1 報告 2026-04-23 Wilcoxon signed-rank W=21 (6 樣本配對最強顯著性), "
        "exact p=0.0156, median gap=0.365; 配 paired control (B3) median gap=0.00003, p=0.578 不顯著 → "
        "確認 LOH-constrained phasing 是 TO-specific; HCC1954 outlier (B2) 解析為 caller FP 背景 84% 非 phasing failure。"
        "(b) HPFineNGroups 重驗中: 評級 ⭐4→⭐3 因 master 無法重現; Phase 2B 計畫 master × flag 機制本質判定。"
        "(c) LOH-constrained phasing 論文主軸候選: Inner same-hap 93-99%, 需獨立驗證。"
        "Q11 預備 (「Thread D 6 樣本可信嗎?」): Wilcoxon 6 樣本配對最強顯著性是 6/6 同方向 W=21 p=0.0156, "
        "本實驗達此上限; evidence grade B (待 Phase 2B 升 A); 多重比較校正未做 (只 1 假設)。"
        "下一張 S22 五大目標銜接 + 研究發展樹。")
    return s


def slide_22_five_goals_tree(prs):
    """S22 (v3): 五大目標 + 樹 (原 v2-26 S23)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "五大目標銜接 + 研究發展樹")
    img = make_s22_research_tree()
    fit_image_within(s, str(img), Inches(0.2), Inches(1.0), Inches(13.0), Inches(5.7))
    add_text_block(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["目標 1/2/4 直接依賴  ·  目標 3/5 部分依賴  ·  Self-phasing 修補是解鎖前提"],
                   font_size=12, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_footer(s, 22)
    set_speaker_note(s,
        "InterSubMod 五大研究目標 (來源 InterSubMod/docs/architecture/20260327_InterSubMod研究願景定錨_01.md): "
        "目標 1 — per-CpG 甲基位點多標籤關聯性評分 (每個 CpG 與 ALT/REF/HP1/HP2 關聯); "
        "目標 2 — clone 結構分析 (sub-clone + 共演化); "
        "目標 3 — 二次打擊與事件順序推論 (依目標 1/2 + LOH/CNV/HP); "
        "目標 4 — TO normal 資訊補強 (無配對 normal 下估計 normal 背景); "
        "目標 5 — 整合 evidence panel 提升 F1 (補充 caller, 保留 TP, 過濾 FP)。"
        "研究發展樹: self-phasing 修補完成 (V5) → 4-bucket 分群可信 → NGroups/HPSig/HP_Ratio 可信 → 五大目標解鎖。"
        "目標 1/2/4 直接依賴 HP tag 正確 (per-CpG 關聯需 hap 分層、clone 結構需 4-bucket、TO normal 補強需 phasing 正確)。"
        "目標 3/5 部分依賴 (3 間接依賴 1/2; 5 部分依賴 + 整合 evidence)。"
        "Self-phasing 修補是 5 大目標解鎖前提, 特別是目標 1/2/4 (V5 修補完才能信任分析結果)。"
        "本張展示研究藍圖, 讓教授看到本工作不只是一次性修補, 而是支撐整個 InterSubMod 的根基。"
        "與 S1 底部小字「修補 self-phasing 是解鎖五大目標的前提」呼應。"
        "下一張 S23 列短期 P0 行動清單。")
    return s


def slide_23_p0_actions(prs):
    """S23 (v3): P0 行動清單 (原 v2-26 S25)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "短期 P0 行動清單")
    rows = [
        ("F1", "Commit V5 working tree (Layer 1.5 + alt guard)", "高",
         "切 2 獨立 commits 後 tag v5.0"),
        ("F3", "7 樣本 V5 BAM 全量重跑", "高",
         "HCC1395 / DORADO / HCC1937 / HCC1954 / H1437 / H2009 / COLO829"),
        ("F4", "master × flag 重驗 HPFineNGroups", "高",
         "對應 S21 機制本質判定"),
    ]
    cy = Inches(1.2)
    rh = Inches(1.4)
    for fid, desc, prio, sub in rows:
        pcolor = C_ACCENT if prio == "高" else (C_AMBER if prio == "中" else C_GRAY)
        add_rect(s, Inches(0.5), cy, Inches(12.4), rh,
                 fill=C_LIGHT, line=pcolor, line_w_pt=2.5)
        add_text_block(s, Inches(0.7), cy + Inches(0.3), Inches(1.0), Inches(0.8),
                       [fid], font_size=28, bold=True, color=pcolor, align=PP_ALIGN.CENTER)
        add_text_block(s, Inches(1.9), cy + Inches(0.2), Inches(8.5), Inches(0.55),
                       [desc], font_size=15, bold=True, color=C_TEXT)
        add_text_block(s, Inches(1.9), cy + Inches(0.8), Inches(8.5), Inches(0.55),
                       [sub], font_size=11, color=C_GRAY)
        add_rect(s, Inches(10.6), cy + Inches(0.35), Inches(2.0), Inches(0.7),
                 fill=pcolor)
        add_text_block(s, Inches(10.6), cy + Inches(0.45), Inches(2.0), Inches(0.5),
                       [prio], font_size=18, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
        cy += rh + Inches(0.12)
    add_text_block(s, Inches(0.5), Inches(6.7), Inches(12.5), Inches(0.4),
                   ["F5-F8 (manifest haplotag_version、cnLOH、trio-phased、100 隨機位點) 入 speaker note"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(s, 23)
    set_speaker_note(s,
        "短期 P0 行動只列 3 條高優先, 其餘 F5-F8 入 speaker note。"
        "**F1 (高) Commit V5 working tree**: 切 2 獨立 commits — "
        "Commit A getVote() Layer 1.5 somatic fallback (行 512-563); "
        "Commit B countSNPHaplotype() alt guard 對稱化 (行 489-494)。"
        "完成後 tag v5.0, 在 InterSubMod manifest.yaml 加 haplotag_version = v5.0 讓下游可追溯。"
        "**F3 (高) 7 樣本 V5 BAM 全量重跑**: 目前只在 HCC1395 5kHz 完整驗證, 跨樣本擴展 R6 未做; "
        "樣本清單 (KDE-corrected): HCC1395_5kHz (已驗證)、HCC1395_DORADO、HCC1937、HCC1954、H1437、H2009、COLO829; "
        "預估 ~10 hr parallel (依 archive TO rerun 規劃)。"
        "**F4 (高) master × flag 重驗 HPFineNGroups**: 對應 S21 機制本質判定, "
        "確認 HPFineNGroups 是 phasing signature 還是 methylation bimodality。"
        "**F5-F8 (中低優先入 speaker note)**: "
        "F5 (中) Manifest 加 haplotag_version 欄位; "
        "F6 (中) cnLOH 區獨立評估方案 (R1 caveat 應對); "
        "F7 (低) Trio-phased 第二 ground truth (R9 應對); "
        "F8 (高) 100 隨機位點 cross-validate (R3 cherry-pick 風險應對)。"
        "Priority 透明展示: 高優先三項 F1 / F3 / F4 是這次 meeting 後優先要做的; "
        "F5-F8 為跟進工作, 視時間與資源安排。"
        "下一張 S24 take-home + Q&A 感謝結語。")
    return s


def slide_24_take_home(prs):
    """S24 (v3): Take-home + Next step + Q&A 感謝 (原 v2-26 S26)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xFA, 0xFB, 0xFC))
    # Top: take-home (big)
    add_rect(s, 0, 0, SLIDE_W, Inches(3.0), fill=C_TITLE_BG)
    add_text_block(s, Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.5),
                   ["Take-home"],
                   font_size=20, bold=True, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.7),
                   ["TO 模式可用  ·  Tag 層問題已解  ·  五大目標解鎖"],
                   font_size=30, bold=True, color=C_TITLE_FG,
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(1.95), Inches(12.5), Inches(0.5),
                   ["read-level concordance +8.3 pp (clean PS, 全基因組)"],
                   font_size=17, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(2.45), Inches(12.5), Inches(0.4),
                   ["V5 不修 self-phasing 本身 (V2b 已處理)  ·  分層處理 phase 與 tag 兩層"],
                   font_size=12, color=RGBColor(0xCF, 0xE2, 0xF3),
                   align=PP_ALIGN.CENTER, italic=True)

    # Middle: 3 next-step commitments
    add_text_block(s, Inches(0.5), Inches(3.20), Inches(12.5), Inches(0.4),
                   ["Next Step (下次 meeting 教授會看到的進度)"],
                   font_size=16, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    items = [
        ("F1", "Commit V5 working tree", "切 2 獨立 commits + tag v5.0"),
        ("F3", "7 樣本擴展 V5 全量重跑", "HCC1395 / DORADO / HCC1937 ... 共 7 樣本"),
        ("F4", "HPFineNGroups master × flag 重驗", "判定 phasing signature 機制本質"),
    ]
    cy = Inches(3.80)
    rh = Inches(0.78)
    colors = [C_ACCENT, C_AMBER, C_GREEN]
    for (fid, desc, sub), color in zip(items, colors):
        add_rect(s, Inches(0.6), cy, Inches(12.2), rh,
                 fill=C_LIGHT, line=color, line_w_pt=2.0)
        add_text_block(s, Inches(0.8), cy + Inches(0.10), Inches(0.9), Inches(0.6),
                       [fid], font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_block(s, Inches(1.9), cy + Inches(0.08), Inches(7.8), Inches(0.4),
                       [desc], font_size=13, bold=True, color=C_TEXT)
        add_text_block(s, Inches(1.9), cy + Inches(0.45), Inches(7.8), Inches(0.35),
                       [sub], font_size=10, color=C_GRAY)
        cy += rh + Inches(0.05)

    # Bottom: Q&A invitation + Thank you (v3 強化結尾)
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.7),
             fill=RGBColor(0xFD, 0xF6, 0xE3), line=C_TITLE_BG, line_w_pt=2.0)
    add_text_block(s, Inches(0.5), Inches(6.62), Inches(12.4), Inches(0.35),
                   ["Q & A 歡迎隨時打斷  ·  詳細數據見 source_materials/  ·  Thank you"],
                   font_size=14, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
    add_text_block(s, Inches(0.5), Inches(6.95), Inches(12.4), Inches(0.3),
                   ["3 source reports + 6 onboarding notes + storyboard_v2.md"],
                   font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
    set_speaker_note(s,
        "結語 v3 強化版 — 倒過來收尾: take-home + next step 承諾 + Q&A 感謝, 不重複 S1 TL;DR。"
        "**Take-home (上半大字)**: TO 模式可用、Tag 層問題已解、五大目標解鎖。"
        "三個重點: "
        "(1) TO 模式可用 — 透過 V5 修補, TO + PON 條件下 tag 品質接近 paired 水準, "
        "read-level concordance +8.3 pp (clean PS blocks, 全基因組)。"
        "(2) Tag 層問題已解 — V3F 解 priority bug + enum mismatch (Bug 1+2), V5 補 Layer 1.5 directional reassignment, "
        "AMB% 17.5%→8.0%, HP:i:33 reads -54%, 介面契約零變動。"
        "(3) 五大目標解鎖 — InterSubMod 目標 1/2/4 直接依賴 HP tag 正確, V5 修補是研究可信度前提。"
        "**口徑校準分層 caveat**: V5 不修 self-phasing 本身 (V2b 已處理 phase scaffold); "
        "self-phasing 問題鏈是被分層處理的 (phase 層由 V2b PON-only, tag 層由 V3F/V5)。"
        "**Next step 3 條承諾** (下次 meeting 教授會看到的進度): "
        "F1 Commit V5 working tree (切 2 獨立 commits + tag v5.0); "
        "F3 7 樣本擴展 V5 全量重跑 (HCC1395_5kHz 已驗證, 待 DORADO / HCC1937 / HCC1954 / H1437 / H2009 / COLO829, "
        "預估 ~10 hr parallel); "
        "F4 HPFineNGroups master × flag 重驗 (判定 phasing signature 機制本質, 對應 Thread D 主軸)。"
        "以承諾收尾比口號收尾更有力, 教授知道下次會看到具體進度。"
        "**v3 結尾強化**: 底部 Q&A 歡迎條 + Thank you, 雙語並列 (CJK + Latin per-char fallback)。"
        "Q&A 隨時歡迎打斷, 詳細數據在 v2/source_materials 三份報告 (01 IGV visual audit、02 V5 audit、"
        "03 longphase TO vs V5 技術報告) 與 v2/notes 6 份 onboarding 文件 (00 background、qa 11+v3 Q12 Q13、"
        "key metrics、code references、industry references、glossary) + 04 Purity calculator failure root cause。"
        "Thank you。")
    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_impact,                # S1 17.3:1 -> 1:1 + 動機小字
        slide_02_pipeline,              # S2 Pipeline + 4 階段
        slide_03_hp_tag_def,            # S3 HP tag 5 整數值 + PS/LOH 兩層
        slide_04_three_layer_evidence,  # S4 證據三層 (合併 v2-26 S5+S6)
        slide_05_paired_vs_to,          # S5 Paired vs TO bar
        slide_06_unlock_phasing_vs_tag, # S6 拆鎖 phasing vs tag
        slide_07_loh_two_layers,        # S7 LOH 兩層精確化
        slide_08_why_important,         # S8 TO 根基 + 4-bucket
        slide_09_3tier_count,           # S9 3-tier 29/14/42
        slide_10_prereq_combined,       # S10 4 函數 + PON + Purity 0.95 (合併 v2-26 S11+S12)
        slide_11_root_cause_v3,         # S11 ★ KEY 根因樹含觸發條件主者
        slide_12_v5_layers_v3,          # S12 ★ V5 三層 + mid-low purity 防禦層
        slide_13_code_diff_v3,          # S13 程式碼 diff + somatic-first 紅標
        slide_14_sanity_5evidence_v3,   # S14 ★ Sanity 4 + 5 證據 + Q1 PON 雙路徑
        slide_15_quant_metrics,         # S15 量化指標 (從 v2-26 S17 拆出)
        slide_16_commit_timeline,       # S16 4-commit timeline (從 v2-26 S17 拆出)
        slide_17_v5max1_climax,         # S17 V5max1 climax + caveat
        slide_18_hp_flip_bar,           # S18 HP flip bar + IGV 縮圖
        slide_19_f1_clarify,            # S19 F1 釐清 + cnLOH 邊界
        slide_20_industry_combined,     # S20 業界家族樹 + 2x2 (吸收 S2 內容)
        slide_21_followup_combined,     # S21 後續可動 + 已初步現 合併
        slide_22_five_goals_tree,       # S22 五大目標 + 樹
        slide_23_p0_actions,            # S23 P0 行動清單
        slide_24_take_home,             # S24 take-home + Next step + Q&A 感謝
    ]

    assert len(builders) == TOTAL_SLIDES, f"Expected {TOTAL_SLIDES} slides, got {len(builders)}"

    for i, fn in enumerate(builders, 1):
        print(f"  [{i:2d}/{TOTAL_SLIDES}] building {fn.__name__} ...")
        fn(prs)

    prs.save(str(OUTPUT_PPTX))
    print(f"\n[OK] PPTX saved to {OUTPUT_PPTX}")
    print(f"     Slides: {len(prs.slides)}")
    return OUTPUT_PPTX


if __name__ == "__main__":
    build()
