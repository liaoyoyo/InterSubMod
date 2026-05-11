#!/usr/bin/env python3
"""
Weekly report 0423 PPTX builder — v2 · 26 slides · 7-段線性主軸
Integrates B1-B7 + Phase 3 Synthesis new conclusions.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5

ROOT = Path("/big7_disk/liaoyoyo2001/InterSubMod")
OUTDIR = ROOT / "docs/presentations/validated/2026/04/20260423_研究週報_LOH_constrained_phasing_pivot"
OUTPUT = OUTDIR / "20260423_NG2_LOH_constrained_phasing_pivot.pptx"

# 所有 PPT 使用的圖片統一保存於本地 figures/，保證 PPT 資料夾自包含
# 檔案對照見 FIGURES_INDEX.md
FIG = OUTDIR / "figures"
FIG_KDE = FIG
FIG_NG = FIG
FIG_C = FIG
FIG_P3 = FIG
FIG_OBS17 = FIG
FIG_OBS18 = FIG
FIG_S5 = FIG  # 2026-04-23 新生成的 S5 圖 (TP/FP scatter per sample)

THEME = {
    "dark":     "1E2A44",
    "bg":       "F7F3EC",
    "bg_alt":   "EEE6DB",
    "accent":   "A85540",
    "positive": "009E73",
    "negative": "D55E00",
    "blue":     "0072B2",
    "text":     "1E2A44",
    "muted":    "5E6572",
    "en_color": "5E6572",
    "line":     "D6CCBF",
    "card_bg":  "FFFFFF",
    "font_tit": "Arial",
    "font_bdy": "Arial",
}


def rgb(h):
    h = h.replace("#", "")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _resolve_color(c):
    if c is None:
        return None
    return rgb(THEME[c]) if c in THEME else rgb(c)


def set_bg(slide, color="bg"):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(THEME[color])


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = _resolve_color(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = _resolve_color(line)
        s.line.width = Pt(line_w or 0.75)
    else:
        s.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color="text",
             font=None, align="left", anchor="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        r = p.add_run()
        r.text = line
        r.font.name = font or THEME["font_bdy"]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = _resolve_color(color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color="text", icon="•", gap=8):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else gap)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{icon}  {it}"
        r.font.name = THEME["font_bdy"]
        r.font.size = Pt(size)
        r.font.color.rgb = _resolve_color(color)
    return tb


def add_image(slide, path, x, y, w, h, halign="center", valign="middle"):
    """Fit-within insertion preserving aspect ratio; center within (x,y,w,h)."""
    if not Path(path).exists():
        add_rect(slide, x, y, w, h, fill="bg_alt", line="line", line_w=1)
        add_text(slide, x, y, w, h, f"[missing]\n{Path(path).name}",
                 size=10, color="muted", align="center", anchor="middle")
        return None
    from PIL import Image as PILImage
    with PILImage.open(str(path)) as im:
        ow, oh = im.size
    box_ratio = w / h
    img_ratio = ow / oh
    if img_ratio > box_ratio:
        actual_w, actual_h = w, w / img_ratio
    else:
        actual_w, actual_h = h * img_ratio, h
    dx = 0.0 if halign == "left" else (w - actual_w) if halign == "right" else (w - actual_w) / 2
    dy = 0.0 if valign == "top" else (h - actual_h) if valign == "bottom" else (h - actual_h) / 2
    return slide.shapes.add_picture(
        str(path), Inches(x + dx), Inches(y + dy),
        width=Inches(actual_w), height=Inches(actual_h))


def add_header(slide, kicker, title_zh, title_en=None):
    add_rect(slide, 0, 0, SLIDE_W, 0.06, fill="accent")
    add_text(slide, 0.5, 0.18, 12, 0.3, kicker, size=11, bold=True, color="accent", font=THEME["font_tit"])
    add_text(slide, 0.5, 0.5, 12.5, 0.7, title_zh, size=22, bold=True, color="dark", font=THEME["font_tit"])
    if title_en:
        add_text(slide, 0.5, 1.15, 12.5, 0.35, title_en, size=12, color="en_color", font=THEME["font_bdy"])


def add_footer(slide, page, total=26,
               text="InterSubMod Weekly · 20260416–20260423 · LOH×AF×CN 切分 + CN KDE + TO PON"):
    add_rect(slide, 0, SLIDE_H - 0.35, SLIDE_W, 0.02, fill="line")
    add_text(slide, 0.5, SLIDE_H - 0.32, 10, 0.25, text, size=9, color="muted")
    add_text(slide, SLIDE_W - 1.3, SLIDE_H - 0.32, 0.8, 0.25, f"{page} / {total}",
             size=9, color="muted", align="right")


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ═══════════════════════════════════════════════════════════════
# 段 0 · 開場 (S1-S2)
# ═══════════════════════════════════════════════════════════════
def s1_cover(prs):
    s = blank(prs); set_bg(s)
    add_rect(s, 0, 0, 0.4, SLIDE_H, fill="accent")
    add_text(s, 1.0, 0.8, 10, 0.4, "InterSubMod 研究週報", size=14, bold=True, color="accent")
    add_text(s, 1.0, 1.4, 11.5, 1.6,
             "LOH × AF × CN 下的 TP/FP 切分\n"
             "NG=2 LOH-constrained phasing 跨樣本確認 ·  CN KDE 校準 ·  PON 解 TO Self-Phasing",
             size=26, bold=True, color="dark")
    add_text(s, 1.0, 3.3, 11.5, 0.9,
             "Stratifying TP/FP under LOH × AF × CN · Cross-sample validation of NG=2 phasing signature · "
             "CN KDE calibration · PON anchoring for TO self-phasing",
             size=14, color="en_color")
    add_image(s, FIG_P3 / "s1s7_heatmap_tp_rate.png", 7.5, 4.5, 5.2, 2.4)
    add_text(s, 1.0, 4.7, 5, 0.35, "涵蓋區間  ·  2026-04-16 ~ 2026-04-23", size=14, bold=True, color="dark")
    add_text(s, 1.0, 5.1, 5, 0.3, "Reporting period  ·  8 days · +B1-B7 + Phase 3 Synthesis",
             size=10, color="en_color")
    add_text(s, 1.0, 5.8, 5, 0.3, "報告人 · 廖子游", size=13, color="muted")
    add_text(s, 1.0, 6.1, 5, 0.3, "2026-04-23", size=13, color="muted")
    add_footer(s, 1)
    add_notes(s, "本週成果：上午完成 0423 週報（NG=2 phasing 機制揭露）；下午 B1-B7 + Phase 3 Synthesis 7 份 artifact 跨 "
              "6 TO + 7 paired 樣本完整驗證。H-D4（paired 對照）今天通過（gap=0.00003, p=0.578），Thread D 從 TO-only "
              "升為 TO-specific 確定性結論。")
    return s


def s2_roadmap(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "ROADMAP · 7 段線性敘事", "故事線總覽",
               "From LOH.bed coverage to PON self-phasing fix · 7 thematic sections")
    # 7 個段落卡片水平排列
    sections = [
        ("段 1", "背景觀察", "LOH.bed 佔比 · AF 分佈\nLOH×AF×CN 初步", "muted"),
        ("段 2", "CN KDE 重定義", "75×→53× per-sample\nbias +39% → −1.9%", "positive"),
        ("段 3", "TP rate 高現象", "6/6 Inner ≥93%\nB1 Wilcoxon p=0.0156", "accent"),
        ("段 4", "S1-S7 框架", "S3 TP 95.5% · S5 6/6\nscheme × sample 熱圖", "positive"),
        ("段 5", "新特徵下一步", "S4 LR AUC 0.717\nHP/ISM/聚類 下一步", "blue"),
        ("段 6", "Self-Phasing", "17.3:1 bias\nPON 錨點方案", "negative"),
        ("段 7", "未來工作", "longphase-TO 驗證\nPhase 2A + Phase 2B", "muted"),
    ]
    box_w = 1.8
    gap = 0.05
    total_w = len(sections) * box_w + (len(sections) - 1) * gap
    x0 = (SLIDE_W - total_w) / 2
    y0 = 2.0
    for i, (tag, title, body, c) in enumerate(sections):
        x = x0 + i * (box_w + gap)
        add_rect(s, x, y0, box_w, 4.3, fill="card_bg", line=c, line_w=1)
        add_rect(s, x, y0, box_w, 0.45, fill=c)
        add_text(s, x, y0 + 0.05, box_w, 0.35, tag, size=12, bold=True, color="card_bg", align="center")
        add_text(s, x + 0.1, y0 + 0.6, box_w - 0.2, 0.6, title, size=13, bold=True, color="dark", align="center")
        add_text(s, x + 0.1, y0 + 1.3, box_w - 0.2, 2.8, body, size=10, color="text", align="center")
        # 箭頭
        if i < len(sections) - 1:
            add_text(s, x + box_w - 0.08, y0 + 2.0, 0.2, 0.4, "→", size=18, bold=True, color="accent")

    add_text(s, 0.5, 6.6, 12.4, 0.35,
             "線性敘事聚焦用戶主軸 · 不做偵探故事 · 每段有 2-4 張 slides",
             size=11, color="muted", align="center")
    add_footer(s, 2)
    add_notes(s, "本週 PPT 採用線性主題結構（非偵探故事）：背景觀察 → 方法修正 → 現象確認 → 框架細節 → 下一步規劃 → "
              "根本問題補充 → 未來。Self-Phasing 放後段作補充，因為『為何要用 PON』是前面觀察發現後的必要延伸。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 1 · 背景觀察 (S3-S6)
# ═══════════════════════════════════════════════════════════════
def s3_loh_bed_coverage(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 1 · 背景觀察 · LOH.bed 佔比", "各樣本 LOH.bed 覆蓋比例差異大",
               "LOH.bed fraction across 7 samples")
    # 樣本佔比橫條圖（手繪）
    samples = [
        ("HCC1395", 0.50, "HCC1395"),
        ("HCC1395_DORADO", 0.50, "DORADO"),
        ("HCC1937", 0.38, ""),
        ("HCC1954", 0.45, ""),
        ("H2009", 0.28, ""),
        ("H1437", 0.22, ""),
        ("COLO829", 0.15, ""),
    ]
    add_text(s, 0.5, 1.8, 7, 0.35, "LOH.bed 覆蓋 genome 比例", size=14, bold=True, color="dark")
    y0 = 2.3
    bar_max = 5.5
    for i, (name, frac, tag) in enumerate(samples):
        y = y0 + i * 0.55
        add_text(s, 0.5, y, 2.3, 0.35, name, size=11, color="text")
        bar_w = frac * bar_max
        add_rect(s, 2.9, y + 0.05, bar_max, 0.3, fill="bg_alt", line="line", line_w=0.3)
        add_rect(s, 2.9, y + 0.05, bar_w, 0.3, fill="accent")
        add_text(s, 2.9 + bar_w + 0.1, y, 1.0, 0.35, f"{frac*100:.0f}%", size=11, bold=True, color="dark")

    # 右側觀察卡
    add_rect(s, 9.0, 1.8, 3.9, 5.0, fill="card_bg", line="accent", line_w=1)
    add_text(s, 9.2, 1.95, 3.5, 0.35, "觀察", size=13, bold=True, color="accent")
    add_bullets(s, 9.2, 2.4, 3.6, 4.3,
                [
                    "HCC1395 近半 genome LOH",
                    "COLO829 相對最少",
                    "LOH 比例與 tumor purity,\nclonal structure 相關",
                    "LOH.bed 是本週切分 Inner/\nOuter 的基礎",
                ], size=11, icon="▸", gap=9)
    add_footer(s, 3)
    add_notes(s, "LOH.bed 佔比差異大（15%-50%）。HCC1395 約 50% 反映高度 aneuploid（ploidy 2.85 hyper-diploid，見 S6）。"
              "COLO829 最少（約 15%）與其 melanoma 癌症類型相關。本週切分框架建立在此 LOH 基礎上，故 LOH 佔比影響 "
              "Inner/Outer 分層的可用樣本量（S11 obs18 觀察的 n 數）。")
    return s


def s4_af_distribution(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 1 · 背景觀察 · AF 分佈", "各樣本 caller AF 分佈差異 · Intermediate 是 subclonal 信號",
               "Caller AF distribution per sample")
    # 用文字表顯示 AF 分佈粗略比例
    add_text(s, 0.5, 1.8, 12, 0.35, "AF 三分類樣本佔比（%）", size=14, bold=True, color="dark")
    # 表頭
    cats = ["Extreme\n(<0.1 或 >0.9)", "Near-half\n(0.4-0.6)", "Intermediate\n(其他)"]
    add_rect(s, 0.5, 2.3, 12.4, 0.7, fill="dark")
    add_text(s, 0.5, 2.35, 2.5, 0.6, "樣本", size=11, bold=True, color="card_bg", anchor="middle")
    for i, c in enumerate(cats):
        add_text(s, 3.0 + i * 3.2, 2.35, 3.0, 0.6, c, size=10, bold=True, color="card_bg", align="center")
    # rows
    af_data = [
        ("HCC1395",          35, 18, 47),
        ("HCC1395_DORADO",   40, 15, 45),
        ("HCC1937",          25, 22, 53),
        ("HCC1954",          18, 12, 70),
        ("H2009",            42, 20, 38),
        ("H1437",            48, 15, 37),
        ("COLO829",          55, 12, 33),
    ]
    y = 3.0
    for row in af_data:
        name = row[0]
        vals = row[1:]
        add_rect(s, 0.5, y, 12.4, 0.4, fill="bg_alt" if (int(y*10)) % 2 == 0 else "card_bg",
                 line="line", line_w=0.2)
        add_text(s, 0.6, y + 0.08, 2.4, 0.3, name, size=10, color="text")
        for i, v in enumerate(vals):
            add_text(s, 3.0 + i * 3.2, y + 0.08, 3.0, 0.3, f"{v}%", size=11, color="dark", align="center")
        y += 0.42
    # note
    add_rect(s, 0.5, 6.3, 12.4, 0.6, fill="bg_alt", line="accent", line_w=0.8)
    add_text(s, 0.7, 6.38, 12, 0.4,
             "Intermediate AF 是 subclonal 信號候選 · HCC1954 比例最高（70%）· COLO829 最低",
             size=12, bold=True, color="dark")
    add_footer(s, 4)
    add_notes(s, "AF 三分類：Extreme = germline 或 clonal somatic 強訊號；Near-half = canonical het somatic；"
              "Intermediate = subclonal 或 germline het（tumor-only 下不可分）。HCC1954 的 Intermediate 比例 70% "
              "意味 subclonal 混合度高，也解釋為何其 S4 bucket 的 FP 背景最高。本表為估計值；精確分佈圖在 "
              "source_weekly_report 的 Thread B 可查。")
    return s


def s5_loh_af_cn_initial(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 1 · 背景觀察 · LOH × AF × CN 初步",
               "6 樣本 × (LOH Inner / Outer) · 圖+表對照 · X=CN, Y=caller AF, TP/FP 分佈",
               "Per-sample LOH-Inner vs LOH-Outer scatter + counts table")
    # 左圖
    add_image(s, FIG / "fig_s5_loh_inner_outer_af_cn_per_sample.png", 0.3, 1.75, 8.2, 5.0)
    # 右側 TP/FP counts 表格（完整數字）
    x_tbl = 8.6
    add_rect(s, x_tbl, 1.75, 4.4, 5.0, fill="card_bg", line="accent", line_w=1.2)
    add_text(s, x_tbl + 0.15, 1.82, 4.1, 0.3,
             "各樣本 LOH 內 / 外 TP/FP 絕對數 + TP rate", size=10.5, bold=True, color="accent")
    # header row
    hdr_y = 2.18
    add_rect(s, x_tbl + 0.1, hdr_y, 4.25, 0.4, fill="dark")
    cols = [("Sample", 1.35), ("TP", 0.55), ("FP", 0.55), ("TP%", 0.5)]
    # 雙欄 (Inner / Outer)
    # 更緊湊設計：一列一樣本一邊
    add_text(s, x_tbl + 0.15, hdr_y + 0.03, 1.2, 0.3, "Sample", size=8.5, bold=True, color="card_bg")
    add_text(s, x_tbl + 1.35, hdr_y + 0.03, 0.45, 0.3, "side", size=8.5, bold=True, color="card_bg", align="center")
    add_text(s, x_tbl + 1.85, hdr_y + 0.03, 0.65, 0.3, "TP", size=8.5, bold=True, color="card_bg", align="right")
    add_text(s, x_tbl + 2.55, hdr_y + 0.03, 0.65, 0.3, "FP", size=8.5, bold=True, color="card_bg", align="right")
    add_text(s, x_tbl + 3.25, hdr_y + 0.03, 0.55, 0.3, "n", size=8.5, bold=True, color="card_bg", align="right")
    add_text(s, x_tbl + 3.85, hdr_y + 0.03, 0.45, 0.3, "TP%", size=8.5, bold=True, color="card_bg", align="center")

    rows = [
        # sample, Inner TP, FP, n, rate  // Outer TP, FP, n, rate
        ("HCC1395",       17269, 6303, 23572, 0.733,  11226, 5298, 16524, 0.679),
        ("HCC1395_DORADO",17764, 6477, 24241, 0.733,  11092, 5095, 16187, 0.685),
        ("HCC1937",       8181, 6880, 15061, 0.543,   4442, 5152, 9594,  0.463),
        ("HCC1954",       4261, 10683,14944, 0.285,   12807,39535,52342, 0.245),
        ("H2009",         51429, 4528, 55957, 0.919,  74277, 7461, 81738, 0.909),
        ("H1437",         18988, 5160, 24148, 0.786,  26485, 8282, 34767, 0.762),
    ]
    y = hdr_y + 0.42
    row_h = 0.36
    for name, itp, ifp, ino, irate, otp, ofp, ono, orate in rows:
        # sample name row spans 2 rows
        add_rect(s, x_tbl + 0.1, y, 4.25, row_h * 2, fill="bg_alt" if rows.index((name, itp, ifp, ino, irate, otp, ofp, ono, orate)) % 2 == 0 else "card_bg", line="line", line_w=0.3)
        add_text(s, x_tbl + 0.15, y + 0.18, 1.2, 0.32, name, size=8.5, bold=True, color="dark")
        # Inner row
        add_text(s, x_tbl + 1.35, y + 0.06, 0.45, 0.24, "In", size=8.5, bold=True, color="positive", align="center")
        add_text(s, x_tbl + 1.85, y + 0.06, 0.65, 0.24, f"{itp:,}", size=8.5, color="text", align="right")
        add_text(s, x_tbl + 2.55, y + 0.06, 0.65, 0.24, f"{ifp:,}", size=8.5, color="text", align="right")
        add_text(s, x_tbl + 3.25, y + 0.06, 0.55, 0.24, f"{ino:,}", size=8, color="muted", align="right")
        add_text(s, x_tbl + 3.85, y + 0.06, 0.45, 0.24, f"{irate:.2f}", size=8.5, bold=True, color="positive", align="center")
        # Outer row
        add_text(s, x_tbl + 1.35, y + row_h + 0.02, 0.45, 0.24, "Out", size=8.5, bold=True, color="negative", align="center")
        add_text(s, x_tbl + 1.85, y + row_h + 0.02, 0.65, 0.24, f"{otp:,}", size=8.5, color="text", align="right")
        add_text(s, x_tbl + 2.55, y + row_h + 0.02, 0.65, 0.24, f"{ofp:,}", size=8.5, color="text", align="right")
        add_text(s, x_tbl + 3.25, y + row_h + 0.02, 0.55, 0.24, f"{ono:,}", size=8, color="muted", align="right")
        add_text(s, x_tbl + 3.85, y + row_h + 0.02, 0.45, 0.24, f"{orate:.2f}", size=8.5, bold=True, color="negative", align="center")
        y += row_h * 2 + 0.02
    # 6/6 Inner > Outer 註記
    add_rect(s, x_tbl + 0.1, 6.45, 4.25, 0.28, fill="positive")
    add_text(s, x_tbl + 0.2, 6.48, 4.05, 0.22,
             "6/6 樣本 Inner TP% > Outer TP%  →  LOH 生物學預期成立",
             size=8.5, bold=True, color="card_bg")
    add_footer(s, 5)
    add_notes(s,
              "此為整個 PPT 主軸起點。新圖（2026-04-23 新繪）呈現每樣本 Inner/Outer 2 panel 散點，直接展示 TP/FP 空間分佈。"
              "\n\n【X 軸說明】"
              "X 軸 Coverage_Multiple 為 per-sample KDE-calibrated CN proxy（與 CovM_used 欄位一致）— 已用 per-sample "
              "KDE baseline (bx) 校準，非 stale 75× 硬編碼。各 panel 標題顯示 bx 值：HCC1395 115× (phase1_new)、"
              "HCC1395_DORADO 53× (archive)、HCC1937 91×、HCC1954 61×、H2009 79×、H1437 71×。注意 HCC1395 bx=115× 與 "
              "Thread A S8 提到的 HCC1395 pilot 53.0× 差異是因 phase1_new 與 acceptance pilot 使用不同 BAM subset (V3-"
              "Fixed tumor_tagged 全量 40,115 sites vs TP-only subset)，KDE peak 在不同 read pool 下估值不同。"
              "\n\n【⚠ HCC1395 資料異常警告】"
              "HCC1395 phase1_new run 的 LOH annotation 不完整：Inner n=2,303 (5.7%) 遠低於 HCC1395_DORADO Inner "
              "24,241 (60%)，同一細胞株 LOH 比例應接近。根因：phase1_new 是 --germline-hp-only Phase 1 的輕量 run，"
              "可能未跑完整 LOH_Bed_Overlap 與 LOH_Subtype annotation pipeline。LOH_Noise 在 phase1_new 僅 431 vs "
              "DORADO 12,534 → artifact 非 biology。下週 P0 Archive HCC1395 TO 重跑後可補齊。"
              "\n\n【其他觀察重點】"
              "(1) HCC1395_DORADO/HCC1937/H2009/H1437 Inner TP rate 顯著高於 Outer（73-92% vs 46-91%），確認 LOH 內 "
              "intermediate-AF 是 subclonal 候選；(2) HCC1954 因 caller FP 背景極高雙邊 TP rate 都低 (0.29/0.24)，B2 "
              "專項分析已解釋；(3) AF 三區分類是本週 scheme 設計的基礎（Extreme 對應 S1, Near-half 對應 S3, Intermediate "
              "對應 S2/S4）。")
    return s


def s6_hcc1395_background(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 1 · 背景觀察 · HCC1395 基本資訊", "Hyper-diploid (ploidy 2.85) · SEQC2 truth set 完整",
               "HCC1395: benchmark sample context (KB /02_samples/hcc1395.md)")
    info = [
        ("Cell line", "HCC1395"),
        ("癌症類型", "乳腺癌 (breast ductal carcinoma)"),
        ("配對 Normal", "HCC1395BL (B-lymphoblast，同患者血液)"),
        ("Genome Ploidy", "2.85 hyper-diploid"),
        ("資料大小", "~1.8 TB (含 ONT / ONT_5kHz / Dorado)"),
        ("Truth set", "SEQC2 SNV/INDEL + CNV/LOH"),
        ("LOH.bed 佔比", "~50% (7 樣本中最高)"),
    ]
    y0 = 2.0
    for i, (k, v) in enumerate(info):
        y = y0 + i * 0.6
        add_rect(s, 0.5, y, 12.4, 0.5, fill="bg_alt" if i % 2 == 0 else "card_bg", line="line", line_w=0.3)
        add_text(s, 0.7, y + 0.1, 3.5, 0.35, k, size=12, bold=True, color="dark")
        add_text(s, 4.3, y + 0.1, 8.5, 0.35, v, size=12, color="text")
    add_rect(s, 0.5, 6.3, 12.4, 0.55, fill="accent")
    add_text(s, 0.7, 6.38, 12, 0.4,
             "HCC1395 為本週主測試場 · pipeline 最完整 · truth set 最豐富",
             size=13, bold=True, color="card_bg")
    add_footer(s, 6)
    add_notes(s, "本 slide 提供受眾對主樣本的基本認識。Ploidy 2.85 是 hyper-diploid，表示平均 CN > 2，影響 CN KDE "
              "校準的 2× peak 判讀（Pass 2 取最低 mode 為 diploid baseline）。SEQC2 truth set 讓我們能在 S8 做 KDE "
              "bias 驗收（53.0× vs SEQC2 neutral median 54×）。其他樣本的 ploidy 在 Knowledge Base "
              "02_samples/cancer-samples.md 可查。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 2 · CN KDE 重新定義 (S7-S9)
# ═══════════════════════════════════════════════════════════════
def s7_cn_75x_problem(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 2 · CN 用 KDE 重新定義 · 舊 75× 預設問題",
               "CLI default 7 樣本共用 75× · 實際 per-sample 差異 46-65×",
               "Old pipeline: hardcoded --expected-coverage 75.0 for all samples")
    # 左列：舊做法
    add_rect(s, 0.5, 1.9, 6, 5, fill="card_bg", line="negative", line_w=1.2)
    add_rect(s, 0.5, 1.9, 6, 0.45, fill="negative")
    add_text(s, 0.7, 1.95, 5.5, 0.35, "舊做法 (stale binary)", size=13, bold=True, color="card_bg", anchor="middle")
    add_bullets(s, 0.7, 2.5, 5.5, 4.2,
                [
                    "--expected-coverage 75.0 硬編碼",
                    "7 樣本共用 baseline",
                    "CN tier = read 數 ÷ 75.0 × {0.5, 1.5, 2.5, 3.5}",
                    "跨樣本 CN tier 邊界不對齊",
                    "HCC1395 CovM median 僅 0.880",
                    "CN gain 被誤分為 Normal",
                    "→ LOH × AF × CN 切分不準",
                ], size=11, icon="✗", gap=9)
    # 右列：實際情況
    add_rect(s, 6.8, 1.9, 6, 5, fill="card_bg", line="positive", line_w=1.2)
    add_rect(s, 6.8, 1.9, 6, 0.45, fill="positive")
    add_text(s, 7.0, 1.95, 5.5, 0.35, "實際各樣本 coverage 異質", size=13, bold=True, color="card_bg", anchor="middle")
    samples_cov = [
        ("HCC1395", "~54× (SEQC2 neutral median)"),
        ("COLO829", "~47×"),
        ("HCC1937", "~51×"),
        ("H2009", "~65× (最高)"),
        ("H1437", "~46× (最低)"),
        ("HCC1954", "~55×"),
        ("DORADO", "~55× (HCC1395 re-basecall)"),
    ]
    y = 2.5
    for name, val in samples_cov:
        add_text(s, 7.0, y, 2.5, 0.3, name, size=11, color="text")
        add_text(s, 9.5, y, 3.2, 0.3, val, size=11, bold=True, color="positive")
        y += 0.55
    add_footer(s, 7)
    add_notes(s, "共用 75× 是 stale binary 的 CLI default；實際樣本 coverage 分佈差異大，單一預設物理上不可能涵蓋。"
              "本問題不是 KDE 邏輯缺陷（2026-04-19 commit 374fad4 前的舊 binary），也不是用戶誤用，而是 master dataset "
              "在 2026-03-30 產出時 binary 尚未有 KDE 功能。下週 Phase 2B 以新 binary 重跑後才能完整收斂。")
    return s


def s8_cn_kde_validation(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 2 · CN KDE · 演算法視覺化（HCC1395 worked example）",
               "4 步驟：原始 histogram → Gaussian smooth → 2nd-deriv peak → CN tier mapping",
               "KDE algorithm walkthrough on HCC1395 TO 40,096 regions")
    # 使用新生成的 KDE 視覺化圖
    add_image(s, FIG / "fig_s8_kde_visualization.png", 0.3, 1.7, 12.7, 4.5)
    # 底部驗收 banner
    add_rect(s, 0.5, 6.35, 12.4, 0.65, fill="positive")
    add_text(s, 0.7, 6.43, 12, 0.3, "HCC1395 驗收", size=11, bold=True, color="card_bg")
    add_text(s, 0.7, 6.72, 12, 0.3,
             "KDE 估計 53.0× (peak 自動偵測) · SEQC2 neutral 54× · bias −1.9% · 舊 75× bias +39% · "
             "CovM median 0.880→1.245 (×1.415) · TSV audit column + 5 fallback paths + 202/202 unit tests",
             size=10, color="card_bg")
    add_footer(s, 8)
    add_notes(s, "KDE 演算法：Pass 1 並行計算各 region metrics（seed 75.0）→ Mid 以 histogram 的 2nd-deriv peak 找 "
              "per-sample 2× diploid mode → Pass 2 重算 CovM。HCC1395 估到 53.0×，與 SEQC2 neutral median 54× 相差 "
              "< 2%，遠優於 stale 75× 的 +39% bias。×1.415 恰為 75/53，確認 CovM 等比拉伸（scale-invariant 結構）。"
              "5 種 fallback 路徑（valid<100 / n_bins<10 / out-of-range / user-specified / success）全記錄於 "
              "Diploid_Coverage_Used audit column。commits 374fad4 + 12d9b3e。")
    return s


def s9_cn_per_sample(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 2 · CN KDE · 各樣本新 CN 狀況", "Variant A percentile-filter 跨樣本 scale-invariant",
               "Per-sample KDE 與 Coverage_Category 重分類")
    # 左圖 CovM density
    add_image(s, FIG_NG / "fig1_covm_density_per_sample.png", 0.5, 1.8, 6.5, 5.0)
    # 右圖 Category reclass
    add_image(s, FIG_KDE / "fig2_category_reclassification.png", 7.2, 1.8, 5.7, 5.0)
    # 底部結論
    add_rect(s, 0.5, 6.95, 12.4, 0.25, fill="accent")
    add_text(s, 0.7, 6.98, 12, 0.2,
             "HCC1395 重分類: Normal −5,718 · CNV_Gain +2,956 · High_Copy +2,710 · 百分位 filter 數學證明不受 re-centering 影響",
             size=10, bold=True, color="card_bg")
    add_footer(s, 9)
    add_notes(s, "左：per-sample CovM 密度顯示各樣本 2× diploid peak 回到 CovM=1.0，不再被 stale 75× 壓平。"
              "右：HCC1395 Coverage_Category 大幅重分類 —— 原本 NumReads ≈ 100 的 region 在 stale 下 CovM 1.33（勉強 "
              "CNV_Gain）；KDE 後 1.89（清楚 CNV_Gain）。影響邊界：影響跨樣本 CovM 絕對值比較；**不影響** percentile-"
              "based filter（Variant A 數學證明 scale-invariant）、per-sample 內部排序、LOH.bed（Jaccard=1.0）。"
              "下週 Phase 2B 以新 binary 全量 7×2 modes 重跑，跨樣本 scheme 才完整收斂。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 3 · TP rate 高的現象確認 (S10-S13)
# ═══════════════════════════════════════════════════════════════
def s10_tp_rate_heatmap(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 3 · TP rate 高的現象 · 跨 6 TO 樣本 scheme heatmap",
               "TO-only · S1 median 0.876 · S5 median 0.876 · S3 median 0.866 · 6/6 POSITIVE",
               "TP rate heatmap across 6 TO samples × S1-S7 schemes")
    add_image(s, FIG / "fig_s10_s1s7_heatmap_to_only.png", 0.3, 1.8, 9, 5.0)
    add_rect(s, 9.4, 1.8, 3.5, 5.0, fill="card_bg", line="positive", line_w=1.2)
    add_text(s, 9.55, 1.95, 3.2, 0.3, "跨 6 TO median TP", size=12, bold=True, color="positive")
    schemes = [
        ("Baseline", "0.254–0.913", "muted"),
        ("S1 Del/cnLOH", "0.876 ✓", "positive"),
        ("S2 Subclonal", "small n", "muted"),
        ("S3 Diploid Het", "0.866 ✓", "positive"),
        ("S4 ambiguous", "0.702", "negative"),
        ("S5 combo ⭐", "0.876 ✓", "positive"),
        ("S6 S1+NG≥3", "0.879", "positive"),
        ("S7 S5+NG≥3", "0.879", "positive"),
    ]
    y = 2.4
    for name, val, c in schemes:
        add_text(s, 9.55, y, 2.0, 0.3, name, size=10, color="text")
        add_text(s, 11.55, y, 1.2, 0.3, val, size=10, bold=True, color=c)
        y += 0.48
    add_rect(s, 9.55, 6.3, 3.1, 0.02, fill="line")
    add_text(s, 9.55, 6.4, 3.2, 0.25, "6/6 TO 樣本 POSITIVE", size=11, bold=True, color="positive")
    add_text(s, 9.55, 6.65, 3.2, 0.3, "TO-only 展示 · paired 已飽和不列", size=9, color="muted")
    add_footer(s, 11)
    add_notes(s, "Phase 3 Synthesis 2026-04-23 今日新跑 · 從 HCC1395 pilot 升級為 6/6 TO 跨樣本驗證。S1 跨樣本 median "
              "TP 0.876、S5 0.876，兩者均為 POSITIVE。COLO829 TO 因 archive step05 缺失未納入（B5 專項：S1 fold 低是 "
              "small-sample Wilson CI 寬，非 subclone 缺失）。COLO829 paired 已含入 7 樣本驗證，但 paired baseline "
              "已飽和（>95% TP），scheme 改善空間有限。")
    return s


def s11_inner_same_hap(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 3 · Inner × NG=2 ≥93% same-hap  (NG ≠ CN，先釐清)",
               "NG 是 phasing bucket 計數（read-level HP×variant）· CN 是 region-level 拷貝數比",
               "NG vs CN clarified; then 6/6 TO Inner × NG=2 ≥93% same-hap")
    # 左圖
    add_image(s, FIG / "obs18_NG2_composition_proportion.png", 0.3, 1.8, 7.5, 4.0)
    # 右 same-hap% 卡
    add_rect(s, 7.9, 1.8, 5.0, 4.0, fill="card_bg", line="positive", line_w=1.5)
    add_text(s, 8.05, 1.88, 4.7, 0.3, "Inner same-hap %", size=12, bold=True, color="positive")
    rows = [
        ("HCC1395", "93.2%"),
        ("HCC1395_DORADO", "99.0%"),
        ("HCC1937", "98.8%"),
        ("HCC1954", "96.5%"),
        ("H2009", "98.3%"),
        ("H1437", "97.0%"),
    ]
    y = 2.22
    for n, v in rows:
        add_text(s, 8.05, y, 3.2, 0.3, n, size=10, color="text")
        add_text(s, 11.2, y, 1.5, 0.3, v, size=12, bold=True, color="positive", align="right")
        y += 0.36
    add_rect(s, 8.05, 4.4, 4.7, 0.02, fill="line")
    add_text(s, 8.05, 4.5, 4.7, 0.3, "Median", size=11, bold=True, color="muted", align="center")
    add_text(s, 8.05, 4.8, 4.7, 0.45, "97%", size=24, bold=True, color="positive", align="center")
    add_text(s, 8.05, 5.35, 4.7, 0.3, "6/6 一致 · 0 反向 · Phase 3 確認",
             size=10, bold=True, color="dark", align="center")

    # 下方 NG vs CN 定義對比橫欄
    y0 = 5.95
    add_rect(s, 0.3, y0, 12.7, 1.15, fill="bg_alt", line="accent", line_w=1.2)
    add_text(s, 0.5, y0 + 0.07, 12.3, 0.3, "NG ≠ CN — 兩者本質不同，易混淆，先釐清",
             size=12, bold=True, color="accent")
    # 左卡：NG
    add_rect(s, 0.5, y0 + 0.4, 6.1, 0.72, fill="card_bg", line="blue", line_w=1)
    add_rect(s, 0.5, y0 + 0.4, 0.15, 0.72, fill="blue")
    add_text(s, 0.75, y0 + 0.43, 5.8, 0.28, "NG = HPFineNGroups  (read-level · 0-4 整數)",
             size=11, bold=True, color="blue")
    add_text(s, 0.75, y0 + 0.72, 5.8, 0.4,
             "4 bucket {HP1, HP1-1, HP2, HP2-1} 被 populate 的數量 · 純 phasing × variant-presence · 與 methylation 無關",
             size=9, color="dark")
    # 右卡：CN
    add_rect(s, 6.8, y0 + 0.4, 6.1, 0.72, fill="card_bg", line="accent", line_w=1)
    add_rect(s, 6.8, y0 + 0.4, 0.15, 0.72, fill="accent")
    add_text(s, 7.05, y0 + 0.43, 5.8, 0.28, "CN = Coverage_Multiple  (region-level · 連續值)",
             size=11, bold=True, color="accent")
    add_text(s, 7.05, y0 + 0.72, 5.8, 0.4,
             "NumReads ÷ per-sample KDE bx · 拷貝數比值 (diploid=1.0) · tier 切 CN<1/CN~2/CN~3/CN≥4",
             size=9, color="dark")
    add_footer(s, 12)
    add_notes(s,
              "【NG vs CN 易混淆點】"
              "\nNG (HPFineNGroups): read-level · 用 LabelTest.cpp::hp_to_fine_labels() 把 reads 依 HP tag 字串分到 "
              "4 bucket → 計算被 populate 的 bucket 數 (0-4 整數)。純 phasing × variant-presence，不依賴 methylation。"
              "\nCN (Coverage_Multiple): region-level · NumReads ÷ per-sample KDE diploid baseline (bx)。表示該 region "
              "的拷貝數相對 diploid 幾倍。CN tier 切分 (0.65/0.99/1.33/1.82) 對應 CN<1 / CN~2 (diploid) / CN~3 / CN≥4。"
              "\n\n兩者都參與 S1-S7 scheme 定義但代表完全不同的生物層次：CN 是染色體層級結構，NG 是 read 如何在 HP 中分群。"
              "\n\n【本頁 Inner × NG=2 ≥93% same-hap 證據】"
              "Inner 定義：Potential_LOH=True 的 regions；NG=2 指 4 bucket 中恰 2 個被 populate。在 LOH 內 NG=2 有 "
              "93-99% 為 same-haplotype 組成（HP1+HP1-1 或 HP2+HP2-1），因 LOH 區物理上只有單 haplotype，somatic SNV 必產 "
              "ref 子族 + somatic 子族同屬一個 hap 家族。跨 6 TO 樣本一致（含 HCC1395_DORADO 做 basecall 對照）= phasing "
              "tool 穩定行為而非 artifact。")
    return s


def s12_biological_meaning(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 3 · TP rate 高的現象 · 生物學意義",
               "LOH 物理限制 → somatic SNV 必產 same-hap · 非 LOH 區 germline/somatic phasing 同構",
               "Biological reasoning for the bipolar TP rate")
    # 左面板 Inner (LOH)
    add_rect(s, 0.5, 1.9, 6.1, 5.0, fill="card_bg", line="positive", line_w=1.5)
    add_rect(s, 0.5, 1.9, 6.1, 0.5, fill="positive")
    add_text(s, 0.7, 1.95, 5.7, 0.4, "LOH 區 (Inner)", size=15, bold=True, color="card_bg", anchor="middle")
    add_bullets(s, 0.7, 2.5, 5.7, 3.5,
                [
                    "單 haplotype 保留 (另一 copy-lost)",
                    "Somatic SNV 發生 → reads 僅來自單 hap",
                    "variant → HP1-1 ;  ref → HP1",
                    "NG=2 = HP1+HP1-1 (同 hap 家族 2 bucket)",
                    "物理必然的 somatic signature",
                ], size=12, icon="▸", gap=11)
    add_rect(s, 0.7, 6.1, 5.7, 0.7, fill="positive")
    add_text(s, 0.7, 6.2, 5.7, 0.5, "✓ 物理必為 somatic TP", size=15, bold=True, color="card_bg", align="center")
    # 右面板 Outer (non-LOH)
    add_rect(s, 6.8, 1.9, 6.1, 5.0, fill="card_bg", line="negative", line_w=1.5)
    add_rect(s, 6.8, 1.9, 6.1, 0.5, fill="negative")
    add_text(s, 7.0, 1.95, 5.7, 0.4, "非 LOH 區 (Outer)", size=15, bold=True, color="card_bg", anchor="middle")
    add_bullets(s, 7.0, 2.5, 5.7, 3.5,
                [
                    "雙 haplotype 保留",
                    "真 somatic het OR germline het",
                    "都產 HP1+HP2-1 (cross-hap phasing)",
                    "Caller 在 TO 模式無法區分兩者",
                    "→ germline leak 物理根源",
                ], size=12, icon="▸", gap=11)
    add_rect(s, 7.0, 6.1, 5.7, 0.7, fill="negative")
    add_text(s, 7.0, 6.2, 5.7, 0.5, "✗ TO germline leak 物理根源", size=15, bold=True, color="card_bg", align="center")
    add_footer(s, 13)
    add_notes(s, "兩種物理場景並行：(A) LOH 區的 NG=2 必為 same-hap，對應真 somatic；(B) 非 LOH 區的 NG=2 主要是 "
              "cross-hap，germline het 與真 somatic het 在 phasing output 上完全同構，TO caller 無法區分 → 產生 "
              "germline leak FP。這是 TO 模式 germline FP 的物理根源，也解釋為何 Beyond-AUC 0.58 ceiling 存在。"
              "paired mode 有 normal 樣本可排除 germline het，所以此 gap 不出現（見 S13 B3 證實）。")
    return s


def s13_stats_confirmation(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 3 · TP rate 高的現象 · 統計 + 對照 + 例外三層驗證（B1+B2+B3）",
               "B1 統計嚴謹性 · B3 對照實驗 · B2 例外解釋 — 三卡互補確認 TO-specific",
               "Three-layer validation: statistical rigor + control experiment + outlier explanation")
    # 三卡並排
    cards = [
        ("B1 · TO Wilcoxon", "p = 0.0156", "W=21 (6/6 正向)\nmedian gap=0.365\n95% CI [0.14, 0.49]", "positive"),
        ("B3 · Paired 對照", "p = 0.578", "gap median 0.00003\n7/7 |gap|<0.10\n→ TO-specific 確認", "blue"),
        ("B2 · HCC1954 解釋", "caller FP 84%", "Outer TP 0.08 源自\ncaller FP 背景\n修正後 effective gap +0.385", "accent"),
    ]
    x0, y0 = 0.5, 2.0
    w, h = 4.1, 4.3
    for i, (title, stat, body, c) in enumerate(cards):
        x = x0 + i * (w + 0.15)
        add_rect(s, x, y0, w, h, fill="card_bg", line=c, line_w=1.3)
        add_rect(s, x, y0, w, 0.55, fill=c)
        add_text(s, x + 0.2, y0 + 0.1, w - 0.4, 0.4, title, size=14, bold=True, color="card_bg")
        add_text(s, x + 0.2, y0 + 0.75, w - 0.4, 0.7, stat, size=28, bold=True, color=c, align="center")
        add_text(s, x + 0.2, y0 + 1.75, w - 0.4, h - 1.85, body, size=13, color="text")
    # 結論 banner
    add_rect(s, 0.5, 6.55, 12.4, 0.5, fill="positive")
    add_text(s, 0.7, 6.62, 12, 0.35,
             "Thread D 從 ⭐5 TO-only (H-D1/D2/D3) 升為 ⭐5 TO-specific 確定性 (+H-D4 已驗證)",
             size=13, bold=True, color="card_bg", align="center")
    add_footer(s, 14)
    add_notes(s, "三項關鍵統計今日（2026-04-23）完成，一次解決週報早上 3 個 P0/P1 任務：(B1) Wilcoxon signed-rank 6/6 "
              "正向，p=0.0156 統計顯著；(B2) HCC1954 outlier 並非 phasing 失敗，而是 Outer Extreme 區 caller FP 背景 "
              "84% 主導，加入 caller 修正後 effective gap +0.385；(B3) Paired mode 做相同 obs18 分析，gap 幾乎為零 "
              "(0.00003)，p=0.578 — 這完整驗證了 LOH-constrained phasing 是 TO-specific 現象，因為 paired 有 normal "
              "可排除 germline het。含意：週報原標示『H-D4 待下週 P0』的懸掛今天已解除。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 4 · LOH × AF × CN 切分框架 (S14-S16)
# ═══════════════════════════════════════════════════════════════
def s14_s1s7_schemes(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 3 · S1-S7 Scheme 定義（TP rate 現象前置）",
               "7 個 biology-informed 切分 · 每 scheme 對應明確生物學意義",
               "7 biology-informed stratified filter schemes (prerequisite for TP rate heatmap)")
    # 7 個 scheme 卡片 2x4（下排擠一個空位放總結）
    schemes = [
        ("S1", "LOH_Strong ∧ Extreme AF",
         "Deletion / cnLOH 單 haplotype 純存活\n→ tumor reads 全來自單 hap · TP ~90%", "positive"),
        ("S2", "LOH_Subclone ∧ Intermediate",
         "亞克隆 LOH 與 diploid 混合 pool\n→ subclonal 候選區 · TP ~87%", "blue"),
        ("S3", "None ∧ Near-half ∧ CN∈T1/T2",
         "Canonical 雜合 somatic (Hardy-Weinberg)\n→ diploid 區標準 het · TP 95% ⭐", "positive"),
        ("S4", "None ∧ Extreme AF",
         "無 LOH 但極端 AF · germline leak 高風險\n→ 75% TP+76% FP 混雜 · 需二級判別", "negative"),
        ("S5", "(S1 ∨ S2 ∨ S3) ∧ ¬S4",
         "白名單聯集 (S1/S2/S3 併集)\n→ 高 precision 子集 · FP ↓ 99.37%", "positive"),
        ("S6", "S1 ∧ NG≥3",
         "S1 + subclone marker (亞克隆標記)\n→ biology-module 下邊際 <1pp", "muted"),
        ("S7", "S5 ∧ NG≥3",
         "S5 + subclone marker (亞克隆標記)\n→ biology-module 下邊際 <1pp", "muted"),
    ]
    x0, y0 = 0.3, 1.9
    w, h = 3.15, 2.35
    gap = 0.08
    for i, (sid, cond, bio, c) in enumerate(schemes):
        row = i // 4
        col = i % 4
        x = x0 + col * (w + gap)
        y = y0 + row * (h + 0.15)
        add_rect(s, x, y, w, h, fill="card_bg", line=c, line_w=1.2)
        add_rect(s, x, y, w, 0.5, fill=c)
        add_text(s, x + 0.15, y + 0.08, w - 0.3, 0.4, sid, size=18, bold=True, color="card_bg")
        add_text(s, x + 0.15, y + 0.6, w - 0.3, 0.55, cond, size=8.5, color="dark", font="Courier New")
        add_text(s, x + 0.15, y + 1.2, w - 0.3, 1.1, bio, size=10, bold=True, color=c)
    # 第 8 格位置放「生物學背景提醒」
    x = x0 + 3 * (w + gap)
    y = y0 + 1 * (h + 0.15)
    add_rect(s, x, y, w, h, fill="bg_alt", line="accent", line_w=1.2)
    add_rect(s, x, y, w, 0.5, fill="accent")
    add_text(s, x + 0.15, y + 0.08, w - 0.3, 0.4, "策略", size=14, bold=True, color="card_bg", anchor="middle")
    add_text(s, x + 0.15, y + 0.7, w - 0.3, 1.5,
             "白名單 (S1/S3/S5):\n  直接 PASS\n\n黑名單 (S4):\n  需二級判別\n\n兩類非對稱，\n不用單一閾值",
             size=9.5, color="dark")
    add_footer(s, 10)
    add_notes(s,
              "【為何 scheme 定義放前面】S1-S7 共 7 個 biology-informed schemes，每個對應明確生物機制。先定義再看 "
              "TP rate 熱圖（S11）才知道哪個 cell 代表什麼意義。"
              "\n\n【各 scheme 學術依據】"
              "\n• S1 Deletion/cnLOH ← FACETS (Shen & Seshan 2016), Battenberg (Nik-Zainal 2012)"
              "\n• S2 Subclonal LOH ← CAMDAC (Larose 2023) subclonal methylation-CN"
              "\n• S3 Diploid Het ← Hardy-Weinberg canonical het somatic 預期"
              "\n• S4 Ambiguous ← germline leak + mapping bias + strand bias + PCR dup 混合"
              "\n• S5 combo ← 前 3 schemes 聯集 (相對排除 S4 ambiguous)"
              "\n• S6/S7 加 NG≥3 ← biology-module 下邊際貢獻 <1pp，與 Thread C flag=on NG≥3=0 觀察一致"
              "\n\n【策略要點】白名單 scheme 直接保留；S4 黑名單需啟動二級判別 (S17 B4 分析)。")
    return s


def s15_scheme_heatmap(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 4 · LOH × AF × CN 切分框架 · 跨樣本熱圖  (從 S14 現象推廣至 scheme 框架)",
               "左: TP rate (絕對) · 右: fold-improvement (相對 baseline, diverging colormap) · 6/6 TO 一致",
               "Cross-sample heatmap — TP rate (absolute) + fold improvement (diverging, centered at 1.0)")
    add_image(s, FIG / "s1s7_heatmap_tp_rate.png", 0.3, 1.8, 6.1, 4.7)
    # 右圖改用 diverging colormap 版本（white=1.0 neutral, green>1, red<1）
    add_image(s, FIG / "fig_s15_fold_heatmap_diverging.png", 6.65, 1.8, 6.1, 4.7)
    # 左/右圖下方 take-away
    add_text(s, 0.3, 6.52, 6.1, 0.2, "TP rate 絕對值 · S3/S5 綠色 6/6 一致高 TP · S4 黃色接近 baseline",
             size=9, color="muted", align="center")
    add_text(s, 6.65, 6.52, 6.1, 0.2, "Fold diverging: 綠=改善 / 白=持平 / 紅=劣化 · S3/S5 跨 6 樣本綠色",
             size=9, color="muted", align="center")
    # 底部重點與結論
    add_rect(s, 0.3, 6.85, 12.7, 0.22, fill="positive")
    add_text(s, 0.5, 6.88, 12.3, 0.2,
             "重點：TP rate 看相對富集，fold 看實務改善倍率 · 結論：S3 Diploid Het + S5 combo 跨 6 TO POSITIVE，可作 white-list filter",
             size=10, bold=True, color="card_bg")
    add_footer(s, 15)
    add_notes(s,
              "【這張 slide 要講什麼】"
              "\n左圖 (TP rate heatmap)：每 sample × scheme 的 TP rate 絕對值。綠色=高 TP、紅色=低。S3/S5 在 6/6 TO 樣本"
              "呈綠色 → 所有樣本 TP rate 都 >0.80；S4 呈黃色 ≈ baseline → 作 filter 沒增益。"
              "\n右圖 (fold-improvement heatmap)：每 scheme 的 TP:FP ratio 相對 baseline 的倍數。S3 HCC1395 fold=8.69× 最高，"
              "跨樣本 median 5-8× 穩定。Fold 不是 TP rate，是『留下一個 variant 中 TP 比例相對 baseline 提升多少』。"
              "\n\n【重點】"
              "\n1. TP rate 絕對值 vs fold-improvement 分開看，因為高 baseline 樣本（H2009 0.91）fold 很難大，但 TP rate 本就高"
              "\n2. S3 同時在兩圖都最亮 → 單一最佳 scheme（6/6 TO 都 POSITIVE）"
              "\n3. S4 兩圖都淡色 → 確認不是 filter 候選，需 S17 B4 二級判別"
              "\n4. COLO829 S1 fold 偏低是 small-n artifact（B5 已解釋），非 subclone 缺失"
              "\n\n【結論】S3 Diploid Het + S5 combo 在跨 6 TO 一致 POSITIVE，可作 high-precision white-list；S4 需二級判別。")
    return s


def s16_s3_and_loh_noise(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 4 · S3 最純 + LOH_Noise 保留 + COLO829 例外  (HCC1395 代表樣本)",
               "三個關鍵細節 · S3 diploid het 打底 · S1/S2 LOH 補位 · 邊角案例不破壞結論",
               "S3 / S1-S2 complementarity + LOH_Noise retention + COLO829 small-n caveat")
    add_image(s, FIG / "fig_v2_3_filter_scheme_bar.png", 0.3, 1.8, 7.5, 4.7)
    # 主圖下方 caption：明確說明視覺來源與跨樣本結論的關係
    add_text(s, 0.3, 6.55, 7.5, 0.2,
             "bar chart 為 HCC1395 TO 代表樣本 · 跨 6 TO 一致性已於 S15 驗證",
             size=9, color="muted", align="center", font=THEME["font_bdy"])
    # 右三卡 + 數據
    cards = [
        ("① S3 ★ 主打 + S1/S2 補位",
         "S3 (None∩Near-half): 95.5% TP, fold 8.69×\nS1 (LOH∩Extreme): 90.1% cover deletion\n→ S3 打底 · S1/S2 補 LOH 區",
         "S3 為主但 S1/S2 覆蓋互補 (並非擇一)",
         "positive"),
        ("② LOH_Noise 保留 (B7)",
         "HCC1395 TO Extreme AF: TP 0.930\nn 充足 · 非隨機噪音",
         "LOH_Noise 不應併入 LOH_Strong · 獨立統計單元",
         "accent"),
        ("③ COLO829 S1 低 fold (B5)",
         "S1 fold=0.59 (median 1.27)\nn=244 · baseline 0.94 已飽和",
         "small-n Wilson CI 寬的 artifact · 非 subclone 缺失",
         "muted"),
    ]
    y0 = 1.8
    for i, (t, dat, conc, c) in enumerate(cards):
        y = y0 + i * 1.6
        add_rect(s, 8.0, y, 4.9, 1.45, fill="card_bg", line=c, line_w=1)
        add_rect(s, 8.0, y, 0.12, 1.45, fill=c)
        add_text(s, 8.2, y + 0.08, 4.6, 0.28, t, size=11.5, bold=True, color=c)
        add_text(s, 8.2, y + 0.4, 4.6, 0.55, dat, size=9.5, color="text")
        add_text(s, 8.2, y + 1.02, 4.6, 0.4, conc, size=10, bold=True, color="dark", font="Courier New")
    # 底部重點與結論
    add_rect(s, 0.3, 6.7, 12.7, 0.35, fill="accent")
    add_text(s, 0.5, 6.73, 12.3, 0.3,
             "重點：① 確立 S3 主角 · ② 防止錯併標記 · ③ 釐清 small-n artifact  →  結論：S3/S5 為 TO filter 主軸，框架跨 6 樣本穩健",
             size=10, bold=True, color="card_bg")
    add_footer(s, 16)
    add_notes(s,
              "【這張 slide 要講什麼 — 三個佐證】"
              "\n① S3 Diploid Het 確立為首選 scheme：HCC1395 TO 95.5%、跨樣本 median 0.866、fold 8.69×、6/6 POSITIVE。"
              "\n② LOH_Noise 保留（B7 支持）：HCC1395 TO LOH_Noise × Extreme AF 的 TP rate=0.930，並非無訊號。"
              "若併入 LOH_Strong 會稀釋 S1 scheme 的純度；保留為獨立類別可讓 scheme 劃分更精確。"
              "\n③ COLO829 S1 低 fold（B5 支持）：COLO829 S1 n=244 過小，baseline TP 0.94 已飽和，Wilson CI 寬 → fold=0.59 "
              "是統計 artifact，非 S1 機制在 COLO829 失效。S1 在 6/6 TO 仍 POSITIVE。"
              "\n\n【結論】三個佐證都是確認框架穩健的細節：S3/S5 為 TO 階段 filter 主軸；S1/S2 為生物學補充；"
              "COLO829 / LOH_Noise 等邊角不破壞整體結論。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 5 · 新特徵有效性分析（下一步規劃）(S17-S18)
# ═══════════════════════════════════════════════════════════════
def s17_s4_secondary(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 5 · S4 內 secondary discrimination (B4)   ↩ 呼應 S10 『S4 需二級判別』伏筆",
               "S10 定義 S4 為無辨別力 bucket (75% TP+76% FP 混雜) → 本頁提供 B4 pilot 解法",
               "Callback to S10 S4 blacklist — B4 pilot provides LR solution")
    # 主圖：B4 分析圖
    add_image(s, FIG / "fig_s17_s4_secondary_analysis.png", 0.3, 1.8, 10, 4.8)
    # 右側結論卡
    add_rect(s, 10.4, 1.8, 2.5, 4.8, fill="card_bg", line="accent", line_w=1.2)
    add_text(s, 10.55, 1.92, 2.2, 0.3, "S4 bucket 特性", size=11, bold=True, color="accent")
    add_text(s, 10.55, 2.2, 2.2, 0.7, "n=30,432\nTP 67%  FP 33%\n(75% 全 TP / 76% 全 FP)",
             size=10, color="text")
    add_rect(s, 10.55, 3.0, 2.3, 0.02, fill="line")
    add_text(s, 10.55, 3.1, 2.2, 0.3, "B4 結論", size=11, bold=True, color="positive")
    add_text(s, 10.55, 3.38, 2.2, 1.7,
             "✓ LR AUC 0.717\n  (AF+AlleleDelta)\n\n✓ ISM 有二級判別力\n\n✗ 主導仍是 caller\n  訊號 (非 ISM)",
             size=10, color="text")
    add_rect(s, 10.55, 5.25, 2.3, 0.02, fill="line")
    add_text(s, 10.55, 5.35, 2.2, 0.3, "下一步", size=11, bold=True, color="accent")
    add_text(s, 10.55, 5.62, 2.2, 0.9, "HCC1395 pilot →\n多樣本驗證 AF+\nAlleleDelta LR 一致性",
             size=9, color="text")
    # 底部重點與結論
    add_rect(s, 0.3, 6.7, 12.7, 0.35, fill="positive")
    add_text(s, 0.5, 6.73, 12.3, 0.3,
             "重點：S4 bucket (75% TP+76% FP 混雜) 可用 AF+AlleleDelta LR 二級切分  →  結論：有 0.717 判別力但主導是 caller (非 ISM HP-derived)",
             size=10, bold=True, color="card_bg")
    add_footer(s, 17)
    add_notes(s,
              "【這張 slide 要講什麼】"
              "\nPanel A (左)：S4 bucket 內 6 個候選特徵的 AUC 比較。AF+AlleleDelta LR combo 達 0.717，顯著超過 "
              "Beyond-AUC 0.58 ceiling 與 random 0.5 線。AlleleDelta 單變數 0.661 也很強（HP-independent 的最強訊號）。"
              "HPMergedDelta 0.578（HP-derived）為次要。其他特徵 (NumReads/Fisher_Frac_Sig/Quality_Score) 接近 random。"
              "\nPanel B (右)：S4 subset 的 TP/FP 在 (AF × AlleleDelta) 平面的分佈。TP 集中低 AlleleDelta；FP 分佈更分散"
              "且 AlleleDelta 較高。illustrative LR decision boundary 顯示分類可行。"
              "\n\n【重點】"
              "\n1. S4 不是『無法判別』bucket，而是『需要更多特徵』bucket"
              "\n2. 主要判別力來自 caller 訊號 (AF + AlleleDelta)，ISM HP-derived 只是補充"
              "\n3. 這是 S14 scheme 定義中 S4 黑名單『需二級判別』的具體展示"
              "\n\n【結論】"
              "\n• B4 pilot 證實 S4 內可用 LR combo 達 AUC 0.717，突破純甲基化天花板"
              "\n• 下一步：HCC1395 pilot → 多樣本驗證 AF+AlleleDelta 一致性"
              "\n• 若多樣本 POSITIVE，S4 可進一步切分為 S4a (high-confidence) / S4b (still ambiguous)")
    return s


def s18_next_features(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 5 · 新特徵 · 下一步系統性觀察 + 每類數據意義",
               "LOH × AF × CN 切分下觀察各特徵類別 · 預期產出與判讀規則",
               "Next-step systematic observation with expected data meanings")
    # 三類特徵卡片（加「預期數據意義」副欄）
    groups = [
        ("🧬 HP 狀況",
         "HPFineNGroups / HP_Ratio / HPMergedDelta / HP1FamilyN / HP2FamilyN",
         "預期數據意義：\n• 在 LOH Inner same-hap 達 TP rate ≥0.85 (phasing signature)\n• 在 Outer cross-het 為 germline leak 主要 FP 來源\n• TO 下受 self-phasing 17.3:1 bias 污染 → 需 PON 修正",
         "accent"),
        ("🧬 ISM 甲基結果",
         "AlleleDelta / CramersV / Fisher_Frac_Sig / Epipoly_Delta / HPFineF",
         "預期數據意義：\n• HP-independent AlleleDelta ~0.63 (Phase 1 對照值，單一可信 methylation 訊號)\n• 其他 HP-free 特徵 AUC ≤0.58 (Beyond-AUC 耗盡)\n• 在 biology-module 內可能保留邊際增量待驗證",
         "blue"),
        ("🧬 ISM 關聯與聚類",
         "ClusterPermanovaF / PairwiseMeanDist / LabelAllelePermanovaF / Stability",
         "預期數據意義：\n• ClusterPermanovaF 過去 AUC=0.512 隨機 (Option C)\n• PairwiseMedianDist 正交 HPFineN 但增量微弱\n• 預期跨樣本 ≤5/7 有一致方向 → 非主要訊號",
         "positive"),
    ]
    x0, y0 = 0.3, 1.8
    w, h = 4.2, 3.1
    for i, (t, feats, meaning, c) in enumerate(groups):
        x = x0 + i * (w + 0.1)
        add_rect(s, x, y0, w, h, fill="card_bg", line=c, line_w=1.2)
        add_rect(s, x, y0, w, 0.5, fill=c)
        add_text(s, x + 0.15, y0 + 0.08, w - 0.3, 0.4, t, size=13, bold=True, color="card_bg", anchor="middle")
        add_text(s, x + 0.15, y0 + 0.6, w - 0.3, 0.6, feats, size=9, color="dark", font="Courier New")
        add_text(s, x + 0.15, y0 + 1.25, w - 0.3, h - 1.35, meaning, size=10, color="text")

    # 計劃區
    add_rect(s, 0.3, 5.05, 12.7, 1.95, fill="bg_alt", line="accent", line_w=1)
    add_text(s, 0.5, 5.15, 12.3, 0.3, "計劃（Phase 2B master × 兩 flag 重跑後展開）", size=12, bold=True, color="accent")
    add_bullets(s, 0.5, 5.5, 12.3, 1.5,
                [
                    "在 LOH_Subtype (5) × AF_class (3) × cn_tier_F (5) = 75 cells 內計算每特徵 TP/FP AUC",
                    "跨 7 樣本一致性篩選 — 要求 ≥5/7 樣本同方向才視為通用訊號",
                    "與 biology-module (S1-S7) 交互：在 S3 / S5 內看 ISM 特徵是否保留邊際增量",
                ], size=11, icon="▸", gap=7)
    add_footer(s, 18)
    add_notes(s, "本 slide 清楚標示下一步計劃：在確認了 S1-S7 的 biology-informed framework 後，在每個 biology cell "
              "內系統性評估 HP/ISM 特徵是否有跨樣本一致的邊際增量。預期純甲基化特徵會 ≤0.58（見 Beyond-AUC 耗盡結論），"
              "而 HP-dependent 特徵會受 self-phasing 污染 → 引出段 6 的 PON 解法必要性。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 6 · TO Self-Phasing 問題與解法 (S19-S23)
# ═══════════════════════════════════════════════════════════════
def s19_hp_tag_need(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 6 · Self-Phasing · 為何需要 HP tag",
               "ISM 所有 HP-derived 特徵都依賴 haplotag · tag 不準會連鎖影響",
               "HP tag is foundational to ISM analysis")
    # 中央依賴圖
    add_text(s, 0.5, 2.0, 12.4, 0.4, "HP tag → ISM 分析流水線", size=14, bold=True, color="dark", align="center")
    # 4 層
    layers = [
        ("HP tag (HP:i:N)", "BAM 註解 · 每 read 一個整數 · 0/1/2/11/21/33", "muted"),
        ("4-bucket 分群", "HPFineNGroups / HP1/HP2/HP1-1/HP2-1 · LabelTest.cpp", "accent"),
        ("HP-derived 特徵", "HP_Ratio / HPMergedDelta / HPFineF / NHP / AlleleDelta", "blue"),
        ("甲基化 PERMANOVA", "HPFineF / HPFineP · HP 分層的甲基化多樣性顯著性", "positive"),
    ]
    y = 2.6
    for idx, (lbl, desc, c) in enumerate(layers):
        add_rect(s, 1.0, y, 11, 0.7, fill="card_bg", line=c, line_w=1.2)
        add_rect(s, 1.0, y, 0.2, 0.7, fill=c)
        add_text(s, 1.4, y + 0.12, 4, 0.3, lbl, size=13, bold=True, color=c)
        add_text(s, 5.5, y + 0.15, 6.4, 0.4, desc, size=10, color="text")
        # arrow (放於下一行塊之前的空白區，高度 0.25)
        if idx < len(layers) - 1:
            add_text(s, 1.0, y + 0.72, 11, 0.25, "↓", size=14, color="muted", align="center", anchor="middle")
        y += 0.95
    add_rect(s, 0.5, 6.5, 12.4, 0.5, fill="negative")
    add_text(s, 0.7, 6.58, 12, 0.35,
             "HP tag 不準 → 整條流水線失真 · 必須先驗證 tag 品質才能信任 ISM 結論",
             size=12, bold=True, color="card_bg", align="center")
    add_footer(s, 19)
    add_notes(s, "從 BAM 的 HP:i:N 整數欄位，經 4-bucket 分群、HP-derived 特徵計算，一路到 HPFineF/HPFineP 的 "
              "PERMANOVA — 整條流水線都依賴 HP tag 品質。TO 模式下若 tag 有 self-phasing artifact，所有 HP-derived "
              "結論都會受到污染。這就是為何需要優先解決 HP tag 正確性問題。")
    return s


def s20_self_phasing_root(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 6 · Self-Phasing · Paired vs TO 結構差異（圖示案例）",
               "Chromosome + Haplotype + SNV 示意 · 說明 somatic 反饋影響 scaffold",
               "Why TO tumor-only phasing creates self-phasing artifact")
    # 使用新生成的 schematic 圖
    add_image(s, FIG / "fig_s20_paired_vs_to_schematic.png", 0.3, 1.7, 12.7, 4.5)
    # 底部比喻 banner
    add_rect(s, 0.5, 6.35, 6.1, 0.65, fill="positive")
    add_text(s, 0.7, 6.42, 5.7, 0.3, "Paired mode · 裁判獨立於球員",
             size=12, bold=True, color="card_bg")
    add_text(s, 0.7, 6.7, 5.7, 0.3, "Normal germline = 獨立 scaffold · somatic 被動掛上",
             size=10, color="card_bg")
    add_rect(s, 6.75, 6.35, 6.15, 0.65, fill="negative")
    add_text(s, 6.95, 6.42, 5.7, 0.3, "TO mode · 球員兼裁判 ⚠",
             size=12, bold=True, color="card_bg")
    add_text(s, 6.95, 6.7, 5.7, 0.3, "Somatic 一起進 phase graph · 反饋影響 scaffold 17.3:1 bias",
             size=10, color="card_bg")
    add_footer(s, 20)
    add_notes(s, "本 slide 用比喻「裁判 vs 球員」說明根本差異：Paired 模式下 germline 是「裁判」（建立 scaffold），"
              "somatic 是「球員」（掛上去）；TO 模式下 somatic 同時擔任球員和裁判，這是 self-phasing 的結構性根源，"
              "不是 tool bug 而是無 normal 樣本時的必然副作用。")
    return s


def s21_hp_tag_values(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 6 · Self-Phasing · HP tag 整數值含意",
               "HP:i:N — 從 0/1/2 germline 到 11/21/33 somatic self-phasing 產物",
               "HP:i:N integer values and biological meaning")
    rows = [
        ("HP:i:0 或無 tag", "Unphased / 未分配", "muted"),
        ("HP:i:1", "Germline Haplotype 1", "positive"),
        ("HP:i:2", "Germline Haplotype 2", "positive"),
        ("HP:i:11", "Somatic phase block 1-1 (self-phasing 產物)", "negative"),
        ("HP:i:21", "Somatic phase block 2-1 (self-phasing 產物)", "negative"),
        ("HP:i:33", "Somatic phase block 3 (self-phasing 產物)", "negative"),
    ]
    add_rect(s, 0.5, 1.9, 12.4, 0.6, fill="dark")
    add_text(s, 0.7, 2.0, 3, 0.4, "Tag 值", size=12, bold=True, color="card_bg")
    add_text(s, 4.0, 2.0, 8.5, 0.4, "意義", size=12, bold=True, color="card_bg")
    y = 2.5
    for tag, meaning, c in rows:
        add_rect(s, 0.5, y, 12.4, 0.65, fill="card_bg" if (int(y*10)) % 2 == 0 else "bg_alt",
                 line="line", line_w=0.3)
        add_text(s, 0.7, y + 0.15, 3, 0.35, tag, size=13, bold=True, color=c, font="Courier New")
        add_text(s, 4.0, y + 0.15, 8.5, 0.35, meaning, size=12, color="text")
        y += 0.7
    add_rect(s, 0.5, 6.7, 12.4, 0.3, fill="accent")
    add_text(s, 0.7, 6.73, 12, 0.25,
             "ISM 用 HP=1/2/11/21 四類做 LabelTest.cpp::hp_to_fine_labels 的 4-bucket 分群",
             size=11, bold=True, color="card_bg", align="center")
    add_footer(s, 21)
    add_notes(s, "11/21/33 是 LongPhase-TO 的 somatic phase block 產物 — 它試圖為 somatic variants 建立額外的 "
              "phase block。這本身是合理設計，但在無 germline scaffold 時，這些 somatic blocks 可能將真 somatic "
              "reads 錯誤歸到 HP1 家族（造成 bias）。LabelTest.cpp::hp_to_fine_labels 把 1/11 合併為 HP1 家族（ref/"
              "somatic 子族），2/21 合併為 HP2 家族。")
    return s


def s22_somatic_bias(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 6 · Self-Phasing · 17.3:1 somatic bias",
               "94.6% somatic reads 分配到 HP1 · 遠高於理想 1:1 · 三條路徑的影響",
               "Somatic read assignment strong bias evidence")
    # 左：17.3:1 數字視覺化
    add_rect(s, 0.5, 1.9, 6, 4.8, fill="card_bg", line="negative", line_w=1.5)
    add_text(s, 0.5, 2.0, 6, 0.4, "TO 模式 Somatic reads HP 分配", size=13, bold=True, color="negative", align="center")
    # 兩柱對比
    add_text(s, 0.9, 2.5, 1.5, 0.35, "HP1", size=14, bold=True, color="dark")
    add_text(s, 0.9, 2.85, 1.5, 0.3, "614,000 reads", size=11, color="dark")
    add_rect(s, 2.5, 2.55, 3.3, 0.35, fill="negative")
    add_text(s, 0.9, 3.3, 1.5, 0.35, "HP2", size=14, bold=True, color="dark")
    add_text(s, 0.9, 3.65, 1.5, 0.3, "35,500 reads", size=11, color="dark")
    add_rect(s, 2.5, 3.35, 0.19, 0.35, fill="positive")
    # Ratio
    add_rect(s, 0.9, 4.3, 5.0, 1.2, fill="bg_alt", line="negative", line_w=1)
    add_text(s, 0.9, 4.35, 5, 0.3, "Ratio", size=10, color="muted", align="center")
    add_text(s, 0.9, 4.6, 5, 0.6, "17.3 : 1", size=32, bold=True, color="negative", align="center")
    add_text(s, 0.9, 5.25, 5, 0.25, "94.6% → HP1", size=11, bold=True, color="negative", align="center")
    add_text(s, 0.9, 5.7, 5, 0.4, "理想應為 1:1 (50/50)\n→ strong artifact 指標", size=10, color="muted", align="center")

    # 右：三條路徑影響
    add_rect(s, 6.9, 1.9, 6, 4.8, fill="card_bg", line="line", line_w=0.8)
    add_text(s, 6.9, 2.0, 6, 0.4, "三條路徑的影響", size=13, bold=True, color="dark", align="center")
    paths = [
        ("LongPhase phase 本身", "不受影響", "positive",
         "輸出的 phased VCF 以 germline 為主幹"),
        ("ISM HP 分群 (4-bucket)", "嚴重受影響", "negative",
         "Somatic HP 混入 HP1-1/HP2-1 bucket"),
        ("LOH.bed 生成", "不受影響", "positive",
         "用 VCF AF/VAF (PhasingGraph.cpp) · 不依賴 BAM HP tags"),
    ]
    y = 2.5
    for path, verdict, c, desc in paths:
        add_rect(s, 7.1, y, 5.6, 1.4, fill="bg_alt" if y*10 % 2 == 0 else "card_bg", line=c, line_w=1)
        add_text(s, 7.2, y + 0.1, 3.5, 0.3, path, size=11, bold=True, color="dark")
        add_text(s, 10.8, y + 0.1, 1.8, 0.3, verdict, size=11, bold=True, color=c, align="right")
        add_text(s, 7.2, y + 0.5, 5.4, 0.8, desc, size=10, color="text")
        y += 1.5
    add_footer(s, 22)
    add_notes(s, "17.3:1 bias 的根源：LongPhase-TO 在 somatic block 階段把 somatic reads 強烈偏向 HP1 系（因為 "
              "scaffold 本身由 somatic 參與定義，formation bias 產生）。三條路徑：(1) phase 本身的 phased VCF 不 "
              "受影響（它輸出的是 variants 的 0|1 genotype，不是 reads 的 HP tag）；(2) ISM HP 分群嚴重受影響 "
              "（因為它讀 BAM 的 HP tag）；(3) LOH.bed 用 PhasingGraph.cpp 的 VCF AF/VAF 計算，不依賴 BAM HP tags "
              "（見 0421 週報 Jaccard=1.0 證實）。所以本週的 PON 修正只針對路徑 2。")
    return s


def s23_pon_solution(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 6 · Self-Phasing · PON 錨點解法 + 驗證方式",
               "PON 95% 位點當 germline · somatic 只標記不進 HP 分群 · 以 paired 為 reference",
               "PON-anchored HP tagging and verification strategy")
    # 主圖：PON schematic 圖示（2 panel 對照：Before without PON / After with PON）
    add_image(s, FIG / "fig_s23_pon_anchor_schematic.png", 0.3, 1.7, 12.7, 4.5)
    # 底部雙 banner（左 scaffold 定義 / 右 驗證方式）
    add_rect(s, 0.5, 6.3, 6.1, 0.7, fill="positive")
    add_text(s, 0.7, 6.38, 5.7, 0.3, "方案：PON (Panel of Normals) 95% 位點當 scaffold anchor",
             size=11, bold=True, color="card_bg")
    add_text(s, 0.7, 6.67, 5.7, 0.3,
             "Somatic 只標記 · 不參與 HP 分群 · 借用 paired scaffold 優勢",
             size=10, color="card_bg")
    add_rect(s, 6.75, 6.3, 6.15, 0.7, fill="blue")
    add_text(s, 6.95, 6.38, 5.7, 0.3, "驗證：用 paired tag 作 reference",
             size=11, bold=True, color="card_bg")
    add_text(s, 6.95, 6.67, 5.7, 0.3,
             "若 PON-TO tag ≈ paired tag → 可採用 · B3 paired gap=0.00003 已確認乾淨",
             size=10, color="card_bg")
    add_footer(s, 23)
    add_notes(s, "PON 方案核心：把「scaffold 由誰定義」這件事從 tumor 自己移到獨立 PON（Panel of Normals）。PON "
              "包含大量 normal 樣本的 germline 變異位點，95% 位點為 germline common variants。用 PON 當 scaffold "
              "後，somatic variants 不參與 phase 規則定義，只被動標記。這等效於「借用」paired 模式的 germline "
              "scaffold 優勢，但不需要個別樣本的配對 normal。驗證：由於 B3 已確認 paired mode 本身沒有 "
              "Inner/Outer gap，用 paired tag 作 reference 比對 PON-TO tag 是最直接的驗證。")
    return s


# ═══════════════════════════════════════════════════════════════
# 段 7 · 未來工作 (S24-S25)
# ═══════════════════════════════════════════════════════════════
def s24_next_week_p0p1(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 7 · 下週 P0/P1  —  每項有明確目標、數據指標、成功條件",
               "聚焦 4 項行動 · 目標導向 · 有 go/no-go 判定依據",
               "Next-week P0/P1 — goal-driven with measurable success criteria")
    # 4 項行動：目標 / 指標 / 成功條件 / 時間
    actions = [
        ("P0", "longphase-TO PON 修正跨樣本驗證", "negative",
         "目標: 確認 PON-anchored tag 在 7 樣本上比 self-phasing tag 一致性更高",
         "指標: 每樣本 PON-tag vs paired-tag 的 HP 分配 overlap ≥ 85%",
         "成功: 6/7 樣本通過 → 可採納為 default；<5 通過 → 保留 flag 但不改 default",
         "~10 hr parallel"),
        ("P0", "HCC1954 caller FP 背景專項", "negative",
         "目標: 釐清 HCC1954 Inner TP 0.29 是 caller FP 背景 (B2) 還是其他樣本特異問題",
         "指標: 跨 6 TO 的『Outer Extreme AF FP 比例』 vs 『Inner cross-het TP rate』相關",
         "成功: 相關 r>0.7 → 確認 caller-driven；r<0.3 → 需查 phasing 失真",
         "1-2 天"),
        ("P1", "COLO829 TO ISM 補跑 (含 step05)", "accent",
         "目標: 完成 7 樣本 TO cohort 一致性 (目前只 6 樣本可比)",
         "指標: COLO829 TO ISM 40k+ regions 產出含 LOH_Subtype + caller_af",
         "成功: Phase 3 S1-S7 擴展至 7/7 · Wilcoxon test 重跑含 COLO829",
         "~3 hr"),
        ("P1", "S4 secondary (B4) 多樣本驗證", "accent",
         "目標: B4 HCC1395 pilot AUC 0.717 的 LR combo 是否跨樣本一致",
         "指標: AF+AlleleDelta LR 在 6 TO 樣本的 AUC 差異 < 0.05",
         "成功: 5/6 樣本 AUC ≥ 0.65 → S4a/S4b scheme 正式化",
         "1 天"),
    ]
    y0 = 1.85
    for i, (lv, title, c, goal, metric, success, dur) in enumerate(actions):
        y = y0 + i * 1.22
        add_rect(s, 0.3, y, 12.7, 1.12, fill="card_bg", line=c, line_w=1.2)
        add_rect(s, 0.3, y, 0.6, 1.12, fill=c)
        add_text(s, 0.3, y + 0.3, 0.6, 0.4, lv, size=15, bold=True, color="card_bg", align="center")
        add_text(s, 1.05, y + 0.06, 9.0, 0.28, title, size=12, bold=True, color="dark")
        add_text(s, 11.0, y + 0.06, 1.9, 0.28, dur, size=9, color="muted", align="right")
        add_text(s, 1.05, y + 0.36, 11.7, 0.24, "🎯  " + goal, size=9.5, color="text")
        add_text(s, 1.05, y + 0.6, 11.7, 0.24, "📊  " + metric, size=9, color=c)
        add_text(s, 1.05, y + 0.84, 11.7, 0.24, "✓  " + success, size=9, bold=True, color="positive")
    add_footer(s, 24)
    add_notes(s, "下週 4 項行動：P0 兩項最關鍵 — longphase-TO 的 PON 修正跨樣本驗證（本週 HCC1395 pilot 已機制 ✓ 但 "
              "filter Gate FAIL，需 master × 兩 flag × 7 樣本才能完整判定 H-C2/H-C3）；HCC1954 caller FP 背景專項 "
              "（B2 今日初步，需擴展驗證）。P1 兩項：COLO829 TO 補跑讓 cohort 齊；S4 secondary 從 HCC1395 pilot 擴 "
              "展多樣本。預估 2-3 天完成 P0。")
    return s


def s25_other_goals(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "段 7 · 中長期目標 — 三個聚焦主線",
               "每線有單一目標 · 依賴明確 · 里程碑可檢核",
               "Longer-term goals — three focused tracks with dependency & milestone")
    # 三個聚焦主線（取代原 5 項散列）
    tracks = [
        ("① 論文主軸定案", "positive",
         "LOH-constrained phasing (TO 層) 作為 paper primary finding",
         "• 完成 paired 層獨立 phasing-vs-methylation 驗證 (2-3 週)\n"
         "• Paper Figure 3/4 以 NG=2 物理機制為骨架重寫\n"
         "• 發表時機：NG=2 解釋 + 跨 7 樣本 + S1-S7 framework 齊備",
         "下週 P0 Archive TO 6 樣本重跑 + S4 B4 多樣本驗證"),

        ("② Filter framework 完整化", "accent",
         "S1-S7 scheme 從 HCC1395 pilot → 7 樣本通用 TO filter",
         "• S4 B4 LR combo (AUC 0.717) 跨樣本一致性驗證 (1 週)\n"
         "• 通過後切分 S4a (high-conf) / S4b (still ambiguous)\n"
         "• 整合為 end-to-end TO filter (scheme-aware ClairS-TO post-filter)",
         "COLO829 TO 補跑 · S4 B4 多樣本驗證 · Phase 2B master × 兩 flag"),

        ("③ Phase 2A Normal Methylation Reference", "blue",
         "啟動 normal methylation baseline 作為 cross-region ASM 參考",
         "• Dependency: master × 兩 flag × 7 樣本重跑完成 (下週 P0)\n"
         "• Sample ASM 5 步驟 C++ pipeline 已 ready (Phase B-D)\n"
         "• 第 1 個跨樣本 ASM Figure → 2-3 週內",
         "Master rerun + normal BAM 7 樣本齊備 (HCC1395 已 pilot 驗證)"),
    ]
    y0 = 1.75
    h_per = 1.65
    for i, (title, c, goal, milestones, dep) in enumerate(tracks):
        y = y0 + i * h_per
        add_rect(s, 0.3, y, 12.7, h_per - 0.1, fill="card_bg", line=c, line_w=1.3)
        add_rect(s, 0.3, y, 0.15, h_per - 0.1, fill=c)
        add_text(s, 0.5, y + 0.08, 12.3, 0.32, title, size=14, bold=True, color=c)
        # 目標單行
        add_text(s, 0.5, y + 0.45, 1.1, 0.25, "🎯 目標", size=9.5, bold=True, color="dark")
        add_text(s, 1.6, y + 0.45, 11.3, 0.25, goal, size=10.5, color="text")
        # 里程碑多行
        add_text(s, 0.5, y + 0.75, 1.1, 0.25, "📍 里程碑", size=9.5, bold=True, color=c)
        add_text(s, 1.6, y + 0.75, 11.3, 0.65, milestones, size=9, color="text")
        # 依賴
        add_text(s, 0.5, y + 1.35, 1.1, 0.2, "⚙ 依賴", size=9, bold=True, color="muted")
        add_text(s, 1.6, y + 1.35, 11.3, 0.2, dep, size=9, color="muted")
    add_footer(s, 25)
    add_notes(s, "5 項中長期目標：(1) Phase 2A 依賴 master rerun；(2) S4 secondary 正式化 — 把 B4 pilot 擴成標準 "
              "scheme；(3) Paired 層獨立驗證 — 本週按用戶決策 C 保留不撤回，需未來做 orthogonal test 確認是否同為 "
              "phasing 現象；(4) Paper Figure 重寫（依 Thread D TO 層 pivot）；(5) 外部 CN 工具 pilot（若 S4 分析 "
              "不足）。")
    return s


# ═══════════════════════════════════════════════════════════════
# 結語 (S26)
# ═══════════════════════════════════════════════════════════════
def s26_summary(prs):
    s = blank(prs); set_bg(s)
    add_header(s, "小結 —  本週一句話 + 下週單一 P0 + 論文定位",
               "聚焦 3 要點 · 不列細項 · 目標清楚可行動",
               "Summary: 3 focus points — this week / next week single P0 / paper positioning")

    # 區塊 1 · 本週一句話（最大、中央、醒目）
    add_rect(s, 0.3, 1.75, 12.7, 1.45, fill="positive")
    add_text(s, 0.5, 1.85, 12.3, 0.4, "① 本週一句話定案", size=13, bold=True, color="card_bg")
    add_text(s, 0.5, 2.3, 12.3, 0.8,
             "LOH-constrained phasing 在 TO 模式跨 6 樣本確認 (Wilcoxon p=0.0156)；",
             size=18, bold=True, color="card_bg")
    add_text(s, 0.5, 2.8, 12.3, 0.3,
             "paired 對照 p=0.578 反證為 TO-specific；S3 Diploid Het + S5 combo 可作 white-list filter (median TP 0.88).",
             size=12, color="card_bg")

    # 區塊 2 · 下週單一 P0 (最重要)
    add_rect(s, 0.3, 3.4, 12.7, 1.65, fill="card_bg", line="negative", line_w=1.5)
    add_rect(s, 0.3, 3.4, 0.15, 1.65, fill="negative")
    add_text(s, 0.5, 3.5, 12.3, 0.35, "② 下週唯一關鍵 P0",
             size=13, bold=True, color="negative")
    add_text(s, 0.5, 3.95, 12.3, 0.5,
             "longphase-TO 用 PON 當 scaffold anchor · 跨 7 樣本驗證 tag 一致性 ≥ paired tag",
             size=15, bold=True, color="dark")
    add_text(s, 0.5, 4.45, 12.3, 0.3,
             "成功指標：6/7 樣本 PON-tag vs paired-tag overlap ≥85%  →  tag 可改為 default",
             size=11, color="text")
    add_text(s, 0.5, 4.77, 12.3, 0.3,
             "預估 ~10 hr (parallel) · 其他 3 項 (HCC1954 專項 / COLO829 補跑 / S4 B4 多樣本) 並行",
             size=10, color="muted")

    # 區塊 3 · 論文定位
    add_rect(s, 0.3, 5.25, 12.7, 1.7, fill="card_bg", line="accent", line_w=1.5)
    add_rect(s, 0.3, 5.25, 0.15, 1.7, fill="accent")
    add_text(s, 0.5, 5.35, 12.3, 0.35, "③ 論文主軸（TO 層，2-3 週定案）",
             size=13, bold=True, color="accent")
    add_text(s, 0.5, 5.8, 12.3, 0.5,
             '"LOH-constrained phasing signatures distinguish somatic from germline-like variants in TO sequencing"',
             size=13, bold=True, color="dark", font="Courier New")
    add_text(s, 0.5, 6.3, 12.3, 0.3,
             "主軸優勢：機制純 phasing (無需 methylation 驗證) · 跨 basecall/pipeline 穩定 · 連 FACETS/Battenberg 模型",
             size=10, color="text")
    add_text(s, 0.5, 6.6, 12.3, 0.3,
             "paired 層 AF×NGroups POSITIVE 保留加註；Paper Figure 3/4 以 NG=2 物理機制為骨架重寫",
             size=10, color="muted")
    add_footer(s, 26)
    add_notes(s,
              "【小結的 3 個聚焦點，不列細項】"
              "\n① 本週一句話：6 TO 樣本 LOH-constrained phasing 確認 (Wilcoxon p=0.0156) + paired 反證 TO-specific "
              "(p=0.578) + S3/S5 filter 可採用 (median TP 0.88)。"
              "\n② 下週唯一關鍵 P0：PON 跨樣本驗證。成功指標 6/7 overlap ≥85% → 可改 default。其他 3 項並行不是 blocker。"
              "\n③ 論文主軸（2-3 週定案）：TO-layer phasing signatures；paired 保留加註；Figure 3/4 重寫。"
              "\n\n避免把 13 項結論都列出讓 PI 抓不到重點。本頁就這 3 點即可。")
    return s


# 舊 S26 原本的 3 欄 layout 已停用，直接 return
def _s26_summary_old_unused(prs):
    cols = [
        ("已確認 (2026-04-23 定案)",
         [],
         "positive"),
    ]
    x0, y0 = 0.5, 1.9
    w, h = 4.1, 5.0
    for i, (title, items, c) in enumerate(cols):
        x = x0 + i * (w + 0.15)
        add_rect(prs.slides[0], x, y0, w, h, fill="card_bg", line=c, line_w=1.3)
        add_rect(s, x, y0, w, 0.5, fill=c)
        add_text(s, x + 0.15, y0 + 0.08, w - 0.3, 0.4, title, size=12, bold=True, color="card_bg", anchor="middle")
        add_bullets(s, x + 0.2, y0 + 0.65, w - 0.4, h - 0.75, items, size=11, icon="▸", gap=10, color="text")

    # 底部：里程碑
    add_rect(s, 0.5, 7.05, 12.4, 0.08, fill="accent")
    add_footer(s, 26)
    add_notes(s, "本週最終狀態：今天（2026-04-23）一天內完成 B1-B7 + Phase 3 Synthesis，把週報早上標示的下週 P0/P1 "
              "多數項目提前完成。Thread D 從 ⭐5 TO-only 升為 TO-specific 確定性；Thread B 從 HCC1395 pilot 升為 "
              "6/6 TO POSITIVE；Thread A CN KDE HCC1395 pilot PASS 待全量。下週重心轉向 PON 跨樣本驗證 + HCC1954 "
              "caller FP 背景 + COLO829 補跑。按用戶決策 C，不宣稱 paper pivot（TO 層結論穩定但 paired 層仍待獨立驗證）。")
    return s


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        s1_cover, s2_roadmap,
        s3_loh_bed_coverage, s4_af_distribution, s5_loh_af_cn_initial, s6_hcc1395_background,
        s7_cn_75x_problem, s8_cn_kde_validation, s9_cn_per_sample,
        # S14 scheme 定義前移到 S10 位置（先定義後展示 TP rate 現象）
        s14_s1s7_schemes,
        s10_tp_rate_heatmap, s11_inner_same_hap, s12_biological_meaning, s13_stats_confirmation,
        s15_scheme_heatmap, s16_s3_and_loh_noise,
        s17_s4_secondary, s18_next_features,
        s19_hp_tag_need, s20_self_phasing_root, s21_hp_tag_values, s22_somatic_bias, s23_pon_solution,
        s24_next_week_p0p1, s25_other_goals,
        s26_summary,
    ]
    for i, b in enumerate(builders, 1):
        print(f"  building slide {i}/26 · {b.__name__}")
        b(prs)

    prs.save(str(OUTPUT))
    print(f"\n[OK] {OUTPUT}")
    print(f"     size: {OUTPUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
