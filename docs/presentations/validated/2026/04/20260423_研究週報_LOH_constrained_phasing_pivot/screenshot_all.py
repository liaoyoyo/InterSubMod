#!/usr/bin/env python3
"""Generate wireframe screenshots of all PPTX slides with CJK support.

Renders all shapes (rectangles, text boxes, images, connectors) from python-pptx
onto PIL canvases. Outputs PNG files to screenshots/ directory.
"""

import os, sys, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

# --- Config ---
PPTX_PATH = "20260423_NG2_LOH_constrained_phasing_pivot.pptx"
OUT_DIR = Path("screenshots")
_ISSUES = []  # collected overflow / img-fail log entries
SCALE = 150  # DPI-ish scale factor
SLIDE_W_INCH = 13.333
SLIDE_H_INCH = 7.5
W_PX = int(SLIDE_W_INCH * SCALE)
H_PX = int(SLIDE_H_INCH * SCALE)

# Font paths — CJK for Chinese, Latin for English/numerals/symbols
CJK_FONT = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
LATIN_FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
LATIN_BOLD = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
MONO_FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf"


def _font(size, bold=False, mono=False, cjk=False):
    """Return a single ImageFont. cjk=True when rendering CJK glyph segment."""
    if mono:
        path = MONO_FONT
    elif cjk:
        path = CJK_FONT
    else:
        path = LATIN_BOLD if bold else LATIN_FONT
    try:
        return ImageFont.truetype(path, size)
    except Exception:
        return ImageFont.load_default()


def _is_cjk(ch):
    cp = ord(ch)
    return (
        0x4E00 <= cp <= 0x9FFF   # CJK Unified
        or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
        or 0x3000 <= cp <= 0x303F  # CJK Symbols and Punctuation
        or 0x3040 <= cp <= 0x30FF  # Hiragana/Katakana
        or 0xFF00 <= cp <= 0xFFEF  # Full-width forms
        or 0x2E80 <= cp <= 0x2EFF  # CJK Radicals
        or 0x2F00 <= cp <= 0x2FDF  # Kangxi Radicals
    )


def _segment(text):
    """Split text into (is_cjk, substring) runs."""
    if not text:
        return []
    runs = []
    cur_cjk = _is_cjk(text[0])
    buf = text[0]
    for ch in text[1:]:
        is_cjk = _is_cjk(ch)
        if is_cjk == cur_cjk:
            buf += ch
        else:
            runs.append((cur_cjk, buf))
            cur_cjk = is_cjk
            buf = ch
    runs.append((cur_cjk, buf))
    return runs


def _draw_mixed(draw, x, y, text, *, size, bold=False, color=(0, 0, 0), mono=False):
    """Draw text supporting mixed Latin/CJK via per-char font fallback."""
    cx = x
    for is_cjk, run in _segment(text):
        f = _font(size, bold=bold, mono=mono, cjk=is_cjk)
        draw.text((cx, y), run, font=f, fill=color)
        bbox = draw.textbbox((0, 0), run, font=f)
        cx += bbox[2] - bbox[0]
    return cx - x  # total width drawn


def _text_w(draw, text, *, size, bold=False, mono=False):
    """Measure width of mixed text with per-segment fonts."""
    w_total = 0
    for is_cjk, run in _segment(text):
        f = _font(size, bold=bold, mono=mono, cjk=is_cjk)
        bbox = draw.textbbox((0, 0), run, font=f)
        w_total += bbox[2] - bbox[0]
    return w_total

def emu_to_px(emu):
    return int(emu / 914400 * SCALE)

def _color_tuple(rgb_obj):
    """Extract (R,G,B) from various pptx color objects."""
    if rgb_obj is None:
        return None
    try:
        # RGBColor object
        return (rgb_obj[0], rgb_obj[1], rgb_obj[2])
    except:
        pass
    try:
        s = str(rgb_obj)
        if len(s) == 6:
            return (int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
    except:
        pass
    return None

def _get_fill_color(shape):
    """Try to get fill color from shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return _color_tuple(fc.rgb)
    except:
        pass
    return None

def _get_shape_bounds(shape):
    """Get (x, y, w, h) in pixels."""
    return (emu_to_px(shape.left), emu_to_px(shape.top),
            emu_to_px(shape.width), emu_to_px(shape.height))

def _extract_text(shape):
    """Extract all text from a shape's text_frame."""
    lines = []
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            line = ""
            for run in para.runs:
                line += run.text
            if line.strip():
                lines.append(line.strip())
    except:
        pass
    return lines

def _get_text_props(shape):
    """Get font size and bold from first run."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                bold = run.font.bold
                return (int(sz / 12700) if sz else 10, bool(bold))
    except:
        pass
    return (10, False)

def _wrap_text(text, size, max_w, draw, *, bold=False, mono=False):
    """Char-level wrapping using mixed-font width measurement."""
    if not text:
        return []
    if _text_w(draw, text, size=size, bold=bold, mono=mono) <= max_w:
        return [text]
    lines, current = [], ""
    for ch in text:
        test = current + ch
        if _text_w(draw, test, size=size, bold=bold, mono=mono) > max_w and current:
            lines.append(current)
            current = ch
        else:
            current = test
    if current:
        lines.append(current)
    return lines

def render_slide(slide, idx, out_dir):
    """Render one slide to PNG."""
    img = Image.new("RGB", (W_PX, H_PX), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Draw border
    draw.rectangle([0, 0, W_PX-1, H_PX-1], outline=(200, 200, 200), width=2)

    # Page number
    _draw_mixed(draw, 10, 5, f"P{idx+1}", size=18, bold=True, color=(100, 100, 100))

    for shape in slide.shapes:
        x, y, w, h = _get_shape_bounds(shape)

        # Skip tiny shapes
        if w < 3 and h < 3:
            continue

        # Get fill color
        fill_color = _get_fill_color(shape)

        # --- Image: real render ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                import io
                pic = Image.open(io.BytesIO(blob)).convert("RGB")
                pic_r = pic.resize((max(1, w), max(1, h)), Image.LANCZOS)
                img.paste(pic_r, (x, y))
                draw.rectangle([x, y, x+w, y+h], outline=(180,180,180), width=1)
            except Exception as e:
                draw.rectangle([x, y, x+w, y+h], outline=(100,150,200), width=2, fill=(230,240,250))
                _draw_mixed(draw, x+4, y+4, f"[IMG FAIL: {type(e).__name__}]",
                            size=12, color=(200, 80, 80))
                _ISSUES.append(f"P{idx+1:02d} · IMG FAIL · {type(e).__name__}")
            continue

        # --- Connector / Line ---
        if hasattr(shape, 'begin_x') or shape.shape_type in (MSO_SHAPE_TYPE.LINE,):
            draw.line([x, y, x+w, y+h], fill=(150,150,150), width=2)
            continue

        # --- Rectangle / TextBox / AutoShape ---
        has_text = bool(_extract_text(shape))

        # Draw fill
        if fill_color:
            draw.rectangle([x, y, x+w, y+h], fill=fill_color, outline=None)
            # Add subtle border
            brightness = sum(fill_color) / 3
            border_c = tuple(max(0, c - 40) for c in fill_color)
            draw.rectangle([x, y, x+w, y+h], outline=border_c, width=1)
        elif has_text and w > 20 and h > 20:
            # Text box without fill - light dashed outline
            draw.rectangle([x, y, x+w, y+h], outline=(220, 220, 220), width=1)

        # Draw text
        if has_text:
            lines = _extract_text(shape)
            pt_size, bold = _get_text_props(shape)
            font_px = max(8, int(pt_size * SCALE / 96))
            font_px = min(font_px, 28)

            is_mono = False
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and ('consol' in run.font.name.lower()
                                              or 'courier' in run.font.name.lower()
                                              or 'mono' in run.font.name.lower()):
                            is_mono = True
                            break
            except:
                pass

            # Determine text color
            txt_color = (0, 0, 0)
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            txt_color = _color_tuple(run.font.color.rgb) or (0, 0, 0)
                            break
                    break
            except:
                pass

            if fill_color and sum(fill_color) / 3 < 100:
                txt_color = (255, 255, 255)

            pad = 4
            max_text_w = w - pad * 2
            ty = y + pad
            overflow = False
            for line in lines:
                wrapped = _wrap_text(line, font_px, max_text_w, draw,
                                     bold=bold, mono=is_mono)
                for wl in wrapped:
                    if ty + font_px > y + h:
                        _draw_mixed(draw, x + w - 70, y + h - 13,
                                    "⚠OVERFLOW", size=10, bold=True, color=(220, 0, 0))
                        snippet = line[:30].replace("\n", " ")
                        _ISSUES.append(
                            f"P{idx+1:02d} · TEXT OVERFLOW · box=({x},{y},{w}x{h}) · '{snippet}...'"
                        )
                        overflow = True
                        break
                    _draw_mixed(draw, x + pad, ty, wl,
                                size=font_px, bold=bold, color=txt_color, mono=is_mono)
                    ty += font_px + 2
                if overflow:
                    break

    # Save
    out_path = out_dir / f"P{idx+1:02d}.png"
    img.save(str(out_path), "PNG")
    return str(out_path)

def main():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation(PPTX_PATH)
    total = len(prs.slides)
    print(f"Rendering {total} slides...")

    for i, slide in enumerate(prs.slides):
        path = render_slide(slide, i, OUT_DIR)
        print(f"  P{i+1:02d} -> {path}")

    print("\n" + "=" * 60)
    if _ISSUES:
        print(f"[ISSUES] {len(_ISSUES)} item(s) detected:")
        for it in _ISSUES:
            print(f"  - {it}")
    else:
        print("[ISSUES] none detected · all text fits · all images loaded ✓")
    print("=" * 60)

    print(f"\nDone. {total} screenshots in {OUT_DIR}/")

if __name__ == "__main__":
    main()
