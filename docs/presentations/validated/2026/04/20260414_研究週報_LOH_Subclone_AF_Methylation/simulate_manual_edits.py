#!/usr/bin/env python3
"""
Simulate the manually edited PPTX by applying all changes from the diff report
to the current code-generated PPTX. This creates a reference file for diff verification.

Usage:
  python3 simulate_manual_edits.py
  → reads current PPTX, applies all 183 target values, saves as simulated version
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import shutil
import sys

EMU_PER_INCH = 914400

# All changes extracted from 20260416_手動編輯差異報告.md
# Format: (slide_1based, shape_name, property, target_value)
# property: "left", "top", "width", "height", "font_size"
# text changes handled separately

POSITION_CHANGES = [
    # P01 — build_slide_1_cover
    (1, "Picture 3", "left", 7.56),
    (1, "Picture 3", "top", 1.08),

    # P04 — build_slide_3_to_motivation
    (4, "TextBox 6", "top", 1.92),
    (4, "TextBox 7", "top", 2.14),

    # P06 — build_slide_ism_loh_impact
    (6, "TextBox 30", "left", 10.12),
    (6, "TextBox 30", "top", 3.57),
    (6, "TextBox 30", "width", 2.56),
    (6, "TextBox 30", "height", 0.79),

    # P07 — build_slide_5_dual_definition
    (7, "TextBox 22", "left", 1.05),
    (7, "TextBox 22", "top", 5.35),
    (7, "TextBox 22", "width", 11.63),
    (7, "TextBox 22", "height", 0.79),

    # P08 — build_slide_6_conclusion_bridge
    (8, "TextBox 18", "width", 6.08),
    (8, "TextBox 18", "height", 0.59),
    (8, "TextBox 21", "top", 3.89),
    (8, "TextBox 21", "height", 0.88),

    # P09 — build_slide_7_cm_cn_proxy
    (9, "Rounded Rectangle 4", "height", 1.92),
    (9, "TextBox 28", "left", 1.02),
    (9, "TextBox 28", "top", 5.52),
    (9, "TextBox 28", "width", 11.66),
    (9, "TextBox 28", "height", 0.54),

    # P10 — build_slide_loh_af_overview
    (10, "Rounded Rectangle 6", "height", 1.62),
    (10, "TextBox 8", "left", 7.33),
    (10, "TextBox 8", "top", 4.52),
    (10, "TextBox 8", "width", 5.65),

    # P11 — build_slide_8_cn_tier_separation
    (11, "Rounded Rectangle 22", "top", 5.60),

    # P14 — build_slide_10_loh_basics
    (14, "Rounded Rectangle 29", "height", 0.65),
    (14, "Rounded Rectangle 56", "top", 5.68),
    (14, "TextBox 57", "top", 5.68),

    # P15 — build_slide_11_clonal_subclonal
    (15, "Rounded Rectangle 14", "height", 1.34),

    # P16 — build_slide_12_af_math
    (16, "Rounded Rectangle 6", "height", 2.89),
    (16, "Rounded Rectangle 11", "height", 2.89),
    (16, "Rounded Rectangle 29", "top", 5.26),
    (16, "Rounded Rectangle 29", "height", 1.36),
    (16, "Rounded Rectangle 30", "top", 5.32),
    (16, "TextBox 31", "top", 5.34),
    (16, "TextBox 32", "top", 5.64),

    # P18 — build_slide_14_cm_calibration
    (18, "Rounded Rectangle 8", "height", 1.77),
    (18, "Rounded Rectangle 12", "top", 4.22),
    (18, "Rounded Rectangle 12", "height", 1.60),
    (18, "Rounded Rectangle 13", "top", 4.27),
    (18, "TextBox 14", "top", 4.29),
    (18, "Table 15", "top", 4.59),
    (18, "Table 15", "height", 1.03),

    # P19 — build_slide_15_camdac
    (19, "Rounded Rectangle 8", "top", 4.65),
    (19, "TextBox 9", "top", 4.67),
    (19, "Rounded Rectangle 14", "top", 4.65),
    (19, "TextBox 15", "top", 4.67),
    (19, "Rounded Rectangle 16", "top", 5.30),
    (19, "Rounded Rectangle 16", "height", 1.38),
    (19, "Rounded Rectangle 17", "top", 5.36),
    (19, "TextBox 18", "top", 5.38),
    (19, "TextBox 19", "top", 5.68),
    (19, "Rounded Rectangle 20", "top", 5.30),
    (19, "Rounded Rectangle 20", "height", 1.38),
    (19, "TextBox 21", "top", 5.36),
    (19, "TextBox 22", "top", 5.70),

    # P20 — build_slide_16_allele_delta
    (20, "TextBox 17", "height", 1.94),
    (20, "Rounded Rectangle 18", "top", 4.93),
    (20, "Rounded Rectangle 19", "top", 4.98),
    (20, "TextBox 20", "top", 5.00),
    (20, "TextBox 21", "top", 5.28),
    (20, "TextBox 21", "height", 0.74),
    (20, "Rounded Rectangle 22", "top", 4.93),
    (20, "TextBox 23", "top", 4.97),
    (20, "TextBox 24", "top", 5.19),

    # P21 — build_slide_17_read_level
    (21, "Rounded Rectangle 7", "top", 2.47),
    (21, "TextBox 8", "top", 2.49),
    (21, "TextBox 9", "top", 2.47),
    (21, "TextBox 10", "top", 2.67),
    (21, "Rounded Rectangle 11", "top", 3.02),
    (21, "TextBox 12", "top", 3.04),
    (21, "TextBox 13", "top", 3.02),
    (21, "TextBox 14", "top", 3.22),
    (21, "Rounded Rectangle 15", "top", 3.57),
    (21, "TextBox 16", "top", 3.59),
    (21, "TextBox 17", "top", 3.57),
    (21, "TextBox 18", "top", 3.77),
    (21, "Rounded Rectangle 19", "top", 4.12),
    (21, "TextBox 20", "top", 4.14),
    (21, "TextBox 21", "top", 4.12),
    (21, "TextBox 22", "top", 4.32),
    (21, "Rounded Rectangle 23", "top", 4.67),
    (21, "TextBox 24", "top", 4.69),
    (21, "TextBox 25", "top", 4.67),
    (21, "TextBox 26", "top", 4.87),
    (21, "Rounded Rectangle 27", "top", 5.22),
    (21, "TextBox 28", "top", 5.24),
    (21, "TextBox 29", "top", 5.22),
    (21, "TextBox 30", "top", 5.42),
    (21, "TextBox 34", "height", 1.05),
    (21, "TextBox 38", "height", 1.05),
    (21, "TextBox 41", "left", 1.40),
    (21, "TextBox 41", "top", 6.10),
    (21, "TextBox 41", "width", 10.95),
    (21, "TextBox 41", "height", 0.54),

    # P22 — build_slide_18_ngroups
    (22, "Rounded Rectangle 4", "width", 1.01),
    (22, "TextBox 5", "top", 1.93),
    (22, "TextBox 5", "width", 0.93),
    (22, "TextBox 13", "top", 2.38),
    (22, "TextBox 14", "top", 2.38),
    (22, "TextBox 15", "left", 10.05),
    (22, "TextBox 15", "top", 2.35),
    (22, "TextBox 15", "width", 1.78),
    (22, "TextBox 15", "height", 0.37),
    (22, "TextBox 17", "top", 3.02),
    (22, "TextBox 18", "top", 3.02),
    (22, "TextBox 19", "left", 10.05),
    (22, "TextBox 19", "top", 3.01),
    (22, "TextBox 19", "width", 2.20),
    (22, "TextBox 19", "height", 0.37),
    (22, "TextBox 21", "top", 3.65),
    (22, "TextBox 22", "top", 3.65),
    (22, "TextBox 23", "left", 10.05),
    (22, "TextBox 23", "top", 3.62),
    (22, "TextBox 23", "width", 2.20),
    (22, "TextBox 23", "height", 0.37),
    (22, "Rounded Rectangle 25", "width", 1.01),
    (22, "TextBox 26", "top", 4.54),
    (22, "TextBox 26", "width", 0.93),
    (22, "TextBox 27", "left", 2.02),
    (22, "TextBox 27", "top", 4.46),
    (22, "TextBox 27", "width", 4.78),
    (22, "TextBox 27", "height", 0.80),
    (22, "TextBox 38", "left", 9.12),
    (22, "TextBox 38", "top", 4.50),
    (22, "TextBox 38", "width", 3.24),
    (22, "TextBox 38", "height", 0.57),

    # P23 — build_slide_19_nr_confound
    (23, "Rounded Rectangle 21", "height", 2.44),
    (23, "Table 22", "top", 3.96),
    (23, "Rounded Rectangle 23", "top", 5.51),
    (23, "Rounded Rectangle 23", "height", 0.61),
    (23, "TextBox 24", "top", 5.53),
    (23, "Rounded Rectangle 25", "height", 1.24),

    # P24 — build_slide_20_pyramid
    (24, "TextBox 9", "top", 4.83),
    (24, "TextBox 12", "top", 4.18),
    (24, "TextBox 15", "top", 3.56),
    (24, "TextBox 18", "top", 2.91),
    (24, "TextBox 21", "left", 3.58),
    (24, "TextBox 21", "top", 2.29),
    (24, "TextBox 21", "width", 1.67),
    (24, "TextBox 21", "height", 0.19),
    (24, "Rounded Rectangle 26", "left", 7.30),
    (24, "Rounded Rectangle 26", "top", 4.91),
    (24, "Rounded Rectangle 26", "width", 5.25),
    (24, "Rounded Rectangle 26", "height", 1.54),
    (24, "TextBox 27", "left", 7.43),
    (24, "TextBox 27", "top", 5.10),
    (24, "TextBox 27", "width", 2.96),
    (24, "TextBox 28", "left", 7.55),
    (24, "TextBox 28", "top", 5.40),
    (24, "TextBox 28", "width", 4.87),
    (24, "TextBox 28", "height", 0.91),

    # P26 — build_slide_21_af_baseline
    (26, "Picture 6", "width", 2.67),
    (26, "Picture 6", "height", 4.70),
    (26, "Rounded Rectangle 10", "left", 4.50),
    (26, "Rounded Rectangle 10", "width", 8.03),
    (26, "TextBox 11", "left", 4.56),
    (26, "TextBox 11", "top", 6.27),
    (26, "TextBox 11", "width", 7.74),
    (26, "TextBox 11", "height", 0.20),

    # P27 — build_slide_22_cn_tier_focus
    (27, "Table 8", "height", 1.03),
    (27, "Rounded Rectangle 9", "top", 5.06),
    (27, "TextBox 10", "top", 5.13),

    # P28 — build_slide_23_ngroups_af
    (28, "TextBox 13", "top", 5.27),

    # P29 — build_slide_24_nr_control
    (29, "Rounded Rectangle 4", "top", 0.69),
    (29, "TextBox 5", "top", 0.71),
    (29, "Picture 7", "top", 1.18),
    (29, "Table 10", "top", 4.97),
    (29, "Table 10", "height", 0.83),
    (29, "TextBox 11", "top", 5.90),
    (29, "Table 14", "top", 4.95),
    (29, "Table 14", "width", 2.77),
    (29, "Table 14", "height", 0.83),
    (29, "TextBox 15", "top", 5.81),

    # P30 — build_slide_25_segment
    (30, "TextBox 11", "left", 2.50),
    (30, "TextBox 11", "top", 4.94),
    (30, "TextBox 11", "width", 3.48),
    (30, "TextBox 11", "height", 1.09),
    (30, "TextBox 15", "left", 9.00),
    (30, "TextBox 15", "top", 4.93),
    (30, "TextBox 15", "width", 3.33),
    (30, "TextBox 15", "height", 1.09),

    # P31 — build_slide_26_per_sample
    (31, "Table 6", "top", 4.64),
    (31, "Table 6", "width", 7.41),
    (31, "Table 6", "height", 1.65),
    (31, "TextBox 9", "left", 9.45),
    (31, "TextBox 9", "top", 4.88),
    (31, "TextBox 9", "width", 2.90),
    (31, "TextBox 9", "height", 1.12),
    (31, "Rounded Rectangle 10", "top", 6.32),
    (31, "TextBox 11", "top", 6.34),

    # P33 — build_slide_27_evidence_summary
    (33, "TextBox 8", "height", 1.55),
    (33, "Rounded Rectangle 9", "height", 1.33),
    (33, "Rounded Rectangle 12", "height", 1.33),

    # P34 — build_slide_28_paired_validation
    (34, "TextBox 5", "top", 2.12),
    (34, "TextBox 6", "left", 1.41),
    (34, "TextBox 6", "top", 2.18),
    (34, "TextBox 6", "height", 0.27),
    (34, "TextBox 7", "top", 3.11),
    (34, "TextBox 7", "height", 0.92),
    (34, "TextBox 9", "top", 2.12),
    (34, "TextBox 10", "left", 5.39),
    (34, "TextBox 10", "top", 2.18),
    (34, "TextBox 10", "height", 0.27),
    (34, "TextBox 11", "top", 3.11),
    (34, "TextBox 11", "height", 0.92),
    (34, "TextBox 13", "top", 2.12),
    (34, "TextBox 14", "left", 9.37),
    (34, "TextBox 14", "top", 2.18),
    (34, "TextBox 14", "height", 0.27),
    (34, "TextBox 15", "top", 3.11),
    (34, "TextBox 15", "height", 0.92),

    # P35 — build_slide_29_normal_asm
    (35, "Rounded Rectangle 27", "height", 0.63),

    # P36 — build_slide_30_tumorlens_table
    (36, "TextBox 5", "top", 2.11),
    (36, "TextBox 6", "top", 2.16),
    (36, "TextBox 7", "top", 2.57),
    (36, "TextBox 7", "width", 4.95),
    (36, "TextBox 7", "height", 1.79),
    (36, "TextBox 9", "top", 2.12),
    (36, "Rounded Rectangle 10", "top", 2.52),
    (36, "TextBox 11", "top", 2.58),
    (36, "TextBox 12", "top", 2.85),
    (36, "Rounded Rectangle 13", "top", 2.52),
    (36, "TextBox 14", "top", 2.58),
    (36, "TextBox 15", "top", 2.85),
    (36, "TextBox 16", "top", 3.88),
    (36, "Rounded Rectangle 17", "height", 1.60),
    (36, "TextBox 20", "width", 11.25),
    (36, "TextBox 20", "height", 0.92),
    (36, "TextBox 21", "height", 0.15),

    # P37 — build_slide_31_positioning_diagram
    (37, "TextBox 8", "left", 7.21),
    (37, "TextBox 8", "top", 2.14),
    (37, "TextBox 8", "width", 3.31),
    (37, "TextBox 8", "height", 0.29),
    (37, "TextBox 9", "left", 7.21),
    (37, "TextBox 9", "width", 3.31),
    (37, "TextBox 9", "height", 2.37),

    # P38 — build_slide_32_ism_positioning
    (38, "Rounded Rectangle 4", "top", 2.01),
    (38, "TextBox 5", "top", 2.07),
    (38, "TextBox 6", "top", 2.46),
    (38, "Rounded Rectangle 7", "top", 4.04),
    (38, "TextBox 8", "top", 4.20),
    (38, "TextBox 9", "top", 4.59),
    (38, "TextBox 10", "top", 5.94),
    (38, "Rounded Rectangle 11", "height", 4.67),

    # P39 — build_slide_33_roadmap
    (39, "Rounded Rectangle 14", "height", 2.11),
    (39, "Rounded Rectangle 16", "height", 2.11),
    (39, "Rounded Rectangle 18", "height", 2.11),
    (39, "Rounded Rectangle 22", "height", 0.95),

    # P40 — build_slide_34_validation_risks
    (40, "Rounded Rectangle 4", "height", 4.43),
    (40, "Table 7", "height", 1.41),
    (40, "TextBox 8", "top", 3.76),
    (40, "TextBox 8", "width", 5.67),
    (40, "TextBox 8", "height", 1.92),
    (40, "Rounded Rectangle 10", "height", 4.43),

    # P41 — build_slide_35_qa
    (41, "TextBox 5", "left", 1.08),
    (41, "TextBox 6", "left", 1.49),
    (41, "TextBox 7", "left", 1.79),
    (41, "TextBox 7", "width", 4.62),
    (41, "TextBox 8", "left", 1.79),
    (41, "TextBox 8", "top", 2.65),
    (41, "TextBox 8", "width", 4.62),
    (41, "TextBox 8", "height", 0.18),
    (41, "TextBox 9", "left", 1.49),
    (41, "TextBox 10", "left", 1.79),
    (41, "TextBox 10", "width", 4.62),
    (41, "TextBox 11", "left", 1.79),
    (41, "TextBox 11", "top", 3.25),
    (41, "TextBox 11", "width", 4.62),
    (41, "TextBox 11", "height", 0.18),
    (41, "TextBox 12", "left", 1.49),
    (41, "TextBox 13", "left", 1.79),
    (41, "TextBox 13", "width", 4.62),
    (41, "TextBox 14", "left", 1.79),
    (41, "TextBox 14", "top", 3.85),
    (41, "TextBox 14", "width", 4.62),
    (41, "TextBox 14", "height", 0.18),
    (41, "TextBox 15", "left", 1.49),
    (41, "TextBox 16", "left", 1.79),
    (41, "TextBox 16", "width", 4.62),
    (41, "TextBox 17", "left", 1.79),
    (41, "TextBox 17", "top", 4.45),
    (41, "TextBox 17", "width", 4.62),
    (41, "TextBox 17", "height", 0.18),
    (41, "TextBox 18", "left", 1.49),
    (41, "TextBox 19", "left", 1.79),
    (41, "TextBox 19", "width", 4.62),
    (41, "TextBox 20", "left", 1.79),
    (41, "TextBox 20", "top", 5.05),
    (41, "TextBox 20", "width", 4.62),
    (41, "TextBox 20", "height", 0.18),
]

# Font size changes
FONT_CHANGES = [
    # P34 — build_slide_28_paired_validation
    (34, "TextBox 6", 14.0),
    (34, "TextBox 7", 14.0),
    (34, "TextBox 10", 14.0),
    (34, "TextBox 11", 14.0),
    (34, "TextBox 14", 14.0),
    (34, "TextBox 15", 14.0),
]


def apply_changes(src_pptx, dst_pptx):
    """Apply all recorded manual edits to create simulated version."""
    shutil.copy2(src_pptx, dst_pptx)
    prs = Presentation(dst_pptx)

    applied = 0
    not_found = 0

    # Build shape index: slide_idx -> {shape_name: shape}
    shape_index = {}
    for si, slide in enumerate(prs.slides):
        shape_map = {}
        for shape in slide.shapes:
            shape_map[shape.name] = shape
        shape_index[si] = shape_map

    # Apply position/size changes
    for (slide_1, shape_name, prop, target) in POSITION_CHANGES:
        si = slide_1 - 1
        if si not in shape_index:
            print(f"  [SKIP] Slide {slide_1} out of range")
            not_found += 1
            continue

        shape = shape_index[si].get(shape_name)
        if shape is None:
            print(f"  [NOT FOUND] Slide {slide_1}, shape '{shape_name}'")
            not_found += 1
            continue

        target_emu = int(target * EMU_PER_INCH)
        if prop == "left":
            shape.left = target_emu
        elif prop == "top":
            shape.top = target_emu
        elif prop == "width":
            shape.width = target_emu
        elif prop == "height":
            shape.height = target_emu

        applied += 1

    # Apply font changes
    for (slide_1, shape_name, target_pt) in FONT_CHANGES:
        si = slide_1 - 1
        shape = shape_index.get(si, {}).get(shape_name)
        if shape is None:
            print(f"  [NOT FOUND] Slide {slide_1}, shape '{shape_name}' (font)")
            not_found += 1
            continue

        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(target_pt)
        except Exception:
            pass
        applied += 1

    prs.save(dst_pptx)
    print(f"\n[OK] Simulated manual edit version saved: {dst_pptx}")
    print(f"     Applied: {applied} changes")
    print(f"     Not found: {not_found} shapes")
    return applied, not_found


if __name__ == "__main__":
    src = "20260414_LOH_Subclone_AF_Methylation.pptx"
    dst = "20260414_simulated_manual_edit.pptx"
    apply_changes(src, dst)
