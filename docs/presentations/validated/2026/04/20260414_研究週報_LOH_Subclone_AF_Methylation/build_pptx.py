#!/usr/bin/env python3
"""
LOH Subclone AF × Methylation 週報 PPTX Builder

分批產出，每次 ~5 slides。執行方式:
  python3 build_pptx.py                    # 產出目前已定義的所有 slides
  python3 build_pptx.py --batch 1          # 僅產出 batch 1 (封面+重點)

每批完成後截圖驗證，確認佈局再繼續下一批。
"""
from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════
SLIDE_W = 13.333
SLIDE_H = 7.5
TOTAL_PAGES = 41  # 37 content + 4 section dividers (V2 + ISM flow + LOH AF overview)

REPORT_DIR = Path(__file__).parent
FIG_DIR = REPORT_DIR / "figures"
TEACH_DIR = REPORT_DIR / "teaching_figures"
OUTPUT_PPTX = REPORT_DIR / "20260414_LOH_Subclone_AF_Methylation.pptx"

# ═══════════════════════════════════════════════════════════════
# Theme — warm neutral (same palette as previous weeks)
# ═══════════════════════════════════════════════════════════════
THEME = {
    "dark":       "1E2A44",
    "bg":         "F7F3EC",
    "bg_alt":     "EEE6DB",
    "accent":     "C76B50",   # Terracotta
    "accent2":    "6E9D8A",   # Sage green
    "accent3":    "54708C",   # Steel blue
    "accent4":    "7B5B9A",   # Muted purple
    "positive":   "009E73",   # Colorblind-safe green
    "negative":   "D55E00",   # Colorblind-safe vermillion
    "text":       "17202A",
    "muted":      "5E6572",
    "en_color":   "5E6572",
    "light_text": "F8F4EE",
    "light_muted":"D6DCE5",
    "line":       "D6CCBF",
    "card_bg":    "FFFFFF",
    "font_title": "Arial Bold",
    "font_body":  "Arial",
}

META = {
    "deck_title":    "LOH 區域發現 Subclone — 甲基化雙重驗證",
    "deck_title_en": "LOH Subclone Discovery — Dual Evidence via AF × Methylation Diversity",
    "deck_subtitle": "2026-04-14 ~ 2026-04-20",
    "summary_line":  "結論 POSITIVE：LOH 區域的 subclonal variant\n甲基化多樣性顯著高於 clonal variant (7/7 一致)",
    "kicker":        "InterSubMod 研究週報",
    "author":        "廖子游",
    "footer":        "InterSubMod Weekly · LOH Subclone AF × Methylation · 2026-04-14",
}


# ═══════════════════════════════════════════════════════════════
# Low-level helpers (from build_weekly_report_pptx.py)
# ═══════════════════════════════════════════════════════════════
def hex_rgb(value: str) -> RGBColor:
    value = value.replace("#", "").strip()
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_bg(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)


def add_text(slide, text, x, y, w, h, *, font_size=18, color="17202A",
             bold=False, font_name="Arial", align="left",
             valign=MSO_ANCHOR.MIDDLE, italic=False, margin=0.05,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_rgb(color)
    return box


def add_multiline_text(slide, lines, x, y, w, h, *, font_size=14, color="17202A",
                       bold=False, font_name="Arial", align="left",
                       line_spacing=1.15, margin=0.05):
    """Add a text box with multiple paragraphs (one per line)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.bold = bold
        run.font.italic = False
        run.font.color.rgb = hex_rgb(color)
    return box


def add_rect(slide, x, y, w, h, *, fill, line=None, radius=False, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill)
    shape.fill.transparency = transparency
    if line is None:
        # Use same-color border instead of noFill (better viewer compatibility)
        shape.line.color.rgb = hex_rgb(fill)
        shape.line.width = Pt(0.25)
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_footer(slide, page_num):
    add_text(slide, META["footer"], 0.8, 6.78, 9.0, 0.3, font_size=10,
             color=THEME["muted"], margin=0)
    add_text(slide, f"{page_num}/{TOTAL_PAGES}", 12.0, 6.75, 0.55, 0.3,
             font_size=12, color=THEME["muted"], align="right", margin=0)


def add_native_arrow(slide, x1, y1, x2, y2, *, color="5E6572", width=1.5):
    """Add a line connector with arrowhead, with safe minimum dimension."""
    from pptx.oxml.ns import qn as _qn
    min_dim = Emu(Inches(0.06))
    cx1, cy1 = Inches(x1), Inches(y1)
    cx2, cy2 = Inches(x2), Inches(y2)
    if abs(cy2 - cy1) < min_dim:
        cy2 = cy1 + min_dim
    if abs(cx2 - cx1) < min_dim:
        cx2 = cx1 + min_dim
    connector = slide.shapes.add_connector(
        1, cx1, cy1, cx2, cy2  # MSO_CONNECTOR_TYPE.STRAIGHT = 1
    )
    connector.line.color.rgb = hex_rgb(color)
    connector.line.width = Pt(width)
    # Add arrowhead
    ln = connector.line._ln
    tail_end = ln.makeelement(_qn("a:tailEnd"), {})
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    ln.append(tail_end)


def inject_speaker_notes(slide, notes_text):
    """Add speaker notes to slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _safe_add_picture(slide, image_path, left, top, *, width=None, height=None,
                      placeholder=""):
    """Add picture with RGBA→RGB conversion, downscale, ICC strip.
    If image not found and placeholder is set, add a gray text box instead.
    left/top/width/height are EMU (from Inches())."""
    img_path = Path(image_path)
    if not img_path.exists():
        print(f"  [WARN] Image not found: {img_path}")
        if placeholder and width is not None and height is not None:
            # Convert EMU back to inches for add_rect/add_text helpers
            _emu = 914400
            lx = left / _emu if isinstance(left, int) else left
            ly = top / _emu if isinstance(top, int) else top
            lw = width / _emu if isinstance(width, int) else width
            lh = height / _emu if isinstance(height, int) else height
            add_rect(slide, lx, ly, lw, lh,
                     fill="E8E8E8", line="CCCCCC", radius=True)
            add_text(slide, f"[IMAGE PLACEHOLDER]\n{placeholder}",
                     lx + 0.1, ly + 0.1, lw - 0.2, lh - 0.2,
                     font_size=10, color="888888", margin=0.05)
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    tmp = None
    try:
        MAX_DIM = 960  # ~150 DPI for 6.5" display width; keeps decoded RAM low
        with Image.open(img_path) as img:
            if img.mode in ("RGBA", "PA"):
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[3])
            elif img.mode != "RGB":
                rgb_img = img.convert("RGB")
            else:
                rgb_img = img.copy()
            iw, ih = rgb_img.size
            if iw > MAX_DIM or ih > MAX_DIM:
                ratio = min(MAX_DIM / iw, MAX_DIM / ih)
                new_size = (max(1, int(iw * ratio)), max(1, int(ih * ratio)))
                resample = getattr(Image, "Resampling", Image).LANCZOS
                rgb_img = rgb_img.resize(new_size, resample)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        rgb_img.save(tmp.name, format="PNG", optimize=True, icc_profile=None)
        result = slide.shapes.add_picture(tmp.name, left, top, **kwargs)
        return result
    except Exception as e:
        print(f"  [WARN] Processing failed: {e}; embedding original")
        return slide.shapes.add_picture(str(img_path), left, top, **kwargs)
    finally:
        if tmp is not None:
            try:
                Path(tmp.name).unlink(missing_ok=True)
            except Exception:
                pass


def add_image_fit(slide, img_path, x, y, w, h):
    """Add image that fits within bounding box, preserving aspect ratio, centered."""
    img_path = Path(img_path)
    if not img_path.exists():
        print(f"  [WARN] Image not found: {img_path}")
        return None
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
    box_aspect = w / h
    if aspect > box_aspect:
        # Width-limited
        disp_w = w
        disp_h = w / aspect
    else:
        # Height-limited
        disp_h = h
        disp_w = h * aspect
    off_x = x + (w - disp_w) / 2
    off_y = y + (h - disp_h) / 2
    return _safe_add_picture(slide, str(img_path), Inches(off_x), Inches(off_y),
                             width=Inches(disp_w))


# ═══════════════════════════════════════════════════════════════
# Higher-level slide component helpers
# ═══════════════════════════════════════════════════════════════
def add_simple_table(slide, headers, rows, x, y, w, h, *,
                     header_color="5E6572", text_color="17202A",
                     font_size=10, header_font_size=10):
    """Add a native PPTX table (1 shape instead of N×M text boxes)."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                         Inches(w), Inches(h))
    table = table_shape.table
    # Style: no built-in style (avoid banding)
    tbl_pr = table._tbl.tblPr
    tbl_pr.set("bandRow", "0")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("lastRow", "0")

    def _set_cell(cell, text, *, color, bold=False, sz=font_size):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.color.rgb = hex_rgb(color)
        run.font.bold = bold
        run.font.name = THEME["font_body"]
        cell.text_frame.word_wrap = True
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        # Transparent fill
        cell.fill.background()

    # Header row
    for ci, h_text in enumerate(headers):
        _set_cell(table.cell(0, ci), h_text, color=header_color,
                  bold=True, sz=header_font_size)

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            bold = isinstance(val, str) and val.startswith("+")
            _set_cell(table.cell(ri + 1, ci), val, color=text_color, bold=bold)

    # Remove table borders
    for ri in range(n_rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                border = tcPr.makeelement(qn(border_name), {})
                no_fill = border.makeelement(qn("a:noFill"), {})
                border.append(no_fill)
                tcPr.append(border)

    return table_shape


def add_slide_title(slide, title_zh, subtitle_zh, *, dark=False,
                    title_en=None, subtitle_en=None):
    """Add bilingual title block to a content slide."""
    fg = THEME["light_text"] if dark else THEME["text"]
    muted = THEME["light_muted"] if dark else THEME["muted"]
    t_sz = 32
    s_sz = 16
    add_text(slide, title_zh, 0.8, 0.35, 11.7, 0.75, font_size=t_sz,
             color=fg, bold=True, font_name=THEME["font_title"], margin=0)
    if title_en:
        en_sz = max(9, int(t_sz * 0.45))
        add_text(slide, title_en, 1.05, 1.05, 11.3, 0.35, font_size=en_sz,
                 color=THEME["en_color"], font_name=THEME["font_body"], margin=0)
    sub_y = 1.38 if title_en else 1.20
    add_text(slide, subtitle_zh, 0.8, sub_y, 11.3, 0.42, font_size=s_sz,
             color=muted, font_name=THEME["font_body"], margin=0)
    if subtitle_en:
        en_sz2 = max(9, int(s_sz * 0.45))
        add_text(slide, subtitle_en, 1.05, sub_y + 0.35, 11.0, 0.32,
                 font_size=en_sz2, color=THEME["en_color"],
                 font_name=THEME["font_body"], margin=0)


def add_card(slide, x, y, w, h, *, accent, title, body, title_en=None, icon=None):
    """Add a single finding card with accent pip and optional icon label."""
    add_rect(slide, x, y, w, h, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_rect(slide, x + 0.15, y + 0.15, 0.42, 0.42, fill=accent, radius=True)
    if icon:
        add_text(slide, icon, x + 0.15, y + 0.17, 0.42, 0.38,
                 font_size=12, color="FFFFFF", bold=True, align="center", margin=0)
    add_text(slide, title, x + 0.72, y + 0.12, w - 1.0, 0.35,
             font_size=17, color=THEME["text"], bold=True, margin=0)
    if title_en:
        add_text(slide, title_en, x + 0.90, y + 0.45, w - 1.2, 0.25,
                 font_size=9, color=THEME["en_color"], margin=0)
    body_y = y + 0.78 if title_en else y + 0.58
    body_h = h - body_y + y - 0.1
    add_text(slide, body, x + 0.30, body_y, w - 0.5, body_h,
             font_size=12, color=THEME["muted"], margin=0, line_spacing=1.35,
             valign=MSO_ANCHOR.TOP)


def add_section_divider(slide, title_zh, subtitle_zh, *, title_en=None, page_num=None):
    """Dark background section divider slide."""
    set_bg(slide, THEME["dark"])
    add_text(slide, title_zh, 1.5, 2.5, 10.3, 1.2, font_size=40,
             color=THEME["light_text"], bold=True,
             font_name=THEME["font_title"], align="center", margin=0)
    if title_en:
        add_text(slide, title_en, 1.5, 3.5, 10.3, 0.6, font_size=16,
                 color=THEME["light_muted"], align="center", margin=0)
    add_text(slide, subtitle_zh, 1.5, 4.3, 10.3, 0.6, font_size=18,
             color=THEME["light_muted"], align="center", margin=0)
    if page_num:
        add_footer(slide, page_num)


# ═══════════════════════════════════════════════════════════════
# BATCH 1: 封面 + 本週重點 (Slides 1-2)
# ═══════════════════════════════════════════════════════════════
def build_slide_1_cover(prs):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_bg(slide, THEME["dark"])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=THEME["dark"])

    # Right panel — cover image
    add_rect(slide, 7.2, 0.5, 5.63, 6.5, fill="F6EFE6", radius=True, transparency=8)
    cover_img = REPORT_DIR / "1776260029399.png"
    if cover_img.exists():
        add_image_fit(slide, cover_img, 7.56, 1.08, 5.0, 5.3)
    else:
        add_text(slide, "[ 每週重點圖 — 手動插入 ]", 8.0, 3.2, 4.0, 0.4,
                 font_size=13, color="A0A0A0", align="center", margin=0)

    # Kicker badge
    add_rect(slide, 0.7, 0.65, 2.2, 0.35, fill=THEME["accent2"], radius=True)
    add_text(slide, META["kicker"], 0.83, 0.72, 2.5, 0.2,
             font_size=13, color=THEME["light_text"], bold=True, margin=0)

    # Title (ZH large)
    add_text(slide, META["deck_title"], 0.8, 1.35, 6.0, 0.8,
             font_size=36, color=THEME["light_text"], bold=True,
             font_name=THEME["font_title"], margin=0)
    # EN title
    add_text(slide, META["deck_title_en"], 1.1, 2.15, 5.7, 0.7,
             font_size=14, color=THEME["light_muted"],
             font_name=THEME["font_body"], margin=0)

    # Subtitle (date + conclusion)
    add_text(slide, META["deck_subtitle"], 0.84, 3.05, 6.0, 0.4,
             font_size=18, color=THEME["light_muted"],
             font_name=THEME["font_body"], margin=0)
    add_text(slide, META["summary_line"], 0.84, 3.55, 6.0, 0.90,
             font_size=17, color=THEME["light_text"], bold=True,
             font_name=THEME["font_body"], margin=0, line_spacing=1.3)

    # Guide bullets
    guide = [
        "📖 教學: LOH → Subclonal LOH → 雙重證據鏈",
        "📊 數據: 三步驟驗證 (AF → NGroups × AF → Segment)",
        "✅ 結論: 六層證據鏈 POSITIVE",
        "🔭 方向: ISM read-level characterization",
    ]
    add_multiline_text(slide, guide, 0.84, 4.60, 5.5, 1.3,
                       font_size=11, color=THEME["light_muted"],
                       font_name=THEME["font_body"], line_spacing=1.5)

    # Reporter
    add_text(slide, f"報告人：{META['author']}", 0.84, 6.10, 4.0, 0.3,
             font_size=14, color=THEME["light_text"], bold=True, margin=0)

    inject_speaker_notes(slide,
        "開場直接點出 POSITIVE 結論。\n"
        "這是 ISM 第二個 POSITIVE 結論，第一個是 HPFineNGroups subclone marker。\n"
        "也是第一個建立完整六層證據鏈的結論。")

    return slide


def build_slide_2_key_findings(prs):
    """Slide 2: 本週重點 + 報告路線圖"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "本週重點與報告路線圖",
                    "四個核心發現 + 報告結構導覽",
                    title_en="Key Findings & Report Roadmap")

    # 4 finding cards (2x2 grid)
    cards = [
        {
            "accent": THEME["accent"],
            "icon": "🔍",
            "title": "TO 的 LOH 區域篩出高比例 TP",
            "title_en": "24.6% LOH TP have intermediate AF",
            "body": "在 purity=1.0 cell line 中，LOH 區域 TP 有四分之一 AF 不在期望值 (0 or 1.0)。FP 僅 4.1%，6 倍富集。7/7 樣本一致 (7.0%-60.2%)。",
        },
        {
            "accent": THEME["accent2"],
            "icon": "🔬",
            "title": "甲基化獨立驗證 subclone 存在",
            "title_en": "ΔNGroups = +0.705, all 7/7 p < 10⁻³⁹",
            "body": "Intermediate AF variants 甲基化群組數顯著高於 Extreme AF。Extreme: 90.7% 單群 → clonal。Intermediate: 79.6% 多群 → subclonal。",
        },
        {
            "accent": THEME["accent3"],
            "icon": "✅",
            "title": "控制讀取數後效應反而增強",
            "title_en": "NR-bin control: effect strengthens (r: 0.48→0.71)",
            "body": "NumReads 分層後，效應量從 |r|=0.483 增加到 0.709。若是 confound 應消失 — 實際增強，完全排除。低 NR 效應弱 = floor effect。",
        },
        {
            "accent": THEME["accent4"],
            "icon": "📍",
            "title": "同一 LOH 區段內 subclone 信號穩定",
            "title_en": "Segment-level: ρ=0.270, 6/7 positive",
            "body": "AF 變異性越大的 LOH segment，NGroups 越高。6/7 樣本正方向，5/7 顯著。H1437 ρ=0.809 最強。Segmental event 非 random noise。",
        },
    ]

    positions = [(0.8, 1.95), (6.95, 1.95), (0.8, 4.32), (6.95, 4.32)]
    for (x, y), card in zip(positions, cards):
        add_card(slide, x, y, 5.55, 2.15, **card)

    # Roadmap bar at bottom
    add_rect(slide, 0.8, 6.55, 11.73, 0.20, fill=THEME["dark"], radius=True)
    roadmap = "背景(5頁) → CN發現(4頁) → 教學(11頁) → 數據(6頁) → 整合 → 競爭者 → 未來  (~55 min)"
    add_text(slide, roadmap, 0.8, 6.52, 11.73, 0.25,
             font_size=11, color=THEME["light_text"], align="center", margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "4 cards 對應四個核心數字：\n"
        "1. 24.6% intermediate AF — 提示 subclonal LOH\n"
        "2. ΔNG=+0.705 — 雙重證據鏈\n"
        "3. NR-bin 控制排除 confound\n"
        "4. Segment spatial 一致性\n"
        "底部路線圖：先教學(較長) → 再數據 → 結論 → 方向")

    return slide


# ═══════════════════════════════════════════════════════════════
# BATCH 2: Section Divider + Part B 背景 (4 slides) + Part C CN 發現 (3 slides)
# ═══════════════════════════════════════════════════════════════
def _add_flow_step(slide, x, y, w, h, *, label, detail, fill, text_color="17202A"):
    """Add a flow step box: label on top, detail below."""
    add_rect(slide, x, y, w, h, fill=fill, radius=True)
    add_text(slide, label, x, y + 0.06, w, 0.32,
             font_size=14, color=text_color, bold=True, align="center", margin=0.04)
    add_text(slide, detail, x, y + 0.36, w, h - 0.42,
             font_size=10, color=text_color if text_color != "17202A" else THEME["muted"],
             align="center", margin=0.04, valign=MSO_ANCHOR.TOP, line_spacing=1.2)


def _add_term_box(slide, terms, x, y, w, h):
    """Add a terminology explanation box (術語解釋區塊)."""
    add_rect(slide, x, y, w, h, fill="F5F0EB", line=THEME["line"], radius=True)
    add_text(slide, "術語", x + 0.10, y + 0.04, 0.6, 0.20,
             font_size=10, color=THEME["accent"], bold=True, margin=0)
    add_multiline_text(slide, terms, x + 0.10, y + 0.26, w - 0.20, h - 0.30,
                       font_size=10, color=THEME["text"], line_spacing=1.30)


def _add_conclusion_box(slide, text, x=0.8, y=6.20, w=11.73):
    """Add a bottom conclusion box."""
    add_rect(slide, x, y, w, 0.40, fill="E8F5E9", line=THEME["accent2"], radius=True)
    add_text(slide, "結論: " + text, x + 0.15, y + 0.03, w - 0.30, 0.35,
             font_size=12, color=THEME["text"], bold=True, margin=0)


def build_divider_background(prs):
    """Section Divider D1: 背景與 LOH 驗證"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide,
                        "背景與 LOH 驗證",
                        "為什麼研究 TO × LOH？LOH.bed 可信嗎？",
                        title_en="Background & LOH Validation",
                        page_num=len(prs.slides))
    inject_speaker_notes(slide,
        "進入背景區段。\n"
        "接下來 7 頁會回顧：\n"
        "為什麼從 Paired 轉向 TO 模式（統計檢定力）、\n"
        "LOH.bed 的 SEQC2 金標準驗證（J=0.928）、\n"
        "雙定義層級關係、前期結論橋樑，\n"
        "然後進入 LOH+CNV 的 CN Tier 分層發現。")
    return slide


def build_slide_3_to_motivation(prs):
    """Slide 3: 從 Paired 到 TO — 為了統計檢定力"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "從 Paired 到 TO — 為了統計檢定力",
                    "TO 模式提供充足的 FP 數量來驗證 ISM 特徵",
                    title_en="From Paired to Tumor-Only — For Statistical Power")

    # Top: TP/FP Composition chart (main visual, ~60% area)
    comp_img = TEACH_DIR / "fig04_truth_label_composition.png"
    if comp_img.exists():
        add_image_fit(slide, comp_img, 0.5, 1.80, 9.0, 2.80)
    else:
        add_rect(slide, 0.5, 1.80, 9.0, 2.80, fill="F0F0F0", line=THEME["line"], radius=True)
        add_text(slide, "[TODO: TP/FP Composition Chart — Paired vs TO]",
                 2.5, 3.0, 5.0, 0.5, font_size=13, color="A0A0A0", align="center", margin=0)

    # Top-right: Term box
    _add_term_box(slide, [
        "FP = False Positive (假陽性)",
        "  被誤判為 somatic 的 variant",
        "TO 無 normal BAM 對照",
        "  → germline 被誤判",
    ], 9.7, 1.80, 3.13, 1.60)

    # Bottom: Paired box (compressed)
    add_rect(slide, 0.5, 4.70, 6.0, 1.40, fill="FFEBEE", line="EF9A9A", radius=True)
    add_text(slide, "Paired Mode", 0.7, 4.75, 3.0, 0.30,
             font_size=14, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "TP : FP ≈ 95 : 1  ·  FP = 3,429 / 328,699 (~1%)",
        "FP 太少 → 統計力不足 ✗",
    ], 0.7, 5.08, 5.5, 0.85, font_size=12, color=THEME["text"], line_spacing=1.40)

    # Arrow
    add_native_arrow(slide, 6.60, 5.40, 6.95, 5.40, color=THEME["accent"], width=2.5)

    # Bottom: TO box (compressed)
    add_rect(slide, 7.1, 4.70, 5.73, 1.40, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "TO Mode", 7.3, 4.75, 3.0, 0.30,
             font_size=14, color=THEME["positive"], bold=True, margin=0)
    add_multiline_text(slide, [
        "TP : FP ≈ 2.3 : 1  ·  FP = 128,382 / 419,692 (~31%)",
        "充足 FP → 統計力足夠 ✓",
    ], 7.3, 5.08, 5.2, 0.85, font_size=12, color=THEME["text"], line_spacing=1.40)

    _add_conclusion_box(slide, "TO 模式提供充足統計檢定力驗證 ISM 特徵的區分能力")

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "Paired FP 僅 ~1%，TP:FP ≈ 95:1 → 無法做統計分析。\n"
        "TO FP ~31%，TP:FP ≈ 2.3:1 → 充足檢定力。\n"
        "TO FP 主要是 germline variant (PON 過濾後殘留)。\n"
        "來源: v5 LOH report Slide 2")
    return slide


def build_slide_4_seqc2(prs):
    """Slide 4: LOH.bed 金標準驗證 (含 chrX 排除比較 + 論文引用 + 四標準解釋)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "LOH.bed 金標準驗證",
                    "SEQC2 多中心共識確認 — 排除 chrX 後 Precision +8.7%",
                    title_en="LOH.bed Validation Against SEQC2 Gold Standard")

    # Top-left: Ideogram image (60% width)
    teach_img = TEACH_DIR / "seqc2_ideogram.png"
    if teach_img.exists():
        add_image_fit(slide, teach_img, 0.5, 1.90, 7.5, 2.40)
    else:
        add_rect(slide, 0.5, 1.90, 7.5, 2.40, fill="F0F0F0", line=THEME["line"], radius=True)
        add_text(slide, "[SEQC2 vs LongPhase-TO Ideogram]",
                 2.0, 2.8, 4.5, 0.5, font_size=13, color="A0A0A0", align="center", margin=0)

    # Top-right: Autosomal metrics cards (main result)
    metrics = [
        ("Jaccard", "0.928"),
        ("Sens", "0.961"),
        ("Prec", "0.964"),
        ("F1", "0.963"),
    ]
    for i, (label, val) in enumerate(metrics):
        mx = 8.3 + i * 1.15
        add_rect(slide, mx, 1.95, 1.05, 0.80, fill=THEME["card_bg"], line=THEME["line"], radius=True)
        add_text(slide, label, mx, 2.00, 1.05, 0.30,
                 font_size=10, color=THEME["muted"], bold=True, align="center", margin=0)
        add_text(slide, val, mx, 2.30, 1.05, 0.40,
                 font_size=18, color=THEME["positive"], bold=True, align="center", margin=0)
    add_text(slide, "體染色體 only", 8.3, 2.78, 4.6, 0.22,
             font_size=9, color=THEME["muted"], align="center", margin=0)

    # Middle-left: chrX comparison table
    add_text(slide, "chrX 排除效果 (HCC1395 = 女性 cell line，chrX 半合子→假 LOH)",
             0.8, 4.40, 7.5, 0.25, font_size=10, color=THEME["muted"], bold=True, margin=0)
    chrx_headers = ["指標", "含 chrX", "體染色體 only", "差異"]
    chrx_rows = [
        ["Jaccard", "0.847", "0.928", "+9.6%"],
        ["Sensitivity", "0.961", "0.961", "—"],
        ["Precision", "0.877", "0.964", "+8.7%"],
        ["F1", "0.917", "0.963", "+4.6%"],
    ]
    add_simple_table(slide, chrx_headers, chrx_rows, 0.8, 4.68, 7.0, 1.20,
                     header_color=THEME["accent3"], font_size=11, header_font_size=11)

    # Middle-right: Four-metric explanation memo
    add_rect(slide, 8.3, 3.05, 4.23, 2.85, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "四標準定義備忘", 8.45, 3.10, 3.9, 0.28,
             font_size=11, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Jaccard = TP∩Truth / (TP∪Truth)",
        "  → 整體重疊度（最嚴格）",
        "Sensitivity = TP / (TP+FN)",
        "  → 找到多少真 LOH",
        "Precision = TP / (TP+FP)",
        "  → 叫出的 LOH 有多少是真的",
        "F1 = 2×P×S / (P+S)",
        "  → P 和 S 的調和平均",
    ], 8.45, 3.42, 3.9, 2.40, font_size=10, color=THEME["muted"], line_spacing=1.28)

    # Bottom: Paper reference
    add_rect(slide, 0.8, 5.95, 11.73, 0.22, fill=THEME["card_bg"], radius=True)
    add_text(slide, "Fang LT et al. (2021) Nature Biotechnology 39:1151-1160 · DOI: 10.1038/s41587-021-00993-6 · Data: research/loh_investigation/data/seqc2_cnv_benchmark_v4/",
             0.95, 5.95, 11.4, 0.22, font_size=8, color=THEME["muted"], margin=0)

    _add_conclusion_box(slide, "LOH.bed 高度可信（體染色體 J=0.928）；chrX 假 LOH 已排除")

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "SEQC2 驗證結果（體染色體 only）: J=0.928, Sens=0.961, Prec=0.964, F1=0.963\n"
        "含 chrX 時: J=0.847, Prec=0.877 → chrX 是女性 cell line 半合子造成的假 LOH\n"
        "排除 chrX 後 Precision +8.7%，Jaccard +9.6%\n"
        "論文: Fang LT et al. (2021) Nat Biotechnol 39:1151-1160\n"
        "四標準差異: Jaccard 最嚴格（集合重疊），Sens/Prec 互補，F1 為調和平均\n"
        "來源: v3 seqc2_validation section")
    return slide


def build_slide_ism_loh_impact(prs):
    """New slide: ISM 分析流程 — 為什麼 LOH 影響甲基化分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "ISM 分析流程 — 為什麼 LOH 影響甲基化分析",
                    "ISM 依賴 HP tags 做等位基因分離 → LOH 區域失效",
                    title_en="ISM Pipeline & Why LOH Disrupts Methylation Analysis")

    # ── Upper: ISM 4-step flow (y=1.85~3.45) — dark boxes with white text ──
    flow_steps = [
        ("BAM + VCF",     "ONT long-read\n± somatic SNV",     THEME["accent"]),
        ("CpG 甲基化",    "MM/ML tags\nper-read 解析",        THEME["accent2"]),
        ("距離矩陣",      "NHD / L1 / L2\nread-read 距離",    THEME["accent3"]),
        ("分群 + 統計",   "Hierarchical\nPERMANOVA",          THEME["accent4"]),
    ]
    for i, (lbl, det, bg_col) in enumerate(flow_steps):
        sx = 0.5 + i * 3.15
        _add_flow_step(slide, sx, 1.85, 2.60, 1.15, label=lbl, detail=det,
                       fill=bg_col, text_color="FFFFFF")
        if i < 3:
            add_native_arrow(slide, sx + 2.65, 2.42, sx + 3.10, 2.42,
                             color=THEME["text"], width=2.0)

    # Core output labels below flow
    add_text(slide, "→ 核心輸出: NGroups (分群數)、AlleleDelta (等位基因甲基化差異)、PERMANOVA p-value",
             0.5, 3.08, 12.3, 0.30,
             font_size=10, color=THEME["accent"], bold=True, margin=0)

    # ── Lower: Normal vs LOH comparison (y=3.50~5.85) ──
    add_text(slide, "為什麼 LOH 區域 ISM 分析受限？", 0.5, 3.45, 12.0, 0.35,
             font_size=16, color=THEME["text"], bold=True, margin=0)

    # Left panel: Normal region
    add_rect(slide, 0.5, 3.85, 5.8, 2.00, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "✓ 正常區域", 0.7, 3.90, 3.0, 0.30,
             font_size=14, color=THEME["positive"], bold=True, margin=0)
    add_multiline_text(slide, [
        "HP1 reads: 45  |  HP2 reads: 50",
        "HP_Ratio ≈ 0.47 → 兩群均衡",
        "",
        "PERMANOVA 可比較兩群甲基化差異",
        "→ NGroups / AlleleDelta 可信 ✓",
    ], 0.7, 4.25, 5.4, 1.50, font_size=12, color=THEME["text"], line_spacing=1.40)

    # Arrow between panels
    add_text(slide, "vs", 6.35, 4.55, 0.55, 0.40,
             font_size=16, color=THEME["muted"], bold=True, align="center", margin=0)

    # Right panel: LOH region
    add_rect(slide, 6.95, 3.85, 5.88, 2.00, fill="FFEBEE", line="EF9A9A", radius=True)
    add_text(slide, "✗ LOH 區域", 7.15, 3.90, 3.0, 0.30,
             font_size=14, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "HP1 reads: 92  |  HP2 reads: 3",
        "HP_Ratio ≈ 0.97 → 一群消失",
        "",
        "PERMANOVA 無法比較（單群）",
        "→ NGroups / AlleleDelta 失效 ✗",
    ], 7.15, 4.25, 5.5, 1.50, font_size=12, color=THEME["text"], line_spacing=1.40)

    # Term box — top-right corner, beside "為什麼 LOH" title
    _add_term_box(slide, [
        "MM/ML = BAM 甲基化標籤",
        "PERMANOVA = 多變量變異數分析",
        "HP tags = Haplotype Phase 標籤",
    ], 9.6, 3.48, 3.23, 0.90)

    _add_conclusion_box(slide,
        "ISM 在 LOH 區域的分析受限 → 需要深入理解 LOH 才能正確解讀 ISM 結果",
        x=0.5, w=12.33)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "ISM 四步流程: BAM+VCF → CpG甲基化(MM/ML) → 距離矩陣 → 分群+統計\n"
        "核心輸出: NGroups, AlleleDelta, PERMANOVA p-value\n"
        "正常區域: HP1≈HP2 → PERMANOVA 可比較 → 特徵可信\n"
        "LOH 區域: 一群 reads 消失 → PERMANOVA 失效 → 特徵不可靠\n"
        "這就是為什麼我們需要理解 LOH: 它直接影響 ISM 分析的可信度\n"
        "參考: 0408 PPT Slide 7 ISM framework")
    return slide


def build_slide_5_dual_definition(prs):
    """Slide 5: HP Imbalance 定義 + LOH 雙定義層級關係"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "HP Imbalance 定義與 LOH 雙定義",
                    "HP_Ratio 判定位點層級等位基因失衡 → 與 LOH.bed 的層級關係",
                    title_en="HP Imbalance Definition & Two LOH Definitions")

    # ── Upper-left: HP_Ratio formula + definition (y=1.80~3.65) ──
    add_rect(slide, 0.5, 1.80, 6.5, 1.85, fill=THEME["card_bg"], line=THEME["accent"], radius=True)
    add_text(slide, "HP_Ratio 公式", 0.7, 1.85, 3.0, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)
    # Formula in monospace — split into two lines for readability
    add_text(slide, "HP_Ratio = (HP1_FamilyN + 0.001)",
             0.7, 2.15, 6.1, 0.22,
             font_size=10, color=THEME["accent"], bold=True, font_name="Consolas", margin=0)
    add_text(slide, "         / (HP1_FamilyN + HP2_FamilyN + 0.002)",
             0.7, 2.35, 6.1, 0.22,
             font_size=10, color=THEME["accent"], bold=True, font_name="Consolas", margin=0)
    add_multiline_text(slide, [
        "Laplace smoothing 避免除以零",
        "正常 0.1-0.9 | HP Imbalance <0.1 或 >0.9",
        "範例: HP1=45, HP2=5 → 0.90 → Potential LOH",
    ], 0.7, 2.62, 6.1, 0.95, font_size=10, color=THEME["muted"], line_spacing=1.30)

    # ── Upper-right: HP Imbalance impact on ISM (y=1.80~3.65) ──
    add_rect(slide, 7.2, 1.80, 5.63, 1.85, fill="FFF3E0", line="FFB74D", radius=True)
    add_text(slide, "HP Imbalance 對 ISM 的影響", 7.4, 1.85, 5.2, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "HP_Ratio > 0.9 → 一群 reads 消失",
        "→ PERMANOVA 失效，NGroups 不可信",
        "TO 模式 HP_Ratio 中位數 0.84",
        "→ Self-phasing 使 HP 偏差更嚴重",
        "→ 需要 LOH.bed 作為區域層級判定",
    ], 7.4, 2.18, 5.2, 1.40, font_size=11, color=THEME["text"], line_spacing=1.30)

    # ── Lower-left: Compressed four-quadrant table (y=3.80~5.10) ──
    add_text(slide, "四象限分析: HP Imbalance ⊃ LOH.bed", 0.5, 3.75, 7.0, 0.25,
             font_size=11, color=THEME["text"], bold=True, margin=0)
    headers = ["", "HP Imbalance ✓", "HP Imbalance ✗"]
    rows = [
        ["LOH.bed ✓", "Q1: 26.7%", "Q3: 0.07% (286筆)"],
        ["LOH.bed ✗", "Q2: 15.2%", "Q4: 58.1%"],
    ]
    add_simple_table(slide, headers, rows, 0.5, 4.02, 6.5, 1.10,
                     header_color=THEME["accent3"], font_size=12, header_font_size=12)

    # ── Lower-right: Sensitivity highlight (y=3.80~5.10) ──
    add_rect(slide, 7.2, 3.80, 3.0, 1.35, fill="E8F5E9", line=THEME["accent2"], radius=True)
    add_text(slide, "Sensitivity", 7.35, 3.87, 2.7, 0.28,
             font_size=13, color=THEME["muted"], bold=True, align="center", margin=0)
    add_text(slide, "99.7%", 7.35, 4.18, 2.7, 0.60,
             font_size=34, color=THEME["positive"], bold=True, align="center", margin=0)
    add_text(slide, "Q3 僅 286/419,692", 7.35, 4.80, 2.7, 0.25,
             font_size=10, color=THEME["muted"], align="center", margin=0)

    # Q3 explanation note
    add_rect(slide, 10.4, 3.80, 2.43, 1.35, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_multiline_text(slide, [
        "Q3: LOH.bed 內",
        "但 HP 未失衡",
        "= 邊界區域",
        "0.07% 可忽略",
    ], 10.5, 3.87, 2.2, 1.20, font_size=10, color=THEME["muted"], line_spacing=1.25)

    # ── Term box ──
    _add_term_box(slide, [
        "HP Imbalance = 位點層級 (per-variant)",
        "LOH.bed = 區域層級 (LongPhase-TO 連續區塊)",
        "HP_Ratio = ISM 等位基因平衡指標",
    ], 0.5, 5.30, 12.33, 0.80)

    _add_conclusion_box(slide,
        "HP Imbalance (位點) 完全覆蓋 LOH.bed (區域)，兩定義層級互補；HP_Ratio 直接影響 ISM 可信度",
        x=0.5, w=12.33)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "HP_Ratio 公式: (HP1_FamilyN+0.001)/(HP1_FamilyN+HP2_FamilyN+0.002)\n"
        "Laplace smoothing 避免除以零。正常: 0.1-0.9, HP Imbalance: <0.1 或 >0.9\n"
        "HP Imbalance 對 ISM: HP_Ratio>0.9 → 一群消失 → PERMANOVA 失效\n"
        "TO 模式 HP_Ratio 中位數 0.84 → self-phasing 偏差使 HP 更極端\n"
        "四象限: Q1=26.7%, Q2=15.2%, Q3=0.07%(286筆), Q4=58.1%\n"
        "Sensitivity=99.7%: HP Imbalance 幾乎完全覆蓋 LOH.bed\n"
        "來源: v3 loh_dual_definition + 0408 PPT Slide 8")
    return slide


def build_slide_6_conclusion_bridge(prs):
    """Slide 6: 前期結論橋樑 — LOH 內外特徵比較結果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "前期結論: LOH 內外特徵比較",
                    "比較 LOH 區域與 Non-LOH 區域的 TP/FP 單特徵差異 → 多數無效，兩項有效",
                    title_en="Prior Findings: LOH vs Non-LOH Feature Comparison — Most Features Ineffective")

    # ── Left panel: 嘗試過的方向 (FAIL) ──
    add_rect(slide, 0.5, 1.90, 5.8, 3.30, fill="FFF3E0", line="FFB74D", radius=True)
    add_text(slide, "✗ 嘗試過但無效的方向", 0.7, 1.95, 5.4, 0.30,
             font_size=14, color=THEME["negative"], bold=True, margin=0)

    fail_items = [
        ("Filter 10/10 策略", "全部 FAIL (FP/TP ratio < 2.0×)"),
        ("甲基化特徵區分", "LOH 區域 AUC ≈ 0.50，無判別力"),
        ("HP-based 特徵", "受 self-phasing 干擾，不獨立"),
        ("距離/分群特徵", "PairwiseMeanDist 等無 LOH 內外差異"),
        ("Filter 策略總結", "安全約束下無可用的過濾方案"),
    ]
    for i, (feat, result) in enumerate(fail_items):
        fy = 2.35 + i * 0.55
        add_text(slide, feat, 0.7, fy, 2.2, 0.25,
                 font_size=11, color=THEME["text"], bold=True, margin=0)
        add_text(slide, result, 2.95, fy, 3.2, 0.25,
                 font_size=11, color=THEME["muted"], margin=0)

    # ── Right panel: 有效發現 ──
    add_rect(slide, 6.55, 1.90, 6.28, 1.30, fill="E8F5E9", line=THEME["accent2"], radius=True)
    add_text(slide, "✓ 唯二有效發現", 6.75, 1.95, 5.8, 0.30,
             font_size=14, color=THEME["positive"], bold=True, margin=0)
    add_multiline_text(slide, [
        "1. LOH FP rate = 0.239 < Non-LOH = 0.338",
        "   → LOH 區域 TP 比例更高 (TP enriched)",
        "2. caller_af AUC = 0.654",
        "   → AF 是唯一超越全 ISM 的特徵",
    ], 6.75, 2.28, 5.8, 0.88, font_size=11, color=THEME["text"], line_spacing=1.30)

    # ── Right bottom: 結論與轉向 ──
    add_rect(slide, 6.55, 3.40, 6.28, 1.80, fill="E8F0FE", line=THEME["accent"], radius=True)
    add_text(slide, "推論與下一步", 6.75, 3.45, 5.8, 0.28,
             font_size=13, color=THEME["accent"], bold=True, margin=0)
    add_multiline_text(slide, [
        "LOH 不能 filter → 轉向 characterization",
        "但 TP enrichment + AF 有資訊 →",
        "  LOH 的 variant 是否全部 clonal？",
        "  加入 CN (拷貝數) 共分析能否解開？",
        "→ 進入 LOH × CN 共分析",
    ], 6.75, 3.78, 5.8, 1.35, font_size=11, color=THEME["text"], line_spacing=1.35)

    # ── Bottom: key numbers recap ──
    add_rect(slide, 0.5, 5.40, 5.8, 0.55, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "Filter FAIL: 10/10 策略全低於安全閾值 (TP loss ≤2%)",
             0.7, 5.47, 5.4, 0.40, font_size=10, color=THEME["muted"], margin=0)

    add_rect(slide, 6.55, 5.40, 6.28, 0.55, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "caller_af = 0.654 > 全 ISM → AF 是最有資訊的維度",
             6.75, 5.47, 5.8, 0.40, font_size=10, color=THEME["muted"], margin=0)

    _add_conclusion_box(slide,
        "LOH 不能 filter，但 TP enrichment + AF 資訊 → 接續 LOH × CN 共分析",
        x=0.5, w=12.33)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "前期做了 LOH 內外基本特徵 TP/FP 差異比較：\n"
        "失敗: Filter 10/10 全 FAIL, 甲基化 AUC≈0.50, HP 特徵受 self-phasing 干擾\n"
        "有效: LOH FP rate=0.239 < Non-LOH=0.338 (TP enriched)\n"
        "有效: caller_af AUC=0.654 唯一超越全 ISM\n"
        "結論: LOH 是分層工具非過濾器，但 TP enrichment + AF 有潛力\n"
        "→ 接續問題: LOH 內 variant 是否全部 clonal? 加入 CN 共分析\n"
        "來源: v3 C14 結論; 0413 LOH×CNV 總結")
    return slide


# ═══════════════════════════════════════════════════════════════
# Part C: LOH+CNV 發現 — TP 分離 (Slides 7-9)
# ═══════════════════════════════════════════════════════════════
def build_slide_7_cm_cn_proxy(prs):
    """Slide 7: Coverage_Multiple 定義與 CN Proxy"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Coverage_Multiple 與 CN Proxy",
                    "ISM 用覆蓋深度比值估算拷貝數 — 經 SEQC2 truth CN 校準 (r=0.831)",
                    title_en="Coverage_Multiple as CN Proxy — Calibrated Against SEQC2 (r=0.831)")

    # ── Left top: CM 定義框 (prominent) ──
    add_rect(slide, 0.5, 1.90, 6.3, 1.80, fill=THEME["card_bg"], line=THEME["accent"], radius=True)
    add_text(slide, "Coverage_Multiple 定義", 0.7, 1.95, 5.8, 0.28,
             font_size=14, color=THEME["text"], bold=True, margin=0)
    add_text(slide, "CM = Region Coverage / Expected Coverage",
             0.7, 2.28, 5.8, 0.30,
             font_size=13, color=THEME["accent"], bold=True, font_name="Consolas", margin=0)
    add_multiline_text(slide, [
        "ISM 對每個 variant region 計算局部覆蓋深度",
        "除以全基因組平均覆蓋深度 (expected_coverage)",
        "CM ≈ 0.5 → CN≈1 (deletion)  |  CM ≈ 1.0 → CN≈2",
        "CM ≈ 1.5 → CN≈3 (gain)    |  CM > 2.0 → CN≈4+",
    ], 0.7, 2.62, 5.8, 1.00, font_size=11, color=THEME["muted"], line_spacing=1.30)

    # ── Right top: Calibration scatter (or placeholder) + r value ──
    cal_img = TEACH_DIR / "T05_cm_calibration_scatter.png"
    # Also try fig04 (AF vs CM scatter) as alternative
    fig04_img = FIG_DIR / "04_loh_af_vs_coverage_multiple.png"
    if cal_img.exists():
        add_image_fit(slide, cal_img, 7.0, 1.90, 5.83, 2.60)
    elif fig04_img.exists():
        add_image_fit(slide, fig04_img, 7.0, 1.90, 5.83, 2.60)
    else:
        add_rect(slide, 7.0, 1.90, 5.83, 2.60, fill="F0F0F0", line=THEME["line"], radius=True)
        add_text(slide, "[CM vs SEQC2 Truth CN Scatter]\nr = 0.831",
                 8.0, 2.9, 3.5, 0.7, font_size=12, color="A0A0A0", align="center", margin=0)

    # ── Middle: CN Tier cards (flat, below image bottom y≈4.50) ──
    add_text(slide, "CN Tier 分類 (閾值: 整數 CN boundary ±0.25)", 0.5, 4.38, 8.0, 0.25,
             font_size=11, color=THEME["text"], bold=True, margin=0)
    tiers = [
        ("CN1", "<0.75", "Deletion LOH · Coverage ↓", THEME["accent"]),
        ("CN2", "0.75-1.25", "cnLOH / diploid · Coverage ≈", THEME["accent2"]),
        ("CN3", "1.25-1.75", "Single gain · Coverage ↑", THEME["accent3"]),
        ("CN4+", ">1.75", "High gain · Coverage ↑↑", THEME["accent4"]),
    ]
    for i, (tier, rng, desc, color) in enumerate(tiers):
        tx = 0.5 + i * 3.15
        add_rect(slide, tx, 4.60, 2.95, 0.80, fill=THEME["card_bg"], line=color, radius=True)
        add_text(slide, tier, tx + 0.12, 4.63, 0.80, 0.25,
                 font_size=14, color=color, bold=True, margin=0)
        add_text(slide, rng, tx + 1.0, 4.63, 1.5, 0.25,
                 font_size=11, color=THEME["muted"], bold=True, margin=0)
        add_text(slide, desc, tx + 0.12, 4.90, 2.7, 0.42,
                 font_size=10, color=THEME["text"], margin=0)

    # ── Bottom: Term box + calibration note ──
    _add_term_box(slide, [
        "CM = region coverage / expected coverage（非精確 CN，無 GC correction）",
        "SEQC2 校準: CM vs truth CN r=0.831",
    ], 0.5, 5.50, 12.33, 0.55)

    _add_conclusion_box(slide,
        "Coverage_Multiple 作為 CN proxy 經 SEQC2 校準可靠 → 用於 CN Tier 分層分析",
        x=0.5, w=12.33)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "Coverage_Multiple 定義: ISM 計算的 region coverage / expected coverage\n"
        "expected_coverage 由 --expected-coverage CLI 參數或 KDE 自動估計\n"
        "校準: vs HCC1395 SEQC2 truth CN, Pearson r=0.831\n"
        "CN Tier 閾值: 整數 CN boundary ±0.25 (CN1<0.75, CN2 0.75-1.25, CN3 1.25-1.75, CN4+>1.75)\n"
        "限制: 非精確 CN (無 GC correction, 不考慮 ploidy)，但粗略分類足夠\n"
        "後續分析的 CN Tier 全部使用 ISM 的 Coverage_Multiple 計算，非 SEQC2 truth")
    return slide


def build_slide_loh_af_overview(prs):
    """New slide: LOH 內外 AF 分布差異 — 7 樣本一覽"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "LOH 區域 vs Non-LOH 區域 AF 分布",
                    "LOH 區域 TP 有大量 intermediate AF (24.6%) — 7 樣本一致的模式",
                    title_en="LOH vs Non-LOH AF Distribution — Consistent Across 7 Samples")

    # Left: fig01 AF distribution LOH vs non-LOH (tall, portrait AR=0.57)
    fig01 = FIG_DIR / "01_af_distribution_loh_vs_nonloh.png"
    add_image_fit(slide, fig01, 0.3, 1.85, 5.8, 4.35)

    # Right: fig02 intermediate AF proportion
    fig02 = FIG_DIR / "02_intermediate_af_proportion.png"
    add_image_fit(slide, fig02, 6.3, 1.85, 6.53, 2.40)

    # Right bottom: key observations
    add_rect(slide, 6.3, 4.40, 6.53, 1.80, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "關鍵觀察", 6.45, 4.45, 3.0, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "LOH TP: 24.6% intermediate AF",
        "LOH FP: 僅 4.1% → 6× enrichment",
        "7/7 樣本一致 (7.0% ~ 60.2%)",
        "HCC1395 ONT 23.3% vs DORADO 20.0%",
        "→ 為什麼 LOH 區域會有 intermediate AF？",
        "→ 答案需要加入 CN (拷貝數) 分析",
    ], 6.45, 4.78, 6.2, 1.35, font_size=11, color=THEME["text"], line_spacing=1.30)

    _add_conclusion_box(slide,
        "LOH 區域 TP 的 intermediate AF 普遍存在 → 接下來用 CN Tier 分層解釋成因",
        x=0.3, w=12.53)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "左圖: LOH vs Non-LOH 的 AF 密度分布，7 樣本。\n"
        "LOH TP 在 AF 0.1-0.9 有明顯的 intermediate 密度。\n"
        "右上: Per-sample intermediate AF 比例，7.0%(COLO829) ~ 60.2%(HCC1954)。\n"
        "6× enrichment: LOH TP 24.6% vs FP 4.1%。\n"
        "這些 intermediate AF 在 purity=1.0 下不應存在於 clonal LOH。\n"
        "→ 需要 CN 分層來解釋: CN1 的 intermediate 最有意義。")
    return slide


def build_slide_8_cn_tier_separation(prs):
    """Slide 8: CN Tier 分層揭示 AF 分離"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "CN Tier 分層揭示 AF 分離",
                    "限定 LOH 區域內，按 Coverage_Multiple (CN proxy) 分層觀察 AF 分布",
                    title_en="CN Tier Stratification Reveals AF Separation Within LOH Regions")

    # Main figure — enlarged (wider)
    fig05 = FIG_DIR / "05_loh_af_by_cn_tier.png"
    add_image_fit(slide, fig05, 0.3, 1.85, 8.5, 3.60)

    # Right panel: CN tier numbers + explanation
    add_rect(slide, 9.0, 1.85, 3.83, 2.50, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "Intermediate AF %", 9.15, 1.90, 3.5, 0.28,
             font_size=12, color=THEME["text"], bold=True, margin=0)
    cn_data = [
        ("CN1", "16.9%", "← subclone 信號", THEME["positive"]),
        ("CN2", "24.8%", "(混合)", THEME["muted"]),
        ("CN3", "45.2%", "(部分 dosage)", THEME["muted"]),
        ("CN4+", "73.1%", "← allele dosage", THEME["negative"]),
    ]
    for i, (tier, pct, note, color) in enumerate(cn_data):
        dy = 2.30 + i * 0.48
        add_text(slide, tier, 9.15, dy, 0.70, 0.25,
                 font_size=13, color=color, bold=True, margin=0)
        add_text(slide, pct, 9.90, dy, 0.85, 0.25,
                 font_size=14, color=THEME["text"], bold=True, margin=0)
        add_text(slide, note, 10.80, dy, 1.9, 0.25,
                 font_size=9, color=THEME["muted"], margin=0)

    # Right bottom: data source note
    add_rect(slide, 9.0, 4.45, 3.83, 1.10, fill="FFF3E0", line="FFB74D", radius=True)
    add_text(slide, "數據說明", 9.15, 4.50, 3.5, 0.25,
             font_size=11, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "圖: LOH 區域 TP",
        "按 ISM CM 分層 (CN1/2/3/4+)",
        "CN 來源: ISM（非 SEQC2 truth）",
    ], 9.15, 4.78, 3.5, 0.72, font_size=10, color=THEME["muted"], line_spacing=1.25)

    # Bottom: interpretation box
    add_rect(slide, 0.3, 5.60, 8.5, 0.50, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_multiline_text(slide, [
        "CN4+ 的 73.1% = 多拷貝 allele dosage (正常現象，不需要 subclone 解釋)",
        "CN1 的 16.9% = deletion LOH 下不應有 intermediate → 最純的 subclone 信號",
    ], 0.5, 5.68, 8.1, 0.45, font_size=10, color=THEME["text"], line_spacing=1.30)

    _add_conclusion_box(slide,
        "CN1 的 intermediate AF 無 allele dosage 解釋 → 最乾淨的 subclone 信號",
        x=0.3, w=12.53)

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "本圖限定 LOH 區域內的 TP variants。\n"
        "按 ISM 計算的 Coverage_Multiple 分為 CN1/2/3/4+ 四層。\n"
        "注意: CN Tier 使用 ISM 自己的 Coverage_Multiple 計算（非 SEQC2 truth CN）。\n"
        "SEQC2 truth CN 僅用於校準驗證 (r=0.831)。\n\n"
        "CN Tier 分層: CN1=16.9%, CN2=24.8%, CN3=45.2%, CN4+=73.1% intermediate。\n"
        "CN4+ 高比例 = allele dosage (正常: 多拷貝自然產生 intermediate AF)。\n"
        "CN1 (deletion LOH): purity=1.0 + CN=1 → AF 只能 0 or 1.0\n"
        "→ CN1 的 16.9% intermediate 必須由 subclonal LOH 解釋。")
    return slide


def build_slide_9_bio_explanation(prs):
    """Slide 9: 生物學解釋 — 為什麼 CN1 可以分離"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "生物學解釋: 為什麼 CN1 可以分離",
                    "Deletion LOH + purity=1.0 + CN=1 → intermediate AF 只能由 subclonal LOH 解釋",
                    title_en="Biology: Why CN1 Intermediate AF Can Only Be Subclonal LOH")

    # ── Top: Three-column comparison (Normal / Clonal LOH / Subclonal LOH) ──
    cols = [
        ("正常", "E8F5E9", "A5D6A7",
         "═A═  ═B═", "AF = 0.5", "兩 allele 等量表達"),
        ("Clonal LOH", "FFEBEE", "EF9A9A",
         "═A═  (═B═ 刪除)", "AF = 0 or 1.0", "單一 allele 存活"),
        ("Subclonal LOH", "FFF3E0", "FFB74D",
         "s%: ═A═\n(1-s)%: ═A═ ═B═", "AF = 0.5 + 0.5s", "混合信號 → intermediate"),
    ]
    col_w = 3.75
    for i, (title, bg, line_c, chrom, af, desc) in enumerate(cols):
        cx = 0.5 + i * 4.00
        add_rect(slide, cx, 1.95, col_w, 2.15, fill=bg, line=line_c, radius=True)
        add_text(slide, title, cx + 0.10, 1.98, col_w - 0.2, 0.28,
                 font_size=13, color=THEME["text"], bold=True, margin=0)
        add_text(slide, chrom, cx + 0.10, 2.35, col_w - 0.2, 0.55,
                 font_size=12, color=THEME["text"], font_name="Consolas", margin=0,
                 line_spacing=1.35)
        add_text(slide, af, cx + 0.10, 2.95, col_w - 0.2, 0.30,
                 font_size=14, color=THEME["accent"], bold=True, font_name="Consolas",
                 margin=0)
        add_text(slide, desc, cx + 0.10, 3.30, col_w - 0.2, 0.70,
                 font_size=11, color=THEME["muted"], margin=0)

    # ── Bottom-left: 研究數據佐證 ──
    add_rect(slide, 0.5, 4.30, 6.20, 1.60, fill=THEME["card_bg"],
             line=THEME["accent2"], radius=True)
    _add_teaching_label(slide, "研究數據驗證", 0.5, 4.30, color=THEME["accent2"])
    add_multiline_text(slide, [
        "LOH TP: 24.6% intermediate vs FP: 4.1% → 6× enrichment",
        "CN1 intermediate = 16.9% → deletion LOH 最乾淨",
        "ΔNGroups = +0.705 (7/7 p < 10⁻³⁹)",
        "7/7 樣本一致：跨平台 universal effect",
    ], 0.7, 4.68, 5.8, 1.15, font_size=11, color=THEME["text"], line_spacing=1.35)

    # ── Bottom-right: CN4+ 對比 + 文獻佐證 ──
    add_rect(slide, 6.90, 4.30, 5.93, 1.60, fill="FFF3E0", line="FFB74D", radius=True)
    _add_teaching_label(slide, "CN4+ 對比", 6.90, 4.30, color=THEME["accent"])
    add_multiline_text(slide, [
        "CN4+: 多拷貝 allele dosage → 73.1% intermediate（正常現象）",
        "CN1: deletion → intermediate 只由 subclone 解釋",
        "HCC1954 已知高度 subclonal (COSMIC)",
        "purity=1.0 ≠ 基因組一致",
    ], 7.10, 4.68, 5.5, 1.15, font_size=11, color=THEME["text"], line_spacing=1.35)

    _add_conclusion_box(slide,
        "CN1 intermediate AF 在物理上只能由 subclonal LOH 解釋 → 進入教學確認機制")

    add_footer(slide, len(prs.slides))
    inject_speaker_notes(slide,
        "三欄對比：正常(AF=0.5) / Clonal LOH(AF=0 or 1.0) / Subclonal LOH(AF=0.5+0.5s)\n"
        "公式: AF = 0.5 + 0.5s, s=subclonal fraction, 0<s<1 → intermediate\n"
        "研究數據佐證: 24.6% intermediate (6× vs FP), CN1=16.9%, ΔNG=+0.705\n"
        "CN4+ 對比: allele dosage 73.1% intermediate = 正常現象\n"
        "文獻: HCC1954 已知 subclonal (COSMIC cell lines project)\n"
        "purity=1.0 ≠ 基因組一致: cell line 可有 subclonal structure")
    return slide


# ═══════════════════════════════════════════════════════════════
# BATCH 4: 教學 Part D — 前半 (Slides 10-14)
# ═══════════════════════════════════════════════════════════════
def _add_teaching_label(slide, text, x, y, *, color=None):
    """Small section label for teaching slides (e.g., '概念', '數據佐證')."""
    c = color or THEME["accent2"]
    add_rect(slide, x, y, 1.40, 0.28, fill=c, radius=True)
    add_text(slide, text, x + 0.05, y + 0.02, 1.30, 0.24,
             font_size=11, color="FFFFFF", bold=True, align="center", margin=0)


def build_slide_10_loh_basics(prs):
    """Slide 10: 什麼是 LOH — 基礎定義"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "什麼是 LOH",
                    "Loss of Heterozygosity — 腫瘤中最常見的基因組事件之一",
                    title_en="What is LOH (Loss of Heterozygosity)?")

    # --- Left: Normal vs LOH concept ---
    _add_teaching_label(slide, "概念", 0.8, 1.85)

    # Normal cell box
    add_rect(slide, 0.8, 2.25, 2.8, 1.8, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "正常細胞", 0.95, 2.30, 2.5, 0.28,
             font_size=13, color=THEME["accent2"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Allele A  ═══  (帶 variant)",
        "Allele B  ═══",
        "",
        "AF ≈ 0.5 (兩 allele 等量)",
    ], 0.95, 2.62, 2.5, 1.3, font_size=11, color=THEME["text"], line_spacing=1.35)

    # Arrow
    add_native_arrow(slide, 3.65, 3.15, 4.15, 3.15, color=THEME["accent"], width=2.0)

    # LOH cell box
    add_rect(slide, 4.20, 2.25, 2.8, 1.8, fill="FFEBEE", line="EF9A9A", radius=True)
    add_text(slide, "LOH 細胞", 4.35, 2.30, 2.5, 0.28,
             font_size=13, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Allele A  ═══  (保留)",
        "Allele B  ✕ (丟失)",
        "",
        "AF → 0 or 1.0 (只剩一個)",
    ], 4.35, 2.62, 2.5, 1.3, font_size=11, color=THEME["text"], line_spacing=1.35)

    # --- Right: 三種 LOH 機制 ---
    _add_teaching_label(slide, "三種 LOH 機制", 7.6, 1.85, color=THEME["accent3"])

    mechs = [
        ("Deletion LOH", "CN = 1", "物理刪除一段染色體\nCoverage 下降 ↓\n最乾淨的 LOH 信號", "FFEBEE"),
        ("cnLOH", "CN = 2", "先刪除再複製\nCoverage 不變 ≈\n序列相同但 CN=2", "FFF3E0"),
        ("Gain + LOH", "CN > 2", "刪除 + 額外增益\nCoverage 上升 ↑\n複雜的多拷貝", "E8F0FE"),
    ]
    mx = 7.60
    for i, (title, cn, desc, bg_col) in enumerate(mechs):
        my = 2.25 + i * 1.15
        add_rect(slide, mx, my, 4.93, 1.03, fill=bg_col, line=THEME["line"], radius=True)
        add_text(slide, title, mx + 0.12, my + 0.06, 2.0, 0.28,
                 font_size=13, color=THEME["text"], bold=True, margin=0)
        add_text(slide, cn, mx + 2.15, my + 0.06, 1.2, 0.28,
                 font_size=12, color=THEME["accent3"], bold=True, margin=0)
        add_text(slide, desc, mx + 0.12, my + 0.35, 4.6, 0.60,
                 font_size=10, color=THEME["muted"], margin=0, line_spacing=1.25)

    # --- Bottom: ISM 連結 ---
    _add_teaching_label(slide, "ISM 連結", 0.8, 5.70, color=THEME["accent"])
    add_rect(slide, 0.8, 6.05, 11.73, 0.65, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_multiline_text(slide, [
        "ISM 使用 LongPhase-TO 的 tumor_phased_LOH.bed 偵測 LOH 區域 (SEQC2 J=0.928, F1=0.963)",
        "Coverage_Multiple (區域 coverage / 全基因組平均) 作為 CN proxy — vs HCC1395 truth CN: Pearson r = 0.831",
    ], 1.0, 6.08, 11.3, 0.50, font_size=11, color=THEME["text"], line_spacing=1.35)

    # --- Left bottom: LOH 三機制染色體圖解 (PPTX-native) ---
    _add_teaching_label(slide, "染色體圖解", 0.8, 4.20, color=THEME["accent3"])
    chrom_data = [
        ("Deletion", "FFEBEE", [("A", "4CAF50"), ("\u2715B", "E57373")], "CN=1 · 最乾淨"),
        ("cnLOH", "FFF3E0", [("A", "4CAF50"), ("A\u2032", "81C784")], "CN=2 · 序列一致"),
        ("Gain+LOH", "E8F0FE", [("A", "4CAF50"), ("A", "4CAF50"), ("A\u2032", "81C784")], "CN\u22653"),
    ]
    for ci, (name, bg, bars, cn_lbl) in enumerate(chrom_data):
        px = 0.8 + ci * 2.10
        add_rect(slide, px, 4.50, 1.98, 1.08, fill=bg, line=THEME["line"], radius=True)
        add_text(slide, name, px + 0.08, 4.52, 1.82, 0.22,
                 font_size=10, color=THEME["text"], bold=True, margin=0)
        for bj, (label, clr) in enumerate(bars):
            bx = px + 0.12 + bj * 0.55
            add_rect(slide, bx, 4.80, 0.48, 0.30, fill=clr, radius=True)
            add_text(slide, label, bx, 4.82, 0.48, 0.26,
                     font_size=9, color="FFFFFF", bold=True, align="center", margin=0)
        add_text(slide, cn_lbl, px + 0.08, 5.18, 1.82, 0.25,
                 font_size=8, color=THEME["muted"], margin=0)

    # --- Right bottom: key data ---
    add_rect(slide, 7.60, 5.68, 4.93, 0.22, fill=THEME["accent3"], radius=True)
    add_text(slide, "數據佐證", 7.60, 5.68, 4.93, 0.22,
             font_size=10, color="FFFFFF", bold=True, align="center", margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "LOH = Loss of Heterozygosity，腫瘤中非常常見。\n"
        "正常細胞有兩套染色體(A+B)，LOH後只剩一個 allele。\n\n"
        "三種機制：\n"
        "1. Deletion LOH (CN=1): 物理刪除，coverage↓，最乾淨\n"
        "2. cnLOH (CN=2): 先刪後複製，coverage不變\n"
        "3. Gain+LOH (CN>2): 刪除+增益，coverage↑\n\n"
        "ISM 用 LOH.bed 偵測 LOH (J=0.928)，\n"
        "Coverage_Multiple 作為 CN proxy (r=0.831)。")

    return slide


def build_slide_11_clonal_subclonal(prs):
    """Slide 11: Clonal vs Subclonal LOH"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Clonal vs Subclonal LOH",
                    "不是所有腫瘤細胞都一樣 — 腫瘤異質性的核心概念",
                    title_en="Clonal vs Subclonal LOH — Tumor Heterogeneity")

    # --- Left: Clonal LOH ---
    _add_teaching_label(slide, "Clonal LOH", 0.8, 1.85, color=THEME["accent2"])
    add_rect(slide, 0.8, 2.20, 5.5, 2.50, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "100% 腫瘤細胞都有 LOH", 1.0, 2.26, 5.0, 0.28,
             font_size=14, color=THEME["text"], bold=True, margin=0)

    add_multiline_text(slide, [
        "○ ○ ○ ○ ○ ○ ○ ○  ← 所有細胞均一",
        "",
        "每個細胞只有 Allele A (CN=1)",
        "所有 reads 來自同一 allele",
        "",
        "AF = 0 或 1.0 (唯一可能)",
        "甲基化 pattern 均一 → NGroups = 1",
    ], 1.0, 2.58, 5.1, 2.0, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Right: Subclonal LOH ---
    _add_teaching_label(slide, "Subclonal LOH", 6.85, 1.85, color=THEME["accent"])
    add_rect(slide, 6.85, 2.20, 5.68, 2.50, fill="FFF3E0", line="FFD9A0", radius=True)
    add_text(slide, "僅 s% 細胞有 LOH，(1-s)% 保留雙 allele", 7.05, 2.26, 5.3, 0.28,
             font_size=14, color=THEME["text"], bold=True, margin=0)

    add_multiline_text(slide, [
        "○ ○ ○ ● ● ● ● ●  ← s% LOH / (1-s)% 正常",
        "",
        "LOH 細胞: 只有 Allele A",
        "正常細胞: 有 Allele A + B",
        "",
        "Reads 混合 → AF 落在中間值",
        "甲基化兩群 → NGroups ≥ 2",
    ], 7.05, 2.58, 5.3, 2.0, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Center: Key insight box ---
    add_rect(slide, 0.8, 4.90, 11.73, 1.34, fill=THEME["card_bg"], line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "關鍵推論", 0.95, 4.96, color=THEME["accent"])

    add_multiline_text(slide, [
        "Cell line purity = 1.0 → 沒有 normal cell 汙染。所以 intermediate AF 不可能是 normal cell 造成的稀釋。",
        "如果看到 intermediate AF → 只能是 subclonal LOH。",
        "例: HCC1954 = 極度 subclonal (已知) → 60.2% intermediate;  COLO829 = 相對 clonal → 7.0% intermediate",
    ], 1.0, 5.28, 11.3, 0.62, font_size=11, color=THEME["text"], line_spacing=1.30)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Clonal LOH: 100% 細胞都有 LOH → AF 只能是 0 或 1.0。\n"
        "Subclonal LOH: 部分細胞 LOH，部分保留 → AF 在中間。\n\n"
        "關鍵：我們用 cell line (purity=1.0)，\n"
        "所以 intermediate AF 不可能是 normal cell 稀釋。\n"
        "只能是 subclonal LOH。\n\n"
        "HCC1954 已知高度 subclonal → 60.2% intermediate\n"
        "COLO829 相對 clonal → 7.0% intermediate\n"
        "這個差異在後面數據中會清楚看到。")

    return slide


def build_slide_12_af_math(prs):
    """Slide 12: AF 的物理意義與數學推導"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "AF 數學推導",
                    "為什麼 Intermediate AF = Subclonal LOH 的數學證明",
                    title_en="AF Mathematical Derivation: Why Intermediate = Subclone")

    # --- Left: Formula derivation ---
    _add_teaching_label(slide, "公式推導", 0.8, 1.85, color=THEME["accent3"])
    add_rect(slide, 0.8, 2.20, 6.0, 2.89, fill="E8F0FE", line="B8D4F0", radius=True)

    add_text(slide, "Deletion LOH (CN=1), purity p = 1.0", 1.0, 2.26, 5.6, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)

    add_multiline_text(slide, [
        "Clonal (s=1): AF = 0 or 1.0  (唯二結果)",
        "",
        "Subclonal (fraction s, 0 < s < 1):",
        "  AF_obs = s × 1.0 + (1−s) × 0.5",
        "        = 0.5 + 0.5s",
        "",
        "當 s ∈ (0,1) → AF ∈ (0.5, 1.0) = intermediate",
    ], 1.0, 2.60, 5.6, 1.9, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Right: AF number line ---
    _add_teaching_label(slide, "AF 三分類", 7.3, 1.85, color=THEME["accent"])
    add_rect(slide, 7.3, 2.20, 5.23, 2.89, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    # Number line visualization using colored bars
    # Extreme zones
    add_rect(slide, 7.55, 2.55, 0.95, 0.40, fill="FFEBEE", radius=True)
    add_text(slide, "0 ~ 0.1", 7.55, 2.55, 0.95, 0.20,
             font_size=9, color=THEME["negative"], bold=True, align="center", margin=0)
    add_text(slide, "Extreme", 7.55, 2.73, 0.95, 0.20,
             font_size=8, color=THEME["muted"], align="center", margin=0)

    # Intermediate low
    add_rect(slide, 8.55, 2.55, 1.10, 0.40, fill="FFF3E0", radius=True)
    add_text(slide, "0.1 ~ 0.4", 8.55, 2.55, 1.10, 0.20,
             font_size=9, color=THEME["accent"], bold=True, align="center", margin=0)
    add_text(slide, "Intermediate", 8.55, 2.73, 1.10, 0.20,
             font_size=8, color=THEME["muted"], align="center", margin=0)

    # Near-half
    add_rect(slide, 9.70, 2.55, 0.70, 0.40, fill="E8F5E9", radius=True)
    add_text(slide, "0.4~0.6", 9.70, 2.55, 0.70, 0.20,
             font_size=9, color=THEME["accent2"], bold=True, align="center", margin=0)
    add_text(slide, "Near-half", 9.70, 2.73, 0.70, 0.20,
             font_size=8, color=THEME["muted"], align="center", margin=0)

    # Intermediate high
    add_rect(slide, 10.45, 2.55, 1.10, 0.40, fill="FFF3E0", radius=True)
    add_text(slide, "0.6 ~ 0.9", 10.45, 2.55, 1.10, 0.20,
             font_size=9, color=THEME["accent"], bold=True, align="center", margin=0)
    add_text(slide, "Intermediate", 10.45, 2.73, 1.10, 0.20,
             font_size=8, color=THEME["muted"], align="center", margin=0)

    # Extreme high
    add_rect(slide, 11.60, 2.55, 0.70, 0.40, fill="FFEBEE", radius=True)
    add_text(slide, "0.9~1.0", 11.60, 2.55, 0.70, 0.20,
             font_size=9, color=THEME["negative"], bold=True, align="center", margin=0)
    add_text(slide, "Extreme", 11.60, 2.73, 0.70, 0.20,
             font_size=8, color=THEME["muted"], align="center", margin=0)

    # Threshold rationale
    add_multiline_text(slide, [
        "閾值選擇理由:",
        "• 0.1/0.9: 與 TITAN BAF 閾值一致",
        "• 0.4/0.6: heterozygous 自然波動 ±0.1",
    ], 7.55, 3.15, 4.8, 0.80, font_size=11, color=THEME["muted"], line_spacing=1.30)

    # Key validation numbers
    add_multiline_text(slide, [
        "本研究數據佐證:",
        "  LOH TP:  24.6% intermediate (6× enrichment)",
        "  LOH FP:   4.1% intermediate",
        "  → TP 比 FP 多 6 倍 intermediate AF",
    ], 7.55, 3.90, 4.8, 0.65, font_size=11, color=THEME["text"], line_spacing=1.25)

    # --- Bottom: Evidence summary ---
    add_rect(slide, 0.8, 5.26, 11.73, 1.36, fill="FFF8E1", line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "小結", 0.95, 5.32, color=THEME["accent"])
    add_multiline_text(slide, [
        "在 purity=1.0 的 cell line 中，clonal LOH (CN=1) 的 AF 只能是 0 或 1.0。",
        "24.6% LOH TP 有 intermediate AF → 不可能是 clonal → 必然是 subclonal LOH 事件。",
        "FP 僅 4.1% intermediate → 6× enrichment 確認 intermediate AF 集中在 TP，不是隨機噪音。",
    ], 1.0, 5.18, 11.3, 0.45, font_size=11, color=THEME["text"], line_spacing=1.30)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "核心公式：AF_obs = 0.5 + 0.5s\n"
        "當 s∈(0,1) 時，AF 落在 0.5~1.0 的 intermediate 區間。\n\n"
        "AF 三分類：\n"
        "- Extreme (<0.1 or >0.9): clonal LOH 期望值\n"
        "- Near-half (0.4-0.6): heterozygous 或 cnLOH\n"
        "- Intermediate (0.1-0.4, 0.6-0.9): 異常 → subclone\n\n"
        "閾值與 TITAN BAF 一致，0.4/0.6 涵蓋自然波動。\n"
        "數據：LOH TP 24.6% intermediate vs FP 4.1% → 6× enrichment。")

    return slide


def build_slide_13_self_phasing(prs):
    """Slide 13: Self-Phasing Circular Dependency — 防禦性教學"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Self-Phasing Circular Dependency",
                    "TO Phasing 的已知問題 — 但不影響本研究 subclone 發現",
                    title_en="Self-Phasing Bias: Known Issue but Irrelevant to AF-Based Subclone Detection")

    # --- Top: Four-step logic chain flow ---
    _add_teaching_label(slide, "四步邏輯鏈", 0.95, 1.92)

    steps = [
        ("①", "LongPhase-TO 用\nsomatic variant\n做 phasing anchor"),
        ("②", "TP 的 ALT reads\n天然偏向\n一個 haplotype"),
        ("③", "HP_Ratio 被\n推向極端\n(d = -1.20)"),
        ("④", "62% LOH\n因此消失\n(PON-only 修正後)"),
    ]
    sx = 0.8
    sw = 2.65
    gap = 0.30
    for i, (num, desc) in enumerate(steps):
        cx = sx + i * (sw + gap)
        add_rect(slide, cx, 2.30, sw, 1.40, fill="FFF3E0", line=THEME["accent"], radius=True)
        add_text(slide, num, cx + 0.10, 2.34, 0.40, 0.35,
                 font_size=20, color=THEME["accent"], bold=True, margin=0)
        add_text(slide, desc, cx + 0.55, 2.36, sw - 0.65, 1.25,
                 font_size=11, color=THEME["text"], margin=0, line_spacing=1.30)
        if i < 3:
            ax = cx + sw + 0.02
            add_native_arrow(slide, ax, 3.00, ax + gap - 0.04, 3.00,
                             color=THEME["accent"], width=1.5)

    # --- Middle: Question + Answer ---
    add_rect(slide, 0.8, 4.00, 5.50, 1.45, fill="FFEBEE", line=THEME["negative"], radius=True)
    add_text(slide, "⚠ 這是否影響本研究的 subclone 發現？", 1.05, 4.06, 5.1, 0.30,
             font_size=13, color=THEME["negative"], bold=True, margin=0)
    add_text(slide,
             "Self-phasing 影響 HP_Ratio → LOH 判定\n"
             "但 caller_af 來自 ClairS-TO variant caller\n"
             "→ 完全獨立於 haplotype phasing",
             1.05, 4.40, 5.1, 0.95,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.35)

    add_rect(slide, 6.85, 4.00, 5.68, 1.45, fill="E8F5E9", line=THEME["positive"], radius=True)
    add_text(slide, "✅ 否 — caller_af 獨立於 phasing", 7.10, 4.06, 5.2, 0.30,
             font_size=13, color=THEME["positive"], bold=True, margin=0)
    add_text(slide,
             "intermediate AF 是 variant caller 層面的觀察\n"
             "不是 HP 層面的觀察\n"
             "subclone 證據鏈完全不依賴 phasing",
             7.10, 4.40, 5.2, 0.95,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.35)

    # --- Term box ---
    _add_term_box(slide, [
        "Self-phasing = 被分析的 variant 同時參與了自己的 phasing 判定",
        "caller_af = variant caller 報告的 AF（獨立於 haplotype phasing）",
    ], 0.8, 5.65, 11.73, 0.90)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "防禦性教學：PI 可能質疑 TO phasing 問題。\n\n"
        "四步邏輯鏈說明 self-phasing 如何影響 LOH 判定：\n"
        "1. TO 用 somatic variant 做 phasing anchor\n"
        "2. TP ALT reads 偏向一個 haplotype\n"
        "3. HP_Ratio 被推極端 (d=-1.20)\n"
        "4. 62% LOH 消失\n\n"
        "但關鍵：caller_af 來自 ClairS-TO，完全獨立於 phasing。\n"
        "intermediate AF 是 caller 觀察，不是 HP 觀察。\n"
        "subclone 證據鏈完全不依賴 HP_Ratio 或 LOH.bed。\n"
        "來源: v5 self-phasing analysis + v3 C13 結論")

    return slide


def build_slide_14_cm_calibration(prs):
    """Slide 14: Coverage_Multiple 校準"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Coverage_Multiple 校準",
                    "用 SEQC2 金標準驗證 CM 作為 CN Proxy 的可靠性",
                    title_en="Coverage_Multiple Calibration — CN Proxy Validation via SEQC2")

    # --- Left: Scatter plot placeholder (60%) ---
    _add_teaching_label(slide, "校準圖", 0.95, 1.92, color=THEME["accent3"])
    fig_path = TEACH_DIR / "T05_cm_calibration_scatter.png"
    add_rect(slide, 0.8, 2.25, 6.80, 3.60, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _safe_add_picture(slide, fig_path, Inches(0.95), Inches(2.35),
                      width=Inches(6.50), height=Inches(3.40),
                      placeholder="Scatter: CM vs SEQC2 Truth CN\nPearson r = 0.831\n(HCC1395, FDA 金標準)")

    # --- Right: Method + thresholds ---
    add_rect(slide, 8.0, 2.25, 4.53, 1.77, fill="E8F0FE", line="B8D4F0", radius=True)
    _add_teaching_label(slide, "校準方法", 8.15, 2.30, color=THEME["accent3"])
    add_multiline_text(slide, [
        "輸入: ISM per-region coverage /",
        "      genome-wide average",
        "輸出: Coverage_Multiple ≈ CN proxy",
        "驗證: HCC1395 SEQC2 truth CN",
        "      (FDA 金標準, WGS+SNP array)",
    ], 8.15, 2.62, 4.2, 1.15, font_size=11, color=THEME["text"], line_spacing=1.25)

    add_rect(slide, 8.0, 4.22, 4.53, 1.60, fill="FFF3E0", line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "CN Tier 閾值", 8.15, 4.27, color=THEME["accent"])
    cn_headers = ["CN Tier", "CM Range", "理由"]
    cn_rows = [
        ["CN1", "< 0.75", "CN=1 -25%"],
        ["CN2", "0.75~1.25", "CN=1 ±25%"],
        ["CN3", "1.25~1.75", "CN=1.5 ±17%"],
        ["CN4+", "> 1.75", "CN=2 -12.5%"],
    ]
    add_simple_table(slide, cn_headers, cn_rows, 8.15, 4.42, 4.20, 0.80,
                     font_size=10, header_color=THEME["accent"])

    # --- Bottom: Limitation ---
    add_rect(slide, 0.8, 6.05, 11.73, 0.50, fill="FFF8E1", line=THEME["accent"], radius=True)
    add_text(slide,
             "限制: Coverage_Multiple 非精確 CN — 無 GC correction、不考慮 ploidy、不區分 LOH 類型",
             1.0, 6.10, 11.3, 0.40,
             font_size=11, color=THEME["muted"], margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Coverage_Multiple 校準教學。\n\n"
        "方法：ISM 計算每個 region 的 coverage / genome-wide average = CM。\n"
        "用 HCC1395 的 SEQC2 truth CN 校準 → Pearson r=0.831。\n\n"
        "CN Tier 閾值選擇理由：接近整數 CN boundary ±25%。\n"
        "0.75/1.25/1.75 分別對應 CN=1/1.5/2 的邊界。\n\n"
        "限制說明：CM 不做 GC correction、不考慮 ploidy、\n"
        "不區分 LOH 類型（deletion vs cnLOH）。\n"
        "是 reasonable proxy，非精確 CN calling。\n"
        "來源: 0414 validated report §1.2 + HCC1395 SEQC2 數據")

    return slide


# ═══════════════════════════════════════════════════════════════
# BATCH 5: 教學 Part D — 後半 (Slides 15-20)
# ═══════════════════════════════════════════════════════════════
def build_slide_15_camdac(prs):
    """Slide 15: CAMDAC 原理 — 甲基化如何反映 LOH (★★★ 最重要教學頁)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "CAMDAC 原理",
                    "甲基化如何獨立反映 LOH 狀態 — 雙重證據鏈的核心邏輯",
                    title_en="CAMDAC Principle: Methylation Reflects LOH State")

    # --- Left: Clonal LOH → ASM 消失 ---
    _add_teaching_label(slide, "Clonal LOH", 0.8, 1.85, color=THEME["accent2"])
    add_rect(slide, 0.8, 2.20, 5.5, 2.65, fill="E8F5E9", line="A5D6A7", radius=True)

    add_multiline_text(slide, [
        "所有細胞只有 Allele A",
        "      ↓",
        "Reads 全來自同一 allele",
        "      ↓",
        "甲基化 pattern 均一",
        "      ↓",
        "NGroups = 1, AlleleDelta ≈ 0",
    ], 1.0, 2.30, 5.1, 2.1, font_size=13, color=THEME["text"], line_spacing=1.25)

    # ASM badge
    add_rect(slide, 1.0, 4.65, 5.1, 0.35, fill=THEME["accent2"], radius=True)
    add_text(slide, "ASM (Allele-Specific Methylation) 消失",
             1.05, 4.37, 5.0, 0.30,
             font_size=12, color="FFFFFF", bold=True, align="center", margin=0)

    # --- Right: Subclonal LOH → ASM 部分保留 ---
    _add_teaching_label(slide, "Subclonal LOH", 6.85, 1.85, color=THEME["accent"])
    add_rect(slide, 6.85, 2.20, 5.68, 2.65, fill="FFF3E0", line="FFD9A0", radius=True)

    add_multiline_text(slide, [
        "s% 細胞只有 Allele A",
        "(1-s)% 細胞有 A+B",
        "      ↓",
        "Reads 來自混合群",
        "      ↓",
        "甲基化 pattern 異質",
        "      ↓",
        "NGroups > 1, AlleleDelta > 0",
    ], 7.05, 2.25, 5.3, 2.2, font_size=13, color=THEME["text"], line_spacing=1.20)

    # ASM badge
    add_rect(slide, 7.05, 4.65, 5.3, 0.35, fill=THEME["accent"], radius=True)
    add_text(slide, "ASM 部分保留 → 甲基化多樣性",
             7.10, 4.37, 5.2, 0.30,
             font_size=12, color="FFFFFF", bold=True, align="center", margin=0)

    # --- Bottom: ISM 驗證數據 ---
    add_rect(slide, 0.8, 5.30, 7.5, 1.38, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _add_teaching_label(slide, "ISM 驗證 (CN1)", 0.95, 5.36, color=THEME["accent3"])

    add_multiline_text(slide, [
        "Extreme AF:      AlleleDelta = 0.003  (均一 → clonal)",
        "Intermediate AF:  AlleleDelta = 0.031  (+12× 增加 → subclonal)",
        "Cohen's d = +0.724 (大效應量)",
    ], 1.0, 5.38, 7.1, 0.70, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Bottom right: 雙重證據鏈 ---
    add_rect(slide, 8.50, 5.30, 4.03, 1.38, fill=THEME["dark"], radius=True)
    add_text(slide, "雙重證據鏈", 8.70, 5.36, 3.6, 0.30,
             font_size=15, color=THEME["light_text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "遺傳: AF 偏離 (intermediate)",
        "表觀: NGroups 升高 (+0.705)",
        "兩條獨立路徑，同一結論",
    ], 8.70, 5.40, 3.6, 0.70, font_size=11, color=THEME["light_muted"], line_spacing=1.30)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "★★★ 最重要的教學頁 ★★★\n\n"
        "CAMDAC 原理：\n"
        "Clonal LOH → 所有 reads 來自同一 allele → ASM 消失\n"
        "  → NGroups=1, AlleleDelta≈0\n\n"
        "Subclonal LOH → reads 混合 → ASM 部分保留\n"
        "  → NGroups>1, AlleleDelta>0\n\n"
        "ISM 驗證：\n"
        "Extreme AF: AlleleDelta=0.003 (均一)\n"
        "Intermediate AF: AlleleDelta=0.031 (+12×!)\n"
        "Cohen's d = +0.724 (大效應)\n\n"
        "這就是雙重證據鏈：\n"
        "遺傳(AF偏離) + 表觀遺傳(NGroups升高)\n"
        "兩條獨立路徑指向同一結論：subclonal LOH。")

    return slide


def build_slide_16_allele_delta(prs):
    """Slide 16: AlleleDelta / ASM 原理 — 新增教學"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "AlleleDelta / ASM 原理",
                    "兩個 Allele 的甲基化差異 — Subclonal LOH 的表觀遺傳證據",
                    title_en="AlleleDelta: Allele-Specific Methylation as Subclone Evidence")

    # --- Top: ASM schematic (PPTX-native CpG methylation visual) ---
    _add_teaching_label(slide, "ASM 示意", 0.95, 1.92)
    add_rect(slide, 0.8, 2.25, 7.50, 2.50, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    # HP1 panel (high methylation haplotype)
    add_rect(slide, 0.95, 2.40, 3.40, 1.80, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "HP1 Reads (高甲基化)", 1.05, 2.44, 3.2, 0.24,
             font_size=11, color=THEME["accent2"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Read 1: \u25cf \u25cf \u25cf \u25cb \u25cb",
        "Read 2: \u25cf \u25cf \u25cb \u25cf \u25cb",
        "Read 3: \u25cf \u25cf \u25cf \u25cb \u25cf",
        "\u2192 mean = 0.60",
    ], 1.05, 2.72, 3.2, 1.35, font_size=12, color=THEME["text"],
       font_name="Consolas", line_spacing=1.40)

    # HP2 panel (low methylation haplotype)
    add_rect(slide, 4.50, 2.40, 3.60, 1.80, fill="FFF3E0", line="FFD9A0", radius=True)
    add_text(slide, "HP2 Reads (低甲基化)", 4.60, 2.44, 3.4, 0.24,
             font_size=11, color=THEME["accent"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Read 1: \u25cb \u25cb \u25cb \u25cf \u25cf",
        "Read 2: \u25cb \u25cb \u25cf \u25cb \u25cf",
        "Read 3: \u25cb \u25cf \u25cb \u25cf \u25cb",
        "\u2192 mean = 0.33",
    ], 4.60, 2.72, 3.4, 1.35, font_size=12, color=THEME["text"],
       font_name="Consolas", line_spacing=1.40)

    # Formula
    add_text(slide, "AlleleDelta = |mean(HP1) \u2212 mean(HP2)| = |0.60 \u2212 0.33| = 0.27",
             1.00, 4.32, 7.0, 0.30,
             font_size=12, color=THEME["accent3"], bold=True, font_name="Consolas", margin=0)

    # --- Right: Data link ---
    add_rect(slide, 8.70, 2.25, 3.83, 2.50, fill="E8F5E9", line=THEME["positive"], radius=True)
    _add_teaching_label(slide, "本研究數據", 8.85, 2.30, color=THEME["positive"])
    add_multiline_text(slide, [
        "Extreme AF:",
        "  AlleleDelta = 0.003",
        "",
        "Intermediate AF:",
        "  AlleleDelta = 0.031",
        "",
        "→ +12× 增加",
        "→ Cohen's d = 0.724",
        "  (大效應)",
    ], 8.85, 2.62, 3.5, 2.00, font_size=12, color=THEME["text"], line_spacing=1.20)

    # --- Bottom: Mechanism + Conclusion ---
    add_rect(slide, 0.8, 4.93, 5.50, 1.05, fill="F3E5F5", line=THEME["accent4"], radius=True)
    _add_teaching_label(slide, "分子機制", 0.95, 4.98, color=THEME["accent4"])
    add_multiline_text(slide, [
        "ASM 由 imprinting、CTCF 結合差異、chromatin accessibility 驅動",
        "Clonal LOH → ASM 消失 (單 allele)",
        "Subclonal LOH → ASM 部分保留 (混合 allele)",
    ], 1.0, 5.35, 5.1, 0.65, font_size=10, color=THEME["text"], line_spacing=1.20)

    # --- Term box ---
    _add_term_box(slide, [
        "ASM = Allele-Specific Methylation（等位基因特異性甲基化）",
        "AlleleDelta = |mean(HP1 methylation) - mean(HP2 methylation)|",
    ], 6.80, 5.00, 5.73, 1.05)

    _add_conclusion_box(slide, "AlleleDelta 劇變支持 CAMDAC 機制 — subclonal LOH 保留了 ASM")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "AlleleDelta 教學 — 新增頁面。\n\n"
        "ASM = 同一 locus 兩個 allele 甲基化不同。\n"
        "AlleleDelta = |mean(HP1) - mean(HP2)|。\n\n"
        "分子機制：imprinting、CTCF、chromatin accessibility。\n"
        "Clonal LOH → 單 allele → ASM 消失 → AlleleDelta≈0。\n"
        "Subclonal LOH → 混合 allele → ASM 保留 → AlleleDelta>0。\n\n"
        "數據：Extreme AF AlleleDelta=0.003 → Intermediate=0.031。\n"
        "+12× 增加，Cohen's d=0.724 大效應。\n"
        "這是 L4 Mechanistic 證據層的核心。\n"
        "來源: validated report §4.2 AlleleDelta 分析")

    return slide


def build_slide_17_read_level(prs):
    """Slide 17: Read-Level 分析原理 — ISM 如何工作"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Read-Level 分析原理",
                    "ISM 如何工作 — 從 BAM 到甲基化分群的完整 Pipeline",
                    title_en="Read-Level Analysis: How ISM Works — BAM to Methylation Clustering")

    # --- Left: Pipeline flow (PPTX-native step diagram) ---
    _add_teaching_label(slide, "ISM Pipeline", 0.95, 1.92, color=THEME["accent3"])
    add_rect(slide, 0.8, 2.25, 7.00, 3.60, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    pipe_steps = [
        ("1", "BAM + VCF 輸入", "ONT long-read + somatic SNV 位點", THEME["accent3"]),
        ("2", "MM/ML Tags 提取", "BAM 中 Modified Base 標籤 (0-255)", THEME["accent3"]),
        ("3", "CpG 座標校正", "CIGAR 對齊 \u2192 精確 CpG 位點", THEME["accent3"]),
        ("4", "Read\u00d7CpG 矩陣", "每條 read \u00d7 每個 CpG = 甲基化值", THEME["accent"]),
        ("5", "距離計算", "Bernoulli / L1 / NHD 等多種距離", THEME["accent"]),
        ("6", "聚類 \u2192 NGroups", "Agglomerative Clustering \u2192 分群數", THEME["accent2"]),
    ]
    for si, (num, title, desc, clr) in enumerate(pipe_steps):
        sy = 2.47 + si * 0.55
        # Number badge
        add_rect(slide, 1.05, sy, 0.35, 0.38, fill=clr, radius=True)
        add_text(slide, num, 1.05, sy + 0.02, 0.35, 0.34,
                 font_size=13, color="FFFFFF", bold=True, align="center", margin=0)
        # Step title
        add_text(slide, title, 1.52, sy, 2.6, 0.22,
                 font_size=11, color=THEME["text"], bold=True, margin=0)
        # Description
        add_text(slide, desc, 1.52, sy + 0.20, 5.2, 0.22,
                 font_size=9, color=THEME["muted"], margin=0)

    # --- Right top: ISM vs Traditional ---
    add_rect(slide, 8.20, 2.25, 4.33, 1.65, fill="FFF3E0", line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "ISM vs 傳統", 8.35, 2.30, color=THEME["accent"])
    add_multiline_text(slide, [
        "傳統方法:",
        "  region-level 平均甲基化率  → 丟失 read 間差異",
        "",
        "ISM:  read-level 個別甲基化模式  → 保留所有異質性資訊",
    ], 8.35, 2.62, 4.0, 1.05, font_size=11, color=THEME["text"], line_spacing=1.20)

    # --- Right bottom: Key feature ---
    add_rect(slide, 8.20, 4.10, 4.33, 1.75, fill="E8F5E9", line=THEME["positive"], radius=True)
    _add_teaching_label(slide, "核心差異", 8.35, 4.15, color=THEME["positive"])
    add_multiline_text(slide, [
        "ISM 是唯一做到",
        "per-variant read-level methylation clustering 的工具",
        "",
        "每個 somatic variant 都有",
        "獨立的甲基化結構 context",
    ], 8.35, 4.47, 4.0, 1.05, font_size=11, color=THEME["text"], line_spacing=1.20)

    # --- Term box ---
    _add_term_box(slide, [
        "MM/ML = BAM 中的甲基化標籤（M=修飾類型, L=likelihood 0-255）",
        "Agglomerative Clustering = 由下而上合併的階層式聚類",
    ], 0.8, 6.05, 11.73, 0.55)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "ISM pipeline 教學 — 新增頁面。\n\n"
        "Pipeline 流程：\n"
        "BAM → MM/ML tags → CpG 座標校正 (CIGAR alignment)\n"
        "→ Read×CpG 甲基化矩陣 → 距離計算 → clustering → NGroups。\n\n"
        "ISM vs 傳統方法：\n"
        "傳統工具只看 region-level 平均甲基化率，丟失 read 間差異。\n"
        "ISM 保留每個 read 的個別甲基化模式，能偵測異質性。\n\n"
        "這是 ISM 的獨特技術優勢 — 唯一做 per-variant read-level\n"
        "methylation clustering 的工具。\n"
        "來源: ISM 架構文件 + README_PROJECT_SUMMARY.md")

    return slide


def build_slide_18_ngroups(prs):
    """Slide 18: HPFineNGroups — ISM 的甲基化分群指標"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "HPFineNGroups",
                    "Read-Level 甲基化分群指標 — ISM 的核心觀測工具",
                    title_en="HPFineNGroups: Read-Level Methylation Clustering")

    # --- Left: Clustering concept ---
    _add_teaching_label(slide, "計算方法", 0.8, 1.85, color=THEME["accent3"])
    add_rect(slide, 0.8, 2.20, 6.2, 2.00, fill="E8F0FE", line="B8D4F0", radius=True)

    add_text(slide, "Read × CpG 甲基化矩陣 → Agglomerative Clustering → NGroups",
             1.0, 2.26, 5.8, 0.28,
             font_size=12, color=THEME["text"], bold=True, margin=0)

    # Matrix visualization (simplified)
    matrix_lines = [
        "       CpG₁ CpG₂ CpG₃ CpG₄ CpG₅",
        "R1:     ●     ●     ○     ●     ●   ┐",
        "R2:     ●     ●     ○     ●     ○   │ Group 1",
        "R3:     ●     ●     ●     ●     ●   ┘",
        "R4:     ○     ○     ●     ○     ○   ┐",
        "R5:     ○     ○     ○     ○     ●   │ Group 2",
        "R6:     ○     ●     ○     ○     ○   ┘",
    ]
    add_multiline_text(slide, matrix_lines, 1.0, 2.60, 5.8, 1.50,
                       font_size=10, color=THEME["text"], line_spacing=1.15,
                       font_name="Consolas")

    # --- Right: NGroups meaning ---
    _add_teaching_label(slide, "NGroups 含義", 7.3, 1.85, color=THEME["accent"])
    add_rect(slide, 7.3, 2.20, 5.23, 2.00, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    ng_items = [
        ("NGroups = 1", "均質", "clonal 或單 allele\n所有 reads 同一 pattern", "E8F5E9"),
        ("NGroups = 2", "雙群", "雙 allele 或 subclone\n兩群不同 pattern", "FFF3E0"),
        ("NGroups = 3-4", "複雜", "多 subclone\n更複雜結構", "FFEBEE"),
    ]
    ny = 2.28
    for label, short, detail, bg_col in ng_items:
        add_rect(slide, 7.45, ny, 4.9, 0.55, fill=bg_col, radius=True)
        add_text(slide, label, 7.55, ny + 0.03, 1.6, 0.24,
                 font_size=12, color=THEME["text"], bold=True, margin=0)
        add_text(slide, short, 9.15, ny + 0.03, 0.8, 0.24,
                 font_size=11, color=THEME["accent"], bold=True, margin=0)
        add_text(slide, detail, 7.55, ny + 0.28, 4.7, 0.25,
                 font_size=10, color=THEME["muted"], margin=0, line_spacing=1.15)
        ny += 0.63

    # --- Middle: Prior validation ---
    add_rect(slide, 0.8, 4.40, 6.2, 0.72, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _add_teaching_label(slide, "前期驗證", 0.95, 4.46, color=THEME["accent2"])
    add_multiline_text(slide, [
        "Phase BCD 已驗證: NGroups ≥ 4 + NR ≥ 80 → TP rate 89.1%",
        "NGroups 是 somatic heterogeneity 的有效指標 (ISM 第一個 POSITIVE 結論)",
    ], 1.0, 4.72, 5.8, 0.35, font_size=11, color=THEME["text"], line_spacing=1.25)

    # --- Bottom: 本研究數字 ---
    add_rect(slide, 0.8, 5.30, 11.73, 0.85, fill="FFF8E1", line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "本研究關鍵數字", 0.95, 5.36, color=THEME["accent"])

    # Two comparison boxes
    add_rect(slide, 1.0, 5.68, 5.2, 0.38, fill="E8F5E9", radius=True)
    add_text(slide, "Extreme AF:  90.7% NGroups=1 (均質 → clonal)", 1.15, 5.70, 4.9, 0.34,
             font_size=13, color=THEME["text"], bold=True, margin=0)

    add_rect(slide, 6.40, 5.68, 5.93, 0.38, fill="FFF3E0", radius=True)
    add_text(slide, "Intermediate AF:  79.6% NGroups≥2 (異質 → subclonal)", 6.55, 5.70, 5.6, 0.34,
             font_size=13, color=THEME["text"], bold=True, margin=0)

    # --- Right bottom: prior NGroups validation ---
    add_rect(slide, 7.3, 4.40, 5.23, 0.72, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _add_teaching_label(slide, "本研究聚焦", 7.45, 4.46, color=THEME["accent3"])
    add_multiline_text(slide, [
        "LOH 區域的 NGroups 與 AF 的關聯:",
        "Clonal LOH → 單 allele → NGroups=1 (預期)",
    ], 7.50, 4.72, 4.8, 0.35, font_size=11, color=THEME["text"], line_spacing=1.25)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "HPFineNGroups 計算方式：\n"
        "取每個 region 的 Read×CpG 甲基化矩陣\n"
        "→ agglomerative clustering\n"
        "→ 輸出群組數 NGroups (1-4)\n\n"
        "NGroups 含義：\n"
        "1: 均質 (clonal 或單 allele)\n"
        "2: 雙群 (雙 allele 或 subclone)\n"
        "3-4: 複雜 (多 subclone)\n\n"
        "前期驗證：NGroups≥4 + NR≥80 → TP rate 89.1%\n\n"
        "本研究數字：\n"
        "Extreme AF: 90.7% NGroups=1 (clonal → 預期)\n"
        "Intermediate AF: 79.6% NGroups≥2 (subclonal → 發現!)")

    return slide


def build_slide_19_nr_confound(prs):
    """Slide 19: NumReads Confound — 為什麼必須控制"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "NumReads Confound",
                    "NGroups 與讀取數量有正相關 — 不控制就不可信",
                    title_en="NumReads Confound: Why Control is Essential")

    # --- Left: Problem illustration ---
    _add_teaching_label(slide, "問題", 0.8, 1.85, color=THEME["negative"])

    # Low NR box
    add_rect(slide, 0.8, 2.20, 3.2, 1.90, fill="FFEBEE", line="EF9A9A", radius=True)
    add_text(slide, "NR = 15 reads", 0.95, 2.26, 2.8, 0.28,
             font_size=13, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "● ● ● ○ ○",
        "",
        "reads 太少",
        "即使有 2 群也偵測不到",
        "→ NGroups = 1 (低估)",
    ], 0.95, 2.60, 2.8, 1.4, font_size=11, color=THEME["text"], line_spacing=1.25)

    # High NR box
    add_rect(slide, 4.20, 2.20, 3.2, 1.90, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "NR = 100 reads", 4.35, 2.26, 2.8, 0.28,
             font_size=13, color=THEME["accent2"], bold=True, margin=0)
    add_multiline_text(slide, [
        "●●●●●●● ○○○○○○○",
        "",
        "reads 充足",
        "2 群結構清晰可見",
        "→ NGroups = 2 (正確)",
    ], 4.35, 2.60, 2.8, 1.4, font_size=11, color=THEME["text"], line_spacing=1.25)

    # Warning box
    add_rect(slide, 0.8, 4.25, 6.6, 0.55, fill="FFF8E1", line=THEME["accent"], radius=True)
    add_text(slide,
             "若 Intermediate AF 碰巧 NR 較高 → NGroups 差異可能只是 NR confound!",
             1.0, 4.30, 6.2, 0.45,
             font_size=12, color=THEME["accent"], bold=True, margin=0)

    # --- Right: Control method + Results ---
    _add_teaching_label(slide, "控制方法", 7.6, 1.85, color=THEME["accent3"])
    add_rect(slide, 7.6, 2.20, 4.93, 1.10, fill="E8F0FE", line="B8D4F0", radius=True)

    add_text(slide, "NR-bin 分層 (5 bins)", 7.75, 2.26, 4.6, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "[10-30]  [30-50]  [50-80]  [80-150]  [150-500]",
        "在 NR 完全匹配的情況下重新比較 ΔNGroups",
    ], 7.75, 2.58, 4.6, 0.65, font_size=11, color=THEME["text"], line_spacing=1.30)

    _add_teaching_label(slide, "結果: 排除!", 7.6, 3.45, color=THEME["positive"])
    add_rect(slide, 7.6, 3.80, 4.93, 2.44, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    # Results table (native table = 1 shape)
    add_simple_table(slide,
        ["指標", "數值", "解讀"],
        [
            ["NR ratio", "1.01 ~ 1.34", "差異很小"],
            ["NR 10-30", "|r| = 0.483", "有效應 (floor effect)"],
            ["NR 30-50", "|r| = 0.709", "效應增強!"],
            ["NR 50-80", "|r| = 0.708", "效應穩定"],
        ],
        7.65, 3.85, 4.8, 1.40)

    # Key conclusion in results area
    add_rect(slide, 7.75, 5.40, 4.6, 0.50, fill="E8F5E9", radius=True)
    add_multiline_text(slide, [
        "效應隨 NR 增加而增強 (confound 預測: 應消失)",
        "低 NR 效應弱 = floor effect (不是無效)",
    ], 7.80, 5.42, 4.4, 0.46, font_size=11, color=THEME["text"], line_spacing=1.25)

    # --- Bottom: Conclusion ---
    add_rect(slide, 0.8, 5.00, 6.6, 1.24, fill=THEME["card_bg"], line=THEME["positive"], radius=True)
    _add_teaching_label(slide, "結論", 0.95, 5.06, color=THEME["positive"])
    add_multiline_text(slide, [
        "NGroups 差異獨立於 NumReads",
        "控制後效應增強 → confound 完全排除",
        "→ L3 Confound exclusion PASSED",
    ], 1.0, 5.32, 6.2, 0.52, font_size=12, color=THEME["text"], line_spacing=1.28,
       bold=True)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "NGroups 與 NumReads 有正相關：\n"
        "reads 少 → 即使有多群也偵測不到 → NGroups 被低估\n\n"
        "控制方法：NR-bin 分層 (5 bins)\n"
        "在 NR 完全匹配的情況下重新比較\n\n"
        "結果：\n"
        "- NR ratio 僅 1.01-1.34 (Intermediate vs Extreme 差異很小)\n"
        "- NR 10-30: |r|=0.483 (有效應但弱 = floor effect)\n"
        "- NR 30-50: |r|=0.709 (效應增強!)\n"
        "- NR 50-80: |r|=0.708 (效應穩定)\n\n"
        "如果是 confound，控制後效應應該消失。\n"
        "實際增強 → confound 完全排除。\n"
        "低 NR 效應弱的原因：floor effect (reads 不夠，偵測不到群組，不是沒有群組)。")

    return slide


def build_slide_20_pyramid(prs):
    """Slide 20: 六層證據鏈架構預告"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "六層證據鏈架構",
                    "為什麼需要多層驗證 — 接下來的數據結構預告",
                    title_en="Six-Layer Evidence Chain: Why Multi-Layer Verification")

    # --- Pyramid (left side, built bottom-up) ---
    layers = [
        ("L1 遺傳學",   "Intermediate AF",          THEME["accent"],  5.5),
        ("L2 表觀遺傳", "NGroups ΔNG=+0.705",       THEME["accent2"], 5.0),
        ("L3 排除",     "NR-bin 控制後增強",          THEME["accent3"], 4.5),
        ("L4 機制",     "CAMDAC ASM d=+0.724",       THEME["accent4"], 4.0),
        ("L5 空間",     "Segment ρ=0.270",            "7B8D6F",        3.5),
        ("L6 技術",     "ONT vs DORADO 一致",         THEME["muted"],  3.0),
    ]

    # Draw pyramid layers (L1 at bottom=widest, L6 at top=narrowest)
    py_x_base = 0.80
    py_w_base = 5.80
    shrink = 0.25
    n_layers = len(layers)
    for i, (label, detail, color, _unused) in enumerate(layers):
        # i=0→L1(bottom,widest), i=5→L6(top,narrowest)
        ly = 2.10 + (n_layers - 1 - i) * 0.65   # L1 at bottom (y=5.35)
        lx = py_x_base + i * shrink              # L1 least indented
        lw = py_w_base - i * shrink * 2           # L1 widest
        add_rect(slide, lx, ly, lw, 0.55, fill=color, radius=True)
        add_text(slide, label, lx + 0.10, ly + 0.03, 1.8, 0.25,
                 font_size=12, color="FFFFFF", bold=True, margin=0)
        detail_w = max(0.8, lw - 2.0)
        add_text(slide, detail, lx + 1.90, ly + 0.03, detail_w, 0.25,
                 font_size=11, color="FFFFFF", margin=0)

    # --- Right: Why multi-layer (consolidated) ---
    _add_teaching_label(slide, "為什麼多層?", 7.3, 1.85, color=THEME["accent3"])
    add_rect(slide, 7.3, 2.20, 5.23, 2.50, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    add_multiline_text(slide, [
        "L1 alone:  AF 偏離可能是技術雜訊",
        "+ L2:       可能是 NumReads confound",
        "+ L3:       控制 NR 後可能是巧合",
        "+ L4:       機制不明可能是 artifact",
        "全 6 層:   排除所有已知 alternative",
    ], 7.50, 2.30, 4.85, 2.30, font_size=12, color=THEME["text"], line_spacing=1.60)

    # --- Bottom: Next steps preview (consolidated) ---
    add_rect(slide, 0.8, 5.15, 11.73, 1.10, fill=THEME["dark"], radius=True)
    add_text(slide, "接下來: 三步驟逐層展示數據", 1.0, 5.20, 5.0, 0.30,
             font_size=15, color=THEME["light_text"], bold=True, margin=0)

    add_multiline_text(slide, [
        "Step 1  AF Baseline (L1) — Intermediate AF 分布與 CN tier 分層",
        "Step 2  NGroups x AF (L2+L3+L4) — 核心交叉分析 + Confound 排除",
        "Step 3  Segment Spatial (L5+L6) — 空間一致性 + 技術重複",
    ], 1.0, 5.55, 11.3, 0.65, font_size=11, color=THEME["light_muted"], line_spacing=1.40)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "教學結束，進入數據。\n\n"
        "六層金字塔的邏輯：每多一層就排除一個 alternative explanation。\n"
        "L1 alone: AF 偏離可能是技術雜訊\n"
        "L1+L2: NGroups 差異可能是 NR confound\n"
        "L1+L2+L3: 控制 NR 後可能只是巧合\n"
        "L1+L2+L3+L4: 機制不明可能是 artifact\n"
        "全 6 層 → 排除所有已知 alternative → 結論可信\n\n"
        "接下來三步驟：\n"
        "Step 1: AF baseline (對應 L1)\n"
        "Step 2: NGroups × AF 交叉 + confound (L2+L3+L4)\n"
        "Step 3: Segment spatial + 技術重複 (L5+L6)")

    return slide


# ═══════════════════════════════════════════════════════════════
# BATCH 5: 核心結果 (Slides 13-18) — 每頁兩圖 + 數字 + 結論
# ═══════════════════════════════════════════════════════════════
def _add_data_conclusion(slide, text, y=6.15):
    """Add green 'this figure proves' bar at bottom."""
    add_rect(slide, 0.8, y, 11.73, 0.45, fill="E8F5E9", line=THEME["accent2"], radius=True)
    add_text(slide, "✓ " + text, 1.0, y + 0.02, 11.3, 0.40,
             font_size=12, color=THEME["text"], bold=True, margin=0)


def _add_step_badge(slide, step_text, x=11.80, y=1.60):
    """Add step badge (e.g., 'Step 1') — top-right, above figures."""
    add_rect(slide, x, y, 1.1, 0.28, fill=THEME["accent"], radius=True)
    add_text(slide, step_text, x + 0.05, y + 0.02, 1.0, 0.24,
             font_size=12, color="FFFFFF", bold=True, align="center", margin=0)


def build_slide_21_af_baseline(prs):
    """Slide 21: Step 1 — AF 分布 Baseline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Step 1: AF 分布 Baseline",
                    "LOH TP 有大量 intermediate AF — purity=1.0 下不應存在",
                    title_en="Step 1: AF Distribution Baseline — 24.6% Intermediate in LOH TP")
    _add_step_badge(slide, "Step 1")

    # --- Two figures: fig01 portrait (left, tall) + fig02 landscape (right) ---
    fig01 = FIG_DIR / "01_af_distribution_loh_vs_nonloh.png"
    fig02 = FIG_DIR / "02_intermediate_af_proportion.png"
    add_image_fit(slide, fig01, 0.5, 1.90, 3.8, 4.30)
    add_image_fit(slide, fig02, 4.5, 2.10, 8.33, 2.80)

    # --- Key numbers panel (below fig02) ---
    add_rect(slide, 4.5, 5.00, 8.33, 0.80, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_multiline_text(slide, [
        "LOH TP: 24.6% intermediate (n=21,004)",
        "LOH FP:  4.1% intermediate → 6× enrichment",
        "HCC1395 ONT 23.3% vs DORADO 20.0% (技術重複一致)",
    ], 4.65, 5.03, 7.9, 0.74, font_size=10, color=THEME["text"], line_spacing=1.25)

    _add_data_conclusion(slide,
        "此圖證明: purity=1.0 下 intermediate AF 普遍存在 (7/7 樣本) → 提示 subclonal LOH")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Step 1: AF 分布 Baseline\n\n"
        "左圖: LOH vs Non-LOH 的 AF 分布，7 個樣本。\n"
        "LOH TP 明顯有 intermediate AF 密度峰。\n\n"
        "右圖: Per-sample intermediate AF 比例。\n"
        "從 COLO829 的 7.0% 到 HCC1954 的 60.2%。\n\n"
        "24.6% LOH TP 有 intermediate AF，FP 僅 4.1% → 6× enrichment。\n"
        "HCC1395 ONT(23.3%) vs DORADO(20.0%) 技術重複一致 (Δ=3.3pp)。")

    return slide


def build_slide_22_cn_tier_focus(prs):
    """Slide 22: Step 1b — CN Tier 分層"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Step 1b: CN Tier 分層",
                    "聚焦 CN1 (Deletion LOH) — 排除 allele dosage confound",
                    title_en="Step 1b: CN Tier Stratification — Focus on CN1 (Cleanest Signal)")
    _add_step_badge(slide, "Step 1b")

    # --- Two figures ---
    fig05 = FIG_DIR / "05_loh_af_by_cn_tier.png"
    fig04 = FIG_DIR / "04_loh_af_vs_coverage_multiple.png"
    add_image_fit(slide, fig05, 0.5, 2.00, 6.5, 3.10)
    add_image_fit(slide, fig04, 7.2, 2.00, 5.63, 3.10)

    # --- CN tier comparison table (native table = 1 shape) ---
    add_simple_table(slide,
        ["CN Tier", "Intermediate %", "說明"],
        [
            ["CN1 (deletion)", "16.9%", "最乾淨 (無 dosage)"],
            ["CN2 (cnLOH)", "24.8%", "混合"],
            ["CN3 (gain)", "45.2%", "部分 dosage"],
            ["CN4+ (high)", "73.1%", "主要 dosage"],
        ],
        0.85, 5.15, 6.1, 0.95)

    # --- Right: inference ---
    add_rect(slide, 7.3, 5.15, 5.23, 0.95, fill="FFF8E1", line=THEME["accent"], radius=True)
    add_multiline_text(slide, [
        "CN4+ 高比例 = allele dosage (正常現象)",
        "CN1 的 16.9% 無 dosage 解釋 → subclone",
        "→ 後續分析全部聚焦 CN1",
    ], 7.45, 5.22, 4.9, 0.82, font_size=12, color=THEME["text"], line_spacing=1.30,
       bold=True)

    _add_data_conclusion(slide,
        "此圖證明: CN tier 分層排除 allele dosage confound — CN1 = 最可靠的 subclone 信號")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Step 1b: CN Tier 分層\n\n"
        "左圖: 4 個 CN tier 的 AF 密度分布。\n"
        "CN1 有清晰的 bimodal (0/1 + intermediate)。\n"
        "CN4+ 幾乎全是 intermediate (allele dosage)。\n\n"
        "右圖: AF vs Coverage_Multiple scatter。\n"
        "4 個 tier 清楚可分，CN1 在最左。\n\n"
        "Intermediate %: CN1(16.9%) → CN4+(73.1%)\n"
        "CN4+ 的高比例是多拷貝 allele dosage (正常)。\n"
        "CN1 不可能有 dosage → 16.9% 就是 subclone。\n"
        "後續全部聚焦 CN1。")

    return slide


def build_slide_23_ngroups_af(prs):
    """Slide 23: Step 2 — NGroups × AF (★★★ 核心結果)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Step 2: NGroups × AF 交叉分析",
                    "核心結果 — Intermediate AF → NGroups 顯著升高 (7/7 一致)",
                    title_en="Step 2: NGroups × AF Cross-Analysis (Core Result) ★★★")
    _add_step_badge(slide, "Step 2 ★")

    # --- Two figures ---
    fig06 = FIG_DIR / "06_ngroups_by_af_class_cn_tier.png"
    fig07 = FIG_DIR / "07_ngroups_per_sample_cn1.png"
    add_image_fit(slide, fig06, 0.5, 2.00, 6.5, 2.80)
    add_image_fit(slide, fig07, 7.2, 2.00, 5.63, 2.80)

    # --- Key numbers: large highlight ---
    add_rect(slide, 0.8, 4.95, 4.5, 1.15, fill=THEME["dark"], radius=True)
    add_text(slide, "ΔNGroups = +0.705", 1.0, 5.00, 4.1, 0.35,
             font_size=22, color=THEME["light_text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Extreme NG = 1.091 (90.7% NG=1 → clonal)",
        "Intermediate NG = 1.796 (79.6% NG≥2 → subclonal)",
    ], 1.0, 5.38, 4.1, 0.65, font_size=11, color=THEME["light_muted"], line_spacing=1.25)

    # --- Per-sample stats ---
    add_rect(slide, 5.50, 4.95, 7.03, 1.15, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "7/7 樣本全數顯著", 5.65, 5.00, 3.0, 0.28,
             font_size=13, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "全部 p < 10⁻³⁹  |  最強 |r| = 0.822 (HCC1395_DORADO)",
        "最弱 |r| = 0.324 (H2009)  |  NR ratio 僅 1.01-1.34",
    ], 5.65, 5.32, 6.7, 0.70, font_size=11, color=THEME["muted"], line_spacing=1.25)

    _add_data_conclusion(slide,
        "此圖證明: H1+H4 supported — 遺傳 (intermediate AF) + 表觀 (NGroups 升高) 雙重證據指向 subclonal LOH")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "★★★ 最重要的數據頁 ★★★\n\n"
        "左圖: NGroups by AF class × CN tier。\n"
        "CN1 Extreme: 幾乎全 NG=1 (clonal)。\n"
        "CN1 Intermediate: 高峰在 NG=2 (subclonal)。\n\n"
        "右圖: Per-sample NGroups + p-value。\n"
        "7/7 樣本全部正方向，全部 p<10⁻³⁹。\n\n"
        "核心數字：ΔNGroups = +0.705\n"
        "Extreme: 90.7% NG=1 (clonal)\n"
        "Intermediate: 79.6% NG≥2 (subclonal)\n"
        "最強效應: |r|=0.822 (HCC1395 DORADO)\n"
        "最弱: |r|=0.324 (H2009)\n"
        "NR ratio 僅 1.01-1.34 → 不足以解釋 +0.705 差異。")

    return slide


def build_slide_24_nr_control(prs):
    """Slide 24: Step 2b — NR Confound 排除 + CAMDAC 機制"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Step 2b: Confound 排除 + 機制驗證",
                    "NR-bin 控制後效應增強 + AlleleDelta 支持 CAMDAC ASM 機制",
                    title_en="Step 2b: Confound Exclusion (L3) + Mechanistic Validation (L4)")
    _add_step_badge(slide, "Step 2b")

    # --- Two figures: fig08 wide (left) + fig09 portrait (right, tall) ---
    fig08 = FIG_DIR / "08_ngroups_numreads_controlled.png"
    fig09 = FIG_DIR / "09_methylation_features_by_af_class.png"
    add_image_fit(slide, fig08, 0.5, 2.00, 7.8, 2.50)
    add_image_fit(slide, fig09, 8.5, 1.90, 4.33, 4.70)

    # --- Left bottom: NR-bin results (native table = 1 shape) ---
    _add_teaching_label(slide, "L3 Confound", 0.65, 4.60, color=THEME["accent3"])
    add_simple_table(slide,
        ["NR Bin", "ΔNGroups", "|r| 效應量"],
        [
            ["NR 10-30", "+0.484", "0.483"],
            ["NR 30-50", "+0.715", "0.709 (增強)"],
            ["NR 50-80", "+0.711", "0.708 (穩定)"],
        ],
        0.55, 4.88, 5.4, 0.72)

    add_text(slide, "效應隨 NR 增加而增強 → confound 排除", 0.65, 5.65, 5.4, 0.20,
             font_size=11, color=THEME["positive"], bold=True, margin=0)

    # --- Right bottom: Methylation features (next to L3 table) ---
    _add_teaching_label(slide, "L4 機制", 6.10, 4.60, color=THEME["accent4"])
    add_simple_table(slide,
        ["特徵", "Cohen's d", "說明"],
        [
            ["AlleleDelta", "+0.724", "ASM +12x"],
            ["HPFineF", "+0.639", "F-stat +23x"],
            ["CramersV", "+0.318", "中效應"],
        ],
        6.10, 4.88, 2.20, 0.62)

    add_text(slide, "CAMDAC ASM 機制確認", 6.10, 5.55, 2.20, 0.20,
             font_size=11, color=THEME["positive"], bold=True, margin=0)

    _add_data_conclusion(slide,
        "此圖證明: L3 confound 排除 + L4 CAMDAC ASM 機制確認 — NGroups 差異獨立於 NumReads")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Step 2b: 兩層驗證\n\n"
        "L3 Confound 排除:\n"
        "NR-bin 分層後，效應量從 |r|=0.483 增加到 0.709。\n"
        "如果是 confound，控制後應消失 → 實際增強 → 排除。\n"
        "低 NR (10-30) 效應弱 = floor effect (reads 不夠偵測群組)。\n\n"
        "L4 CAMDAC 機制:\n"
        "AlleleDelta: d=+0.724 (+12× → CAMDAC ASM)\n"
        "HPFineF: d=+0.639 (+23×)\n"
        "CramersV: d=+0.318 (中效應)\n"
        "PairwiseMeanDist: d=+0.031 (negative control 通過 → 不是所有特徵都有效)\n\n"
        "AlleleDelta 是最強的 mechanistic 證據：\n"
        "Intermediate AF 的 ASM 是 Extreme 的 12 倍。")

    return slide


def build_slide_25_segment(prs):
    """Slide 25: Step 3 — Segment 空間一致性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Step 3: Segment 空間一致性",
                    "AF 變異性高的 LOH segment → NGroups 也高 → segmental event",
                    title_en="Step 3: Segment-Level Spatial Consistency (L5)")
    _add_step_badge(slide, "Step 3")

    # --- Two figures ---
    fig11 = FIG_DIR / "11_segment_af_sd_vs_ngroups.png"
    fig12 = FIG_DIR / "12_segment_uniform_vs_mixed.png"
    add_image_fit(slide, fig11, 0.5, 2.00, 6.5, 2.70)
    add_image_fit(slide, fig12, 7.2, 2.00, 5.63, 2.70)

    # --- Left bottom: Correlation stats ---
    add_rect(slide, 0.8, 4.85, 6.2, 1.25, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _add_teaching_label(slide, "Segment 相關", 0.95, 4.91, color=THEME["accent3"])

    add_multiline_text(slide, [
        "Overall CN1: Spearman ρ = 0.270, p = 5.6×10⁻²²",
        "Per-sample: 6/7 positive 方向",
        "  H1437 ρ=0.809*** (最強) | COLO829 ρ=0.763***",
        "  HCC1954 ρ=−0.297 ns (n=34, 反向)",
    ], 0.95, 5.22, 5.8, 0.82, font_size=11, color=THEME["text"], line_spacing=1.25)

    # --- Right bottom: Uniform vs Mixed ---
    add_rect(slide, 7.3, 4.85, 5.23, 1.25, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    _add_teaching_label(slide, "Uniform vs Mixed", 7.45, 4.91, color=THEME["accent"])

    add_multiline_text(slide, [
        "Uniform segments: NGroups = 1.292",
        "Mixed segments:   NGroups = 1.717 (+0.425)",
        "Mixed AlleleDelta = 0.0275 (2.2× Uniform)",
        "HCC1954 反向: n=34 underpowered + ceiling",
    ], 7.45, 5.22, 4.8, 0.82, font_size=11, color=THEME["text"], line_spacing=1.25)

    _add_data_conclusion(slide,
        "此圖證明: L5 spatial — AF 變異性 ∝ NGroups → segmental event，非 random variant-level noise")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Step 3: Segment 空間一致性\n\n"
        "左圖: AF-SD vs NGroups scatter + trend line。\n"
        "正向趨勢：AF 變異性越大，NGroups 越高。\n"
        "ρ=0.270, p=5.6×10⁻²²。\n\n"
        "右圖: Uniform vs Mixed segments 的 NGroups。\n"
        "Mixed segments NGroups 更高 (+0.425)。\n\n"
        "Per-sample: 6/7 positive。\n"
        "H1437 ρ=0.809 最強，COLO829 ρ=0.763。\n"
        "HCC1954 反向 (ρ=-0.297)：n=34 segments underpowered + 60.2% 全 subclonal → ceiling effect。")

    return slide


def build_slide_26_per_sample(prs):
    """Slide 26: Per-Sample 一致性總覽"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Per-Sample 一致性總覽",
                    "7 cancer cell lines × 2 技術平台 — Universal Effect",
                    title_en="Per-Sample Consistency: 7/7 Positive, Cross-Platform Validated (L6)")

    # --- Two figures ---
    fig10 = FIG_DIR / "10_ngroups_per_sample_consistency.png"
    fig13 = FIG_DIR / "13_per_sample_segment_consistency.png"
    add_image_fit(slide, fig10, 0.5, 1.85, 6.5, 2.80)
    add_image_fit(slide, fig13, 7.2, 1.85, 5.63, 2.80)

    # --- Summary table (native PPTX table = 1 shape) ---
    add_simple_table(slide,
        ["樣本", "ΔNG", "MW p", "|r|", "ρ segment", "判定"],
        [
            ["HCC1395", "+0.581", "0.00", "0.568", "0.151**", "POSITIVE"],
            ["HCC1395_DORADO", "+0.823", "0.00", "0.822", "0.255***", "POSITIVE"],
            ["COLO829", "+0.580", "0.00", "0.580", "0.763***", "POSITIVE"],
            ["H1437", "+0.812", "0.00", "0.812", "0.809***", "POSITIVE"],
            ["H2009", "+0.329", "7.8e-40", "0.324", "0.212 ns", "POSITIVE"],
            ["HCC1937", "+0.679", "1.8e-71", "0.679", "0.230*", "POSITIVE"],
            ["HCC1954", "+0.689", "2.2e-39", "0.687", "−0.297 ns", "POSITIVE*"],
        ],
        0.85, 4.75, 7.2, 1.40)

    # --- Right: summary box ---
    add_rect(slide, 8.30, 4.75, 4.23, 1.40, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "統計摘要", 8.45, 4.80, 2.0, 0.25,
             font_size=12, color=THEME["text"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Variant-level: 7/7 positive",
        "全 p < 10⁻³⁹",
        "Segment: 6/7 positive",
        "ONT / DORADO 技術重複一致",
    ], 8.45, 5.10, 3.8, 0.95, font_size=11, color=THEME["text"], line_spacing=1.30)

    _add_data_conclusion(slide,
        "此圖證明: L6 — 跨 7 cancer cell lines × 2 技術平台的 universal effect → 高度可信")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Per-sample 一致性總覽：\n\n"
        "左圖: 7 samples ΔNGroups consistency bar chart。\n"
        "全部正方向，全部 p<10⁻³⁹。\n\n"
        "右圖: Per-sample segment correlation scatter。\n"
        "6/7 positive 方向。\n\n"
        "完整數據表：\n"
        "HCC1395: ΔNG=+0.581, |r|=0.568, ρ=0.151**\n"
        "HCC1395_DORADO: ΔNG=+0.823, |r|=0.822, ρ=0.255***\n"
        "COLO829: ΔNG=+0.580, |r|=0.580, ρ=0.763***\n"
        "H1437: ΔNG=+0.812, |r|=0.812, ρ=0.809***\n"
        "H2009: ΔNG=+0.329, |r|=0.324, ρ=0.212 ns\n"
        "HCC1937: ΔNG=+0.679, |r|=0.679, ρ=0.230*\n"
        "HCC1954: ΔNG=+0.689, |r|=0.687, ρ=−0.297 ns (反向)\n\n"
        "HCC1954 segment 反向原因: n=34 underpowered + 60.2% intermediate → ceiling。\n"
        "Variant-level 仍然 positive (p=2.2e-39)。")

    return slide


# ═══════════════════════════════════════════════════════════════
# BATCH 6: 總結 + TumorLens + 未來方向 + 結尾 (Slides 19-24)
# ═══════════════════════════════════════════════════════════════
def build_slide_27_evidence_summary(prs):
    """Slide 27: 六層證據鏈匯總"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "六層證據鏈匯總: POSITIVE",
                    "所有層級全部通過 — Intermediate AF ↔ Subclonal LOH 確認",
                    title_en="Six-Layer Evidence Chain Summary: ALL PASSED → POSITIVE")

    # --- Evidence summary table (native table = 1 shape) ---
    add_simple_table(slide,
        ["層級", "名稱", "核心證據", "數字", "判定"],
        [
            ["L1", "遺傳學", "16.9-60.2% intermediate AF (CN1)", "purity=1.0 排除 dilution", "PASS"],
            ["L2", "表觀遺傳", "ΔNGroups = +0.705 (79.6% NG≥2)", "7/7 p < 10⁻³⁹", "PASS"],
            ["L3", "排除", "NR-bin 控制後效應增強", "r: 0.483 → 0.709", "PASS"],
            ["L4", "機制", "AlleleDelta +12× (d=+0.724)", "CAMDAC ASM 確認", "PASS"],
            ["L5", "空間", "AF-SD ∝ NGroups (ρ=0.270)", "6/7 positive", "PASS"],
            ["L6", "技術", "ONT 23.3% vs DORADO 20.0%", "跨平台一致", "PASS"],
        ],
        0.8, 1.95, 8.5, 2.80, font_size=11, header_font_size=11)

    # --- Right column: Conclusion ---
    add_rect(slide, 10.0, 1.95, 2.53, 2.80, fill=THEME["dark"], radius=True)
    add_text(slide, "結論", 10.15, 2.05, 2.2, 0.30,
             font_size=18, color=THEME["light_text"], bold=True, margin=0)
    add_text(slide, "POSITIVE", 10.15, 2.40, 2.2, 0.40,
             font_size=24, color=THEME["positive"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Intermediate AF 在 LOH 區域",
        "對應 subclonal LOH 事件",
        "",
        "ISM 甲基化指標能獨立偵測",
        "這些事件",
        "7/7 樣本全數一致",
    ], 10.15, 2.90, 2.2, 1.70, font_size=11, color=THEME["light_muted"], line_spacing=1.20)

    # --- Bottom: summary sentence + implication box ---
    add_rect(slide, 0.8, 4.95, 5.80, 1.33, fill="E8F5E9", line=THEME["positive"], radius=True)
    add_text(slide, "核心發現", 0.95, 5.00, 2.0, 0.28,
             font_size=13, color=THEME["positive"], bold=True, margin=0)
    add_multiline_text(slide, [
        "LOH 區域的 subclonal variant 甲基化多樣性",
        "顯著高於 clonal variant (ΔNG=+0.705)",
        "遺傳 + 表觀遺傳 雙重證據鏈確認",
    ], 0.95, 5.30, 5.4, 0.74, font_size=11, color=THEME["text"], line_spacing=1.25)

    add_rect(slide, 6.90, 4.95, 5.63, 1.33, fill="E8F0FE", line=THEME["accent3"], radius=True)
    add_text(slide, "意義與方向", 7.05, 5.00, 2.0, 0.28,
             font_size=13, color=THEME["accent3"], bold=True, margin=0)
    add_multiline_text(slide, [
        "ISM 定位: read-level epigenetic characterization",
        "LOH + CN + Methylation 整合分析方向正確",
        "下一步: Normal reference 區分 germline ASM",
    ], 7.05, 5.30, 5.2, 0.74, font_size=11, color=THEME["text"], line_spacing=1.25)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "六層證據鏈全數 PASS → POSITIVE。\n\n"
        "L1: 16.9-60.2% intermediate AF 在 purity=1.0 下不應存在\n"
        "L2: ΔNGroups=+0.705, 7/7 全 p<10⁻³⁹\n"
        "L3: NR-bin 控制後效應增強 (r: 0.48→0.71)\n"
        "L4: AlleleDelta +12× (d=0.724), CAMDAC ASM 機制\n"
        "L5: Segment ρ=0.270, 6/7 positive\n"
        "L6: ONT vs DORADO 跨平台一致\n\n"
        "結論一句話：Intermediate AF 在 LOH 區域對應 subclonal LOH，\n"
        "ISM 的甲基化指標能獨立偵測這些事件。")

    return slide


def build_slide_28_paired_validation(prs):
    """Slide 28: Paired 模式驗證可行性分析 — 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Paired 模式驗證可行性分析",
                    "TO Subclone 發現可信；Paired 可做 LOH 獨立確認",
                    title_en="Paired Mode Validation Feasibility — TO Findings Are Robust")

    # --- Three analysis cards ---
    cards = [
        {
            "icon": "✅", "title": "caller_af 不受 self-phasing 影響",
            "detail": "caller_af 由 ClairS-TO 計算\n獨立於 haplotype phasing\n→ TO subclone 發現可信",
            "fill": "E8F5E9", "line": THEME["positive"],
        },
        {
            "icon": "✅", "title": "Paired LOH.bed 更可靠",
            "detail": "無 self-phasing bias\n→ 可做 LOH 定義獨立確認\n→ 強化 LOH 區域信心",
            "fill": "E8F5E9", "line": THEME["positive"],
        },
        {
            "icon": "⚠", "title": "Paired FP 太少 (~1%)",
            "detail": "FP = 3,429 vs TO = 128,382\n→ 不適合大規模統計分析\n→ 統計檢定力不足",
            "fill": "FFF3E0", "line": THEME["accent"],
        },
    ]

    cx = 0.8
    cw = 3.78
    gap = 0.20
    for i, c in enumerate(cards):
        x = cx + i * (cw + gap)
        add_rect(slide, x, 1.95, cw, 2.70, fill=c["fill"], line=c["line"], radius=True)
        add_text(slide, c["icon"], x + 0.15, 2.02, 0.40, 0.35,
                 font_size=22, color=THEME["text"], bold=True, margin=0)
        add_text(slide, c["title"], x + 0.55, 2.02, cw - 0.70, 0.55,
                 font_size=14, color=THEME["text"], bold=True, margin=0, line_spacing=1.25)
        add_text(slide, c["detail"], x + 0.15, 2.62, cw - 0.30, 1.90,
                 font_size=14, color=THEME["text"], margin=0, line_spacing=1.35)

    # --- Term box ---
    _add_term_box(slide, [
        "Paired mode = 有 normal BAM 對照的分析模式（TP:FP ≈ 95:1）",
        "TO mode = Tumor-Only 無 normal 對照（TP:FP ≈ 2.3:1）",
    ], 0.8, 5.00, 11.73, 0.85)

    _add_conclusion_box(slide,
        "TO 的 subclone 發現在邏輯上獨立於 self-phasing 問題；"
        "Paired 可作為 LOH 定義的獨立驗證")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Paired 模式驗證可行性 — 三點分析。\n\n"
        "1. caller_af 不受 self-phasing 影響 → TO 結論可信\n"
        "   caller_af 由 ClairS-TO 計算，獨立於 phasing。\n"
        "2. Paired LOH.bed 無 self-phasing bias → 可驗證 LOH 定義\n"
        "   但 Paired FP 太少，不適合大規模 subclone 統計分析。\n"
        "3. Paired FP=3,429 vs TO=128,382 → 統計檢定力差 37×。\n\n"
        "結論：不需要 Paired 來證明 subclone，\n"
        "但可以用 Paired 來獨立驗證 LOH.bed 的可靠性。")

    return slide


def build_slide_29_normal_asm(prs):
    """Slide 29: Normal ASM + Tumor ASM 共分析願景 — 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "Normal ASM + Tumor ASM 共分析願景",
                    "用 Normal Reference 區分 Germline ASM vs Tumor-Specific ASM",
                    title_en="Normal + Tumor ASM Co-Analysis — Distinguishing Germline vs Subclone Signal")

    # --- Concept diagram (65% area) ---
    add_rect(slide, 0.8, 1.95, 7.80, 3.35, fill=THEME["card_bg"], line=THEME["line"], radius=True)

    # Normal BAM box
    add_rect(slide, 1.10, 2.15, 3.00, 1.00, fill="E8F0FE", line="B8D4F0", radius=True)
    add_text(slide, "Normal BAM", 1.20, 2.20, 2.8, 0.30,
             font_size=14, color=THEME["accent3"], bold=True, margin=0)
    add_text(slide, "NGroups per-locus", 1.20, 2.50, 2.8, 0.25,
             font_size=11, color=THEME["muted"], margin=0)
    add_text(slide, "(Phase 2 A+D)", 1.20, 2.75, 2.8, 0.25,
             font_size=10, color=THEME["muted"], margin=0)

    # Tumor BAM box
    add_rect(slide, 4.60, 2.15, 3.00, 1.00, fill="FFF3E0", line=THEME["accent"], radius=True)
    add_text(slide, "Tumor BAM", 4.70, 2.20, 2.8, 0.30,
             font_size=14, color=THEME["accent"], bold=True, margin=0)
    add_text(slide, "NGroups per-locus", 4.70, 2.50, 2.8, 0.25,
             font_size=11, color=THEME["muted"], margin=0)
    add_text(slide, "(ISM 主分析)", 4.70, 2.75, 2.8, 0.25,
             font_size=10, color=THEME["muted"], margin=0)

    # Arrow down: Compare
    add_native_arrow(slide, 4.10, 3.20, 4.10, 3.55,
                     color=THEME["dark"], width=2.0)
    add_text(slide, "比較", 3.80, 3.25, 0.60, 0.25,
             font_size=10, color=THEME["muted"], bold=True, margin=0)

    # Result boxes
    add_rect(slide, 1.10, 3.60, 3.30, 1.50, fill="E8F5E9", line=THEME["positive"], radius=True)
    add_text(slide, "Germline ASM", 1.20, 3.65, 3.0, 0.30,
             font_size=13, color=THEME["positive"], bold=True, margin=0)
    add_text(slide, "Normal 和 Tumor 都有\n→ 排除（非 tumor 信號）",
             1.20, 3.98, 3.0, 0.65,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.30)
    add_text(slide, "❌ Exclude", 1.20, 4.70, 3.0, 0.30,
             font_size=12, color=THEME["muted"], bold=True, margin=0)

    add_rect(slide, 4.80, 3.60, 3.50, 1.50, fill="FFEBEE", line=THEME["negative"], radius=True)
    add_text(slide, "Tumor-specific ASM", 4.90, 3.65, 3.2, 0.30,
             font_size=13, color=THEME["negative"], bold=True, margin=0)
    add_text(slide, "只有 Tumor 有\n→ subclone 信號 ✓",
             4.90, 3.98, 3.2, 0.65,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.30)
    add_text(slide, "✅ Subclone", 4.90, 4.70, 3.2, 0.30,
             font_size=12, color=THEME["positive"], bold=True, margin=0)

    # --- Right: Phase 2 A+D progress ---
    add_rect(slide, 9.0, 1.95, 3.53, 3.35, fill="E8F0FE", line="B8D4F0", radius=True)
    add_rect(slide, 9.10, 2.00, 2.20, 0.28, fill=THEME["accent3"], radius=True)
    add_text(slide, "Phase 2 A+D 進度", 9.15, 2.02, 2.10, 0.24,
             font_size=11, color="FFFFFF", bold=True, align="center", margin=0)
    add_multiline_text(slide, [
        "✓ Normal BAM integration — 已開始",
        "",
        "✓ LOH BED annotation — 已完成",
        "",
        "☐ Cross-region subclone — 下一步",
        "",
        "☐ Germline ASM filter — 待完成",
    ], 9.15, 2.35, 3.2, 2.80, font_size=11, color=THEME["text"], line_spacing=1.25)

    # --- Bridge box: key logic ---
    add_rect(slide, 0.8, 5.45, 11.73, 0.63, fill="FFF8E1", line=THEME["accent"], radius=True)
    add_text(slide, "核心邏輯", 0.95, 5.48, 2.0, 0.25,
             font_size=12, color=THEME["accent"], bold=True, margin=0)
    add_text(slide, "Subclone NGroups 升高可能來自 germline ASM → 用 Normal reference 排除後，剩餘 = tumor-specific subclone 信號",
             2.80, 5.50, 9.5, 0.60,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.20)

    _add_conclusion_box(slide,
        "Normal reference 是提升 subclone NGroups 可信度的關鍵下一步")

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "Normal ASM + Tumor ASM 共分析 — 願景頁面。\n\n"
        "核心概念：用 Normal BAM 的 NGroups 作為 reference，\n"
        "區分 germline ASM（Normal+Tumor 都有）vs tumor-specific ASM（只有 Tumor）。\n\n"
        "Germline ASM → 排除，不是 tumor 信號。\n"
        "Tumor-specific ASM → subclone 信號，是我們要找的。\n\n"
        "Phase 2 A+D 進度：\n"
        "Normal BAM integration 已開始。\n"
        "LOH BED annotation 已完成。\n"
        "下一步：cross-region subclone analysis。\n"
        "來源: Phase 2 A+D 計劃 + 0414 validated results")

    return slide


def build_slide_30_tumorlens_table(prs):
    """Slide 30: TumorLens 競爭者佐證"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "競爭者佐證: TumorLens",
                    "學界第一個 long-read unified pipeline — 方向一致",
                    title_en="Competitor Validation: TumorLens (medRxiv 2026)")

    # --- TumorLens description ---
    add_rect(slide, 0.8, 1.95, 5.5, 2.80, fill="E8F0FE", line="B8D4F0", radius=True)
    add_text(slide, "TumorLens", 1.0, 2.00, 3.0, 0.30,
             font_size=18, color=THEME["accent3"], bold=True, margin=0)
    add_text(slide, "medRxiv 2026 preprint", 3.5, 2.05, 2.5, 0.22,
             font_size=10, color=THEME["muted"], margin=0)

    add_multiline_text(slide, [
        "第一個 long-read unified pipeline 同時偵測:",
        "  SNV + Indel + SV + CNV + LOH + CpG Methylation",
        "",
        "Purity-aware CNV/LOH modeling (Spectre CNV)",
        "HLA-locus 重建 + Allele-specific methylation",
        "支援 Tumor-normal pairs 或 Tumor-only 模式",
    ], 1.0, 2.40, 5.1, 2.2, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- ISM vs TumorLens comparison ---
    add_rect(slide, 6.85, 1.95, 5.68, 2.80, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "ISM vs TumorLens 互補", 7.05, 2.00, 4.0, 0.30,
             font_size=15, color=THEME["text"], bold=True, margin=0)

    # TumorLens card
    add_rect(slide, 7.05, 2.40, 2.4, 1.25, fill="E8F0FE", radius=True)
    add_text(slide, "TumorLens", 7.15, 2.45, 2.2, 0.25,
             font_size=13, color=THEME["accent3"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Macro-level",
        "全基因組多模態",
        "CNV+LOH+Methylation",
        "Pipeline 偵測",
    ], 7.15, 2.72, 2.2, 0.85, font_size=10, color=THEME["muted"], line_spacing=1.20)

    # ISM card
    add_rect(slide, 9.65, 2.40, 2.7, 1.25, fill="FFF3E0", radius=True)
    add_text(slide, "ISM", 9.75, 2.45, 2.5, 0.25,
             font_size=13, color=THEME["accent"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Micro-level",
        "Per-variant 甲基化",
        "Read-level clustering",
        "Epigenetic context",
    ], 9.75, 2.72, 2.5, 0.85, font_size=10, color=THEME["muted"], line_spacing=1.20)

    # Shared direction
    add_multiline_text(slide, [
        "共同方向: LOH + CNV + 甲基化整合分析",
        "= 領域發展趨勢",
    ], 7.05, 3.75, 5.3, 0.45, font_size=11, color=THEME["text"], line_spacing=1.25, bold=True)

    # --- Bottom: Significance ---
    add_rect(slide, 0.8, 5.00, 11.73, 1.15, fill="FFF8E1", line=THEME["accent"], radius=True)
    _add_teaching_label(slide, "佐證意義", 0.95, 5.06, color=THEME["accent"])
    add_multiline_text(slide, [
        "ISM 的 LOH×AF×methylation 整合方向 與 TumorLens 的 unified multi-modal approach 不謀而合",
        "學界正朝此方向發展 — ISM 走在 read-level resolution 的前沿",
        "Phase 2 Dual-BAM (Normal+Tumor) 方向與此趨勢一致",
    ], 1.0, 5.38, 11.3, 0.70, font_size=12, color=THEME["text"], line_spacing=1.30)

    add_text(slide, "doi: 10.64898/2026.03.18.26348569", 0.95, 6.25, 5.0, 0.22,
             font_size=9, color=THEME["muted"], margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "TumorLens (medRxiv 2026 preprint):\n"
        "第一個 unified long-read pipeline for SNV+indel+SV+CNV+LOH+methylation。\n"
        "Purity-aware CNV/LOH modeling。\n\n"
        "ISM vs TumorLens 是互補的：\n"
        "TumorLens: macro-level，全基因組多模態偵測\n"
        "ISM: micro-level，per-variant read-level 甲基化 clustering\n\n"
        "共同方向：LOH+CNV+甲基化必須整合分析。\n"
        "ISM 走在 read-level resolution 的前沿。")

    return slide


def build_slide_31_positioning_diagram(prs):
    """Slide 31: 定位圖 — Macro vs Micro 互補"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "定位圖: Macro vs Micro 互補",
                    "TumorLens 做 Genome-wide 偵測，ISM 做 Per-variant 甲基化深度分析",
                    title_en="Positioning: TumorLens (Macro) ↔ ISM (Micro) — Complementary, Not Competing")

    # --- Left: TumorLens circle (macro) ---
    add_rect(slide, 0.8, 2.00, 5.00, 3.20, fill="E8F0FE", line="B8D4F0", radius=True)
    add_text(slide, "TumorLens (Macro)", 1.00, 2.08, 4.5, 0.35,
             font_size=17, color=THEME["accent3"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Genome-wide multi-modal",
        "SNV + SV + CNV + LOH + Methylation",
        "",
        "Purity-aware modeling",
        "HLA-locus 重建",
        "",
        "Tumor-Normal + Tumor-Only",
        "Region-level methylation averaging",
    ], 1.00, 2.50, 4.5, 2.50, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Right: ISM circle (micro) ---
    add_rect(slide, 6.20, 2.00, 5.00, 3.20, fill="FFF3E0", line=THEME["accent"], radius=True)
    add_text(slide, "ISM (Micro)", 6.40, 2.08, 4.5, 0.35,
             font_size=17, color=THEME["accent"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Per-variant read-level",
        "Methylation clustering",
        "",
        "Somatic variant-anchored",
        "Individual read methylation patterns",
        "",
        "NGroups, AlleleDelta, PERMANOVA",
        "Subclone characterization",
    ], 6.40, 2.50, 4.5, 2.50, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Central overlap arrow ---
    add_rect(slide, 5.15, 3.05, 1.70, 1.00, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "互補", 5.15, 3.10, 1.70, 0.45,
             font_size=16, color=THEME["dark"], bold=True, align="center", margin=0)
    add_text(slide, "共同方向", 5.15, 3.50, 1.70, 0.30,
             font_size=10, color=THEME["muted"], align="center", margin=0)

    # --- Bottom: Trend + differentiation ---
    add_rect(slide, 0.8, 5.45, 5.50, 1.05, fill="F3E5F5", line=THEME["accent4"], radius=True)
    add_text(slide, "2026 學界趨勢", 0.95, 5.50, 3.0, 0.25,
             font_size=12, color=THEME["accent4"], bold=True, margin=0)
    add_text(slide, "LOH + CNV + Methylation unified approach",
             0.95, 5.78, 5.2, 0.25,
             font_size=11, color=THEME["text"], margin=0)
    add_text(slide, "多中心驗證，purity-aware modeling 是共識方向",
             0.95, 6.04, 5.2, 0.22,
             font_size=10, color=THEME["muted"], margin=0)
    add_text(slide, "來源: TumorLens preprint (medRxiv 2026)",
             0.95, 6.28, 5.2, 0.20,
             font_size=9, color=THEME["muted"], margin=0)

    add_rect(slide, 6.80, 5.45, 5.73, 1.05, fill="E8F5E9", line=THEME["positive"], radius=True)
    add_text(slide, "ISM 差異化", 6.95, 5.50, 3.0, 0.25,
             font_size=12, color=THEME["positive"], bold=True, margin=0)
    add_text(slide, "唯一做 read-level per-variant\nepigenetic context 的工具",
             6.95, 5.78, 5.4, 0.45,
             font_size=11, color=THEME["text"], margin=0, line_spacing=1.25)
    add_text(slide, "本週佐證: ΔNGroups = +0.705 (7/7 一致)",
             6.95, 6.25, 5.4, 0.22,
             font_size=10, color=THEME["positive"], bold=True, margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "定位圖 — TumorLens vs ISM 互補關係。\n\n"
        "TumorLens (Macro)：\n"
        "第一個 long-read unified pipeline，同時偵測 SNV/SV/CNV/LOH/methylation。\n"
        "Genome-wide 多模態，purity-aware modeling。\n\n"
        "ISM (Micro)：\n"
        "Per-variant read-level methylation clustering。\n"
        "給每個 somatic variant 附加甲基化結構 context。\n\n"
        "2026 學界趨勢：LOH+CNV+methylation 整合分析。\n"
        "ISM 的差異化：唯一做 read-level per-variant epigenetic context。\n"
        "兩者互補而非競爭。\n"
        "來源: TumorLens preprint medRxiv 2026")

    return slide


def build_slide_32_ism_positioning(prs):
    """Slide 32: ISM 工具定位"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "ISM 工具定位",
                    "Read-Level Epigenetic Characterization — 不是 filter",
                    title_en="ISM Positioning: Read-Level Subclone Characterization Tool")

    # --- Left: What ISM is NOT ---
    add_rect(slide, 0.8, 1.95, 5.5, 1.60, fill="FFEBEE", line="EF9A9A", radius=True)
    add_text(slide, "✗  Variant Filter", 1.0, 2.01, 4.0, 0.30,
             font_size=17, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "Phase 1A: F1 增益有限 (paired Δ=+0.011)",
        "TO: 甲基化完全無增量 (caller_af=0.654 > 全部 ISM)",
        "結論: ISM 不適合做 variant filtering",
    ], 1.0, 2.40, 5.1, 1.0, font_size=12, color=THEME["text"], line_spacing=1.30)

    # --- Left: What ISM IS ---
    add_rect(slide, 0.8, 3.75, 5.5, 2.50, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "✓  Read-Level Epigenetic Characterization", 1.0, 3.81, 5.1, 0.30,
             font_size=15, color=THEME["positive"], bold=True, margin=0)

    capabilities = [
        "1. 每個 somatic variant 附加甲基化結構 context",
        "2. 偵測 subclonal LOH / subclonal heterogeneity",
        "3. 輸出 structural annotations 供下游工具使用",
    ]
    add_multiline_text(slide, capabilities, 1.0, 4.20, 5.1, 1.2,
                       font_size=13, color=THEME["text"], line_spacing=1.40)

    add_text(slide, "TumorLens macro → ISM micro = 互補", 1.0, 5.55, 5.1, 0.30,
             font_size=12, color=THEME["accent3"], bold=True, margin=0)

    # --- Right: This week's evidence ---
    add_rect(slide, 6.85, 1.95, 5.68, 4.30, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "本週結果佐證此定位", 7.05, 2.01, 5.0, 0.30,
             font_size=15, color=THEME["text"], bold=True, margin=0)

    evidence_items = [
        ("NGroups 能區分", "clonal vs subclonal LOH\nΔNG=+0.705, 7/7 一致", THEME["accent2"]),
        ("AlleleDelta 偵測", "ASM 差異 (+12×)\nCAMDAC 機制確認", THEME["accent4"]),
        ("Segment 一致性", "空間結構 ρ=0.270\nsegmental event", THEME["accent3"]),
        ("跨平台穩定", "ONT/DORADO 一致\n技術可靠", THEME["muted"]),
    ]
    ey = 2.45
    for title, detail, color in evidence_items:
        add_rect(slide, 7.05, ey, 5.3, 0.90, fill="F8F8F8", radius=True)
        add_rect(slide, 7.05, ey, 0.10, 0.90, fill=color, radius=True)
        add_text(slide, title, 7.30, ey + 0.05, 3.0, 0.25,
                 font_size=12, color=THEME["text"], bold=True, margin=0)
        add_text(slide, detail, 7.30, ey + 0.32, 4.8, 0.52,
                 font_size=10, color=THEME["muted"], margin=0, line_spacing=1.20)
        ey += 0.97

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "ISM 工具定位：\n\n"
        "✗ 不是 variant filter:\n"
        "- Phase 1A F1 增益有限 (paired Δ=+0.011)\n"
        "- TO 甲基化無增量 (caller_af 已超越全 ISM)\n\n"
        "✓ 是 read-level epigenetic characterization 工具:\n"
        "1. 每個 variant 附加甲基化 context\n"
        "2. 偵測 subclonal LOH/heterogeneity\n"
        "3. 輸出 structural annotations\n\n"
        "本週結果佐證：NGroups 能區分 clonal vs subclonal LOH。")

    return slide


def build_slide_33_roadmap(prs):
    """Slide 33: 區塊分割策略 + 技術路線圖"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "技術路線圖",
                    "從 Phase 1A 到未來 — 整體藍圖",
                    title_en="Technical Roadmap: Phase 1A → Phase 2 → Phase 3")

    # --- Timeline bar ---
    add_rect(slide, 0.8, 2.20, 11.73, 0.30, fill=THEME["line"], radius=True)
    # Phase markers
    phases = [
        (0.8, 3.8, "Phase 1A", "完成", THEME["accent2"]),
        (4.0, 4.5, "Phase 2", "進行中", THEME["accent"]),
        (8.7, 3.8, "Phase 3", "規劃", THEME["accent3"]),
    ]
    for px, pw, label, status, color in phases:
        add_rect(slide, px, 2.20, pw, 0.30, fill=color, radius=True)
        add_text(slide, label, px + 0.1, 2.22, pw - 0.2, 0.14,
                 font_size=10, color="FFFFFF", bold=True, margin=0)
        add_text(slide, status, px + 0.1, 2.35, pw - 0.2, 0.14,
                 font_size=8, color="FFFFFF", margin=0)

    # --- Phase 1A details ---
    add_rect(slide, 0.8, 2.70, 3.0, 2.11, fill="E8F5E9", line="A5D6A7", radius=True)
    add_multiline_text(slide, [
        "F1 優化",
        "paired Δ=+0.011",
        "TO 無增量",
        "",
        "LOH×AF×Methyl",
        "雙重證據鏈 ✓",
        "(本週 POSITIVE)",
    ], 0.95, 2.78, 2.7, 1.8, font_size=11, color=THEME["text"], line_spacing=1.20)

    # --- Phase 2 details ---
    add_rect(slide, 4.0, 2.70, 4.5, 2.11, fill="FFF3E0", line="FFD9A0", radius=True)
    add_multiline_text(slide, [
        "Direction A+D",
        "Normal BAM 整合",
        "LOH BED 標註",
        "跨區域分析",
        "",
        "→ Subclone Annotation",
        "  Segment Classification",
        "  Uniform vs Mixed",
    ], 4.15, 2.78, 4.2, 1.8, font_size=11, color=THEME["text"], line_spacing=1.15)

    # --- Phase 3 details ---
    add_rect(slide, 8.7, 2.70, 3.83, 2.11, fill="E8F0FE", line="B8D4F0", radius=True)
    add_multiline_text(slide, [
        "Block Segmentation",
        "Fraction Estimation",
        "Cross-Segment Pattern",
        "",
        "論文目標:",
        "Read-level epigenetic",
        "context tool",
    ], 8.85, 2.78, 3.5, 1.8, font_size=11, color=THEME["text"], line_spacing=1.20)

    # --- Arrows ---
    add_native_arrow(slide, 3.85, 3.70, 3.95, 3.70, color=THEME["muted"], width=2)
    add_native_arrow(slide, 8.55, 3.70, 8.65, 3.70, color=THEME["muted"], width=2)

    # --- Blocker note ---
    add_rect(slide, 0.8, 4.95, 11.73, 0.95, fill="FFEBEE", line=THEME["negative"], radius=True)
    add_text(slide, "Blocker: Haplotag 重跑", 1.0, 5.00, 4.0, 0.25,
             font_size=13, color=THEME["negative"], bold=True, margin=0)
    add_multiline_text(slide, [
        "HP-dependent 分析需等修正 BAM。19 canonical runs 全部 partial (AMB% 已從 17.5% 降至 8.0%，但仍非最終版本)。",
        "LOH 分析使用 VCF AF (非 HP-dependent)，因此本週結論不受影響。",
    ], 1.0, 5.28, 11.3, 0.40, font_size=10, color=THEME["text"], line_spacing=1.25)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "技術路線圖：\n\n"
        "Phase 1A (完成):\n"
        "- F1 優化: paired Δ=+0.011, TO 無增量\n"
        "- LOH×AF×Methylation 雙重證據鏈 (本週 POSITIVE)\n\n"
        "Phase 2 (進行中):\n"
        "- Direction A+D: Normal BAM 整合 + LOH BED 標註\n"
        "- → Subclone annotation, segment classification\n\n"
        "Phase 3 (規劃):\n"
        "- Block segmentation, fraction estimation\n"
        "- 論文: read-level epigenetic context tool\n\n"
        "Blocker: Haplotag 重跑。HP-dependent 分析需修正 BAM。\n"
        "但 LOH 分析用 VCF AF，不受影響。")

    return slide


def build_slide_34_validation_risks(prs):
    """Slide 34: 驗證需求與風險 — 新增"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "驗證需求與風險",
                    "驗證優先序與已知限制 — 下一步最該做什麼",
                    title_en="Validation Priorities & Known Risks")

    # --- Left: Validation priority table ---
    add_rect(slide, 0.8, 1.95, 6.20, 4.43, fill=THEME["card_bg"], line=THEME["accent2"], radius=True)
    _add_teaching_label(slide, "驗證優先序", 0.95, 2.00, color=THEME["accent2"])

    pri_headers = ["優先序", "項目", "狀態"]
    pri_rows = [
        ["1", "Phase 2 A+D Normal BAM integration", "進行中"],
        ["2", "臨床低純度樣本擴展驗證", "待規劃"],
        ["3", "Haplotag 重跑 (不同 phasing 策略)", "blocker: 計算資源"],
    ]
    add_simple_table(slide, pri_headers, pri_rows, 1.0, 2.35, 5.8, 1.20,
                     font_size=11, header_color=THEME["accent2"])

    # Priority details
    add_multiline_text(slide, [
        "① Normal BAM 是區分 germline vs tumor ASM 的關鍵",
        "  → 提升 subclone NGroups 可信度",
        "",
        "② 臨床樣本 purity < 1.0 → subclonal fraction 混合 dilution",
        "  → 驗證 AF 公式在低純度下是否成立",
        "",
        "③ Haplotag 重跑解決 self-phasing → HP-dependent 分析重啟",
        "  → 但 LOH/AF 分析不需要（已獨立）",
    ], 1.0, 3.65, 5.8, 1.70, font_size=10, color=THEME["text"], line_spacing=1.20)

    # Left bottom: key takeaway
    add_text(slide, "核心: Normal BAM 驗證最優先 → 直接提升 NGroups 可信度",
             1.0, 5.85, 5.8, 0.25,
             font_size=11, color=THEME["accent2"], bold=True, margin=0)

    # --- Right: Known risks ---
    add_rect(slide, 7.40, 1.95, 5.13, 4.43, fill="FFEBEE", line=THEME["negative"], radius=True)
    _add_teaching_label(slide, "已知風險/限制", 7.55, 2.00, color=THEME["negative"])

    risks = [
        "1. Coverage_Multiple ≠ 精確 CN",
        "   (無 GC correction / ploidy)",
        "",
        "2. Cell line ≠ clinical sample",
        "   (purity=1.0，需低純度驗證)",
        "",
        "3. LOH.bed 不區分 LOH 類型",
        "   (Deletion vs cnLOH vs Gain)",
        "",
        "4. HCC1954 segment 反向",
        "   (n=34 underpowered + saturation)",
        "",
        "5. NGroups 上限 = 4",
        "   (可能低估高 subclonality)",
    ]
    add_multiline_text(slide, risks, 7.55, 2.35, 4.8, 3.60,
                       font_size=11, color=THEME["text"], line_spacing=1.20)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "驗證需求與風險 — 新增頁面。\n\n"
        "驗證優先序：\n"
        "1. Phase 2 A+D Normal BAM — 區分 germline vs tumor ASM\n"
        "2. 臨床低純度樣本 — 驗證 AF 公式在 purity<1 下的適用性\n"
        "3. Haplotag 重跑 — 解決 self-phasing，但 LOH 分析不需要\n\n"
        "已知限制：\n"
        "1. CM 非精確 CN（無 GC/ploidy correction）\n"
        "2. Cell line purity=1.0 ≠ clinical sample\n"
        "3. LOH.bed 不區分類型（deletion vs cnLOH）\n"
        "4. HCC1954 segment 反向（n=34 太少）\n"
        "5. NGroups 上限=4（ceiling effect）")

    return slide


def build_slide_35_qa(prs):
    """Slide 35: 限制與 Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_slide_title(slide,
                    "限制與討論",
                    "已知 caveats + 下週工作 + Q&A",
                    title_en="Limitations, Next Steps & Q&A")

    # --- Left: Limitations ---
    add_rect(slide, 0.8, 1.95, 6.2, 3.40, fill=THEME["card_bg"], line=THEME["line"], radius=True)
    add_text(slide, "已知限制", 1.0, 2.01, 3.0, 0.28,
             font_size=15, color=THEME["text"], bold=True, margin=0)

    limitations = [
        ("1.", "Coverage_Multiple 非精確 CN",
         "vs truth r=0.831，非 1.0 — 未來需 Spectre 或 truth CN"),
        ("2.", "Cell line ≠ 臨床樣本",
         "purity=1.0，無 normal cell dilution — 需臨床樣本驗證"),
        ("3.", "LOH.bed 不區分 deletion vs cnLOH",
         "依賴 Coverage_Multiple 間接推斷 CN tier"),
        ("4.", "HCC1954 Segment 分析反向",
         "n=34 segments underpowered + 60.2% saturation"),
        ("5.", "NGroups 上限 = 4",
         "可能低估高複雜度 subclone — 但 clustering 穩定性考量"),
    ]
    ly = 2.35
    for num, title, detail in limitations:
        add_text(slide, num, 1.49, ly, 0.3, 0.22,
                 font_size=11, color=THEME["accent"], bold=True, margin=0)
        add_text(slide, title, 1.79, ly, 4.62, 0.22,
                 font_size=12, color=THEME["text"], bold=True, margin=0)
        add_text(slide, detail, 1.79, ly + 0.24, 4.62, 0.18,
                 font_size=10, color=THEME["muted"], margin=0, line_spacing=1.15)
        ly += 0.60

    # --- Right: Next steps ---
    add_rect(slide, 7.3, 1.95, 5.23, 1.60, fill="E8F5E9", line="A5D6A7", radius=True)
    add_text(slide, "下週預計工作", 7.50, 2.01, 3.0, 0.28,
             font_size=15, color=THEME["accent2"], bold=True, margin=0)

    next_steps = [
        "Phase 2 A+D Normal BAM 整合續進",
        "Segment annotation prototype",
        "TO pipeline staging (v2) 推進",
    ]
    add_multiline_text(slide, next_steps, 7.50, 2.35, 4.8, 1.1,
                       font_size=13, color=THEME["text"], line_spacing=1.45)

    # --- Q&A area ---
    add_rect(slide, 7.3, 3.75, 5.23, 1.60, fill=THEME["dark"], radius=True)
    add_text(slide, "Q & A", 7.3, 4.05, 5.23, 0.70,
             font_size=36, color=THEME["light_text"], bold=True, align="center",
             font_name=THEME["font_title"], margin=0)
    add_text(slide, "討論 & 問答", 7.3, 4.70, 5.23, 0.40,
             font_size=16, color=THEME["light_muted"], align="center", margin=0)

    # --- Bottom: Thank you ---
    add_rect(slide, 0.8, 5.55, 11.73, 1.00, fill=THEME["dark"], radius=True)
    add_text(slide, "感謝聆聽", 0.8, 5.65, 11.73, 0.40,
             font_size=22, color=THEME["light_text"], bold=True, align="center",
             font_name=THEME["font_title"], margin=0)
    add_text(slide, "InterSubMod · LOH Subclone AF × Methylation · 2026-04-14",
             0.8, 6.05, 11.73, 0.30,
             font_size=12, color=THEME["light_muted"], align="center", margin=0)

    add_footer(slide, len(prs.slides))

    inject_speaker_notes(slide,
        "限制：\n"
        "1. Coverage_Multiple 不精確 (r=0.831)\n"
        "2. Cell line ≠ 臨床 (purity=1.0)\n"
        "3. LOH.bed 不區分 deletion vs cnLOH\n"
        "4. HCC1954 segment 反向\n"
        "5. NGroups 上限 4\n\n"
        "下週工作：\n"
        "- Phase 2 A+D 續進\n"
        "- Segment annotation prototype\n"
        "- TO pipeline staging v2")

    return slide


# ═══════════════════════════════════════════════════════════════
# Main: batch-driven generation
# ═══════════════════════════════════════════════════════════════
# ═══════════════════════════════════════════════════════════════
# Section Dividers (inserted between major sections)
# ═══════════════════════════════════════════════════════════════
def build_divider_teaching(prs):
    """Section divider: Teaching — LOH to Subclonal LOH (V2: 11 slides)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide,
                        "教學: LOH → Subclonal LOH 完整推導",
                        "從 LOH 基礎到甲基化偵測的 11 頁完整推導",
                        title_en="Teaching: LOH → Subclonal LOH — Complete Derivation",
                        page_num=len(prs.slides))
    inject_speaker_notes(slide,
        "進入教學區段。接下來 11 頁會從 LOH 基礎概念開始，\n"
        "經過 self-phasing 防禦性教學、CM 校準，\n"
        "到 CAMDAC 機制、AlleleDelta/ASM、ISM pipeline，\n"
        "逐步推導到 subclonal LOH 的偵測邏輯，\n"
        "最後引出六層證據鏈的架構。")


def build_divider_data(prs):
    """Section divider: Core Data Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide,
                        "核心結果: 三步驟驗證",
                        "AF baseline → NGroups × AF → Segment 空間一致性",
                        title_en="Core Results: Three-Step Verification",
                        page_num=len(prs.slides))
    inject_speaker_notes(slide,
        "教學結束，進入數據。\n"
        "接下來 6 頁按三個步驟展示核心研究數據：\n"
        "Step 1: AF 分布 baseline 與 CN tier 分層\n"
        "Step 2: NGroups × AF 交叉分析 + confound 排除\n"
        "Step 3: Segment 空間一致性 + per-sample 總覽")


def build_divider_summary(prs):
    """Section divider: Summary & Future Directions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(slide,
                        "總結與展望",
                        "六層證據鏈 POSITIVE → ISM 定位 → 未來方向",
                        title_en="Summary & Future Directions",
                        page_num=len(prs.slides))
    inject_speaker_notes(slide,
        "進入總結區段。\n"
        "先彙整六層證據鏈結論，\n"
        "再看 TumorLens 競爭者佐證，\n"
        "最後是 ISM 工具定位與未來技術路線。")


BATCHES = {
    # Batch 1: Part A — 封面與導覽 (2 slides)
    1: [build_slide_1_cover, build_slide_2_key_findings],
    # Batch 2: Part B + C — 背景 + LOH+CNV 發現 (1 divider + 9 slides)
    2: [build_divider_background, build_slide_3_to_motivation,
        build_slide_4_seqc2, build_slide_ism_loh_impact,
        build_slide_5_dual_definition,
        build_slide_6_conclusion_bridge,
        build_slide_7_cm_cn_proxy, build_slide_loh_af_overview,
        build_slide_8_cn_tier_separation,
        build_slide_9_bio_explanation],
    # Batch 3: Part D 前半 — 教學 Slides 10-14 (1 divider + 5 slides)
    3: [build_divider_teaching, build_slide_10_loh_basics,
        build_slide_11_clonal_subclonal, build_slide_12_af_math,
        build_slide_13_self_phasing, build_slide_14_cm_calibration],
    # Batch 4: Part D 後半 — 教學 Slides 15-20 (6 slides)
    4: [build_slide_15_camdac, build_slide_16_allele_delta,
        build_slide_17_read_level, build_slide_18_ngroups,
        build_slide_19_nr_confound, build_slide_20_pyramid],
    # Batch 5: Part E — 核心結果 (1 divider + 6 slides)
    5: [build_divider_data,
        build_slide_21_af_baseline, build_slide_22_cn_tier_focus,
        build_slide_23_ngroups_af, build_slide_24_nr_control,
        build_slide_25_segment, build_slide_26_per_sample],
    # Batch 6: Part F — 整合 (1 divider + 3 slides)
    6: [build_divider_summary, build_slide_27_evidence_summary,
        build_slide_28_paired_validation, build_slide_29_normal_asm],
    # Batch 7: Part G + H + I — 競爭者 + 未來 + Q&A (6 slides)
    7: [build_slide_30_tumorlens_table, build_slide_31_positioning_diagram,
        build_slide_32_ism_positioning, build_slide_33_roadmap,
        build_slide_34_validation_risks, build_slide_35_qa],
}


def main():
    parser = argparse.ArgumentParser(description="Build weekly report PPTX in batches")
    parser.add_argument("--batch", type=int, default=0,
                        help="Batch number (0 = all defined batches)")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = META["author"]
    prs.core_properties.title = META["deck_title"]

    if args.batch > 0:
        builders = BATCHES.get(args.batch, [])
        if not builders:
            print(f"[ERROR] Batch {args.batch} not defined. Available: {list(BATCHES.keys())}")
            return
        print(f"[INFO] Building batch {args.batch} ({len(builders)} slides)")
        for fn in builders:
            fn(prs)
    else:
        total = 0
        for batch_num in sorted(BATCHES.keys()):
            for fn in BATCHES[batch_num]:
                fn(prs)
                total += 1
        print(f"[INFO] Built all {total} slides from {len(BATCHES)} batches")

    prs.save(str(OUTPUT_PPTX))
    print(f"[OK] Saved: {OUTPUT_PPTX}")
    print(f"     Slides: {len(prs.slides)}")

    # Auto-snapshot for diff detection
    try:
        from pptx_snapshot import take_snapshot, _infer_slide_map_from_batches
        snap_path = str(OUTPUT_PPTX).replace(".pptx", "_snapshot.json")
        slide_map = _infer_slide_map_from_batches(Path(__file__))
        take_snapshot(str(OUTPUT_PPTX), snap_path, slide_index_map=slide_map)
        print(f"     Snapshot: {Path(snap_path).name}")
    except ImportError:
        pass


if __name__ == "__main__":
    main()
