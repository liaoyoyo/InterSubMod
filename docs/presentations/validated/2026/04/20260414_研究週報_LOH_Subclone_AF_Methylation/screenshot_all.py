#!/usr/bin/env python3
"""Generate wireframe screenshots of all PPTX slides with CJK support.

Renders all shapes (rectangles, text boxes, images, connectors) from python-pptx
onto PIL canvases. Outputs PNG files to screenshots/ directory.
"""

import os, sys, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

# --- Config ---
PPTX_PATH = "20260414_LOH_Subclone_AF_Methylation.pptx"
OUT_DIR = Path("screenshots")
SCALE = 150  # DPI-ish scale factor
SLIDE_W_INCH = 13.333
SLIDE_H_INCH = 7.5
W_PX = int(SLIDE_W_INCH * SCALE)
H_PX = int(SLIDE_H_INCH * SCALE)

# Font paths
CJK_FONT = "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"
MONO_FONT = "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf"

def _font(size, bold=False, mono=False):
    path = MONO_FONT if mono else CJK_FONT
    try:
        f = ImageFont.truetype(path, size)
        return f
    except:
        return ImageFont.load_default()

def emu_to_px(emu):
    return int(emu / 914400 * SCALE)

def _color_tuple(rgb_obj):
    """Extract (R,G,B) from various pptx color objects."""
    if rgb_obj is None:
        return None
    try:
        # RGBColor object
        return (rgb_obj[0], rgb_obj[1], rgb_obj[2])
    except:
        pass
    try:
        s = str(rgb_obj)
        if len(s) == 6:
            return (int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
    except:
        pass
    return None

def _get_fill_color(shape):
    """Try to get fill color from shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return _color_tuple(fc.rgb)
    except:
        pass
    return None

def _get_shape_bounds(shape):
    """Get (x, y, w, h) in pixels."""
    return (emu_to_px(shape.left), emu_to_px(shape.top),
            emu_to_px(shape.width), emu_to_px(shape.height))

def _extract_text(shape):
    """Extract all text from a shape's text_frame."""
    lines = []
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            line = ""
            for run in para.runs:
                line += run.text
            if line.strip():
                lines.append(line.strip())
    except:
        pass
    return lines

def _get_text_props(shape):
    """Get font size and bold from first run."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                bold = run.font.bold
                return (int(sz / 12700) if sz else 10, bool(bold))
    except:
        pass
    return (10, False)

def _wrap_text(text, font, max_w, draw):
    """Simple word/char wrapping."""
    if not text:
        return []
    # Try full line first
    bbox = draw.textbbox((0,0), text, font=font)
    if (bbox[2] - bbox[0]) <= max_w:
        return [text]

    # Wrap by character for CJK-heavy text
    lines = []
    current = ""
    for ch in text:
        test = current + ch
        bbox = draw.textbbox((0,0), test, font=font)
        if (bbox[2] - bbox[0]) > max_w and current:
            lines.append(current)
            current = ch
        else:
            current = test
    if current:
        lines.append(current)
    return lines

def render_slide(slide, idx, out_dir):
    """Render one slide to PNG."""
    img = Image.new("RGB", (W_PX, H_PX), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Draw border
    draw.rectangle([0, 0, W_PX-1, H_PX-1], outline=(200, 200, 200), width=2)

    # Page number
    pg_font = _font(18, bold=True)
    draw.text((10, 5), f"P{idx+1}", fill=(100,100,100), font=pg_font)

    for shape in slide.shapes:
        x, y, w, h = _get_shape_bounds(shape)

        # Skip tiny shapes
        if w < 3 and h < 3:
            continue

        # Get fill color
        fill_color = _get_fill_color(shape)

        # --- Image placeholder ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            draw.rectangle([x, y, x+w, y+h], outline=(100,150,200), width=2, fill=(230,240,250))
            img_font = _font(12)
            draw.text((x+4, y+4), "[IMG]", fill=(100,150,200), font=img_font)
            continue

        # --- Connector / Line ---
        if hasattr(shape, 'begin_x') or shape.shape_type in (MSO_SHAPE_TYPE.LINE,):
            draw.line([x, y, x+w, y+h], fill=(150,150,150), width=2)
            continue

        # --- Rectangle / TextBox / AutoShape ---
        has_text = bool(_extract_text(shape))

        # Draw fill
        if fill_color:
            draw.rectangle([x, y, x+w, y+h], fill=fill_color, outline=None)
            # Add subtle border
            brightness = sum(fill_color) / 3
            border_c = tuple(max(0, c - 40) for c in fill_color)
            draw.rectangle([x, y, x+w, y+h], outline=border_c, width=1)
        elif has_text and w > 20 and h > 20:
            # Text box without fill - light dashed outline
            draw.rectangle([x, y, x+w, y+h], outline=(220, 220, 220), width=1)

        # Draw text
        if has_text:
            lines = _extract_text(shape)
            pt_size, bold = _get_text_props(shape)
            # Scale font size
            font_px = max(8, int(pt_size * SCALE / 96))
            font_px = min(font_px, 28)  # cap

            is_mono = False
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and 'consol' in run.font.name.lower():
                            is_mono = True
                            break
            except:
                pass

            txt_font = _font(font_px, bold=bold, mono=is_mono)

            # Determine text color
            txt_color = (0, 0, 0)
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            txt_color = _color_tuple(run.font.color.rgb) or (0,0,0)
                            break
                    break
            except:
                pass

            # If dark fill, use white text
            if fill_color and sum(fill_color)/3 < 100:
                txt_color = (255, 255, 255)

            # Render text with wrapping
            pad = 4
            max_text_w = w - pad * 2
            ty = y + pad
            for line in lines:
                wrapped = _wrap_text(line, txt_font, max_text_w, draw)
                for wl in wrapped:
                    if ty + font_px > y + h:
                        # OVERFLOW indicator
                        of_font = _font(10, bold=True)
                        draw.text((x + w - 60, y + h - 14), "⚠OVERFLOW", fill=(255, 0, 0), font=of_font)
                        break
                    draw.text((x + pad, ty), wl, fill=txt_color, font=txt_font)
                    ty += font_px + 2
                else:
                    continue
                break

    # Save
    out_path = out_dir / f"P{idx+1:02d}.png"
    img.save(str(out_path), "PNG")
    return str(out_path)

def main():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation(PPTX_PATH)
    total = len(prs.slides)
    print(f"Rendering {total} slides...")

    for i, slide in enumerate(prs.slides):
        path = render_slide(slide, i, OUT_DIR)
        print(f"  P{i+1:02d} -> {path}")

    print(f"\nDone. {total} screenshots in {OUT_DIR}/")

if __name__ == "__main__":
    main()
