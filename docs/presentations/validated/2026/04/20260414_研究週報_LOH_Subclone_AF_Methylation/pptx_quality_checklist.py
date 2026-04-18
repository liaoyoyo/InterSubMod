#!/usr/bin/env python3
"""
PPTX Quality Checklist — Automated quality checks for weekly report PPTX.

Usage:
  python3 pptx_quality_checklist.py <pptx_file>
  python3 pptx_quality_checklist.py  # defaults to *.pptx in current dir

Checks C01-C14 based on memory rules and rendering troubleshoot handbook.
"""
from __future__ import annotations

import io
import sys
import unicodedata
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════
ALLOWED_FONTS = {"Arial", "Arial Bold", "Consolas"}
MAX_CN_TITLE_CHARS = 15
MAX_FOCUS_POINTS = 3
FOOTER_SAFETY_Y = Inches(6.68)
SLIDE_HEIGHT = Inches(7.5)
MIN_CONNECTOR_EMU = 36576
MAX_IMAGE_WIDTH = int(1920 * 1.1)  # 2112
MIN_NOTES_LEN = 20

# Namespaces for XML parsing
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _count_display_chars(text: str) -> int:
    """Count display characters (CJK = 1 char each, ignore whitespace)."""
    count = 0
    for ch in text.strip():
        if ch.isspace():
            continue
        count += 1
    return count


def _is_cjk(ch: str) -> bool:
    """Check if character is CJK."""
    try:
        name = unicodedata.name(ch, "")
        return "CJK" in name or "HANGUL" in name or "HIRAGANA" in name or "KATAKANA" in name
    except ValueError:
        return False


def _count_cjk_chars(text: str) -> int:
    """Count CJK characters in text."""
    return sum(1 for ch in text if _is_cjk(ch))


def _has_cjk(text: str) -> bool:
    """Check if text contains any CJK characters."""
    return any(_is_cjk(ch) for ch in text)


def _is_english_text(text: str) -> bool:
    """Heuristic: text is primarily English if <20% CJK chars."""
    if not text.strip():
        return False
    cjk = _count_cjk_chars(text)
    total = len(text.strip())
    return total > 3 and cjk / total < 0.2


class Issue:
    """Single quality issue."""
    def __init__(self, check_id: str, level: str, slide_num: int, message: str):
        self.check_id = check_id
        self.level = level  # FAIL, WARN, INFO
        self.slide_num = slide_num
        self.message = message

    def __repr__(self):
        icon = {"FAIL": "\u274c", "WARN": "\u26a0", "INFO": "\u2139\ufe0f"}.get(self.level, "?")
        return f"  {icon} {self.check_id} {self.level}: {self.message}"


def check_c01_fonts(prs: Presentation) -> list[Issue]:
    """C01: All fonts must be in ALLOWED_FONTS."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        bad_fonts = set()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        fname = run.font.name
                        if fname and fname not in ALLOWED_FONTS:
                            bad_fonts.add(fname)
        for f in sorted(bad_fonts):
            issues.append(Issue("C01", "FAIL", idx, f'font "{f}" not in {ALLOWED_FONTS}'))
    return issues


def check_c02_title_length(prs: Presentation) -> list[Issue]:
    """C02: Slide main titles should be <= 15 CJK characters.
    Only checks shapes in the title zone (y < 1.5") with font_size >= 28pt."""
    issues = []
    title_zone_max_y = Inches(1.5)
    min_title_pt = 28 * 12700  # 28pt in EMU
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.top is not None and shape.top > title_zone_max_y:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if not run.font.bold:
                        continue
                    fs = run.font.size
                    if fs and fs >= min_title_pt:
                        text = run.text.strip()
                        if "\n" in text:
                            text = text.split("\n")[0]
                        cjk_count = _count_cjk_chars(text)
                        if cjk_count > MAX_CN_TITLE_CHARS:
                            issues.append(Issue("C02", "WARN", idx,
                                f'title "{text[:30]}..." has {cjk_count} CJK chars (max {MAX_CN_TITLE_CHARS})'))
    return issues


def check_c04_en_translation(prs: Presentation) -> list[Issue]:
    """C04: Each slide should have at least 1 English text element."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        has_english = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text
            if _is_english_text(full_text):
                has_english = True
                break
        if not has_english:
            issues.append(Issue("C04", "WARN", idx, "no English text found on this slide"))
    return issues


def check_c05_notes(prs: Presentation) -> list[Issue]:
    """C05: Each slide should have speaker notes >= MIN_NOTES_LEN chars."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        try:
            ns = slide.notes_slide
            text = ns.notes_text_frame.text.strip() if ns and ns.notes_text_frame else ""
        except Exception:
            text = ""
        if len(text) < MIN_NOTES_LEN:
            issues.append(Issue("C05", "FAIL", idx,
                f"notes too short ({len(text)} chars, min {MIN_NOTES_LEN})"))
    return issues


def check_c06_footer(prs: Presentation) -> list[Issue]:
    """C06: Each slide (except cover) should have a footer element (y > 6.5")."""
    issues = []
    footer_y_min = Inches(6.5)
    for idx, slide in enumerate(prs.slides, 1):
        if idx == 1:  # Skip cover
            continue
        has_footer = False
        for shape in slide.shapes:
            if shape.top and shape.top >= footer_y_min:
                has_footer = True
                break
        if not has_footer:
            issues.append(Issue("C06", "WARN", idx, "no element found in footer zone (y > 6.5\")"))
    return issues


def check_c07_safety_line(prs: Presentation) -> list[Issue]:
    """C07: All non-footer shapes must have bottom <= 6.68".
    Skips cover slide (1) and section dividers (dark bg, few shapes)."""
    issues = []
    footer_y = Inches(6.70)  # Slightly above footer to exclude footer elements
    for idx, slide in enumerate(prs.slides, 1):
        if idx == 1:  # Skip cover page (full-bleed design)
            continue
        # Skip section dividers (dark bg, few shapes)
        if len(slide.shapes) <= 12 and "1E2A44" in slide._element.xml:
            continue
        for shape in slide.shapes:
            if shape.top is None or shape.height is None:
                continue
            bottom = shape.top + shape.height
            # Skip footer elements themselves
            if shape.top >= footer_y:
                continue
            if bottom > FOOTER_SAFETY_Y + Inches(0.1):  # Small tolerance
                bottom_in = bottom / Inches(1)
                issues.append(Issue("C07", "WARN", idx,
                    f"shape bottom at {bottom_in:.2f}\" exceeds safety line 6.68\""))
    return issues


def check_c08_background(prs: Presentation) -> list[Issue]:
    """C08: Each slide should have <p:bg> element."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_xml = slide._element.xml
        if "<p:bg>" not in slide_xml and "<p:bg " not in slide_xml:
            issues.append(Issue("C08", "FAIL", idx, "missing <p:bg> background element"))
    return issues


def check_c09_zero_size(prs: Presentation) -> list[Issue]:
    """C09: No shape should have width or height = 0."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.width is not None and shape.width == 0:
                issues.append(Issue("C09", "FAIL", idx,
                    f"shape '{shape.name}' has width=0"))
            if shape.height is not None and shape.height == 0:
                issues.append(Issue("C09", "FAIL", idx,
                    f"shape '{shape.name}' has height=0"))
    return issues


def check_c10_small_connector(prs: Presentation) -> list[Issue]:
    """C10: Connector shapes must have extent >= 36576 EMU."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            sp_xml = shape._element.xml
            if "cxnSp" not in sp_xml:
                continue
            w = shape.width or 0
            h = shape.height or 0
            if w < MIN_CONNECTOR_EMU and h < MIN_CONNECTOR_EMU:
                issues.append(Issue("C10", "FAIL", idx,
                    f"connector '{shape.name}' too small: {w}x{h} EMU"))
    return issues


def check_c11_rgba_images(prs: Presentation) -> list[Issue]:
    """C11: Embedded images must not be RGBA/PA mode."""
    issues = []
    try:
        pptx_path = prs._part.package._pkg_file
        with zipfile.ZipFile(pptx_path) as zf:
            for name in zf.namelist():
                if name.startswith("ppt/media/") and name.lower().endswith((".png", ".jpg", ".jpeg")):
                    data = zf.read(name)
                    with Image.open(io.BytesIO(data)) as img:
                        if img.mode in ("RGBA", "PA"):
                            issues.append(Issue("C11", "FAIL", 0,
                                f"embedded image {name} is {img.mode} mode"))
    except Exception as e:
        issues.append(Issue("C11", "WARN", 0, f"cannot check images: {e}"))
    return issues


def check_c12_large_images(prs: Presentation) -> list[Issue]:
    """C12: Embedded images width should be <= 1920*1.1 px."""
    issues = []
    try:
        pptx_path = prs._part.package._pkg_file
        with zipfile.ZipFile(pptx_path) as zf:
            for name in zf.namelist():
                if name.startswith("ppt/media/") and name.lower().endswith((".png", ".jpg", ".jpeg")):
                    data = zf.read(name)
                    with Image.open(io.BytesIO(data)) as img:
                        if img.size[0] > MAX_IMAGE_WIDTH:
                            issues.append(Issue("C12", "WARN", 0,
                                f"image {name} width={img.size[0]}px > {MAX_IMAGE_WIDTH}px"))
    except Exception as e:
        issues.append(Issue("C12", "WARN", 0, f"cannot check images: {e}"))
    return issues


def check_c13_italic(prs: Presentation) -> list[Issue]:
    """C13: No italic fonts (except proper nouns in small text)."""
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.italic:
                        text = run.text.strip()[:40]
                        issues.append(Issue("C13", "WARN", idx,
                            f'italic text found: "{text}"'))
    return issues


def check_c14_section_dividers(prs: Presentation) -> list[Issue]:
    """C14: Section dividers should appear every 5-8 content slides."""
    issues = []
    # Detect dark-bg slides as section dividers
    divider_indices = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_xml = slide._element.xml
        # Check for dark background (1E2A44 = our dark theme)
        if "1E2A44" in slide_xml:
            # Check if background fill is dark
            bg_xml = slide_xml
            if "<p:bg>" in bg_xml or "<p:bg " in bg_xml:
                # Simple heuristic: if dark color in bg section and few shapes
                shape_count = len(slide.shapes)
                if shape_count <= 12:  # Dividers typically have fewer shapes
                    divider_indices.append(idx)

    # Check gaps between dividers
    if len(divider_indices) < 2:
        issues.append(Issue("C14", "WARN", 0,
            f"only {len(divider_indices)} section divider(s) found; recommend every 5-8 slides"))
    else:
        for i in range(1, len(divider_indices)):
            gap = divider_indices[i] - divider_indices[i-1]
            if gap > 10:
                issues.append(Issue("C14", "WARN", divider_indices[i],
                    f"gap of {gap} slides since last divider (recommend 5-8)"))

    return issues


def run_all_checks(pptx_path: str) -> dict:
    """Run all C01-C14 checks and return results."""
    prs = Presentation(pptx_path)
    # Store the file path for image checks
    prs._part.package._pkg_file = pptx_path

    checks = [
        ("C01 Fonts",             check_c01_fonts),
        ("C02 Title Length",      check_c02_title_length),
        ("C04 EN Translation",    check_c04_en_translation),
        ("C05 Notes",             check_c05_notes),
        ("C06 Footer",            check_c06_footer),
        ("C07 Safety Line",       check_c07_safety_line),
        ("C08 Background",        check_c08_background),
        ("C09 Zero Size",         check_c09_zero_size),
        ("C10 Small Connector",   check_c10_small_connector),
        ("C11 RGBA Images",       check_c11_rgba_images),
        ("C12 Large Images",      check_c12_large_images),
        ("C13 Italic",            check_c13_italic),
        ("C14 Section Dividers",  check_c14_section_dividers),
    ]

    all_issues: list[Issue] = []
    for name, fn in checks:
        issues = fn(prs)
        all_issues.extend(issues)

    return {
        "pptx_path": pptx_path,
        "slide_count": len(prs.slides),
        "issues": all_issues,
        "checks_run": len(checks),
    }


def print_report(results: dict):
    """Print formatted quality report."""
    issues = results["issues"]
    slide_count = results["slide_count"]

    print(f"\n{'=' * 60}")
    print(f"  PPTX Quality Report")
    print(f"  File: {Path(results['pptx_path']).name}")
    print(f"  Slides: {slide_count}")
    print(f"  Checks: {results['checks_run']}")
    print(f"{'=' * 60}\n")

    # Group by slide
    by_slide: dict[int, list[Issue]] = {}
    for issue in issues:
        by_slide.setdefault(issue.slide_num, []).append(issue)

    # Per-slide report
    fail_count = sum(1 for i in issues if i.level == "FAIL")
    warn_count = sum(1 for i in issues if i.level == "WARN")

    for slide_num in range(1, slide_count + 1):
        slide_issues = by_slide.get(slide_num, [])
        slide_fails = [i for i in slide_issues if i.level == "FAIL"]
        slide_warns = [i for i in slide_issues if i.level == "WARN"]

        if slide_fails:
            status = "\u274c FAIL"
        elif slide_warns:
            status = "\u26a0  WARN"
        else:
            status = "\u2705 PASS"

        print(f"Slide {slide_num:2d}: {status}")
        for issue in slide_issues:
            print(f"  {issue}")

    # Global issues (slide_num = 0)
    global_issues = by_slide.get(0, [])
    if global_issues:
        print(f"\nGlobal:")
        for issue in global_issues:
            print(f"  {issue}")

    # Summary
    pass_slides = slide_count - len(set(i.slide_num for i in issues if i.slide_num > 0))
    print(f"\n{'=' * 60}")
    print(f"  SUMMARY: {pass_slides} PASS / {fail_count} FAIL / {warn_count} WARN")
    if fail_count == 0:
        print(f"  \u2705 ALL CRITICAL CHECKS PASSED")
    else:
        print(f"  \u274c {fail_count} CRITICAL ISSUE(S) NEED FIXING")
    print(f"{'=' * 60}\n")

    return fail_count


def main():
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        # Find *.pptx in current directory
        pptx_files = list(Path(".").glob("*.pptx"))
        pptx_files = [f for f in pptx_files if not f.name.startswith("test_")]
        if not pptx_files:
            print("[ERROR] No .pptx files found in current directory")
            sys.exit(1)
        pptx_path = str(pptx_files[0])
        print(f"[INFO] Using: {pptx_path}")

    if not Path(pptx_path).exists():
        print(f"[ERROR] File not found: {pptx_path}")
        sys.exit(1)

    results = run_all_checks(pptx_path)
    fail_count = print_report(results)
    sys.exit(1 if fail_count > 0 else 0)


if __name__ == "__main__":
    main()
