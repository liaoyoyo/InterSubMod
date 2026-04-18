#!/usr/bin/env python3
"""
PPTX Snapshot & Diff Tool — detect manual edits and map back to build_pptx.py.

Usage:
  python3 pptx_snapshot.py snapshot <file.pptx> [-o snapshot.json]
  python3 pptx_snapshot.py diff <edited.pptx> [--baseline snapshot.json] [--threshold 0.05]
  python3 pptx_snapshot.py diff <edited.pptx> --precise          # threshold=0.01", catches tiny moves
  python3 pptx_snapshot.py diff <edited.pptx> --no-archive       # skip auto-archive
  python3 pptx_snapshot.py inspect <file.pptx> --slide 7

Tracks per shape:
  - Position (left, top), size (width, height), rotation
  - Full text + per-paragraph alignment/spacing + per-run font detail
  - Fill color, line color, font color/size/bold/italic/name
  - Table: dimensions + per-cell text + per-cell fill color
  - Images: content hash for replacement detection

Auto-archive (default on): saves edited snapshot + JSON diff to diff_archive/
"""
from __future__ import annotations

import argparse
import ast
import hashlib
import json
import math
import re
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════
EMU_PER_INCH = 914400
PRECISION = 2  # decimal places for inches

DEFAULT_THRESHOLD = 0.05  # inches — position/size changes below this are ignored
PRECISE_THRESHOLD = 0.01  # inches — for --precise mode (catches almost all moves)
FONT_SIZE_THRESHOLD = 0.5  # pt — font size changes below this are ignored
PROXIMITY_LIMIT = 2.0  # inches — max distance for Tier 2 matching

# Helper call patterns to search in builder functions
HELPER_PATTERNS = [
    "add_text(", "add_rect(", "add_multiline_text(", "add_image_fit(",
    "add_simple_table(", "_add_term_box(", "_add_teaching_label(",
    "_add_step_badge(", "add_section_divider(",
]


# ═══════════════════════════════════════════════════════════════
# Shape property extraction
# ═══════════════════════════════════════════════════════════════
def _emu_to_inches(emu_val):
    """Convert EMU to inches, rounded to PRECISION decimals."""
    if emu_val is None:
        return None
    return round(emu_val / EMU_PER_INCH, PRECISION)


def _get_fill_color(shape):
    """Extract fill color hex string, or None if no solid fill."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc and fc.rgb:
                return str(fc.rgb)
    except Exception:
        pass
    return None


def _get_line_color(shape):
    """Extract line/border color hex string, or None."""
    try:
        ln = shape.line
        if ln and ln.color and ln.color.rgb:
            return str(ln.color.rgb)
    except Exception:
        pass
    return None


def _get_font_props(shape):
    """Extract font properties from the first text run (legacy compat)."""
    font_size = None
    font_bold = None
    font_color = None
    font_name = None
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    font_size = round(run.font.size.pt, 1)
                font_bold = run.font.bold
                if run.font.color and run.font.color.rgb:
                    font_color = str(run.font.color.rgb)
                if run.font.name:
                    font_name = run.font.name
                return font_size, font_bold, font_color, font_name
    except Exception:
        pass
    return font_size, font_bold, font_color, font_name


def _get_paragraphs_detail(shape):
    """Extract per-paragraph, per-run detail for precise comparison.

    Returns list of paragraph dicts, each containing:
      - alignment, line_spacing, space_before, space_after, indent
      - runs: list of {text, font_size, font_bold, font_color, font_name, font_italic}
    """
    paragraphs = []
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            p_info = {
                "alignment": str(para.alignment) if para.alignment else None,
                "level": para.level,
            }
            # Paragraph format properties
            try:
                pf = para.paragraph_format
                if pf.line_spacing is not None:
                    p_info["line_spacing"] = float(pf.line_spacing)
                if pf.space_before is not None:
                    p_info["space_before"] = round(pf.space_before.pt, 1) if hasattr(pf.space_before, 'pt') else float(pf.space_before)
                if pf.space_after is not None:
                    p_info["space_after"] = round(pf.space_after.pt, 1) if hasattr(pf.space_after, 'pt') else float(pf.space_after)
            except Exception:
                pass

            runs = []
            for run in para.runs:
                r_info = {"text": run.text}
                f = run.font
                if f.size is not None:
                    r_info["font_size"] = round(f.size.pt, 1)
                if f.bold is not None:
                    r_info["font_bold"] = f.bold
                if f.italic is not None:
                    r_info["font_italic"] = f.italic
                try:
                    if f.color and f.color.rgb:
                        r_info["font_color"] = str(f.color.rgb)
                except Exception:
                    pass
                if f.name:
                    r_info["font_name"] = f.name
                runs.append(r_info)

            p_info["runs"] = runs
            paragraphs.append(p_info)
    except Exception:
        pass
    return paragraphs if paragraphs else None


def _get_text(shape):
    """Extract all text from shape, joining paragraphs with newline."""
    try:
        return shape.text_frame.text
    except Exception:
        return None


def _get_table_info(shape):
    """Extract table dimensions, cell text, and cell fills."""
    try:
        tbl = shape.table
        rows = len(tbl.rows)
        cols = len(tbl.columns)
        cells = []
        for r in range(rows):
            row_cells = []
            for c in range(cols):
                cell = tbl.cell(r, c)
                cell_info = {"text": cell.text}
                # Cell fill color
                try:
                    fill = cell.fill
                    if fill.type is not None:
                        fc = fill.fore_color
                        if fc and fc.rgb:
                            cell_info["fill"] = str(fc.rgb)
                except Exception:
                    pass
                row_cells.append(cell_info)
            cells.append(row_cells)
        return {"rows": rows, "cols": cols, "cells": cells}
    except Exception:
        return None


def _get_image_hash(shape):
    """Get MD5 hash of image blob for change detection."""
    try:
        blob = shape.image.blob
        return hashlib.md5(blob).hexdigest()
    except Exception:
        return None


def _shape_type_name(shape):
    """Human-readable shape type."""
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "TEXT_BOX"
        elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "AUTO_SHAPE"
        elif st == MSO_SHAPE_TYPE.PICTURE:
            return "PICTURE"
        elif st == MSO_SHAPE_TYPE.TABLE:
            return "TABLE"
        elif st == MSO_SHAPE_TYPE.GROUP:
            return "GROUP"
        elif st == MSO_SHAPE_TYPE.FREEFORM:
            return "FREEFORM"
        elif st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "PLACEHOLDER"
        else:
            return str(st)
    except Exception:
        return "UNKNOWN"


def extract_shape_props(shape):
    """Extract all relevant properties from a single shape."""
    props = {
        "name": shape.name,
        "shape_type": _shape_type_name(shape),
        "left": _emu_to_inches(shape.left),
        "top": _emu_to_inches(shape.top),
        "width": _emu_to_inches(shape.width),
        "height": _emu_to_inches(shape.height),
        "rotation": round(shape.rotation, 1) if shape.rotation else 0,
    }

    # Text and font
    text = _get_text(shape)
    if text is not None:
        props["text"] = text
        fs, fb, fc, fn = _get_font_props(shape)
        props["font_size"] = fs
        props["font_bold"] = fb
        props["font_color"] = fc
        props["font_name"] = fn
        # Per-paragraph, per-run detail for precise comparison
        paras = _get_paragraphs_detail(shape)
        if paras:
            props["paragraphs"] = paras

    # Fill and line
    props["fill_color"] = _get_fill_color(shape)
    props["line_color"] = _get_line_color(shape)

    # Type-specific
    stype = props["shape_type"]
    if stype == "TABLE":
        props["table"] = _get_table_info(shape)
    elif stype == "PICTURE":
        props["image_hash"] = _get_image_hash(shape)

    return props


# ═══════════════════════════════════════════════════════════════
# Snapshot generation
# ═══════════════════════════════════════════════════════════════
def take_snapshot(pptx_path, output_path=None, slide_index_map=None):
    """Read a PPTX and produce a JSON snapshot of all shape properties."""
    prs = Presentation(pptx_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            shapes_data.append(extract_shape_props(shape))

        slide_entry = {
            "index": idx,
            "shapes": shapes_data,
        }
        if slide_index_map and idx in slide_index_map:
            slide_entry["builder_function"] = slide_index_map[idx]

        slides_data.append(slide_entry)

    snapshot = {
        "source_file": Path(pptx_path).name,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "generator": "pptx_snapshot.py",
        "total_slides": len(prs.slides),
        "slides": slides_data,
    }

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(snapshot, f, ensure_ascii=False, indent=2)

    return snapshot


# ═══════════════════════════════════════════════════════════════
# Shape matching (4-tier algorithm)
# ═══════════════════════════════════════════════════════════════
def _l2_distance(s1, s2):
    """L2 distance between shape centers."""
    cx1 = (s1.get("left", 0) or 0) + (s1.get("width", 0) or 0) / 2
    cy1 = (s1.get("top", 0) or 0) + (s1.get("height", 0) or 0) / 2
    cx2 = (s2.get("left", 0) or 0) + (s2.get("width", 0) or 0) / 2
    cy2 = (s2.get("top", 0) or 0) + (s2.get("height", 0) or 0) / 2
    return math.sqrt((cx1 - cx2) ** 2 + (cy1 - cy2) ** 2)


def match_shapes(baseline_shapes, edited_shapes):
    """Match shapes between baseline and edited using 4-tier algorithm.

    Returns list of (baseline_shape, edited_shape, match_tier).
    Unmatched baseline shapes have edited_shape=None.
    Unmatched edited shapes have baseline_shape=None.
    """
    matched = []
    unmatched_base = list(range(len(baseline_shapes)))
    unmatched_edit = list(range(len(edited_shapes)))

    # Tier 1: Exact name match
    for bi in list(unmatched_base):
        for ei in list(unmatched_edit):
            if baseline_shapes[bi]["name"] == edited_shapes[ei]["name"]:
                matched.append((baseline_shapes[bi], edited_shapes[ei], 1))
                unmatched_base.remove(bi)
                unmatched_edit.remove(ei)
                break

    # Tier 2: Same type + proximity
    for bi in list(unmatched_base):
        best_ei = None
        best_dist = PROXIMITY_LIMIT
        for ei in unmatched_edit:
            if baseline_shapes[bi]["shape_type"] == edited_shapes[ei]["shape_type"]:
                d = _l2_distance(baseline_shapes[bi], edited_shapes[ei])
                if d < best_dist:
                    best_dist = d
                    best_ei = ei
        if best_ei is not None:
            matched.append((baseline_shapes[bi], edited_shapes[best_ei], 2))
            unmatched_base.remove(bi)
            unmatched_edit.remove(best_ei)

    # Tier 3: Text content match (first 30 chars)
    for bi in list(unmatched_base):
        b_text = baseline_shapes[bi].get("text", "") or ""
        if not b_text:
            continue
        b_prefix = b_text[:30]
        for ei in list(unmatched_edit):
            e_text = edited_shapes[ei].get("text", "") or ""
            if e_text[:30] == b_prefix and b_prefix:
                matched.append((baseline_shapes[bi], edited_shapes[ei], 3))
                unmatched_base.remove(bi)
                unmatched_edit.remove(ei)
                break

    # Tier 4: Unmatched
    for bi in unmatched_base:
        matched.append((baseline_shapes[bi], None, 4))  # REMOVED
    for ei in unmatched_edit:
        matched.append((None, edited_shapes[ei], 4))  # ADDED

    return matched


# ═══════════════════════════════════════════════════════════════
# Diff computation
# ═══════════════════════════════════════════════════════════════
def _compare_value(old, new, threshold=0):
    """Compare two values with optional numeric threshold."""
    if old is None and new is None:
        return False
    if old is None or new is None:
        return True
    if isinstance(old, (int, float)) and isinstance(new, (int, float)):
        return abs(old - new) > threshold
    return old != new


def diff_shapes(baseline_shape, edited_shape, pos_threshold, font_threshold):
    """Compare two matched shapes and return list of changes."""
    changes = []

    # Position
    for prop in ["left", "top"]:
        old_v = baseline_shape.get(prop)
        new_v = edited_shape.get(prop)
        if _compare_value(old_v, new_v, pos_threshold):
            changes.append({
                "change_type": "MOVED",
                "property": prop,
                "old_value": old_v,
                "new_value": new_v,
                "delta": round((new_v or 0) - (old_v or 0), PRECISION),
            })

    # Size
    for prop in ["width", "height"]:
        old_v = baseline_shape.get(prop)
        new_v = edited_shape.get(prop)
        if _compare_value(old_v, new_v, pos_threshold):
            changes.append({
                "change_type": "RESIZED",
                "property": prop,
                "old_value": old_v,
                "new_value": new_v,
                "delta": round((new_v or 0) - (old_v or 0), PRECISION),
            })

    # Text (full, no truncation — detailed per-run diff is in paragraph section)
    old_text = baseline_shape.get("text")
    new_text = edited_shape.get("text")
    if old_text is not None and new_text is not None and old_text != new_text:
        changes.append({
            "change_type": "TEXT",
            "property": "text",
            "old_value": old_text,
            "new_value": new_text,
        })

    # Font size
    old_fs = baseline_shape.get("font_size")
    new_fs = edited_shape.get("font_size")
    if _compare_value(old_fs, new_fs, font_threshold):
        changes.append({
            "change_type": "FONT",
            "property": "font_size",
            "old_value": old_fs,
            "new_value": new_fs,
        })

    # Font bold
    old_fb = baseline_shape.get("font_bold")
    new_fb = edited_shape.get("font_bold")
    if old_fb is not None and new_fb is not None and old_fb != new_fb:
        changes.append({
            "change_type": "FONT",
            "property": "font_bold",
            "old_value": old_fb,
            "new_value": new_fb,
        })

    # Font color
    old_fc = baseline_shape.get("font_color")
    new_fc = edited_shape.get("font_color")
    if old_fc is not None and new_fc is not None and old_fc != new_fc:
        changes.append({
            "change_type": "COLOR",
            "property": "font_color",
            "old_value": old_fc,
            "new_value": new_fc,
        })

    # Fill color
    old_fill = baseline_shape.get("fill_color")
    new_fill = edited_shape.get("fill_color")
    if _compare_value(old_fill, new_fill):
        changes.append({
            "change_type": "COLOR",
            "property": "fill_color",
            "old_value": old_fill,
            "new_value": new_fill,
        })

    # Line color
    old_line = baseline_shape.get("line_color")
    new_line = edited_shape.get("line_color")
    if _compare_value(old_line, new_line):
        changes.append({
            "change_type": "COLOR",
            "property": "line_color",
            "old_value": old_line,
            "new_value": new_line,
        })

    # Rotation
    old_rot = baseline_shape.get("rotation", 0)
    new_rot = edited_shape.get("rotation", 0)
    if abs((old_rot or 0) - (new_rot or 0)) > 0.5:
        changes.append({
            "change_type": "ROTATED",
            "property": "rotation",
            "old_value": old_rot,
            "new_value": new_rot,
        })

    # Image hash
    old_img = baseline_shape.get("image_hash")
    new_img = edited_shape.get("image_hash")
    if old_img and new_img and old_img != new_img:
        changes.append({
            "change_type": "IMAGE_REPLACED",
            "property": "image_hash",
            "old_value": old_img[:12],
            "new_value": new_img[:12],
        })

    # Font name
    old_fn = baseline_shape.get("font_name")
    new_fn = edited_shape.get("font_name")
    if old_fn and new_fn and old_fn != new_fn:
        changes.append({
            "change_type": "FONT",
            "property": "font_name",
            "old_value": old_fn,
            "new_value": new_fn,
        })

    # Paragraph-level detail diff
    old_paras = baseline_shape.get("paragraphs")
    new_paras = edited_shape.get("paragraphs")
    if old_paras and new_paras:
        para_changes = _diff_paragraphs(old_paras, new_paras)
        changes.extend(para_changes)

    # Table cell-level diff
    old_tbl = baseline_shape.get("table")
    new_tbl = edited_shape.get("table")
    if old_tbl and new_tbl:
        tbl_changes = _diff_table_cells(old_tbl, new_tbl)
        changes.extend(tbl_changes)

    return changes


def _diff_paragraphs(old_paras, new_paras):
    """Compare paragraph-level properties: alignment, spacing, per-run fonts."""
    changes = []
    n = min(len(old_paras), len(new_paras))

    if len(old_paras) != len(new_paras):
        changes.append({
            "change_type": "PARA_COUNT",
            "property": "paragraph_count",
            "old_value": len(old_paras),
            "new_value": len(new_paras),
        })

    for i in range(n):
        op, np_ = old_paras[i], new_paras[i]
        prefix = f"para[{i}]"

        # Alignment
        if op.get("alignment") != np_.get("alignment"):
            changes.append({
                "change_type": "PARA_ALIGN",
                "property": f"{prefix}.alignment",
                "old_value": op.get("alignment"),
                "new_value": np_.get("alignment"),
            })

        # Line spacing
        if op.get("line_spacing") != np_.get("line_spacing"):
            old_ls = op.get("line_spacing")
            new_ls = np_.get("line_spacing")
            if old_ls is not None or new_ls is not None:
                changes.append({
                    "change_type": "PARA_SPACING",
                    "property": f"{prefix}.line_spacing",
                    "old_value": old_ls,
                    "new_value": new_ls,
                })

        # Per-run font changes
        old_runs = op.get("runs", [])
        new_runs = np_.get("runs", [])
        rn = min(len(old_runs), len(new_runs))
        for j in range(rn):
            o_r, n_r = old_runs[j], new_runs[j]
            run_prefix = f"{prefix}.run[{j}]"

            # Run text
            if o_r.get("text") != n_r.get("text"):
                changes.append({
                    "change_type": "RUN_TEXT",
                    "property": f"{run_prefix}.text",
                    "old_value": o_r.get("text", ""),
                    "new_value": n_r.get("text", ""),
                })

            # Run font size
            ofs = o_r.get("font_size")
            nfs = n_r.get("font_size")
            if ofs is not None and nfs is not None and abs(ofs - nfs) > 0.5:
                changes.append({
                    "change_type": "RUN_FONT",
                    "property": f"{run_prefix}.font_size",
                    "old_value": ofs,
                    "new_value": nfs,
                })

            # Run bold
            if o_r.get("font_bold") != n_r.get("font_bold"):
                changes.append({
                    "change_type": "RUN_FONT",
                    "property": f"{run_prefix}.font_bold",
                    "old_value": o_r.get("font_bold"),
                    "new_value": n_r.get("font_bold"),
                })

            # Run italic
            if o_r.get("font_italic") != n_r.get("font_italic"):
                changes.append({
                    "change_type": "RUN_FONT",
                    "property": f"{run_prefix}.font_italic",
                    "old_value": o_r.get("font_italic"),
                    "new_value": n_r.get("font_italic"),
                })

            # Run color
            if o_r.get("font_color") != n_r.get("font_color"):
                ofc = o_r.get("font_color")
                nfc = n_r.get("font_color")
                if ofc is not None or nfc is not None:
                    changes.append({
                        "change_type": "RUN_FONT",
                        "property": f"{run_prefix}.font_color",
                        "old_value": ofc,
                        "new_value": nfc,
                    })

            # Run font name
            if o_r.get("font_name") != n_r.get("font_name"):
                ofn = o_r.get("font_name")
                nfn = n_r.get("font_name")
                if ofn is not None or nfn is not None:
                    changes.append({
                        "change_type": "RUN_FONT",
                        "property": f"{run_prefix}.font_name",
                        "old_value": ofn,
                        "new_value": nfn,
                    })

    return changes


def _diff_table_cells(old_tbl, new_tbl):
    """Compare table cells: text and fill color."""
    changes = []
    if old_tbl["rows"] != new_tbl["rows"] or old_tbl["cols"] != new_tbl["cols"]:
        changes.append({
            "change_type": "TABLE_STRUCTURE",
            "property": "dimensions",
            "old_value": f"{old_tbl['rows']}x{old_tbl['cols']}",
            "new_value": f"{new_tbl['rows']}x{new_tbl['cols']}",
        })
        return changes

    for r in range(old_tbl["rows"]):
        for c in range(old_tbl["cols"]):
            old_cell = old_tbl["cells"][r][c]
            new_cell = new_tbl["cells"][r][c]

            # Handle both old format (string) and new format (dict)
            old_text = old_cell["text"] if isinstance(old_cell, dict) else old_cell
            new_text = new_cell["text"] if isinstance(new_cell, dict) else new_cell

            if old_text != new_text:
                changes.append({
                    "change_type": "TABLE_CELL_TEXT",
                    "property": f"cell[{r},{c}].text",
                    "old_value": old_text,
                    "new_value": new_text,
                })

            # Cell fill (only if dict format)
            if isinstance(old_cell, dict) and isinstance(new_cell, dict):
                old_fill = old_cell.get("fill")
                new_fill = new_cell.get("fill")
                if old_fill != new_fill and (old_fill or new_fill):
                    changes.append({
                        "change_type": "TABLE_CELL_FILL",
                        "property": f"cell[{r},{c}].fill",
                        "old_value": old_fill,
                        "new_value": new_fill,
                    })

    return changes


# ═══════════════════════════════════════════════════════════════
# Builder function mapping via AST
# ═══════════════════════════════════════════════════════════════
def _parse_builder_lines(build_pptx_path):
    """Parse build_pptx.py to get function_name → line_number mapping."""
    source = Path(build_pptx_path).read_text(encoding="utf-8")
    tree = ast.parse(source)
    func_lines = {}
    for node in ast.walk(tree):
        if isinstance(node, ast.FunctionDef):
            func_lines[node.name] = node.lineno
    return func_lines


def _find_helper_calls_in_range(source_lines, start_line, end_line, old_values):
    """Search for helper function calls in a line range that contain old values.

    Returns list of (line_number, line_text) that likely need updating.
    """
    hits = []
    for i in range(start_line - 1, min(end_line, len(source_lines))):
        line = source_lines[i]
        # Check if line contains a helper call
        has_helper = any(pat in line for pat in HELPER_PATTERNS)
        if not has_helper:
            continue
        # Check if any old value appears in the line
        for val in old_values:
            if val is None:
                continue
            val_str = str(val)
            if val_str in line:
                hits.append((i + 1, line.rstrip()))
                break
    return hits


def locate_in_builder(build_pptx_path, func_name, changes):
    """Find the specific lines in build_pptx.py that need updating."""
    source = Path(build_pptx_path).read_text(encoding="utf-8")
    source_lines = source.split("\n")
    func_lines = _parse_builder_lines(build_pptx_path)

    if func_name not in func_lines:
        return None

    start = func_lines[func_name]

    # Find function end (next function at same indent level or end of file)
    sorted_funcs = sorted(func_lines.items(), key=lambda x: x[1])
    end = len(source_lines)
    for fn, ln in sorted_funcs:
        if ln > start:
            end = ln
            break

    # Collect old values from changes to search for
    old_values = set()
    for ch in changes:
        old_v = ch.get("old_value")
        if old_v is not None:
            old_values.add(old_v)

    return _find_helper_calls_in_range(source_lines, start, end, old_values)


# ═══════════════════════════════════════════════════════════════
# Full diff pipeline
# ═══════════════════════════════════════════════════════════════
def run_diff(edited_pptx, baseline_json, threshold=DEFAULT_THRESHOLD,
             build_pptx_path=None, output_json=False):
    """Run full diff between baseline snapshot and edited PPTX."""
    # Load baseline
    with open(baseline_json, "r", encoding="utf-8") as f:
        baseline = json.load(f)

    # Take snapshot of edited
    edited_snap = take_snapshot(edited_pptx)

    # Parse builder line numbers if build_pptx.py available
    builder_lines = {}
    if build_pptx_path and Path(build_pptx_path).exists():
        builder_lines = _parse_builder_lines(build_pptx_path)

    all_changes = []
    total_modified = 0
    total_added = 0
    total_removed = 0
    slides_changed = 0

    n_slides = min(len(baseline["slides"]), len(edited_snap["slides"]))

    for si in range(n_slides):
        b_slide = baseline["slides"][si]
        e_slide = edited_snap["slides"][si]

        builder_func = b_slide.get("builder_function", f"slide_{si}")
        builder_line = builder_lines.get(builder_func)

        matches = match_shapes(b_slide["shapes"], e_slide["shapes"])

        slide_changes = []
        for b_shape, e_shape, tier in matches:
            if b_shape is None:
                # ADDED
                slide_changes.append({
                    "shape_name": e_shape["name"],
                    "match_tier": tier,
                    "changes": [{"change_type": "ADDED", "property": "shape",
                                 "new_value": e_shape["shape_type"]}],
                })
                total_added += 1
            elif e_shape is None:
                # REMOVED
                slide_changes.append({
                    "shape_name": b_shape["name"],
                    "match_tier": tier,
                    "changes": [{"change_type": "REMOVED", "property": "shape",
                                 "old_value": b_shape["shape_type"]}],
                })
                total_removed += 1
            else:
                # Compare
                changes = diff_shapes(b_shape, e_shape, threshold, FONT_SIZE_THRESHOLD)
                if changes:
                    shape_entry = {
                        "shape_name": b_shape["name"],
                        "match_tier": tier,
                        "changes": changes,
                    }
                    # Try to locate in builder
                    if build_pptx_path and builder_func in builder_lines:
                        locs = locate_in_builder(build_pptx_path, builder_func, changes)
                        if locs:
                            shape_entry["source_hints"] = [
                                {"line": ln, "code": code} for ln, code in locs
                            ]
                    slide_changes.append(shape_entry)
                    total_modified += 1

        if slide_changes:
            slides_changed += 1
            slide_entry = {
                "slide_index": si,
                "slide_display": si + 1,
                "builder_function": builder_func,
            }
            if builder_line:
                slide_entry["builder_line"] = builder_line
            slide_entry["shape_changes"] = slide_changes
            all_changes.append(slide_entry)

    # Handle extra slides in edited
    if len(edited_snap["slides"]) > len(baseline["slides"]):
        for si in range(n_slides, len(edited_snap["slides"])):
            all_changes.append({
                "slide_index": si,
                "slide_display": si + 1,
                "builder_function": None,
                "shape_changes": [{"change_type": "SLIDE_ADDED"}],
            })

    result = {
        "baseline_file": baseline["source_file"],
        "edited_file": Path(edited_pptx).name,
        "threshold_inches": threshold,
        "summary": {
            "slides_changed": slides_changed,
            "shapes_modified": total_modified,
            "shapes_added": total_added,
            "shapes_removed": total_removed,
        },
        "changes": all_changes,
    }

    return result


# ═══════════════════════════════════════════════════════════════
# Output formatting
# ═══════════════════════════════════════════════════════════════
def format_diff_human(result):
    """Format diff result as human-readable text."""
    lines = []
    summary = result["summary"]

    if summary["shapes_modified"] == 0 and summary["shapes_added"] == 0 and summary["shapes_removed"] == 0:
        lines.append("No changes detected.")
        lines.append(f"  Baseline: {result['baseline_file']}")
        lines.append(f"  Edited:   {result['edited_file']}")
        lines.append(f"  Threshold: {result['threshold_inches']}\"")
        return "\n".join(lines)

    for slide_ch in result["changes"]:
        si = slide_ch["slide_display"]
        func = slide_ch.get("builder_function", "?")
        bline = slide_ch.get("builder_line", "")
        header = f"=== Slide {si} ({func}"
        if bline:
            header += f", line {bline}"
        header += ") ==="
        lines.append(header)

        for shape_ch in slide_ch.get("shape_changes", []):
            if isinstance(shape_ch, dict) and shape_ch.get("change_type") == "SLIDE_ADDED":
                lines.append("  [SLIDE_ADDED]")
                continue

            sname = shape_ch.get("shape_name", "?")
            tier = shape_ch.get("match_tier", "?")
            lines.append(f"  Shape: \"{sname}\" (match tier {tier})")

            for ch in shape_ch.get("changes", []):
                ct = ch["change_type"]
                prop = ch.get("property", "")
                old_v = ch.get("old_value", "")
                new_v = ch.get("new_value", "")
                delta = ch.get("delta", "")

                if ct in ("MOVED", "RESIZED"):
                    delta_str = f" ({'+' if delta > 0 else ''}{delta}\")" if delta else ""
                    lines.append(f"    [{ct}] {prop}: {old_v} -> {new_v}{delta_str}")
                elif ct == "TEXT":
                    lines.append(f"    [{ct}] {prop}:")
                    lines.append(f"      BEFORE: \"{old_v}\"")
                    lines.append(f"      AFTER:  \"{new_v}\"")
                elif ct == "FONT":
                    lines.append(f"    [{ct}] {prop}: {old_v} -> {new_v}")
                elif ct == "COLOR":
                    lines.append(f"    [{ct}] {prop}: {old_v} -> {new_v}")
                elif ct in ("ADDED", "REMOVED"):
                    lines.append(f"    [{ct}] {old_v or new_v}")
                elif ct == "IMAGE_REPLACED":
                    lines.append(f"    [{ct}] hash: {old_v}... -> {new_v}...")
                elif ct == "ROTATED":
                    lines.append(f"    [{ct}] {old_v} -> {new_v} deg")
                elif ct == "PARA_COUNT":
                    lines.append(f"    [PARA_COUNT] paragraphs: {old_v} -> {new_v}")
                elif ct == "PARA_ALIGN":
                    lines.append(f"    [PARA_ALIGN] {prop}: {old_v} -> {new_v}")
                elif ct == "PARA_SPACING":
                    lines.append(f"    [PARA_SPACING] {prop}: {old_v} -> {new_v}")
                elif ct == "RUN_TEXT":
                    lines.append(f"    [RUN_TEXT] {prop}: \"{old_v}\" -> \"{new_v}\"")
                elif ct == "RUN_FONT":
                    lines.append(f"    [RUN_FONT] {prop}: {old_v} -> {new_v}")
                elif ct == "TABLE_STRUCTURE":
                    lines.append(f"    [TABLE] dimensions: {old_v} -> {new_v}")
                elif ct == "TABLE_CELL_TEXT":
                    lines.append(f"    [TABLE_CELL] {prop}: \"{old_v}\" -> \"{new_v}\"")
                elif ct == "TABLE_CELL_FILL":
                    lines.append(f"    [TABLE_CELL_FILL] {prop}: {old_v} -> {new_v}")
                else:
                    lines.append(f"    [{ct}] {prop}: {old_v} -> {new_v}")

            # Source hints
            hints = shape_ch.get("source_hints", [])
            for hint in hints:
                lines.append(f"    -> build_pptx.py:{hint['line']}  {hint['code']}")

        lines.append("")

    lines.append(f"=== Summary: {summary['slides_changed']} slides changed, "
                 f"{summary['shapes_modified']} modified, "
                 f"{summary['shapes_added']} added, "
                 f"{summary['shapes_removed']} removed ===")
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════
# Inspect command
# ═══════════════════════════════════════════════════════════════
def inspect_slide(pptx_path, slide_num):
    """Print all shape properties for a single slide."""
    prs = Presentation(pptx_path)
    idx = slide_num - 1  # 1-based to 0-based
    if idx < 0 or idx >= len(prs.slides):
        print(f"[ERROR] Slide {slide_num} out of range (1-{len(prs.slides)})")
        return

    slide = prs.slides[idx]
    print(f"=== Slide {slide_num} ({len(slide.shapes)} shapes) ===\n")

    for i, shape in enumerate(slide.shapes):
        props = extract_shape_props(shape)
        print(f"  [{i}] \"{props['name']}\" ({props['shape_type']})")
        print(f"      pos:  ({props['left']}, {props['top']})  size: ({props['width']} x {props['height']})")

        text = props.get("text")
        if text:
            display_text = text[:60].replace("\n", "\\n")
            if len(text) > 60:
                display_text += "..."
            print(f"      text: \"{display_text}\"")

        fs = props.get("font_size")
        fb = props.get("font_bold")
        fc = props.get("font_color")
        fn = props.get("font_name")
        if any(v is not None for v in [fs, fb, fc, fn]):
            parts = []
            if fs is not None:
                parts.append(f"size={fs}pt")
            if fb:
                parts.append("bold")
            if fc:
                parts.append(f"color={fc}")
            if fn:
                parts.append(f"font={fn}")
            print(f"      font: {', '.join(parts)}")

        fill = props.get("fill_color")
        line = props.get("line_color")
        if fill or line:
            parts = []
            if fill:
                parts.append(f"fill={fill}")
            if line:
                parts.append(f"line={line}")
            print(f"      style: {', '.join(parts)}")

        tbl = props.get("table")
        if tbl:
            print(f"      table: {tbl['rows']}x{tbl['cols']}")

        img = props.get("image_hash")
        if img:
            print(f"      image: hash={img[:12]}...")

        print()


# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════
def main():
    parser = argparse.ArgumentParser(
        description="PPTX Snapshot & Diff — detect manual edits, map to build_pptx.py"
    )
    sub = parser.add_subparsers(dest="command")

    # snapshot
    p_snap = sub.add_parser("snapshot", help="Generate JSON snapshot of a PPTX")
    p_snap.add_argument("pptx", help="Path to PPTX file")
    p_snap.add_argument("-o", "--output", help="Output JSON path (default: <pptx>_snapshot.json)")
    p_snap.add_argument("--builder-map", help="Optional: builder function index map JSON")

    # diff
    p_diff = sub.add_parser("diff", help="Compare edited PPTX against baseline snapshot")
    p_diff.add_argument("pptx", help="Path to edited PPTX file")
    p_diff.add_argument("--baseline", help="Baseline snapshot JSON (default: <pptx>_snapshot.json)")
    p_diff.add_argument("--threshold", type=float, default=DEFAULT_THRESHOLD,
                        help=f"Position/size change threshold in inches (default: {DEFAULT_THRESHOLD})")
    p_diff.add_argument("--precise", action="store_true",
                        help=f"Use precise threshold ({PRECISE_THRESHOLD}\") to catch even tiny moves")
    p_diff.add_argument("--build-pptx", help="Path to build_pptx.py for source line hints")
    p_diff.add_argument("--json", action="store_true", help="Output as JSON instead of text")
    p_diff.add_argument("--archive", action="store_true", default=True,
                        help="Auto-save edited snapshot + JSON diff (default: enabled)")
    p_diff.add_argument("--no-archive", dest="archive", action="store_false",
                        help="Disable auto-archive")

    # inspect
    p_insp = sub.add_parser("inspect", help="Show all shapes on a specific slide")
    p_insp.add_argument("pptx", help="Path to PPTX file")
    p_insp.add_argument("--slide", type=int, required=True, help="Slide number (1-based)")

    args = parser.parse_args()

    if args.command == "snapshot":
        out = args.output or args.pptx.replace(".pptx", "_snapshot.json")
        # Try to build slide index map from BATCHES if build_pptx.py is available
        slide_map = None
        build_pptx = Path(args.pptx).parent / "build_pptx.py"
        if build_pptx.exists():
            slide_map = _infer_slide_map_from_batches(build_pptx)
        snap = take_snapshot(args.pptx, out, slide_index_map=slide_map)
        print(f"[OK] Snapshot saved: {out}")
        print(f"     Slides: {snap['total_slides']}")
        total_shapes = sum(len(s["shapes"]) for s in snap["slides"])
        print(f"     Shapes: {total_shapes}")

    elif args.command == "diff":
        baseline = args.baseline
        if not baseline:
            baseline = args.pptx.replace(".pptx", "_snapshot.json")
        if not Path(baseline).exists():
            print(f"[ERROR] Baseline not found: {baseline}")
            print("  Run 'snapshot' first, or specify --baseline")
            sys.exit(1)

        # --precise overrides --threshold
        threshold = PRECISE_THRESHOLD if args.precise else args.threshold

        build_pptx = args.build_pptx
        if not build_pptx:
            candidate = Path(args.pptx).parent / "build_pptx.py"
            if candidate.exists():
                build_pptx = str(candidate)

        result = run_diff(args.pptx, baseline, threshold, build_pptx)

        if args.json:
            print(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            print(format_diff_human(result))

        # Auto-archive: save edited snapshot + diff JSON
        if args.archive and result["summary"]["shapes_modified"] + result["summary"]["shapes_added"] + result["summary"]["shapes_removed"] > 0:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_stem = Path(args.pptx).stem
            archive_dir = Path(args.pptx).parent / "diff_archive"
            archive_dir.mkdir(exist_ok=True)

            # Save edited snapshot
            edited_snap_path = archive_dir / f"{ts}_{pptx_stem}_edited_snapshot.json"
            slide_map = _infer_slide_map_from_batches(Path(build_pptx)) if build_pptx else None
            take_snapshot(args.pptx, str(edited_snap_path), slide_index_map=slide_map)

            # Save diff result JSON
            diff_json_path = archive_dir / f"{ts}_{pptx_stem}_diff.json"
            with open(diff_json_path, "w", encoding="utf-8") as f:
                json.dump(result, ensure_ascii=False, indent=2, fp=f)

            print(f"\n[Archive] Saved to {archive_dir.name}/")
            print(f"  Edited snapshot: {edited_snap_path.name}")
            print(f"  Diff result:     {diff_json_path.name}")

    elif args.command == "inspect":
        inspect_slide(args.pptx, args.slide)

    else:
        parser.print_help()


def _infer_slide_map_from_batches(build_pptx_path):
    """Parse BATCHES dict from build_pptx.py source to infer slide → function mapping.

    This uses regex to extract function names from BATCHES definition,
    since importing build_pptx.py would have side effects.
    """
    source = Path(build_pptx_path).read_text(encoding="utf-8")

    # Find BATCHES block
    match = re.search(r"^BATCHES\s*=\s*\{", source, re.MULTILINE)
    if not match:
        return None

    # Extract all function references in order from BATCHES
    start = match.start()
    # Find the closing brace
    depth = 0
    end = start
    for i in range(start, len(source)):
        if source[i] == "{":
            depth += 1
        elif source[i] == "}":
            depth -= 1
            if depth == 0:
                end = i + 1
                break

    batches_text = source[start:end]

    # Parse batch numbers and their function lists
    # Pattern: number: [func1, func2, ...]
    slide_map = {}
    idx = 0
    # Find all list entries in order
    func_names = re.findall(r"\b(build_\w+)\b", batches_text)
    for fname in func_names:
        slide_map[idx] = fname
        idx += 1

    return slide_map


if __name__ == "__main__":
    main()
