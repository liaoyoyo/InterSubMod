#!/usr/bin/env python3
"""
build_pptx.py — Build Self-Phasing PPT 25-slide deck from markdown specs.

Reads:
- 03_slide_layout_script.md (per-slide spec)
- 04_speaker_script.md (Tier 2 speaker notes)
- figures/*.png (G1-G4 supplementary + master report figures)

Writes:
- output.pptx

Usage:
    python3 build_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

HERE = Path(__file__).parent
FIGURES = HERE / "figures"
MASTER  = FIGURES / "master"        # F1-F7 from 5/8 整合報告 + paired audit
IGV     = FIGURES / "igv"           # D_SP1/2/3 from PI 報告 by_HP_4ver
OUT = HERE / "output.pptx"

# ─── Color palette (對齊 R-G4 / preview.html) ──────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)
ORANGE      = RGBColor(0x9A, 0x34, 0x12)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
YELLOW      = RGBColor(0xFB, 0xBF, 0x24)
RED         = RGBColor(0xDC, 0x26, 0x26)
DARK_RED    = RGBColor(0x7F, 0x1D, 0x1D)
DARK_GREEN  = RGBColor(0x16, 0x65, 0x34)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY   = RGBColor(0x37, 0x41, 0x51)
OFF_WHITE   = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_YEL   = RGBColor(0xFF, 0xF3, 0xCD)
LIGHT_RED   = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# Fonts (CJK fallback chain handled by python-pptx via font.name)
FONT_CJK = "Noto Sans CJK TC"      # primary CJK
FONT_LATIN = "DejaVu Sans"          # primary Latin

# ─── Setup 16:9 presentation ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════

def new_slide():
    return prs.slides.add_slide(BLANK)

def add_textbox(slide, left, top, width, height, text="", *,
                font=FONT_CJK, size=14, bold=False, italic=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT, fill=None, border=None):
    """General-purpose text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(1)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if isinstance(text, str):
        text_lines = [text]
    else:
        text_lines = list(text)
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text_lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_title_block(slide, zh, en=""):
    """Slide title (assertion-evidence) + EN subtitle (60% size, indented)."""
    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7),
                zh, size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if en:
        add_textbox(slide, Inches(0.65), Inches(0.85), Inches(12.2), Inches(0.4),
                    en, size=14, italic=True, color=GRAY, align=PP_ALIGN.LEFT)
    # Title underline
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(1.32),
                                       Inches(12.93), Inches(1.32))
    line.line.color.rgb = NAVY
    line.line.width = Pt(2)

def add_glossary_footer(slide, items):
    """Bottom glossary box (R-G4 中等術語 + 簡單術語 footnote)."""
    if not items:
        return
    text = "  ·  ".join(items)
    add_textbox(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
                text, size=10, color=GRAY, italic=True)

def add_image_fit(slide, img_path, left, top, max_width, max_height):
    """Add image with fit-within (no aspect distortion)."""
    img_path = Path(img_path)
    if not img_path.exists():
        # placeholder if missing
        add_textbox(slide, left, top, max_width, max_height,
                    f"[missing: {img_path.name}]",
                    size=11, color=RED, fill=LIGHT_RED, border=RED,
                    align=PP_ALIGN.CENTER)
        return None
    with PILImage.open(img_path) as im:
        w, h = im.size
    aspect = w / h
    target_aspect = max_width / max_height
    if aspect > target_aspect:
        out_w = max_width
        out_h = int(max_width / aspect)
    else:
        out_h = max_height
        out_w = int(max_height * aspect)
    out_left = left + (max_width - out_w) // 2
    out_top = top + (max_height - out_h) // 2
    return slide.shapes.add_picture(str(img_path), out_left, out_top,
                                    width=out_w, height=out_h)

def add_table_box(slide, left, top, width, height, headers, rows,
                  highlight_rows=None):
    """Simple data table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    # Headers
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = FONT_CJK
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_CJK
            run.font.size = Pt(11)
            run.font.color.rgb = DARK_GRAY
            if highlight_rows and i in highlight_rows:
                cell.fill.solid()
                color = highlight_rows[i]
                cell.fill.fore_color.rgb = color
    return tbl

def add_caveat_box(slide, left, top, width, height, label, text, en_note=""):
    """Yellow caveat box."""
    box = add_textbox(slide, left, top, width, height,
                      "", fill=LIGHT_YEL, border=YELLOW)
    tf = box.text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = label + " "
    r1.font.bold = True
    r1.font.color.rgb = DARK_RED
    r1.font.size = Pt(13)
    r1.font.name = FONT_CJK
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(13)
    r2.font.color.rgb = DARK_RED
    r2.font.name = FONT_CJK
    if en_note:
        p2 = tf.add_paragraph()
        r3 = p2.add_run()
        r3.text = en_note
        r3.font.size = Pt(10)
        r3.font.italic = True
        r3.font.color.rgb = DARK_RED
        r3.font.name = FONT_LATIN
    return box

def add_red_frame(slide, left, top, width, height, label, text, en_note=""):
    """Red frame for V5 caveat (slide 16 only)."""
    box = add_textbox(slide, left, top, width, height,
                      "", fill=LIGHT_RED, border=RED)
    tf = box.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = label + " "
    r1.font.bold = True
    r1.font.color.rgb = RED
    r1.font.size = Pt(14)
    r1.font.name = FONT_CJK
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(13)
    r2.font.color.rgb = DARK_RED
    r2.font.name = FONT_CJK
    if en_note:
        p2 = tf.add_paragraph()
        r3 = p2.add_run()
        r3.text = en_note
        r3.font.size = Pt(10)
        r3.font.italic = True
        r3.font.color.rgb = DARK_RED
        r3.font.name = FONT_LATIN
    return box

def add_stat_box(slide, left, top, width, height, number, label, color=NAVY):
    """Big-number stat box for TL;DR."""
    box = add_textbox(slide, left, top, width, height, "",
                      fill=OFF_WHITE, border=GRAY)
    tf = box.text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = number
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = color
    r1.font.name = FONT_LATIN
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = GRAY
    r2.font.name = FONT_CJK
    return box

def add_cliffhanger_box(slide, left, top, width, height, line1, line2):
    """Yellow dashed cliffhanger box for slide 14."""
    box = add_textbox(slide, left, top, width, height,
                      "", fill=LIGHT_YEL, border=YELLOW)
    tf = box.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = line1
    r1.font.size = Pt(13)
    r1.font.italic = True
    r1.font.color.rgb = DARK_RED
    r1.font.name = FONT_CJK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = line2
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = DARK_RED
    r2.font.name = FONT_CJK
    return box

def add_speaker_note(slide, t2_text, t3_text=""):
    """Add speaker note (Tier 2) + Tier 3 oral-optional appendix."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    p1 = notes.paragraphs[0]
    r1 = p1.add_run()
    r1.text = t2_text
    r1.font.size = Pt(12)
    r1.font.name = FONT_CJK
    if t3_text:
        p2 = notes.add_paragraph()
        r2 = p2.add_run()
        r2.text = "[ORAL-OPTIONAL] " + t3_text
        r2.font.size = Pt(10)
        r2.font.italic = True
        r2.font.color.rgb = GRAY
        r2.font.name = FONT_CJK


# ═══════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_cover():
    s = new_slide()
    # Background fill
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE
    bg.line.fill.background()
    # Title
    add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.9),
                "Self-Phasing 整合觀察",
                size=42, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.6),
                "Self-Phasing Integration Synthesis",
                size=22, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5),
                "── V5 Layer 1.5 設計缺陷揭露 ──",
                size=20, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
                "longphase-to-mod 5 commits 修補成熟 + 5/9 paired cross-ref 新發現",
                size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    # Date + audience
    add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                "2026-05-10  ·  PI / lab meeting",
                size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
                "Source: 5/8 整合報告 + 5/9 errata + paired Step D",
                size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_speaker_note(s,
        "謝謝各位。今天 20 分鐘的報告主題是 self-phasing 整合觀察與 V5 Layer 1.5 設計缺陷揭露。"
        "整合 longphase-to-mod 5 commits 的修補成熟度，以及 5/9 paired cross-ref 新發現。"
        "目的：協助 PI 決策 V5 是否作 production tag baseline，以及是否啟動 F-paired-D3 follow-up。",
        "5/8 主報告 1,211 行 + 7 figures commit hash")

def slide_02_tldr():
    s = new_slide()
    add_title_block(s,
        "TL;DR — 修補主線確立，但 V5 Layer 1.5 germline-absent 區仍待補強",
        "Main fix line established; V5 Layer 1.5 still gaps in germline-absent regions")
    # Two stat groups
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(0.4),
                "修補主線 (V3F + V5)", size=13, color=GRAY)
    add_stat_box(s, Inches(0.4), Inches(2.0), Inches(2.9), Inches(1.3),
                 "17.3:1", "HP1 偏 baseline")
    add_stat_box(s, Inches(3.5), Inches(2.0), Inches(2.9), Inches(1.3),
                 "34,855", "read-level victims")

    add_textbox(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.4),
                "關鍵驗證", size=13, color=GRAY)
    add_stat_box(s, Inches(6.9), Inches(2.0), Inches(2.9), Inches(1.3),
                 "+13.3 pp", "paired GT (V5)", color=GREEN)
    add_stat_box(s, Inches(10.0), Inches(2.0), Inches(2.9), Inches(1.3),
                 "100%", "V3F+V5 修正率", color=GREEN)

    # Overall impact
    add_textbox(s, Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.5),
                "整體影響: 20/0 指標 no regression (caller F1 三版完全相同)",
                size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                fill=LIGHT_BLUE, border=NAVY)

    # Caveat
    add_caveat_box(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(1.5),
        "⚠ Caveat — 5/9 新發現:",
        "V5 Layer 1.5 germline-absent 區與 baseline 4.19:1 偏 HP1 完全相同; "
        "V3F 標 hp=33 反而更穩健 (~5% germline-absent 區，占比小不阻擋整體)",
        "V5 Layer 1.5 retains priority bug bias in germline-absent regions (~5% chr19 events)")

    add_glossary_footer(s, [
        "📖 sub-clone: 腫瘤內基因型相同細胞群; 不同 sub-clone 帶不同 somatic"])

    add_speaker_note(s,
        "今天 thesis 兩個焦點。正面: 修補主線確立。baseline LongPhase-TO 全基因組 HP1:HP2 = 17.3:1 systematic bias; "
        "V3F + V5 兩層修補在 read-level 對 chr19 752 + 全基因組 34,855 victims 修正率 100%; "
        "paired GT +13.3 pp; 20 指標 0 regression; caller F1 三版完全相同。"
        "反面: V5 Layer 1.5 germline-absent 區仍未修對 — 與 baseline 4.19:1 完全相同; V3F hp=33 更穩健。"
        "整體 12 維度 10 ✅ / 1 ⚠️ / 2 待跑; V5 仍可作 production baseline。",
        "V3F = 41ff147 / V5 = d0bcd8c + 938f0df / paired Step D = 766ec5f")

def slide_03_genome_173():
    s = new_slide()
    add_title_block(s,
        "全基因組 HP1:HP2 = 17.3:1 是 systematic bias，非樣本性質",
        "Genome-wide 17.3:1 = systematic engineering artifact, not sample variation")

    # Stat table
    add_table_box(s, Inches(0.4), Inches(1.6), Inches(6.5), Inches(2.2),
        ["指標", "baseline", "隨機預期", "偏離"],
        [
            ["HP1 reads (somatic ALT)", "614,000", "~325,000", "1.89×"],
            ["HP2 reads (somatic ALT)", "35,500", "~325,000", "0.11×"],
            ["HP1:HP2 ratio", "17.3:1", "1:1", "17.3×"],
            ["HP1 占比", "94.6%", "~50%", "+44.6 pp"],
        ],
        highlight_rows={3: LIGHT_YEL})

    # Visual contrast box
    add_textbox(s, Inches(7.2), Inches(1.6), Inches(5.7), Inches(0.5),
        "HP1 占比 94.6%", size=20, bold=True, color=ORANGE,
        align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.2), Inches(2.2), Inches(5.7), Inches(0.4),
        "↓↓↓↓↓↓↓↓↓", size=18, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.2), Inches(2.7), Inches(5.7), Inches(0.5),
        "隨機預期 50%", size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # 3 arguments
    add_textbox(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.4),
        "三條獨立論證 (3 independent arguments):", size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.5),
        "①  生物學: tumor sub-clone 跨 23 條染色體不該系統偏 HP1",
        size=13, color=DARK_GRAY)
    add_textbox(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5),
        "②  跨 chr 一致: cnLOH artifact 只影響單一 chr; 94.6% 跨 chr 一致",
        size=13, color=DARK_GRAY)
    add_textbox(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.5),
        "③  paired 對照: paired tumor-normal 同 reads HP1:HP2 ≈ 1:1",
        size=13, color=DARK_GRAY)

    # Conclusion
    add_textbox(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.5),
        "→ 17.3:1 不是 sample 性質，是 LongPhase-TO 的 systematic engineering artifact",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)

    add_glossary_footer(s, [
        "📖 germline het: 個人遺傳變異中雜合位點 (A/G)",
        "📖 sub-clone: 腫瘤內基因型相同細胞群",
        "📖 haplotype: 父系/母系兩條染色體之一"])

    add_speaker_note(s,
        "baseline LongPhase-TO 全基因組: HP1 reads 614K vs HP2 35.5K, 比例 17.3:1 vs 隨機 1:1。"
        "94.6% 占比是 systematic bias 的硬證據, 三條獨立論證: "
        "(1) 生物學 tumor sub-clone 不該跨 23 chr 系統偏 HP1; "
        "(2) cnLOH artifact 只影響單 chr 但這偏移跨 23 chr 一致; "
        "(3) paired pipeline HP1:HP2 ≈ 1:1。"
        "三條互相獨立不可能同時都是 sample 效應 → engineering artifact。",
        "cnLOH 機制細節 / 23 chr 一致性表 / paired pipeline 程式碼差異")

def slide_04a_sp1():
    s = new_slide()
    add_title_block(s,
        "SP1 chr19:17,565,944 — baseline 113:0 → V5 翻轉至 HP2 ✅ paired",
        "SP1 — baseline 113:0 → V5 flips to HP2, aligned with paired")

    add_table_box(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.7),
        ["位點", "baseline HP1:HP2", "V5", "paired GT", "對齊?"],
        [["SP1 chr19:17,565,944", "113 : 0", "翻轉至 HP2 主導", "HP2", "✅"]],
        highlight_rows={0: LIGHT_RED})

    # IGV image (D_SP1 by_HP_4ver: baseline / V2b / V3F / V5 / paired_T / paired_N)
    add_image_fit(s, IGV / "D_SP1_chr19_17565944.png",
                  Inches(0.4), Inches(2.5), Inches(8.5), Inches(4.0))

    add_textbox(s, Inches(9.2), Inches(2.5), Inches(3.7), Inches(4.0),
        "為何能排除\n是噪音/caller/alignment?\n\n"
        "•  baseline 與 paired\n"
        "    方向相反\n"
        "    (不是衰減而是翻轉)\n\n"
        "•  V5 修正後與 paired\n"
        "    ground truth 重合\n\n"
        "→  read assignment\n"
        "    強制集中的鐵證",
        size=11, color=DARK_GRAY, fill=LIGHT_BLUE, border=NAVY)

    add_glossary_footer(s, [
        "📖 haplotype: 父系/母系兩條染色體之一; germline het 的方向標"])

    add_speaker_note(s,
        "全基因組 17.3:1 是平均值; IGV 6-BAM 並列篩到 chr19 三個近 100% 失衡位點。"
        "先看 SP1 chr19:17,565,944: baseline panel 113 reads 全部 HP1, HP2 stack 為 0; "
        "V5 翻轉至 HP2 主導, 與 paired tumor 一致。"
        "排除噪音 / caller / alignment: baseline 與 paired 完全反向 (不是衰減而是翻轉), "
        "V5 修正後與 paired ground truth 重合 → read assignment 強制集中的鐵證。next: SP2/SP3。",
        "6-BAM 並列順序 / V2b 中間階段意義")

def slide_04b_sp2_sp3():
    s = new_slide()
    add_title_block(s,
        "SP2 + SP3 並列 — 同模式: baseline 全 HP1 / V5 翻 HP2 / 3/3 對齊 paired",
        "SP2 + SP3 same pattern; 3/3 aligned with paired ground truth")

    # Two side-by-side
    add_table_box(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.7),
        ["位點", "baseline", "V5"],
        [["SP2 chr19:12,452,332", "109 : 1", "HP2 主導"]],
        highlight_rows={0: LIGHT_RED})
    add_table_box(s, Inches(6.9), Inches(1.6), Inches(6.0), Inches(0.7),
        ["位點", "baseline", "V5"],
        [["SP3 chr19:12,467,180", "108 : 0", "HP2 主導"]],
        highlight_rows={0: LIGHT_RED})

    # IGV images side-by-side
    add_image_fit(s, IGV / "D_SP2_chr19_12452332.png",
                  Inches(0.4), Inches(2.6), Inches(6.0), Inches(2.0))
    add_image_fit(s, IGV / "D_SP3_chr19_12467180.png",
                  Inches(6.9), Inches(2.6), Inches(6.0), Inches(2.0))

    add_textbox(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(0.5),
        "→ 三 SP 都在 chr19:12-17M 區段 → 對齊 slide 09 chr19 752 victims hotspot 區",
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        fill=GREEN)

    add_textbox(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.2),
        "引出四問:\n"
        "①  為何全集中一邊?  (→ S2 機制)\n"
        "②  read 層級多少個案?  (→ S3 量化)\n"
        "③  三版各修什麼?  (→ S4 修補)\n"
        "④  是否都修對?  (→ S5 驗證)",
        size=11, color=DARK_GRAY)

    add_speaker_note(s,
        "SP2 chr19:12,452,332 baseline 109:1; SP3 chr19:12,467,180 baseline 108:0; "
        "與 SP1 同模式: baseline 全 HP1, V5 翻 HP2 對齊 paired ground truth, 3/3 全對齊。"
        "三 SP 都在 chr19:12-17M, 與 slide 09 chr19 752 victims 1Mb hotspot 散點圖最高 enrichment 區段對齊 — "
        "read-level 個案與 IGV 截圖屬同一機制不同層級觀察。"
        "三 SP 都 V5 修對引出四問機制 / 量化 / 修補 / 驗證。",
        "paired_T 與 paired_N 對照細節")

def slide_05_player_referee():
    s = new_slide()
    add_title_block(s,
        "phasing 層球員兼裁判 — somatic 100% 共現蓋過 germline 50/50",
        "Phasing layer player-as-referee — somatic overrules germline")

    # G1 figure
    add_image_fit(s, FIGURES / "G1_player_as_referee.png",
                  Inches(0.4), Inches(1.5), Inches(8.0), Inches(4.5))

    # Right-side text
    add_textbox(s, Inches(8.7), Inches(1.6), Inches(4.4), Inches(0.4),
        "TO 模式致命 (4 條):", size=12, bold=True, color=DARK_RED)
    add_textbox(s, Inches(8.7), Inches(2.0), Inches(4.4), Inches(2.0),
        "•  TO 沒 paired normal → 用 PoN\n"
        "    (somatic 不在 PoN → 當 germline)\n"
        "•  somatic 進 graph 後 edge\n"
        "    weight 暴漲 (100% > 50% 共現)\n"
        "•  自我增強: 3 somatic 共現 →\n"
        "    該 haplotype 強度 ×3\n"
        "•  germline 真實訊號被 overrule",
        size=10, color=DARK_GRAY, fill=LIGHT_RED, border=RED)

    add_textbox(s, Inches(8.7), Inches(4.2), Inches(4.4), Inches(0.8),
        "解法: PON-only flag\n(commit 8b8c1fd)\nphasing graph 只放 PoN germline",
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        fill=GREEN)

    add_glossary_footer(s, [
        "📖 PoN: Panel of Normals — 多正常樣本建構的 germline reference set",
        "📖 phasing graph: het 位點當 node, read 共現當 edge"])

    add_speaker_note(s,
        "機制兩層 bug, 先 phasing 層 — 球員兼裁判隱喻。"
        "正常 phasing graph 的物理基礎: germline het 在 HP1/HP2 兩 stream 50/50 隨機分佈。"
        "TO 沒 paired normal 區分 germline / somatic, 只能用 PoN; 未在 PoN 內的位點被當 germline → "
        "somatic 進 graph 後 edge weight 暴漲 (100% > 50% 共現), graph 自我增強迴圈。"
        "球員兼裁判: somatic 應該被 phase 反過來主導 graph。"
        "修法 PON-only flag (commit 8b8c1fd) phasing 階段只放 PoN germline。但這只解 phasing 層, tag 還壞。",
        "TO vs paired PoN 對照 / 自我增強迴圈 / Pass 1 vs Pass 2 預告")

def slide_06_priority_bug():
    s = new_slide()
    add_title_block(s,
        "tagging 層 getVote priority bug — 1 票 somatic 觸發誤標",
        "tagging layer getVote priority bug — 1 somatic vote triggers mislabel")

    # F1 mechanism figure (top-left)
    add_image_fit(s, MASTER / "F1_priority_bug_mechanism.png",
                  Inches(0.4), Inches(1.5), Inches(7.0), Inches(3.2))

    # Right: 5-vote real read example
    add_textbox(s, Inches(7.7), Inches(1.5), Inches(5.4), Inches(0.4),
        "Real read 範例 (752 同模式):", size=11, bold=True, color=DARK_RED)
    add_textbox(s, Inches(7.7), Inches(2.0), Inches(5.4), Inches(2.7),
        "germline HP1 = 0\n"
        "germline HP2 = 5  ← 主導\n"
        "somatic HP1_1 = 1  ← 1 票觸發\n"
        "somatic HP2_1 = 0\n"
        "──────────────────────\n"
        "baseline: → hp=11 ❌ break\n"
        "          (germline 5 票被忽略)\n"
        "正確答案: hp=21\n"
        "(germline HP2=5 主導 + 標 21)",
        size=10, font=FONT_LATIN, color=DARK_GRAY, fill=OFF_WHITE, border=GRAY)

    add_textbox(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.3),
        "→ tumor sub-clone somatic 100% 同方向 → priority bug 把所有受影響 reads 標 HP:i:11 系列\n"
        "→ 17.3:1 偏移在 tag layer 形成\n"
        "(注意: tag layer 與 §slide 05 phasing layer 是不同層 bug, 必須分別修補)",
        size=12, color=NAVY, fill=LIGHT_BLUE, border=NAVY)

    add_glossary_footer(s, [
        "📖 sub-clone: 不同 sub-clone 帶不同 somatic",
        "ⓘ priority bug: vector 順序檢查 + break early 導致前面條目蓋過後面",
        "ⓘ break early: for 迴圈第一個非空處 break"])

    add_speaker_note(s,
        "tagging 層 priority bug: baseline getVote() vector 順序 ① somatic ② mixed ③ germline; "
        "for 迴圈第一個非空 break early, germline 永遠看不到。"
        "範例 read: germline HP2=5 主導, somatic HP1_1=1 觸發, baseline 標 hp=11 錯, 應 hp=21。"
        "全 752 chr19 victims 同模式。"
        "tumor sub-clone somatic 100% 同方向 → priority bug 把所有受影響 reads 標 HP:i:11 系列 → 17.3:1 形成。"
        "tag layer 與 phasing layer 是不同層 bug, 不能合併單 commit 修。",
        "enum HAPLOTYPE1_1=2 vs HP tag int=11 比較 bug / 5-vote countMap 結構")

def slide_07_two_layer_table():
    s = new_slide()
    add_title_block(s,
        "兩層 bug 兩層修補對應 — phasing PON-only + tagging V3F+V5 缺一不可",
        "Two-layer bugs: PON-only (phasing) + V3F+V5 (tagging) both required")

    # Mapping table
    add_table_box(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.0),
        ["Layer", "Bug", "修補 commit", "§5 章節"],
        [
            ["phasing 層", "球員兼裁判 (somatic 進 graph)",
             "8b8c1fd PON-only flag", "§5.2"],
            ["tagging 層", "priority bug (vector 順序錯)",
             "V3F: 41ff147 + 380e8d2",  "§5.3"],
            ["", "", "V5: + d0bcd8c + 938f0df", "§5.4"],
        ],
        highlight_rows={0: LIGHT_BLUE, 1: LIGHT_GREEN, 2: LIGHT_GREEN})

    # Two why questions
    add_textbox(s, Inches(0.4), Inches(4.0), Inches(6.0), Inches(2.0),
        "為何不能只用 PON-only?\n\n"
        "解 phasing 但 tag 仍壞\n"
        "→ 99.9% reads 仍標 HP:i:11 系列",
        size=12, color=DARK_GRAY, fill=LIGHT_YEL, border=YELLOW)

    add_textbox(s, Inches(6.9), Inches(4.0), Inches(6.0), Inches(2.0),
        "為何 V3F 不夠還要 V5?\n\n"
        "V3F Layer 1 only → germline 缺席區\n"
        "reads 全 untagged; V5 Layer 1.5 補 fallback\n"
        "(但有 germline-absent caveat, slide 16 詳述)",
        size=12, color=DARK_GRAY, fill=LIGHT_YEL, border=YELLOW)

    add_glossary_footer(s, [
        "📖 PoN: Panel of Normals — 多正常樣本建構的 germline reference set"])

    add_speaker_note(s,
        "兩層 bug 是獨立問題: phasing 層球員兼裁判 → 8b8c1fd PON-only; "
        "tagging 層 priority bug → V3F two-layer (41ff147 + 380e8d2 INDEL guard) → "
        "V5 (+ d0bcd8c Pass 2 ploidy fix + bundled Layer 1.5 + 938f0df threshold)。"
        "任一單修不夠必須 stacking。"
        "PON-only 不夠: 解 phasing 但 tag 仍偏 99.9% reads 標 HP:i:11; "
        "V3F 不夠: Layer 1 only germline 缺席區 reads 全 untagged。"
        "V5 Layer 1.5 有自己 caveat (slide 16 詳述)。",
        "5 commit 順序 / V3F 命名來歷")

def slide_08_chr19_752():
    s = new_slide()
    add_title_block(s,
        "chr19 752 victims — 100% 單向 baseline=11 → V3F/V5=21",
        "chr19 752 victims — 100% unidirectional fix")

    # Scale + fix rate
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(0.4),
        "規模:", size=12, bold=True, color=NAVY)
    add_textbox(s, Inches(0.4), Inches(2.0), Inches(6.0), Inches(1.4),
        "Dump rows:        549,206\n"
        "3-way merged:     1,069,832\n"
        "Priority bug victims: 752",
        size=12, font=FONT_LATIN, color=DARK_GRAY, fill=OFF_WHITE, border=GRAY)

    add_textbox(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.4),
        "修正率:", size=12, bold=True, color=GREEN)
    add_textbox(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(1.4),
        "V3F 修正率:  100.00% ✅\n"
        "V5  修正率:  100.00% ✅\n"
        "全 752 條無一條反向",
        size=12, font=FONT_LATIN, color=DARK_GREEN, fill=LIGHT_GREEN, border=GREEN)

    # 4-path verification (compressed to make room for F4)
    add_table_box(s, Inches(0.4), Inches(3.6), Inches(7.0), Inches(2.4),
        ["路徑", "結果", "判定"],
        [
            ["①  個案 trace ≥10",     "752 條",                     "✅"],
            ["②  1Mb 區域聚集",        "30M(215)+27M(133) 46%",      "⚠ PARTIAL"],
            ["③  Somatic density",    "high≥5=0 受害, 低票觸發",     "🔄 反向"],
            ["④  修正後消失",          "V3F/V5 100%",                "✅"],
        ],
        highlight_rows={0: LIGHT_GREEN, 3: LIGHT_GREEN, 1: LIGHT_YEL})

    # F4 hotspot scatter (right side)
    add_image_fit(s, MASTER / "F4_chr19_752_victims_scatter.png",
                  Inches(7.7), Inches(3.4), Inches(5.4), Inches(2.7))

    add_textbox(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.6),
        "→ 3 PASS + 1 PARTIAL = priority bug 機制因果確立 (chr19 only scope)",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)

    add_glossary_footer(s, ["ⓘ scope: chr19 only"])

    add_speaker_note(s,
        "chr19 read-level audit: dump 549,206 rows × 3 binary versions JOIN merged 1.07M events, "
        "篩 germline_majority ≠ somatic_majority 雙向矛盾 = 752 victims。"
        "V3F + V5 修正率 100%, 全 752 條 baseline=11 → V3F=21 單向。"
        "4-path 驗證: ① 752 條 PASS; ② 1Mb 30M 215 + 27M 133 共佔 46% PARTIAL; "
        "③ density 共變 high vote ≥5 = 0 受害, 低票觸發 — 反向有意義 (high 票 sub-clone 一致已對齊); "
        "④ V3F/V5 修正後消失 PASS。3 PASS + 1 PARTIAL 機制因果確立。",
        "4-path detail / read_name 真實 case 5 條 / SP1/2/3 對應 chr19 hotspot")

def slide_09_genome_34855():
    s = new_slide()
    add_title_block(s,
        "全基因組 34,855 victims (46×) — chr19 占 2.16% rank 19",
        "Genome-wide 34,855 victims; main hotspots chr7/chr2/chr1")

    # Scale comparison
    add_table_box(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(1.4),
        ["", "chr19 pilot", "Genome F1", "倍數"],
        [
            ["Dump rows", "549,206", "29,973,253", "54.6×"],
            ["Priority bug victims", "752", "34,855", "46.4×"],
            ["V3F / V5 修正率", "100% / 100%", "100% / 100%", "一致 ✅"],
        ],
        highlight_rows={1: LIGHT_YEL, 2: LIGHT_GREEN})

    # Per-chr table
    add_table_box(s, Inches(0.4), Inches(3.2), Inches(8.0), Inches(2.4),
        ["chr", "victims", "占全基因組", "rank"],
        [
            ["chr7",  "3,508", "10.1%", "1"],
            ["chr2",  "2,792", "8.0%",  "2"],
            ["chr1",  "2,674", "7.7%",  "3"],
            ["chr16", "2,584", "7.4%",  "4"],
            ["chr19", "752",   "2.16%", "19  ★ 不是主 hotspot"],
            ["chr8",  "666",   "1.9%",  "21  ★ 冷區"],
        ],
        highlight_rows={4: LIGHT_RED, 5: LIGHT_BLUE})

    # F2 per-chr enrichment (right side)
    add_image_fit(s, MASTER / "F2_priority_bug_per_chr_enrichment.png",
                  Inches(8.7), Inches(3.2), Inches(4.4), Inches(2.4))

    add_glossary_footer(s, [
        "ⓘ scope: 全基因組 (T1.2-F1)",
        "ⓘ chr8 LOH+HPSig 是 ISM 下游 hotspot (§8.2)"])

    add_speaker_note(s,
        "chr19 752 不是局部 artifact: 全基因組擴展 34,855 victims (46.4×)。"
        "V3F/V5 修正率仍 100% 方向一致。"
        "per-chr 分佈推翻原 chr19 pilot 結論: 主要 hotspot 是 chr7 (3,508) / chr2 / chr1 / chr16 / chr20, "
        "chr19 占比僅 2.16% rank 19。"
        "chr8 priority bug enrichment 0.34× (rank 21 冷區), 與 chr8 LOH+HPSig hotspot 是不同 layer "
        "(後者 ISM 下游 false-positive 富集 7.4×, 與 priority bug 無直接關聯)。",
        "全 chr enrichment ‰ 表 / chrY 小 N 高 ‰ 解釋")

def slide_10_5_commits():
    s = new_slide()
    add_title_block(s,
        "5 commits 兩層三版 stacking — baseline → V3F → V5",
        "5 commits two-layer three-version stacking")

    # F3 timeline figure
    add_image_fit(s, MASTER / "F3_binary_commit_timeline.png",
                  Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.3))

    # Stacking
    add_textbox(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.5),
        "V3-Fixed = baseline + 41ff147 + 380e8d2",
        size=13, bold=True, color=DARK_GREEN)
    add_textbox(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.5),
        "V5 = V3F + d0bcd8c + 938f0df",
        size=13, bold=True, color=DARK_GREEN)

    # Note
    add_textbox(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.2),
        "★ 41ff147 是修偏移的關鍵 commit (tagging 層 priority bug fix)\n"
        "累計 ~155 lines tagging-layer + ~40 lines phasing-layer\n"
        "HaplotagProcess.h:66-68 介面契約零變動\n"
        "d0bcd8c 是唯一跨兩層 commit (Pass 2 ploidy fix + Layer 1.5 fallback)",
        size=11, color=DARK_GRAY, fill=OFF_WHITE, border=GRAY)

    add_glossary_footer(s, [
        "ⓘ INDEL guard: 補 OOB UB (HAPLOTYPE_UNDEFINED 檢查)",
        "ⓘ threshold 0.95 → 0.9: Pass 2 second round 觸發 purity 閾值"])

    add_speaker_note(s,
        "5 commits 漸進完成 self-phasing 修補。"
        "8b8c1fd (4-09) PON-only flag 解 phasing 層 (藍); "
        "41ff147 (4-10) two-layer getVote 解 tagging 層 priority bug (綠) ★ 是修偏移關鍵; "
        "380e8d2 (4-25) INDEL guard 補 OOB UB (綠); "
        "d0bcd8c (4-30) Pass 2 ploidy fix + bundled Layer 1.5 + countSNP guard (紫 跨兩層); "
        "938f0df (4-30) threshold 0.95→0.9 (藍 Pass 2 觸發)。"
        "V3F = baseline + 41ff147 + 380e8d2; V5 = V3F + d0bcd8c + 938f0df。"
        "累計 155 行 tagging + 40 行 phasing; 介面契約零變動。",
        "各 commit 對應 layer 細節 / 為何不合併 / cherry-pick 自 upstream zhenyu")

def slide_11_getvote_3versions():
    s = new_slide()
    add_title_block(s,
        "getVote 三版差異 — baseline 順序錯 → V3F 兩層 → V5 +Layer 1.5",
        "getVote 3-version diff — baseline ordered → V3F two-layer → V5 +1.5")

    # G3 figure
    add_image_fit(s, FIGURES / "G3_getVote_three_layer.png",
                  Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.6))

    # Caveat
    add_caveat_box(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.6),
        "⚠ Layer 1.5 caveat:",
        "V5 Layer 1.5 在 germline-absent 區會繼承 priority bug 偏移 → (slide 16 詳述)")

    add_glossary_footer(s, [
        "📖 germline het: 個人遺傳變異中雜合位點",
        "ⓘ two-layer / Layer 1.5 fallback / phased votes"])

    add_speaker_note(s,
        "三版 code side-by-side。baseline 紅底: vector keys 順序 ① somatic FIRST, for 迴圈第一個非空 break, priority bug 在這。"
        "V3F 綠底: 拆 explicit Layer 1 germline only 決方向 (gR=1 或 2); Layer 2 somatic annotation 標 hp=11/21/33。"
        "Layer 1 germline 永不被 somatic overrule。"
        "V5 綠底加黃 highlight: 保留 Layer 1, 新增 Layer 1.5 fallback — germline 缺席時用 somatic phased votes 決方向; Layer 2 同 V3F。"
        "黃 highlight 是 caveat 預告: Layer 1.5 在 germline-absent 區繼承 priority bug 偏移, slide 16 詳述。",
        "enum HAPLOTYPE1_1=2 / HP tag int=11 比較 / V3F bonus 修")

def slide_12_sp_fixed():
    s = new_slide()
    add_title_block(s,
        "個案層 V5 修正 3/3 + 全基因組 HP1:HP2 17.3:1 → ~1:1",
        "Site-level V5 fixes 3/3 + genome HP1:HP2 17.3:1 → ~1:1")

    # Site-level table
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4),
        "個案層: SP1/2/3 修正後對齊 paired", size=13, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.5),
        ["位點", "baseline", "V5 翻轉", "paired", "對齊?"],
        [
            ["SP1 chr19:17,565,944", "113:0",  "HP2 主導", "HP2", "✅"],
            ["SP2 chr19:12,452,332", "109:1",  "HP2 主導", "HP2", "✅"],
            ["SP3 chr19:12,467,180", "108:0",  "HP2 主導", "HP2", "✅"],
        ],
        highlight_rows={0: LIGHT_GREEN, 1: LIGHT_GREEN, 2: LIGHT_GREEN})

    # Genome-level table
    add_textbox(s, Inches(0.4), Inches(3.7), Inches(12.5), Inches(0.4),
        "全基因組層: 17.3:1 → ~1:1 消除偏移", size=13, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(1.4),
        ["指標", "baseline", "V5", "Δ"],
        [
            ["HP1:HP2 ratio",       "17.3:1", "~1:1",      "消除偏移"],
            ["94.6% somatic→HP1",   "是",     "~50%",      "balanced"],
            ["15-site Problem PS",  "48.5%",  "52.0%",     "+3.5 pp"],
        ],
        highlight_rows={0: LIGHT_GREEN, 1: LIGHT_GREEN, 2: LIGHT_GREEN})

    # F5 zero-sum 4-quadrant figure (compact at bottom)
    add_image_fit(s, MASTER / "F5_layer15_zero_sum_4quadrant.png",
                  Inches(0.4), Inches(5.7), Inches(8.5), Inches(1.3))
    add_textbox(s, Inches(9.0), Inches(5.7), Inches(4.0), Inches(1.3),
        "F5 zero-sum:\n"
        "germline=0  +560,881\n"
        "germline>0  -560,881\n"
        "總和         =0",
        size=11, color=DARK_GREEN, fill=LIGHT_GREEN, border=GREEN)

    add_glossary_footer(s, [
        "📖 haplotype: 父系/母系兩條染色體之一",
        "ⓘ scope: chr19 個案 + 全基因組統計"])

    add_speaker_note(s,
        "個案層: SP1/2/3 baseline 113:0 / 109:1 / 108:0, V5 翻 HP2 對齊 paired 3/3 ✅。"
        "全基因組層: HP1:HP2 17.3:1 → ~1:1 消除偏移; 94.6% somatic→HP1 → ~50% balanced; "
        "15-site Problem PS (含 SP1/2/3) 48.5% → 52.0% +3.5 pp 看似小但機制顯著。"
        "F5 zero-sum 重分配: germline=0 +560,881 reads (V5 Layer 1.5 觸發) / germline>0 -560,881 / 總和 0。"
        "Pass 2 reclassify 104K germline het 為 somatic / ./. 是因果。",
        "V2b / V3F 中間版本 / Layer 1.5 zero-sum 詳細機制")

def slide_13_20_metrics():
    s = new_slide()
    add_title_block(s,
        "20 指標 no regression — 6 項 ⭐ 顯著改善 +8.3 ~ +99.7%",
        "20 metrics no regression; 6 significant improvements")

    # 5 categories
    cats = [
        ("①  ISM aggregate",  "TP_rate +0.005 / HP_Ratio 0.788→0.574 / Potential_LOH +3.5 pp", LIGHT_GREEN),
        ("②  HP_Ratio AUC",   "All -0.005 (隨機區間) / Inner +0.002",                          LIGHT_GREEN),
        ("③  Methylation 6 feat", "全 ±0.01 內持平",                                            LIGHT_GREEN),
        ("④  Paired GT concord. ⭐", "clean PS +8.3 pp / 15-Aggr +6.65 pp / 15-Clean PS +13.3 pp", LIGHT_YEL),
        ("⑤  HP / LOH 結構 ⭐",   "N50 +99.7% / Phased +23.6 pp / 1.36× 快 / LOH 完全相同",      LIGHT_YEL),
    ]
    for i, (cat, content, fill) in enumerate(cats):
        y = Inches(1.5 + i * 0.65)
        add_textbox(s, Inches(0.4), y, Inches(3.0), Inches(0.55),
            cat, size=12, bold=True, color=NAVY, fill=fill, border=GRAY)
        add_textbox(s, Inches(3.5), y, Inches(9.4), Inches(0.55),
            content, size=12, color=DARK_GRAY, fill=fill, border=GRAY)

    # 6 highlights
    add_textbox(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5),
        "6 顯著改善 ⭐: N50 +99.7% / Phased rate +23.6 pp / 1.36× 速度 / 15-Clean PS +13.3 pp / 全基因組 clean PS +8.3 pp / 15-Aggr +6.65 pp",
        size=11, bold=True, color=WHITE, fill=GREEN)

    # Conclusion
    add_textbox(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.6),
        "→ 20/0 指標 no regression — V5 全面 production-ready",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)

    add_textbox(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.4),
        "(HP_Ratio 0.788→0.574 是 tag bias 修正非變差 — pre-registered metrics, 無 cherry-picking)",
        size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_glossary_footer(s, [
        "📖 LOH: 雜合性丟失 (Loss of Heterozygosity)",
        "ⓘ scope: HCC1395 5kHz @ 0.93 purity (PI 報告 V5 = Pass 1 only BAM)"])

    add_speaker_note(s,
        "5 大類別 20 指標全綠 (pre-registered metrics 無 cherry-picking): "
        "① ISM aggregate 3 項 — HP_Ratio 0.788→0.574 是 tag bias 修正非變差; "
        "② HP_Ratio AUC 兩項在隨機區間; ③ methylation 6 全持平; "
        "④ Paired GT 4 項 ⭐ +6.65 ~ +13.3 pp; "
        "⑤ HP/LOH 5 項 ⭐ N50 +99.7% / Phased +23.6 pp / 1.36× 快 / LOH Jaccard=1.0。"
        "20/0 no regression — V5 全面 production-ready。",
        "methylation 6 feat 列表 / HP_Ratio 詳解 / LOH bed Jaccard=1.0 機制")

def slide_14_caller_f1_cliffhanger():
    s = new_slide()
    add_title_block(s,
        "Caller F1 vs SEQC2 三版完全相同; purity 0.6 完整對照 0 critical regression",
        "Caller F1 identical across versions; purity 0.6 fully verified")

    # 0.93 purity table
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
        "HCC1395 5kHz @ 0.93 purity:",
        size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(1.85), Inches(12.5), Inches(1.6),
        ["版本", "TP", "FP", "FN", "Precision", "Recall", "F1"],
        [
            ["A1 baseline", "28,509", "11,606", "10,938", "0.7107", "0.7227", "0.7166"],
            ["A3 V3F",      "28,509", "11,606", "10,938", "0.7107", "0.7227", "0.7166"],
            ["A5 V5",       "28,509", "11,606", "10,938", "0.7107", "0.7227", "0.7166"],
        ],
        highlight_rows={0: LIGHT_GREEN, 1: LIGHT_GREEN, 2: LIGHT_GREEN})

    # 0.6 purity row
    add_textbox(s, Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.4),
        "HCC1395 t30_n20 @ 0.6 purity:", size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.7),
        ["版本", "TP", "FP", "FN", "F1"],
        [["B1/B3/B5 三版", "24,190", "13,487", "15,257", "0.6273"]],
        highlight_rows={0: LIGHT_GREEN})

    # Causal chain
    add_textbox(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(0.7),
        "因果鏈: ClairS-TO PASS set 由 caller 決定 → V5 改 GT/PS/GT2/GT3 不改 FILTER\n"
        "         → PASS set 不變 → TP/FP/FN 完全相同 → F1 完全相同",
        size=11, color=DARK_GRAY, fill=LIGHT_BLUE, border=NAVY)

    # Cliffhanger
    add_cliffhanger_box(s, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.0),
        "→ V5 不改 caller; ΔF1 (0.93→0.6) = -0.0893 為 ClairS-TO 性質",
        "→ 但 5/9 paired cross-ref 揭露另一面...  〔next slide 15-16〕")

    add_glossary_footer(s, [
        "📖 purity: 樣本中腫瘤細胞占比",
        "ⓘ PASS set: ClairS-TO snv.vcf FILTER=PASS 集合"])

    add_speaker_note(s,
        "Caller F1 vs SEQC2 v1.2.1。0.93 purity: A1/A3/A5 三版 TP=28,509 / FP=11,606 / FN=10,938 完全相同, F1=0.7166。"
        "0.6 purity: B1/B3/B5 三版 TP=24,190 / FP=13,487 / FN=15,257 完全相同, F1=0.6273。"
        "因果鏈: ClairS-TO PASS set 由 caller 決定, longphase-to phase 改 GT/PS/GT2/GT3 不改 FILTER → PASS set 不變 → TP/FP/FN 不變 → F1 不變。"
        "V5 不改 caller; ΔF1 (0.93→0.6)=-0.0893 是 ClairS-TO 在低 purity 本身偵測下降。"
        "正確 metric 是 read-level tag concordance (paired GT +13.3 pp)。"
        "turning point — 但 5/9 paired cross-ref 揭露另一面, 我們將進入 V5 Layer 1.5 設計缺陷。",
        "PASS set / FILTER 機制細節 / purity 0.6 N50 微差")

def slide_15_paired_mode():
    s = new_slide()
    add_title_block(s,
        "paired mode 整體無偏移 — HP1:HP2 = 1:1.275; som_ratio mean 0.462",
        "paired mode no systematic bias")

    # paired chr19 distribution
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(0.4),
        "paired chr19 HP:Z: 分布 (354,919 tagged):", size=11, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(1.95), Inches(6.0), Inches(2.0),
        ["HP:Z:", "reads", "%"],
        [
            ["HP:Z:2",   "183,309", "51.6%"],
            ["HP:Z:1",   "143,760", "40.5%"],
            ["HP:Z:2-1", "14,504",  "4.1%"],
            ["HP:Z:1-1", "12,401",  "3.5%"],
            ["HP:Z:3",   "1,145",   "0.3%"],
        ])

    # paired vs TO ratio
    add_textbox(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.4),
        "paired vs TO ratio:", size=11, bold=True, color=NAVY)
    add_table_box(s, Inches(6.9), Inches(1.95), Inches(6.0), Inches(1.4),
        ["", "paired", "TO baseline"],
        [
            ["germline", "1:1.275 ✅", "17.3:1 ❌ priority bug"],
            ["somatic",  "1:1.169 ✅", "全偏 HP1"],
        ],
        highlight_rows={0: LIGHT_GREEN, 1: LIGHT_GREEN})

    # 57 windows
    add_textbox(s, Inches(6.9), Inches(3.6), Inches(6.0), Inches(1.0),
        "57 chr19 1Mb windows:\n"
        "som_ratio mean 0.462 / median 0.494 / stdev 0.332",
        size=11, color=DARK_GRAY, fill=LIGHT_BLUE, border=NAVY)

    # F6 paired vs TO HP distribution (bottom-left)
    add_image_fit(s, MASTER / "F6_paired_vs_TO_HP_distribution.png",
                  Inches(0.4), Inches(4.3), Inches(7.0), Inches(1.9))

    # Real signal cases (bottom-right)
    add_textbox(s, Inches(7.7), Inches(4.3), Inches(5.4), Inches(0.3),
        "真實 sub-clone signal cases:", size=11, bold=True, color=NAVY)
    add_textbox(s, Inches(7.7), Inches(4.6), Inches(5.4), Inches(1.6),
        "•  3M 全 HP2-1 (755/0) LOH\n"
        "•  0M 全 HP1-1 (330/1) 反向\n"
        "•  17M 對稱 0.500 (265/265)\n"
        "    → SP1 附近 paired 認雙\n"
        "    sub-clone (vs TO 113:0)",
        size=10, color=DARK_GRAY, fill=OFF_WHITE, border=GRAY)

    # paired uses different binary
    add_textbox(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.5),
        "paired = longphase-s (獨立 codebase, 非 longphase-to fork); HP tag 用字串 HP:Z: 編碼",
        size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_glossary_footer(s, [
        "📖 HP:i: vs HP:Z:: longphase-to=整數 (1/2/11/21/33), longphase-s=字串",
        "ⓘ som_ratio: HP1-1 / (HP1-1+HP2-1) 票數"])

    add_speaker_note(s,
        "paired mode 用不同 binary longphase-s (非 longphase-to fork), HP tag 用 HP:Z: 字串編碼。"
        "chr19 paired 分布 HP:Z:2 51.6% / HP:Z:1 40.5% germline 1:1.275 接近隨機; "
        "somatic HP:Z:2-1 4.1% / HP:Z:1-1 3.5% 也 1:1.169。"
        "對比 TO baseline 17.3:1 priority bug — paired mode 整體無 systematic bias。"
        "57 windows som_ratio mean 0.462 / median 0.494 / stdev 0.332 跨 0-1 全範圍 = 真實 sub-clone signal。"
        "chr19:3M LOH 方向; chr19:0M 反向; chr19:17M 對稱 0.500 SP1 附近 paired 認雙 sub-clone (vs TO 113:0 失衡)。",
        "longphase-s codebase / SomaticHaplotagProcess.cpp:533 / paired 軸對齊")

def slide_16_v5_caveat():
    s = new_slide()
    add_title_block(s,
        "V5 Layer 1.5 設計缺陷: germline-absent 區與 baseline 4.19:1 完全相同",
        "V5 Layer 1.5 design caveat — identical to baseline in germline-absent")
    # Override title color
    # (already set via add_title_block, but emphasize via red)

    # G4 figure (centerpiece)
    add_image_fit(s, FIGURES / "G4_germline_absent_three_versions.png",
                  Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.4))

    # Conclusion red frame
    add_red_frame(s, Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.85),
        "★ 結論:",
        "V5 Layer 1.5 = priority bug 的 feature 化非修補 — V3F 標 hp=33 反而更穩健",
        "V5 Layer 1.5 makes priority bug a feature, not a fix.")

    add_glossary_footer(s, [
        "📖 sub-clone: 腫瘤內基因型相同細胞群",
        "ⓘ scope: paired chr19 germline-absent (5,789 events ≈ 5% chr19)"])

    add_speaker_note(s,
        "整份報告最重要的一張 slide。"
        "對 paired chr19 read_name × T1.2 baseline/V3F/V5 vote dump JOIN, 篩 cnt_HP1+HP2=0 且 somatic>0 = 5,789 chr19 events。"
        "cross-tab: baseline hp=11 (HP1) 3,312 / hp=21 (HP2) 791 → 4.19:1 priority bug 次峰。"
        "V3F 全 5,789 標 hp=33 保守不選邊 ✅。"
        "V5 hp=11=3,313 / hp=21=790 → 4.19:1 與 baseline 完全相同！"
        "機制詮釋: V5 Layer 1.5 germline=0 用 sHP1 vs sHP2 票數決方向; sub-clone somatic 100% 共現 → graph 偏向同一 haplotype → 投票偏向同邊 → Layer 1.5 結果繼承 priority bug 偏移。"
        "V5 Layer 1.5 = priority bug 的 feature 化非修補; V3F hp=33 是更穩健選擇。"
        "V5 設計時未對 paired ground truth 做 germline-absent cross-ref, 5/9 paired audit Step D 才補上發現。",
        "cross-binary axis alignment / phasing graph §3.2 / Layer 1.5 改回 V3F default / F-paired-D3 工作量")

def slide_17_errata():
    s = new_slide()
    add_title_block(s,
        "PI 報告 4-29 5 條 errata 已 patch; 主結論不撤回",
        "5 errata patched; main conclusions retained")

    erratas = [
        ("E1", "§3.3.3 chr19 SP1/2/3 解讀降級",
         "原: 主要 hotspot → 新: 可重現案例 (chr19 占 priority bug 2.16% rank 19)",
         OFF_WHITE),
        ("E2", "§5.2 V5 working tree commit 狀態",
         "原: 未 commit → 新: ✅ 2026-04-30 已 commit (d0bcd8c + 938f0df)",
         OFF_WHITE),
        ("E3", "§5.2 priority bug 證據強度升級",
         "原: commit msg + 3 IGV → 新: + 34,855 read-level 鐵證 V3F+V5 修正率 100%",
         OFF_WHITE),
        ("★ E4", "§6.4/§6.5 V5 數值歸因精確化 (最重要 errata)",
         "原: V5 four-commit 整體效益 → 新: V5 BAM = Pass 1 only (ploidy bug); "
         "主要 V3F + Layer 1.5; Pass 2 二次效益尚未量化",
         LIGHT_YEL),
        ("★ E5", "§5.2 V5 Layer 1.5 設計 (5/10 加)",
         "原: Layer 1.5 = germline 缺席 fallback 隱含修補 → 新: "
         "germline-absent 區與 baseline 4.19:1 完全相同, priority bug feature 化非修補; V3F hp=33 反而穩健",
         LIGHT_YEL),
    ]
    for i, (eid, desc, content, fill) in enumerate(erratas):
        y = Inches(1.5 + i * 0.95)
        add_textbox(s, Inches(0.4), y, Inches(2.0), Inches(0.85),
            f"{eid}", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
            fill=fill, border=GRAY)
        add_textbox(s, Inches(2.5), y, Inches(10.4), Inches(0.85),
            f"{desc}\n{content}",
            size=10, color=DARK_GRAY, fill=fill, border=GRAY)

    add_glossary_footer(s, ["ⓘ commit chain: f17754f → 2553e96 → 71d21bd"])

    add_speaker_note(s,
        "PI 報告 4-29 5 條 errata patched, 主結論不撤回。"
        "E1 chr19 SP1/2/3 從主要 hotspot 改為可重現案例 (占 2.16% rank 19); "
        "E2 V5 已 commit d0bcd8c + 938f0df 2026-04-30; "
        "E3 priority bug 證據從 commit msg + IGV 升級為 34,855 read-level 鐵證 100%; "
        "★ E4 V5 數值歸因精確化 — V5 BAM = Pass 1 only (ploidy bug 讓 purity=0), 主要 V3F + Layer 1.5, Pass 2 二次效益尚未量化; "
        "★ E5 5/10 加 — V5 Layer 1.5 germline-absent 與 baseline 4.19:1 完全相同, priority bug feature 化非修補, V3F hp=33 反而穩健。"
        "E5 來源 5/9 paired audit Step D, 5/10 amend。E1-E3 表述精確化, E4+E5 核心。",
        "各 errata commit hash / errata patch 工作量 / 為何不撤回 vs 補 banner")

def slide_18_followup():
    s = new_slide()
    add_title_block(s,
        "整體成熟度 + 5 follow-up — V5 仍可作 production baseline",
        "Maturity status + 5 follow-up cycles")

    # Maturity grid (12 dimensions, abbreviated)
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4),
        "整體成熟度: 10 ✅ / 1 ⚠️ / 2 ⏸",
        size=13, bold=True, color=NAVY)

    maturity = [
        "✅ 機制因果", "✅ 修補設計", "✅ chr19 SP 對齊",
        "✅ 全基因組擴展", "✅ V5/V3F zero-sum",  "✅ 20 指標 0 regression",
        "✅ Caller F1 三版相同", "✅ purity 0.6 對照", "✅ 三路徑算法",
        "✅ Pass 2 量化 +3.51%", "✅ 版本對齊 938f0df", "✅ Paired Step A+C NEGATIVE",
        "⚠ V5 Layer 1.5 設計缺陷 (E5)",
        "⏸ Pass 2 second 獨立貢獻 (T1.3)",
        "⏸ 跨樣本 (T3)",
    ]
    for i, m in enumerate(maturity):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.2)
        y = Inches(2.0 + row * 0.45)
        if "✅" in m:
            fill = LIGHT_GREEN
        elif "⚠" in m:
            fill = LIGHT_YEL
        else:
            fill = OFF_WHITE
        add_textbox(s, x, y, Inches(4.0), Inches(0.4),
            m, size=10, color=DARK_GRAY, fill=fill, border=GRAY)

    # Follow-up table
    add_textbox(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.4),
        "5 follow-up cycle (依 ROI 排序):", size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.5),
        ["ID", "內容", "預估"],
        [
            ["★ F-paired-D3", "Layer 1.5 改 V3F ISM 影響評估 (決定 V5 vs V3F default)", "1-2 day"],
            ["F-paired-D1",   "germline-absent 全基因組擴展", "0.5 day"],
            ["F-paired-D2",   "phase block 內 axis-aligned 分析", "1 day"],
            ["T3",            "7 樣本跨樣本擴展", "1-2 day"],
            ["T1.3",          "4-cell ablation (Pass 2 second 獨立貢獻)", "3 day"],
        ],
        highlight_rows={0: LIGHT_YEL})

    add_textbox(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
        "→ V5 仍可作 production baseline; F-paired-D3 量化後決定是否回 V3F default",
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)

    add_glossary_footer(s, ["📖 LOH: 雜合性丟失 / cnLOH: 拷貝中性 LOH"])

    add_speaker_note(s,
        "整體成熟度 12 維度: 10 ✅ + 1 ⚠️ V5 Layer 1.5 germline-absent 設計缺陷 + 2 ⏸ 待跑 (T1.3 Pass 2 second 獨立貢獻 + T3 跨樣本)。"
        "5 follow-up 排序按 ROI: ★ F-paired-D3 1-2 day 最重要 (決定 V5 vs V3F default); "
        "F-paired-D1 0.5 day; F-paired-D2 1 day; T3 1-2 day; T1.3 3 day。"
        "V5 仍可作 production tag baseline; germline-absent 區占 ~5% 不阻擋整體; "
        "F-paired-D3 ISM 影響量化後決定是否回歸 V3F default — 待 PI 決策觸發。",
        "T1.3 4-cell ablation / 7 樣本 binary patch / cnLOH 雙親同源")

# ─── Q&A Backup ─────────────────────────────────────────────────────────

def slide_b1_pass2():
    s = new_slide()
    add_title_block(s,
        "Pass 2 = 只重跑 2-point edgeConnectResult; 高 purity 才觸發",
        "Pass 2 = re-run 2-point only; high purity (>0.9) gate")

    # G2 figure
    add_image_fit(s, FIGURES / "G2_pass_two_flow.png",
                  Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.0))

    add_glossary_footer(s, [
        "📖 somaticCalling vs edgeConnectResult: 兩 phasing graph 演算法",
        "📖 purity: 樣本中腫瘤細胞占比"])

    add_speaker_note(s,
        "Q: Pass 2 為什麼只跑 2-point 不重跑 3-point?"
        "兩個分別函式: somaticCalling 用 3-point patternMining (與 purity 無關); edgeConnectResult 用 2-point 永遠跑。"
        "低 purity ≤ 0.9: Pass 1 跑 2+3-point; Pass 2 跳過。"
        "高 purity > 0.9: Pass 1 跑 2+3-point; Pass 2 只重跑 2-point (Pass 1 已產出穩定 origin 分類)。"
        "Pass 2 incremental: phased var -2.90 pp / blocks -9.79% / N50 +3.51%。Pass 2 是 polish/merge。"
        "常見誤解: '低 purity 用 3-point' 倒過來; 高 purity 才多做事 = Pass 2 多 2-point。",
        "patternMining first/second/third path / Pass 2 不重跑 somaticCalling 原因")

def slide_b2_purity06():
    s = new_slide()
    add_title_block(s,
        "purity 0.6 樣本 baseline vs V5 完整對照 — 0 critical regression",
        "purity 0.6 sample fully verified")

    # 6 caller F1
    add_textbox(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
        "6 Caller F1 (vs SEQC2 v1.2.1) — 三版完全相同:",
        size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(1.85), Inches(12.5), Inches(1.6),
        ["指標", "baseline 0.6", "V5 0.6", "Δ"],
        [
            ["TP",         "24,190", "24,190", "0 ✅"],
            ["FP",         "13,487", "13,487", "0 ✅"],
            ["FN",         "15,257", "15,257", "0 ✅"],
            ["F1",         "0.6273", "0.6273", "0 ✅"],
        ],
        highlight_rows={i: LIGHT_GREEN for i in range(4)})

    # 9 structure
    add_textbox(s, Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.4),
        "9 結構指標 — 4 改善 + 1 微差 + 1 持平:",
        size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.4),
        ["指標", "baseline", "V5", "Δ", "eval"],
        [
            ["phased%",   "61.82",  "65.83",  "+4.01 pp",  "✅"],
            ["n_blocks",  "9,748",  "11,514", "+18.1%",    "✅"],
            ["N50 (bp)",  "798,903","683,296","-14.5%",    "微差 (≥600K)"],
            ["HP:i:33",   "0",      "20",     "+20",       "✅ conservative"],
            ["AMB%",      "0.00",   "3.12",   "+3.12 pp",  "✅"],
            ["purity",    "0.607",  "0.634",  "+0.027",    "✅ closer to 0.6"],
        ],
        highlight_rows={2: LIGHT_YEL})

    add_glossary_footer(s, ["📖 purity: 樣本中腫瘤細胞占比"])

    add_speaker_note(s,
        "Q: 低純度樣本是否變差? "
        "purity 0.6 樣本 baseline vs V5 完整對照 (HCC1395 t30_n20)。"
        "6 caller F1 全完全相同 (V5 不改 caller)。"
        "9 結構指標: phased% +4.01 pp / n_blocks +18.1% / HP:i:33 +20 (V5 conservative) / "
        "AMB% +3.12 pp / HP1/HP2 0.48→0.38 修正 / purity 0.607→0.634 更接近真實 0.6 / "
        "Pass 2 兩者都不觸發 / LOH 完全相同; 唯一 N50 -14.5% 但仍 ≥600 K。"
        "0 critical regression; ploidy bug 在低 purity 自我治癒。",
        "V3F vs V5 0.6 對比 / N50 為何接受 / Pass 2 不針對低 purity")

def slide_b3_cross_sample():
    s = new_slide()
    add_title_block(s,
        "跨樣本擴展 (T3) + cnLOH 雙親同源待開放",
        "Cross-sample expansion (T3) + cnLOH bi-parental")

    # T3 sample list
    add_textbox(s, Inches(0.4), Inches(1.5), Inches(8.0), Inches(0.4),
        "T3 跨樣本擴展工作:", size=12, bold=True, color=NAVY)
    add_table_box(s, Inches(0.4), Inches(2.0), Inches(8.0), Inches(3.0),
        ["樣本", "狀態", "預估"],
        [
            ["HCC1395 5kHz",        "✅ 已驗證 (本報告)", "完成"],
            ["HCC1395_DORADO",     "⏸ 待跑",          "~3 hr"],
            ["HCC1937",             "⏸ 待跑",          "~3 hr"],
            ["HCC1954",             "⏸ 待跑",          "~3 hr"],
            ["H1437",               "⏸ 待跑",          "~3 hr"],
            ["H2009",               "⏸ 待跑",          "~3 hr"],
            ["COLO829",             "⏸ 待跑",          "~3 hr"],
        ],
        highlight_rows={0: LIGHT_GREEN})

    # cnLOH
    add_textbox(s, Inches(8.7), Inches(1.5), Inches(4.4), Inches(0.4),
        "cnLOH 雙親同源待開放:", size=12, bold=True, color=NAVY)
    add_textbox(s, Inches(8.7), Inches(2.0), Inches(4.4), Inches(3.0),
        "情境: parent 1 vs parent 2 染色體\n在 LOH 區難 phase\n\n"
        "影響: Layer 1.5 在 cnLOH 區的設計選擇待量化\n\n"
        "關係: 與 V5 Layer 1.5 設計缺陷 (E5) 連動",
        size=11, color=DARK_GRAY, fill=LIGHT_BLUE, border=NAVY)

    # Priority sorting
    add_textbox(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.5),
        "Follow-up 排序: F-paired-D3 (1-2 d, actionable) > T3 (1-2 d, generalizability) > T1.3 (3 d, ablation)",
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)
    add_textbox(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.4),
        "T3 總工作: 1-2 day (含 binary patch + dump 6 × 3 版本)",
        size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_glossary_footer(s, [
        "📖 LOH: 雜合性丟失",
        "📖 cnLOH: 拷貝中性 LOH (保留 2 套但同源)"])

    add_speaker_note(s,
        "Q: 跨樣本一致性? cnLOH 區? "
        "T3 跨樣本擴展: HCC1395 5kHz 已驗證為本報告主案例; "
        "6 樣本待跑 — HCC1395_DORADO / HCC1937 / HCC1954 / H1437 / H2009 / COLO829 — 每樣本 ~3 hr; T3 共 1-2 day。"
        "cnLOH 雙親同源: parent 1 vs parent 2 染色體在 LOH 區因物理上同源難 phase; Layer 1.5 設計選擇對 cnLOH 區影響待量化, 與 E5 連動。"
        "Follow-up 排序: F-paired-D3 1-2 d 最重要 actionable (V5 vs V3F default); T3 1-2 d generalizability check; T1.3 3 d ablation。",
        "7 樣本各自特性 / cnLOH 機制 / F-paired-D3 詳細")


# ═══════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════

def main():
    print("[build] Self-Phasing PPT — 25 slide build_pptx.py")
    print(f"[build] Output: {OUT}")
    print(f"[build] Figures dir: {FIGURES}")
    print()

    # Main slides 1-18
    slide_01_cover()
    slide_02_tldr()
    slide_03_genome_173()
    slide_04a_sp1()
    slide_04b_sp2_sp3()
    slide_05_player_referee()
    slide_06_priority_bug()
    slide_07_two_layer_table()
    slide_08_chr19_752()
    slide_09_genome_34855()
    slide_10_5_commits()
    slide_11_getvote_3versions()
    slide_12_sp_fixed()
    slide_13_20_metrics()
    slide_14_caller_f1_cliffhanger()
    slide_15_paired_mode()
    slide_16_v5_caveat()
    slide_17_errata()
    slide_18_followup()

    # Backup B1-B3
    slide_b1_pass2()
    slide_b2_purity06()
    slide_b3_cross_sample()

    n = len(prs.slides)
    print(f"[build] Total slides: {n}")
    prs.save(OUT)
    sz = OUT.stat().st_size
    print(f"[ok] Saved: {OUT} ({sz / 1024:.1f} KB)")


if __name__ == "__main__":
    main()
