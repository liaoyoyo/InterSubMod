#!/usr/bin/env python3
"""extract_pptx.py — 從 PPTX 抽「名詞/物件/排版/字型」供 methods-example 慣例學習.

輸出 per-deck:
  <out>/<deck>/slides.json   每 slide: 尺寸 + shapes(type, bbox inch, text runs[text/pt/bold/color])
  <out>/<deck>/text.md       人類可讀 per-slide 文字 + 字型摘要 + 圖清單
  <out>/<deck>/media/        unzip 出的嵌入圖（PNG/SVG/…）
不渲染整頁（無 libreoffice）；嵌入媒體圖直接抽出供 vision/SVG 分析.
用法: python3 extract_pptx.py "<a.pptx>" "<b.pptx>" -o <out_dir>
"""
import argparse, json, os, re, subprocess, sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0


def slug(p):
    b = os.path.splitext(os.path.basename(p))[0]
    return re.sub(r"[^A-Za-z0-9]+", "_", b).strip("_").lower()[:40] or "deck"


def color_of(run):
    try:
        c = run.font.color
        if c and c.type is not None:
            if str(c.type) == "MSO_THEME_COLOR.NOT_THEME_COLOR" or c.rgb is not None:
                return "#" + str(c.rgb)
            return f"theme:{c.theme_color}"
    except Exception:
        pass
    return None


def emu(v):
    return round(v / EMU_IN, 2) if v is not None else None


def shape_dict(sh):
    d = {"name": getattr(sh, "name", ""), "type": str(getattr(sh, "shape_type", "")),
         "bbox_in": [emu(sh.left), emu(sh.top), emu(sh.width), emu(sh.height)]}
    runs = []
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.text.strip():
                    runs.append({"t": r.text, "pt": (r.font.size.pt if r.font.size else None),
                                 "b": bool(r.font.bold), "color": color_of(r)})
    if runs:
        d["runs"] = runs
    if sh.shape_type == 13:  # PICTURE
        try:
            d["image"] = os.path.basename(sh.image.filename or "") or f"blob.{sh.image.ext}"
            d["image_ext"] = sh.image.ext
        except Exception:
            d["image"] = "?"
    if getattr(sh, "has_table", False) and sh.has_table:
        d["table"] = [[c.text for c in row.cells] for row in sh.table.rows]
    return d


def extract(pptx_path, out):
    sl = slug(pptx_path)
    ddir = os.path.join(out, sl)
    os.makedirs(os.path.join(ddir, "media"), exist_ok=True)
    prs = Presentation(pptx_path)
    W, H = emu(prs.slide_width), emu(prs.slide_height)
    slides = []
    md = [f"# {os.path.basename(pptx_path)}", f"slide size: {W} x {H} in · {len(prs.slides)} slides\n"]
    for i, slide in enumerate(prs.slides, 1):
        shapes = [shape_dict(s) for s in slide.shapes]
        slides.append({"n": i, "shapes": shapes})
        md.append(f"\n## slide {i}")
        for s in shapes:
            if "runs" in s:
                for r in s["runs"]:
                    flag = "**" if r["b"] else ""
                    pt = f"{r['pt']}pt" if r["pt"] else "·"
                    md.append(f"- [{pt}{' B' if r['b'] else ''}{(' '+r['color']) if r['color'] else ''}] {flag}{r['t']}{flag}")
            if "image" in s:
                md.append(f"  - 🖼 IMAGE: {s['image']}  bbox(in)={s['bbox_in']}")
            if "table" in s:
                md.append(f"  - 📋 TABLE {len(s['table'])}x{len(s['table'][0]) if s['table'] else 0}: {s['table']}")
    json.dump({"deck": os.path.basename(pptx_path), "slide_size_in": [W, H], "slides": slides},
              open(os.path.join(ddir, "slides.json"), "w"), ensure_ascii=False, indent=1)
    open(os.path.join(ddir, "text.md"), "w").write("\n".join(md))
    # unzip media
    subprocess.run(["unzip", "-o", "-j", pptx_path, "ppt/media/*", "-d", os.path.join(ddir, "media")],
                   capture_output=True)
    nmedia = len(os.listdir(os.path.join(ddir, "media")))
    # font size histogram
    fonts = {}
    for s in slides:
        for sh in s["shapes"]:
            for r in sh.get("runs", []):
                if r["pt"]:
                    fonts[r["pt"]] = fonts.get(r["pt"], 0) + 1
    print(f"[OK] {sl}: {len(slides)} slides, {nmedia} media → {ddir}")
    print(f"     font pt histogram: {dict(sorted(fonts.items(), reverse=True))}")
    return {"slug": sl, "dir": ddir, "slides": len(slides), "media": nmedia, "fonts": fonts, "size_in": [W, H]}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", nargs="+")
    ap.add_argument("-o", "--out", required=True)
    a = ap.parse_args()
    os.makedirs(a.out, exist_ok=True)
    summary = [extract(p, a.out) for p in a.pptx]
    json.dump(summary, open(os.path.join(a.out, "_summary.json"), "w"), ensure_ascii=False, indent=1)
    print(f"\n[DONE] summary → {os.path.join(a.out, '_summary.json')}")


if __name__ == "__main__":
    main()
