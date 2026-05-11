#!/usr/bin/env python3
"""build_pptx.py — Self-phasing V5 provenance audit weekly report (18 slides).

Source: master_draft.md §10 (18-slide architecture, improvement_report template).
Audience: PI (advisor), 25 min target.
Main thesis (≤30 字): V5 為 5 commits, PI 數據為 Pass 1 only, Pass 2 重驗 P0.

Output:
    output.pptx                 — 18-slide deck
    speaker_script.md           — full speaker note dump for review

Run:
    cd .../pptx_build
    python3 build_pptx.py
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

# Resolve repo root for ppt_toolkit + style_library
HERE = Path(__file__).resolve().parent
REPO_ROOT = HERE
for _ in range(8):
    if (REPO_ROOT / "tools" / "ppt_toolkit").is_dir():
        break
    REPO_ROOT = REPO_ROOT.parent
sys.path.insert(0, str(REPO_ROOT / "tools"))

# Point loader to pptx-build style library (the live one as of 2026-05-05).
os.environ["MYPPT_STYLE_LIBRARY"] = str(
    REPO_ROOT / ".claude" / "skills" / "pptx-build" / "style_library"
)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ppt_toolkit import (
    PALETTE,
    add_text_with_fallback,
    fit_image_within,
    set_speaker_note,
)

# Figure assets directory
FIGURES_DIR = HERE / "figures"

# ---------------------------------------------------------------------------
# Canvas constants (16:9)
# ---------------------------------------------------------------------------
CANVAS_W = Inches(13.333)
CANVAS_H = Inches(7.5)
TITLE_Y = Inches(0.4)
TITLE_H = Inches(0.85)
TITLE_X = Inches(0.5)
TITLE_W = Inches(12.333)

# Tier S/A/B/C colour map (used for focal-point markers)
TIER_S_COLOR = PALETTE["accent"]   # red
TIER_A_COLOR = PALETTE["green"]
TIER_B_COLOR = PALETTE["blue"]
TIER_NEUTRAL = PALETTE["gray"]


# ---------------------------------------------------------------------------
# Reusable primitives
# ---------------------------------------------------------------------------
def add_assertion_title(slide, text, *, color_bg=None):
    """Thesis-style banner across the top — heading IS the conclusion."""
    bg = color_bg or PALETTE["title_bg"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, text, font_size=24, bold=True, color=PALETTE["title_fg"], paragraph_index=0
    )
    return bar


def add_textbox(slide, x, y, w, h, lines, *, font_size=16, bold=False, color=None,
                align=PP_ALIGN.LEFT, fill=None, border=None):
    """A simple wrapper that draws a (possibly multi-line) textbox."""
    if fill is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if border is not None:
            shp.line.color.rgb = border
            shp.line.width = Pt(1.0)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
    else:
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.paragraphs[0].text = ""
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        add_text_with_fallback(
            tf, line, font_size=font_size, bold=bold, color=color or PALETTE["text"],
            align=align, paragraph_index=i,
        )


def add_caveat_strip(slide, x, y, w, text):
    """Red-banded caveat strip (Tier S/A — surface caveat ON slide)."""
    h = Inches(0.55)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    shp.line.color.rgb = TIER_S_COLOR
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, text, font_size=14, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C),
        paragraph_index=0,
    )


def add_insight_card(slide, x, y, w, h, lines, *, title=None):
    """Green-banded insight card (Tier A — supports thesis)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shp.line.color.rgb = TIER_A_COLOR
    shp.line.width = Pt(1.0)
    tf = shp.text_frame
    tf.margin_left = Emu(100000)
    tf.margin_right = Emu(100000)
    tf.margin_top = Emu(60000)
    tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    idx = 0
    if title:
        add_text_with_fallback(
            tf, title, font_size=15, bold=True, color=RGBColor(0x1B, 0x5E, 0x20),
            paragraph_index=idx,
        )
        idx += 1
    for line in lines:
        add_text_with_fallback(
            tf, line, font_size=13, color=PALETTE["text"], paragraph_index=idx,
        )
        idx += 1


def add_neutral_box(slide, x, y, w, h, lines, *, title=None, accent=None, body_font_size=12):
    """Neutral grey box (Tier B — context / definition)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shp.line.color.rgb = accent or TIER_NEUTRAL
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Emu(100000)
    tf.margin_right = Emu(100000)
    tf.margin_top = Emu(60000)
    tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    idx = 0
    if title:
        add_text_with_fallback(
            tf, title, font_size=14, bold=True, color=PALETTE["title_bg"],
            paragraph_index=idx,
        )
        idx += 1
    for line in lines:
        add_text_with_fallback(
            tf, line, font_size=body_font_size, color=PALETTE["text"], paragraph_index=idx,
        )
        idx += 1


def add_method_box(slide, x, y, w, h, lines, *, title=None):
    """Blue-banded method/mechanism box (Tier A — supports thesis with method/info)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shp.line.color.rgb = TIER_B_COLOR
    shp.line.width = Pt(1.0)
    tf = shp.text_frame
    tf.margin_left = Emu(100000)
    tf.margin_right = Emu(100000)
    tf.margin_top = Emu(60000)
    tf.margin_bottom = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    idx = 0
    if title:
        add_text_with_fallback(
            tf, title, font_size=14, bold=True, color=RGBColor(0x0D, 0x47, 0xA1),
            paragraph_index=idx,
        )
        idx += 1
    for line in lines:
        add_text_with_fallback(
            tf, line, font_size=13, color=PALETTE["text"], paragraph_index=idx,
        )
        idx += 1


def add_kpi_block(slide, x, y, w, h, label, value, *, value_color=None, sub=None):
    """Big-number KPI block."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = PALETTE["light"]
    shp.line.color.rgb = PALETTE["border"]
    shp.line.width = Pt(0.5)
    tf = shp.text_frame
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(80000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, label, font_size=12, color=PALETTE["gray"], align=PP_ALIGN.CENTER,
        paragraph_index=0,
    )
    add_text_with_fallback(
        tf, value, font_size=28, bold=True, color=value_color or PALETTE["title_bg"],
        align=PP_ALIGN.CENTER,
    )
    if sub:
        add_text_with_fallback(
            tf, sub, font_size=11, color=PALETTE["gray"], align=PP_ALIGN.CENTER,
        )


def add_footer_citation(slide, text):
    """Bottom-right citation footnote, 9pt grey (myPPT §16)."""
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.35)
    )
    tf = box.text_frame
    tf.margin_left = Emu(20000)
    tf.margin_right = Emu(20000)
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, text, font_size=9, color=PALETTE["gray"], align=PP_ALIGN.RIGHT,
        paragraph_index=0,
    )


def add_focal_marker(slide, text):
    """Top-right focal point chip — confirms slide passed §20.B."""
    w, h = Inches(3.5), Inches(0.32)
    # 距右緣 0.1 in（從 0.5 in 縮，符合 §16 對齊規則）
    x = CANVAS_W - w - Inches(0.1)
    y = Inches(0.05)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shp.line.color.rgb = PALETTE["amber"]
    shp.line.width = Pt(0.5)
    tf = shp.text_frame
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, f"[Focal] {text}", font_size=10, bold=True,
        color=RGBColor(0xE6, 0x5C, 0x00), align=PP_ALIGN.CENTER, paragraph_index=0,
    )


# ---------------------------------------------------------------------------
# Per-slide builders
# ---------------------------------------------------------------------------
def slide_01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Big title — 上移 0.3 in 讓 Main Thesis 框中心落在 ~50%
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(12.0), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf,
        "Self-phasing V5 provenance 校正",
        font_size=44, bold=True, color=PALETTE["title_bg"], paragraph_index=0,
    )
    add_text_with_fallback(
        tf,
        "與 Pass 2 重驗 P0",
        font_size=44, bold=True, color=PALETTE["title_bg"],
    )
    # Thesis bar — 上移配合 title 上移
    thesis_bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.7), Inches(11.9), Inches(1.9)
    )
    thesis_bg.shadow.inherit = False
    thesis_bg.fill.solid()
    thesis_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    thesis_bg.line.color.rgb = PALETTE["amber"]
    thesis_bg.line.width = Pt(2.0)
    tf = thesis_bg.text_frame
    tf.margin_left = Emu(160000)
    tf.margin_right = Emu(160000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, "Main Thesis", font_size=14, bold=True, color=PALETTE["amber"],
        paragraph_index=0,
    )
    add_text_with_fallback(
        tf,
        "V5 為 5 commits，PI 數據為 Pass 1 only，Pass 2 重驗 P0",
        font_size=22, bold=True, color=PALETTE["text"],
    )
    # Meta + Mini-glossary
    box = s.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf,
        "週報 · 2026-04-26 ~ 2026-05-05 · 受眾: PI · 預設時長: 25 min · 18 slides",
        font_size=14, color=PALETTE["gray"], paragraph_index=0,
    )
    # Mini-glossary：PI 第一次看會碰到的術語
    add_text_with_fallback(
        tf,
        "術語：Pass 2 = phasing 的 second-round 計算（4/30 才能觸發）｜ P0 = 最高優先級",
        font_size=11, color=PALETTE["gray"],
    )
    set_speaker_note(
        s,
        "開場 30 秒抓 main thesis：本週最關鍵的發現是 5/05 provenance audit 揭露"
        "4/29 PI 報告的全部 V5 數值都源自 4/12 產的 BAM，那份 BAM 因 ploidy bug 讓 "
        "purity=0、Pass 2 從未觸發。所以「V5 為當前最佳 phasing baseline」的主結論"
        "需暫停為「Pass 1 only 條件下觀察」，Pass 2 完整重驗是本週最高優先級 (P0)。"
        "同時 5/01 三條路 audit + force_path2only ablation 實證「路 3 second round 抵消"
        "路 2 反轉」，新 baseline self-phasing 反而更甚。本次彙報聚焦兩件 critical findings"
        " 與下一步決策。\n\n"
        "[ORAL-OPTIONAL] 若教授問 deliverable 數量，可補：本週共 9 份 deliverable，"
        "包含 4/28 V5 audit、4/29 技術報告、4/30 V3F ablation × 2、5/01 三條路 audit、"
        "5/01 force_path2only ablation、5/05 provenance audit。",
        min_chars=300,
    )
    add_focal_marker(s, "Main thesis 鎖定")
    return s


def slide_02_background(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Thread D 主軸切換已完成，本週聚焦 V5 audit chain 校正")
    # Two columns
    add_neutral_box(
        s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5),
        [
            "• 4/26 Thread D (LOH-constrained phasing) 確認為主軸",
            "• 4/26 Thread B (NG=2 LOH-AF interaction) 撤回",
            "• 4/26-5/05 預定推進 V5 audit chain",
            "• 25 hr Pass 2 重驗成為本週唯一決定性投入",
        ],
        title="上週起點 (Memory + 4/23 週報)",
    )
    add_insight_card(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(2.5),
        [
            "[OK] 4/28 V5 audit haplotag sanity 通過",
            "[OK] 4/29 longphase TO vs V5 技術報告",
            "[OK] 4/30 V3F ablation × 2 跨 purity 確認",
            "[!] 5/05 critical finding: PI 數據 Pass 1 only",
            "[!] Pass 2 重驗尚未啟動 (P0)",
        ],
        title="本週進展 + 紅色 caveat",
    )
    # Sub-thread context
    add_method_box(
        s, Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.5),
        [
            "主線 (problem)：4/29 PI 報告 V5 結論基於 Pass 1 only 條件，需暫停為「Pass 1 only 觀察」",
            "Sub-thread (progress)：4/26-5/05 共 9 份 deliverable 完成",
            "教授視角優先序：主線追問 (caveat + Pass 2 重驗) > Sub-thread 進度 (commit 鏈 + 三條路機制)",
        ],
        title="本週敘事框架 (混合主線)",
    )
    add_footer_citation(s, "Source: master_draft.md §1, §2, Memory project_v5_somatic_fallback_verification")
    set_speaker_note(
        s,
        "本週上下文先架構好。上週週報明確設定三件事：Thread D 主軸切換、Thread B 撤回、"
        "V5 audit chain 推進。本週前 5 天大致按計畫完成，9 份 deliverable 涵蓋 V5 audit / "
        "技術報告 / V3F ablation / 三條路 audit / force_path2only ablation。但 5/05 做"
        "provenance audit 時發現一個 critical issue：4/29 PI 報告引用的 V5 數據實際來自"
        "4/12 那份 BAM，那份 BAM 是 Pass 1 only 結果，所以「V5 為當前最佳」這個結論的"
        "現實基礎需要重新確認。\n\n"
        "敘事框架是混合主線：主線是 problem (V5 caveat)，但 sub-thread 是 progress "
        "(9 份 deliverable)，所以結尾仍能 actionable。教授視角優先序我建議先聽完 critical "
        "finding 與 Pass 2 重驗計畫，再回看其他進度。\n\n"
        "[ORAL-OPTIONAL] 若教授追問為何 4/29 報告未及時發現 Pass 1 only，可解釋：當時"
        "log 只看「flag 是否觸發」沒看「purity 是否 > 0.9」，這是 5/05 audit 才補上的盲點。",
        min_chars=300,
    )
    add_focal_marker(s, "上週背景 + 本週主線")
    return s


def slide_03_v5_commit_chain(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "V5 修補鏈為 5 commits，4/30 補 ploidy fix + threshold cherry-pick")
    # Embed pre-rendered timeline figure (matplotlib)
    fit_image_within(
        s,
        FIGURES_DIR / "fig_s03_v5_5commits_timeline.png",
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.4),
        border=False,
    )
    # Working tree caveat resolved insight (still on-slide as Tier 1 takeaway)
    add_insight_card(
        s, Inches(0.5), Inches(6.05), Inches(12.333), Inches(0.95),
        [
            "git diff --stat HEAD = 0 byte (working tree clean) · 舊 caveat R1 已隨 938f0df 解決",
        ],
        title="[已解決] 4/30 working tree caveat R1",
    )
    add_footer_citation(s, "Source: 20260505_self_phasing_V5_data_provenance_audit_01.md §1.1")
    set_speaker_note(
        s,
        "V5 不是 4 commits 是 5 commits，這是本週第一個校正點。原本記錄的 V5 是前三個 commit"
        " (PON-only flag、tag layer fix、somatic fallback)，4/30 那天補了兩個關鍵 commit："
        "d0bcd8c 修 ploidy bug，938f0df 從上游 cherry-pick threshold 0.95→0.9 調整。"
        "這兩個都是紅色節點，因為 d0bcd8c 解開 Pass 2 觸發機制，938f0df 引入路 3 抵消效應。\n\n"
        "之前 audit 留有一個 caveat R1：working tree 有未 commit 改動。5/05 確認 git diff --stat HEAD"
        " 為空，全部邏輯都已 commit、可追溯。這個 caveat 隨 938f0df 一起被解決，是本週副效益。\n\n"
        "[ORAL-OPTIONAL] threshold cherry-pick 細節：上游 zhenyu 把 highPurity 觸發閾值從 0.95 降到 0.9，"
        "讓更多 case 進入 Pass 2 second round。這個改動在我們 case 是雙刃劍 — 後面 slide 6 會看到"
        "路 3 second round 反而抵消路 2 反轉效果。",
        min_chars=300,
    )
    add_focal_marker(s, "5 commits + 4/30 兩個關鍵節點")
    return s


def slide_04_pi_pass1_only(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(
        s, "* 5/05 Critical Finding: PI 報告 V5 數據限 Pass 1 only 條件",
        color_bg=PALETTE["accent"],
    )
    # Caveat banner
    add_caveat_strip(
        s, Inches(0.5), Inches(1.5), Inches(12.333),
        "[!] 4/29 PI 報告引用的所有 V5 數值都是 Pass 1 only 條件（Pass 2 從未觸發）",
    )
    # 1 行白話因果鋪墊（PI 還沒看到 S05 之前先給）
    add_method_box(
        s, Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.85),
        [
            "因果一句話：4/12 BAM 因 ploidy bug 算出 purity=0 → V5 跳過 second-round → 數值來自單輪修補（非完整版）",
        ],
        title="為何 Pass 1 only",
    )
    # Evidence 合併單欄（簡化原雙欄）
    add_method_box(
        s, Inches(0.5), Inches(3.15), Inches(12.333), Inches(1.6),
        [
            "BAM = pononly_v2b/tumor_tagged.bam（4/12 產）",
            "log: purity=0 + 無 'second round phasing' 字串",
            "→ Pass 2 second round 跳過 → 數值 = Pass 1 + tag layer fix",
        ],
        title="Evidence",
    )
    # Mini-table: PI report numbers (精簡為 3 行)
    add_neutral_box(
        s, Inches(0.5), Inches(4.85), Inches(12.333), Inches(1.85),
        [
            "clean PS V5=88.2% / BL=74.9% (+13.3pp)  ·  AMB% 17.5→8.0% (-54.3%)",
            "HP:i:33 read 計數 239,679 → 110,197 (-54%, whole genome, PI 口徑)",
            "→ 結論「V5 為當前最佳 phasing baseline」需重驗於 Pass 2 真實觸發版",
        ],
        title="PI 報告關鍵數字（皆 Pass 1 only）",
    )
    add_footer_citation(s, "Source: master_draft.md Top Findings #2, audit_01.md §1.2")
    set_speaker_note(
        s,
        "這張是本週最重要的 slide。5/05 做 provenance audit 時發現一個被忽視的盲點："
        "4/29 PI 報告所有 V5 數值都引用 4/12 產的那份 BAM，但那份 BAM 因為 ploidy bug 讓 purity 算出 0，"
        "整個 highPurity 條件 false，Pass 2 second round 從未觸發。所以那些好看的數字 — "
        "clean PS +13.3pp、AMB% 從 17.5% 降到 8%、HP:i:33 完全消失 — 全部是 "
        "「PON-only Pass 1 + tag layer fix」的功勞，不是 V5 完整版的功勞。\n\n"
        "Evidence 兩條：A 是 BAM provenance，log 顯示 purity:0、無 second round 字串；"
        "B 是機制因果，purity=0 → highPurity=false → Pass 2 跳過。兩條交叉確認。\n\n"
        "為什麼這很 critical？因為「V5 為當前最佳 phasing baseline」這個主結論在現有 audit chain"
        "（4/29 技術報告、4/28 audit、Memory project_v5_somatic_fallback_verification）裡到處被引用，"
        "但實際上這個結論的現實基礎是 Pass 1 only 條件。我們需要 Pass 2 真實觸發版重驗才能維持。\n\n"
        "[ORAL-OPTIONAL] highPurity 三次多項式公式在 PhasingProcess.cpp:142-220，"
        "purity 從 read coverage / depth 算出來，ploidy bug 讓計算 chain 在 ploidy=0 時短路。",
        min_chars=400,
    )
    add_focal_marker(s, "Tier S — 主結論需暫停")
    return s


def slide_05_mechanism(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "ploidy bug → purity=0 → highPurity=false → Pass 2 從未觸發")
    # Embed pre-rendered causal-chain figure (matplotlib)
    fit_image_within(
        s,
        FIGURES_DIR / "fig_s05_ploidy_causal_chain.png",
        Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6),
        border=False,
    )
    add_footer_citation(s, "Source: PhasingProcess.cpp:142-220, audit_01.md §2.3")
    set_speaker_note(
        s,
        "把 critical finding 視覺化成 5 階段因果鏈。從左到右：4/12 BAM → ploidy=0 → purity=0 → "
        "highPurity=false → Pass 1 only result。中間三步 (紅色) 是 ploidy bug 的影響鏈，"
        "右邊的橘色節點是最終結果，也就是 4/29 PI 報告引用的數值。\n\n"
        "4/30 d0bcd8c 解了這個 bug，purity 計算重新正常，現在 Pass 2 second round 會真的觸發。"
        "但問題是 4/29 PI 報告寫的時候用的還是 4/12 那份 BAM，所以數字維持原樣。要更新就必須"
        "在 4/30 修補後 binary 下整套重跑：HCC1395 5kHz sanity、concordance、ISM benchmark，"
        "估計 ~25 hr 含 7 樣本平行。\n\n"
        "Decisive next step P1 就是這個重跑。如果不做，整條 V5 audit chain 都卡在 Pass 1 only 條件。\n\n"
        "[ORAL-OPTIONAL] ploidy 計算分支在 PhasingProcess.cpp:142-220 那個區段，d0bcd8c 改了 5 行 — "
        "詳細 diff 在 audit_01.md §2.3。可以追加討論為何上游沒一起 cherry-pick。",
        min_chars=300,
    )
    add_focal_marker(s, "因果鏈 + 解法 + P0")
    return s


def slide_06_three_paths(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "5/01 三條路 audit：路 3 second round 抵消路 2 反轉效果")
    # Embed pre-rendered three-paths codeflow figure (matplotlib)
    fit_image_within(
        s,
        FIGURES_DIR / "fig_s06_three_paths_codeflow.png",
        Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6),
        border=False,
    )
    add_footer_citation(s, "Source: 20260501_latest_longphase_to_3paths_audit_01.md §3, PhasingProcess.cpp:142-220")
    set_speaker_note(
        s,
        "本週第二個 critical finding：5/01 三條路 audit。我們先把 PhasingProcess 拆成三條路徑來理解。"
        "路 1 是原始 baseline，沒有任何 flag，single pass，self-phasing 17.3:1 殘留。"
        "路 2 是 V5 flag PON-only，4 步流程含 somaticCalling 重跑，這是反轉 self-phasing 的核心 — "
        "舊 V5 在 ploidy bug 條件下走純路 2，HP1:HP2 = 0.735 [OK] 反轉。"
        "路 3 是新加的 highPurity > 0.9 second round，只重 phase 不重 call，self-phasing 偏移殘留。\n\n"
        "問題是新 binary 下 V5 flag 會走「路 2 + 路 3」，路 3 的 second round 把路 2 反轉的東西"
        "重新偏掉，最終 HP1:HP2 = 1.400，比新 baseline (路 1+3) 1.400 還高、比舊 baseline (路 1) 1.328 更糟。\n\n"
        "教授視角的 implication：V5 flag 在新 binary 下「機制成立」(路 2 仍能反轉) 但「最終效果消失」"
        "(路 3 抵消)。這直接影響「V5 flag 還要不要留」的決策，後面 slide 11 會回到這個問題。\n\n"
        "[ORAL-OPTIONAL] HaplotagProcess.cpp:484-563 是 somaticCalling 重跑的程式碼，"
        "如果想看具體哪幾行讓路 2 反轉，可以打開該段。",
        min_chars=300,
    )
    add_focal_marker(s, "Tier S — 路 3 抵消路 2")
    return s


def slide_07_5version_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "5 版本 HP1:HP2 對比 (Ground truth ≈ 1:1)：NEW noPath3 復現 OLD V5 反轉")
    # Table 5 rows × 4 cols（合併走的路+反轉?，HP_33 移 footer）
    rows = [
        ("版本", "走的路 + 結果", "HP1:HP2", "HP_33"),
        ("OLD baseline", "路 1：基準", "1.328", "2,640"),
        ("OLD V5 flag", "路 2 (bug)：[OK] 反轉", "0.735", "14,524"),
        ("NEW baseline", "路 1+3：[X] 更糟", "1.400", "3,468"),
        ("NEW V5 flag", "路 2+3：[X] 抵消", "1.400", "3,468"),
        ("NEW V5 noPath3", "路 2 only (forced)：[OK] 反轉*", "1.127*", "28*"),
    ]
    base_x = Inches(0.5)
    base_y = Inches(1.5)
    col_widths = [Inches(2.8), Inches(4.5), Inches(2.5), Inches(2.533)]
    row_h = Inches(0.6)
    for r_idx, row in enumerate(rows):
        x = base_x
        for c_idx, val in enumerate(row):
            w = col_widths[c_idx]
            y = base_y + row_h * r_idx
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            cell.shadow.inherit = False
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALETTE["title_bg"]
            elif row[0] in ("OLD V5 flag", "NEW V5 noPath3"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # green tint
            elif row[0] in ("NEW baseline", "NEW V5 flag"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)  # red tint
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            cell.line.color.rgb = PALETTE["border"]
            cell.line.width = Pt(0.5)
            tf = cell.text_frame
            tf.margin_left = Emu(60000)
            tf.margin_right = Emu(60000)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            color = PALETTE["title_fg"] if r_idx == 0 else PALETTE["text"]
            bold = (r_idx == 0) or (c_idx == 0)
            add_text_with_fallback(
                tf, val, font_size=12, bold=bold, color=color,
                align=PP_ALIGN.CENTER, paragraph_index=0,
            )
            x = x + w
    # Insight strip
    add_insight_card(
        s, Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.4),
        [
            "OLD V5 (路 2 only, ploidy bug 副作用) HP1:HP2 = 0.735，反轉成功",
            "NEW noPath3 (強制跳過路 3) HP1:HP2 = 1.127*，反轉成功（* 15-site cherry-pick）",
            "→ 兩者都反轉，方向一致 → 路 3 抵消假說 PASS（noPath3 全 BAM 預期 ≈ 0.735，待 monitor [U]）",
        ],
        title="[OK] 證據 — 路 3 抵消假說 PASS",
    )
    add_footer_citation(s, "* noPath3 數字來自 force_path2only_ablation §3.4 的 15-site cherry-picked sample；其他 4 列為 threshold_compare 全 BAM。Source: §3.4 + 3paths_audit §3.1")
    set_speaker_note(
        s,
        "這張表是本週最重要的數值對照。5 個版本：OLD baseline / OLD V5 / NEW baseline / NEW V5 / NEW noPath3。"
        "前兩個是 4/12 BAM 結果（ploidy bug 條件），後三個是 4/30 修補後 binary 結果。\n\n"
        "綠色行：成功反轉 (HP1:HP2 < 1)。OLD V5 (0.735, 全 BAM) + NEW noPath3 (1.127*, 15-site cherry) 都成功。"
        "紅色行：抵消或更糟。NEW baseline (1.400) + NEW V5 (1.400) 都失敗。\n\n"
        "關鍵 caveat — noPath3 表格中的 1.127 / HP_33=28 是 15-site cherry-picked sample 的數字"
        "（force_path2only §3.4），不是全 BAM 觀測值。全 BAM monitor 還在跑，預期值與舊 V5 全 BAM "
        "（0.735, HP_33=14,524）接近。當前用 cherry-pick 是因為它能快速證明「路 3 抵消」假說的方向 "
        "（反轉成立），但量化等價性需等全 BAM 完成才能下結論 [U]。\n\n"
        "教授可能追問：為何不等全 BAM？答：全 BAM 跑 ~3 hr × 7 樣本要 21 hr，與 P1 重驗共用 binary 才"
        "能效率最大化；當前 15-site 已足以證明假說方向，量化由 P2 補。\n\n"
        "教授視角的 implication：如果我們本地維護 longphase-to-noPath3 變體 binary，"
        "等於回到舊 V5 的反轉效果但無 ploidy bug。這是 slide 11 三選項中的 (b)。\n\n"
        "[ORAL-OPTIONAL] noPath3 1.127 vs OLD V5 0.735 的差距：兩者 scope 不同（noPath3 是 15-site，"
        "OLD V5 是全 BAM），不能直接比較絕對值，只能比較方向（都反轉）。全 BAM 跑完才能量化等價程度。",
        min_chars=300,
    )
    add_focal_marker(s, "等價性證明：NEW noPath3 ≈ OLD V5")
    return s


def slide_08_force_path2only(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "force_path2only ablation：強制 noPath3 復現舊 V5 反轉（負控制 PASS）")
    # Setup vs Result two-column
    add_neutral_box(
        s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5),
        [
            "Hypothesis：「路 3 second round 抵消路 2 反轉」",
            "",
            "Method：",
            "  • patch PhasingProcess.cpp 一行：",
            "  • bool highPurity = false;",
            "  • 強制跳過路 3，不論實際 purity",
            "",
            "Sample：HCC1395 5kHz",
            "Phase 時間：813s (vs 路 2+3 的 2881s, 節省 71%)",
            "預期：HP1:HP2 接近 OLD V5 (0.735)",
        ],
        title="Setup",
    )
    add_insight_card(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(4.5),
        [
            "Result (15-site cherry-picked):",
            "  • HP1:HP2 = 1.127 (反轉成立)",
            "  • HP_33 = 28",
            "",
            "→ 反轉方向與舊 V5 全 BAM (0.735) 一致",
            "→ 路 3 抵消假說 PASS（方向）",
            "→ 全 BAM 等價量化待 monitor 完成",
            "",
            "副產物：",
            "  • Phase 時間節省 71% (813s vs 2881s)",
            "  • 同一 commit 938f0df 下可重現",
        ],
        title="Result + Verdict",
    )
    add_method_box(
        s, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.85),
        [
            "結論：路 3 是 NEW V5 失效的 single root cause；如未來決定維護 noPath3 變體，22.55 MB binary 已建立。",
        ],
    )
    add_footer_citation(s, "Source: 20260501_v5_force_path2only_ablation_01.md")
    set_speaker_note(
        s,
        "這張是負控制實驗 (negative control)。我們想驗證「路 3 抵消路 2 反轉」這個假說。"
        "做法：直接在 PhasingProcess.cpp patch 一行，bool highPurity = false，"
        "強制讓 second round 條件失敗、跳過路 3，不論實際 purity 計算結果。\n\n"
        "如果路 3 真的是抵消來源，跳過路 3 應該復現舊 V5 反轉。結果：HP1:HP2 = 1.127 反轉成功"
        "（15-site cherry-picked sample，HP_33 = 28）。反轉方向與舊 V5 全 BAM (0.735) 一致，"
        "假說 PASS（方向）。量化等價性待全 BAM monitor 完成。\n\n"
        "副產物：phase 時間從 2881s 降到 813s，節省 71%。這個沒有實際生產意義 (因為 noPath3 binary "
        "並非 default)，但顯示路 3 是計算密集區段。\n\n"
        "教授視角的 implication：這個實驗證明「機制完全清楚」— V5 flag 本身沒壞，問題只在路 3。"
        "如果決定維護本地 longphase-to-noPath3 變體 binary，22.55 MB 可直接 deploy。\n\n"
        "[ORAL-OPTIONAL] 為什麼選 force_path2only 這種 hack 而非改 threshold？因為 threshold 改動"
        "(0.95 vs 0.9) 是上游 zhenyu 那邊維護的，本地改動會 diverge；強制 highPurity=false 是一行 patch，"
        "可作為 ablation 工具，不影響上游。",
        min_chars=300,
    )
    add_focal_marker(s, "Tier A — 假說 PASS")
    return s


def slide_09_caller_f1_unchanged(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Caller F1 全 6 版本相同：FILTER 不變，需用 ISM 下游 F1 評估 phasing")
    # KPI cells
    add_kpi_block(
        s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.0),
        "Caller F1 @ Pass=0.93", "0.7166",
        value_color=PALETTE["title_bg"],
        sub="6 版本完全相同 (OLD/NEW × baseline/V5/noPath3)",
    )
    add_kpi_block(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(2.0),
        "Caller F1 @ Pass=0.6", "0.6273",
        value_color=PALETTE["title_bg"],
        sub="3 版本完全相同 (4/30 V3F ablation 首次發布)",
    )
    add_method_box(
        s, Inches(0.5), Inches(3.7), Inches(12.333), Inches(1.6),
        [
            "longphase-to 改動只動 phasing graph，不改 FILTER 欄位",
            "Caller F1 (caller 直接輸出) 對 phasing 改動不敏感",
            "→ phasing 品質必須用 ISM SuggestFilter 下游 F1 評估，不是 caller F1",
        ],
        title="為何 Caller F1 不能評估 phasing 改動",
    )
    add_caveat_strip(
        s, Inches(0.5), Inches(5.5), Inches(12.333),
        "[!] 4/29 PI 報告引用的「F1 不變」結論正確，但需附註：phasing 品質需另用 ISM F1 衡量",
    )
    add_footer_citation(s, "Source: 20260430_V3F_ablation_purity06_results_01.md (首次發布), 20260430_V3F_getVote_only_ablation_results_01.md")
    set_speaker_note(
        s,
        "這張是 supporting evidence — 確認 caller F1 完全不受 phasing 改動影響。"
        "0.93 pass quality 6 個版本都是 0.7166，0.6 pass quality 3 個版本都是 0.6273。"
        "0.6273 這個數字是本週首次發布 (4/30 V3F purity ablation)，之前只有 0.93 點數據。\n\n"
        "為什麼 caller F1 對 phasing 改動完全不敏感？因為 longphase-to 只動 phasing graph "
        "(decide HP tag 怎麼分配)，不動 FILTER 欄位 (PASS/LowQual/...)。Caller F1 看的是 FILTER，"
        "所以 phasing 改動 invisible。\n\n"
        "教授視角 implication：4/29 PI 報告寫「caller F1 不變」是對的，但要附註 — 這個不變不代表"
        "「phasing 沒改動」，只代表「caller 沒受傷」。phasing 品質要用 ISM SuggestFilter 下游 F1 衡量"
        "(下游 F1 才會看到 HP tag 分配差異)。\n\n"
        "[ORAL-OPTIONAL] 0.6273 這個數字之前未公布的原因：4/30 V3F ablation 第一次跑出 0.6 pass quality "
        "點，作為 cross-purity 確認。詳見該 ablation 報告。",
        min_chars=300,
    )
    add_focal_marker(s, "FILTER 不變 → 需 ISM F1 評估")
    return s


def slide_10_caveat_banner(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(
        s, "[!] V5 結論 caveat 總結：暫停為「Pass 1 only 條件下觀察」，待 Pass 2 重驗",
        color_bg=PALETTE["accent"],
    )
    # Big caveat block — 加高至 3.0 in 增加視覺比例
    big = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(3.0)
    )
    big.shadow.inherit = False
    big.fill.solid()
    big.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    big.line.color.rgb = PALETTE["accent"]
    big.line.width = Pt(2.5)
    tf = big.text_frame
    tf.margin_left = Emu(160000)
    tf.margin_right = Emu(160000)
    tf.margin_top = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    add_text_with_fallback(
        tf, "[!] Caveat Banner — V5 結論狀態",
        font_size=18, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C),
        paragraph_index=0,
    )
    add_text_with_fallback(
        tf,
        "「V5 為當前最佳 phasing baseline」此結論基於 Pass 1 only 條件 (4/12 BAM, ploidy bug)。",
        font_size=14, color=PALETTE["text"],
    )
    add_text_with_fallback(
        tf,
        "4/30 修補後 Pass 2 真實觸發但 sanity / concordance / ISM benchmark 尚未重驗。",
        font_size=14, color=PALETTE["text"],
    )
    add_text_with_fallback(
        tf,
        "→ 暫停為「Pass 1 only 條件下觀察」，等 P0 重驗 (~25 hr) 完成才能恢復為「驗證後結論」。",
        font_size=14, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C),
    )
    # Affected artifacts — 加高至 2.7 in + caveat strip 距離縮 0.1 in
    add_neutral_box(
        s, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.4),
        [
            "• 4/29 longphase TO vs V5 技術報告 → 補 caveat banner",
            "• 4/28 V5 audit haplotag sanity → 標記 BAM = 4/12 (Pass 1 only)",
            "• Memory project_v5_somatic_fallback_verification → status 從 concluded → needs_rerun",
            "• 引用 V5 為 baseline 的 Thread D 主軸切換文件 → 加 footnote",
        ],
        title="受影響 artifacts (需更新 caveat)",
        accent=PALETTE["amber"],
    )
    add_footer_citation(s, "Source: master_draft.md §11 R1, audit_01.md §3")
    set_speaker_note(
        s,
        "這張是把所有需要更新的 artifacts 列出來，作為 action item 的視覺整合。"
        "不是新發現，是把 critical finding 對 audit chain 的影響具體化。\n\n"
        "Banner 內文：「V5 為當前最佳 phasing baseline」這個結論基於 Pass 1 only 條件，"
        "目前需暫停為「Pass 1 only 條件下觀察」。等 Pass 2 重驗 ~25 hr 完成才能恢復為「驗證後結論」。\n\n"
        "受影響 artifacts 4 個：4/29 longphase TO vs V5 技術報告、4/28 V5 audit、"
        "Memory project_v5_somatic_fallback_verification (status 改 needs_rerun)、"
        "Thread D 主軸切換文件 (因為它引用 V5 為 baseline)。每個都要補 caveat 或更新狀態。\n\n"
        "教授視角的 implication：要決定是「立即補 caveat」還是「等 Pass 2 重驗完一次更新」。"
        "兩個選擇 trade-off 不一樣 — 立即補 caveat 透明但會讓上游讀者困惑；等重驗完一次更新清晰但"
        "需 1 週透明度落差。這是 Top Asks #1。\n\n"
        "[ORAL-OPTIONAL] 也可以考慮折衷：先在 Memory 改 needs_rerun (內部記錄)，"
        "等重驗完才更新外部技術報告。",
        min_chars=300,
    )
    add_focal_marker(s, "Caveat 影響 4 個 artifacts")
    return s


def slide_11_decision_tree(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "V5 留存決策三選項：a) 接受新 default · b) 維護 noPath3 · c) 上游 PR")
    options = [
        (
            "(a) 接受新 default",
            PALETTE["gray"],
            [
                "走路 1+3，self-phasing 1.40 殘留",
                "與上游 zhenyu 同步維護成本最低",
                "下游 ISM 必須容忍 self-phasing artifact",
                "估時：0 hr (現狀)",
            ],
        ),
        (
            "(b) 維護 noPath3 binary",
            PALETTE["green"],
            [
                "本地 longphase-to-noPath3 (22.55 MB 已建)",
                "復現舊 V5 反轉效果 (HP1:HP2 = 1.127)",
                "需週期性同步上游修補",
                "估時：1 hr/week 維護",
            ],
        ),
        (
            "(c) 上游 PR",
            PALETTE["title_bg"],
            [
                "提案 zhenyu 修改路 3 邏輯",
                "讓 second round 也重 call somatic",
                "best long-term, 但時程不可控",
                "估時：~ 2 週 (含溝通)",
            ],
        ),
    ]
    n = len(options)
    box_w = Inches(4.0)
    gap = (Inches(12.333) - box_w * n) / max(n - 1, 1)
    base_x = Inches(0.5)
    base_y = Inches(1.5)
    box_h = Inches(4.0)
    for i, (label, color, lines) in enumerate(options):
        x = base_x + (box_w + gap) * i
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, base_y, box_w, box_h)
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        shp.line.color.rgb = color
        shp.line.width = Pt(2.5)
        tf = shp.text_frame
        tf.margin_left = Emu(120000)
        tf.margin_right = Emu(120000)
        tf.margin_top = Emu(100000)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        add_text_with_fallback(
            tf, label, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER,
            paragraph_index=0,
        )
        add_text_with_fallback(
            tf, "─" * 12, font_size=10, color=PALETTE["border"], align=PP_ALIGN.CENTER,
        )
        for line in lines:
            add_text_with_fallback(
                tf, "• " + line, font_size=12, color=PALETTE["text"],
            )
    add_method_box(
        s, Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.4),
        [
            "決策建議 (待教授判斷)：",
            "短期 (本週) 採 (b) 維護 noPath3 binary 作為 ISM 重驗備援；",
            "中期 (本月) 推 (c) 上游 PR；(a) 僅在 (b)(c) 都失敗時 fallback。",
        ],
        title="My recommendation",
    )
    add_footer_citation(s, "Source: master_draft.md §16 Priority 4, Top Asks #2")
    set_speaker_note(
        s,
        "這是 Top Asks 的第二個必判斷決策點。新 binary 下 V5 flag 不再有效反轉 self-phasing，"
        "我們有三個選項。\n\n"
        "(a) 接受新 default：什麼都不做。優點是與上游同步無維護成本，缺點是下游 ISM 必須容忍 "
        "self-phasing 1.40 殘留。0 hr。\n\n"
        "(b) 維護 longphase-to-noPath3 變體 binary：22.55 MB 已建好。可復現舊 V5 反轉效果。"
        "缺點是上游每次更新需週期性同步，~1 hr/week。\n\n"
        "(c) 上游 PR：提案 zhenyu 修路 3 邏輯讓 second round 重 call。最徹底但時程不可控，~2 週。\n\n"
        "我的建議：短期 (b) 作為 ISM 重驗備援，中期推 (c)，(a) 作 fallback。但這要教授判斷 — "
        "因為 (c) 涉及上游溝通的政治成本，(b) 涉及 binary divergence 風險。\n\n"
        "[ORAL-OPTIONAL] 22.55 MB binary 大小其實很合理，longphase 主執行檔本來就 ~22 MB。"
        "已 staged 在本地 build/，未 commit；commit 與否視 (b)(c) 決策。",
        min_chars=300,
    )
    add_focal_marker(s, "三選項 — 待教授判斷")
    return s


def slide_12_cross_sample_risk(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "938f0df 跨樣本影響 [U]：「新 baseline 偏移」是否 systematic？")
    # Sample status table
    samples = [
        ("HCC1395 5kHz", "1.328 → 1.400", "+5.4%", "[OK] 已測"),
        ("HCC1395 DORADO", "?", "[U]", "⏳ Priority 2"),
        ("HCC1937", "?", "[U]", "⏳ Priority 2"),
        ("HCC1954", "?", "[U]", "⏳ Priority 2"),
        ("H2009", "?", "[U]", "⏳ Priority 2"),
        ("H1437", "?", "[U]", "⏳ Priority 2"),
        ("COLO829", "?", "[U]", "⏳ Priority 2"),
    ]
    base_x = Inches(0.5)
    base_y = Inches(1.5)
    col_widths = [Inches(3.5), Inches(3.5), Inches(2.5), Inches(2.833)]
    headers = ["樣本", "OLD → NEW baseline HP1:HP2", "變化", "狀態"]
    row_h = Inches(0.55)
    rows = [headers] + samples
    for r_idx, row in enumerate(rows):
        x = base_x
        for c_idx, val in enumerate(row):
            w = col_widths[c_idx]
            y = base_y + row_h * r_idx
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            cell.shadow.inherit = False
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALETTE["title_bg"]
            elif r_idx == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            cell.line.color.rgb = PALETTE["border"]
            cell.line.width = Pt(0.5)
            tf = cell.text_frame
            tf.margin_left = Emu(60000)
            tf.margin_right = Emu(60000)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            color = PALETTE["title_fg"] if r_idx == 0 else PALETTE["text"]
            bold = r_idx == 0
            add_text_with_fallback(
                tf, val, font_size=12, bold=bold, color=color,
                align=PP_ALIGN.CENTER, paragraph_index=0,
            )
            x = x + w
    add_caveat_strip(
        s, Inches(0.5), Inches(6.4), Inches(12.333),
        "[!] 若 systematic：考慮 (b) noPath3 / (c) 上游 PR；若樣本特異：可暫接受新 default",
    )
    add_footer_citation(s, "Source: master_draft.md R2 + Risk R2")
    set_speaker_note(
        s,
        "這張是 risk surface — 把 cross-sample 不確定性視覺化。"
        "目前 HCC1395 5kHz 一個樣本的 baseline HP1:HP2 從 1.328 (OLD) 升到 1.400 (NEW)，+5.4%。"
        "其他 6 個樣本還沒測。\n\n"
        "為什麼 cross-sample 重要？因為「新 baseline 比舊 baseline 更偏 HP1」可能是 systematic "
        "(每個樣本都受 938f0df 影響) 或樣本特異 (只是 HCC1395 5kHz 的 noise)。"
        "若 systematic，slide 11 的決策 (b) 或 (c) 都需要走；若樣本特異，可暫接受新 default。\n\n"
        "Priority 2 預估 4 hr (跨 6 樣本平行)，是 Pass 2 重驗 (P1, 25 hr) 的子集。可以 P1+P2 一起做。\n\n"
        "[ORAL-OPTIONAL] HCC1395 DORADO 是 Dorado basecaller 重 call 過的版本，"
        "與 5kHz 是同一樣本不同 basecaller。這兩個算 1.5 個樣本 — 若兩者一致，仍需要其他樣本確認。",
        min_chars=300,
    )
    add_focal_marker(s, "[U] 跨樣本待驗")
    return s


def slide_13_hpfinengroups_revalidation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "HPFineNGroups marker 在新 baseline 下重驗 [U]（與 R1 平行）")
    # 1 行術語 footnote（PI 不熟）
    add_method_box(
        s, Inches(0.5), Inches(1.45), Inches(12.333), Inches(0.55),
        [
            "術語：marker = 候選生物標記｜phasing signature = 因 phasing 演算法產生的副作用模式（非真實生物訊號）",
        ],
    )
    add_neutral_box(
        s, Inches(0.5), Inches(2.1), Inches(6.0), Inches(3.9),
        [
            "4/28 audit 結果 (4/12 BAM, Pass 1 only)：",
            "  - marker tier 4-stars → 3-stars 降級",
            "  - HCC1395 TO ClairS-TO raw split 無法重現",
            "  - Fisher odds 0.913 反向 p=3.5e-3",
            "  - flag=on NG>=3 = 0",
            "",
            "重詮釋為 phasing signature：",
            "  - 不是 methylation marker",
            "  - 是 self-phasing 的副產物",
        ],
        title="既有狀態 (4/28)",
    )
    add_method_box(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(4.5),
        [
            "新 baseline (路 1+3) 下需重新驗證：",
            "",
            "  1. 是否仍見 NG≥3 = 0？",
            "    → 若是：phasing signature 假說 ✓",
            "    → 若否：機制需重新理解",
            "",
            "  2. flag=on (NEW V5 路 2+3) 是否變化？",
            "    → 預期類似 (因為仍走路 3)",
            "",
            "  3. noPath3 binary 是否復現舊行為？",
            "    → 若是：完全對稱 paired check",
            "",
            "Priority 3 估時：0.5 hr (與 R1 共用 binary)",
        ],
        title="R3：待驗 hypothesis",
    )
    add_caveat_strip(
        s, Inches(0.5), Inches(6.2), Inches(12.333),
        "[!] 4/28 audit 是 Pass 1 only 結果；新 baseline 結論需 R3 重驗才能 conclude",
    )
    add_footer_citation(s, "Source: Memory project_hpfinengroups_subclone_marker, Top Asks #6")
    set_speaker_note(
        s,
        "這張是 audit chain 的另一條延伸 — HPFineNGroups subclone marker。"
        "4/23 週報把它從 4-stars 降到 3-stars，重詮釋為 phasing signature。但所有 audit 都是 4/12 BAM "
        "(Pass 1 only)。新 baseline 下需要重驗。\n\n"
        "需要驗的三件事：(1) NG≥3 = 0 是否仍成立 (若是 → phasing signature 假說 ✓)；"
        "(2) flag=on (NEW V5 路 2+3) 是否變化 (預期類似因為仍走路 3)；"
        "(3) noPath3 binary 是否復現舊行為 (若是 → 完全對稱 paired check)。\n\n"
        "Priority 3 估時 0.5 hr，與 R1 (Pass 2 重驗) 共用 binary，所以可一起做。\n\n"
        "教授視角 implication：4/28 audit 結論「marker tier 3-stars」目前暫定，等 R3 驗完才能確認。"
        "若 R3 結果與 4/28 不一致，可能要再降一級或重新評估機制。\n\n"
        "[ORAL-OPTIONAL] HPFineNGroups 是 src/include/HPFineNGroups.hpp 定義的特徵，"
        "計算 HP tag 內部精細分群的群數。詳細定義已在 Memory feedback_feature_name_vs_definition_rule.",
        min_chars=300,
    )
    add_focal_marker(s, "R3 — Marker 在新 baseline 下重驗")
    return s


def slide_14_thread_b_context(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Thread B 撤回：機制 pivot 為 phasing artifact，TO 撤回 / Paired 保留")
    # 1 行術語 footnote (PI 不熟)
    add_method_box(
        s, Inches(0.5), Inches(1.45), Inches(12.333), Inches(0.55),
        [
            "術語：paired = 含 normal 對照樣本｜TO = tumor-only 無對照｜NG = NGroups (HPFineNGroups subclone marker)",
        ],
    )
    # 整段時間軸
    add_neutral_box(
        s, Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.7),
        [
            "原假說 (4/19): NG=2 LOH-AF interaction → methylation 訊號",
            "重審 (4/23-26): self-phasing 偏移先污染 HP_Ratio LOH → 訊號 root cause 是 phasing artifact",
            "現狀 (5/01): 機制 pivot 為 phasing signature",
        ],
        title="Thread B 演進時間軸",
    )
    # 兩層處置合併 (簡化原雙欄)
    add_method_box(
        s, Inches(0.5), Inches(3.95), Inches(12.333), Inches(2.7),
        [
            "TO 層：撤回 NG=2 methylation 機制；重詮釋為 phasing signature；等 Pass 2 重驗才能再下結論",
            "Paired 層：保留 POSITIVE — Inter AF→NGroups +0.705~+0.787 (7/7 樣本)；",
            "             paired 不受 self-phasing 影響，但需獨立 phasing-vs-methylation 驗證",
            "→ Thread B 不是完全失敗：TO 層需重驗，Paired 層仍有訊號",
        ],
        title="兩層處置 (TO 撤回 / Paired 保留)",
    )
    add_footer_citation(s, "Source: Memory project_loh_subclone_af_methylation_positive, project_hpfinengroups_subclone_marker")
    set_speaker_note(
        s,
        "這張是把 Thread B 的演進與目前狀態講清楚，因為它與本週主軸 (Thread D 切換) 直接相關。\n\n"
        "Thread B 原本是 4/19 提的假說 — NG=2 LOH-AF interaction 產生 methylation 訊號。"
        "4/23-26 重新審視時發現 self-phasing 會先污染 HP_Ratio LOH，所以「methylation 訊號」其實是"
        "phasing artifact。5/01 機制 pivot 為 phasing signature。\n\n"
        "TO 層：撤回 methylation 機制，重詮釋為 phasing signature，等 Pass 2 重驗。\n"
        "Paired 層：保留 POSITIVE 但加 caveat。Inter AF→NGroups +0.705 ~ +0.787 (7/7) 是 paired 層"
        "的觀察，paired 不受 self-phasing 影響所以保留，但需要獨立的 phasing-vs-methylation 驗證才能"
        "確認 mechanism。\n\n"
        "教授視角 implication：Thread B 不是完全失敗 — paired 層仍有訊號，TO 層需重驗。"
        "這也說明 self-phasing artifact 對下游分析的污染範圍很廣。\n\n"
        "[ORAL-OPTIONAL] Memory 兩條：project_loh_subclone_af_methylation_positive (paired 保留) 與 "
        "project_hpfinengroups_subclone_marker (TO 降級)。詳情可追問。",
        min_chars=300,
    )
    add_focal_marker(s, "TO 撤回, Paired 保留")
    return s


def slide_15_future_priorities(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Future Priorities：5 件事 + 工時，總 ~33 hr (P1+P2 平行可壓 25 hr)")
    priorities = [
        ("P1", "[P1]", "Pass 2 重驗", "25 hr", PALETTE["accent"], "決定性 — 若不做主結論失基礎"),
        ("P2", "[P2]", "938f0df 跨樣本", "4 hr", PALETTE["amber"], "與 P1 共用 binary 並行"),
        ("P3", "[P2]", "PI 報告 caveat", "2 hr", PALETTE["amber"], "視 P1 完成度決定立即/延後"),
        ("P4", "[P2]", "V5 留存決策文件", "1 hr", PALETTE["title_bg"], "等教授判斷後撰寫"),
        ("P5", "[!]", "Memory 校正", "0.5 hr", PALETTE["gray"], "5-05 校正完成"),
    ]
    base_x = Inches(0.5)
    base_y = Inches(1.5)
    card_w = Inches(2.4)
    gap = (Inches(12.333) - card_w * 5) / 4
    card_h = Inches(3.8)
    for i, (pid, stars, title, hours, color, note) in enumerate(priorities):
        x = base_x + (card_w + gap) * i
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, base_y, card_w, card_h)
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        shp.line.color.rgb = color
        shp.line.width = Pt(2.0)
        tf = shp.text_frame
        tf.margin_left = Emu(80000)
        tf.margin_right = Emu(80000)
        tf.margin_top = Emu(80000)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        add_text_with_fallback(
            tf, pid, font_size=22, bold=True, color=color,
            align=PP_ALIGN.CENTER, paragraph_index=0,
        )
        add_text_with_fallback(
            tf, stars, font_size=14, color=color, align=PP_ALIGN.CENTER,
        )
        add_text_with_fallback(
            tf, title, font_size=14, bold=True, color=PALETTE["text"],
            align=PP_ALIGN.CENTER,
        )
        add_text_with_fallback(
            tf, hours, font_size=20, bold=True, color=PALETTE["title_bg"],
            align=PP_ALIGN.CENTER,
        )
        add_text_with_fallback(
            tf, note, font_size=11, color=PALETTE["gray"], align=PP_ALIGN.CENTER,
        )
    add_method_box(
        s, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.4),
        [
            "本週執行計畫：P1 + P2 平行 (週日-週四 25 hr wall time)；P3 待 P1 完成；",
            "P4 待教授決策 (slide 11 三選項)；P5 0.5 hr 收尾。",
            "預期本週交付：Pass 2 重驗報告 + cross-sample 938f0df 評估 + caveat 補丁。",
        ],
        title="本週工時規劃",
    )
    add_footer_citation(s, "Source: master_draft.md §16")
    set_speaker_note(
        s,
        "把 Future 5 件事視覺化成優先序卡片。P1 Pass 2 重驗是決定性 25 hr — 不做的話「V5 為當前最佳」"
        "結論失去現實基礎，所有引用 V5 的 audit chain 都卡住。P2 跨樣本 4 hr 與 P1 共用 binary 並行。"
        "P3 PI 報告 caveat 2 hr 看 P1 完成度決定立即補還是等重驗完一起。P4 V5 留存決策文件 1 hr "
        "等教授判斷三選項後再寫。P5 Memory 校正 0.5 hr 收尾。\n\n"
        "本週執行計畫：P1+P2 並行 ~25 hr wall time，週日到週四完成。P3 視 P1 結果決定。\n\n"
        "預期交付：Pass 2 重驗報告、cross-sample 938f0df 評估、caveat 補丁。"
        "下週週報就能基於 Pass 2 真實數據更新 V5 結論。\n\n"
        "[ORAL-OPTIONAL] 25 hr 不是純計算時間，是 wall time。實際運算 7 樣本平行 ~21 hr，"
        "分析整理 ~4 hr。如果 cluster busy 可能延到 30 hr，但仍週內。",
        min_chars=300,
    )
    add_focal_marker(s, "5 priorities + 工時")
    return s


def slide_16_take_home(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Take-home: 3 件事 — Commit 鏈 / Pass 1 only / Pass 2 重驗 P0", color_bg=PALETTE["green"])
    items = [
        (
            "1. Commit 鏈完整化",
            "V5 是 5 commits（含 4/30 ploidy fix + threshold cherry-pick）",
            "[OK] 5/05 audit 完成 — working tree caveat R1 解決",
            PALETTE["green"],
        ),
        (
            "2. Pass 1 only Caveat",
            "4/29 PI 報告所有 V5 數值基於 Pass 1 only 條件 (ploidy bug 副作用)",
            "[!] 主結論暫停為「Pass 1 only 條件下觀察」",
            PALETTE["accent"],
        ),
        (
            "3. Pass 2 重驗 P0",
            "4/30 修補後 binary 下重跑 sanity / concordance / ISM benchmark",
            "🎯 ~25 hr 含 7 樣本平行；本週執行；下週可恢復「驗證後結論」",
            PALETTE["title_bg"],
        ),
    ]
    base_x = Inches(0.5)
    base_y = Inches(1.5)
    card_w = Inches(12.333)
    # 加高至 1.7 in（+0.2 in × 3 = +0.6 in），視覺比例 ~46% → ~54%
    card_h = Inches(1.7)
    gap = Inches(0.15)
    for i, (label, body, status, color) in enumerate(items):
        y = base_y + (card_h + gap) * i
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x, y, card_w, card_h)
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        shp.line.color.rgb = color
        shp.line.width = Pt(2.5)
        tf = shp.text_frame
        tf.margin_left = Emu(160000)
        tf.margin_right = Emu(160000)
        tf.margin_top = Emu(80000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        add_text_with_fallback(
            tf, label, font_size=20, bold=True, color=color, paragraph_index=0,
        )
        add_text_with_fallback(
            tf, body, font_size=14, color=PALETTE["text"],
        )
        add_text_with_fallback(
            tf, status, font_size=12, bold=True, color=color,
        )
    add_footer_citation(s, "Source: master_draft.md §0 Highlights, §7 Priority order")
    set_speaker_note(
        s,
        "Take-home 3 件事，對應 main thesis 的三個面向。"
        "第一件事：commit 鏈完整化。V5 不是 4 commits 是 5 commits，4/30 補了兩個關鍵 commit "
        "(d0bcd8c ploidy fix + 938f0df threshold cherry-pick)。working tree caveat R1 同時解決。"
        "這部分是 [OK] 已完成。\n\n"
        "第二件事：Pass 1 only caveat。4/29 PI 報告所有 V5 數值都基於 Pass 1 only 條件，"
        "因為 4/12 BAM 的 ploidy bug 讓 Pass 2 從未觸發。這部分是 [!] 主結論暫停為「Pass 1 only 條件下觀察」。\n\n"
        "第三件事：Pass 2 重驗 P0。在 4/30 修補後 binary 下重跑 sanity/concordance/ISM benchmark，"
        "~25 hr 含 7 樣本平行。本週執行，下週週報就能恢復為「驗證後結論」。\n\n"
        "三件事按 main thesis 順序：commit chain (5 commits) → Pass 1 only → Pass 2 重驗 P0。\n\n"
        "[ORAL-OPTIONAL] 強調本週的 critical finding 不是 V5 失敗，是「V5 結論基礎需校正」— "
        "更謹慎的態度而非結論翻轉。",
        min_chars=300,
    )
    add_focal_marker(s, "3 件事 = main thesis 結構")
    return s


def slide_17_qa_backup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Q&A 預備：教授必問 3 + 可能問 4 — evidence 已整備")
    must = [
        ("Q1 PI 結論還能用嗎？", "暫停為 Pass 1 only 觀察；補 caveat banner 而非全面撤回；P0 重驗 (~25 hr)"),
        ("Q2 新 baseline 偏移 systematic？", "目前只有 HCC1395 5kHz [O]；P2 跨樣本評估後才下結論 [U]"),
        ("Q3 V5 flag 還要留嗎？", "三選項：(a) 接受 default / (b) 維護 noPath3 / (c) 上游 PR — 待教授判斷"),
    ]
    may = [
        ("Q4 caller F1 為何重要？", "FILTER 不變 → caller 不傷；phasing 品質需 ISM 下游 F1 評估"),
        ("Q5 4 → 5 commit 為何？", "4/30 補 d0bcd8c (ploidy fix) + 938f0df (threshold cherry-pick); R1 同步解決"),
        ("Q6 HPFineNGroups 還有效？", "4/28 audit 是 4/12 BAM；新 binary R3 待驗 (與 R1 平行 0.5 hr)"),
        ("Q7 下週時間夠？", "P1+P2 ~29 hr 平行可壓 25 hr wall time；週內可完成"),
    ]
    # Two columns — 字級提升至 14pt 下限（PI 25 min 內需快速辨讀）
    add_neutral_box(
        s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5),
        [item for q, a in must for item in [f"[P1] {q}", f"  → {a}", ""]],
        title="必問 (Must-Answer, 3 個)",
        accent=PALETTE["accent"],
        body_font_size=14,
    )
    add_neutral_box(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(5.5),
        [item for q, a in may for item in [f"[P2] {q}", f"  → {a}", ""]],
        title="可能問 (May-Ask, 4 個)",
        accent=PALETTE["amber"],
        body_font_size=13,
    )
    add_footer_citation(s, "Source: master_draft.md §17")
    set_speaker_note(
        s,
        "Q&A 預備頁是 backup — 不一定講，看時間與教授提問方向。"
        "必問 3 個是直接挑戰主結論的：PI 結論還能用嗎、新 baseline 偏移 systematic 嗎、V5 flag 還要留嗎。"
        "三個都已在前面 slide 處理過，這裡只是 quick-reference。\n\n"
        "可能問 4 個是技術細節：caller F1 為何重要、4→5 commits 怎麼回事、"
        "HPFineNGroups 還有效嗎、下週時間夠嗎。每個都有 1 句 evidence-backed 回答。\n\n"
        "[ORAL-OPTIONAL] 如果教授問其他不在這 7 個的問題，預設回應：「我把 evidence 整理到 audit_01.md "
        "之後 follow up」— 不要硬答超出已驗證範圍的東西。",
        min_chars=300,
    )
    add_focal_marker(s, "7 Q&A — backup")
    return s


def slide_18_acknowledgments(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_assertion_title(s, "Acknowledgments + References — 9 份 deliverable + Memory + 上游 zhenyu")
    add_neutral_box(
        s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0),
        [
            "本週 9 份 deliverable：",
            "  • 4/26 Thread D 主軸切換",
            "  • 4/28 V5 audit haplotag sanity",
            "  • 4/29 longphase TO vs V5 技術報告",
            "  • 4/29 supplement getVote design intent",
            "  • 4/30 V3F purity ablation × 2",
            "  • 5/01 三條路 audit",
            "  • 5/01 v5_force_path2only ablation",
            "  • 5/05 V5 provenance audit * critical",
            "",
            "Memory 校正：",
            "  • project_v5_somatic_fallback_verification → needs_rerun",
        ],
        title="Deliverables",
    )
    add_method_box(
        s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(5.0),
        [
            "上游 zhenyu：",
            "  • 938f0df threshold cherry-pick",
            "  • d0bcd8c ploidy bug fix collaboration",
            "",
            "Reference 清單 (10 項)：",
            "  • 20260505_self_phasing_V5_data_provenance_audit_01.md",
            "  • 20260428_Self_Phasing_Baseline_V5_Audit_01.md",
            "  • 20260429_longphase_TO_vs_V5_技術報告_01.md",
            "  • 20260429_supplement_getVote_design_intent_QA_01.md",
            "  • 20260430_V3F_ablation_purity06_results_01.md",
            "  • 20260430_V3F_getVote_only_ablation_results_01.md",
            "  • 20260501_latest_longphase_to_3paths_audit_01.md",
            "  • 20260501_v5_force_path2only_ablation_01.md",
            "  • 20260426_Thread_D_LOH_constrained_phasing_01.md",
            "  • research_landscape/02_Self_Phasing根因.md",
        ],
        title="Acknowledgments + References",
    )
    add_footer_citation(s, "完整 master_draft.md → InterSubMod/docs/reports/validated/2026/05/20260505_週報.../")
    set_speaker_note(
        s,
        "結尾頁面 — Acknowledgments + References。本週 9 份 deliverable 完成，包括 4/26 Thread D 主軸切換、"
        "4/28 V5 audit、4/29 兩份 (技術報告 + supplement)、4/30 兩份 V3F ablation、5/01 兩份 (三條路 + force_path2only)、"
        "5/05 critical finding provenance audit。Memory 同步校正一條。\n\n"
        "上游 zhenyu collaboration 兩個 commit：938f0df threshold cherry-pick 與 d0bcd8c ploidy fix。"
        "兩個都是 4/30 同一天落地，讓本週 audit 正好卡在 caveat 邊界。\n\n"
        "Reference 10 項涵蓋本週與上週相關文件，完整 master_draft.md 在 InterSubMod/docs/reports/validated/...\n\n"
        "[ORAL-OPTIONAL] 上週 4/23 週報還有 3 份 (Thread B 撤回、HPFineNGroups 降級、LOH-AF 保留)，"
        "與本週 9 份合計 ~12 份。Memory 與證據鏈索引可以追問。",
        min_chars=300,
    )
    add_focal_marker(s, "結尾 — 9 份 deliverable")
    return s


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    builders = [
        slide_01_cover,
        slide_02_background,
        slide_03_v5_commit_chain,
        slide_04_pi_pass1_only,
        slide_05_mechanism,
        slide_06_three_paths,
        slide_07_5version_table,
        slide_08_force_path2only,
        slide_09_caller_f1_unchanged,
        slide_10_caveat_banner,
        slide_11_decision_tree,
        slide_12_cross_sample_risk,
        slide_13_hpfinengroups_revalidation,
        slide_14_thread_b_context,
        slide_15_future_priorities,
        slide_16_take_home,
        slide_17_qa_backup,
        slide_18_acknowledgments,
    ]
    for b in builders:
        b(prs)

    out = HERE / "output.pptx"
    prs.save(str(out))
    print(f"[OK] PPTX saved: {out}")
    print(f"     Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
